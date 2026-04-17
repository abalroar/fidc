from __future__ import annotations

import io
import importlib.util
from pathlib import Path
import tempfile
import unittest
from datetime import datetime, timezone

import pandas as pd

from services.fundonet_dashboard import _build_return_history, build_dashboard_data
from services.fundonet_pdf_export import build_dashboard_pdf_bytes


class FundonetDashboardTests(unittest.TestCase):
    def test_build_dashboard_data_computes_summary_returns_and_event_filters(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        self.assertEqual("01/2026", dashboard.latest_competencia)
        self.assertAlmostEqual(20_500.0, dashboard.summary["pl_total"] or 0.0)
        self.assertAlmostEqual(48.7804878, dashboard.summary["subordinacao_pct"] or 0.0, places=5)
        self.assertIsNone(dashboard.summary["direitos_creditorios"])
        self.assertIsNone(dashboard.summary["alocacao_pct"])
        self.assertEqual("01/2026", dashboard.composition_latest_df["competencia"].iloc[0])
        self.assertAlmostEqual(2_000.0, dashboard.summary["emissao_mes"] or 0.0)
        self.assertAlmostEqual(500.0, dashboard.summary["amortizacao_mes"] or 0.0)
        self.assertAlmostEqual(250.0, dashboard.summary["resgate_solicitado_mes"] or 0.0)
        self.assertEqual(2, len(dashboard.event_history_df))

        amort_row = dashboard.event_history_df[dashboard.event_history_df["event_type"] == "amortizacao"].iloc[0]
        self.assertAlmostEqual(-500.0, amort_row["valor_total_assinado"])
        self.assertAlmostEqual(-2.43902439, amort_row["valor_total_pct_pl"], places=6)

        event_summary = dashboard.event_summary_latest_df.set_index("event_type")
        self.assertAlmostEqual(-500.0, event_summary.loc["amortizacao", "valor_total_assinado"])
        self.assertAlmostEqual(250.0, event_summary.loc["resgate_solicitado", "valor_total"])
        self.assertEqual("reported_value", event_summary.loc["resgate_solicitado", "source_status"])

        senior_row = dashboard.return_summary_df[dashboard.return_summary_df["label"] == "Sênior · Série 1"].iloc[0]
        self.assertAlmostEqual(2.0, senior_row["retorno_mes_pct"], places=6)
        self.assertAlmostEqual(2.0, senior_row["retorno_ano_pct"], places=6)
        self.assertAlmostEqual(3.02, senior_row["retorno_12m_pct"], places=2)

        benchmark_row = dashboard.performance_vs_benchmark_latest_df[
            dashboard.performance_vs_benchmark_latest_df["label"] == "Sênior · Série 1"
        ].iloc[0]
        self.assertAlmostEqual(1.8, benchmark_row["desempenho_esperado_pct"], places=6)
        self.assertAlmostEqual(2.0, benchmark_row["desempenho_real_pct"], places=6)
        self.assertAlmostEqual(20.0, benchmark_row["gap_bps"], places=6)

    def test_build_return_history_preserves_multiple_classes_with_same_tipo(self) -> None:
        wide_lookup = pd.DataFrame(columns=["tag_path"]).set_index("tag_path", drop=False)
        listas_df = pd.DataFrame(
            [
                {
                    "competencia": "01/2026",
                    "list_group_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SUBORD",
                    "list_index": 1,
                    "tag": "TIPO",
                    "valor_excel": "Mezz",
                },
                {
                    "competencia": "01/2026",
                    "list_group_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SUBORD",
                    "list_index": 1,
                    "tag": "SERIE",
                    "valor_excel": "Série A",
                },
                {
                    "competencia": "01/2026",
                    "list_group_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SUBORD",
                    "list_index": 1,
                    "tag": "PR_APURADA",
                    "valor_excel": 1.0,
                },
                {
                    "competencia": "01/2026",
                    "list_group_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SUBORD",
                    "list_index": 2,
                    "tag": "TIPO",
                    "valor_excel": "Mezz",
                },
                {
                    "competencia": "01/2026",
                    "list_group_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SUBORD",
                    "list_index": 2,
                    "tag": "SERIE",
                    "valor_excel": "Série B",
                },
                {
                    "competencia": "01/2026",
                    "list_group_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SUBORD",
                    "list_index": 2,
                    "tag": "PR_APURADA",
                    "valor_excel": 2.0,
                },
            ]
        )

        history = _build_return_history(
            wide_lookup=wide_lookup,
            listas_df=listas_df,
            competencias=["01/2026"],
        )

        self.assertEqual(2, len(history))
        self.assertEqual(
            {"Mezz · Série A", "Mezz · Série B"},
            set(history["class_label"].tolist()),
        )
        self.assertEqual(2, history["class_key"].nunique())
        self.assertEqual("RENT_MES.PR_APURADA", history["return_source"].iloc[0])

    def test_build_dashboard_data_exposes_fund_header_information(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        self.assertEqual("12345678000190", dashboard.fund_info["cnpj_fundo"])
        self.assertEqual("Teste FIDC RL", dashboard.fund_info["nome_fundo"])
        self.assertEqual("FECHADO", dashboard.fund_info["condominio"])
        self.assertEqual("Oliveira Trust", dashboard.fund_info["nome_administrador"])
        self.assertEqual("Banco Daycoval", dashboard.fund_info["nome_custodiante"])
        self.assertEqual("JGP", dashboard.fund_info["nome_gestor"])
        self.assertEqual("12/2025 a 01/2026", dashboard.fund_info["periodo_analisado"])

    def test_build_dashboard_data_uses_dicred_total_and_exposes_cvm_tables(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_dicred_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        self.assertAlmostEqual(16_000.0, dashboard.summary["direitos_creditorios"] or 0.0)
        self.assertAlmostEqual(80.0, dashboard.summary["alocacao_pct"] or 0.0)

        segment_row = dashboard.segment_latest_df.iloc[0]
        self.assertEqual("Serviços", segment_row["segmento"])
        self.assertAlmostEqual(16_000.0, segment_row["valor"])

        holder_total = dashboard.holder_latest_df[
            dashboard.holder_latest_df["categoria"] == "Total de cotistas"
        ].iloc[0]
        self.assertAlmostEqual(3.0, holder_total["quantidade"])

        rate_row = dashboard.rate_negotiation_latest_df.iloc[0]
        self.assertEqual("Com aquisição", rate_row["grupo"])
        self.assertEqual("Desconto compra", rate_row["operacao"])
        self.assertAlmostEqual(2.0, rate_row["taxa_media"])

        tracking_row = dashboard.tracking_latest_df[
            dashboard.tracking_latest_df["indicador"] == "Alocação em direitos creditórios"
        ].iloc[0]
        self.assertAlmostEqual(80.0, tracking_row["valor"])

        bucket_labels = dashboard.default_buckets_latest_df.sort_values("ordem")["faixa"].tolist()
        self.assertEqual(
            [
                "Até 30 dias",
                "31 a 60 dias",
                "61 a 90 dias",
                "91 a 120 dias",
                "121 a 150 dias",
                "151 a 180 dias",
                "181 a 360 dias",
                "361 a 720 dias",
                "721 a 1080 dias",
                "Acima de 1080 dias",
            ],
            bucket_labels,
        )
        default_361 = dashboard.default_buckets_latest_df[
            dashboard.default_buckets_latest_df["faixa"] == "361 a 720 dias"
        ].iloc[0]
        self.assertAlmostEqual(7.0, default_361["valor"])
        self.assertEqual("reported_value", default_361["source_status"])
        self.assertAlmostEqual(
            6.25,
            dashboard.default_history_df.iloc[0]["somatorio_inadimplentes_aux_validacao_pct_dcs"],
        )
        self.assertAlmostEqual(
            1_000.0,
            dashboard.default_history_df.iloc[0]["somatorio_inadimplentes_aux_validacao"],
        )

        maturity_361 = dashboard.maturity_latest_df[dashboard.maturity_latest_df["faixa"] == "361 a 720 dias"].iloc[0]
        self.assertAlmostEqual(13.0, maturity_361["valor"])
        self.assertEqual("reported_value", maturity_361["source_status"])

    def test_pdf_export_generates_report_bytes(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        pdf_bytes = build_dashboard_pdf_bytes(dashboard)

        self.assertTrue(pdf_bytes.startswith(b"%PDF"))
        self.assertGreater(len(pdf_bytes), 5_000)

    def test_build_dashboard_data_exposes_risk_monitoring_artifacts(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_dicred_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        risk_lookup = dashboard.risk_metrics_df.set_index("metric_id")
        self.assertEqual(
            {"Risco de crédito", "Risco estrutural"},
            set(dashboard.risk_metrics_df["risk_block"]),
        )
        self.assertAlmostEqual(4.375, risk_lookup.loc["inadimplencia_pct", "value"])
        self.assertEqual("calculado", risk_lookup.loc["inadimplencia_pct", "state"])
        self.assertAlmostEqual(100.0, risk_lookup.loc["concentracao_segmento_proxy", "value"])
        self.assertEqual("critico", risk_lookup.loc["subordinacao_pct", "criticality"])
        self.assertIn("Índice de cobertura", dashboard.coverage_gap_df["tema"].tolist())
        self.assertIn("Subordinação reportada (IME)", dashboard.mini_glossary_df["termo"].tolist())
        self.assertIn(
            "summary.inadimplencia_pct",
            dashboard.current_dashboard_inventory_df["nome_variavel"].tolist(),
        )
        self.assertEqual(
            "agregado_direitos_creditorios_item3",
            dashboard.dc_canonical_history_df.iloc[0]["dc_total_fonte_efetiva"],
        )
        self.assertIn(
            "Over regulatório da inadimplência",
            dashboard.executive_memory_df["componente"].tolist(),
        )
        self.assertAlmostEqual(2000.0, dashboard.liquidity_history_df.iloc[0]["liquidez_imediata"])
        self.assertAlmostEqual(100.0, dashboard.default_buckets_latest_df["percentual"].sum())
        self.assertAlmostEqual(16_000.0, dashboard.summary["inadimplencia_denominador"] or 0.0)

    def test_build_dashboard_data_preserves_missing_history_as_nan(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_missing_history_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        first_row = dashboard.asset_history_df.sort_values("competencia_dt").iloc[0]
        self.assertTrue(pd.isna(first_row["ativos_totais"]))
        self.assertTrue(pd.isna(first_row["carteira"]))
        self.assertTrue(pd.isna(first_row["direitos_creditorios"]))
        self.assertTrue(pd.isna(dashboard.liquidity_history_df.sort_values("competencia_dt").iloc[0]["liquidez_imediata"]))

    def test_build_dashboard_data_prefers_maturity_over_zero_aggregates(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_overdue_priority_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        self.assertAlmostEqual(500.0, dashboard.summary["inadimplencia_total"] or 0.0)
        self.assertAlmostEqual(2_000.0, dashboard.summary["inadimplencia_denominador"] or 0.0)
        self.assertAlmostEqual(25.0, dashboard.summary["inadimplencia_pct"] or 0.0)
        self.assertAlmostEqual(12.5, dashboard.summary["provisao_pct_direitos"] or 0.0)
        self.assertAlmostEqual(50.0, dashboard.summary["cobertura_pct"] or 0.0)
        risk_lookup = dashboard.risk_metrics_df.set_index("metric_id")
        self.assertAlmostEqual(50.0, risk_lookup.loc["provisao_pct_inadimplencia", "value"])

    def test_build_dashboard_pptx_bytes(self) -> None:
        if importlib.util.find_spec("pptx") is None:
            self.skipTest("python-pptx não instalado no ambiente local")
        from pptx import Presentation
        from services.fundonet_ppt_export import build_dashboard_pptx_bytes

        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_dicred_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        pptx_bytes = build_dashboard_pptx_bytes(
            dashboard,
            generated_at=datetime(2026, 4, 14, 15, 0, tzinfo=timezone.utc),
        )
        self.assertTrue(pptx_bytes.startswith(b"PK"))
        self.assertGreater(len(pptx_bytes), 10_000)
        presentation = Presentation(io.BytesIO(pptx_bytes))
        self.assertEqual(4, len(presentation.slides))
        slide_1_text = " ".join(
            shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
        )
        self.assertIn("Gerado em: 14/04/2026 12:00 GMT-3", slide_1_text)
        structure_charts = [shape.chart for shape in presentation.slides[0].shapes if getattr(shape, "has_chart", False)]
        self.assertEqual(2, len(structure_charts))
        self.assertTrue(any("<c:legend>" in chart._chartSpace.xml for chart in structure_charts))

        term_slide_text = " ".join(
            shape.text for shape in presentation.slides[1].shapes if hasattr(shape, "text")
        )
        self.assertIn("Rentabilidade e prazo", term_slide_text)
        term_charts = [shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False)]
        self.assertGreaterEqual(len(term_charts), 2)
        self.assertTrue(any("Total" in chart._chartSpace.xml for chart in term_charts))
        self.assertTrue(all("Vencidos" not in chart._chartSpace.xml for chart in term_charts if len(chart.series) == 3))

        credit_charts = [shape.chart for shape in presentation.slides[2].shapes if getattr(shape, "has_chart", False)]
        self.assertGreaterEqual(len(credit_charts), 2)
        self.assertTrue(any('<c:axPos val="r"' in chart._chartSpace.xml for chart in credit_charts))
        self.assertGreaterEqual(sum("<c:legend>" in chart._chartSpace.xml for chart in credit_charts), 2)

        deterioration_text = " ".join(
            shape.text for shape in presentation.slides[3].shapes if hasattr(shape, "text")
        )
        self.assertIn("Deterioração acumulada", deterioration_text)
        deterioration_charts = [shape.chart for shape in presentation.slides[3].shapes if getattr(shape, "has_chart", False)]
        self.assertGreaterEqual(len(deterioration_charts), 1)

    def test_build_dashboard_pptx_bytes_sanitizes_nan_and_inf_series(self) -> None:
        if importlib.util.find_spec("pptx") is None:
            self.skipTest("python-pptx não instalado no ambiente local")
        from services.fundonet_ppt_export import build_dashboard_pptx_bytes

        with tempfile.TemporaryDirectory() as tmpdir:
            workspace = Path(tmpdir)
            self._write_dicred_fixture_csvs(workspace)

            dashboard = build_dashboard_data(
                wide_csv_path=workspace / "informes_wide.csv",
                listas_csv_path=workspace / "estruturas_lista.csv",
                docs_csv_path=workspace / "documentos_filtrados.csv",
            )

        if not dashboard.return_history_df.empty:
            dashboard.return_history_df.loc[dashboard.return_history_df.index[0], "retorno_mensal_pct"] = float("nan")
        if not dashboard.default_history_df.empty:
            dashboard.default_history_df.loc[dashboard.default_history_df.index[0], "cobertura_pct"] = float("inf")

        pptx_bytes = build_dashboard_pptx_bytes(
            dashboard,
            generated_at=datetime(2026, 4, 14, 15, 0, tzinfo=timezone.utc),
        )
        self.assertTrue(pptx_bytes.startswith(b"PK"))
        self.assertGreater(len(pptx_bytes), 10_000)

    @staticmethod
    def _write_fixture_csvs(workspace: Path) -> None:
        wide_rows = [
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "NR_CNPJ_FUNDO",
                "tag_path": "DOC_ARQ/CAB_INFORM/NR_CNPJ_FUNDO",
                "descricao": "CNPJ fundo",
                "12/2025": "12345678000190",
                "01/2026": "12345678000190",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "NR_CNPJ_CLASSE",
                "tag_path": "DOC_ARQ/CAB_INFORM/NR_CNPJ_CLASSE",
                "descricao": "CNPJ classe",
                "12/2025": "12345678000190",
                "01/2026": "12345678000190",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "NR_CNPJ_ADM",
                "tag_path": "DOC_ARQ/CAB_INFORM/NR_CNPJ_ADM",
                "descricao": "CNPJ administrador",
                "12/2025": "99887766000155",
                "01/2026": "99887766000155",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "TP_CONDOMINIO",
                "tag_path": "DOC_ARQ/CAB_INFORM/TP_CONDOMINIO",
                "descricao": "Condomínio",
                "12/2025": "FECHADO",
                "01/2026": "FECHADO",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "CLASS_UNICA",
                "tag_path": "DOC_ARQ/CAB_INFORM/CLASS_UNICA",
                "descricao": "Classe única",
                "12/2025": "SIM",
                "01/2026": "SIM",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "tag": "SERIE",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/SERIE",
                "descricao": "Série",
                "12/2025": "Série 1",
                "01/2026": "Série 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "tag": "QT_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/QT_COTAS",
                "descricao": "Qt cotas",
                "12/2025": "100",
                "01/2026": "100",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "tag": "VL_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/VL_COTAS",
                "descricao": "Vl cotas",
                "12/2025": "100",
                "01/2026": "105",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "tag": "TIPO",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/TIPO",
                "descricao": "Tipo",
                "12/2025": "Subordinada 1",
                "01/2026": "Subordinada 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "tag": "QT_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/QT_COTAS",
                "descricao": "Qt cotas",
                "12/2025": "500",
                "01/2026": "500",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "tag": "VL_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/VL_COTAS",
                "descricao": "Vl cotas",
                "12/2025": "18",
                "01/2026": "20",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "RENT_MES/RENT_CLASSE_SENIOR",
                "tag": "SERIE",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SENIOR/SERIE",
                "descricao": "Série",
                "12/2025": "Série 1",
                "01/2026": "Série 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "RENT_MES/RENT_CLASSE_SENIOR",
                "tag": "PR_APURADA",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SENIOR/PR_APURADA",
                "descricao": "Retorno",
                "12/2025": "1.00",
                "01/2026": "2.00",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "RENT_MES/RENT_CLASSE_SUBORD",
                "tag": "TIPO",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SUBORD/TIPO",
                "descricao": "Tipo",
                "12/2025": "Subordinada 1",
                "01/2026": "Subordinada 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "RENT_MES/RENT_CLASSE_SUBORD",
                "tag": "PR_APURADA",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/RENT_MES/RENT_CLASSE_SUBORD/PR_APURADA",
                "descricao": "Retorno",
                "12/2025": "0.50",
                "01/2026": "1.00",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESEMP/CLASSE_SENIOR",
                "tag": "SERIE",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESEMP/CLASSE_SENIOR/SERIE",
                "descricao": "Série",
                "12/2025": "Série 1",
                "01/2026": "Série 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESEMP/CLASSE_SENIOR",
                "tag": "DESEMP_ESP",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESEMP/CLASSE_SENIOR/DESEMP_ESP",
                "descricao": "Benchmark",
                "12/2025": "0.80",
                "01/2026": "1.80",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESEMP/CLASSE_SENIOR",
                "tag": "DESEMP_REAL",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESEMP/CLASSE_SENIOR/DESEMP_REAL",
                "descricao": "Realizado",
                "12/2025": "1.00",
                "01/2026": "2.00",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESEMP/CLASSE_SUBORD",
                "tag": "TIPO",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESEMP/CLASSE_SUBORD/TIPO",
                "descricao": "Tipo",
                "12/2025": "Subordinada 1",
                "01/2026": "Subordinada 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESEMP/CLASSE_SUBORD",
                "tag": "DESEMP_ESP",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESEMP/CLASSE_SUBORD/DESEMP_ESP",
                "descricao": "Benchmark",
                "12/2025": "0.40",
                "01/2026": "0.90",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESEMP/CLASSE_SUBORD",
                "tag": "DESEMP_REAL",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESEMP/CLASSE_SUBORD/DESEMP_REAL",
                "descricao": "Realizado",
                "12/2025": "0.50",
                "01/2026": "1.00",
            },
            {
                "bloco": "APLIC_ATIVO",
                "sub_bloco": "",
                "tag": "VL_SOM_APLIC_ATIVO",
                "tag_path": "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/VL_SOM_APLIC_ATIVO",
                "descricao": "Ativos",
                "12/2025": "20000",
                "01/2026": "21000",
            },
            {
                "bloco": "APLIC_ATIVO",
                "sub_bloco": "",
                "tag": "VL_CARTEIRA",
                "tag_path": "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/VL_CARTEIRA",
                "descricao": "Carteira",
                "12/2025": "15000",
                "01/2026": "18000",
            },
            {
                "bloco": "APLIC_ATIVO",
                "sub_bloco": "CRED_EXISTE",
                "tag": "VL_CRED_EXISTE_VENC_ADIMPL",
                "tag_path": "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/CRED_EXISTE/VL_CRED_EXISTE_VENC_ADIMPL",
                "descricao": "Direitos",
                "12/2025": "12000",
                "01/2026": "0",
            },
            {
                "bloco": "APLIC_ATIVO",
                "sub_bloco": "CRED_EXISTE",
                "tag": "VL_CRED_TOTAL_VENC_INAD",
                "tag_path": "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/CRED_EXISTE/VL_CRED_TOTAL_VENC_INAD",
                "descricao": "Inadimplência",
                "12/2025": "300",
                "01/2026": "0",
            },
            {
                "bloco": "APLIC_ATIVO",
                "sub_bloco": "CRED_EXISTE",
                "tag": "VL_PROVIS_REDUC_RECUP",
                "tag_path": "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/CRED_EXISTE/VL_PROVIS_REDUC_RECUP",
                "descricao": "Provisão",
                "12/2025": "100",
                "01/2026": "0",
            },
            {
                "bloco": "NEGOC_DICRED_MES",
                "sub_bloco": "AQUISICOES",
                "tag": "VL_DICRED_AQUIS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/NEGOC_DICRED_MES/AQUISICOES/VL_DICRED_AQUIS",
                "descricao": "Aquisições",
                "12/2025": "1000",
                "01/2026": "800",
            },
            {
                "bloco": "NEGOC_DICRED_MES",
                "sub_bloco": "DICRED_MES_ALIEN",
                "tag": "VL_DICRED_ALIEN",
                "tag_path": "DOC_ARQ/LISTA_INFORM/NEGOC_DICRED_MES/DICRED_MES_ALIEN/VL_DICRED_ALIEN",
                "descricao": "Alienações",
                "12/2025": "200",
                "01/2026": "100",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "LIQUIDEZ",
                "tag": "VL_ATIV_LIQDEZ",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/LIQUIDEZ/VL_ATIV_LIQDEZ",
                "descricao": "Liquidez",
                "12/2025": "500",
                "01/2026": "600",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "LIQUIDEZ",
                "tag": "VL_ATIV_LIQDEZ_30",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/LIQUIDEZ/VL_ATIV_LIQDEZ_30",
                "descricao": "Liquidez 30",
                "12/2025": "300",
                "01/2026": "350",
            },
            {
                "bloco": "COMPMT_DICRED_AQUIS",
                "sub_bloco": "",
                "tag": "VL_PRAZO_VENC_30",
                "tag_path": "DOC_ARQ/LISTA_INFORM/COMPMT_DICRED_AQUIS/VL_PRAZO_VENC_30",
                "descricao": "Prazo 30",
                "12/2025": "2000",
                "01/2026": "0",
            },
            {
                "bloco": "COMPMT_DICRED_AQUIS",
                "sub_bloco": "",
                "tag": "VL_INAD_VENC_30",
                "tag_path": "DOC_ARQ/LISTA_INFORM/COMPMT_DICRED_AQUIS/VL_INAD_VENC_30",
                "descricao": "Inad 30",
                "12/2025": "100",
                "01/2026": "0",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/CAPT_MES/CLASSE_SUBORD",
                "tag": "TIPO",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/CAPT_MES/CLASSE_SUBORD/TIPO",
                "descricao": "Tipo",
                "12/2025": "Subordinada 1",
                "01/2026": "Subordinada 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/CAPT_MES/CLASSE_SUBORD",
                "tag": "VL_TOTAL",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/CAPT_MES/CLASSE_SUBORD/VL_TOTAL",
                "descricao": "Emissão",
                "12/2025": "0",
                "01/2026": "2000",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/CAPT_MES/CLASSE_SUBORD",
                "tag": "QT_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/CAPT_MES/CLASSE_SUBORD/QT_COTAS",
                "descricao": "Qt emissão",
                "12/2025": "0",
                "01/2026": "100",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/AMORT/CLASSE_SENIOR",
                "tag": "SERIE",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/AMORT/CLASSE_SENIOR/SERIE",
                "descricao": "Série",
                "12/2025": "Série 1",
                "01/2026": "Série 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/AMORT/CLASSE_SENIOR",
                "tag": "VL_TOTAL",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/AMORT/CLASSE_SENIOR/VL_TOTAL",
                "descricao": "Amortização",
                "12/2025": "0",
                "01/2026": "500",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/AMORT/CLASSE_SENIOR",
                "tag": "VL_COTA",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/AMORT/CLASSE_SENIOR/VL_COTA",
                "descricao": "Amortização/cota",
                "12/2025": "0",
                "01/2026": "5",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/RESG_SOLIC/CLASSE_SENIOR",
                "tag": "SERIE",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/RESG_SOLIC/CLASSE_SENIOR/SERIE",
                "descricao": "Série",
                "12/2025": "Série 1",
                "01/2026": "Série 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/RESG_SOLIC/CLASSE_SENIOR",
                "tag": "QT_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/RESG_SOLIC/CLASSE_SENIOR/QT_COTAS",
                "descricao": "Qt resgate solicitado",
                "12/2025": "0",
                "01/2026": "10",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "CAPTA_RESGA_AMORTI/RESG_SOLIC/CLASSE_SENIOR",
                "tag": "VL_PAGO",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/CAPTA_RESGA_AMORTI/RESG_SOLIC/CLASSE_SENIOR/VL_PAGO",
                "descricao": "Resgate solicitado",
                "12/2025": "0",
                "01/2026": "250",
            },
        ]
        pd.DataFrame(wide_rows).to_csv(workspace / "informes_wide.csv", index=False)

        listas_columns = [
            "competencia",
            "list_group_path",
            "list_index",
            "tag",
            "valor_excel",
        ]
        pd.DataFrame(columns=listas_columns).to_csv(workspace / "estruturas_lista.csv", index=False)

        docs_rows = [
            {
                "documento_id": "1",
                "competencia": "12/2025",
                "data_entrega": "20/01/2026 09:00",
                "fundo_ou_classe": "Classe",
                "nome_fundo": "Teste FIDC RL",
                "nome_administrador": "Oliveira Trust",
                "nome_custodiante": "Banco Daycoval",
                "nome_gestor": "JGP",
                "processamento": "ok",
            },
            {
                "documento_id": "2",
                "competencia": "01/2026",
                "data_entrega": "20/02/2026 09:00",
                "fundo_ou_classe": "Classe",
                "nome_fundo": "Teste FIDC RL",
                "nome_administrador": "Oliveira Trust",
                "nome_custodiante": "Banco Daycoval",
                "nome_gestor": "JGP",
                "processamento": "ok",
            },
        ]
        pd.DataFrame(docs_rows).to_csv(workspace / "documentos_filtrados.csv", index=False)

    @staticmethod
    def _write_dicred_fixture_csvs(workspace: Path) -> None:
        competencia = "02/2026"

        def row(bloco: str, sub_bloco: str, tag: str, tag_path: str, value: object) -> dict[str, object]:
            return {
                "bloco": bloco,
                "sub_bloco": sub_bloco,
                "tag": tag,
                "tag_path": tag_path,
                "descricao": tag,
                competencia: value,
            }

        wide_rows = [
            row("CAB_INFORM", "", "NR_CNPJ_FUNDO", "DOC_ARQ/CAB_INFORM/NR_CNPJ_FUNDO", "50473039000102"),
            row("CAB_INFORM", "", "NR_CNPJ_CLASSE", "DOC_ARQ/CAB_INFORM/NR_CNPJ_CLASSE", "50473039000102"),
            row("CAB_INFORM", "", "NR_CNPJ_ADM", "DOC_ARQ/CAB_INFORM/NR_CNPJ_ADM", "36113876000191"),
            row("CAB_INFORM", "", "NM_CLASSE", "DOC_ARQ/CAB_INFORM/NM_CLASSE", "SELLER FIDC"),
            row("CAB_INFORM", "", "TP_CONDOMINIO", "DOC_ARQ/CAB_INFORM/TP_CONDOMINIO", "FECHADO"),
            row("CAB_INFORM", "", "CLASS_UNICA", "DOC_ARQ/CAB_INFORM/CLASS_UNICA", "NAO"),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "SERIE",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/SERIE",
                "Série 1",
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "QT_COTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/QT_COTAS",
                100,
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "VL_COTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/VL_COTAS",
                100,
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "TIPO",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/TIPO",
                "Subordinada",
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "QT_COTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/QT_COTAS",
                100,
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "VL_COTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/VL_COTAS",
                50,
            ),
            row("APLIC_ATIVO", "", "VL_SOM_APLIC_ATIVO", "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/VL_SOM_APLIC_ATIVO", 21000),
            row("APLIC_ATIVO", "", "VL_CARTEIRA", "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/VL_CARTEIRA", 20000),
            row("APLIC_ATIVO", "DICRED", "VL_DICRED", "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/DICRED/VL_DICRED", 16000),
            row(
                "APLIC_ATIVO",
                "DICRED",
                "VL_DICRED_TOTAL_VENC_INAD",
                "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/DICRED/VL_DICRED_TOTAL_VENC_INAD",
                700,
            ),
            row(
                "APLIC_ATIVO",
                "DICRED",
                "VL_DICRED_EXISTE_INAD",
                "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/DICRED/VL_DICRED_EXISTE_INAD",
                200,
            ),
            row(
                "APLIC_ATIVO",
                "DICRED",
                "VL_DICRED_VENC_PEND",
                "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/DICRED/VL_DICRED_VENC_PEND",
                100,
            ),
            row(
                "NEGOC_DICRED_MES",
                "AQUISICOES",
                "VL_DICRED_AQUIS",
                "DOC_ARQ/LISTA_INFORM/NEGOC_DICRED_MES/AQUISICOES/VL_DICRED_AQUIS",
                600,
            ),
            row(
                "NEGOC_DICRED_MES",
                "DICRED_MES_ALIEN",
                "VL_DICRED_ALIEN",
                "DOC_ARQ/LISTA_INFORM/NEGOC_DICRED_MES/DICRED_MES_ALIEN/VL_DICRED_ALIEN",
                200,
            ),
            row(
                "CART_SEGMT",
                "SEGMT_SERV",
                "VL_SOM_SEGMT_SERV",
                "DOC_ARQ/LISTA_INFORM/CART_SEGMT/SEGMT_SERV/VL_SOM_SEGMT_SERV",
                16000,
            ),
            row(
                "OUTRAS_INFORM",
                "LIQUIDEZ",
                "VL_ATIV_LIQDEZ",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/LIQUIDEZ/VL_ATIV_LIQDEZ",
                2000,
            ),
            row(
                "OUTRAS_INFORM",
                "LIQUIDEZ",
                "VL_ATIV_LIQDEZ_30",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/LIQUIDEZ/VL_ATIV_LIQDEZ_30",
                3500,
            ),
            row(
                "OUTRAS_INFORM",
                "NUM_COTISTAS",
                "QT_TOTAL_COTISTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/NUM_COTISTAS/QT_TOTAL_COTISTAS",
                3,
            ),
            row(
                "OUTRAS_INFORM",
                "NUM_COTISTAS",
                "QT_TOTAL_COTISTAS_SENIOR",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/NUM_COTISTAS/QT_TOTAL_COTISTAS_SENIOR",
                2,
            ),
            row(
                "OUTRAS_INFORM",
                "NUM_COTISTAS",
                "QT_TOTAL_COTISTAS_SUBORD",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/NUM_COTISTAS/QT_TOTAL_COTISTAS_SUBORD",
                1,
            ),
            row(
                "TAXA_NEGOC_DICRED_MES",
                "TAXA_NEGOC_DICRED_MES_AQUIS/TAXA_NEGOC_DICRED_MES_AQUIS_DESC_COMPRA",
                "TX_MIN",
                "DOC_ARQ/LISTA_INFORM/TAXA_NEGOC_DICRED_MES/TAXA_NEGOC_DICRED_MES_AQUIS/TAXA_NEGOC_DICRED_MES_AQUIS_DESC_COMPRA/TX_MIN",
                1,
            ),
            row(
                "TAXA_NEGOC_DICRED_MES",
                "TAXA_NEGOC_DICRED_MES_AQUIS/TAXA_NEGOC_DICRED_MES_AQUIS_DESC_COMPRA",
                "TX_MEDIO",
                "DOC_ARQ/LISTA_INFORM/TAXA_NEGOC_DICRED_MES/TAXA_NEGOC_DICRED_MES_AQUIS/TAXA_NEGOC_DICRED_MES_AQUIS_DESC_COMPRA/TX_MEDIO",
                2,
            ),
            row(
                "TAXA_NEGOC_DICRED_MES",
                "TAXA_NEGOC_DICRED_MES_AQUIS/TAXA_NEGOC_DICRED_MES_AQUIS_DESC_COMPRA",
                "TX_MAXIMO",
                "DOC_ARQ/LISTA_INFORM/TAXA_NEGOC_DICRED_MES/TAXA_NEGOC_DICRED_MES_AQUIS/TAXA_NEGOC_DICRED_MES_AQUIS_DESC_COMPRA/TX_MAXIMO",
                3,
            ),
            row(
                "COMPMT_DICRED_AQUIS",
                "",
                "VL_INAD_VENC_361_720",
                "DOC_ARQ/LISTA_INFORM/COMPMT_DICRED_AQUIS/VL_INAD_VENC_361_720",
                7,
            ),
            row(
                "COMPMT_DICRED_SEM_AQUIS",
                "",
                "VL_INAD_VENC_361_720",
                "DOC_ARQ/LISTA_INFORM/COMPMT_DICRED_SEM_AQUIS/VL_INAD_VENC_361_720",
                0,
            ),
            row(
                "COMPMT_DICRED_AQUIS",
                "",
                "VL_PRAZO_VENC_361_720",
                "DOC_ARQ/LISTA_INFORM/COMPMT_DICRED_AQUIS/VL_PRAZO_VENC_361_720",
                13,
            ),
            row(
                "COMPMT_DICRED_SEM_AQUIS",
                "",
                "VL_PRAZO_VENC_361_720",
                "DOC_ARQ/LISTA_INFORM/COMPMT_DICRED_SEM_AQUIS/VL_PRAZO_VENC_361_720",
                0,
            ),
        ]
        pd.DataFrame(wide_rows).to_csv(workspace / "informes_wide.csv", index=False)

        listas_rows = [
            {
                "competencia": competencia,
                "list_group_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/NUM_COTISTAS/CLASSE_SENIOR",
                "list_index": 1,
                "tag": "SERIE",
                "valor_excel": "Série 1",
            },
            {
                "competencia": competencia,
                "list_group_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/NUM_COTISTAS/CLASSE_SENIOR",
                "list_index": 1,
                "tag": "QT_COTISTAS",
                "valor_excel": 2,
            },
        ]
        pd.DataFrame(listas_rows).to_csv(workspace / "estruturas_lista.csv", index=False)

        docs_rows = [
            {
                "documento_id": "1136117",
                "competencia": competencia,
                "data_entrega": "20/03/2026 09:00",
                "fundo_ou_classe": "Classe",
                "nome_fundo": "Seller FIDC",
                "nome_administrador": "Oliveira Trust",
                "nome_custodiante": "Banco Daycoval",
                "nome_gestor": "JGP",
                "processamento": "ok",
            }
        ]
        pd.DataFrame(docs_rows).to_csv(workspace / "documentos_filtrados.csv", index=False)

    @staticmethod
    def _write_missing_history_fixture_csvs(workspace: Path) -> None:
        wide_rows = [
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "NR_CNPJ_FUNDO",
                "tag_path": "DOC_ARQ/CAB_INFORM/NR_CNPJ_FUNDO",
                "descricao": "CNPJ fundo",
                "12/2025": "11111111000111",
                "01/2026": "11111111000111",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "NR_CNPJ_CLASSE",
                "tag_path": "DOC_ARQ/CAB_INFORM/NR_CNPJ_CLASSE",
                "descricao": "CNPJ classe",
                "12/2025": "11111111000111",
                "01/2026": "11111111000111",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "NR_CNPJ_ADM",
                "tag_path": "DOC_ARQ/CAB_INFORM/NR_CNPJ_ADM",
                "descricao": "CNPJ administrador",
                "12/2025": "22222222000122",
                "01/2026": "22222222000122",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "NM_CLASSE",
                "tag_path": "DOC_ARQ/CAB_INFORM/NM_CLASSE",
                "descricao": "Classe",
                "12/2025": "Classe Teste",
                "01/2026": "Classe Teste",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "TP_CONDOMINIO",
                "tag_path": "DOC_ARQ/CAB_INFORM/TP_CONDOMINIO",
                "descricao": "Condomínio",
                "12/2025": "FECHADO",
                "01/2026": "FECHADO",
            },
            {
                "bloco": "CAB_INFORM",
                "sub_bloco": "",
                "tag": "CLASS_UNICA",
                "tag_path": "DOC_ARQ/CAB_INFORM/CLASS_UNICA",
                "descricao": "Classe única",
                "12/2025": "SIM",
                "01/2026": "SIM",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "tag": "SERIE",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/SERIE",
                "descricao": "Série",
                "12/2025": "Série 1",
                "01/2026": "Série 1",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "tag": "QT_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/QT_COTAS",
                "descricao": "Qt cotas",
                "12/2025": "100",
                "01/2026": "100",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "tag": "VL_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/VL_COTAS",
                "descricao": "Vl cotas",
                "12/2025": "100",
                "01/2026": "101",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "tag": "TIPO",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/TIPO",
                "descricao": "Tipo",
                "12/2025": "Subordinada",
                "01/2026": "Subordinada",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "tag": "QT_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/QT_COTAS",
                "descricao": "Qt cotas",
                "12/2025": "50",
                "01/2026": "50",
            },
            {
                "bloco": "OUTRAS_INFORM",
                "sub_bloco": "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "tag": "VL_COTAS",
                "tag_path": "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/VL_COTAS",
                "descricao": "Vl cotas",
                "12/2025": "20",
                "01/2026": "20",
            },
            {
                "bloco": "APLIC_ATIVO",
                "sub_bloco": "",
                "tag": "VL_SOM_APLIC_ATIVO",
                "tag_path": "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/VL_SOM_APLIC_ATIVO",
                "descricao": "Ativos",
                "12/2025": "",
                "01/2026": "1000",
            },
            {
                "bloco": "APLIC_ATIVO",
                "sub_bloco": "",
                "tag": "VL_CARTEIRA",
                "tag_path": "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/VL_CARTEIRA",
                "descricao": "Carteira",
                "12/2025": "",
                "01/2026": "900",
            },
            {
                "bloco": "APLIC_ATIVO",
                "sub_bloco": "DICRED",
                "tag": "VL_DICRED",
                "tag_path": "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/DICRED/VL_DICRED",
                "descricao": "Direitos",
                "12/2025": "",
                "01/2026": "800",
            },
        ]
        pd.DataFrame(wide_rows).to_csv(workspace / "informes_wide.csv", index=False)

        pd.DataFrame(columns=["competencia", "list_group_path", "list_index", "tag", "valor_excel"]).to_csv(
            workspace / "estruturas_lista.csv",
            index=False,
        )

        docs_rows = [
            {
                "documento_id": "10",
                "competencia": "12/2025",
                "data_entrega": "20/01/2026 09:00",
                "fundo_ou_classe": "Classe",
                "nome_fundo": "Teste Missing FIDC",
                "processamento": "ok",
            },
            {
                "documento_id": "11",
                "competencia": "01/2026",
                "data_entrega": "20/02/2026 09:00",
                "fundo_ou_classe": "Classe",
                "nome_fundo": "Teste Missing FIDC",
                "processamento": "ok",
            },
        ]
        pd.DataFrame(docs_rows).to_csv(workspace / "documentos_filtrados.csv", index=False)

    @staticmethod
    def _write_overdue_priority_fixture_csvs(workspace: Path) -> None:
        competencia = "03/2026"

        def row(bloco: str, sub_bloco: str, tag: str, tag_path: str, value: object) -> dict[str, object]:
            return {
                "bloco": bloco,
                "sub_bloco": sub_bloco,
                "tag": tag,
                "tag_path": tag_path,
                "descricao": tag,
                competencia: value,
            }

        wide_rows = [
            row("CAB_INFORM", "", "NR_CNPJ_FUNDO", "DOC_ARQ/CAB_INFORM/NR_CNPJ_FUNDO", "41970012000126"),
            row("CAB_INFORM", "", "NR_CNPJ_CLASSE", "DOC_ARQ/CAB_INFORM/NR_CNPJ_CLASSE", "41970012000126"),
            row("CAB_INFORM", "", "NR_CNPJ_ADM", "DOC_ARQ/CAB_INFORM/NR_CNPJ_ADM", "36113876000191"),
            row("CAB_INFORM", "", "NM_CLASSE", "DOC_ARQ/CAB_INFORM/NM_CLASSE", "FIDC Prioridade Vencidos"),
            row("CAB_INFORM", "", "TP_CONDOMINIO", "DOC_ARQ/CAB_INFORM/TP_CONDOMINIO", "FECHADO"),
            row("CAB_INFORM", "", "CLASS_UNICA", "DOC_ARQ/CAB_INFORM/CLASS_UNICA", "SIM"),
            row("APLIC_ATIVO", "", "VL_SOM_APLIC_ATIVO", "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/VL_SOM_APLIC_ATIVO", 2_200),
            row("APLIC_ATIVO", "", "VL_CARTEIRA", "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/VL_CARTEIRA", 2_100),
            row("APLIC_ATIVO", "DICRED", "VL_DICRED", "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/DICRED/VL_DICRED", 0),
            row("APLIC_ATIVO", "CRED_EXISTE", "VL_CRED_TOTAL_VENC_INAD", "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/CRED_EXISTE/VL_CRED_TOTAL_VENC_INAD", 0),
            row("APLIC_ATIVO", "CRED_EXISTE", "VL_PROVIS_REDUC_RECUP", "DOC_ARQ/LISTA_INFORM/APLIC_ATIVO/CRED_EXISTE/VL_PROVIS_REDUC_RECUP", 250),
            row("COMPMT_DICRED_AQUIS", "", "VL_SOM_INAD_VENC", "DOC_ARQ/LISTA_INFORM/COMPMT_DICRED_AQUIS/VL_SOM_INAD_VENC", 500),
            row("COMPMT_DICRED_AQUIS", "", "VL_PRAZO_VENC_30", "DOC_ARQ/LISTA_INFORM/COMPMT_DICRED_AQUIS/VL_PRAZO_VENC_30", 1_500),
            row(
                "OUTRAS_INFORM",
                "NUM_COTISTAS",
                "QT_TOTAL_COTISTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/NUM_COTISTAS/QT_TOTAL_COTISTAS",
                2,
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "SERIE",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/SERIE",
                "Série 1",
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "QT_COTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/QT_COTAS",
                100,
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR",
                "VL_COTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SENIOR/VL_COTAS",
                10,
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "TIPO",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/TIPO",
                "Subordinada",
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "QT_COTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/QT_COTAS",
                100,
            ),
            row(
                "OUTRAS_INFORM",
                "DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD",
                "VL_COTAS",
                "DOC_ARQ/LISTA_INFORM/OUTRAS_INFORM/DESC_SERIE_CLASSE/DESC_SERIE_CLASSE_SUBORD/VL_COTAS",
                10,
            ),
        ]
        pd.DataFrame(wide_rows).to_csv(workspace / "informes_wide.csv", index=False)
        pd.DataFrame(columns=["competencia", "list_group_path", "list_index", "tag", "valor_excel"]).to_csv(
            workspace / "estruturas_lista.csv",
            index=False,
        )
        pd.DataFrame(
            [
                {
                    "documento_id": "20",
                    "competencia": competencia,
                    "data_entrega": "20/04/2026 09:00",
                    "fundo_ou_classe": "Classe",
                    "nome_fundo": "FIDC Prioridade Vencidos",
                    "nome_administrador": "Oliveira Trust",
                    "processamento": "ok",
                }
            ]
        ).to_csv(workspace / "documentos_filtrados.csv", index=False)


if __name__ == "__main__":
    unittest.main()
