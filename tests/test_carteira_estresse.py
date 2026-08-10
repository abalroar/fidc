"""O teste de estresse: reconhecer o buraco de provisão e ver quem desenquadra."""

from __future__ import annotations

from pathlib import Path
import sys

import pandas as pd
import pytest


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from services.carteira_estresse import (  # noqa: E402
    CAPITAL_CONSUMIDO,
    DESENQUADRADO,
    ENQUADRADO,
    estressar,
    nao_reportantes,
    sob_estresse,
)


def fundo(**campos):
    base = {
        "cnpj": "11111111000191",
        "fundo": "FIDC Teste",
        "rotulo": "FIDC Teste",
        "categoria_estrutural": "Financeiro",
        "vl_cotas_subordinadas": 30e6,
        "vl_cotas_total": 100e6,
        "referencia_pct": 20.0,
        "carteira_dc": 90e6,
        "dc_inadimplentes": 20e6,
        "pdd_brl": 10e6,
        "cobertura_pct": 50.0,
    }
    base.update(campos)
    return base


def test_o_deficit_e_inadimplencia_menos_pdd() -> None:
    """E não uma fração da carteira: a carteira reportada já é líquida de PDD."""

    linha = estressar(pd.DataFrame([fundo()])).iloc[0]

    assert linha["deficit_brl"] == pytest.approx(10e6)


def test_o_deficit_desce_do_subordinado_e_do_total() -> None:
    """A perda consome a classe subordinada; o sênior está protegido."""

    linha = estressar(pd.DataFrame([fundo()])).iloc[0]

    assert linha["sub_antes_pct"] == pytest.approx(30.0)
    # (30 - 10) / (100 - 10) = 22,22%
    assert linha["sub_pos_pct"] == pytest.approx(20e6 / 90e6 * 100)
    assert linha["folga_pos_pp"] == pytest.approx(2.2222, abs=1e-3)
    assert linha["estresse_status"] == ENQUADRADO
    assert linha["aporte_brl"] == pytest.approx(0.0)


def test_o_investcred_confere_com_a_conta_a_mao() -> None:
    """O exemplo que motivou o teste, com os números do Informe de jun-26."""

    linha = estressar(
        pd.DataFrame(
            [
                fundo(
                    vl_cotas_subordinadas=950_612_296.48,
                    vl_cotas_total=1_193_589_922.61,
                    referencia_pct=67.4267,
                    carteira_dc=992_925_206.78,
                    dc_inadimplentes=1_668_944_312.63,
                    pdd_brl=1_443_843_834.23,
                    cobertura_pct=86.5124,
                )
            ]
        )
    ).iloc[0]

    assert linha["deficit_brl"] == pytest.approx(225_100_478.40, abs=1.0)
    assert linha["sub_antes_pct"] == pytest.approx(79.64, abs=0.01)
    assert linha["sub_pos_pct"] == pytest.approx(74.91, abs=0.01)
    assert linha["folga_pos_pp"] == pytest.approx(7.48, abs=0.02)
    assert linha["estresse_status"] == ENQUADRADO


def test_o_aporte_reenquadra_exatamente() -> None:
    """A propriedade que define o número: aportado, o fundo bate o mínimo."""

    linha = estressar(
        pd.DataFrame([fundo(vl_cotas_subordinadas=12e6, referencia_pct=20.0)])
    ).iloc[0]

    assert linha["estresse_status"] == DESENQUADRADO
    aporte = linha["aporte_brl"]
    assert aporte > 0

    sub = linha["sub_pos_brl"] + aporte
    total = linha["total_pos_brl"] + aporte
    assert sub / total * 100 == pytest.approx(20.0)


def test_um_deficit_maior_que_o_subordinado_deixa_o_senior_exposto() -> None:
    """Subordinação negativa é informação: a classe júnior foi consumida."""

    linha = estressar(
        pd.DataFrame([fundo(vl_cotas_subordinadas=5e6, dc_inadimplentes=30e6)])
    ).iloc[0]

    assert linha["deficit_brl"] == pytest.approx(20e6)
    assert linha["sub_pos_brl"] < 0
    assert linha["sub_pos_pct"] < 0
    assert linha["estresse_status"] == DESENQUADRADO


def test_deficit_maior_que_o_patrimonio_e_capital_consumido() -> None:
    linha = estressar(
        pd.DataFrame([fundo(dc_inadimplentes=200e6, pdd_brl=0.0)])
    ).iloc[0]

    assert linha["total_pos_brl"] < 0
    assert linha["estresse_status"] == CAPITAL_CONSUMIDO


def test_cobertura_plena_nao_gera_estresse() -> None:
    """Provisionado até o fim, Δ é zero e a subordinação não se move."""

    linha = estressar(
        pd.DataFrame([fundo(pdd_brl=25e6, cobertura_pct=125.0)])
    ).iloc[0]

    assert linha["deficit_brl"] == pytest.approx(0.0)
    assert linha["sub_pos_pct"] == pytest.approx(linha["sub_antes_pct"])


def test_o_teste_alcanca_so_quem_esta_abaixo_de_cem() -> None:
    """Cobertura plena ocuparia linha sem mudar nada."""

    frame = pd.DataFrame(
        [
            fundo(cnpj="11111111000191", cobertura_pct=50.0),
            fundo(cnpj="22222222000172", cobertura_pct=125.0, pdd_brl=25e6),
            fundo(cnpj="33333333000153", cobertura_pct=40.0, referencia_pct=None),
        ]
    )

    assert set(sob_estresse(frame)["cnpj"]) == {"11111111000191"}


def test_a_apuracao_separa_o_zero_legitimo_do_silencio() -> None:
    frame = pd.DataFrame(
        [
            fundo(cnpj="11111111000191", dc_inadimplentes=0.0, pdd_brl=0.0),
            fundo(cnpj="22222222000172", dc_inadimplentes=0.0, pdd_brl=5e6),
            fundo(cnpj="33333333000153", dc_inadimplentes=8e6, pdd_brl=0.0),
            fundo(cnpj="44444444000134", carteira_dc=0.0, dc_inadimplentes=0.0, pdd_brl=0.0),
            # Este reporta os dois e não entra na apuração.
            fundo(cnpj="55555555000115", dc_inadimplentes=8e6, pdd_brl=4e6),
        ]
    )

    apuracao = nao_reportantes(frame).set_index("cnpj")

    assert "55555555000115" not in apuracao.index
    assert apuracao.at["11111111000191", "caso"].startswith("nem PDD")
    assert apuracao.at["22222222000172", "caso"].startswith("provisiona")
    assert apuracao.at["33333333000153", "caso"].startswith("declara")
    assert apuracao.at["44444444000134", "caso"].startswith("sem carteira")


def test_a_lamina_de_estresse_entra_no_deck() -> None:
    """Duas lâminas: método, resultado e apuração fora da área projetada."""

    from pptx import Presentation
    from pptx.util import Inches

    from services.bba_deck import SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN
    from services.carteira_estresse_deck import append_stress_slide

    apresentacao = Presentation()
    apresentacao.slide_width = Inches(SLIDE_WIDTH_IN)
    apresentacao.slide_height = Inches(SLIDE_HEIGHT_IN)
    antes = len(apresentacao.slides._sldIdLst)

    append_stress_slide(apresentacao, ROOT / "data" / "industry_study")

    assert len(apresentacao.slides._sldIdLst) == antes + 2
    metodologia = apresentacao.slides[-2]
    slide = apresentacao.slides[-1]
    assert any(
        forma.has_text_frame and "Stress Test | Metodologia" in forma.text_frame.text
        for forma in metodologia.shapes
    )
    assert not any(forma.has_table for forma in metodologia.shapes)
    assert not any(forma.shape_type == 13 for forma in slide.shapes)
    tabelas = [shape for shape in slide.shapes if shape.has_table]
    assert len(tabelas) == 2

    dentro = [t for t in tabelas if t.left < Inches(SLIDE_WIDTH_IN)]
    fora = [t for t in tabelas if t.left >= Inches(SLIDE_WIDTH_IN)]
    assert len(dentro) == 1 and len(fora) == 1
    # Cabeçalho e os nove desenquadramentos da triagem.
    assert len(dentro[0].table.rows) == 10
    assert [cell.text for cell in dentro[0].table.rows[0].cells] == [
        "FIDC",
        "Cob.",
        "PL R$ mm",
        "Sub Mín.",
        "Sub/PL",
        "Sub+Mez/PL",
        "Folga",
        "Aporte R$ mm",
        "Aporte/PL",
    ]

    # Nada da lâmina invade o rodapé nem transborda a margem direita.
    for tabela in dentro:
        assert tabela.top + tabela.height <= Inches(6.96)
        assert tabela.left + tabela.width <= Inches(12.68)

    # E a tabela é nativa, não uma imagem.
    assert fora[0].table.cell(0, 0).text == "FIDC"


def test_a_lamina_entra_logo_apos_o_bloco_de_subordinacao() -> None:
    """A lâmina fecha o bloco estrutural, e a numeração adiante acompanha."""

    from pptx import Presentation
    from pptx.util import Inches

    from services.bba_deck import Deck, SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN
    from services.carteira_deck import last_structural_slide
    from services.carteira_estresse_deck import _caixa_de_pagina, append_stress_slide

    apresentacao = Presentation()
    apresentacao.slide_width = Inches(SLIDE_WIDTH_IN)
    apresentacao.slide_height = Inches(SLIDE_HEIGHT_IN)
    deck = Deck("BLOCO")
    deck.prs = apresentacao
    deck.page = 0
    data_dir = ROOT / "data" / "industry_study"
    ultimo_slide_estrutural = last_structural_slide(data_dir)
    for numero in range(1, ultimo_slide_estrutural + 4):
        deck.footer(deck.slide(f"Lâmina {numero}"), "fonte")

    total = len(apresentacao.slides._sldIdLst)
    append_stress_slide(apresentacao, data_dir)

    assert len(apresentacao.slides._sldIdLst) == total + 2
    metodologia = apresentacao.slides[ultimo_slide_estrutural]
    posicionada = apresentacao.slides[ultimo_slide_estrutural + 1]
    assert any(
        forma.has_text_frame and "Stress Test | Metodologia" in forma.text_frame.text
        for forma in metodologia.shapes
    )
    assert any(
        forma.has_text_frame
        and "Stress Test | Nove FIDCs c/ cobertura < 100%" in forma.text_frame.text
        for forma in posicionada.shapes
    )
    # A numeração continua sequencial de ponta a ponta.
    numeros = [_caixa_de_pagina(s)[1] for s in apresentacao.slides]
    assert numeros == list(range(1, len(numeros) + 1))
