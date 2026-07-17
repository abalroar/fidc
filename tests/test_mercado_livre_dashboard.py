from __future__ import annotations

from datetime import date
from io import BytesIO
import json
from pathlib import Path
import tempfile
from types import SimpleNamespace
import unittest
from unittest.mock import patch
import zipfile

from openpyxl import load_workbook
import pandas as pd

from services.fidc_model.b3_cdi import B3CdiMonthlyRate
from services.fund_return_matrix import RETURN_YTD_COLUMN
from services.mercado_livre_dashboard import (
    build_consolidated_monthly_base,
    build_consolidated_snapshot_excel_bytes,
    build_excel_export_bytes,
    build_full_variable_csv_zip_bytes,
    build_full_variable_excel_export_bytes,
    build_full_variable_export_matrix,
    build_fund_monthly_base,
    build_mercado_livre_outputs,
    build_wide_table,
    load_outputs_from_cache,
    MercadoLivreOutputs,
    order_period_columns_desc,
    portfolio_identity_key,
    save_outputs_to_cache,
)
from services.mercado_livre_ppt_export import build_pptx_export_bytes
from services.meli_credit_monitor import build_meli_monitor_outputs
from services.portfolio_store import PortfolioFund, PortfolioRecord
from services.mercado_livre_visuals import npl_coverage_chart, pl_subordination_chart
from tabs.tab_dashboard_meli import (
    _build_fund_return_table,
    _render_stacked_funds_view,
    _resolve_fund_return_cdi_rates,
    render_dashboard_meli_analysis,
    resolve_fund_return_export_inputs,
)
from tabs.tab_mercado_livre import (
    CONSOLIDATED_SCOPE_VALUE,
    _MERCADO_LIVRE_UI_CSS,
    _build_base_scope_options,
    _build_mercado_livre_guide_markdown,
    _resolve_default_base_scope,
    _dense_wide_value,
    _display_window_bounds,
    _display_window_months,
    _build_credit_monitor_for_display,
    _period_with_yoy_lookback,
    _tag_outputs_requested_period,
    _filter_outputs_by_competencia_months,
    _display_wide_table,
    _render_wide_table_html,
    _resolve_existing_portfolio_for_save,
)
from services.ime_period import build_custom_period


class MercadoLivreDashboardTests(unittest.TestCase):
    def test_build_fund_monthly_base_calculates_accumulated_npl_and_ex360(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={
                1: 100.0,
                2: 50.0,
                3: 30.0,
                4: 20.0,
                5: 10.0,
                6: 5.0,
                7: 15.0,
                8: 7.0,
                9: 3.0,
                10: 0.0,
            },
        )

        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)
        row = monthly.iloc[0]

        self.assertAlmostEqual(25.0, row["subordinacao_total_pct"])
        self.assertAlmostEqual(60.0, row["npl_over90"])
        self.assertAlmostEqual(10.0, row["npl_over360"])
        self.assertAlmostEqual(50.0, row["npl_over90_ex360"])
        self.assertAlmostEqual(990.0, row["carteira_ex360"])
        self.assertAlmostEqual(50.0 / 990.0 * 100.0, row["npl_over90_ex360_pct"], places=6)
        self.assertTrue(bool(row["pdd_ex360_calculavel"]))
        self.assertAlmostEqual(10.0, row["baixa_over360_carteira"])
        self.assertAlmostEqual(10.0, row["baixa_over360_pdd"])
        self.assertAlmostEqual(0.0, row["baixa_over360_pl"])
        self.assertAlmostEqual(90.0, row["pdd_ex360"])
        self.assertAlmostEqual(900.0, row["carteira_liquida_ex360"])
        self.assertNotIn("PDD ex-360 não calculável", row["warnings"])

    def test_ex360_writeoff_reduces_pl_when_pdd_is_insufficient(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=5.0,
            buckets={8: 10.0},
        )

        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)
        row = monthly.iloc[0]

        self.assertAlmostEqual(10.0, row["baixa_over360_carteira"])
        self.assertAlmostEqual(5.0, row["baixa_over360_pdd"])
        self.assertAlmostEqual(5.0, row["baixa_over360_pl"])
        self.assertAlmostEqual(0.0, row["pdd_ex360"])
        self.assertAlmostEqual(95.0, row["pl_total_ex360"])
        self.assertAlmostEqual(20.0, row["pl_subordinada_mezz_ex360"])
        self.assertIn("PDD menor que Over 360", row["warnings"])

    def test_pl_chart_uses_ex360_subordination_series(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=5.0,
            buckets={8: 10.0},
        )
        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)

        chart_payload = json.dumps(pl_subordination_chart(monthly).to_dict(), ensure_ascii=False)

        self.assertIn("Subordinada + Mez ex-360", chart_payload)
        self.assertIn("% Subordinação Total ex-360", chart_payload)
        self.assertIn("21,1%", chart_payload)
        for legacy_color in ("#1f77b4", "#8ecae6", "#d1495b", "#1b998b"):
            self.assertNotIn(legacy_color, chart_payload)

    def test_npl_chart_uses_meli_palette_and_last_point_labels(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 8: 10.0},
        )
        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)

        chart_payload = json.dumps(npl_coverage_chart(monthly).to_dict(), ensure_ascii=False)

        self.assertIn("#000000", chart_payload)
        self.assertIn("#E47811", chart_payload)
        self.assertIn("Séries", chart_payload)
        self.assertIn("NPL Over 90d ex-360 / Carteira", chart_payload)
        self.assertIn("PDD Ex / NPL Over 90d ex-360", chart_payload)
        self.assertIn("4,0%", chart_payload)
        self.assertIn("225,0%", chart_payload)
        for legacy_color in ("#1f77b4", "#8ecae6", "#d1495b", "#1b998b"):
            self.assertNotIn(legacy_color, chart_payload)

    def test_official_pl_overrides_class_sum_and_keeps_reconciliation(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        official_pl = pd.DataFrame(
            [{"competencia": "01/2026", "pl_total_oficial": 110.0, "pl_total_oficial_source_status": "reported_value"}]
        )

        monthly = build_fund_monthly_base(
            cnpj="11111111000111",
            fund_name="FIDC A",
            dashboard=dashboard,
            official_pl_history_df=official_pl,
        )
        row = monthly.iloc[0]

        self.assertAlmostEqual(110.0, row["pl_total"])
        self.assertAlmostEqual(100.0, row["pl_total_classes"])
        self.assertAlmostEqual(10.0, row["pl_reconciliacao_delta"])
        self.assertEqual("PATRLIQ/VL_PATRIM_LIQ", row["pl_total_usado_fonte"])
        self.assertIn("PL oficial diverge", row["warnings"])

    def test_consolidated_base_sums_absolute_values_and_recalculates_ratios(self) -> None:
        dashboard_a = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        dashboard_b = _dashboard(
            fund_name="FIDC B",
            cnpj="22222222000122",
            pl_total=200.0,
            pl_senior=160.0,
            pl_mezz=0.0,
            pl_sub=40.0,
            carteira=2_000.0,
            pdd=200.0,
            buckets={4: 80.0, 7: 20.0, 8: 0.0},
        )
        monthly_a = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard_a)
        monthly_b = build_fund_monthly_base(cnpj="22222222000122", fund_name="FIDC B", dashboard=dashboard_b)

        consolidated = build_consolidated_monthly_base(
            portfolio_name="Carteira",
            fund_monthly_frames={"11111111000111": monthly_a, "22222222000122": monthly_b},
        )
        row = consolidated.iloc[0]

        self.assertAlmostEqual(300.0, row["pl_total"])
        self.assertAlmostEqual(65.0, row["pl_subordinada_mezz"])
        self.assertAlmostEqual(65.0 / 300.0 * 100.0, row["subordinacao_total_pct"], places=6)
        self.assertAlmostEqual(160.0, row["npl_over90"])
        self.assertAlmostEqual(300.0 / 160.0 * 100.0, row["pdd_npl_over90_pct"], places=6)

    def test_roll_rate_uses_current_bucket_over_previous_month_base(self) -> None:
        competencias = ["01/2026", "02/2026"]
        dashboard = SimpleNamespace(
            competencias=competencias,
            fund_info={"nome_fundo": "FIDC A", "cnpj_fundo": "11111111000111", "nome_classe": "Classe única"},
            subordination_history_df=pd.DataFrame(
                [
                    {
                        "competencia": competencia,
                        "competencia_dt": pd.Timestamp(f"2026-{month:02d}-01"),
                        "pl_total": 100.0,
                        "pl_senior": 75.0,
                        "pl_mezzanino": 15.0,
                        "pl_subordinada_strict": 10.0,
                        "pl_subordinada": 25.0,
                    }
                    for month, competencia in enumerate(competencias, start=1)
                ]
            ),
            default_history_df=pd.DataFrame(
                [
                    {
                        "competencia": competencia,
                        "competencia_dt": pd.Timestamp(f"2026-{month:02d}-01"),
                        "direitos_creditorios": 1_000.0,
                        "direitos_creditorios_fonte": "teste",
                        "provisao_total": 100.0,
                    }
                    for month, competencia in enumerate(competencias, start=1)
                ]
            ),
            dc_canonical_history_df=pd.DataFrame(),
            default_buckets_history_df=pd.DataFrame(
                [
                    {"competencia": "01/2026", "competencia_dt": pd.Timestamp("2026-01-01"), "ordem": 1, "faixa": "Até 30 dias", "valor": 100.0},
                    {"competencia": "01/2026", "competencia_dt": pd.Timestamp("2026-01-01"), "ordem": 2, "faixa": "31 a 60 dias", "valor": 20.0},
                    {"competencia": "02/2026", "competencia_dt": pd.Timestamp("2026-02-01"), "ordem": 1, "faixa": "Até 30 dias", "valor": 80.0},
                    {"competencia": "02/2026", "competencia_dt": pd.Timestamp("2026-02-01"), "ordem": 2, "faixa": "31 a 60 dias", "valor": 49.0},
                ]
            ),
        )

        monthly = build_fund_monthly_base(cnpj="11111111000111", fund_name="FIDC A", dashboard=dashboard)
        jan = monthly.loc[monthly["competencia"] == "01/2026"].iloc[0]
        feb = monthly.loc[monthly["competencia"] == "02/2026"].iloc[0]

        self.assertAlmostEqual(980.0, jan["carteira_em_dia_mais_ate30"])
        self.assertTrue(pd.isna(jan["roll_rate_base_t_minus_1"]))
        self.assertAlmostEqual(980.0, feb["roll_rate_base_t_minus_1"])
        self.assertAlmostEqual(5.0, feb["roll_rate_31_60_pct"])

        wide = build_wide_table(monthly, scope_name="FIDC A")
        base_row = wide.loc[wide["Métrica"] == "Carteira em dia + atrasada até 30d (t-1)"].iloc[0]
        roll_row = wide.loc[wide["Métrica"] == "Roll Rate"].iloc[0]
        self.assertEqual("R$ 980,00", base_row["fev/26"])
        self.assertEqual("N/D", base_row["jan/26"])
        self.assertEqual("5,00%", roll_row["fev/26"])

    def test_wide_table_and_excel_export_include_required_blocks(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        wide = build_wide_table(outputs.fund_monthly["11111111000111"], scope_name="FIDC A")

        self.assertIn("jan/26", wide.columns)
        self.assertIn("7. Visão Ex-Vencidos > 360d", set(wide["Bloco"]))
        self.assertIn("NPL Over 90d", set(wide["Métrica"]))

        excel_bytes = build_excel_export_bytes(outputs)
        workbook = load_workbook(BytesIO(excel_bytes), data_only=True)
        self.assertIn("Consolidado", workbook.sheetnames)
        self.assertIn("Auditoria", workbook.sheetnames)
        self.assertTrue(zipfile.is_zipfile(BytesIO(excel_bytes)))
        ws = workbook["Consolidado"]
        header = [cell.value for cell in ws[1]]
        self.assertEqual("Métrica", header[0])
        self.assertEqual("Memória / fórmula", header[-1])
        pl_row = next(row for row in ws.iter_rows(min_row=2, values_only=False) if row[0].value == "PL FIDC total")
        self.assertIsInstance(pl_row[1].value, (int, float))
        self.assertNotIsInstance(pl_row[1].value, str)
        sub_row = next(row for row in ws.iter_rows(min_row=2, values_only=False) if row[0].value == "% Subordinação Total")
        self.assertAlmostEqual(0.25, sub_row[1].value, places=8)
        self.assertIn("%", sub_row[1].number_format)

        snapshot_bytes = build_consolidated_snapshot_excel_bytes(outputs)
        snapshot_workbook = load_workbook(BytesIO(snapshot_bytes), data_only=True)
        self.assertIn("Resumo exibido", snapshot_workbook.sheetnames)
        self.assertIn("Dados gráficos", snapshot_workbook.sheetnames)
        self.assertIn("Gráficos", snapshot_workbook.sheetnames)
        self.assertGreaterEqual(len(snapshot_workbook["Gráficos"]._charts), 2)

        pptx_bytes = build_pptx_export_bytes(outputs)
        self.assertTrue(pptx_bytes.startswith(b"PK"))
        self.assertTrue(zipfile.is_zipfile(BytesIO(pptx_bytes)))
        from pptx import Presentation

        presentation = Presentation(BytesIO(pptx_bytes))
        self.assertEqual(len(outputs.fund_monthly) + 1, len(presentation.slides))
        for slide in presentation.slides:
            chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
            self.assertEqual(4, len(chart_shapes))
            self.assertTrue(all(shape.chart.has_title for shape in chart_shapes))
        with zipfile.ZipFile(BytesIO(pptx_bytes)) as archive:
            names = archive.namelist()
            xml_payload = "\n".join(
                archive.read(name).decode("utf-8", errors="ignore")
                for name in names
                if name.endswith(".xml")
            )
            chart_xml = "\n".join(
                archive.read(name).decode("utf-8")
                for name in names
                if name.startswith("ppt/charts/chart") and name.endswith(".xml")
            )
        self.assertTrue(any(name.startswith("ppt/charts/chart") for name in names))
        self.assertIn("<c:dLbls>", chart_xml)
        self.assertIn("<c:dLblPos", chart_xml)
        self.assertIn("<c:title>", chart_xml)
        self.assertIn("PL total:", xml_payload)
        self.assertIn('sz="1000"', xml_payload)

    def test_full_variable_export_has_equivalent_numeric_sheets_and_xml_source_names(self) -> None:
        dashboard_a = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={1: 25.0, 2: 10.0, 4: 40.0, 7: 10.0, 8: 10.0},
        )
        dashboard_b = _dashboard(
            fund_name="FIDC B",
            cnpj="22222222000122",
            pl_total=200.0,
            pl_senior=150.0,
            pl_mezz=25.0,
            pl_sub=25.0,
            carteira=2_000.0,
            pdd=150.0,
            buckets={1: 35.0, 2: 20.0, 4: 80.0, 7: 20.0, 8: 20.0},
        )
        return_months = pd.date_range("2025-07-01", "2026-06-01", freq="MS")
        return_values = [1.25, pd.NA, 0.0, -0.5, 1.0, 1.1, 1.2, 1.3, 1.4, 1.5, 1.6, 1.7]
        dashboard_a.return_history_df = pd.DataFrame(
            {
                "competencia": [value.strftime("%m/%Y") for value in return_months],
                "competencia_dt": return_months,
                "class_kind": ["senior"] * len(return_months),
                "class_key": ["senior:1"] * len(return_months),
                "class_label": ["Sênior 1"] * len(return_months),
                "retorno_mensal_pct": return_values,
            }
        )
        dashboard_a.return_summary_df = pd.DataFrame(
            {
                "class_kind": ["senior"],
                "class_key": ["senior:1"],
                "class_label": ["Sênior 1"],
                "latest_competencia": ["06/2026"],
                "retorno_12m_pct": [pd.NA],
                "retorno_ano_pct": [9.0],
                "ytd_status": ["completo"],
                "ytd_competencias_ausentes": [""],
            }
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={
                "11111111000111": ("FIDC A", dashboard_a),
                "22222222000122": ("FIDC B", dashboard_b),
            },
            period_label="01/2026 a 01/2026",
        )

        matrix = build_full_variable_export_matrix(outputs.consolidated_monthly)
        self.assertIn("Nome", matrix.columns)
        self.assertIn("Nome original da variável", matrix.columns)
        self.assertEqual("PL FIDC total", matrix.loc[matrix["Nome"] == "PL FIDC total", "Nome"].iloc[0])
        self.assertIn(
            "DOC_ARQ/LISTA_INFORM/PATRLIQ/VL_PATRIM_LIQ",
            matrix.loc[matrix["Nome"] == "PL FIDC oficial", "Nome original da variável"].iloc[0],
        )

        excel_bytes = build_full_variable_excel_export_bytes(outputs)
        workbook = load_workbook(BytesIO(excel_bytes), data_only=True)
        self.assertIn("Consolidado", workbook.sheetnames)
        self.assertEqual(
            ["Consolidado", "FIDC A", "Rent - FIDC A", "FIDC B", "Rent - FIDC B"],
            workbook.sheetnames,
        )
        consolidated = workbook["Consolidado"]
        fund_a = workbook["FIDC A"]
        fund_b = workbook["FIDC B"]
        self.assertEqual(
            [cell.value for cell in consolidated[1]],
            [cell.value for cell in fund_a[1]],
        )
        self.assertEqual(
            [row[1] for row in consolidated.iter_rows(min_row=2, values_only=True)],
            [row[1] for row in fund_b.iter_rows(min_row=2, values_only=True)],
        )
        pl_row = next(row for row in consolidated.iter_rows(min_row=2, values_only=False) if row[1].value == "PL FIDC total")
        carteira_row = next(row for row in consolidated.iter_rows(min_row=2, values_only=False) if row[1].value == "Carteira Bruta total")
        sub_row = next(row for row in consolidated.iter_rows(min_row=2, values_only=False) if row[1].value == "% Subordinação Total")
        self.assertAlmostEqual(300.0, pl_row[4].value)
        self.assertAlmostEqual(3_000.0, carteira_row[4].value)
        self.assertAlmostEqual(0.25, sub_row[4].value)
        self.assertIn("R$", pl_row[4].number_format)
        self.assertIn("%", sub_row[4].number_format)

        returns_a = workbook["Rent - FIDC A"]
        self.assertEqual(15, returns_a.max_column)
        self.assertEqual("Série", returns_a.cell(row=1, column=1).value)
        self.assertEqual("jul/25", returns_a.cell(row=1, column=2).value)
        self.assertEqual("Ac. Últ. 12m (%)", returns_a.cell(row=1, column=14).value)
        self.assertEqual("Acumulado YTD (%)", returns_a.cell(row=1, column=15).value)
        self.assertEqual("Sênior 1", returns_a.cell(row=2, column=1).value)
        self.assertAlmostEqual(0.0125, returns_a.cell(row=2, column=2).value)
        self.assertIsNone(returns_a.cell(row=2, column=3).value)
        self.assertEqual("0.00%", returns_a.cell(row=2, column=2).number_format)
        self.assertEqual("0.00%", returns_a.cell(row=2, column=15).number_format)
        self.assertEqual("solid", returns_a.cell(row=1, column=1).fill.fill_type)
        self.assertEqual("001F2937", returns_a.cell(row=1, column=1).fill.fgColor.rgb)
        self.assertEqual("B2", returns_a.freeze_panes)
        self.assertFalse(returns_a.sheet_view.showGridLines)
        self.assertEqual(1, len(returns_a.tables))
        self.assertEqual("A1:O2", next(iter(returns_a.tables.values())).ref)

        returns_b = workbook["Rent - FIDC B"]
        self.assertEqual(1, returns_b.max_row)
        self.assertEqual(1, len(returns_b.tables))
        self.assertEqual("A1:C1", next(iter(returns_b.tables.values())).ref)
        self.assertNotEqual(
            next(iter(returns_a.tables.values())).displayName,
            next(iter(returns_b.tables.values())).displayName,
        )

        comparison_export = build_full_variable_excel_export_bytes(
            outputs,
            monitor_outputs=build_meli_monitor_outputs(outputs),
        )
        comparison_workbook = load_workbook(BytesIO(comparison_export), data_only=True)
        self.assertIn("Comparativo crédito", comparison_workbook.sheetnames)
        comparison = comparison_workbook["Comparativo crédito"]
        self.assertEqual(4, comparison.max_row)
        self.assertEqual("Consolidado", comparison.cell(row=2, column=1).value)
        self.assertEqual("Fundo", comparison.cell(row=3, column=1).value)
        self.assertEqual("01/2026", comparison.cell(row=2, column=4).value)
        self.assertEqual("0.00%", comparison.cell(row=2, column=8).number_format)
        self.assertEqual(1, len(comparison.tables))

        csv_zip = build_full_variable_csv_zip_bytes(outputs)
        with zipfile.ZipFile(BytesIO(csv_zip)) as archive:
            self.assertIn("consolidado.csv", archive.namelist())
            self.assertTrue(any(name.endswith("fidc_a.csv") for name in archive.namelist()))
            csv_payload = archive.read("consolidado.csv").decode("utf-8-sig")
        self.assertIn("Nome original da variável", csv_payload)
        self.assertIn("PL FIDC total", csv_payload)

    def test_pptx_export_uses_one_2x2_slide_per_fund_plus_consolidated(self) -> None:
        dashboard_a = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        dashboard_b = _dashboard(
            fund_name="FIDC B",
            cnpj="22222222000122",
            pl_total=200.0,
            pl_senior=160.0,
            pl_mezz=20.0,
            pl_sub=20.0,
            carteira=2_000.0,
            pdd=200.0,
            buckets={4: 80.0, 7: 20.0, 8: 20.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-2",
            portfolio_name="Carteira",
            dashboards_by_cnpj={
                "11111111000111": ("FIDC A", dashboard_a),
                "22222222000122": ("FIDC B", dashboard_b),
            },
            period_label="01/2026 a 01/2026",
        )
        from pptx import Presentation

        presentation = Presentation(BytesIO(build_pptx_export_bytes(outputs)))

        self.assertEqual(3, len(presentation.slides))
        for slide in presentation.slides:
            self.assertEqual(4, sum(1 for shape in slide.shapes if getattr(shape, "has_chart", False)))

    def test_wide_table_orders_periods_newest_to_oldest(self) -> None:
        monthly = pd.DataFrame(
            [
                {
                    "fund_name": "FIDC A",
                    "cnpj": "11111111000111",
                    "tipo_classe": "Classe",
                    "competencia": "11/2025",
                    "competencia_dt": pd.Timestamp("2025-11-01"),
                    "pl_total": 100.0,
                    "pl_senior": 75.0,
                    "pl_subordinada_mezz": 25.0,
                    "subordinacao_total_pct": 25.0,
                },
                {
                    "fund_name": "FIDC A",
                    "cnpj": "11111111000111",
                    "tipo_classe": "Classe",
                    "competencia": "03/2026",
                    "competencia_dt": pd.Timestamp("2026-03-01"),
                    "pl_total": 110.0,
                    "pl_senior": 80.0,
                    "pl_subordinada_mezz": 30.0,
                    "subordinacao_total_pct": 27.2727,
                },
                {
                    "fund_name": "FIDC A",
                    "cnpj": "11111111000111",
                    "tipo_classe": "Classe",
                    "competencia": "01/2026",
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "pl_total": 105.0,
                    "pl_senior": 78.0,
                    "pl_subordinada_mezz": 27.0,
                    "subordinacao_total_pct": 25.7143,
                },
            ]
        )
        wide = build_wide_table(monthly, scope_name="FIDC A")
        period_columns = order_period_columns_desc(wide.columns)

        self.assertEqual(["mar/26", "jan/26", "nov/25"], period_columns)
        self.assertLess(wide.columns.get_loc("mar/26"), wide.columns.get_loc("jan/26"))
        self.assertEqual("Memória / fórmula", wide.columns[-1])

    def test_consolidated_wide_html_has_section_rows_and_dense_formatting(self) -> None:
        wide = pd.DataFrame(
            [
                {
                    "Bloco": "2. PL FIDC",
                    "Métrica": "PL FIDC total",
                    "Memória / fórmula": "PATRLIQ/VL_PATRIM_LIQ",
                    "jan/26": "R$ mm 1.500,00",
                    "fev/26": "N/D",
                },
                {
                    "Bloco": "2. PL FIDC",
                    "Métrica": "% Subordinação Total",
                    "Memória / fórmula": "(Subordinada + Mezanino) / PL FIDC total.",
                    "jan/26": "27,25%",
                    "fev/26": "28,75%",
                },
                {
                    "Bloco": "7. Visão Ex-Vencidos > 360d",
                    "Métrica": "NPL Over 90d Ex 360 / Carteira Ex 360",
                    "Memória / fórmula": "NPL Over 90d Ex 360 / Carteira Ex 360.",
                    "jan/26": "12,50%",
                    "fev/26": "13,75%",
                },
                {
                    "Bloco": "7. Visão Ex-Vencidos > 360d",
                    "Métrica": "PDD / NPL Over 90d Ex 360",
                    "Memória / fórmula": "PDD Ex Over 360d / NPL Over 90d Ex 360.",
                    "jan/26": "100,00%",
                    "fev/26": "110,00%",
                },
            ]
        )

        html = _render_wide_table_html(wide)

        self.assertIn("wide-table-wrapper", html)
        self.assertIn("<details class='wide-section' style=", html)
        self.assertNotIn("<details class='wide-section' open", html)
        self.assertIn("<col class='label-col-width' style='width: 280px;'>", html)
        self.assertEqual(4, html.count("<col class='period-col-width' style='width: 96px;'>"))
        self.assertIn("<col class='formula-col-width' style='width: 340px;'>", html)
        self.assertIn("<summary>PL FIDC</summary>", html)
        self.assertIn("PL FIDC", html)
        self.assertIn("1.500,0 MM", html)
        self.assertIn("27,2%", html)
        self.assertIn("<summary>Visão Ex-Vencidos &gt; 360d</summary>", html)
        self.assertIn("<tr class='destaque'>\n<td class='label'>NPL Over 90d Ex 360 / Carteira Ex 360</td>", html)
        self.assertIn("<tr class='destaque'>\n<td class='label'>PDD / NPL Over 90d Ex 360</td>", html)
        self.assertNotIn(">N/D<", html)

    def test_dense_wide_value_uses_empty_for_missing_cells_and_br_format(self) -> None:
        self.assertEqual("", _dense_wide_value("N/D"))
        self.assertEqual("178,0%", _dense_wide_value("178,04%"))
        self.assertEqual("5.380,0 MM", _dense_wide_value("R$ mm 5.380,00"))

    def test_base_scope_options_show_consolidated_only_for_multi_fund_portfolios(self) -> None:
        portfolio = PortfolioRecord(
            id="portfolio-1",
            name="Carteira Teste",
            funds=(
                PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),
                PortfolioFund(cnpj="22222222000122", display_name="FIDC B"),
            ),
            created_at="2026-05-28T00:00:00Z",
            updated_at="2026-05-28T00:00:00Z",
        )
        outputs = SimpleNamespace(
            fund_monthly={
                "11111111000111": pd.DataFrame({"fund_name": ["FIDC A"]}),
                "22222222000122": pd.DataFrame({"fund_name": ["FIDC B"]}),
            },
            fund_wide={
                "11111111000111": pd.DataFrame({"Métrica": ["PL"]}),
                "22222222000122": pd.DataFrame({"Métrica": ["PL"]}),
            },
            consolidated_monthly=pd.DataFrame({"fund_name": ["Carteira Teste"]}),
            consolidated_wide=pd.DataFrame({"Métrica": ["PL"]}),
        )

        options = _build_base_scope_options(display_outputs=outputs, selected_portfolio=portfolio)

        self.assertEqual(CONSOLIDATED_SCOPE_VALUE, options[0].value)
        self.assertEqual("consolidated", options[0].kind)
        self.assertEqual(["fund::11111111000111", "fund::22222222000122"], [option.value for option in options[1:]])

    def test_base_scope_options_hide_consolidated_for_single_fund_portfolios(self) -> None:
        portfolio = PortfolioRecord(
            id="portfolio-1",
            name="Carteira Teste",
            funds=(PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),),
            created_at="2026-05-28T00:00:00Z",
            updated_at="2026-05-28T00:00:00Z",
        )
        outputs = SimpleNamespace(
            fund_monthly={"11111111000111": pd.DataFrame({"fund_name": ["FIDC A"]})},
            fund_wide={"11111111000111": pd.DataFrame({"Métrica": ["PL"]})},
            consolidated_monthly=pd.DataFrame({"fund_name": ["Carteira Teste"]}),
            consolidated_wide=pd.DataFrame({"Métrica": ["PL"]}),
        )

        options = _build_base_scope_options(display_outputs=outputs, selected_portfolio=portfolio)

        self.assertEqual(1, len(options))
        self.assertEqual("fund::11111111000111", options[0].value)
        self.assertEqual("fund", options[0].kind)
        self.assertEqual("FIDC A · 11111111000111", options[0].label)

    def test_default_base_scope_uses_consolidated_only_when_multiple_funds_exist(self) -> None:
        portfolio = PortfolioRecord(
            id="portfolio-1",
            name="Carteira Teste",
            funds=(
                PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),
                PortfolioFund(cnpj="22222222000122", display_name="FIDC B"),
            ),
            created_at="2026-05-28T00:00:00Z",
            updated_at="2026-05-28T00:00:00Z",
        )
        outputs = SimpleNamespace(
            fund_monthly={
                "11111111000111": pd.DataFrame({"fund_name": ["FIDC A"]}),
                "22222222000122": pd.DataFrame({"fund_name": ["FIDC B"]}),
            },
            fund_wide={
                "11111111000111": pd.DataFrame({"Métrica": ["PL"]}),
                "22222222000122": pd.DataFrame({"Métrica": ["PL"]}),
            },
            consolidated_monthly=pd.DataFrame({"fund_name": ["Carteira Teste"]}),
            consolidated_wide=pd.DataFrame({"Métrica": ["PL"]}),
        )

        option = _resolve_default_base_scope(display_outputs=outputs, selected_portfolio=portfolio)

        self.assertIsNotNone(option)
        self.assertEqual(CONSOLIDATED_SCOPE_VALUE, option.value)
        self.assertEqual("consolidated", option.kind)

    def test_default_base_scope_keeps_single_fund_direct_without_consolidated_label(self) -> None:
        portfolio = PortfolioRecord(
            id="portfolio-1",
            name="Carteira Teste",
            funds=(PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),),
            created_at="2026-05-28T00:00:00Z",
            updated_at="2026-05-28T00:00:00Z",
        )
        outputs = SimpleNamespace(
            fund_monthly={"11111111000111": pd.DataFrame({"fund_name": ["FIDC A"]})},
            fund_wide={"11111111000111": pd.DataFrame({"Métrica": ["PL"]})},
            consolidated_monthly=pd.DataFrame({"fund_name": ["Carteira Teste"]}),
            consolidated_wide=pd.DataFrame({"Métrica": ["PL"]}),
        )

        option = _resolve_default_base_scope(display_outputs=outputs, selected_portfolio=portfolio)

        self.assertIsNotNone(option)
        self.assertEqual("fund::11111111000111", option.value)
        self.assertEqual("fund", option.kind)

    def test_display_wide_table_removes_audit_blocks_but_keeps_formula_memory(self) -> None:
        wide = pd.DataFrame(
            [
                {"Bloco": "2. PL FIDC", "Métrica": "PL FIDC total", "Memória / fórmula": "PATRLIQ"},
                {"Bloco": "9. Campos auxiliares de auditoria", "Métrica": "Fonte", "Memória / fórmula": "debug"},
            ]
        )

        display = _display_wide_table(wide)

        self.assertEqual(["PL FIDC total"], display["Métrica"].tolist())
        self.assertIn("Memória / fórmula", display.columns)

    def test_wide_table_uses_money_scale_by_metric_for_readability(self) -> None:
        monthly = pd.DataFrame(
            [
                {
                    "scope": "consolidado",
                    "fund_name": "Carteira",
                    "cnpj": "CONSOLIDADO",
                    "competencia": "01/2026",
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "carteira_bruta": 6_200_000_000.0,
                    "pdd_total": 320_000_000.0,
                }
            ]
        )

        wide = build_wide_table(monthly, scope_name="Carteira")
        carteira = wide.loc[wide["Métrica"] == "Carteira Bruta total", "jan/26"].iloc[0]
        pdd = wide.loc[wide["Métrica"] == "PDD total", "jan/26"].iloc[0]

        self.assertEqual("R$ mm 6.200,0", carteira)
        self.assertEqual("R$ mm 320,0", pdd)

    def test_wide_table_uses_billion_scale_only_above_one_trillion(self) -> None:
        monthly = pd.DataFrame(
            [
                {
                    "scope": "consolidado",
                    "fund_name": "Carteira",
                    "cnpj": "CONSOLIDADO",
                    "competencia": "01/2026",
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "carteira_bruta": 1_200_000_000_000.0,
                }
            ]
        )

        wide = build_wide_table(monthly, scope_name="Carteira")
        carteira = wide.loc[wide["Métrica"] == "Carteira Bruta total", "jan/26"].iloc[0]

        self.assertEqual("R$ bi 1.200,0", carteira)

    def test_wide_table_text_columns_wrap_in_fixed_grid(self) -> None:
        self.assertIn("table-layout: fixed;", _MERCADO_LIVRE_UI_CSS)
        self.assertIn("vertical-align: top;", _MERCADO_LIVRE_UI_CSS)
        self.assertIn("overflow-wrap: anywhere;", _MERCADO_LIVRE_UI_CSS)
        self.assertIn("white-space: normal;", _MERCADO_LIVRE_UI_CSS)
        self.assertIn(".meli-kpi-grid", _MERCADO_LIVRE_UI_CSS)
        self.assertIn(".meli-kpi-card", _MERCADO_LIVRE_UI_CSS)
        td_block = _MERCADO_LIVRE_UI_CSS.split(".wide-table td {", 1)[1].split("}", 1)[0]
        self.assertNotIn("white-space: nowrap;", td_block)
        self.assertIn("white-space: normal;", td_block)
        self.assertIn("overflow-wrap: anywhere;", td_block)
        self.assertIn("line-height: 1.25;", td_block)

    def test_mercado_livre_guide_documents_usage_and_mechanics(self) -> None:
        guide = _build_mercado_livre_guide_markdown()

        self.assertIn("Como usar", guide)
        self.assertIn("Base analítica", guide)
        self.assertNotIn("Tabela Completa", guide)
        self.assertIn("Carteira de Crédito", guide)
        self.assertIn("Mecânica essencial", guide)
        self.assertIn("não há média simples de percentuais", guide)
        self.assertIn("NPL Over é acumulado", guide)
        self.assertIn("ex-360 remove vencidos acima de 360 dias", guide)
        self.assertIn("filtro visual", guide.lower())
        self.assertNotIn("tabelas wide", guide.lower())

    def test_display_window_full_option_keeps_loaded_36_months_by_default(self) -> None:
        available = [
            date(2023 + (idx + 3) // 12, ((idx + 3) % 12) + 1, 1)
            for idx in range(36)
        ]
        start, end = _display_window_bounds(
            selected="Todo período carregado",
            loaded_start=available[0],
            loaded_end=available[-1],
        )

        self.assertEqual(available[0], start)
        self.assertEqual(available[-1], end)

    def test_display_window_12m_filters_only_when_explicitly_selected(self) -> None:
        available = [
            date(2023 + (idx + 3) // 12, ((idx + 3) % 12) + 1, 1)
            for idx in range(36)
        ]
        start, end = _display_window_bounds(
            selected="12M",
            loaded_start=available[0],
            loaded_end=available[-1],
        )

        self.assertEqual(date(2025, 3, 1), start)
        self.assertEqual(available[-1], end)

    def test_display_window_12m_anchors_to_latest_available_when_requested_end_is_missing(self) -> None:
        available = [
            date(2024 + (idx + 1) // 12, ((idx + 1) % 12) + 1, 1)
            for idx in range(26)
            if date(2024 + (idx + 1) // 12, ((idx + 1) % 12) + 1, 1) <= date(2026, 3, 1)
        ]

        months = _display_window_months(selected="12M", available=available)

        self.assertEqual(date(2025, 3, 1), months[0])
        self.assertEqual(date(2026, 3, 1), months[-1])
        self.assertEqual(13, len(months))

    def test_fund_return_table_anchors_ytd_months_independently_by_fund(self) -> None:
        def return_frames(*, latest_month: int) -> tuple[pd.DataFrame, pd.DataFrame]:
            competencias = pd.date_range("2026-01-01", periods=latest_month, freq="MS")
            history = pd.DataFrame(
                {
                    "competencia": [value.strftime("%m/%Y") for value in competencias],
                    "competencia_dt": competencias,
                    "class_kind": ["senior"] * len(competencias),
                    "class_key": ["senior:1"] * len(competencias),
                    "class_label": ["Sênior"] * len(competencias),
                    "retorno_mensal_pct": [1.0] * len(competencias),
                }
            )
            summary = pd.DataFrame(
                {
                    "class_kind": ["senior"],
                    "class_key": ["senior:1"],
                    "class_label": ["Sênior"],
                    "latest_competencia": [f"{latest_month:02d}/2026"],
                    "retorno_ano_pct": [((1.01**latest_month) - 1.0) * 100.0],
                    "ytd_status": ["completo"],
                    "ytd_competencias_ausentes": [""],
                }
            )
            return history, summary

        history_a, summary_a = return_frames(latest_month=6)
        history_b, summary_b = return_frames(latest_month=5)
        outputs = SimpleNamespace(
            fund_return_history={"A": history_a, "B": history_b},
            fund_return_summary={"A": summary_a, "B": summary_b},
        )

        table_a = _build_fund_return_table(outputs=outputs, cnpj="A")
        table_b = _build_fund_return_table(outputs=outputs, cnpj="B")

        self.assertIn("jun/26", table_a.columns)
        self.assertNotIn("jun/26", table_b.columns)
        self.assertIn("mai/26", table_b.columns)
        self.assertEqual("6,15%", table_a.iloc[0][RETURN_YTD_COLUMN])
        self.assertEqual("5,10%", table_b.iloc[0][RETURN_YTD_COLUMN])

    def test_resolve_fund_return_cdi_rates_uses_latest_competencia_window_without_network(self) -> None:
        summary = pd.DataFrame(
            {
                "class_kind": ["senior"],
                "class_key": ["senior:1"],
                "class_label": ["Sênior"],
                "latest_competencia": ["06/2026"],
            }
        )
        outputs = SimpleNamespace(
            fund_return_history={"fund": pd.DataFrame()},
            fund_return_summary={"fund": summary},
        )
        rates = (
            B3CdiMonthlyRate(
                mes="2026-06",
                cdi_mensal=0.01,
                dias_uteis=21,
                data_inicio=date(2026, 6, 1),
                data_fim=date(2026, 6, 30),
                source="fixture",
                expected_dias_uteis=21,
            ),
        )

        with patch(
            "tabs.tab_dashboard_meli._fetch_fund_return_cdi_rates",
            return_value=rates,
        ) as fetch_rates:
            resolved, error = _resolve_fund_return_cdi_rates(outputs=outputs, cnpj="fund")

        self.assertEqual(rates, resolved)
        self.assertIsNone(error)
        fetch_rates.assert_called_once_with("2025-07-01", "2026-06-30")

    def test_resolve_fund_return_export_inputs_reuses_subtab_resolvers(self) -> None:
        outputs = SimpleNamespace()
        rates = (
            B3CdiMonthlyRate(
                mes="2026-06",
                cdi_mensal=0.01,
                dias_uteis=21,
                data_inicio=date(2026, 6, 1),
                data_fim=date(2026, 6, 30),
                source="fixture",
            ),
        )
        with (
            patch(
                "tabs.tab_dashboard_meli._resolve_fund_return_cdi_rates",
                return_value=(rates, None),
            ) as resolve_cdi,
            patch(
                "tabs.tab_dashboard_meli._resolve_fund_return_benchmark",
                return_value=SimpleNamespace(spreads_by_class_key={"senior:1": 0.035}),
            ) as resolve_benchmark,
        ):
            cdi_by_fund, benchmark_by_fund = resolve_fund_return_export_inputs(
                outputs=outputs,
                cnpjs=["fund"],
            )

        self.assertEqual({"fund": rates}, cdi_by_fund)
        self.assertEqual({"fund": {"senior:1": 0.035}}, benchmark_by_fund)
        resolve_cdi.assert_called_once_with(outputs=outputs, cnpj="fund")
        resolve_benchmark.assert_called_once_with(outputs=outputs, cnpj="fund")

    def test_stacked_single_fund_renders_return_table_without_duplicate_charts(self) -> None:
        outputs = SimpleNamespace()
        monitor_outputs = SimpleNamespace(fund_monitor={"11111111000111": pd.DataFrame()})
        portfolio = SimpleNamespace(id="portfolio-1")

        with (
            patch("tabs.tab_dashboard_meli._render_fund_return_table") as render_return_table,
            patch("tabs.tab_dashboard_meli._render_fund_dashboards") as render_fund_dashboards,
        ):
            _render_stacked_funds_view(
                outputs=outputs,
                monitor_outputs=monitor_outputs,
                selected_portfolio=portfolio,
            )

        render_return_table.assert_called_once_with(outputs=outputs, cnpj="11111111000111")
        render_fund_dashboards.assert_not_called()

    def test_inline_portfolio_renders_fund_returns_before_consolidated_charts(self) -> None:
        events: list[str] = []
        outputs = SimpleNamespace()
        monitor_outputs = SimpleNamespace(
            consolidated_monitor=pd.DataFrame(),
            fund_monitor={"11111111000111": pd.DataFrame()},
        )
        portfolio = SimpleNamespace(id="portfolio-1", name="Carteira")

        with (
            patch("tabs.tab_dashboard_meli._render_kpis"),
            patch(
                "tabs.tab_dashboard_meli._render_stacked_funds_view",
                side_effect=lambda **_kwargs: events.append("fundos"),
            ),
            patch(
                "tabs.tab_dashboard_meli._render_consolidated_dashboard",
                side_effect=lambda *_args, **_kwargs: events.append("consolidado"),
            ),
            patch("tabs.tab_dashboard_meli._render_audit"),
            patch("tabs.tab_dashboard_meli._render_methodology"),
            patch("tabs.tab_dashboard_meli.st.expander"),
        ):
            render_dashboard_meli_analysis(
                outputs=outputs,
                selected_portfolio=portfolio,
                monitor_outputs=monitor_outputs,
                research_outputs=SimpleNamespace(),
                verification_report=pd.DataFrame(),
                pptx_bytes=b"pptx",
                use_tabs=False,
                show_guide=False,
                show_downloads=False,
            )

        self.assertEqual(["fundos", "consolidado"], events)

    def test_requested_period_is_loaded_with_yoy_lookback_metadata(self) -> None:
        requested = build_custom_period(start_month=date(2025, 5, 1), end_month=date(2026, 4, 1))
        calculation = _period_with_yoy_lookback(requested)
        outputs = MercadoLivreOutputs(
            fund_monthly={},
            fund_wide={},
            consolidated_monthly=pd.DataFrame(),
            consolidated_wide=pd.DataFrame(),
            warnings_df=pd.DataFrame(),
            metadata={},
        )

        tagged = _tag_outputs_requested_period(outputs, requested_period=requested, calculation_period=calculation)

        self.assertEqual(date(2024, 5, 1), calculation.start_month)
        self.assertEqual("12M", tagged.metadata["requested_window_option"])
        self.assertEqual("2025-05-01", tagged.metadata["requested_period_start"])
        self.assertEqual("2024-05-01", tagged.metadata["calculation_period_start"])

    def test_credit_monitor_keeps_yoy_values_after_display_filter(self) -> None:
        months = pd.date_range("2025-01-01", "2026-12-01", freq="MS")
        monthly = pd.DataFrame(
            {
                "competencia": [f"{item.month:02d}/{item.year}" for item in months],
                "competencia_dt": months,
                "fund_name": "FIDC A",
                "carteira_ex360": [100.0 + idx for idx in range(len(months))],
            }
        )
        outputs = MercadoLivreOutputs(
            fund_monthly={"1": monthly},
            fund_wide={"1": pd.DataFrame()},
            consolidated_monthly=monthly,
            consolidated_wide=pd.DataFrame(),
            warnings_df=pd.DataFrame(),
            metadata={
                "display_period_months": [value.date().isoformat() for value in months[-12:]],
                "display_period_start": months[-12].date().isoformat(),
                "display_period_end": months[-1].date().isoformat(),
            },
        )

        monitor = _build_credit_monitor_for_display(outputs=outputs, display_outputs=outputs)

        self.assertEqual(12, len(monitor.consolidated_monitor))
        self.assertFalse(monitor.consolidated_monitor["carteira_ex360_yoy_pct"].isna().any())
        self.assertAlmostEqual((112.0 / 100.0 - 1.0) * 100.0, monitor.consolidated_monitor.iloc[0]["carteira_ex360_yoy_pct"])

    def test_display_window_decembers_plus_current_year_uses_non_contiguous_months(self) -> None:
        available = [
            date(2023, 11, 1),
            date(2023, 12, 1),
            date(2024, 11, 1),
            date(2025, 12, 1),
            date(2026, 1, 1),
            date(2026, 2, 1),
            date(2026, 3, 1),
            date(2026, 4, 1),
        ]

        months = _display_window_months(selected="Dezembros + Ano Atual", available=available)

        self.assertEqual(
            [
                date(2023, 12, 1),
                date(2025, 12, 1),
                date(2026, 1, 1),
                date(2026, 2, 1),
                date(2026, 3, 1),
                date(2026, 4, 1),
            ],
            months,
        )

    def test_filter_outputs_by_explicit_months_keeps_same_months_in_wide_and_metadata(self) -> None:
        monthly = pd.DataFrame(
            [
                {"competencia": "11/2023", "competencia_dt": pd.Timestamp("2023-11-01"), "fund_name": "FIDC A", "pl_total": 1.0},
                {"competencia": "12/2023", "competencia_dt": pd.Timestamp("2023-12-01"), "fund_name": "FIDC A", "pl_total": 2.0},
                {"competencia": "01/2026", "competencia_dt": pd.Timestamp("2026-01-01"), "fund_name": "FIDC A", "pl_total": 3.0},
                {"competencia": "04/2026", "competencia_dt": pd.Timestamp("2026-04-01"), "fund_name": "FIDC A", "pl_total": 4.0},
            ]
        )
        return_history = pd.DataFrame({"competencia": ["01/2026"], "retorno_mensal_pct": [1.0]})
        outputs = MercadoLivreOutputs(
            fund_monthly={"1": monthly},
            fund_wide={"1": pd.DataFrame()},
            consolidated_monthly=monthly,
            consolidated_wide=pd.DataFrame(),
            warnings_df=pd.DataFrame(),
            metadata={"portfolio_name": "Carteira"},
            fund_return_history={"1": return_history},
        )

        filtered = _filter_outputs_by_competencia_months(
            outputs,
            months=[date(2023, 12, 1), date(2026, 1, 1), date(2026, 4, 1)],
            label="dez/23 + jan/26 → abr/26",
            mode="Dezembros + Ano Atual",
        )

        self.assertEqual(["12/2023", "01/2026", "04/2026"], filtered.consolidated_monthly["competencia"].tolist())
        self.assertEqual(
            ["2023-12-01", "2026-01-01", "2026-04-01"],
            filtered.metadata["display_period_months"],
        )
        self.assertIn("abr/26", filtered.consolidated_wide.columns)
        self.assertIn("dez/23", filtered.consolidated_wide.columns)
        self.assertTrue(filtered.fund_return_history["1"].equals(return_history))

    def test_compact_wide_table_hides_operational_identification_rows_only_in_main_view(self) -> None:
        monthly = pd.DataFrame(
            [
                {
                    "competencia": "03/2026",
                    "competencia_dt": pd.Timestamp("2026-03-01"),
                    "fund_name": "Carteira Teste",
                    "cnpj": "CONSOLIDADO",
                    "pl_total": 100.0,
                    "carteira_bruta": 80.0,
                    "pdd_total": 5.0,
                    "npl_over90": 2.0,
                }
            ]
        )
        wide = build_wide_table(monthly, scope_name="Carteira Teste")

        compact = _display_wide_table(wide, compact=True)
        full = _display_wide_table(wide, compact=False)

        self.assertNotIn("Nome do fundo", compact["Métrica"].tolist())
        self.assertNotIn("Período final", compact["Métrica"].tolist())
        self.assertIn("PL FIDC total", compact["Métrica"].tolist())
        self.assertIn("NPL Over 90d / Carteira", compact["Métrica"].tolist())
        self.assertIn("Nome do fundo", full["Métrica"].tolist())
        self.assertIn("Período final", full["Métrica"].tolist())

    def test_snapshot_export_uses_all_displayed_months_instead_of_tail_six(self) -> None:
        months = pd.date_range("2023-12-01", "2026-04-01", freq="MS")
        monthly = pd.DataFrame(
            {
                "competencia": [f"{item.month:02d}/{item.year}" for item in months],
                "competencia_dt": months,
                "pl_total": range(len(months)),
                "pl_senior": range(len(months)),
                "pl_subordinada_mezz_ex360": range(len(months)),
                "subordinacao_total_ex360_pct": 10.0,
                "carteira_bruta": range(len(months)),
                "carteira_ex360": range(len(months)),
                "pdd_ex360": range(len(months)),
                "npl_over90_ex360": range(len(months)),
                "npl_over90_ex360_pct": 1.0,
                "pdd_npl_over90_ex360_pct": 100.0,
                "roll_rate_31_60_pct": 2.0,
            }
        )
        outputs = SimpleNamespace(consolidated_monthly=monthly)

        snapshot_bytes = build_consolidated_snapshot_excel_bytes(outputs)
        workbook = load_workbook(BytesIO(snapshot_bytes), data_only=True)
        summary_header = [cell.value for cell in workbook["Resumo exibido"][1]]
        chart_rows = list(workbook["Dados gráficos"].iter_rows(min_row=2, values_only=True))

        self.assertIn("dez/23", summary_header)
        self.assertIn("abr/26", summary_header)
        self.assertEqual(len(months), len(chart_rows))

    def test_outputs_cache_roundtrip_uses_deterministic_identity(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        dashboard.return_history_df = pd.DataFrame(
            {
                "competencia": ["01/2026"],
                "competencia_dt": pd.to_datetime(["2026-01-01"]),
                "class_kind": ["senior"],
                "class_key": ["senior:1"],
                "class_label": ["Sênior"],
                "retorno_mensal_pct": [1.25],
            }
        )
        dashboard.return_summary_df = pd.DataFrame(
            {
                "class_kind": ["senior"],
                "class_key": ["senior:1"],
                "class_label": ["Sênior"],
                "latest_competencia": ["01/2026"],
                "retorno_ano_pct": [1.25],
                "ytd_status": ["completo"],
                "ytd_competencias_ausentes": [""],
            }
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        funds = (PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),)

        with tempfile.TemporaryDirectory() as tmp_dir:
            root = save_outputs_to_cache(
                outputs,
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )
            loaded = load_outputs_from_cache(
                portfolio_id="outro-id-visual",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )
            self.assertTrue((root / "metadata.json").exists())
            metadata = json.loads((root / "metadata.json").read_text(encoding="utf-8"))

        self.assertIsNotNone(loaded)
        assert loaded is not None
        self.assertIn("11111111000111", loaded.fund_monthly)
        self.assertEqual(1.25, loaded.fund_return_history["11111111000111"].iloc[0]["retorno_mensal_pct"])
        self.assertEqual("01/2026", loaded.fund_return_summary["11111111000111"].iloc[0]["latest_competencia"])
        self.assertEqual(outputs.metadata["loaded_period_label"], loaded.metadata["loaded_period_label"])
        self.assertEqual(["11111111000111"], metadata["requested_funds"])
        self.assertEqual(["11111111000111"], metadata["loaded_funds"])
        self.assertTrue(metadata["cache_complete"])

    def test_outputs_cache_rejects_partial_portfolio_cache(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        requested_funds = (
            PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),
            PortfolioFund(cnpj="22222222000122", display_name="FIDC B"),
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            root = save_outputs_to_cache(
                outputs,
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=requested_funds,
                base_dir=Path(tmp_dir),
            )
            metadata = json.loads((root / "metadata.json").read_text(encoding="utf-8"))
            loaded = load_outputs_from_cache(
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=requested_funds,
                base_dir=Path(tmp_dir),
            )

        self.assertFalse(metadata["cache_complete"])
        self.assertEqual(["11111111000111", "22222222000122"], metadata["requested_funds"])
        self.assertEqual(["11111111000111"], metadata["loaded_funds"])
        self.assertIsNone(loaded)

    def test_outputs_cache_treats_empty_optional_warnings_as_empty_dataframe(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        funds = (PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),)

        with tempfile.TemporaryDirectory() as tmp_dir:
            root = save_outputs_to_cache(
                outputs,
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )
            (root / "warnings.csv").write_text("", encoding="utf-8")

            loaded = load_outputs_from_cache(
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )

        self.assertIsNotNone(loaded)
        assert loaded is not None
        self.assertTrue(loaded.warnings_df.empty)

    def test_outputs_cache_invalidates_empty_required_csvs(self) -> None:
        dashboard = _dashboard(
            fund_name="FIDC A",
            cnpj="11111111000111",
            pl_total=100.0,
            pl_senior=75.0,
            pl_mezz=15.0,
            pl_sub=10.0,
            carteira=1_000.0,
            pdd=100.0,
            buckets={4: 40.0, 7: 10.0, 8: 10.0},
        )
        outputs = build_mercado_livre_outputs(
            portfolio_id="portfolio-1",
            portfolio_name="Carteira",
            dashboards_by_cnpj={"11111111000111": ("FIDC A", dashboard)},
            period_label="01/2026 a 01/2026",
        )
        funds = (PortfolioFund(cnpj="11111111000111", display_name="FIDC A"),)

        with tempfile.TemporaryDirectory() as tmp_dir:
            root = save_outputs_to_cache(
                outputs,
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )
            (root / "monthly_consolidado.csv").write_text("", encoding="utf-8")

            loaded = load_outputs_from_cache(
                portfolio_id="portfolio-1",
                period_key="2026-01-01::2026-01-01",
                portfolio_funds=funds,
                base_dir=Path(tmp_dir),
            )

        self.assertIsNone(loaded)

    def test_portfolio_identity_key_is_deterministic_and_duplicate_save_reuses_same_basket(self) -> None:
        funds_a = (
            PortfolioFund(cnpj="11111111000111", display_name="A"),
            PortfolioFund(cnpj="22222222000122", display_name="B"),
        )
        funds_b = (
            PortfolioFund(cnpj="22222222000122", display_name="B"),
            PortfolioFund(cnpj="11111111000111", display_name="A"),
        )

        self.assertEqual(
            portfolio_identity_key(funds_a, fallback="p1"),
            portfolio_identity_key(funds_b, fallback="p2"),
        )

        existing = PortfolioRecord(
            id="portfolio-1",
            name="MELI (Todos)",
            funds=funds_a,
            created_at="2026-04-29T00:00:00Z",
            updated_at="2026-04-29T00:00:00Z",
        )
        action = _resolve_existing_portfolio_for_save(
            portfolios=[existing],
            target=None,
            name="Outro Nome",
            funds=list(funds_b),
        )

        self.assertEqual("reuse", action["action"])
        self.assertEqual(existing.id, action["portfolio"].id)


def _dashboard(
    *,
    fund_name: str,
    cnpj: str,
    pl_total: float,
    pl_senior: float,
    pl_mezz: float,
    pl_sub: float,
    carteira: float,
    pdd: float,
    buckets: dict[int, float],
) -> SimpleNamespace:
    competencia = "01/2026"
    bucket_labels = {
        1: "Até 30 dias",
        2: "31 a 60 dias",
        3: "61 a 90 dias",
        4: "91 a 120 dias",
        5: "121 a 150 dias",
        6: "151 a 180 dias",
        7: "181 a 360 dias",
        8: "361 a 720 dias",
        9: "721 a 1080 dias",
        10: "Acima de 1080 dias",
    }
    return SimpleNamespace(
        competencias=[competencia],
        fund_info={"nome_fundo": fund_name, "cnpj_fundo": cnpj, "nome_classe": "Classe única"},
        subordination_history_df=pd.DataFrame(
            [
                {
                    "competencia": competencia,
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "pl_total": pl_total,
                    "pl_senior": pl_senior,
                    "pl_mezzanino": pl_mezz,
                    "pl_subordinada_strict": pl_sub,
                    "pl_subordinada": pl_mezz + pl_sub,
                }
            ]
        ),
        default_history_df=pd.DataFrame(
            [
                {
                    "competencia": competencia,
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "direitos_creditorios": carteira,
                    "direitos_creditorios_fonte": "teste",
                    "provisao_total": pdd,
                }
            ]
        ),
        dc_canonical_history_df=pd.DataFrame(
            [
                {
                    "competencia": competencia,
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "dc_total_canonico": carteira,
                    "dc_total_fonte_efetiva": "teste",
                }
            ]
        ),
        default_buckets_history_df=pd.DataFrame(
            [
                {
                    "competencia": competencia,
                    "competencia_dt": pd.Timestamp("2026-01-01"),
                    "ordem": ordem,
                    "faixa": bucket_labels[ordem],
                    "valor": valor,
                    "source_status": "reported_value",
                }
                for ordem, valor in buckets.items()
            ]
        ),
    )


if __name__ == "__main__":
    unittest.main()
