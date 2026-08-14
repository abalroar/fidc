from io import BytesIO

from openpyxl import load_workbook
from pptx import Presentation

from scripts.build_deep_dive_package import get_portfolio_by_id
from services.fidc_analytical_slide import (
    build_fidc_analytical_pptx_bytes,
    build_fidc_analytical_xlsx_bytes,
    build_slide_data,
)


def _portfolio():
    portfolio = get_portfolio_by_id("meutudo_ago26_202608")
    assert portfolio is not None
    return portfolio


def test_meutudo_portfolio_has_expected_scope_and_class_substitution():
    portfolio = _portfolio()
    cnpjs = [fund.cnpj for fund in portfolio.funds]
    assert len(cnpjs) == 10
    assert "63953620000165" in cnpjs
    assert "63953619000130" not in cnpjs


def test_slide_data_reconciles_latest_snapshot_and_rdb_evidence():
    data = build_slide_data(_portfolio().funds)
    assert len(data) == 10
    assert {item.competence for item in data} == {"06/2026"}
    explicit = next(item for item in data if item.cnpj == "54464892000100")
    assert explicit.rdb_status == "Explícito"
    assert "RDB Parati" in explicit.rdb_finding
    for item in data:
        if item.senior_pct_pl is not None and item.subordinated_pct_pl is not None:
            assert abs(item.senior_pct_pl + item.subordinated_pct_pl - 1) < 0.01


def test_outputs_are_native_and_auditable():
    portfolio = _portfolio()
    pptx = build_fidc_analytical_pptx_bytes(portfolio.funds, portfolio_name=portfolio.name)
    prs = Presentation(BytesIO(pptx))
    assert len(prs.slides) >= 11
    assert sum(shape.has_table for slide in prs.slides for shape in slide.shapes) >= 11
    assert all(not (shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height) for slide in prs.slides for shape in slide.shapes)

    xlsx = build_fidc_analytical_xlsx_bytes(portfolio.funds, portfolio_name=portfolio.name)
    wb = load_workbook(BytesIO(xlsx), data_only=False)
    assert wb.sheetnames == ["Snapshot", "Emissões", "RDB", "Validações", "Metodologia"]
    assert len(wb["Snapshot"].tables) == 1
    assert len(wb["RDB"].tables) == 1
