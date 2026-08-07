"""A carteira de subordinação: registro, resolução e troca de slides."""

from __future__ import annotations

from pathlib import Path
import sys

import pandas as pd
import pytest


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from services.carteira_subordinacao import (  # noqa: E402
    ORIGEM_CARTEIRA_101,
    ORIGEM_MANUAL,
    RegistryError,
    dumbbell_figure,
    format_cnpj,
    load_registry,
    normalize_cnpj,
    remove_entry,
    resolve_portfolio,
    save_entry,
    seed_registry,
    short_fund_name,
)

FUNDO_A = "11111111000191"
FUNDO_B = "22222222000172"
FUNDO_MORTO = "33333333000153"
FUNDO_QUEBRADO = "55555555000115"
FUNDO_UNICO = "66666666000196"


@pytest.fixture()
def data_dir(tmp_path: Path) -> Path:
    """Uma base mínima com o que ``resolve_portfolio`` lê."""

    def linha(competencia, cnpj, nome, pl, sub, total=None):
        total = pl if total is None else total
        return {
            "competencia": competencia, "cnpj": cnpj, "denominacao": nome,
            "pl": pl, "subordinacao_pct": sub,
            "vl_cotas_total": total,
            "vl_cotas_subordinadas": (sub or 0.0) * total,
        }

    monthly = pd.DataFrame(
        [
            # O fundo A reportou até julho; o B parou em junho.  A resolução
            # tem de dar a cada um a sua própria competência mais recente.
            linha("2026-06", FUNDO_A, "FIDC A", 100e6, 0.20),
            linha("2026-07", FUNDO_A, "FIDC A", 120e6, 0.25),
            linha("2026-06", FUNDO_B, "FIDC B", 80e6, 0.08),
            # Sem PL a linha não conta como reporte.
            linha("2026-07", FUNDO_B, "FIDC B", 0.0, None),
            # Parou de reportar bem antes do corte.
            linha("2025-01", FUNDO_MORTO, "FIDC Morto", 10e6, 0.30),
            # O fundo QUEBRADO tem patrimônio em julho, mas o quadro de cotas
            # daquele mês veio zerado: junho é o dado bom.
            linha("2026-06", FUNDO_QUEBRADO, "FIDC Quebrado", 50e6, 0.30),
            linha("2026-07", FUNDO_QUEBRADO, "FIDC Quebrado", 48e6, 0.0, total=0.0),
            # O fundo CLASSE ÚNICA reporta o quadro de cotas todo mês, sempre
            # sem cota subordinada: 0% ali é o dado, não uma falha.
            linha("2026-05", FUNDO_UNICO, "FIDC Único", 30e6, 0.0),
            linha("2026-06", FUNDO_UNICO, "FIDC Único", 31e6, 0.0),
            linha("2026-07", FUNDO_UNICO, "FIDC Único", 32e6, 0.0),
        ]
    )
    monthly.to_csv(tmp_path / "vehicle_monthly.csv.gz", index=False, compression="gzip")
    pd.DataFrame(
        [{"cnpj_fundo": FUNDO_A, "cnpj_classe": FUNDO_A, "tipo_anbima": "Financeiro",
          "foco_anbima": "Consignado"},
         {"cnpj_fundo": FUNDO_B, "cnpj_classe": FUNDO_B, "tipo_anbima": "Outros",
          "foco_anbima": ""}]
    ).to_csv(tmp_path / "industry_anbima_classification.csv.gz", index=False, compression="gzip")
    return tmp_path


# ---------------------------------------------------------------------------
# Registro
# ---------------------------------------------------------------------------

def test_o_cnpj_entra_com_ou_sem_mascara(data_dir: Path) -> None:
    save_entry(cnpj="11.111.111/0001-91", subordinacao_minima_pct=15, data_dir=data_dir)

    registry = load_registry(data_dir)

    assert registry["cnpj"].tolist() == [FUNDO_A]
    assert format_cnpj(FUNDO_A) == "11.111.111/0001-91"


def test_um_cnpj_incompleto_e_recusado(data_dir: Path) -> None:
    with pytest.raises(RegistryError, match="14 dígitos"):
        save_entry(cnpj="111", subordinacao_minima_pct=15, data_dir=data_dir)


def test_salvar_sem_nenhum_minimo_e_recusado(data_dir: Path) -> None:
    """Sem mínimo não existe comparação — a linha não teria função."""

    with pytest.raises(RegistryError, match="ao menos um"):
        save_entry(cnpj=FUNDO_A, data_dir=data_dir)


def test_um_minimo_fora_da_escala_e_recusado(data_dir: Path) -> None:
    with pytest.raises(RegistryError, match="entre 0 e 100"):
        save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=150, data_dir=data_dir)


def test_o_estrutural_nao_pode_ser_menor_que_o_junior(data_dir: Path) -> None:
    """O estrutural soma o mezanino ao júnior; menor que o júnior é erro de digitação."""

    with pytest.raises(RegistryError, match="não pode ser menor"):
        save_entry(
            cnpj=FUNDO_A,
            subordinacao_minima_pct=20,
            subordinacao_estrutural_pct=10,
            data_dir=data_dir,
        )


def test_salvar_o_mesmo_cnpj_duas_vezes_substitui_a_linha(data_dir: Path) -> None:
    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=15, data_dir=data_dir)
    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=18, fonte="Regulamento novo",
               data_dir=data_dir)

    registry = load_registry(data_dir)

    assert len(registry) == 1
    assert registry.iloc[0]["subordinacao_minima_pct"] == 18
    assert registry.iloc[0]["fonte"] == "Regulamento novo"


def test_remover_tira_o_cnpj_do_registro(data_dir: Path) -> None:
    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=15, data_dir=data_dir)
    save_entry(cnpj=FUNDO_B, subordinacao_minima_pct=5, data_dir=data_dir)

    remove_entry(FUNDO_A, data_dir)

    assert load_registry(data_dir)["cnpj"].tolist() == [FUNDO_B]


def test_a_semeadura_preserva_o_que_o_analista_editou(tmp_path: Path) -> None:
    """A Carteira 101 e a inclusão manual moram no mesmo arquivo; semear de
    novo não pode apagar a decisão de quem editou à mão."""

    pd.DataFrame(
        [{"cnpj_fundo": FUNDO_A, "documento_id_regulamento": "123",
          "documento_data_regulamento": "2026-01-01", "pagina_clausula": "10",
          "subordinacao_minima_junior_pct": "12.5", "suporte_estrutural_minimo_pct": "",
          "status_curadoria_documental": "documentado"}]
    ).to_csv(tmp_path / "industry_carteira_1_document_curation.csv", index=False)
    pd.DataFrame([{"cnpj_fundo": FUNDO_A, "nome_foto": "FIDC A"}]).to_csv(
        tmp_path / "industry_carteira_1_scope.csv", index=False
    )
    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=99, origem=ORIGEM_MANUAL,
               data_dir=tmp_path)

    seeded = seed_registry(tmp_path)

    linha = seeded[seeded["cnpj"].eq(FUNDO_A)].iloc[0]
    assert linha["subordinacao_minima_pct"] == 99
    assert linha["origem"] == ORIGEM_MANUAL

    forcado = seed_registry(tmp_path, overwrite=True)

    assert forcado[forcado["cnpj"].eq(FUNDO_A)].iloc[0]["origem"] == ORIGEM_CARTEIRA_101


# ---------------------------------------------------------------------------
# Resolução
# ---------------------------------------------------------------------------

def test_cada_fundo_entra_com_a_sua_competencia_mais_recente(data_dir: Path) -> None:
    """Julho para um, junho para outro: a carteira mistura sem preferir um mês."""

    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=15, data_dir=data_dir)
    save_entry(cnpj=FUNDO_B, subordinacao_minima_pct=10, data_dir=data_dir)

    frame = resolve_portfolio(data_dir).frame.set_index("cnpj")

    assert frame.at[FUNDO_A, "competencia"] == "2026-07"
    assert frame.at[FUNDO_B, "competencia"] == "2026-06"
    # E o valor é o da competência escolhida, não o da anterior.
    assert frame.at[FUNDO_A, "sub_atual_pct"] == pytest.approx(25.0)
    assert frame.at[FUNDO_A, "pl_mm"] == pytest.approx(120.0)


def test_a_fracao_do_informe_vira_ponto_percentual(data_dir: Path) -> None:
    """O Informe reporta 0,08; o mínimo do registro está em 10 p.p."""

    save_entry(cnpj=FUNDO_B, subordinacao_minima_pct=10, data_dir=data_dir)

    linha = resolve_portfolio(data_dir).frame.iloc[0]

    assert linha["sub_atual_pct"] == pytest.approx(8.0)
    assert linha["folga_pp"] == pytest.approx(-2.0)
    assert bool(linha["abaixo_do_minimo"]) is True


def test_o_fundo_que_parou_de_reportar_sai_do_grafico(data_dir: Path) -> None:
    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=15, data_dir=data_dir)
    save_entry(cnpj=FUNDO_MORTO, subordinacao_minima_pct=15, data_dir=data_dir)

    ativos = resolve_portfolio(data_dir).frame
    todos = resolve_portfolio(data_dir, somente_ativos=False).frame

    assert ativos["cnpj"].tolist() == [FUNDO_A]
    assert set(todos["cnpj"]) == {FUNDO_A, FUNDO_MORTO}


def test_o_estrutural_manda_quando_existe(data_dir: Path) -> None:
    """Havendo mezanino, a comparação é contra Sub+Mez, e o slide diz isso."""

    save_entry(
        cnpj=FUNDO_A,
        subordinacao_minima_pct=10,
        subordinacao_estrutural_pct=22,
        data_dir=data_dir,
    )

    linha = resolve_portfolio(data_dir).frame.iloc[0]

    assert linha["referencia_pct"] == pytest.approx(22.0)
    assert linha["referencia_tipo"] == "Estrutural (Sub+Mez)"
    assert linha["folga_pp"] == pytest.approx(3.0)


def test_sem_subordinacao_no_informe_o_fundo_nao_e_comparavel(data_dir: Path) -> None:
    save_entry(cnpj="44444444000134", subordinacao_minima_pct=10, data_dir=data_dir)

    frame = resolve_portfolio(data_dir, somente_ativos=False).frame
    linha = frame[frame["cnpj"].eq("44444444000134")].iloc[0]

    assert bool(linha["comparavel"]) is False


# ---------------------------------------------------------------------------
# Gráfico
# ---------------------------------------------------------------------------

def test_o_verde_marca_sempre_o_maior_dos_dois(data_dir: Path) -> None:
    """O fundo A está acima do mínimo e o B abaixo; o ponto verde troca de lado.

    É essa troca que carrega o alerta: o verde no lugar do mínimo significa que
    o fundo não alcança o que o regulamento exige.
    """

    import numpy as np

    from services.carteira_subordinacao import COLOR_HIGHER

    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=15, data_dir=data_dir)
    save_entry(cnpj=FUNDO_B, subordinacao_minima_pct=10, data_dir=data_dir)

    frame = resolve_portfolio(data_dir).frame
    figure = dumbbell_figure(frame)
    axes = figure.axes[0]

    ordenado = frame[frame["comparavel"]].sort_values("sub_atual_pct", ascending=False)
    esperado = np.maximum(
        ordenado["sub_atual_pct"].to_numpy(), ordenado["referencia_pct"].to_numpy()
    )
    verdes = [
        colecao
        for colecao in axes.collections
        if getattr(colecao, "get_facecolor", None) is not None
        and len(colecao.get_facecolor())
        and matplotlib_hex(colecao.get_facecolor()[0]) == COLOR_HIGHER.lower()
    ]
    assert len(verdes) == 1
    alturas = verdes[0].get_offsets()[:, 1]
    assert np.allclose(sorted(alturas), sorted(esperado))


def matplotlib_hex(rgba) -> str:
    from matplotlib.colors import to_hex

    return to_hex(rgba)


def test_a_haste_de_quem_esta_abaixo_do_minimo_e_destacada(data_dir: Path) -> None:
    from services.carteira_subordinacao import COLOR_GAP

    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=15, data_dir=data_dir)
    save_entry(cnpj=FUNDO_B, subordinacao_minima_pct=10, data_dir=data_dir)

    axes = dumbbell_figure(resolve_portfolio(data_dir).frame).axes[0]

    vermelhas = [
        colecao
        for colecao in axes.collections
        if getattr(colecao, "get_color", None) is not None
        and len(getattr(colecao, "get_color")())
        and matplotlib_hex(colecao.get_color()[0]) == COLOR_GAP.lower()
    ]
    assert vermelhas, "a haste do fundo em falta precisa aparecer em vermelho"


def test_o_grafico_vazio_nao_quebra(data_dir: Path) -> None:
    figure = dumbbell_figure(resolve_portfolio(data_dir).frame)

    assert figure.axes


def test_o_nome_do_fundo_perde_a_boilerplate_registral() -> None:
    """Corta-se **na** frase registral; subtraí-la deixaria restos sem sentido."""

    assert short_fund_name(
        "ATLANTA FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS DE RESPONSABILIDADE LIMITADA"
    ) == "Atlanta"
    assert short_fund_name("CLASSE ÚNICA DO FIDC CRÉDITO NITRO AGRO") == "Crédito Nitro Agro"
    # Subtrair a boilerplate produziria "Cobuccio de Fechada de Re…".
    assert short_fund_name(
        "COBUCCIO FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS DE CLASSE ÚNICA "
        "FECHADA DE RESPONSABILIDADE LIMITADA"
    ) == "Cobuccio"
    # O nome próprio pode vir depois da frase registral.
    assert short_fund_name(
        "FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS PAGSEGURO I"
    ) == "Pagseguro I"
    assert normalize_cnpj("11.111.111/0001-91") == FUNDO_A


def test_a_caixa_do_rotulo_separa_sigla_de_palavra() -> None:
    """BLUE, CASH e CLUB são palavras; VTK e II são siglas."""

    assert short_fund_name("RESIDENCE CLUB FIDC") == "Residence Club"
    assert short_fund_name("CD CASH FIDC RESPONSABILIDADE LIMITADA") == "CD Cash"
    assert short_fund_name("VTK FIDC RESPONSABILIDADE LIMITADA") == "VTK"
    assert short_fund_name("VIA INVEST II FIDC") == "Via Invest II"
    assert short_fund_name("TMAQ 21 FIDC DE ARRANJOS DE PAGAMENTO") == "Tmaq 21"


# ---------------------------------------------------------------------------
# Deck
# ---------------------------------------------------------------------------

def test_a_troca_de_slides_nao_duplica_partes_no_pacote() -> None:
    """Trocar slides removendo e recriando emitiria dois ``slideN.xml``.

    O ``next_partname`` do python-pptx devolve um nome já ocupado quando a
    numeração das partes tem buraco — que é o que remover um slide do meio do
    deck produz.  Por isso a troca reescreve no lugar, e este teste guarda
    exatamente essa propriedade: nenhum nome de parte repetido depois de
    reescrever a carteira **e** acrescentar o ranking.
    """

    import warnings
    from io import BytesIO

    from pptx import Presentation

    from services.anbima_executive_export import append_anbima_slides
    from services.carteira_deck import REPLACED_SLIDE_RANGE, replace_structural_slides
    from services.industry_ppt_export import build_industry_pptx_bytes

    data_dir = ROOT / "data" / "industry_study"
    presentation = Presentation(BytesIO(build_industry_pptx_bytes(data_dir)))
    antes = len(presentation.slides)

    escritos = replace_structural_slides(presentation, data_dir)
    append_anbima_slides(presentation, data_dir)

    # Sete gráficos — o consolidado mais as seis categorias estruturais — para
    # seis slots; o excedente é acrescentado, nunca conquistado por remoção.
    assert escritos >= REPLACED_SLIDE_RANGE[1] - REPLACED_SLIDE_RANGE[0] + 1
    with warnings.catch_warnings(record=True) as capturados:
        warnings.simplefilter("always")
        buffer = BytesIO()
        presentation.save(buffer)
    duplicadas = [w for w in capturados if "Duplicate name" in str(w.message)]
    assert duplicadas == []

    reaberto = Presentation(BytesIO(buffer.getvalue()))
    assert len(reaberto.slides) > antes
    primeiro, _ = REPLACED_SLIDE_RANGE
    for numero in range(primeiro, primeiro + escritos):
        slide = list(reaberto.slides)[numero - 1]
        textos = " ".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        )
        assert "CARTEIRA 101" in textos
        assert any(shape.shape_type == 13 for shape in slide.shapes)


def test_a_aba_carteira_esta_no_menu_da_industria() -> None:
    from tabs.tab_industry_study import INDUSTRY_VIEW_TABS

    assert "Carteira" in INDUSTRY_VIEW_TABS


def test_o_registro_participa_da_chave_de_cache_das_exportacoes(tmp_path: Path) -> None:
    """Salvar um fundo tem de invalidar o deck; sem isso o site serve o antigo."""

    import tabs.tab_industry_study as painel
    from services.carteira_subordinacao import registry_path

    antes = painel._industry_export_signature()
    caminho = registry_path(painel._DATA_DIR)
    original = caminho.read_bytes()
    marca = caminho.stat().st_mtime
    try:
        caminho.write_bytes(original + b"\n99999999000199,,10,,False,manual,,,,\n")
        depois = painel._industry_export_signature()
    finally:
        caminho.write_bytes(original)
        import os

        os.utime(caminho, (marca, marca))

    assert depois != antes


# ---------------------------------------------------------------------------
# Integridade do quadro de cotas
# ---------------------------------------------------------------------------

def test_o_mes_com_quadro_de_cotas_zerado_cede_lugar_ao_anterior(data_dir: Path) -> None:
    """Patrimônio positivo com total de cotas zerado não é fundo sem subordinação.

    É a quebra do quadro de cotas naquele mês, e ela produziria 0% onde o mês
    anterior traz o número real.
    """

    save_entry(cnpj=FUNDO_QUEBRADO, subordinacao_minima_pct=10, data_dir=data_dir)

    linha = resolve_portfolio(data_dir).frame.iloc[0]

    assert linha["competencia"] == "2026-06"
    assert linha["sub_atual_pct"] == pytest.approx(30.0)
    assert bool(linha["quadro_de_cotas_integro"]) is True


def test_zero_por_cento_reportado_todo_mes_permanece_zero(data_dir: Path) -> None:
    """Um fundo de classe única reporta 0% de verdade; não se corrige isso.

    A distinção entre o 0% real e o 0% de falha está em ``meses_sem_subordinada``:
    um mês isolado é ruído, uma sequência descreve a estrutura do fundo.
    """

    save_entry(cnpj=FUNDO_UNICO, subordinacao_minima_pct=10, data_dir=data_dir)

    linha = resolve_portfolio(data_dir).frame.iloc[0]

    assert linha["competencia"] == "2026-07"
    assert linha["sub_atual_pct"] == pytest.approx(0.0)
    assert bool(linha["quadro_de_cotas_integro"]) is True
    assert linha["meses_sem_subordinada"] == 3
    assert bool(linha["abaixo_do_minimo"]) is True


# ---------------------------------------------------------------------------
# Taxonomia dos slides estruturais
# ---------------------------------------------------------------------------

def test_a_carteira_usa_a_taxonomia_dos_slides_estruturais() -> None:
    """O corte é o dos slides 18–23, e não o tipo ANBIMA."""

    from services.carteira_deck import slide_plans
    from services.carteira_subordinacao import resolve_portfolio as resolver

    esperadas = {
        "Financeiro",
        "Adquirência",
        "Agro / Revenda",
        "Risco Corporativo",
        "Consignado INSS e FGTS",
        "Factoring",
    }
    data_dir = ROOT / "data" / "industry_study"
    frame = resolver(data_dir).frame
    presentes = set(frame.loc[frame["comparavel"], "categoria_estrutural"])

    assert presentes.issubset(esperadas | {"Não classificado"})
    assert esperadas.issubset(presentes)

    categorias = {str(plano["categoria"]) for plano in slide_plans(data_dir)}
    assert esperadas == categorias


def test_um_cnpj_fora_da_taxonomia_nao_inventa_categoria(data_dir: Path) -> None:
    from services.carteira_subordinacao import NAO_CLASSIFICADO

    save_entry(cnpj=FUNDO_A, subordinacao_minima_pct=15, data_dir=data_dir)

    assert resolve_portfolio(data_dir).frame.iloc[0]["categoria_estrutural"] == (
        NAO_CLASSIFICADO
    )


# ---------------------------------------------------------------------------
# Tabela ao lado do gráfico
# ---------------------------------------------------------------------------

def test_a_ordenacao_da_tabela_privilegia_materialidade_e_depois_risco() -> None:
    """PL do maior para o menor; a folga desempata pelo mais perto do limite."""

    import pandas as pd

    from services.carteira_deck import order_for_table

    bruto = pd.DataFrame(
        [
            {"cnpj": "a", "pl_mm": 100.0, "folga_pp": 20.0},
            {"cnpj": "b", "pl_mm": 500.0, "folga_pp": 9.0},
            {"cnpj": "c", "pl_mm": 500.0, "folga_pp": 1.0},
            {"cnpj": "d", "pl_mm": None, "folga_pp": -3.0},
        ]
    )

    assert order_for_table(bruto)["cnpj"].tolist() == ["c", "b", "a", "d"]


def test_a_folga_e_atual_menos_minimo_com_uma_casa() -> None:
    import pandas as pd

    from services.carteira_deck import table_rows

    from services.carteira_deck import TABLE_HEADER

    frame = pd.DataFrame(
        [
            {"cnpj": "a", "fundo": "FIDC ALFA", "pl_mm": 1234.5, "multi_flag": "",
             "referencia_pct": 15.0, "sub_atual_pct": 21.34, "folga_pp": 6.34},
            {"cnpj": "b", "fundo": "FIDC BETA", "pl_mm": 5.0, "multi_flag": "",
             "referencia_pct": 10.0, "sub_atual_pct": 7.5, "folga_pp": -2.5},
        ]
    )
    rows, _fills, _colors = table_rows(frame)

    assert rows[0] == list(TABLE_HEADER)
    assert rows[1][1:] == ["1.234,5", "15,0", "21,3", "+6,3"]
    assert rows[2][1:] == ["5,0", "10,0", "7,5", "-2,5"]


def test_a_folga_ganha_cor_por_banda_de_risco() -> None:
    """Vermelho só no desenquadramento; amarelo até +1,5; verde acima disso."""

    import pandas as pd

    from services.carteira_deck import (
        COL_FOLGA,
        FILL_AMARELO,
        FILL_VERDE,
        FILL_VERMELHO,
        table_rows,
    )

    frame = pd.DataFrame(
        [
            {"cnpj": "a", "fundo": "A", "pl_mm": 5.0, "multi_flag": "",
             "referencia_pct": 10.0, "sub_atual_pct": 30.0, "folga_pp": 20.0},
            {"cnpj": "b", "fundo": "B", "pl_mm": 4.0, "multi_flag": "",
             "referencia_pct": 10.0, "sub_atual_pct": 11.6, "folga_pp": 1.6},
            {"cnpj": "c", "fundo": "C", "pl_mm": 3.0, "multi_flag": "",
             "referencia_pct": 10.0, "sub_atual_pct": 11.5, "folga_pp": 1.5},
            {"cnpj": "d", "fundo": "D", "pl_mm": 2.0, "multi_flag": "",
             "referencia_pct": 10.0, "sub_atual_pct": 10.0, "folga_pp": 0.0},
            {"cnpj": "e", "fundo": "E", "pl_mm": 1.0, "multi_flag": "",
             "referencia_pct": 10.0, "sub_atual_pct": 9.0, "folga_pp": -1.0},
        ]
    )
    _rows, fills, _colors = table_rows(frame)

    assert fills[(1, COL_FOLGA)] == FILL_VERDE
    assert fills[(2, COL_FOLGA)] == FILL_VERDE
    # A borda de +1,5 pertence ao amarelo, e +0,0 também.
    assert fills[(3, COL_FOLGA)] == FILL_AMARELO
    assert fills[(4, COL_FOLGA)] == FILL_AMARELO
    assert fills[(5, COL_FOLGA)] == FILL_VERMELHO


def test_todo_fidc_do_grafico_aparece_na_tabela_do_slide() -> None:
    """A razão de a tabela existir: o gráfico rotula poucos, ela nomeia todos."""

    from services.carteira_deck import slide_plans, table_rows
    from services.carteira_subordinacao import short_fund_name

    data_dir = ROOT / "data" / "industry_study"
    por_categoria: dict[str, set[str]] = {}
    graficos: dict[str, set[str]] = {}
    for plano in slide_plans(data_dir):
        categoria = str(plano["categoria"])
        rows, _f, _c = table_rows(plano["tabela"])
        por_categoria.setdefault(categoria, set()).update(linha[0] for linha in rows[1:])
        graficos.setdefault(categoria, set()).update(
            short_fund_name(str(nome), limite=26) for nome in plano["grafico"]["fundo"]
        )

    for categoria, nomes in graficos.items():
        # O rótulo da tabela pode carregar a marca de multicedente; o nome
        # curto tem de aparecer nele.
        for nome in nomes:
            assert any(nome in linha for linha in por_categoria[categoria]), (
                categoria,
                nome,
            )


def test_o_grafico_nomeia_todos_os_veiculos_no_eixo() -> None:
    """Um rótulo por ponto, na vertical: nenhum FIDC fica anônimo no gráfico."""

    from services.carteira_deck import slide_plans
    from services.carteira_subordinacao import dumbbell_figure

    plano = slide_plans(ROOT / "data" / "industry_study")[0]
    grafico = plano["grafico"]
    axes = dumbbell_figure(grafico, nomear_todos=True, figsize=(12.05, 3.1)).axes[0]

    marcas = [marca.get_text() for marca in axes.get_xticklabels()]
    assert len(marcas) == len(grafico)
    assert all(texto.strip() for texto in marcas)


# ---------------------------------------------------------------------------
# Top 100 para revisão do universo Middle
# ---------------------------------------------------------------------------

def test_o_top100_entra_no_deck_em_tabelas_nativas_ordenadas_por_pl() -> None:
    from io import BytesIO

    from pptx import Presentation

    from services.bba_deck import SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN
    from services.top100_middle_deck import (
        COLUMNS,
        ROWS_PER_SLIDE,
        append_top100_slides,
        load_review,
    )

    data_dir = ROOT / "data" / "industry_study"
    review = load_review(data_dir)
    presentation = Presentation()
    presentation.slide_width = int(SLIDE_WIDTH_IN * 914400)
    presentation.slide_height = int(SLIDE_HEIGHT_IN * 914400)

    blocos = append_top100_slides(presentation, data_dir)

    assert blocos == -(-len(review) // ROWS_PER_SLIDE)
    buffer = BytesIO()
    presentation.save(buffer)
    reaberto = Presentation(BytesIO(buffer.getvalue()))

    total = 0
    for slide in reaberto.slides:
        tabelas = [shape for shape in slide.shapes if shape.has_table]
        assert len(tabelas) == 1
        tabela = tabelas[0].table
        assert len(tabela.columns) == len(COLUMNS)
        total += len(tabela.rows) - 1
    assert total == len(review)

    # A ordenação por PL é o que dá materialidade à revisão.
    valores = review["pl_num"].dropna().tolist()
    assert valores == sorted(valores, reverse=True)


def test_a_planilha_do_top100_e_uma_tabela_de_excel_filtravel(tmp_path: Path) -> None:
    """Tabela nativa, não intervalo formatado: filtro e ordenação vêm com ela."""

    from openpyxl import load_workbook

    sys.path.insert(0, str(ROOT / "scripts"))
    from build_top100_middle_review_xlsx import write_workbook  # noqa: E402

    from services.top100_middle_deck import load_review

    destino = tmp_path / "top100.xlsx"
    write_workbook(load_review(ROOT / "data" / "industry_study"), destino)
    workbook = load_workbook(destino)
    sheet = workbook["Top 100"]

    assert list(sheet.tables) == ["Top100Middle"]
    assert sheet.freeze_panes == "C2"
    cabecalho = [cell.value for cell in sheet[1]]
    assert cabecalho[0] == "#" and cabecalho[-1] == "MIDDLE"
    # A coluna de revisão só aceita Sim ou Não.
    assert len(sheet.data_validations.dataValidation) == 1
    assert sheet.data_validations.dataValidation[0].formula1 == '"Sim,Não"'


def test_as_tabelas_do_deck_cabem_na_largura_do_conteudo() -> None:
    """A soma das larguras precisa fechar na margem; passando, a última coluna
    sai pela borda da lâmina — e no Top 100 a última é a que o revisor preenche."""

    from services.bba_deck import CONTENT_WIDTH_IN, MARGIN_IN, SLIDE_WIDTH_IN
    from services.top100_middle_deck import COLUMNS

    largura_top100 = sum(largura for _k, _t, largura, _a in COLUMNS)

    assert largura_top100 <= CONTENT_WIDTH_IN
    assert MARGIN_IN + largura_top100 <= SLIDE_WIDTH_IN - MARGIN_IN


def test_a_tabela_da_carteira_termina_na_margem_direita() -> None:
    from io import BytesIO

    from pptx import Presentation
    from pptx.util import Emu

    from services.bba_deck import MARGIN_IN, SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN
    from services.carteira_deck import draw_carteira_slide, slide_plans

    presentation = Presentation()
    presentation.slide_width = int(SLIDE_WIDTH_IN * 914400)
    presentation.slide_height = int(SLIDE_HEIGHT_IN * 914400)
    from services.bba_deck import Deck

    deck = Deck("teste", presentation)
    plano = slide_plans(ROOT / "data" / "industry_study")[0]
    slide = deck.blank()
    draw_carteira_slide(deck, slide, plano)

    buffer = BytesIO()
    presentation.save(buffer)
    reaberto = list(Presentation(BytesIO(buffer.getvalue())).slides)[0]
    imagens = [shape for shape in reaberto.shapes if shape.shape_type == 13]
    tabelas = sorted(
        (shape for shape in reaberto.shapes if shape.has_table), key=lambda s: s.left
    )

    assert imagens and len(tabelas) == 2
    # O gráfico toma a largura inteira do conteúdo e a tabela vem abaixo dele,
    # em duas tranches que terminam na margem direita.
    assert Emu(imagens[0].width).inches == pytest.approx(12.05, abs=0.01)
    assert Emu(tabelas[0].top).inches > Emu(imagens[0].top + imagens[0].height).inches
    assert Emu(tabelas[1].left).inches > Emu(
        tabelas[0].left + tabelas[0].width
    ).inches
    direita = Emu(tabelas[1].left + tabelas[1].width).inches
    assert direita <= SLIDE_WIDTH_IN - MARGIN_IN + 0.01
