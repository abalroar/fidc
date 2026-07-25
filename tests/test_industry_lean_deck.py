"""Contrato do recorte enxuto do deck executivo."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from services.industry_lean_deck import (
    LEAN_SLIDE_HEADINGS,
    LeanDeckUnavailable,
    build_lean_pptx_bytes,
)


ROOT = Path(__file__).resolve().parents[1]
PUBLISHED_PPTX = (
    ROOT / "data" / "industry_study" / "generated_revision" / "industry_executive_revised.pptx"
)


def _published_bytes() -> bytes:
    if not PUBLISHED_PPTX.exists():
        pytest.skip("bundle publicado indisponível")
    return PUBLISHED_PPTX.read_bytes()


def _headings(pptx_bytes: bytes) -> list[str]:
    presentation = Presentation(BytesIO(pptx_bytes))
    headings: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                headings.append(shape.text_frame.text.strip().splitlines()[0].strip())
                break
    return headings


def _synthetic_deck(headings: list[str]) -> bytes:
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for heading in headings:
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Emu(457200), Emu(457200), Emu(4572000), Emu(457200))
        box.text_frame.text = heading
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_lean_deck_keeps_only_requested_headings_in_deck_order() -> None:
    source = _synthetic_deck(["CAPA", "MEIO", "FIM", "EXTRA"])

    lean = build_lean_pptx_bytes(source, headings=("FIM", "CAPA"))

    assert _headings(lean) == ["CAPA", "FIM"]


def test_lean_deck_fails_closed_when_a_heading_is_missing() -> None:
    source = _synthetic_deck(["CAPA", "MEIO"])

    with pytest.raises(LeanDeckUnavailable) as excinfo:
        build_lean_pptx_bytes(source, headings=("CAPA", "AUSENTE"))

    assert "AUSENTE" in str(excinfo.value)


def test_lean_deck_matches_headings_ignoring_accents_and_separators() -> None:
    source = _synthetic_deck(["OFERTAS ENCERRADAS · RENDA FIXA"])

    lean = build_lean_pptx_bytes(source, headings=("Ofertas encerradas — Renda Fixa",))

    assert len(_headings(lean)) == 1


def test_lean_selection_covers_the_dcm_narrative() -> None:
    assert LEAN_SLIDE_HEADINGS[0] == "INDÚSTRIA DE FIDCs"
    assert LEAN_SLIDE_HEADINGS[-1] == "CONCLUSÕES ESTRATÉGICAS · DCM"
    assert "PRINCIPAIS CONCLUSÕES" in LEAN_SLIDE_HEADINGS
    assert len(set(LEAN_SLIDE_HEADINGS)) == len(LEAN_SLIDE_HEADINGS)


def test_lean_deck_from_published_bundle_preserves_native_charts() -> None:
    published = _published_bytes()
    available = [
        heading
        for heading in LEAN_SLIDE_HEADINGS
        if heading in set(_headings(published))
    ]

    lean = build_lean_pptx_bytes(published, headings=tuple(available))

    presentation = Presentation(BytesIO(lean))
    assert len(presentation.slides) == len(available)
    charts = sum(
        1
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_chart
    )
    assert charts > 0
    assert len(lean) < len(published)
