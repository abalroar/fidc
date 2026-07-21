from __future__ import annotations

from io import BytesIO
from pathlib import Path
from types import SimpleNamespace
from xml.etree import ElementTree as ET
from zipfile import ZipFile

import pandas as pd

from scripts.classify_large_fidcs_from_documents import (
    classify_large_fund_precise,
    extract_anbima_labels,
)
from services.industry_intelligence import (
    build_competence_status,
    build_competitive_position,
    build_offer_annual,
)
from services.industry_ppt_export import (
    _annual_history,
    build_industry_pptx_bytes,
    build_industry_xlsx_bytes,
)
from services.industry_revision_export import (
    EXPECTED_SLIDES,
    get_revision_export_status,
    validate_revision_pptx,
    validate_revision_xlsx,
)


CHART = "http://schemas.openxmlformats.org/drawingml/2006/chart"
SHEET = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"


def test_annual_history_keeps_full_2015_window_and_uses_latest_complete_month() -> None:
    industry = pd.DataFrame(
        {
            "competencia": [
                "2014-12",
                "2015-11",
                "2015-12",
                "2016-12",
                "2019-12",
                "2025-12",
                "2026-04",
                "2026-05",
                "2026-06",
            ],
            "pl_total": [80, 90, 100, 120, 180, 300, 310, 330, 120],
            "pl_fic_fidc": [8, 9, 10, 12, 18, 30, 31, 33, 12],
        }
    )
    pack = SimpleNamespace(
        competences=SimpleNamespace(latest_complete="2026-05")
    )

    history = _annual_history(industry, pack)

    assert history["competencia"].tolist() == [
        "2015-12",
        "2016-12",
        "2019-12",
        "2025-12",
        "2026-05",
    ]
    assert history["period_label"].tolist()[-1] == "Mai/26"
    assert history.iloc[0]["pl_ex_fic"] == 90
    assert history.iloc[-1]["pl_ex_fic"] == 297


def test_competence_status_keeps_partial_tail_out_of_consolidated_snapshot() -> None:
    industry = pd.DataFrame(
        {
            "competencia": ["2026-04", "2026-05", "2026-06"],
            "n_veiculos": [4_200, 4_224, 1_990],
            "pl_total": [950e9, 958e9, 337e9],
        }
    )

    status = build_competence_status(industry)

    assert status.iloc[-1]["publication_status"] == "preliminar"
    assert status.iloc[-2]["publication_status"] == "completa"
    assert status.iloc[-1]["vehicle_ratio_vs_previous"] < 0.5


def test_offer_annual_separates_total_initial_and_relevant_ticket_scopes() -> None:
    offers = pd.DataFrame(
        {
            "year": [2026, 2026, 2026],
            "period": ["2026YTD"] * 3,
            "offer_id": ["a", "b", "c"],
            "issuer_cnpj": ["1", "1", "2"],
            "registered_volume_brl": [500e6, 200e6, 400e6],
            "valid_offer": [True, True, False],
            "closed_offer": [True, True, False],
            "initial_offer": [True, False, True],
            "ticket_relevant": [True, False, True],
            "investor_data_available": [True, True, False],
            "single_investor": [True, False, False],
            "investor_count": [1, 4, 0],
            "placed_volume_proxy_brl": [500e6, 180e6, 0],
        }
    )

    annual = build_offer_annual(offers).iloc[0]

    assert annual["valid_offers"] == 2
    assert annual["initial_offers"] == 1
    assert annual["relevant_ticket_offers"] == 1
    assert annual["valid_registered_volume_brl"] == 700e6


def test_competitive_position_exposes_structuring_to_servicing_conversion_gap() -> None:
    offers = pd.DataFrame(
        {
            "year": [2026, 2026, 2026],
            "period": ["2026YTD"] * 3,
            "offer_id": ["a", "b", "c"],
            "registered_volume_brl": [600e6, 500e6, 400e6],
            "valid_offer": [True] * 3,
            "ticket_relevant": [True] * 3,
            "leader_group": ["Itaú", "Itaú", "Bradesco"],
            "administrator_group": ["Oliveira Trust", "Itaú", "Bradesco"],
            "manager_group": ["Oliveira Trust", "Tercon", "Bradesco"],
            "custodian_group": ["Oliveira Trust", "Itaú", "Bradesco"],
            "same_platform_admin_manager_custodian": [True, False, True],
            "closed_offer": [True, True, True],
            "investor_data_available": [True, True, True],
            "single_investor": [True, False, False],
            "investor_count": [1, 8, 5],
            "fund_investor_present": [False, True, True],
            "bank_investor_present": [True, True, False],
        }
    )

    row = build_competitive_position(offers).iloc[0]

    assert row["itau_coordinator_rank"] == 1
    assert row["itau_coordinator_volume_brl"] == 1.1e9
    assert row["itau_administrator_volume_brl"] == 500e6
    assert row["itau_coordinator_share"] > row["itau_administrator_share"]


def test_large_fund_classifier_prioritizes_specific_lastro_over_generic_terms() -> None:
    text = (
        "Índice de preços ao consumidor. Os direitos creditórios são oriundos de litígios "
        "contra pessoas jurídicas de direito público e podem estar representados por precatórios."
    )
    result = classify_large_fund_precise("Alternative Assets III", text, "Financeiro")

    assert result["n1"] == "Judicial/Precatórios/NPL"
    assert result["confidence"] == "alta"


def test_anbima_labels_remain_separate_from_economic_classification() -> None:
    text = (
        'O fundo é classificado como FIDC, tipo "Agro, Indústria e Comércio" e '
        'foco de atuação em "Crédito Corporativo", conforme diretriz ANBIMA.'
    )

    anbima_type, focus, evidence = extract_anbima_labels(text)

    assert "AGRO" in anbima_type
    assert "CREDITO CORPORATIVO" in focus
    assert evidence


def test_industry_exports_are_valid_office_files() -> None:
    from pptx import Presentation

    status = get_revision_export_status()
    assert status.bundle_valid, status.validation_error
    assert status.bundle_id
    assert Path(status.pptx_path).name == "industry_executive_revised.pptx"
    assert Path(status.xlsx_path).name == "industry_data_revised.xlsx"

    pptx = build_industry_pptx_bytes()
    xlsx = build_industry_xlsx_bytes()
    validate_revision_pptx(pptx)
    validate_revision_xlsx(xlsx)

    presentation = Presentation(BytesIO(pptx))
    assert len(presentation.slides) == EXPECTED_SLIDES == 55
    slide_texts: list[str] = []
    for slide in presentation.slides:
        visible_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                visible_parts.append(shape.text)
            if getattr(shape, "has_table", False):
                visible_parts.extend(cell.text for row in shape.table.rows for cell in row.cells)
        slide_texts.append("\n".join(visible_parts))
    visible_text = "\n".join(slide_texts)
    for expected in (
        "INDÚSTRIA DE FIDCs",
        "SÍNTESE EXECUTIVA",
        "ESCALA DA INDÚSTRIA",
        "BASE INVESTIDORA",
        "OBSERVABILIDADE DA INADIMPLÊNCIA",
        "PRESTADORES · RANKING E CONCENTRAÇÃO",
        "MARKET SHARE · ADMINISTRAÇÃO",
        "MARKET SHARE · GESTÃO",
        "MARKET SHARE · CUSTÓDIA",
        "PRESTADORES · EVOLUÇÃO DO RANKING",
        "PRESTADORES INDEPENDENTES · EVOLUÇÃO",
        "FIDCs DOS CINCO BANCOS · COORTE ATUAL",
        "PRESTADORES · LIDERANÇA EXPLICADA",
        "CBSF / REAG · DESTINO DOS FUNDOS",
        "PRESTADORES · MIGRAÇÃO EM ADMINISTRAÇÃO",
        "PRESTADORES · MIGRAÇÃO EM GESTÃO",
        "PRESTADORES · MIGRAÇÃO EM CUSTÓDIA",
        "INADIMPLÊNCIA · COORTE ATUAL POR RECEBÍVEL",
        "OFERTAS ENCERRADAS · DISTRIBUIÇÃO DO TICKET",
        "OFERTAS ENCERRADAS · ORIGINADORES NOMINÁVEIS",
        "PRINCIPAIS CONCLUSÕES",
            "RANKING · TOP 20 FIDCs",
        "RANKING · TOP 20 OUTROS",
        "APÊNDICE · CURADORIA TOP 20",
    ):
        assert expected in visible_text
    profiles = [
        slide_text
        for slide_text in slide_texts
        if "APÊNDICE · CURADORIA TOP 20" in slide_text
    ]
    assert len(profiles) == 20
    for rank, slide_text in enumerate(profiles, start=1):
        assert f"#{rank} " in slide_text

    with ZipFile(BytesIO(pptx)) as archive:
        office_xml = b"".join(
            archive.read(name).upper()
            for name in archive.namelist()
            if name.endswith(".xml")
            and (
                name.startswith("ppt/slides/")
                or name.startswith("ppt/theme/")
                or "/charts/chart" in name
            )
        )
        assert b"EC7000" in office_xml
        assert b"172A3A" not in office_xml

        chart_names = [
            name
            for name in archive.namelist()
            if "/charts/chart" in name and name.endswith(".xml")
        ]
        assert chart_names
        chart_xml = b"".join(archive.read(name) for name in chart_names)
        assert b'<c:axId val="-' not in chart_xml
        assert b'<c:crossAx val="-' not in chart_xml
        for name in chart_names:
            root = ET.fromstring(archive.read(name))
            for smooth in root.iter(f"{{{CHART}}}smooth"):
                assert smooth.attrib.get("val", "0").lower() not in {"1", "true"}
            for marker in root.iter(f"{{{CHART}}}marker"):
                symbol = marker.find(f"{{{CHART}}}symbol")
                assert symbol is not None
                assert symbol.attrib.get("val") == "none"

    with ZipFile(BytesIO(xlsx)) as archive:
        workbook = ET.fromstring(archive.read("xl/workbook.xml"))
        sheet_names = {
            sheet.attrib["name"]
            for sheet in workbook.findall(f".//{{{SHEET}}}sheet")
        }
    assert {
        "QA Inadimplência",
        "Base competência-CNPJ",
        "Base por fundo-CNPJ",
        "Concentração de monoestruturas",
        "Market share por subtipo",
        "Top 20 FIDCs",
        "Top 20 Outros",
        "Curadoria Top 20",
        "Ranking prestadores",
        "Taxonomia adquirência",
        "Atribuição prestadores",
        "Fluxos prestadores",
        "Migração CBSF",
        "Checks revisão",
    }.issubset(sheet_names)
