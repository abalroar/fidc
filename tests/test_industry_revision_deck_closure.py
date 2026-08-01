from __future__ import annotations

from collections import Counter
import json
import os
from pathlib import Path
import re
import zipfile

import pandas as pd
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
REVISION_DIR = ROOT / "data" / "industry_study" / "generated_revision"
PPTX = Path(
    os.environ.get("FIDC_TEST_PPTX", REVISION_DIR / "industry_executive_revised.pptx")
)
PAYLOAD = Path(
    os.environ.get("FIDC_TEST_PAYLOAD", REVISION_DIR / "artifact_payload.json")
)


def _presentation() -> Presentation:
    return Presentation(PPTX)


def _slide_text(slide) -> str:
    return "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def _find_slide_index(slides: list[str], *tokens: str) -> int:
    normalized = [token.casefold() for token in tokens]
    return next(
        index
        for index, text in enumerate(slides)
        if all(token in text.casefold() for token in normalized)
    )


def test_deck_has_no_duplicate_text_box_at_same_position() -> None:
    for slide_number, slide in enumerate(_presentation().slides, start=1):
        keys = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = re.sub(r"\s+", " ", shape.text).strip()
            if not text:
                continue
            keys.append((int(shape.left), int(shape.top), text))
        duplicates = [key for key, count in Counter(keys).items() if count > 1]
        assert not duplicates, (
            f"slide {slide_number} contém caixa duplicada em posição e texto: "
            f"{duplicates}"
        )


def test_deck_has_no_truncated_fragments_or_visible_technical_slugs() -> None:
    visible = "\n".join(_slide_text(slide) for slide in _presentation().slides)
    assert "…" not in visible
    assert "..." not in visible
    assert "TOS CREDITORIOS CEDIDOS SAO ORIUNDOS DO" not in visible.upper()
    assert "ESENTE ANEXO, A CLASSE PODERA ADQUIRIR DO S" not in visible.upper()
    assert not re.search(r"(?<=[a-záéíóúãõç])_[a-záéíóúãõç]", visible)
    assert not re.search(
        r"(?:^|\n)\s*[•\-]?\s*(Vale destacar|É importante notar|Cabe ressaltar)",
        visible,
        flags=re.IGNORECASE,
    )
    assert not re.search(
        r"\b(robusto|robusta|robustos|robustas|expressivo|expressiva|"
        r"expressivos|expressivas|significativo|significativa|"
        r"significativos|significativas)\b",
        visible,
        flags=re.IGNORECASE,
    )


def test_removed_market_share_sections_remain_available_in_the_payload() -> None:
    slides = [_slide_text(slide) for slide in _presentation().slides]
    provider_ranking = _find_slide_index(
        slides,
        "QI lidera administração; BTG lidera gestão e custódia",
    )
    provider_concentration = _find_slide_index(
        slides,
        "PRESTADORES · RANKING E CONCENTRAÇÃO",
    )
    top20 = _find_slide_index(slides, "RANKING · TOP 20 FIDCs")
    top20_other = _find_slide_index(
        slides, "Outros: o único bloco que encolheu"
    )
    flagship = _find_slide_index(slides, "RISCO ESTRUTURAL · COBERTURA POR TAXONOMIA")
    carteira_1 = _find_slide_index(slides, "RISCO ESTRUTURAL · CARTEIRA VS. PARES")
    investor_base = _find_slide_index(slides, "Quase todo o volume vai para o investidor profissional")
    assert top20 < top20_other < flagship < carteira_1 < provider_ranking
    assert provider_ranking < provider_concentration < investor_base
    visible = "\n".join(slides)
    for removed in (
        "MARKET SHARE · ADMINISTRAÇÃO",
        "MARKET SHARE · GESTÃO",
        "MARKET SHARE · CUSTÓDIA",
        "ADMINISTRAÇÃO POR SUBTIPO",
        "GESTÃO POR SUBTIPO",
        "CUSTÓDIA POR SUBTIPO",
    ):
        assert removed not in visible
    payload = json.loads(PAYLOAD.read_text(encoding="utf-8"))
    assert payload["market_share"]
    assert payload["market_share_top10_fixed"]


def test_top20_rankings_remain_in_payload_and_profiles_are_omitted_from_deck() -> None:
    payload = json.loads(PAYLOAD.read_text(encoding="utf-8"))
    assert len(payload["top20_fidcs"]) == 20
    assert len(payload["top20_outros"]) == 20
    assert [int(row["rank"]) for row in payload["top20_fidcs"]] == list(
        range(1, 21)
    )
    assert [int(row["rank_outros"]) for row in payload["top20_outros"]] == list(
        range(1, 21)
    )
    assert all(
        row.get("anbima_tipo") not in {None, "", "N/D"}
        and row.get("anbima_foco") not in {None, "", "N/D"}
        for row in payload["top20_fidcs"]
    )

    assert len(payload["profiles"]) == 20
    visible = "\n".join(_slide_text(slide) for slide in _presentation().slides)
    assert "APÊNDICE · CURADORIA TOP 20" not in visible


def test_hhi_uses_antitrust_points_scale() -> None:
    frame = pd.read_csv(REVISION_DIR / "monoestrutura_concentracao.csv")
    hhi = pd.to_numeric(frame["hhi_fundos"], errors="raise")
    assert hhi.between(0, 10_000).all()
    assert hhi.max() == 10_000
    visible = "\n".join(_slide_text(slide) for slide in _presentation().slides)
    assert "CONCENTRAÇÃO DAS MONOESTRUTURAS" not in visible


def test_aging_reconciles_to_full_table_i_and_ex360_is_published() -> None:
    qa = pd.read_csv(REVISION_DIR / "qa_inadimplencia_competencia.csv")
    latest = qa.sort_values("competencia").iloc[-1]
    assert latest["aging_publication_status"] == "publicável"
    assert bool(latest["inadimplencia_ex_360d_publicavel"])
    assert abs(float(latest["aging_gap_vs_tabela_i_completa_brl"])) < 1.0
    assert float(latest["aging_parcelas_inadimplentes_brl"]) > 6_000_000_000


def test_table_i_ii_reconciliation_and_documentary_classification_are_published() -> None:
    reconciliation = pd.read_csv(
        REVISION_DIR / "reconciliacao_tabelas_i_ii_resumo.csv"
    )
    assert set(reconciliation["competencia"]) == {"2023-12", "2026-06"}
    latest = reconciliation.set_index("competencia").loc["2026-06"]
    assert int(latest["fundos_sem_abertura_tabela_ii"]) >= 0
    assert 0 <= float(latest["gap_positivo_top20_share"]) <= 1

    payload = json.loads(PAYLOAD.read_text(encoding="utf-8"))
    evidence = next(
        row
        for row in payload["classification_coverage"]
        if row["categoria"] == "Evidência documental"
    )
    assert float(evidence["pl"]) > 0
    mt_global = next(
        row
        for row in payload["top20_fidcs"]
        if row["cnpj_fundo"] == "63953619000130"
    )
    assert mt_global["anbima_tipo"] == "Financeiro"
    assert mt_global["anbima_foco"] == "Crédito Consignado"


def test_native_line_charts_keep_markers_and_smoothing_disabled() -> None:
    with zipfile.ZipFile(PPTX) as archive:
        charts = b"".join(
            archive.read(name)
            for name in archive.namelist()
            if "/charts/chart" in name and name.endswith(".xml")
        )
    assert b'<c:smooth val="1"' not in charts
    assert b'<c:smooth val="true"' not in charts
    for chunk in charts.replace(b" />", b"/>").split(b"<c:marker>")[1:]:
        marker = chunk.split(b"</c:marker>", 1)[0]
        assert b'<c:symbol val="none"' in marker
