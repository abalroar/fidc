from __future__ import annotations

from datetime import date
from io import BytesIO
from pathlib import Path

import pandas as pd
import pytest

from services.cloudwalk_financial_cost import (
    CostRunConfig,
    build_financial_cost_outputs,
    funding_lines_from_frame,
    load_spread_overrides,
)
from services.cloudwalk_financial_cost_exports import (
    build_cloudwalk_financial_cost_pptx_bytes,
    build_cloudwalk_financial_cost_xlsx_bytes,
)
from services.financial_cost_scope import (
    SCOPE_CLOUDWALK,
    SCOPE_CNPJS,
    SCOPE_PORTFOLIO,
    build_cloudwalk_scope,
    curation_data_signature,
    parse_manual_cnpj_selection,
    resolve_scope_curation,
    scope_from_cnpjs,
    scope_from_portfolio,
)
from services.portfolio_store import PortfolioFund, PortfolioRecord
from tabs import tab_cloudwalk_financial_cost as cost_tab


ROOT = Path(__file__).resolve().parents[1]
CLOUDWALK_EMISSIONS = ROOT / "data/regulatory_profiles/cloudwalk_cotas_emissoes_pagamentos.csv"
CLOUDWALK_CONFIG = ROOT / "config/cloudwalk_financial_cost_inputs.json"


def _emission_row(
    *,
    fund: str,
    cnpj: str,
    series: str = "1ª série sênior",
    spread: str = "Taxa DI + 1,00% a.a.",
    volume: str = "R$ 1.000.000,00",
) -> dict[str, str]:
    return {
        "Fundo": fund,
        "CNPJ": cnpj,
        "Cota/Classe": series,
        "Tipo": "Sênior",
        "Data deliberação": "01/01/2025",
        "Data emissão / 1ª integralização": "01/01/2025",
        "Data encerramento/oferta": "",
        "Quantidade": "",
        "Volume": volume,
        "VNU": "",
        "Remuneração": spread,
        "Juros/remuneração": spread,
        "Amortização principal": "Sem calendário fixo identificado",
        "Status/evidência": "evento de cota",
        "Fonte": "documento.pdf · p.10",
        "Status curadoria": "curado documental",
    }


def test_manual_cnpj_parser_validates_deduplicates_and_preserves_order() -> None:
    parsed = parse_manual_cnpj_selection(
        "08.417.544/0001-65, 34408539000104\n08.417.544/0001-65; 11.111.111/1111-11"
    )

    assert parsed.cnpjs == ("08417544000165", "34408539000104")
    assert parsed.duplicates == ("08417544000165",)
    assert parsed.invalid == ("11.111.111/1111-11",)


def test_scope_builders_keep_cloudwalk_as_system_preset_and_portfolio_identity() -> None:
    cloudwalk = build_cloudwalk_scope(CLOUDWALK_EMISSIONS)
    portfolio = PortfolioRecord(
        id="portfolio-1",
        name="Carteira teste",
        funds=(PortfolioFund(cnpj="08417544000165", display_name="FIDC Teste"),),
        created_at="2026-01-01T00:00:00Z",
        updated_at="2026-01-01T00:00:00Z",
    )

    assert cloudwalk.kind == SCOPE_CLOUDWALK
    assert cloudwalk.label == "CloudWalk"
    assert len(cloudwalk.cnpjs) == 11
    assert scope_from_portfolio(portfolio).kind == SCOPE_PORTFOLIO
    assert scope_from_portfolio(portfolio).cnpjs == ("08417544000165",)
    assert scope_from_cnpjs(["08417544000165"]).kind == SCOPE_CNPJS


def test_saved_portfolio_default_finds_cloudwalk_basket_instead_of_first_item() -> None:
    cloudwalk = build_cloudwalk_scope(CLOUDWALK_EMISSIONS)
    first = PortfolioRecord(
        id="alphabetical-first",
        name="A primeira carteira",
        funds=(PortfolioFund(cnpj="08417544000165", display_name="FIDC A"),),
        created_at="2026-01-01T00:00:00Z",
        updated_at="2026-01-01T00:00:00Z",
    )
    cloudwalk_portfolio = PortfolioRecord(
        id="cloudwalk-saved",
        name="Cloudwalk salva",
        funds=tuple(PortfolioFund(cnpj=cnpj, display_name=name) for cnpj, name in cloudwalk.fund_names),
        created_at="2026-01-01T00:00:00Z",
        updated_at="2026-01-01T00:00:00Z",
    )

    assert cost_tab._default_portfolio_id([first, cloudwalk_portfolio], cloudwalk) == "cloudwalk-saved"


def test_scope_curation_prefers_dedicated_profile_over_all_fidcs(tmp_path: Path) -> None:
    curated_dir = tmp_path / "profiles"
    knowledge_dir = tmp_path / "knowledge"
    curated_dir.mkdir()
    knowledge_dir.mkdir()
    dedicated = pd.DataFrame(
        [_emission_row(fund="FIDC Dedicado", cnpj="08.417.544/0001-65", spread="CDI + 1,25% a.a.")]
    )
    triage = pd.DataFrame(
        [_emission_row(fund="FIDC Triagem", cnpj="08.417.544/0001-65", spread="CDI + 9,00% a.a.")]
    )
    dedicated.to_csv(curated_dir / "dedicado_cotas_emissoes_pagamentos.csv", index=False)
    triage.to_csv(curated_dir / "all_fidcs_cotas_emissoes_pagamentos.csv", index=False)

    result = resolve_scope_curation(
        scope_from_cnpjs(["08417544000165"]),
        curated_dir=curated_dir,
        knowledge_dir=knowledge_dir,
    )
    lines = funding_lines_from_frame(result.emissions_df)

    assert len(lines) == 1
    assert lines[0].fund_name == "FIDC Dedicado"
    assert lines[0].spread_aa == pytest.approx(0.0125)
    assert result.coverage_df.iloc[0]["profile_type"] == "curado"


def test_ambiguous_series_are_blocked_instead_of_double_counted(tmp_path: Path) -> None:
    curated_dir = tmp_path / "profiles"
    knowledge_dir = tmp_path / "knowledge"
    curated_dir.mkdir()
    knowledge_dir.mkdir()
    frame = pd.DataFrame(
        [
            _emission_row(fund="FIDC A", cnpj="08.417.544/0001-65", volume="R$ 1.000.000,00"),
            _emission_row(fund="FIDC A", cnpj="08.417.544/0001-65", volume="R$ 2.000.000,00"),
        ]
    )
    frame.to_csv(curated_dir / "all_fidcs_cotas_emissoes_pagamentos.csv", index=False)

    result = resolve_scope_curation(
        scope_from_cnpjs(["08417544000165"]),
        curated_dir=curated_dir,
        knowledge_dir=knowledge_dir,
    )

    assert result.emissions_df.empty
    assert int(result.coverage_df.iloc[0]["ambiguous_rows_blocked"]) == 2
    assert result.coverage_df.iloc[0]["status"] == "Curadoria ambígua; séries bloqueadas"


def test_manual_override_is_exact_per_fund_and_keeps_curated_value() -> None:
    frame = pd.DataFrame(
        [
            _emission_row(fund="FIDC A", cnpj="08.417.544/0001-65", spread="CDI + 1,00% a.a."),
            _emission_row(fund="FIDC B", cnpj="34.408.539/0001-04", spread="CDI + 2,00% a.a."),
        ]
    )
    lines = funding_lines_from_frame(
        frame,
        spread_overrides={"08417544000165|1ª série sênior": 0.03},
    )

    assert lines[0].curated_spread_aa == pytest.approx(0.01)
    assert lines[0].manual_spread_aa == pytest.approx(0.03)
    assert lines[0].spread_aa == pytest.approx(0.03)
    assert lines[1].curated_spread_aa == pytest.approx(0.02)
    assert lines[1].manual_spread_aa is None
    assert lines[1].spread_aa == pytest.approx(0.02)


def test_cloudwalk_regression_keeps_15_curated_and_4_manual_spreads() -> None:
    scope = build_cloudwalk_scope(CLOUDWALK_EMISSIONS)
    curation = resolve_scope_curation(scope)
    overrides = load_spread_overrides(CLOUDWALK_CONFIG)
    lines = funding_lines_from_frame(curation.emissions_df, spread_overrides=overrides)
    active = [line for line in lines if line.included]

    assert len(curation.emissions_df) == 27
    assert len(active) == 19
    assert sum(line.curated_spread_aa is not None for line in active) == 15
    assert sum(line.manual_spread_aa is not None for line in active) == 4
    assert sum(line.spread_aa is None for line in active) == 0


def test_clearing_manual_value_restores_curated_spread_in_editor() -> None:
    frame = pd.DataFrame(
        [_emission_row(fund="FIDC A", cnpj="08.417.544/0001-65", spread="CDI + 1,00% a.a.")]
    )
    curated_lines = funding_lines_from_frame(frame)
    with_manual = cost_tab._spread_input_table(
        curated_lines,
        {"08417544000165|1ª série sênior": 0.03},
    )
    restored = cost_tab._spread_input_table(curated_lines, {})

    assert with_manual.iloc[0]["CDI+ efetivo (% a.a.)"] == pytest.approx(3.0)
    assert with_manual.iloc[0]["Origem efetiva"] == "Manual"
    assert restored.iloc[0]["CDI+ efetivo (% a.a.)"] == pytest.approx(1.0)
    assert restored.iloc[0]["Origem efetiva"] == "Documento/curadoria"


def test_curation_signature_changes_when_source_file_changes(tmp_path: Path) -> None:
    path = tmp_path / "emissions.csv"
    pd.DataFrame([_emission_row(fund="FIDC A", cnpj="08.417.544/0001-65")]).to_csv(path, index=False)
    scope = build_cloudwalk_scope(path)
    before = curation_data_signature(scope)
    path.write_text(path.read_text(encoding="utf-8") + "\n", encoding="utf-8")
    after = curation_data_signature(scope)

    assert before != after


def test_generic_export_prefix_does_not_keep_cloudwalk_name() -> None:
    assert cost_tab._scope_export_prefix(SCOPE_CLOUDWALK, "CloudWalk") == "cloudwalk_financial_cost"
    assert cost_tab._scope_export_prefix(SCOPE_PORTFOLIO, "Crédito Ágil 1T26") == "credito_agil_1t26_financial_cost"


def test_generic_exports_use_scope_label_instead_of_cloudwalk() -> None:
    from openpyxl import load_workbook
    from pptx import Presentation

    lines = funding_lines_from_frame(
        pd.DataFrame([_emission_row(fund="FIDC A", cnpj="08.417.544/0001-65")])
    )
    outputs = build_financial_cost_outputs(
        lines=lines,
        snapshots=[],
        config=CostRunConfig(
            start_date=date(2026, 1, 1),
            end_date=date(2026, 12, 31),
            snapshot_date=date(2026, 5, 14),
            cdi_aa=0.13,
            cdi_source="fixture",
            scope_label="Carteira Teste",
            scope_kind=SCOPE_PORTFOLIO,
        ),
    )

    workbook = load_workbook(BytesIO(build_cloudwalk_financial_cost_xlsx_bytes(outputs, scope_label="Carteira Teste")))
    presentation = Presentation(BytesIO(build_cloudwalk_financial_cost_pptx_bytes(outputs, scope_label="Carteira Teste")))
    slide_text = " ".join(shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text_frame"))

    assert workbook["Resumo"]["A1"].value == "Carteira Teste - custo financeiro de FIDCs"
    assert "Carteira Teste" in slide_text
    assert "CloudWalk" not in slide_text
