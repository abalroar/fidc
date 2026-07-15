from __future__ import annotations

from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import pytest

from services.pptx_merge import PptxMergeError, merge_pptx_bytes


def _native_deck_bytes(
    *,
    slide_title: str,
    chart_title: str,
    chart_values: tuple[float, float],
    include_table: bool = False,
    width_inches: float = 13.333,
) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(width_inches)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(5), Inches(0.4)).text = slide_title

    chart_data = CategoryChartData()
    chart_data.categories = ["jan/26", "fev/26"]
    chart_data.add_series(chart_title, chart_values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4),
        Inches(0.9),
        Inches(6),
        Inches(4),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title

    if include_table:
        table = slide.shapes.add_table(
            2,
            2,
            Inches(7.0),
            Inches(0.9),
            Inches(5.5),
            Inches(1.2),
        ).table
        table.cell(0, 0).text = "Escopo"
        table.cell(0, 1).text = "Valor"
        table.cell(1, 0).text = "Fundo B"
        table.cell(1, 1).text = "12,34%"

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_merge_pptx_bytes_preserves_native_charts_tables_and_embedded_workbooks() -> None:
    primary = _native_deck_bytes(
        slide_title="Deck principal",
        chart_title="Gráfico principal",
        chart_values=(1.0, 2.0),
    )
    appendix = _native_deck_bytes(
        slide_title="Deck complementar",
        chart_title="Gráfico complementar",
        chart_values=(10.0, 20.0),
        include_table=True,
    )

    merged = merge_pptx_bytes(primary, appendix)

    presentation = Presentation(BytesIO(merged))
    assert len(presentation.slides) == 2
    assert "Deck principal" in "\n".join(
        shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
    )
    assert "Deck complementar" in "\n".join(
        shape.text for shape in presentation.slides[1].shapes if hasattr(shape, "text")
    )
    assert sum(bool(getattr(shape, "has_chart", False)) for shape in presentation.slides[0].shapes) == 1
    assert sum(bool(getattr(shape, "has_chart", False)) for shape in presentation.slides[1].shapes) == 1
    tables = [shape.table for shape in presentation.slides[1].shapes if getattr(shape, "has_table", False)]
    assert len(tables) == 1
    assert tables[0].cell(1, 0).text == "Fundo B"

    with ZipFile(BytesIO(merged)) as archive:
        chart_parts = sorted(
            name
            for name in archive.namelist()
            if name.startswith("ppt/charts/chart") and name.endswith(".xml")
        )
        workbook_parts = sorted(
            name
            for name in archive.namelist()
            if name.startswith("ppt/embeddings/") and name.endswith(".xlsx")
        )
        slide_two_relationships = archive.read("ppt/slides/_rels/slide2.xml.rels").decode("utf-8")
        workbook_values = [
            load_workbook(BytesIO(archive.read(name)), data_only=False).active["B2"].value
            for name in workbook_parts
        ]

    assert len(chart_parts) == 2
    assert len(workbook_parts) == 2
    assert sorted(workbook_values) == [1.0, 10.0]
    assert "/chart" in slide_two_relationships
    assert "slideLayout" in slide_two_relationships

    # A second Office-package round trip must keep both native objects intact.
    roundtrip = BytesIO()
    presentation.save(roundtrip)
    reopened = Presentation(BytesIO(roundtrip.getvalue()))
    assert len(reopened.slides) == 2
    assert any(getattr(shape, "has_table", False) for shape in reopened.slides[1].shapes)
    assert any(getattr(shape, "has_chart", False) for shape in reopened.slides[1].shapes)


def test_merge_pptx_bytes_rejects_different_slide_sizes() -> None:
    primary = _native_deck_bytes(
        slide_title="Principal",
        chart_title="A",
        chart_values=(1.0, 2.0),
    )
    appendix = _native_deck_bytes(
        slide_title="Apêndice",
        chart_title="B",
        chart_values=(3.0, 4.0),
        width_inches=10.0,
    )

    with pytest.raises(PptxMergeError, match="tamanho de slide diferente"):
        merge_pptx_bytes(primary, appendix)


def test_merge_pptx_bytes_rejects_incompatible_theme_chain() -> None:
    primary = _native_deck_bytes(
        slide_title="Principal",
        chart_title="A",
        chart_values=(1.0, 2.0),
    )
    appendix = _native_deck_bytes(
        slide_title="Apêndice",
        chart_title="B",
        chart_values=(3.0, 4.0),
    )
    mutated = BytesIO()
    with ZipFile(BytesIO(appendix), "r") as source, ZipFile(mutated, "w", compression=ZIP_DEFLATED) as target:
        for info in source.infolist():
            payload = source.read(info.filename)
            if info.filename == "ppt/theme/theme1.xml":
                payload = payload.replace(b"Office Theme", b"Merged Theme")
            target.writestr(info, payload)

    with pytest.raises(PptxMergeError, match="layout/master/tema incompatível"):
        merge_pptx_bytes(primary, mutated.getvalue())


def test_merge_pptx_bytes_returns_primary_unchanged_without_appendices() -> None:
    primary = _native_deck_bytes(
        slide_title="Principal",
        chart_title="A",
        chart_values=(1.0, 2.0),
    )

    assert merge_pptx_bytes(primary) == primary


def test_merge_pptx_bytes_rejects_invalid_payload() -> None:
    primary = _native_deck_bytes(
        slide_title="Principal",
        chart_title="A",
        chart_values=(1.0, 2.0),
    )

    with pytest.raises(PptxMergeError, match="não é um PPTX válido"):
        merge_pptx_bytes(primary, b"not-a-pptx")
