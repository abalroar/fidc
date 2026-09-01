from __future__ import annotations

import csv
import json
from pathlib import Path

import pandas as pd
import pytest

from services.industry_requested_revision import (
    PF_PJ_CATEGORY, build_category_cnpj_ledger,
    build_issuance_by_display_category, build_provider_comparison,
    build_stock_scenarios, build_credit_screen, load_pfpj_ledger, prepare_funds,
)


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "data" / "industry_study" / "generated_revision" / "directors_update"


def test_carteira_101_one_senior_account_reconciliation() -> None:
    payload = json.loads((OUTPUT / "fidc_directors_update_data.json").read_text())
    summary = payload["carteira_101"]["summary"]
    funds = payload["carteira_101"]["one_senior_account_funds"]

    assert summary["portfolio_funds"] == 101
    assert summary["funds_with_cvm_data"] == 78
    assert summary["one_senior_account_funds"] == 19
    assert summary["pl_coverage"] > 0.98
    assert summary["primary_competence"] == "2026-07"
    assert summary["fallback_competence"] == "2026-06"
    assert len(funds) == 19
    assert all(row["contas_senior_reportadas"] == 1 for row in funds)
    assert [row["pl_publicado_brl"] for row in funds] == sorted(
        (row["pl_publicado_brl"] for row in funds), reverse=True
    )

    with (OUTPUT / "carteira_101_cotistas_senior_202607.csv").open(
        newline="", encoding="utf-8"
    ) as handle:
        assert sum(1 for _ in csv.DictReader(handle)) == 101


def test_financeiro_decomposition_closes_and_keeps_unsegmented_buckets() -> None:
    payload = json.loads((OUTPUT / "fidc_directors_update_data.json").read_text())
    summary = payload["financeiro"]["summary"]
    rows = payload["financeiro"]["decomposition"]

    assert summary["financeiro_funds"] == 1103
    assert round(sum(row["pl_brl"] for row in rows), 2) == round(
        summary["financeiro_pl_brl"], 2
    )
    assert abs(sum(row["share"] for row in rows) - 1) < 1e-8
    assert sum(row["fundos"] for row in rows) == summary["financeiro_funds"]
    assert summary["tapso_inside_financeiro"] is True
    assert summary["tapso_type_applied"] == "Financeiro"
    assert summary["tapso_pl_brl"] > 41_000_000_000
    assert any("sem segregacao" in row["bucket_financeiro"] for row in rows)


def test_director_provider_top_five_preserved_and_both_comparators_appended():
    rows = [dict(competencia="2026-06", papel="gestor", participante=name,
                 pl_brl=pl, fundos=1, denominador_pl_brl=1000)
            for name, pl in [("A",100),("B",90),("C",80),("D",70),("E",60),
                             ("Itaú",10),("Kinea",5),("Kanastra Administracao De Recursos",8),("Limine Trust",7)]]
    ranking, lineage = build_provider_comparison(pd.DataFrame(rows), "2026-06")
    selected = ranking.dropna(subset=["ordem_slide"]).sort_values("ordem_slide")
    assert selected.participante.tolist() == ["A","B","C","D","E","Itaú","Kanastra (incl. Limine)"]
    assert selected.set_index("participante").loc["Itaú", "pl_brl"] == 15
    assert selected.set_index("participante").loc["Kanastra (incl. Limine)", "pl_brl"] == 15
    assert ranking.pl_brl.sum() == lineage.pl_brl.sum()


def test_director_exclusion_changes_numerator_and_denominator():
    f = pd.DataFrame([dict(competencia="2026-06",categoria_slide="Financeiro",pl=40,excluido_cenario=True),
                      dict(competencia="2026-06",categoria_slide="Financeiro",pl=10,excluido_cenario=False),
                      dict(competencia="2026-06",categoria_slide="Fomento Mercantil",pl=50,excluido_cenario=False)])
    r = build_stock_scenarios(f,["2026-06"])
    s = r[r.cenario.eq("sem_tapso_petrobras")]
    assert s.pl_brl.sum() == 60
    assert s.loc[s.categoria.eq("Financeiro"),"share"].iloc[0] == pytest.approx(1/6)
    assert s.share.sum() == pytest.approx(1)


def test_credit_screen_does_not_infer_pulverization_or_classify_from_name():
    f = pd.DataFrame([
        dict(denominacao="PULVERIZADO PF",anbima_tipo_curado="Financeiro",anbima_foco_curado="Multicarteira Financeiro",taxonomia_funcional_n1_curada="",taxonomia_funcional_n2_curada="",categoria_slide="Financeiro"),
        dict(denominacao="X",anbima_tipo_curado="Financeiro",anbima_foco_curado="Crédito Pessoal",taxonomia_funcional_n1_curada="",taxonomia_funcional_n2_curada="",categoria_slide="Financeiro"),
        dict(denominacao="Y",anbima_tipo_curado="Fomento Mercantil",anbima_foco_curado="Fomento Mercantil",taxonomia_funcional_n1_curada="Crédito PJ",taxonomia_funcional_n2_curada="",categoria_slide="Fomento Mercantil"),
    ]).assign(competencia="2026-06",pl=10,segmento_financeiro_principal="")
    r = build_credit_screen(f,"2026-06")
    assert r.recorte_credito.tolist() == ["Financeiro sem segregação","PF pessoal / estudantil / BNPL (triagem)","Fomento Mercantil"]
    assert r.pulverizacao_validada.str.startswith("N/D").all()
    assert r.share_total.sum() == pytest.approx(1)


def test_requested_revision_downloads_match_the_approved_artifacts():
    from io import BytesIO
    from zipfile import ZipFile
    from pptx import Presentation
    from services.industry_requested_revision_export import (
        DOWNLOADS, RELEASE_DIR, load_requested_revision_downloads,
    )

    data_dir = ROOT / "data/industry_study"
    downloads = load_requested_revision_downloads(data_dir)
    assert set(downloads) == {"complete", "slides", "package"}
    assert downloads["complete"] == (data_dir / RELEASE_DIR / DOWNLOADS["complete"]).read_bytes()
    deck = Presentation(BytesIO(downloads["complete"]))
    assert len(deck.slides) == 33
    deck_text = [
        " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in deck.slides
    ]
    assert "Posição competitiva do Itaú BBA" in deck_text[13]
    assert "Em FIDC o Itaú BBA lidera com 45,7%" in deck_text[20]
    assert "Premissas, limitações e fontes" in deck_text[25]
    assert "Volume por prestador" in deck_text[27]
    provider_xml = "".join(
        rel.target_part.blob.decode()
        for rel in deck.slides[27].part.rels.values()
        if rel.reltype.endswith("/chart")
    )
    for color in ("FF5500", "7030A0", "2456D6", "1D4080", "7A1F3D", "30353A"):
        assert color in provider_xml
    with ZipFile(BytesIO(downloads["package"])) as archive:
        assert archive.read(DOWNLOADS["slides"]) == downloads["slides"]
        stock = pd.read_csv(archive.open("bases/saldo_cenarios.csv"))
        latest = stock[stock.competencia.eq("2026-06")]
        assert latest.loc[latest.cenario.eq("sem_tapso_petrobras"), "pl_brl"].sum() == pytest.approx(718610541429.85)
        report = archive.read("Relatorio_Revisao_Diretoria_v2.md").decode()
        assert "R$ 7.102.640.442,08" in report
        assert "Número total de devedores: **N/D**" in report
        assert "33 slides" in report
        assert "ANBIMA" in report


@pytest.mark.parametrize("broken", ["complete", "package"])
def test_requested_revision_rejects_corrupt_materialized_downloads(tmp_path, broken):
    import shutil
    from services.industry_requested_revision_export import (
        DOWNLOADS, RELEASE_DIR, load_requested_revision_downloads,
    )

    shutil.copytree(ROOT / "data/industry_study" / RELEASE_DIR, tmp_path / RELEASE_DIR)
    path = tmp_path / RELEASE_DIR / DOWNLOADS[broken]
    path.write_bytes(path.read_bytes() + b"corrupt")
    with pytest.raises(ValueError, match="divergente do manifesto"):
        load_requested_revision_downloads(tmp_path)


def test_requested_revision_rejects_partial_release(tmp_path):
    import shutil
    from services.industry_requested_revision_export import RELEASE_DIR, load_requested_revision_downloads

    shutil.copytree(ROOT / "data/industry_study" / RELEASE_DIR, tmp_path / RELEASE_DIR)
    (tmp_path / RELEASE_DIR / "release.json").unlink()
    with pytest.raises(FileNotFoundError):
        load_requested_revision_downloads(tmp_path)


def test_pfpj_ledger_preserves_26_decisions_and_excludes_11_and_26():
    path = ROOT / "data/industry_study/director_pfpj_ledger_20260901.csv"
    ledger = load_pfpj_ledger(path)
    included = ledger[ledger.incluir_pfpj]
    excluded = ledger[~ledger.incluir_pfpj]
    assert len(included) == 24
    assert set(excluded.ordem.astype(int)) == {11, 26}
    assert included.pl_brl.sum() == pytest.approx(7_102_640_442.08)
    assert set(excluded.nome_exibicao) == {"Sólido", "BizCapital Finpass PME"}


def test_pfpj_overlay_moves_only_selected_financeiro_rows():
    ledger = pd.DataFrame([
        {"cnpj_fundo": "00000000000001", "incluir_pfpj": True},
        {"cnpj_fundo": "00000000000002", "incluir_pfpj": False},
    ])
    funds = pd.DataFrame([
        {"competencia": "2026-06", "cnpj_fundo": "00000000000001", "cnpj_classe": "", "is_fic_fidc": False, "pl": 7, "anbima_tipo": "Financeiro", "anbima_foco": "Crédito Pessoal"},
        {"competencia": "2026-06", "cnpj_fundo": "00000000000002", "cnpj_classe": "", "is_fic_fidc": False, "pl": 3, "anbima_tipo": "Financeiro", "anbima_foco": "Crédito Pessoal"},
    ])
    actions = pd.DataFrame(columns=["cnpj_fundo", "status"])
    frame = prepare_funds(funds, actions, ledger)
    assert frame.set_index("cnpj_fundo").loc["00000000000001", "categoria_slide"] == PF_PJ_CATEGORY
    assert frame.set_index("cnpj_fundo").loc["00000000000002", "categoria_slide"] == "Financeiro"


def test_issuance_uses_fund_then_class_and_keeps_unmatched_as_nd():
    ledger = pd.DataFrame([
        {"cnpj_fundo": "00000000000001", "cnpj_classe": "00000000000011", "categoria_slide": PF_PJ_CATEGORY},
        {"cnpj_fundo": "00000000000002", "cnpj_classe": "", "categoria_slide": "Financeiro"},
    ])
    offers = pd.DataFrame([
        {"cnpj_emissor": "00000000000001", "data_encerramento": "2023-01-10", "registered_volume_brl": 10},
        {"cnpj_emissor": "00000000000011", "data_encerramento": "2024-01-10", "registered_volume_brl": 20},
        {"cnpj_emissor": "00000000000099", "data_encerramento": "2025-01-10", "registered_volume_brl": 30},
        {"cnpj_emissor": "00000000000002", "data_encerramento": "2025-02-10", "registered_volume_brl": 40},
        {"cnpj_emissor": "00000000000002", "data_encerramento": "2026-02-10", "registered_volume_brl": 50},
    ])
    summary, detail, audit = build_issuance_by_display_category(
        offers, ledger, anbima_2023_brl=10
    )
    assert detail.loc[detail.cnpj_n.eq("00000000000001"), "match_categoria"].iloc[0] == "cnpj_fundo"
    assert detail.loc[detail.cnpj_n.eq("00000000000011"), "match_categoria"].iloc[0] == "cnpj_classe"
    assert detail.loc[detail.cnpj_n.eq("00000000000099"), "categoria"].iloc[0] == "N/D"
    assert summary.groupby("period_key").share.sum().eq(1).all()
    assert audit.loc[audit.period_key.eq("2025"), "nao_localizado_emissores"].iloc[0] == 1
