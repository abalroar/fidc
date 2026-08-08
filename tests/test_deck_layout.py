"""Reordenar lâminas sem quebrar o pacote nem a numeração."""

from __future__ import annotations

from pathlib import Path
import sys

import pytest


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from services.deck_layout import (  # noqa: E402
    move_slides,
    page_box,
    renumber_pages,
)


def deck(quantas: int):
    from pptx import Presentation
    from pptx.util import Inches

    from services.bba_deck import Deck, SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN

    apresentacao = Presentation()
    apresentacao.slide_width = Inches(SLIDE_WIDTH_IN)
    apresentacao.slide_height = Inches(SLIDE_HEIGHT_IN)
    montador = Deck("BLOCO")
    montador.prs = apresentacao
    montador.page = 0
    for numero in range(1, quantas + 1):
        montador.footer(montador.slide(f"Lâmina {numero}"), "fonte")
    return apresentacao


def titulos(apresentacao) -> list[str]:
    saida = []
    for slide in apresentacao.slides:
        for forma in slide.shapes:
            texto = forma.text_frame.text.strip() if forma.has_text_frame else ""
            if texto.startswith("Lâmina "):
                saida.append(texto)
                break
    return saida


def test_o_bloco_vai_para_antes_da_ancora() -> None:
    """Mover 9–17 para 52 põe o bloco logo antes da lâmina que hoje é a 52ª."""

    apresentacao = deck(56)
    move_slides(apresentacao, 9, 17, 52)

    nomes = titulos(apresentacao)
    assert len(nomes) == 56
    # O que vinha antes do bloco não se move.
    assert nomes[:8] == [f"Lâmina {n}" for n in range(1, 9)]
    # O bloco reaparece inteiro e na ordem, terminando na 51ª posição.
    assert nomes[42:51] == [f"Lâmina {n}" for n in range(9, 18)]
    # E a lâmina 52 continua sendo a 52ª.
    assert nomes[51] == "Lâmina 52"


def test_mover_nao_duplica_nem_perde_lamina() -> None:
    apresentacao = deck(20)
    move_slides(apresentacao, 3, 6, 15)

    nomes = titulos(apresentacao)
    assert sorted(nomes, key=lambda t: int(t.split()[-1])) == [
        f"Lâmina {n}" for n in range(1, 21)
    ]


def test_mover_para_dentro_do_proprio_bloco_e_no_op() -> None:
    apresentacao = deck(10)
    antes = titulos(apresentacao)

    move_slides(apresentacao, 3, 6, 5)

    assert titulos(apresentacao) == antes


def test_um_bloco_fora_do_deck_e_recusado() -> None:
    apresentacao = deck(10)

    with pytest.raises(IndexError, match="não cabe"):
        move_slides(apresentacao, 8, 12, 3)


def test_a_renumeracao_segue_a_posicao() -> None:
    apresentacao = deck(20)
    move_slides(apresentacao, 3, 6, 15)

    assert renumber_pages(apresentacao) == 20
    numeros = [int(page_box(s).text_frame.text) for s in apresentacao.slides]
    assert numeros == list(range(1, 21))


def test_o_pacote_nao_ganha_parte_repetida(tmp_path: Path) -> None:
    """A razão de mover em vez de remover e recriar: partnames duplicados."""

    import zipfile

    apresentacao = deck(30)
    move_slides(apresentacao, 5, 12, 28)
    renumber_pages(apresentacao)
    destino = tmp_path / "deck.pptx"
    apresentacao.save(destino)

    with zipfile.ZipFile(destino) as pacote:
        partes = [n for n in pacote.namelist() if n.startswith("ppt/slides/slide")]
    assert len(partes) == len(set(partes)) == 30
