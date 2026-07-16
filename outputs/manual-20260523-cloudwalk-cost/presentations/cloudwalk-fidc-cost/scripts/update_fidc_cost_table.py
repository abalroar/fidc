from __future__ import annotations

import shutil
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from datetime import date

ROOT = Path("/Users/matheusjprates/fidc")
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from services.cloudwalk_financial_cost import (
    _balance_at,
    load_amortization_convention_overrides,
    load_funding_lines,
    load_spread_overrides,
)
from services.waterfall_schedule import only_digits

DATA = ROOT / "reports/cloudwalk_financial_cost_deck_data/line_2025.csv"
CASH_DATA = ROOT / "reports/cloudwalk_financial_cost_deck_data/ime_snapshot_2025.csv"
EMISSIONS_CSV = ROOT / "data/regulatory_profiles/cloudwalk_cotas_emissoes_pagamentos.csv"
CONFIG_JSON = ROOT / "config/cloudwalk_financial_cost_inputs.json"
PL_WATERFALL_SUMMARY = ROOT / "reports/cloudwalk_financial_cost_deck_data/pl_waterfall_summary_2025.csv"
IMPLIED_CDI_CSV = ROOT / "reports/cloudwalk_financial_cost_deck_data/implied_cdi_plus.csv"
CAPITAL_STACK_CSV = ROOT / "reports/cloudwalk_financial_cost_deck_data/capital_stack_ime_by_class.csv"
CAPITAL_STACK_OFFICIAL_CSV = ROOT / "reports/cloudwalk_financial_cost_deck_data/capital_stack_official_pl.csv"
PL_MONTHLY_CSV = ROOT / "reports/cloudwalk_financial_cost_deck_data/pl_monthly_2025_2026q1.csv"
EVENTS_2025_2026Q1_CSV = ROOT / "reports/cloudwalk_financial_cost_deck_data/events_2025_2026q1.csv"
WATERFALL_AMORT_EVENTS_CSV = ROOT / "reports/cloudwalk_financial_cost_deck_data/pl_waterfall_amortization_events_2025.csv"
DECK_MAIN = ROOT / "reports/cloudwalk_bba_deck/cloudwalk-fidc-custo-financeiro-2025.pptx"
DECK_CONDENSED = ROOT / "reports/cloudwalk_bba_deck/cloudwalk-fidc-custo-financeiro-2025-condensado.pptx"
DECK_OUTPUT = (
    ROOT
    / "outputs/manual-20260523-cloudwalk-cost/presentations/cloudwalk-fidc-cost/output/cloudwalk-fidc-custo-financeiro-2025.pptx"
)

NAVY = RGBColor(0, 35, 78)
BLUE = RGBColor(0, 87, 166)
ORANGE = RGBColor(255, 111, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(17, 24, 39)
GRID = RGBColor(214, 221, 230)
LIGHT = RGBColor(246, 248, 251)
BG = RGBColor(248, 247, 243)
GREEN = RGBColor(15, 157, 88)
RED = RGBColor(217, 48, 37)
GREY = RGBColor(120, 132, 145)
TODAY = date(2026, 5, 23)
SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)


def short_name(name: str) -> str:
    upper = name.upper()
    if "BELA" in upper:
        return "Bela"
    if "PI FUNDO" in upper:
        return "PI"
    if "A.I." in upper:
        return "A.I."
    if "AKIRA II" in upper:
        return "Akira II"
    if "AKIRA I" in upper:
        return "Akira I"
    if "KICK ASS I" in upper:
        return "Kick Ass I"
    if "BIG PICTURE III" in upper:
        return "Big Picture III"
    if "BIG PICTURE II" in upper:
        return "Big Picture II"
    if "BIG PICTURE IV" in upper:
        return "Big Picture IV"
    if "BIG PICTURE I" in upper:
        return "Big Picture I"
    return name[:24]


def fmt_bi(value: float) -> str:
    return f"{value / 1e9:.2f}".replace(".", ",") + " bi"


def fmt_bi_money(value: float) -> str:
    return "R$" + fmt_bi(value)


def fmt_mm(value: float) -> str:
    return f"R${value / 1e6:.0f} mi"


def fmt_pct(value: float) -> str:
    return f"{value * 100:.2f}%".replace(".", ",")


def fmt_pct_one(value: float) -> str:
    return f"{value * 100:.1f}%".replace(".", ",")


def fmt_bps(value: float) -> str:
    return f"{value * 10000:+.0f} bps".replace(".", ",")


def fmt_pp(value: float) -> str:
    return f"{value * 100:+.1f} p.p.".replace(".", ",")


def spread_source_label(value: object) -> str:
    text = str(value or "")
    if text.startswith("override:"):
        return "Manual"
    if text == "curadoria:remuneracao":
        return "Documento"
    if text == "pendente":
        return "Pendente"
    return text


def class_label(value: object) -> str:
    return {"senior": "Sênior", "mezzanino": "Mezanino"}.get(str(value or ""), str(value or ""))


def build_table_rows() -> list[list[str]]:
    frame = pd.read_csv(DATA)
    included = frame[frame["included_in_cost"] == True].copy()
    included["short"] = included["fund_name"].map(short_name)

    rows: list[dict[str, object]] = []
    for fund, group in included.groupby("short"):
        saldo = float(group["saldo_medio_programado"].sum())
        custo = float(group["custo_programado_bruto"].sum())
        spread_ponderado = float((group["spread_cdi_plus_aa"] * group["saldo_medio_programado"]).sum() / saldo)
        spreads = sorted(float(x) for x in group["spread_cdi_plus_aa"].dropna().unique())
        if len(spreads) <= 1:
            faixa = fmt_pct(spread_ponderado)
        else:
            faixa = f"{fmt_pct(spreads[0])}-{fmt_pct(spreads[-1])}"
        rows.append(
            {
                "fidc": fund,
                "saldo": saldo,
                "spread": spread_ponderado,
                "faixa": faixa,
                "custo": custo,
            }
        )

    rows.sort(key=lambda item: float(item["custo"]), reverse=True)
    table = [["FIDC", "Saldo méd.", "CDI+ pond.", "Faixa CDI+", "Custo 25"]]
    for row in rows:
        table.append(
            [
                str(row["fidc"]),
                fmt_bi(float(row["saldo"])),
                fmt_pct(float(row["spread"])),
                str(row["faixa"]),
                fmt_mm(float(row["custo"])),
            ]
        )
    return table


def build_series_rows() -> list[list[str]]:
    frame = pd.read_csv(DATA)
    included = frame[frame["included_in_cost"] == True].copy()
    included["short"] = included["fund_name"].map(short_name)
    included["source_label"] = included["spread_source"].map(spread_source_label)
    included = included.sort_values(["short", "classe"])
    rows = [["FIDC", "Série", "Classe", "CDI+ a.a.", "Fonte", "Volume", "Custo"]]
    for _, row in included.iterrows():
        rows.append(
            [
                row["short"],
                str(row["classe"]),
                class_label(row["class_macro"]),
                fmt_pct(float(row["spread_cdi_plus_aa"])),
                row["source_label"],
                fmt_bi(float(row["volume_emitido"])),
                fmt_mm(float(row["custo_programado_bruto"])),
            ]
        )
    return rows


def cash_lft_totals() -> tuple[float, float]:
    frame = pd.read_csv(CASH_DATA)
    cash_base = pd.to_numeric(frame.get("cash_like_caixa_lft"), errors="coerce").fillna(0.0).sum()
    cash_yield = pd.to_numeric(frame.get("rendimento_estimado_caixa_lft"), errors="coerce").fillna(0.0).sum()
    return float(cash_base), float(cash_yield)


def funding_lines():
    return load_funding_lines(
        EMISSIONS_CSV,
        spread_overrides=load_spread_overrides(CONFIG_JSON),
        amortization_convention_overrides=load_amortization_convention_overrides(CONFIG_JSON),
    )


def build_structure_rows(lines) -> pd.DataFrame:
    if CAPITAL_STACK_CSV.exists() and CAPITAL_STACK_OFFICIAL_CSV.exists():
        stack = pd.read_csv(CAPITAL_STACK_CSV)
        official = pd.read_csv(CAPITAL_STACK_OFFICIAL_CSV)
        class_map = {"senior": "Sênior", "mezzanino": "Mezanino", "subordinada": "Subordinada"}
        grouped = (
            stack.groupby(["Data", "class_macro"], as_index=False)
            .agg(Saldo=("PL classe", "sum"), FIDCs=("cnpj", "nunique"), Séries=("class_label", "count"))
        )
        records = []
        for data_label in ["31/12/25", "Hoje (23/05/26)"]:
            subset = grouped[grouped["Data"].eq(data_label)].copy()
            for macro in ["senior", "mezzanino", "subordinada"]:
                row = subset[subset["class_macro"].eq(macro)]
                if row.empty:
                    saldo, fidcs, series = 0.0, 0, 0
                else:
                    saldo = float(row["Saldo"].iloc[0])
                    fidcs = int(row["FIDCs"].iloc[0])
                    series = int(row["Séries"].iloc[0])
                records.append(
                    {
                        "Data": data_label,
                        "Classe": class_map[macro],
                        "Saldo": saldo,
                        "FIDCs": fidcs,
                        "Séries": series,
                        "Nota": "PL por classe no IME",
                    }
                )
            class_total = float(stack.loc[stack["Data"].eq(data_label), "PL classe"].sum())
            official_total = float(official.loc[official["Data"].eq(data_label), "PL oficial"].sum())
            diff = official_total - class_total
            if abs(diff) > 1_000_000:
                records.append(
                    {
                        "Data": data_label,
                        "Classe": "Não reconcil.",
                        "Saldo": diff,
                        "FIDCs": "n.d.",
                        "Séries": "n.d.",
                        "Nota": "PL oficial - soma por classe",
                    }
                )
            records.append(
                {
                    "Data": data_label,
                    "Classe": "Total",
                    "Saldo": official_total,
                    "FIDCs": int(stack.loc[stack["Data"].eq(data_label), "cnpj"].nunique()),
                    "Séries": int(stack.loc[stack["Data"].eq(data_label), "class_label"].count()),
                    "Nota": "PL oficial do IME",
                }
            )
        return pd.DataFrame(records)

    pl_2025 = float(pd.read_csv(PL_WATERFALL_SUMMARY)["pl_final"].iloc[0])
    records = []
    for label, ref_date, pl_total in [("31/12/25", date(2025, 12, 31), pl_2025), ("Hoje (23/05/26)", TODAY, None)]:
        class_rows = []
        active_total = set()
        active_series_total = 0
        for macro in ["senior", "mezzanino"]:
            macro_lines = [line for line in lines if line.included and line.class_macro == macro]
            balance = sum(_balance_at(line, ref_date) for line in macro_lines)
            active = [line for line in macro_lines if _balance_at(line, ref_date) > 0]
            active_total.update(only_digits(line.cnpj) for line in active)
            active_series_total += len(active)
            class_rows.append(
                {
                    "Data": label,
                    "Classe": "Sênior" if macro == "senior" else "Mezanino",
                    "Saldo": balance,
                    "FIDCs": len({only_digits(line.cnpj) for line in active}),
                    "Séries": len(active),
                    "Nota": "principal vivo documental",
                }
            )
        senior_mezz = sum(row["Saldo"] for row in class_rows)
        sub_residual = max((pl_total or 0.0) - senior_mezz, 0.0) if pl_total is not None else 0.0
        class_rows.append(
            {
                "Data": label,
                "Classe": "Subordinada",
                "Saldo": sub_residual,
                "FIDCs": "n.d." if sub_residual <= 0 else str(len(active_total)),
                "Séries": "n.d.",
                "Nota": "PL final - Sr - Mez" if pl_total is not None else "não reconciliado com IME recente",
            }
        )
        class_rows.append(
            {
                "Data": label,
                "Classe": "Total",
                "Saldo": senior_mezz + sub_residual,
                "FIDCs": len(active_total),
                "Séries": active_series_total,
                "Nota": "Sr + Mez + Sub/Residual",
            }
        )
        records.extend(class_rows)
    return pd.DataFrame(records)


def build_amortization_2026_rows(lines) -> pd.DataFrame:
    records = []
    for line in lines:
        if not line.included:
            continue
        events = [(item_date, amount) for item_date, amount in line.amortizations if date(2026, 1, 1) <= item_date <= date(2026, 12, 31)]
        if not events:
            continue
        records.append(
            {
                "FIDC": short_name(line.fund_name),
                "Série": line.classe,
                "Classe": class_label(line.class_macro),
                "Total": sum(amount for _, amount in events),
                "Eventos": len(events),
                "Primeira": min(item_date for item_date, _ in events).strftime("%d/%m/%y"),
                "Última": max(item_date for item_date, _ in events).strftime("%d/%m/%y"),
                "Datas": "; ".join(f"{item_date:%d/%m}: {fmt_mm(amount)}" for item_date, amount in events),
            }
        )
    return pd.DataFrame(records).sort_values("Total", ascending=False)


def build_implied_cdi_rows() -> pd.DataFrame:
    return pd.read_csv(IMPLIED_CDI_CSV)


def build_monthly_pl_summary() -> pd.DataFrame:
    frame = pd.read_csv(PL_MONTHLY_CSV)
    frame["pl"] = pd.to_numeric(frame["pl"], errors="coerce").fillna(0.0)
    output = frame.groupby("competencia", as_index=False).agg(pl=("pl", "sum"), fundos=("cnpj", "nunique"))
    output["sort_key"] = output["competencia"].map(lambda label: int(str(label).split("/")[1]) * 100 + int(str(label).split("/")[0]))
    return output.sort_values("sort_key").reset_index(drop=True)


def build_2025_amortization_by_fund() -> pd.DataFrame:
    frame = pd.read_csv(WATERFALL_AMORT_EVENTS_CSV)
    frame["valor_total"] = pd.to_numeric(frame["valor_total"], errors="coerce").fillna(0.0)
    frame["display_name"] = frame["short_name"].map(short_name)
    grouped = (
        frame.groupby("display_name", as_index=False)
        .agg(valor=("valor_total", "sum"), eventos=("valor_total", "size"))
        .sort_values("valor", ascending=False)
    )
    return grouped


def build_2026_q1_captations() -> pd.DataFrame:
    frame = pd.read_csv(EVENTS_2025_2026Q1_CSV)
    frame["valor_total"] = pd.to_numeric(frame["valor_total"], errors="coerce").fillna(0.0)
    mask = frame["event_type"].eq("emissao") & frame["competencia"].isin(["01/2026", "02/2026", "03/2026"])
    output = frame[mask].copy()
    output = output.sort_values(["FIDC", "competencia", "class_kind", "class_label"])
    return output[["FIDC", "competencia", "class_label", "valor_total"]]


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def set_textbox_text(shape, text: str, font_size: float = 7.0) -> None:
    shape.text = text
    frame = shape.text_frame
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    for paragraph in frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = BLACK


def add_metric_card(slide, x: float, y: float, label: str, value: str, detail: str) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.12), Inches(0.90))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(225, 220, 212)
    tf = card.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = "Arial"
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = BLACK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Arial"
    r2.font.size = Pt(7.2)
    r2.font.bold = True
    r2.font.color.rgb = BLACK
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = detail
    r3.font.name = "Arial"
    r3.font.size = Pt(6.6)
    r3.font.color.rgb = RGBColor(93, 104, 116)


def style_cell(cell, text: str, *, header: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT, shade: bool = False) -> None:
    cell.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.015)
    cell.margin_bottom = Inches(0.015)

    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = BLACK if header else (LIGHT if shade else WHITE)

    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.bold = header
            run.font.size = Pt(6.6 if header else 6.25)
            run.font.color.rgb = WHITE if header else BLACK


def remove_matching_slides(prs: Presentation, needle: str) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for idx in range(len(slide_ids) - 1, -1, -1):
        slide = prs.slides[idx]
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        if any(needle in text for text in texts):
            prs.slides._sldIdLst.remove(slide_ids[idx])


def update_cover(slide, cash_base: float, cash_yield: float) -> None:
    for shape in list(slide.shapes):
        if (
            shape.top >= Inches(4.30)
            and shape.top <= Inches(5.35)
            and shape.left >= Inches(0.45)
            and shape.left <= Inches(10.2)
        ):
            remove_shape(shape)

    metrics = [
        ("2025 bruto", "R$1,277 bi", "gross-up / margem"),
        ("2025 líquido", "R$1,049 bi", "após caixa/LFT"),
        ("Aplicação caixa/LFT", fmt_bi_money(cash_base), f"rend. {fmt_mm(cash_yield)}"),
        ("Saldo médio 2025", "R$8,624 bi", "cotas remuneradas"),
        ("Run-rate 2026", "R$2,09 bi", "jan-abr anualizado"),
    ]
    for idx, (label, value, detail) in enumerate(metrics):
        add_metric_card(slide, 0.62 + idx * 2.34, 4.48, label, value, detail)


def add_price_table(slide, rows: list[list[str]], x, y, w, h) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    col_count = len(rows[0])
    if col_count == 7:
        widths = [1.14, 1.46, 0.57, 0.62, 0.62, 0.83, 0.64]
        numeric_cols = {3, 5, 6}
    elif col_count == 5:
        widths = [1.05, 1.22, 0.86, 0.58, 0.92]
        numeric_cols = {2, 3, 4}
    elif col_count == 6:
        widths = [1.40, 0.92, 0.86, 0.95, 1.08, 0.82]
        numeric_cols = {1, 2, 3, 4, 5}
    elif col_count == 2:
        widths = [1.25, 0.82]
        numeric_cols = {1}
    else:
        widths = [1.0] * col_count
        numeric_cols = set(range(1, col_count))
    scale = float(w) / float(Inches(sum(widths)))
    for col, width in zip(table.columns, widths):
        col.width = int(Inches(width) * scale)
    table.rows[0].height = Inches(0.25)
    for idx in range(1, len(table.rows)):
        table.rows[idx].height = Inches(0.235)
    for r, values in enumerate(rows):
        for c, value in enumerate(values):
            align = PP_ALIGN.RIGHT if c in numeric_cols else PP_ALIGN.LEFT
            style_cell(table.cell(r, c), value, header=(r == 0), align=align, shade=(r > 0 and r % 2 == 0))


def month_label(competencia: str, with_year: bool = False) -> str:
    month, year = str(competencia).split("/")
    labels = {
        "01": "Jan",
        "02": "Fev",
        "03": "Mar",
        "04": "Abr",
        "05": "Mai",
        "06": "Jun",
        "07": "Jul",
        "08": "Ago",
        "09": "Set",
        "10": "Out",
        "11": "Nov",
        "12": "Dez",
    }
    suffix = f"/{year[-2:]}" if with_year else ""
    return f"{labels.get(month, month)}{suffix}"


def add_side_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, amount_cols: set[int] | None = None) -> None:
    amount_cols = amount_cols or set()
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    table.rows[0].height = Inches(0.28)
    row_h = min(0.33, max(0.20, (h - 0.28) / max(len(rows) - 1, 1)))
    for idx in range(1, len(table.rows)):
        table.rows[idx].height = Inches(row_h)

    if len(rows[0]) == 3:
        widths = [1.64, 1.02, 0.72]
        numeric_cols = {1, 2}
    elif len(rows[0]) == 4:
        widths = [1.14, 0.68, 1.72, 1.05]
        numeric_cols = {3}
    else:
        widths = [1.0] * len(rows[0])
        numeric_cols = set(range(1, len(rows[0])))
    scale = float(Inches(w)) / float(Inches(sum(widths)))
    for col, width in zip(table.columns, widths):
        col.width = int(Inches(width) * scale)

    for r, values in enumerate(rows):
        for c, value in enumerate(values):
            align = PP_ALIGN.RIGHT if c in numeric_cols else PP_ALIGN.LEFT
            style_cell(table.cell(r, c), value, header=(r == 0), align=align, shade=(r > 0 and r % 2 == 0))
            if r > 0 and c in amount_cols:
                color = GREEN if not str(value).startswith("(") else RED
                for paragraph in table.cell(r, c).text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = color


def add_chart_caption(slide, text: str, x: float, y: float, w: float) -> None:
    caption = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.18))
    caption.text = text
    for paragraph in caption.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(7.2)
            run.font.bold = True
            run.font.color.rgb = RGBColor(93, 104, 116)


def add_pl_evolution_2025_slide(prs: Presentation, monthly_pl: pd.DataFrame, amort_2025: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    data_2025 = monthly_pl[(monthly_pl["sort_key"] >= 202501) & (monthly_pl["sort_key"] <= 202512)].copy()
    jan_pl = float(data_2025.loc[data_2025["competencia"].eq("01/2025"), "pl"].iloc[0])
    dec_pl = float(data_2025.loc[data_2025["competencia"].eq("12/2025"), "pl"].iloc[0])
    total_amort = float(amort_2025["valor"].sum())
    dec_funds = int(data_2025.loc[data_2025["competencia"].eq("12/2025"), "fundos"].iloc[0])
    add_slide_header(slide, "PL 2025", "Evolução mensal do PL e amortizações reportadas por fundo")
    add_metric_card(slide, 0.78, 1.58, "PL jan/25*", fmt_bi_money(jan_pl), "base IME disponível")
    add_metric_card(slide, 3.17, 1.58, "PL dez/25", fmt_bi_money(dec_pl), f"{dec_funds} FIDCs no mês")
    add_metric_card(slide, 5.39, 1.58, "Amortizações 2025", f"({fmt_bi_money(total_amort)})", "eventos do IME/CVM")

    add_chart_caption(slide, "PL mensal agregado, R$ bi", 0.88, 2.66, 2.5)
    chart_data = CategoryChartData()
    categories = [month_label(label) + ("*" if label == "01/2025" else "") for label in data_2025["competencia"].tolist()]
    chart_data.categories = categories
    chart_data.add_series("PL agregado", (data_2025["pl"] / 1e9).round(3).tolist())
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.82),
        Inches(2.92),
        Inches(6.40),
        Inches(2.85),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 16
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.category_axis.tick_labels.font.size = Pt(7)
    set_series_color(chart.series[0], ORANGE)

    rows = [["FIDC", "Amort. 2025", "Eventos"]]
    for _, row in amort_2025.iterrows():
        rows.append([row["display_name"], f"({fmt_mm(float(row['valor']))})", str(int(row["eventos"]))])
    rows.append(["Total", f"({fmt_bi_money(total_amort)})", str(int(amort_2025["eventos"].sum()))])
    add_side_table(slide, rows, 7.55, 1.70, 4.65, 4.02, amount_cols={1})

    add_note_box(
        slide,
        "* Jan/25 tem universo parcial no cache do IME; dez/25 reflete a base com 10 FIDCs. Amortizações acima são eventos efetivos de amortização reportados no IME/CVM, usados na janela reconciliada do waterfall de PL, não cronograma linear estimado.",
        0.82,
        6.02,
        11.38,
        0.62,
    )
    add_footer(slide, len(prs.slides))


def add_pl_evolution_2026_q1_slide(prs: Presentation, monthly_pl: pd.DataFrame, q1_captations: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    data_2026 = monthly_pl[(monthly_pl["sort_key"] >= 202601) & (monthly_pl["sort_key"] <= 202603)].copy()
    jan_pl = float(data_2026.loc[data_2026["competencia"].eq("01/2026"), "pl"].iloc[0])
    mar_pl = float(data_2026.loc[data_2026["competencia"].eq("03/2026"), "pl"].iloc[0])
    total_cap = float(q1_captations["valor_total"].sum())
    add_slide_header(slide, "PL JAN-MAR/26", "Captação de Bela explica a recomposição do PL no 1T26")
    add_metric_card(slide, 0.78, 1.58, "PL jan/26", fmt_bi_money(jan_pl), "10 FIDCs no IME")
    add_metric_card(slide, 3.17, 1.58, "PL mar/26", fmt_bi_money(mar_pl), "8 FIDCs no IME")
    add_metric_card(slide, 5.39, 1.58, "Captações 1T26", fmt_bi_money(total_cap), "eventos reportados")

    add_chart_caption(slide, "PL mensal agregado, R$ bi", 0.88, 2.66, 2.5)
    chart_data = CategoryChartData()
    chart_data.categories = [month_label(label, with_year=True) for label in data_2026["competencia"].tolist()]
    chart_data.add_series("PL agregado", (data_2026["pl"] / 1e9).round(3).tolist())
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.82),
        Inches(2.92),
        Inches(6.40),
        Inches(2.85),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 18
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.category_axis.tick_labels.font.size = Pt(8)
    set_series_color(chart.series[0], ORANGE)

    rows = [["FIDC", "Mês", "Série / classe", "Captação"]]
    for _, row in q1_captations.iterrows():
        rows.append(
            [
                short_name(row["FIDC"]),
                month_label(row["competencia"], with_year=True),
                str(row["class_label"]).replace("Sênior · ", "Sr ").replace("Mezanino 1 · ", "Mez "),
                fmt_mm(float(row["valor_total"])),
            ]
        )
    rows.append(["Total", "", "", fmt_bi_money(total_cap)])
    add_side_table(slide, rows, 7.55, 1.70, 4.65, 2.20, amount_cols={3})

    add_note_box(
        slide,
        "* A base IME abre Bela em três linhas de classe/série em mar/26: Sênior série 2 e Mezanino séries 3/4. Isso é compatível com captações concentradas em Bela; a consolidação em duas emissões deve ser validada contra os documentos de oferta/tesouraria.",
        7.55,
        4.22,
        4.65,
        1.00,
    )
    add_note_box(
        slide,
        "Leitura: o PL fica praticamente estável em jan-fev e salta em mar/26, quando Bela adiciona R$3,53 bi de cotas no IME. A cor verde marca entrada de capital; o gráfico preserva a paleta Itaú BBA em laranja/preto/cinza.",
        0.82,
        6.02,
        11.38,
        0.62,
    )
    add_footer(slide, len(prs.slides))


def add_slide_header(slide, kicker: str, title_text: str) -> None:
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.43), Inches(0.08), Inches(0.22))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    kicker_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.39), Inches(3.4), Inches(0.25))
    kicker_box.text = kicker
    for paragraph in kicker_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(7.5)
            run.font.bold = True
            run.font.color.rgb = BLACK
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.88), Inches(11.2), Inches(0.72))
    title.text = title_text
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = BLACK


def add_footer(slide, number: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.16), Inches(9.3), Inches(0.16))
    footer.text = "Fonte: TomaConta FIDCs / IME CVM; documentos regulatórios locais; B3/Cetip MediaCDI. Valores em R$ nominais."
    for paragraph in footer.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(5.8)
            run.font.color.rgb = RGBColor(93, 104, 116)

    page = slide.shapes.add_textbox(Inches(12.48), Inches(7.12), Inches(0.35), Inches(0.18))
    page.text = f"{number:02d}"
    for paragraph in page.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(6.5)
            run.font.color.rgb = RGBColor(93, 104, 116)


def add_note_box(slide, text: str, x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(225, 220, 212)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.10)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(8)
    r.font.color.rgb = BLACK


def set_series_color(series, color: RGBColor) -> None:
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.color.rgb = color


def rebuild_waterfall_slide(slide) -> None:
    for shape in list(slide.shapes):
        remove_shape(shape)

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    add_slide_header(slide, "PL WATERFALL", "PL dos fundos cai R$1,53 bi em 2025, puxado por amortizações reportadas")

    steps = [
        ("PL inicial", 14.29681344964, "total"),
        ("Captações", 0.0444, "up"),
        ("Amortizações", -2.41004129267, "down"),
        ("Accrual", 0.83980257679, "up"),
        ("PL final", 12.77097473376, "total"),
    ]
    running = steps[0][1]
    base = [0.0]
    inc = [0.0]
    dec = [0.0]
    totals = [steps[0][1]]
    for label, value, kind in steps[1:-1]:
        previous = running
        running += value
        base.append(min(previous, running))
        inc.append(value if value > 0 else 0.0)
        dec.append(abs(value) if value < 0 else 0.0)
        totals.append(0.0)
    base.append(0.0)
    inc.append(0.0)
    dec.append(0.0)
    totals.append(steps[-1][1])

    chart_data = CategoryChartData()
    chart_data.categories = [item[0] for item in steps]
    chart_data.add_series("Base invisível", base)
    chart_data.add_series("Aumentos", inc)
    chart_data.add_series("Reduções", dec)
    chart_data.add_series("Totais", totals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.78),
        Inches(1.72),
        Inches(7.25),
        Inches(4.10),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 16
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.category_axis.tick_labels.font.size = Pt(7)
    set_series_color(chart.series[0], BG)
    set_series_color(chart.series[1], GREEN)
    set_series_color(chart.series[2], RED)
    set_series_color(chart.series[3], BLACK)

    table_rows = [
        ["Etapa", "R$ bi"],
        ["PL inicial", "14,297"],
        ["+ Captações", "0,044"],
        ["- Amortizações", "(2,410)"],
        ["+ Accrual", "0,840"],
        ["PL final", "12,771"],
    ]
    add_price_table(slide, table_rows, Inches(8.45), Inches(1.82), Inches(3.15), Inches(1.72))
    add_note_box(
        slide,
        "Leitura: PL cai R$1,53 bi no ano; o accrual positivo não compensou R$2,41 bi de amortizações. O gráfico é nativo Office/Excel, com base invisível, aumentos em verde, reduções em vermelho e totais em preto.",
        8.45,
        3.85,
        3.95,
        1.10,
    )
    add_note_box(
        slide,
        "Maior contribuição: Akira II amortizou R$1,15 bi nos eventos IME/CVM que alimentam o waterfall; o detalhamento fundo a fundo está no slide seguinte.",
        0.78,
        5.98,
        7.25,
        0.58,
    )
    add_footer(slide, 4)


def add_structure_slide(prs: Presentation, structure: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    add_slide_header(slide, "ESTRUTURA", "Capital stack agregado: subordinada aparece explicitamente no IME")

    chart_structure = structure[~structure["Classe"].isin(["Total", "Não reconcil."])].copy()
    pivot = chart_structure.pivot(index="Data", columns="Classe", values="Saldo").fillna(0.0) / 1e9
    chart_data = CategoryChartData()
    chart_data.categories = list(pivot.index)
    for label, color in [("Sênior", BLACK), ("Mezanino", ORANGE), ("Subordinada", GREY)]:
        chart_data.add_series(label, pivot.get(label, pd.Series([0.0] * len(pivot), index=pivot.index)).tolist())
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.75),
        Inches(1.78),
        Inches(6.35),
        Inches(4.25),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.category_axis.tick_labels.font.size = Pt(8)
    for series, color in zip(chart.series, [BLACK, ORANGE, GREY]):
        set_series_color(series, color)

    rows = [["Data", "Classe", "Saldo", "FIDCs", "Séries"]]
    for _, row in structure.iterrows():
        fidcs = str(row["FIDCs"])
        series = str(row["Séries"])
        if isinstance(row["FIDCs"], float) and row["FIDCs"].is_integer():
            fidcs = str(int(row["FIDCs"]))
        if isinstance(row["Séries"], float) and row["Séries"].is_integer():
            series = str(int(row["Séries"]))
        rows.append([row["Data"], row["Classe"], fmt_bi_money(float(row["Saldo"])), fidcs, series])
    add_price_table(slide, rows, Inches(7.45), Inches(1.72), Inches(4.65), Inches(3.20))
    add_note_box(
        slide,
        "Correção metodológica: Sr/Mez/Sub agora vêm do PL por classe no IME mais recente de cada fundo, não de principal documental. A diferença pequena contra o PL oficial fica em 'Não reconcil.'.",
        7.45,
        5.15,
        4.65,
        0.95,
    )
    add_footer(slide, len(prs.slides))


def add_amortization_2026_slide(prs: Presentation, amort_df: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    total = amort_df["Total"].sum()
    add_slide_header(slide, "AMORTIZAÇÕES 2026", f"Há R${total / 1e9:.2f} bi de amortizações programadas em 2026, todas em cotas seniores".replace(".", ","))

    chart_data = CategoryChartData()
    chart_data.categories = amort_df["FIDC"].tolist()
    chart_data.add_series("Amortizações 2026", (amort_df["Total"] / 1e6).tolist())
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.72),
        Inches(1.72),
        Inches(6.50),
        Inches(4.65),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.category_axis.tick_labels.font.size = Pt(7)
    set_series_color(chart.series[0], ORANGE)

    rows = [["FIDC", "Série", "Total", "Eventos", "Janela"]]
    for _, row in amort_df.iterrows():
        rows.append([row["FIDC"], row["Série"].replace("1ª série ", "1ª "), fmt_mm(row["Total"]), str(int(row["Eventos"])), f"{row['Primeira']} a {row['Última']}"])
    add_price_table(slide, rows, Inches(7.55), Inches(1.72), Inches(4.80), Inches(3.10))
    top3 = amort_df.head(3)["Total"].sum() / total if total else 0.0
    add_note_box(
        slide,
        f"Top 3 representam {top3:.0%} das amortizações de 2026: Akira I, Big Picture I e Kick Ass I. Big Pictures concentram pagamentos em nov/dez; Akira I e Kick Ass I são mais distribuídos ao longo do ano.",
        7.55,
        5.05,
        4.80,
        1.00,
    )
    add_footer(slide, len(prs.slides))


def add_implied_cdi_slide(prs: Presentation, implied: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    spread_2025 = float(implied.loc[implied["cenario"].str.contains("2025"), "spread_implicito_bruto_aa"].iloc[0])
    spread_2026 = float(implied.loc[implied["cenario"].str.contains("2026"), "spread_implicito_bruto_aa"].iloc[0])
    pct_cdi_2025 = float(implied.loc[implied["cenario"].str.contains("2025"), "pct_cdi_bruto"].iloc[0])
    pct_cdi_2026 = float(implied.loc[implied["cenario"].str.contains("2026"), "pct_cdi_bruto"].iloc[0])
    delta = spread_2026 - spread_2025
    delta_pct_cdi = pct_cdi_2026 - pct_cdi_2025
    add_slide_header(
        slide,
        "CUSTO TESOURARIA",
        "Custo unitário melhora em 2026 nas duas leituras: CDI+ e % do CDI",
    )

    add_metric_card(slide, 0.78, 1.72, "2025 unitário", f"CDI+{fmt_pct(spread_2025)}", f"{fmt_pct_one(pct_cdi_2025)} do CDI")
    add_metric_card(slide, 3.17, 1.72, "2026 unitário", f"CDI+{fmt_pct(spread_2026)}", f"{fmt_pct_one(pct_cdi_2026)} do CDI")
    add_metric_card(slide, 5.56, 1.72, "Variação", fmt_bps(delta), f"{fmt_pp(delta_pct_cdi)} do CDI")

    axis_label = slide.shapes.add_textbox(Inches(0.85), Inches(2.83), Inches(2.1), Inches(0.18))
    axis_label.text = "CDI+ bruto, em bps"
    for paragraph in axis_label.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(7.2)
            run.font.bold = True
            run.font.color.rgb = RGBColor(93, 104, 116)

    chart_data = CategoryChartData()
    chart_data.categories = ["2025", "2026 run-rate"]
    chart_data.add_series("CDI+ implícito bruto (bps)", [spread_2025 * 10000.0, spread_2026 * 10000.0])
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.85),
        Inches(3.05),
        Inches(5.85),
        Inches(2.65),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 160
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.category_axis.tick_labels.font.size = Pt(8)
    set_series_color(chart.series[0], ORANGE)

    table_rows = [["Cenário", "Custo bruto", "CDI puro", "% CDI", "Spread R$", "CDI+ impl."]]
    for _, row in implied.iterrows():
        table_rows.append(
            [
                "2025" if "2025" in row["cenario"] else "2026 run-rate",
                fmt_bi_money(float(row["custo_bruto_anualizado"])),
                fmt_bi_money(float(row["custo_cdi_puro_anualizado"])),
                fmt_pct_one(float(row["pct_cdi_bruto"])),
                fmt_mm(float(row["parcela_spread_anualizada"])),
                fmt_pct(float(row["spread_implicito_bruto_aa"])),
            ]
        )
    add_price_table(slide, table_rows, Inches(7.35), Inches(1.72), Inches(5.10), Inches(1.35))

    add_note_box(
        slide,
        "Memória CDI+: resolve s com CDI_mês composto e spread aditivo por dias úteis. Memória % CDI: resolve k em Custo = soma Saldo_i x [(1+CDI_trecho)^k - 1].",
        7.35,
        3.40,
        5.10,
        1.08,
    )
    add_note_box(
        slide,
        "Leitura para comitê: o custo absoluto sobe porque a base utilizada cresce, especialmente com Bela, mas o preço unitário bruto melhora levemente. Isso sugere gestão eficiente de spread, ainda dependente de benchmark externo e confirmação da Tesouraria.",
        7.35,
        4.78,
        5.10,
        1.10,
    )
    add_note_box(
        slide,
        f"Resultado: % CDI bruto cai de {fmt_pct_one(pct_cdi_2025)} para {fmt_pct_one(pct_cdi_2026)}; CDI+ bruto cai de {fmt_pct(spread_2025)} para {fmt_pct(spread_2026)}.",
        0.85,
        5.98,
        5.85,
        0.58,
    )
    add_footer(slide, len(prs.slides))


def move_last_slide_to(prs: Presentation, new_index: int) -> None:
    slide_ids = prs.slides._sldIdLst
    last = slide_ids[-1]
    slide_ids.remove(last)
    slide_ids.insert(new_index, last)


def refresh_page_numbers(prs: Presentation) -> None:
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            if shape.left > Inches(11.8) and shape.top > Inches(6.8) and shape.text.strip().isdigit():
                shape.text = f"{idx:02d}"
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.RIGHT
                    for run in paragraph.runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(6.5)
                        run.font.color.rgb = RGBColor(93, 104, 116)


def add_series_price_slide(prs: Presentation, rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.43), Inches(0.08), Inches(0.22))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    kicker = slide.shapes.add_textbox(Inches(0.72), Inches(0.39), Inches(2.5), Inches(0.25))
    kicker.text = "PREÇO POR SÉRIE"
    for paragraph in kicker.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(7.5)
            run.font.bold = True
            run.font.color.rgb = BLACK

    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.88), Inches(11.2), Inches(0.55))
    title.text = "CDI+ exato usado no motor, incluindo preços inputados manualmente"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = BLACK

    midpoint = (len(rows) + 1) // 2
    left_rows = [rows[0]] + rows[1:midpoint]
    right_rows = [rows[0]] + rows[midpoint:]
    add_price_table(slide, left_rows, Inches(0.55), Inches(1.62), Inches(6.05), Inches(4.75))
    add_price_table(slide, right_rows, Inches(6.78), Inches(1.62), Inches(5.95), Inches(4.75))

    note = slide.shapes.add_textbox(Inches(0.68), Inches(6.36), Inches(10.4), Inches(0.18))
    note.text = "Fonte Manual = input externo usado quando o documento baixado não trazia CDI+; demais preços vieram de regulamento/ata/oferta curados."
    for paragraph in note.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(6.7)
            run.font.color.rgb = RGBColor(93, 104, 116)

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.16), Inches(9.3), Inches(0.16))
    footer.text = "Fonte: TomaConta FIDCs / IME CVM; documentos regulatórios locais; B3/Cetip MediaCDI. Valores em R$ nominais."
    for paragraph in footer.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(5.8)
            run.font.color.rgb = RGBColor(93, 104, 116)

    page = slide.shapes.add_textbox(Inches(12.48), Inches(7.12), Inches(0.35), Inches(0.18))
    page.text = f"{len(prs.slides):02d}"
    for paragraph in page.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(6.5)
            run.font.color.rgb = RGBColor(93, 104, 116)


def update_deck(
    path: Path,
    rows: list[list[str]],
    series_rows: list[list[str]],
    cash_base: float,
    cash_yield: float,
    structure: pd.DataFrame,
    amort_df: pd.DataFrame,
    implied: pd.DataFrame,
    monthly_pl: pd.DataFrame,
    amort_2025: pd.DataFrame,
    q1_captations: pd.DataFrame,
) -> None:
    prs = Presentation(path)
    remove_matching_slides(prs, "CDI+ exato usado no motor")
    remove_matching_slides(prs, "Capital stack agregado")
    remove_matching_slides(prs, "amortizações programadas em 2026")
    remove_matching_slides(prs, "CDI+ implícito bruto")
    remove_matching_slides(prs, "Evolução mensal do PL e amortizações")
    remove_matching_slides(prs, "Captação de Bela explica")
    update_cover(prs.slides[0], cash_base, cash_yield)
    rebuild_waterfall_slide(prs.slides[3])
    slide = prs.slides[2]

    for shape in list(slide.shapes):
        if getattr(shape, "has_table", False):
            remove_shape(shape)

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text.startswith("PI, A.I. e Bela concentram"):
                shape.text = "PI, A.I. e Bela concentram a despesa; custo contratado fica em CDI+0,75% a CDI+4,99%"
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(22)
                        run.font.bold = True
                        run.font.color.rgb = BLACK
            elif text == "Top 3 = 63% do custo bruto" or text.startswith("Tabela cobre todos"):
                set_textbox_text(shape, "Tabela cobre todos os FIDCs ativos: CDI+ ponderado pelo saldo médio e faixa por tranche.")
            elif text.startswith("CDI mensal B3 elevou"):
                set_textbox_text(shape, "Top 3 = 63% do custo bruto; PI e A.I. combinam senior barato com mezzanine mais caro.")
            elif text.startswith("Run-rate 2026 antigo"):
                set_textbox_text(shape, "Run-rate 2026 antigo de R$2,2 bi era outro período/base; não é 2025.")

    table_shape = slide.shapes.add_table(
        len(rows),
        len(rows[0]),
        Inches(7.33),
        Inches(1.72),
        Inches(4.80),
        Inches(2.98),
    )
    table = table_shape.table
    widths = [1.48, 0.86, 0.78, 0.98, 0.70]
    for col, width in zip(table.columns, widths):
        col.width = Inches(width)
    table.rows[0].height = Inches(0.24)
    for idx in range(1, len(table.rows)):
        table.rows[idx].height = Inches(0.245)

    for r, values in enumerate(rows):
        for c, value in enumerate(values):
            align = PP_ALIGN.RIGHT if c in (1, 2, 4) else PP_ALIGN.LEFT
            style_cell(table.cell(r, c), value, header=(r == 0), align=align, shade=(r > 0 and r % 2 == 0))

    add_pl_evolution_2025_slide(prs, monthly_pl, amort_2025)
    move_last_slide_to(prs, 4)
    add_pl_evolution_2026_q1_slide(prs, monthly_pl, q1_captations)
    move_last_slide_to(prs, 5)
    add_structure_slide(prs, structure)
    move_last_slide_to(prs, 6)
    add_amortization_2026_slide(prs, amort_df)
    move_last_slide_to(prs, 7)
    add_implied_cdi_slide(prs, implied)
    move_last_slide_to(prs, 10)
    add_series_price_slide(prs, series_rows)
    refresh_page_numbers(prs)
    prs.save(path)


def main() -> None:
    rows = build_table_rows()
    series_rows = build_series_rows()
    cash_base, cash_yield = cash_lft_totals()
    lines = funding_lines()
    structure = build_structure_rows(lines)
    amort_df = build_amortization_2026_rows(lines)
    implied = build_implied_cdi_rows()
    monthly_pl = build_monthly_pl_summary()
    amort_2025 = build_2025_amortization_by_fund()
    q1_captations = build_2026_q1_captations()
    backup = DECK_MAIN.with_name("cloudwalk-fidc-custo-financeiro-2025-pre-custo-fidcs.pptx")
    if not backup.exists():
        shutil.copy2(DECK_MAIN, backup)

    for path in [DECK_MAIN, DECK_CONDENSED, DECK_OUTPUT]:
        shutil.copy2(backup, path)
        update_deck(
            path,
            rows,
            series_rows,
            cash_base,
            cash_yield,
            structure,
            amort_df,
            implied,
            monthly_pl,
            amort_2025,
            q1_captations,
        )


if __name__ == "__main__":
    main()
