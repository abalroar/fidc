# -*- coding: utf-8 -*-
from __future__ import annotations

import json
import math
import os
import sys
from collections import defaultdict
from datetime import date
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/Users/matheusjprates/fidc")
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from services.cloudwalk_financial_cost import (
    load_amortization_convention_overrides,
    load_funding_lines,
    load_spread_overrides,
)

OUT = ROOT / "reports/cloudwalk_bba_deck/reverse_target_output/cloudwalk-fidcs-bba-shortdeck.pptx"
ROOT_COPY = ROOT / "Cloudwalk_FIDC_BBA_ShortDeck.pptx"
PROMPT_OUT = ROOT / "reports/cloudwalk_bba_deck/reverse_target_output/cloudwalk-bba-style-prompt.md"

DATA_DIR = ROOT / "reports/cloudwalk_financial_cost_deck_data"
CONFIG_JSON = ROOT / "config/cloudwalk_financial_cost_inputs.json"
EMISSIONS_CSV = ROOT / "data/regulatory_profiles/cloudwalk_cotas_emissoes_pagamentos.csv"

W = Inches(13.333333)
H = Inches(7.5)

BLACK = RGBColor(14, 18, 27)
TEXT = RGBColor(28, 32, 39)
GREY = RGBColor(104, 111, 122)
LIGHT = RGBColor(244, 245, 247)
GRID = RGBColor(202, 206, 214)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(232, 92, 0)
GREEN = RGBColor(29, 143, 84)
RED = RGBColor(174, 45, 43)
TAUPE = RGBColor(247, 246, 242)


def short_name(name: str) -> str:
    upper = str(name).upper()
    pairs = [
        ("Bela", "BELA"),
        ("PI", "PI FUNDO"),
        ("A.I.", "A.I."),
        ("Akira II", "AKIRA II"),
        ("Akira I", "AKIRA I"),
        ("Kick Ass I", "KICK ASS I"),
        ("Big Picture III", "BIG PICTURE III"),
        ("Big Picture II", "BIG PICTURE II"),
        ("Big Picture IV", "BIG PICTURE IV"),
        ("Big Picture I", "BIG PICTURE I"),
    ]
    for label, token in pairs:
        if token in upper:
            return label
    return str(name)[:24]


def fmt_bi(value: float, digits: int = 1, money: bool = True) -> str:
    txt = f"{value / 1e9:.{digits}f}".replace(".", ",")
    return f"R${txt} bi" if money else f"{txt} bi"


def fmt_mm(value: float, digits: int = 0, money: bool = True) -> str:
    txt = f"{value / 1e6:.{digits}f}".replace(".", ",")
    return f"R${txt} mi" if money else f"{txt} mi"


def fmt_pct(value: float, digits: int = 2) -> str:
    return f"{value * 100:.{digits}f}%".replace(".", ",")


def box(slide, x, y, w, h, radius=True, fill=WHITE, line=BLACK, lw=0.65):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(lw)
    if radius and hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.035
    return shape


def text(slide, x, y, w, h, value, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.text = str(value)
    tf = shape.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = valign
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return shape


def rich_title(slide, title_left: str, title_right: str):
    logo = box(slide, 0.47, 0.36, 0.42, 0.30, radius=True, fill=BLACK, line=BLACK, lw=0)
    logo.text = "itaú BBA"
    for p in logo.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(5.6)
            r.font.bold = True
            r.font.color.rgb = WHITE
    title = slide.shapes.add_textbox(Inches(1.08), Inches(0.30), Inches(11.5), Inches(0.42))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_left
    r.font.name = "Arial"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = BLACK
    r2 = p.add_run()
    r2.text = " | " + title_right
    r2.font.name = "Arial"
    r2.font.size = Pt(17)
    r2.font.bold = False
    r2.font.color.rgb = TEXT


def footer(slide, page: int):
    text(slide, 0.50, 7.04, 8.9, 0.12, "Fonte: TomaConta FIDCs / IME CVM; documentos regulatórios locais; B3/Cetip MediaCDI. Valores em R$ nominais.", 5.2, color=GREY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.47), Inches(7.20), Inches(11.95), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    text(slide, 0.50, 7.23, 2.1, 0.18, "CloudWalk |", 8.5, color=GREY)
    text(slide, 0.50, 7.39, 2.4, 0.10, "Corporativo | Interno", 5.8, bold=True, color=TEXT)
    logo = box(slide, 11.95, 7.27, 0.47, 0.25, radius=True, fill=BLACK, line=BLACK, lw=0)
    logo.text = "itaú BBA"
    for p in logo.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(5.3)
            r.font.bold = True
            r.font.color.rgb = WHITE
    text(slide, 12.55, 7.27, 0.28, 0.14, f"{page:02d}", 6.3, color=GREY, align=PP_ALIGN.RIGHT)


def panel(slide, x, y, w, h, title):
    box(slide, x, y, w, h, radius=True, fill=TAUPE, line=BLACK, lw=0.55)
    text(slide, x + 0.20, y - 0.06, w * 0.55, 0.25, title, 13, bold=True, color=BLACK)
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + min(w * 0.34, 2.4)), Inches(y + 0.02), Inches(max(w - min(w * 0.34, 2.4) - 0.25, 0.1)), Inches(0.006))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLACK
    rule.line.fill.background()


def style_cell(cell, header=False, total=False, negative=False, positive=False, align=PP_ALIGN.LEFT, size=8.8):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLACK if (header or total) else TAUPE
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(size)
            r.font.bold = header or total
            if header or total:
                r.font.color.rgb = WHITE
            elif negative:
                r.font.color.rgb = RED
            elif positive:
                r.font.color.rgb = GREEN
            else:
                r.font.color.rgb = TEXT


def add_table(slide, rows, x, y, w, h, widths=None, numeric_cols=None, negative_cols=None, positive_cols=None, size=8.4):
    numeric_cols = numeric_cols or set()
    negative_cols = negative_cols or set()
    positive_cols = positive_cols or set()
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        scale = w / sum(widths)
        for idx, col_w in enumerate(widths):
            table.columns[idx].width = Inches(col_w * scale)
    header_h = min(0.34, h / max(len(rows), 1))
    table.rows[0].height = Inches(header_h)
    if len(rows) > 1:
        row_h = max(0.16, (h - header_h) / (len(rows) - 1))
        for r in range(1, len(rows)):
            table.rows[r].height = Inches(row_h)
    for r, row in enumerate(rows):
        is_total = r == len(rows) - 1 and str(row[0]).lower() in {"total", "saldo total"}
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            style_cell(
                cell,
                header=(r == 0),
                total=is_total,
                negative=(r > 0 and c in negative_cols),
                positive=(r > 0 and c in positive_cols),
                align=PP_ALIGN.RIGHT if c in numeric_cols else PP_ALIGN.LEFT,
                size=size,
            )
    return table_shape


def set_series(series, color):
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.color.rgb = color


def clean_chart(chart, has_legend=False):
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    try:
        chart.value_axis.major_gridlines.format.line.color.rgb = GRID
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.category_axis.tick_labels.font.size = Pt(7)
    except Exception:
        pass


def load_metrics():
    metrics = json.loads((DATA_DIR / "deck_metrics.json").read_text())
    line = pd.read_csv(DATA_DIR / "line_2025.csv")
    included = line[line["included_in_cost"] == True].copy()
    included["short"] = included["fund_name"].map(short_name)

    cost_rows = []
    for fund, group in included.groupby("short"):
        saldo = float(group["saldo_medio_programado"].sum())
        custo = float(group["custo_programado_bruto"].sum())
        spread = float((group["spread_cdi_plus_aa"] * group["saldo_medio_programado"]).sum() / saldo)
        spreads = sorted(float(x) for x in group["spread_cdi_plus_aa"].dropna().unique())
        faixa = fmt_pct(spread) if len(spreads) <= 1 else f"{fmt_pct(spreads[0])}-{fmt_pct(spreads[-1])}"
        cost_rows.append((fund, saldo, spread, faixa, custo))
    cost_rows.sort(key=lambda row: row[4], reverse=True)

    amort_events = pd.read_csv(DATA_DIR / "pl_waterfall_amortization_events_2025.csv")
    amort_events["valor_total"] = pd.to_numeric(amort_events["valor_total"], errors="coerce").fillna(0.0)
    amort_events["short"] = amort_events["short_name"].map(short_name)
    amort_2025 = amort_events.groupby("short", as_index=False)["valor_total"].sum().sort_values("valor_total", ascending=False)

    events = pd.read_csv(DATA_DIR / "events_2025_2026q1.csv")
    events["valor_total"] = pd.to_numeric(events["valor_total"], errors="coerce").fillna(0.0)
    q1_caps = events[
        events["event_type"].eq("emissao") & events["competencia"].isin(["01/2026", "02/2026", "03/2026"])
    ].copy()
    q1_caps["short"] = q1_caps["FIDC"].map(short_name)

    stack = pd.read_csv(DATA_DIR / "capital_stack_ime_by_class.csv")
    stack["PL classe"] = pd.to_numeric(stack["PL classe"], errors="coerce").fillna(0.0)
    stack_pivot = stack.groupby(["Data", "class_macro"])["PL classe"].sum().unstack(fill_value=0.0)
    official = pd.read_csv(DATA_DIR / "capital_stack_official_pl.csv").groupby("Data")["PL oficial"].sum()

    dur_month = pd.read_csv(DATA_DIR / "duration_monthly_2026_jan_apr.csv")
    dur_fund = pd.read_csv(DATA_DIR / "duration_by_fund_2026_jan_apr.csv")
    dur_fund["short"] = dur_fund["fund_name"].map(short_name)
    dur_fund = dur_fund.sort_values("avg_stock_recebiveis_a_vencer", ascending=False)

    monthly_2026 = pd.read_csv(ROOT / "reports/cloudwalk_financial_cost/cloudwalk_financial_cost_monthly.csv")
    monthly_2026["short"] = monthly_2026["fund_name"].map(short_name)
    ytd = monthly_2026[monthly_2026["mes"].between("2026-01", "2026-04")].copy()
    month_raw = ytd.groupby("mes")["custo_programado_bruto"].sum()
    fund_raw = ytd.groupby("short")["custo_programado_bruto"].sum().sort_values(ascending=False)
    runrate = pd.read_csv(DATA_DIR / "prognostico_2026_runrate.csv")
    jan_apr = runrate[runrate["periodo"].eq("jan-abr/26")].iloc[0]
    scale = float(jan_apr["custo_bruto"]) / float(month_raw.sum())
    month_cost = (month_raw * scale).reindex(["2026-01", "2026-02", "2026-03", "2026-04"])
    fund_cost = fund_raw * scale

    implied = pd.read_csv(DATA_DIR / "implied_cdi_plus.csv")

    lines = load_funding_lines(
        EMISSIONS_CSV,
        spread_overrides=load_spread_overrides(CONFIG_JSON),
        amortization_convention_overrides=load_amortization_convention_overrides(CONFIG_JSON),
    )
    amort_2026 = defaultdict(float)
    for line in lines:
        if not line.included:
            continue
        for event_date, amount in line.amortizations:
            if date(2026, 1, 1) <= event_date <= date(2026, 12, 31):
                amort_2026[short_name(line.fund_name)] += amount
    amort_2026 = sorted(amort_2026.items(), key=lambda item: item[1])

    return {
        "metrics": metrics,
        "cost_rows": cost_rows,
        "amort_2025": amort_2025,
        "q1_caps": q1_caps,
        "stack": stack_pivot,
        "official": official,
        "dur_month": dur_month,
        "dur_fund": dur_fund,
        "month_cost": month_cost,
        "fund_cost": fund_cost,
        "implied": implied,
        "amort_2026": amort_2026,
        "runrate": jan_apr,
    }


def add_horizontal_bar(slide, data, x, y, w, h, title):
    chart_data = CategoryChartData()
    chart_data.categories = [item[0] for item in data]
    chart_data.add_series(title, [round(item[1]) for item in data])
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
    chart = chart_shape.chart
    clean_chart(chart, False)
    chart.value_axis.minimum_scale = 0
    max_value = max(item[1] for item in data) if data else 1
    chart.value_axis.maximum_scale = math.ceil(max_value / 100) * 100
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.number_format = "0"
    chart.plots[0].data_labels.font.size = Pt(8)
    set_series(chart.series[0], BLACK)
    return chart_shape


def add_column_chart(slide, labels, values, x, y, w, h, color=BLACK, max_scale=None, label_format="0"):
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("Valor", values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
    chart = chart_shape.chart
    clean_chart(chart, False)
    chart.value_axis.minimum_scale = 0
    if max_scale:
        chart.value_axis.maximum_scale = max_scale
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.number_format = label_format
    chart.plots[0].data_labels.font.size = Pt(8)
    set_series(chart.series[0], color)
    return chart_shape


def slide_methodology(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rich_title(slide, "Despesa Financeira FIDCs", "estimativa recompõe R$1,0 bi líquido e ~R$26 bi de TPV em 2025")
    panel(slide, 0.55, 1.05, 6.05, 5.92, "Motor de cálculo")
    panel(slide, 6.90, 1.05, 5.85, 5.92, "Quebra por FIDC")

    text(slide, 0.82, 1.38, 5.25, 0.58, "A DRE da IP não explicita a despesa dos FIDCs;\no estudo recompõe o custo econômico.", 13.5, color=TEXT)
    text(slide, 1.10, 2.15, 4.70, 0.26, "Custo FIDC = custo bruto contratado - carry estimado do caixa/LFT", 11.5, bold=True)

    steps = [
        ("Informe Mensal CVM", "volume, classe e preço de cota"),
        ("Saldos mensais", "saldo médio por fundo/série"),
        ("Tipo de cota", "Sr/Mez entram; sub já retorna na IP"),
        ("Caixa / LFT", "carry reduz custo financeiro líquido"),
        ("CDI B3", "mensal composto + dias úteis"),
    ]
    x0, y0 = 0.95, 2.78
    for i, (label, detail) in enumerate(steps):
        yy = y0 + i * 0.62
        box(slide, x0, yy, 1.65, 0.44, radius=True, fill=WHITE, line=GRID, lw=0.5)
        text(slide, x0 + 0.08, yy + 0.07, 1.49, 0.13, label, 8.1, bold=True, align=PP_ALIGN.CENTER)
        text(slide, x0 + 0.08, yy + 0.22, 1.49, 0.10, detail, 6.5, align=PP_ALIGN.CENTER, color=GREY)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(x0 + 0.73), Inches(yy + 0.43), Inches(0.12), Inches(0.18))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ORANGE
            arrow.line.fill.background()
    comments = [
        "Curva mensal evita repetir um saldo anual único.",
        "Amortizações e captações entram por data/evento.",
        "Spread CDI+ é aplicado por série remunerada.",
        "Caixa aplicado em LFTs atenua a despesa.",
        "Resultado combina CDI realizado e estoque exposto.",
    ]
    for i, value in enumerate(comments):
        text(slide, 3.06, y0 + i * 0.62 + 0.09, 2.90, 0.18, value, 9.3, color=TEXT)

    metrics = data["metrics"]["cost_summary"][2]
    text(slide, 0.82, 6.12, 1.2, 0.18, "Resultado", 10.5, bold=True)
    text(slide, 0.82, 6.43, 5.20, 0.20, f"• Custo bruto estimado: {fmt_bi(metrics['despesa_financeira_bruta'], 3)}", 10.6, bold=True)
    text(slide, 0.82, 6.68, 5.20, 0.20, f"• Custo líquido após caixa/LFT: {fmt_bi(metrics['despesa_financeira_liquida'], 3)}", 10.6, bold=True)

    rows = [["FIDC", "Saldo méd.", "CDI+ Pond.", "Faixa CDI+", "Custo 25"]]
    for fund, saldo, spread, faixa, custo in data["cost_rows"]:
        rows.append([fund, fmt_bi(saldo, 2, money=False), fmt_pct(spread), faixa, fmt_mm(custo)])
    rows.append(["Total", fmt_bi(metrics["saldo_base"], 2, money=False), "", "", fmt_bi(metrics["despesa_financeira_bruta"], 3)])
    add_table(slide, rows, 7.02, 1.42, 5.52, 3.18, widths=[1.05, 0.92, 0.86, 1.22, 0.92], numeric_cols={1, 2, 3, 4}, size=8.0)

    dur = data["metrics"]["duration_2025"]
    text(slide, 7.15, 5.25, 2.5, 0.22, "TPV antecipado via FIDCs", 11.2, bold=True)
    text(slide, 8.05, 5.70, 3.80, 0.22, "(365 / prazo médio) x estoque médio de recebíveis", 10.0, color=TEXT, align=PP_ALIGN.CENTER)
    text(slide, 7.18, 6.18, 5.0, 0.24, f"• Duration estimada: ~{dur['duration_flow_implied']:.0f} dias", 10.5, bold=True)
    text(slide, 7.18, 6.52, 5.0, 0.24, f"• TPV antecipado 2025: ~{fmt_bi(dur['tpv_period'], 0)}", 10.5, bold=True)
    footer(slide, 1)


def slide_evolution(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rich_title(slide, "Evolução FIDCs", "PL cai R$1,5 bi em 2025 e recompõe para R$16,1 bi após Bela")
    panel(slide, 0.55, 1.02, 12.20, 2.70, "PL 2025 e amortizações")
    panel(slide, 0.55, 3.94, 7.65, 2.83, "Capital stack agregado")
    panel(slide, 8.35, 3.94, 4.40, 2.83, "Captações 2026")

    wf = data["metrics"]["pl_waterfall_summary"][0]
    steps = [
        ("PL inicial", wf["pl_inicial"] / 1e9, "total"),
        ("Captações", wf["captacoes"] / 1e9, "up"),
        ("Amortizações", -wf["amortizacoes"] / 1e9, "down"),
        ("Accrual", wf["accrual_rentabilidade_residual"] / 1e9, "up"),
        ("PL final", wf["pl_final"] / 1e9, "total"),
    ]
    running = steps[0][1]
    base, inc, dec, total = [0], [0], [0], [steps[0][1]]
    for _, value, kind in steps[1:-1]:
        previous = running
        running += value
        base.append(min(previous, running))
        inc.append(value if value > 0 else 0)
        dec.append(abs(value) if value < 0 else 0)
        total.append(0)
    base.append(0)
    inc.append(0)
    dec.append(0)
    total.append(steps[-1][1])
    chart_data = CategoryChartData()
    chart_data.categories = [s[0] for s in steps]
    chart_data.add_series("Base", base)
    chart_data.add_series("Aumentos", inc)
    chart_data.add_series("Reduções", dec)
    chart_data.add_series("Totais", total)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.92), Inches(1.53), Inches(6.30), Inches(1.88), chart_data)
    chart = chart_shape.chart
    clean_chart(chart, True)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 16
    set_series(chart.series[0], TAUPE)
    set_series(chart.series[1], GREEN)
    set_series(chart.series[2], RED)
    set_series(chart.series[3], BLACK)
    for x, y, label, color in [
        (1.55, 2.45, "14,3", WHITE),
        (3.70, 1.45, "2,4", RED),
        (5.05, 1.50, "0,84", BLACK),
        (6.55, 2.53, "12,8", WHITE),
    ]:
        text(slide, x, y, 0.55, 0.18, label, 11, bold=True, color=color, align=PP_ALIGN.CENTER)

    amort_rows = [["FIDC", "Amort. 2025"]]
    for _, row in data["amort_2025"].iterrows():
        amort_rows.append([row["short"], f"({fmt_mm(float(row['valor_total']))})"])
    amort_rows.append(["Total", f"({fmt_bi(float(data['amort_2025']['valor_total'].sum()), 2)})"])
    add_table(slide, amort_rows, 8.55, 1.34, 3.82, 2.15, widths=[1.8, 1.6], numeric_cols={1}, negative_cols={1}, size=7.8)

    stack = data["stack"]
    labels = ["31/12/25", "23/05/26"]
    senior = [stack.loc["31/12/25", "senior"] / 1e9, stack.loc["Hoje (23/05/26)", "senior"] / 1e9]
    mezz = [stack.loc["31/12/25", "mezzanino"] / 1e9, stack.loc["Hoje (23/05/26)", "mezzanino"] / 1e9]
    sub = [stack.loc["31/12/25", "subordinada"] / 1e9, stack.loc["Hoje (23/05/26)", "subordinada"] / 1e9]
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("Sênior", senior)
    chart_data.add_series("Mezanino", mezz)
    chart_data.add_series("Subordinada", sub)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(1.10), Inches(4.35), Inches(6.35), Inches(2.15), chart_data)
    chart = chart_shape.chart
    clean_chart(chart, True)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 20
    for series, color in zip(chart.series, [BLACK, ORANGE, GREY]):
        set_series(series, color)
    totals = [float(data["official"].loc["31/12/25"]) / 1e9, float(data["official"].loc["Hoje (23/05/26)"]) / 1e9]
    text(slide, 2.72, 4.12, 0.70, 0.18, f"{totals[0]:.1f}".replace(".", ","), 12, bold=True, align=PP_ALIGN.CENTER)
    text(slide, 6.30, 4.12, 0.70, 0.18, f"{totals[1]:.1f}".replace(".", ","), 12, bold=True, align=PP_ALIGN.CENTER)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(4.05), Inches(4.26), Inches(1.98), Inches(0.20))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN
    arrow.line.fill.background()
    text(slide, 6.96, 4.48, 0.96, 0.40, "Subordinação\n~3,1%", 7.8, color=TEXT, align=PP_ALIGN.CENTER)

    cap_rows = [["FIDC", "Mês", "Série / Classe", "Captação"]]
    q1 = data["q1_caps"].copy()
    q1["sort_class"] = q1["class_macro"].map({"senior": 0, "mezzanino": 1, "subordinada": 2}).fillna(9)
    q1 = q1.sort_values(["competencia", "sort_class", "valor_total"], ascending=[True, True, False])
    for _, row in q1.iterrows():
        cap_rows.append([
            row["short"],
            row["competencia"].replace("/2026", "-26"),
            str(row["class_label"]).replace("Sênior · ", "Sr ").replace("Mezanino 1 · ", "Mez "),
            fmt_mm(float(row["valor_total"])),
        ])
    cap_rows.append(["Total", "", "", fmt_bi(float(q1["valor_total"].sum()), 2)])
    add_table(slide, cap_rows, 8.55, 4.39, 3.84, 2.13, widths=[0.75, 0.55, 1.10, 0.95], numeric_cols={3}, positive_cols={3}, size=8.0)
    footer(slide, 2)


def slide_2026(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    implied = data["implied"]
    spread25 = float(implied.loc[implied["cenario"].str.contains("2025"), "spread_implicito_bruto_aa"].iloc[0])
    spread26 = float(implied.loc[implied["cenario"].str.contains("2026"), "spread_implicito_bruto_aa"].iloc[0])
    runrate = data["runrate"]
    rich_title(
        slide,
        "Estimativa 2026",
        f"run-rate bruto cai para CDI+{fmt_pct(spread26)} e anualiza perto de {fmt_bi(float(runrate['custo_bruto_anualizado']), 1)}",
    )
    panel(slide, 0.55, 1.02, 5.70, 2.45, f"Amortizações 2026: {fmt_bi(sum(v for _, v in data['amort_2026']), 1)} programados")
    panel(slide, 6.50, 1.02, 6.25, 2.45, "Duration 2026")
    panel(slide, 0.55, 3.70, 5.70, 3.02, f"Custo financeiro bruto YTD: ~{fmt_mm(float(runrate['custo_bruto']))}")
    panel(slide, 6.50, 3.70, 6.25, 3.02, "CDI+ implícito bruto")

    amort = [(name, value / 1e6) for name, value in data["amort_2026"]]
    add_horizontal_bar(slide, amort, 0.86, 1.42, 4.95, 1.72, "R$ mi")

    dm = data["dur_month"]
    add_column_chart(
        slide,
        ["01", "02", "03", "04"],
        [round(x) for x in dm["duration_days"].tolist()],
        7.02,
        1.35,
        4.85,
        0.92,
        color=BLACK,
        max_scale=150,
    )
    dur_rows = [["FIDC", "Jan-abr/26", "Estoque méd."]]
    for _, row in data["dur_fund"].head(6).iterrows():
        dur_rows.append([row["short"], f"{row['duration_days_stock_day_weighted']:.0f}d", fmt_bi(float(row["avg_stock_recebiveis_a_vencer"]), 1, money=False)])
    add_table(slide, dur_rows, 6.80, 2.34, 5.60, 1.02, widths=[2.55, 1.25, 1.28], numeric_cols={1, 2}, size=7.1)

    months = ["Jan", "Fev", "Mar", "Abr"]
    add_column_chart(
        slide,
        months,
        [round(v / 1e6) for v in data["month_cost"].tolist()],
        1.02,
        4.12,
        4.65,
        1.10,
        color=BLACK,
        max_scale=260,
    )
    top = data["fund_cost"].sort_values(ascending=False)
    demais = top.iloc[4:].sum()
    rows = [["FIDC", "YTD Abr-26", "Anualizado"]]
    for fund, value in top.head(4).items():
        rows.append([fund, fmt_mm(float(value)), fmt_bi(float(value) * 3, 2)])
    rows.append(["Demais", fmt_mm(float(demais)), fmt_bi(float(demais) * 3, 2)])
    rows.append(["Total", fmt_mm(float(top.sum())), fmt_bi(float(top.sum()) * 3, 1)])
    add_table(slide, rows, 1.02, 5.55, 4.65, 0.95, widths=[1.65, 1.05, 1.15], numeric_cols={1, 2}, size=7.3)

    add_column_chart(
        slide,
        ["2025", "2026 run-rate"],
        [round(spread25 * 10000, 1), round(spread26 * 10000, 1)],
        7.20,
        4.40,
        4.70,
        1.72,
        color=BLACK,
        max_scale=160,
        label_format="0.0",
    )
    text(slide, 7.42, 4.08, 2.0, 0.18, f"2025FY: CDI+{fmt_pct(spread25)}", 10.2, bold=True)
    text(slide, 10.02, 4.08, 2.2, 0.18, f"2026 YTD: CDI+{fmt_pct(spread26)}", 10.2, bold=True)
    text(slide, 7.40, 6.25, 4.78, 0.22, "Cálculo pondera CDI mensal composto e spread por saldo médio exposto em cada trecho.", 8.3, color=TEXT)
    footer(slide, 3)


def write_prompt():
    prompt = """# Prompt reverso para gerar o short deck Itaú BBA

Você é um analista sênior de crédito estruturado e designer de apresentações executivas. Gere um PPTX 16:9, em português, com no máximo 3 slides, usando apenas dados fornecidos em CSV/JSON e sem criar capa.

Objetivo: transformar uma memória de cálculo de custo financeiro de FIDCs em um short deck executivo no estilo Itaú BBA: denso, limpo, institucional, com tabelas e gráficos compactos.

Design obrigatório:
- Fundo off-white claro; tipografia Arial/Calibri; títulos no topo com logo pequeno à esquerda.
- Paleta: preto para barras/totais/headers, laranja Itaú para destaque/mezanino, cinza para subordinada/apoio, verde para aumentos/captações, vermelho para reduções/amortizações.
- Cada slide deve ter 2 a 4 caixas grandes com borda fina arredondada; o título de cada caixa fica encostado na borda superior, com uma linha horizontal continuando à direita.
- Tabelas: header preto com fonte branca; linha Total preta com fonte branca; valores negativos em vermelho; valores positivos em verde.
- Gráficos: simples, nativos/editáveis quando possível, com rótulos diretos e pouca legenda. Não usar cards decorativos, hero, gradientes ou ilustrações.
- Rodapé corporativo em todas as páginas: fonte, CloudWalk | Corporativo | Interno, número da página e logo pequeno.

Estrutura de slides:
1. Metodologia + breakdown por FIDC:
   - Título com a conclusão: custo financeiro FIDC recomposto, custo bruto, custo líquido e TPV estimado.
   - Painel esquerdo: fluxo metodológico em 5 etapas: Informe Mensal CVM; saldos mensais; tipo de cota; caixa/LFT; CDI B3. Ao lado de cada etapa, explicar em uma linha como entra no motor.
   - Resultado no rodapé do painel esquerdo: custo bruto e custo líquido após carry de caixa/LFT.
   - Painel direito: tabela fundo a fundo com FIDC, saldo médio, CDI+ ponderado, faixa CDI+ e custo do ano.
   - Bloco final: fórmula de TPV antecipado = (365 / prazo médio) x estoque médio de recebíveis, com duration e TPV.

2. Evolução FIDCs:
   - Painel superior: waterfall do PL de 2025: PL inicial + captações - amortizações + accrual/rentabilidade = PL final. Usar preto para totais, verde para aumento, vermelho para redução.
   - Ao lado: tabela de amortizações 2025 por FIDC, ordenada do maior para o menor, com total.
   - Painel inferior esquerdo: capital stack agregado comparando 31/12/25 vs data atual, empilhado por Sênior, Mezanino e Subordinada; incluir seta verde mostrando aumento de PL.
   - Painel inferior direito: captações 2026 por FIDC/mês/série/classe, com total.

3. Prognóstico / run-rate 2026:
   - Painel superior esquerdo: amortizações programadas em 2026 por fundo em barras horizontais.
   - Painel superior direito: duration jan-abr/26, com gráfico mensal e tabela por FIDC mostrando duration e estoque médio.
   - Painel inferior esquerdo: custo financeiro bruto YTD jan-abr/26 em barras mensais e tabela por FIDC com YTD e anualizado.
   - Painel inferior direito: CDI+ implícito bruto comparando 2025FY vs 2026 run-rate, em bps e com a taxa CDI+ textual acima de cada barra.

Regras analíticas:
- Citar quando a fonte é IME/CVM, documentos regulatórios locais, B3/Cetip MediaCDI ou input manual.
- CDI sempre mensal composto quando houver dado mensal.
- Não usar saldo médio anual único para custo mensal se houver captação/amortização intraperíodo; usar saldos por trecho/mês.
- Amortizações 2025 devem vir dos eventos efetivos reportados no IME/CVM quando o objetivo for reconciliar PL; amortizações futuras podem vir de cronogramas documentais, claramente rotuladas como programadas.
- Não inventar dado faltante: colocar asterisco e nota metodológica curta.
"""
    PROMPT_OUT.write_text(prompt, encoding="utf-8")


def main():
    data = load_metrics()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_methodology(prs, data)
    slide_evolution(prs, data)
    slide_2026(prs, data)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    prs.save(ROOT_COPY)
    write_prompt()
    for path in [OUT, ROOT_COPY, PROMPT_OUT]:
        try:
            os.system(f"xattr -c {path!s}")
        except Exception:
            pass
    print(OUT)
    print(ROOT_COPY)
    print(PROMPT_OUT)


if __name__ == "__main__":
    main()
