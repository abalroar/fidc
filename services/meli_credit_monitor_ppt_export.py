from __future__ import annotations

from io import BytesIO
from typing import Any

import pandas as pd

from services.mercado_livre_dashboard import PT_MONTH_ABBR


MELI_BLACK = "000000"
MELI_ORANGE = "E47811"
MELI_DARK_GRAY = "3F3F3F"
MELI_MEDIUM_GRAY = "8C8C8C"
MELI_WHITE = "FFFFFF"


def build_dashboard_meli_pptx_bytes(monitor_outputs: Any, research_outputs: Any | None = None) -> bytes:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - environment guard
        raise RuntimeError("Dependência python-pptx não instalada.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    _add_consolidated_slide(
        prs=prs,
        layout=layout,
        monitor_outputs=monitor_outputs,
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    _add_consolidated_detail_slide(
        prs=prs,
        layout=layout,
        monitor_outputs=monitor_outputs,
        research_outputs=research_outputs,
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    for cnpj, monitor in getattr(monitor_outputs, "fund_monitor", {}).items():
        fund_name = _fund_name(monitor, fallback=str(cnpj))
        _add_fund_slide(
            prs=prs,
            layout=layout,
            title=fund_name,
            monitor_df=monitor,
            cohort_df=getattr(monitor_outputs, "fund_cohorts", {}).get(cnpj, pd.DataFrame()),
            CategoryChartData=CategoryChartData,
            RGBColor=RGBColor,
            XL_CHART_TYPE=XL_CHART_TYPE,
            XL_LABEL_POSITION=XL_LABEL_POSITION,
            XL_LEGEND_POSITION=XL_LEGEND_POSITION,
            XL_MARKER_STYLE=XL_MARKER_STYLE,
            Inches=Inches,
            Pt=Pt,
        )

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _add_consolidated_slide(
    *,
    prs,
    layout,
    monitor_outputs: Any,
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    Inches,
    Pt,
) -> None:
    slide = prs.slides.add_slide(layout)
    _style_slide(slide, RGBColor)
    _add_header(slide, "Análise Crédito - Consolidado", RGBColor, Inches, Pt)
    slots = _grid_2x2_slots()
    df = _chart_monthly(getattr(monitor_outputs, "consolidated_monitor", pd.DataFrame()))
    categories = _category_labels(df)
    _add_multi_line_chart(
        slide=slide,
        slot=slots[0],
        title="Roll rates",
        df=df,
        categories=categories,
        series_specs=[
            ("Roll 61-90 M-3", "roll_61_90_m3_pct", MELI_BLACK),
            ("Roll 151-180 M-6", "roll_151_180_m6_pct", MELI_ORANGE),
        ],
        y_axis_title="%",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    _add_stacked_npl_chart(
        slide=slide,
        slot=slots[1],
        df=df,
        categories=categories,
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        Inches=Inches,
        Pt=Pt,
    )
    _add_money_bar_chart(
        slide=slide,
        slot=slots[2],
        title="Carteira ex-360",
        df=df,
        categories=categories,
        column="carteira_ex360",
        series_name="Carteira ex-360",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        Inches=Inches,
        Pt=Pt,
    )
    _add_multi_line_chart(
        slide=slide,
        slot=slots[3],
        title="Crescimento YoY carteira ex-360",
        df=df,
        categories=categories,
        series_specs=[("Crescimento YoY", "carteira_ex360_yoy_pct", MELI_ORANGE)],
        y_axis_title="%",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )


def _add_fund_slide(
    *,
    prs,
    layout,
    title: str,
    monitor_df: pd.DataFrame,
    cohort_df: pd.DataFrame,
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    Inches,
    Pt,
) -> None:
    slide = prs.slides.add_slide(layout)
    _style_slide(slide, RGBColor)
    _add_header(slide, title, RGBColor, Inches, Pt)
    slots = _grid_2x2_slots()
    df = _chart_monthly(monitor_df)
    categories = _category_labels(df)
    _add_multi_line_chart(
        slide=slide,
        slot=slots[0],
        title="Roll rates",
        df=df,
        categories=categories,
        series_specs=[
            ("Roll 61-90 M-3", "roll_61_90_m3_pct", MELI_BLACK),
            ("Roll 151-180 M-6", "roll_151_180_m6_pct", MELI_ORANGE),
        ],
        y_axis_title="%",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    _add_stacked_npl_chart(
        slide=slide,
        slot=slots[1],
        df=df,
        categories=categories,
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        Inches=Inches,
        Pt=Pt,
    )
    _add_money_bar_chart(
        slide=slide,
        slot=slots[2],
        title="Carteira ex-360",
        df=df,
        categories=categories,
        column="carteira_ex360",
        series_name="Carteira ex-360",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        Inches=Inches,
        Pt=Pt,
    )
    _add_cohort_chart(
        slide=slide,
        slot=slots[3],
        title="Cohorts recentes",
        cohort_df=cohort_df,
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )


def _add_consolidated_detail_slide(
    *,
    prs,
    layout,
    monitor_outputs: Any,
    research_outputs: Any | None,
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    Inches,
    Pt,
) -> None:
    slide = prs.slides.add_slide(layout)
    _style_slide(slide, RGBColor)
    _add_header(slide, "Análise Crédito - Consolidado (continuação)", RGBColor, Inches, Pt)
    slots = _grid_2x2_slots()
    duration_df = _duration_frame(getattr(monitor_outputs, "consolidated_monitor", pd.DataFrame()), getattr(monitor_outputs, "fund_monitor", {}))
    duration_series = [column for column in duration_df.columns if column not in {"competencia_dt", "competencia"}]
    _add_multi_line_chart(
        slide=slide,
        slot=slots[0],
        title="Duration por FIDC",
        df=duration_df,
        categories=_category_labels(duration_df),
        series_specs=[(serie, serie, color) for serie, color in _series_colors(duration_series)],
        y_axis_title="meses",
        value_is_percent=False,
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    _add_cohort_chart(
        slide=slide,
        slot=slots[1],
        title="Cohorts recentes",
        cohort_df=getattr(monitor_outputs, "consolidated_cohorts", pd.DataFrame()),
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    roll_df = getattr(research_outputs, "roll_seasonality", pd.DataFrame()) if research_outputs is not None else pd.DataFrame()
    roll_61 = _research_roll_wide(roll_df, metric_id="roll_61_90_m3")
    roll_151 = _research_roll_wide(roll_df, metric_id="roll_151_180_m6")
    _add_multi_line_chart(
        slide=slide,
        slot=slots[2],
        title="Roll 61-90 por mês do ano",
        df=roll_61,
        categories=_wide_categories(roll_61),
        series_specs=[(serie, serie, color) for serie, color in _series_colors(_wide_series(roll_61))],
        y_axis_title="%",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    _add_multi_line_chart(
        slide=slide,
        slot=slots[3],
        title="Roll 151-180 por mês do ano",
        df=roll_151,
        categories=_wide_categories(roll_151),
        series_specs=[(serie, serie, color) for serie, color in _series_colors(_wide_series(roll_151))],
        y_axis_title="%",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )


def _add_multi_line_chart(
    *,
    slide,
    slot: tuple[float, float, float, float],
    title: str,
    df: pd.DataFrame,
    categories: list[str],
    series_specs: list[tuple[str, str, str]],
    y_axis_title: str,
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    Inches,
    Pt,
    value_is_percent: bool = True,
) -> None:
    valid_specs = [(name, column, color) for name, column, color in series_specs if column in df.columns or name in set(df.get("serie", pd.Series(dtype="object")))]
    if df.empty or not categories or not valid_specs:
        _add_empty(slide, slot, title, "Sem dados calculáveis.", RGBColor, Inches, Pt)
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    series_values: list[list[float | None]] = []
    for name, column, color in valid_specs:
        if column in df.columns:
            values = [_chart_value(value, percent=value_is_percent) for value in df[column]]
        else:
            series_df = df[df["serie"].eq(name)].copy()
            values = [_chart_value(value, percent=value_is_percent) for value in series_df.get("valor", pd.Series(dtype="float"))]
        chart_data.add_series(name, values)
        series_values.append(values)
    chart = _add_chart(slide, slot, XL_CHART_TYPE.LINE_MARKERS, chart_data, Inches)
    _style_common_chart(chart, title, "Competência", y_axis_title, "0.0%" if value_is_percent else "#,##0.0", RGBColor, Pt)
    chart.has_legend = True
    _style_legend(chart, XL_LEGEND_POSITION, RGBColor, Pt)
    for series, (_, _, color) in zip(chart.series, valid_specs, strict=False):
        _set_series_line(series, _rgb(color, RGBColor), XL_MARKER_STYLE.CIRCLE)
    _apply_last_point_labels(
        chart,
        series_values,
        [color for _, _, color in valid_specs],
        percent=value_is_percent,
        decimals=1,
        RGBColor=RGBColor,
        Pt=Pt,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
    )


def _add_stacked_npl_chart(
    *,
    slide,
    slot: tuple[float, float, float, float],
    df: pd.DataFrame,
    categories: list[str],
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    Inches,
    Pt,
) -> None:
    if df.empty or not categories:
        _add_empty(slide, slot, "NPL ex-360 por severidade", "Sem dados calculáveis.", RGBColor, Inches, Pt)
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    series_values = [
        [_chart_value(value, percent=True) for value in df.get("npl_1_90_pct", pd.Series(dtype="float"))],
        [_chart_value(value, percent=True) for value in df.get("npl_91_360_pct", pd.Series(dtype="float"))],
    ]
    chart_data.add_series("NPL 1-90d", series_values[0])
    chart_data.add_series("NPL 91-360d", series_values[1])
    chart = _add_chart(slide, slot, XL_CHART_TYPE.COLUMN_STACKED, chart_data, Inches)
    _style_common_chart(chart, "NPL ex-360 por severidade", "Competência", "% da carteira ex-360", "0.0%", RGBColor, Pt)
    chart.has_legend = True
    _style_legend(chart, XL_LEGEND_POSITION, RGBColor, Pt)
    for series, color in zip(chart.series, [MELI_BLACK, MELI_ORANGE], strict=False):
        _set_series_fill(series, _rgb(color, RGBColor))
    _apply_last_point_labels(
        chart,
        series_values,
        [MELI_WHITE, MELI_BLACK],
        percent=True,
        decimals=1,
        RGBColor=RGBColor,
        Pt=Pt,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        positions=[XL_LABEL_POSITION.CENTER, XL_LABEL_POSITION.CENTER],
    )


def _add_money_bar_chart(
    *,
    slide,
    slot: tuple[float, float, float, float],
    title: str,
    df: pd.DataFrame,
    categories: list[str],
    column: str,
    series_name: str,
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    Inches,
    Pt,
) -> None:
    if df.empty or not categories or column not in df.columns:
        _add_empty(slide, slot, title, "Sem dados calculáveis.", RGBColor, Inches, Pt)
        return
    scale, scale_label = _money_scale(df[column])
    values = [_scaled_value(value, scale) for value in df[column]]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = _add_chart(slide, slot, XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data, Inches)
    _style_common_chart(chart, title, "Competência", scale_label, "#,##0.0", RGBColor, Pt)
    chart.has_legend = False
    if chart.series:
        _set_series_fill(chart.series[0], _rgb(MELI_BLACK, RGBColor))
    _apply_last_point_labels(
        chart,
        [values],
        [MELI_BLACK],
        percent=False,
        decimals=0,
        RGBColor=RGBColor,
        Pt=Pt,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        positions=[XL_LABEL_POSITION.OUTSIDE_END],
    )


def _add_cohort_chart(
    *,
    slide,
    slot: tuple[float, float, float, float],
    title: str,
    cohort_df: pd.DataFrame,
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    Inches,
    Pt,
) -> None:
    if cohort_df is None or cohort_df.empty:
        _add_empty(slide, slot, title, "Sem dados calculáveis.", RGBColor, Inches, Pt)
        return
    df = cohort_df.copy()
    recent = df[["cohort", "cohort_dt"]].drop_duplicates().sort_values("cohort_dt").tail(6)["cohort"].tolist()
    df = df[df["cohort"].isin(recent)].copy()
    categories = ["M1", "M2", "M3", "M4", "M5", "M6"]
    colors = _cohort_ppt_colors(len(recent))
    chart_data = CategoryChartData()
    chart_data.categories = categories
    series_values: list[list[float | None]] = []
    for idx, cohort in enumerate(recent):
        group = df[df["cohort"].eq(cohort)].set_index("mes_ciclo")
        values = [_chart_value(group.loc[label, "valor_pct"], percent=True) if label in group.index else None for label in categories]
        chart_data.add_series(cohort, values)
        series_values.append(values)
    chart = _add_chart(slide, slot, XL_CHART_TYPE.LINE_MARKERS, chart_data, Inches)
    _style_common_chart(chart, title, "Mês de maturação", "% do saldo a vencer em 30d", "0.0%", RGBColor, Pt)
    chart.has_legend = True
    _style_legend(chart, XL_LEGEND_POSITION, RGBColor, Pt)
    for series, color in zip(chart.series, colors, strict=False):
        _set_series_line(series, _rgb(color, RGBColor), XL_MARKER_STYLE.CIRCLE)
    _apply_last_point_labels(
        chart,
        series_values,
        colors,
        percent=True,
        decimals=1,
        RGBColor=RGBColor,
        Pt=Pt,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
    )


def _research_roll_wide(roll_df: pd.DataFrame, *, metric_id: str) -> pd.DataFrame:
    if roll_df is None or roll_df.empty:
        return pd.DataFrame()
    df = roll_df[roll_df.get("scope", pd.Series(dtype="object")).astype(str).eq("consolidado") & roll_df["metric_id"].eq(metric_id)].copy()
    if df.empty:
        return pd.DataFrame()
    df["month"] = pd.to_numeric(df.get("month"), errors="coerce")
    df = df.sort_values("month")
    wide = df.pivot_table(index=["month", "month_label"], columns="series_name", values="value_pct", aggfunc="last").reset_index()
    wide.columns.name = None
    return wide.sort_values("month").rename(columns={"month_label": "competencia"}).reset_index(drop=True)


def _research_table_metric_wide(table_df: pd.DataFrame, *, metric_id: str) -> pd.DataFrame:
    if table_df is None or table_df.empty:
        return pd.DataFrame()
    df = table_df[table_df.get("scope", pd.Series(dtype="object")).astype(str).eq("consolidado") & table_df["metric_id"].eq(metric_id)].copy()
    if df.empty:
        return pd.DataFrame()
    df["competencia_dt"] = pd.to_datetime(df.get("competencia_dt"), errors="coerce")
    df = df.sort_values("competencia_dt")
    out = pd.DataFrame(
        {
            "competencia_dt": df["competencia_dt"],
            "competencia": [_format_month_label(ts) for ts in df["competencia_dt"]],
            "NPL 1-360": pd.to_numeric(df["value"], errors="coerce"),
        }
    )
    return out.reset_index(drop=True)


def _research_cohort_wide(cohort_df: pd.DataFrame) -> pd.DataFrame:
    if cohort_df is None or cohort_df.empty:
        return pd.DataFrame()
    df = cohort_df[cohort_df.get("scope", pd.Series(dtype="object")).astype(str).eq("consolidado")].copy()
    if df.empty:
        return pd.DataFrame()
    df["ordem"] = pd.to_numeric(df.get("ordem"), errors="coerce")
    wide = df.pivot_table(index=["ordem", "mes_ciclo"], columns="series_name", values="value_pct", aggfunc="last").reset_index()
    wide.columns.name = None
    return wide.sort_values("ordem").rename(columns={"mes_ciclo": "competencia"}).reset_index(drop=True)


def _wide_categories(df: pd.DataFrame) -> list[str]:
    if df is None or df.empty or "competencia" not in df.columns:
        return []
    return df["competencia"].astype(str).tolist()


def _wide_series(df: pd.DataFrame) -> list[str]:
    if df is None or df.empty:
        return []
    return [column for column in df.columns if column not in {"competencia_dt", "competencia", "month", "ordem"}]


def _add_chart(slide, slot, chart_type, chart_data, Inches):  # noqa: ANN001
    left, top, width, height = slot
    return slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data).chart


def _grid_2x2_slots() -> list[tuple[float, float, float, float]]:
    margin_x = 0.48
    top = 0.72
    gap_x = 0.26
    gap_y = 0.28
    width = (13.333 - margin_x * 2 - gap_x) / 2
    height = (7.5 - top - 0.24 - gap_y) / 2
    return [
        (margin_x, top, width, height),
        (margin_x + width + gap_x, top, width, height),
        (margin_x, top + height + gap_y, width, height),
        (margin_x + width + gap_x, top + height + gap_y, width, height),
    ]


def _style_common_chart(chart, title: str, x_title: str, y_title: str, number_format: str, RGBColor, Pt) -> None:  # noqa: ANN001
    chart.has_title = True
    _set_text(chart.chart_title.text_frame, title, Pt(11), True, MELI_BLACK, RGBColor)
    chart.category_axis.has_title = True
    _set_text(chart.category_axis.axis_title.text_frame, x_title, Pt(9), False, MELI_DARK_GRAY, RGBColor)
    chart.value_axis.has_title = True
    _set_text(chart.value_axis.axis_title.text_frame, y_title, Pt(9), False, MELI_DARK_GRAY, RGBColor)
    chart.category_axis.tick_labels.font.name = "Calibri"
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.color.rgb = _rgb(MELI_DARK_GRAY, RGBColor)
    chart.value_axis.tick_labels.font.name = "Calibri"
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.color.rgb = _rgb(MELI_DARK_GRAY, RGBColor)
    chart.value_axis.tick_labels.number_format = number_format
    chart.value_axis.has_major_gridlines = True


def _set_text(text_frame, text: str, font_size, bold: bool, color: str, RGBColor) -> None:  # noqa: ANN001
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = _rgb(color, RGBColor)


def _add_header(slide, title: str, RGBColor, Inches, Pt) -> None:  # noqa: ANN001
    box = slide.shapes.add_textbox(Inches(0.48), Inches(0.18), Inches(12.35), Inches(0.32))
    _set_text(box.text_frame, title, Pt(16), True, MELI_BLACK, RGBColor)


def _add_empty(slide, slot, title: str, message: str, RGBColor, Inches, Pt) -> None:  # noqa: ANN001
    left, top, width, height = slot
    title_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.04), Inches(width), Inches(0.24))
    _set_text(title_box.text_frame, title, Pt(10.5), True, MELI_BLACK, RGBColor)
    msg_box = slide.shapes.add_textbox(Inches(left), Inches(top + height / 2 - 0.1), Inches(width), Inches(0.28))
    _set_text(msg_box.text_frame, message, Pt(10), False, MELI_MEDIUM_GRAY, RGBColor)


def _style_legend(chart, XL_LEGEND_POSITION, RGBColor, Pt) -> None:  # noqa: ANN001
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = True
    chart.legend.font.name = "Calibri"
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = _rgb(MELI_DARK_GRAY, RGBColor)


def _apply_last_point_labels(
    chart,
    series_values: list[list[float | None]],
    colors: list[str],
    *,
    percent: bool,
    decimals: int,
    RGBColor,
    Pt,
    XL_LABEL_POSITION,
    positions: list | None = None,
) -> None:
    default_positions = [XL_LABEL_POSITION.RIGHT, XL_LABEL_POSITION.ABOVE, XL_LABEL_POSITION.BELOW, XL_LABEL_POSITION.LEFT]
    label_positions = positions or default_positions
    for series_idx, values in enumerate(series_values):
        point_idx = _last_non_null_index(values)
        if point_idx is None or series_idx >= len(chart.series):
            continue
        try:
            point = chart.series[series_idx].points[point_idx]
        except (IndexError, TypeError):
            continue
        label = _format_ppt_value(values[point_idx], percent=percent, decimals=decimals)
        _set_point_label(
            point,
            label,
            colors[series_idx % len(colors)],
            label_positions[series_idx % len(label_positions)],
            RGBColor,
            Pt,
        )


def _set_point_label(point, text: str, color: str, position, RGBColor, Pt) -> None:  # noqa: ANN001
    label = point.data_label
    label.position = position
    label.has_text_frame = True
    text_frame = label.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = _rgb(color, RGBColor)


def _style_slide(slide, RGBColor) -> None:  # noqa: ANN001
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def _set_series_fill(series, color) -> None:  # noqa: ANN001
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.color.rgb = color


def _set_series_line(series, color, marker_style) -> None:  # noqa: ANN001
    series.format.line.color.rgb = color
    series.format.line.width = 19050
    series.marker.style = marker_style
    series.marker.size = 5
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = color
    series.marker.format.line.color.rgb = color


def _duration_frame(consolidated_monitor: pd.DataFrame, fund_monitor: dict[str, pd.DataFrame]) -> pd.DataFrame:
    rows = []
    if consolidated_monitor is not None and not consolidated_monitor.empty:
        rows.extend(_duration_rows(consolidated_monitor, "Consolidado"))
    for cnpj, frame in fund_monitor.items():
        rows.extend(_duration_rows(frame, _fund_name(frame, fallback=str(cnpj))))
    if not rows:
        return pd.DataFrame()
    long = pd.DataFrame(rows)
    wide = long.pivot_table(index=["competencia_dt", "competencia"], columns="serie", values="duration_months", aggfunc="last").reset_index()
    wide.columns.name = None
    return wide.sort_values("competencia_dt").reset_index(drop=True)


def _duration_rows(frame: pd.DataFrame, label: str) -> list[dict[str, object]]:
    if frame is None or frame.empty:
        return []
    df = _chart_monthly(frame)
    values = pd.to_numeric(df.get("duration_months"), errors="coerce")
    return [
        {
            "competencia_dt": row.get("competencia_dt"),
            "competencia": row.get("competencia"),
            "serie": label,
            "duration_months": values.iloc[idx],
        }
        for idx, row in df.iterrows()
    ]


def _series_colors(series: pd.Series | list[str] | None) -> list[tuple[str, str]]:
    if series is None:
        return []
    colors = [MELI_BLACK, MELI_ORANGE, MELI_DARK_GRAY, MELI_MEDIUM_GRAY]
    names = pd.Series(series).dropna().drop_duplicates().astype(str).tolist()
    return [(name, colors[idx % len(colors)]) for idx, name in enumerate(names)]


def _cohort_ppt_colors(count: int) -> list[str]:
    colors = ["D9D9D9", "BDBDBD", "A0A0A0", "838383", "666666", "4A4A4A", "242424", MELI_BLACK]
    return colors[-count:] if count > 0 else colors


def _chart_monthly(monthly_df: pd.DataFrame) -> pd.DataFrame:
    if not isinstance(monthly_df, pd.DataFrame) or monthly_df.empty:
        return pd.DataFrame()
    df = monthly_df.copy()
    df["competencia_dt"] = pd.to_datetime(df.get("competencia_dt", df.get("competencia")), errors="coerce")
    return df.sort_values("competencia_dt").reset_index(drop=True)


def _category_labels(df: pd.DataFrame) -> list[str]:
    labels = []
    for _, row in df.iterrows():
        ts = pd.to_datetime(row.get("competencia_dt"), errors="coerce")
        labels.append(f"{PT_MONTH_ABBR[int(ts.month)]}/{str(int(ts.year))[-2:]}" if pd.notna(ts) else str(row.get("competencia") or ""))
    return labels


def _format_month_label(value: object) -> str:
    ts = pd.to_datetime(value, errors="coerce")
    return f"{PT_MONTH_ABBR[int(ts.month)]}/{str(int(ts.year))[-2:]}" if pd.notna(ts) else str(value or "")


def _money_scale(values: pd.Series) -> tuple[float, str]:
    max_value = pd.to_numeric(values, errors="coerce").abs().max()
    if pd.isna(max_value):
        max_value = 0.0
    if max_value >= 1_000_000_000_000:
        return 1_000_000_000.0, "R$ bi"
    if max_value >= 1_000_000:
        return 1_000_000.0, "R$ mm"
    if max_value >= 1_000:
        return 1_000.0, "R$ mil"
    return 1.0, "R$"


def _chart_value(value: object, *, percent: bool) -> float | None:
    numeric = _num(value)
    if numeric is None:
        return None
    return numeric / 100.0 if percent else numeric


def _scaled_value(value: object, divisor: float) -> float | None:
    numeric = _num(value)
    return None if numeric is None else numeric / divisor


def _last_non_null_index(values: list[object]) -> int | None:
    for idx in range(len(values) - 1, -1, -1):
        value = values[idx]
        if value is not None and not pd.isna(value):
            return idx
    return None


def _format_ppt_value(value: object, *, percent: bool, decimals: int = 1) -> str:
    numeric = _num(value)
    if numeric is None:
        return "N/D"
    if percent:
        return f"{numeric * 100.0:.{decimals}f}%".replace(".", ",")
    return f"{numeric:,.{decimals}f}".replace(",", "_").replace(".", ",").replace("_", ".")


def _num(value: object) -> float | None:
    parsed = pd.to_numeric(pd.Series([value]), errors="coerce").iloc[0]
    if pd.isna(parsed):
        return None
    return float(parsed)


def _fund_name(frame: pd.DataFrame, *, fallback: str) -> str:
    if frame is not None and not frame.empty and "fund_name" in frame.columns and frame["fund_name"].notna().any():
        return str(frame["fund_name"].dropna().iloc[0])
    return fallback


def _rgb(hex_color: str, RGBColor):  # noqa: ANN001
    clean = str(hex_color).strip().lstrip("#")
    return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))
