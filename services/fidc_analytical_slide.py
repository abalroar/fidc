"""Native IBBA-style FIDC analytical slides and auditable Excel workbook."""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
import re
from typing import Iterable

import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx.util import Pt

from services.bba_deck import GRAY_100, GRAY_500, GRAY_700, ORANGE, Deck, fmt_mm
from services.regulatory_profiles import load_regulatory_profile


ROOT = Path(__file__).resolve().parents[1]
VEHICLE_MONTHLY = ROOT / "data/industry_study/vehicle_monthly.csv.gz"
RDB_EVIDENCE = ROOT / "data/regulatory_profiles/meutudo_ago26_rdb_evidence.csv"


@dataclass
class Issuance:
    date: str
    series: str
    quota_class: str
    volume_mm: float | None
    remuneration: str
    amortization: str
    source: str


@dataclass
class FidcSlideData:
    name: str
    cnpj: str
    competence: str = "N/D"
    classification: str = "N/D"
    pl_mm: float | None = None
    receivables_pct_pl: float | None = None
    senior_pct_pl: float | None = None
    subordinated_pct_pl: float | None = None
    identified_issuance_volume_mm: float | None = None
    first_issuance: str = "N/D"
    last_issuance: str = "N/D"
    minimum_subordination: str = "N/D"
    issuances: list[Issuance] = field(default_factory=list)
    rdb_status: str = "Não localizado"
    rdb_finding: str = "N/D"
    rdb_source: str = "N/D"
    validations: list[str] = field(default_factory=list)


def normalize_cnpj(value: object) -> str:
    return re.sub(r"\D", "", str(value or "")).zfill(14)


def _display(value: object) -> str:
    text = str(value or "").strip()
    return text if text and text.lower() != "nan" else "N/D"


def _number(value: object) -> float | None:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None
    return None if pd.isna(number) else number


def _money_mm(value: object) -> float | None:
    text = str(value or "").strip()
    if not text:
        return None
    match = re.search(r"(?:R\$\s*)?([\d\.]+(?:,\d+)?)", text)
    if not match:
        return None
    number = float(match.group(1).replace(".", "").replace(",", "."))
    return number / 1_000_000


def _date(value: object) -> pd.Timestamp | None:
    parsed = pd.to_datetime(str(value or "").strip(), dayfirst=True, errors="coerce")
    return None if pd.isna(parsed) else parsed


def _short_name(name: str) -> str:
    cleaned = re.sub(r"\b(FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS|RESPONSABILIDADE LIMITADA|RESP LIMITADA)\b", "", name, flags=re.I)
    return re.sub(r"\s+", " ", cleaned).strip(" -")[:56]


def _clean_remuneration(value: object) -> str:
    text = _display(value)
    supported = re.search(r"\b(CDI|DI\s*\+|Taxa DI|IPCA|IGP-M|prefixad|residual|conforme (?:boletim|suplemento|regulamento)|não aplicável|N/D)\b", text, re.I)
    return text if supported and len(text) <= 180 else "N/D — texto extraído sem associação confiável"


def build_slide_data(funds: Iterable[object]) -> list[FidcSlideData]:
    monthly = pd.read_csv(VEHICLE_MONTHLY, low_memory=False) if VEHICLE_MONTHLY.exists() else pd.DataFrame()
    if not monthly.empty:
        monthly["_cnpj"] = monthly["cnpj"].map(normalize_cnpj)
        monthly["_date"] = pd.to_datetime(monthly["competencia"], errors="coerce")
    evidence = pd.read_csv(RDB_EVIDENCE, dtype=str, keep_default_na=False) if RDB_EVIDENCE.exists() else pd.DataFrame()
    if not evidence.empty:
        evidence["_cnpj"] = evidence["CNPJ"].map(normalize_cnpj)
    result: list[FidcSlideData] = []
    for fund in funds:
        cnpj = normalize_cnpj(getattr(fund, "cnpj", ""))
        name = _display(getattr(fund, "display_name", ""))
        data = FidcSlideData(name=_short_name(name), cnpj=cnpj)
        rows = monthly[monthly["_cnpj"].eq(cnpj)].sort_values("_date") if not monthly.empty else pd.DataFrame()
        if not rows.empty:
            row = rows.iloc[-1]
            if data.name == "N/D":
                data.name = _short_name(_display(row.get("denominacao")))
            data.competence = row["_date"].strftime("%m/%Y") if not pd.isna(row["_date"]) else _display(row.get("competencia"))
            data.classification = _display(row.get("classificacao_anbima"))
            pl = _number(row.get("pl")); dc = _number(row.get("carteira_dc")); sub = _number(row.get("vl_cotas_subordinadas"))
            data.pl_mm = pl / 1_000_000 if pl is not None else None
            data.receivables_pct_pl = dc / pl if dc is not None and pl else None
            data.subordinated_pct_pl = sub / pl if sub is not None and pl else None
            data.senior_pct_pl = 1 - data.subordinated_pct_pl if data.subordinated_pct_pl is not None else None
        profile = load_regulatory_profile(cnpj)
        if profile is not None and not profile.emissions_df.empty:
            seen: set[tuple[str, str, str, str]] = set()
            for _, row in profile.emissions_df.iterrows():
                date_value = _date(row.get("Data emissão / 1ª integralização")) or _date(row.get("Data deliberação"))
                if profile.profile_type in {"heurístico", "triagem estruturada"}:
                    continue
                series = _display(row.get("Cota/Classe")); volume_text = _display(row.get("Volume")); source = _display(row.get("Fonte"))
                key = (date_value.strftime("%Y-%m-%d") if date_value is not None else "", series.lower(), volume_text, source.split(";")[0])
                if key in seen:
                    continue
                seen.add(key)
                data.issuances.append(Issuance(
                    date=date_value.strftime("%m/%Y") if date_value is not None else "N/D",
                    series=series,
                    quota_class=_display(row.get("Tipo")),
                    volume_mm=_money_mm(row.get("Volume")),
                    remuneration=_clean_remuneration(row.get("Remuneração")),
                    amortization=_display(row.get("Amortização principal")),
                    source=source,
                ))
            dated = [(_date(i.date), i) for i in data.issuances if i.date != "N/D"]
            data.issuances.sort(key=lambda item: (_date(item.date) or pd.Timestamp.max, item.series))
            economic_issuances = [i for i in data.issuances if "evento" not in i.quota_class.lower() and "incorpora" not in i.series.lower() and "cisão" not in i.series.lower()]
            volumes = [i.volume_mm for i in economic_issuances if i.volume_mm is not None]
            data.identified_issuance_volume_mm = sum(volumes) if volumes else None
            valid_dates = sorted(_date(i.date) for i in economic_issuances if i.date != "N/D" and _date(i.date) is not None)
            if valid_dates:
                data.first_issuance, data.last_issuance = valid_dates[0].strftime("%m/%Y"), valid_dates[-1].strftime("%m/%Y")
        if profile is not None and profile.profile_type in {"curado", "curado parcial"} and not profile.criteria_df.empty:
            criteria = profile.criteria_df
            mask = criteria.apply(lambda row: "subord" in " ".join(map(str, row.values)).lower(), axis=1)
            if mask.any():
                data.minimum_subordination = _display(criteria.loc[mask].iloc[0].get("Limite/regra"))
        ev = evidence[evidence["_cnpj"].eq(cnpj)] if not evidence.empty else pd.DataFrame()
        if not ev.empty:
            row = ev.iloc[0]
            data.rdb_status = _display(row.get("Status")); data.rdb_finding = _display(row.get("Achado regulamento"))
            data.rdb_source = f"{_display(row.get('Documento'))} · p. {_display(row.get('Página'))}"
        _validate(data)
        result.append(data)
    return result


def _validate(data: FidcSlideData) -> None:
    if data.senior_pct_pl is not None and data.subordinated_pct_pl is not None:
        if abs(data.senior_pct_pl + data.subordinated_pct_pl - 1) > 0.01:
            data.validations.append("Divergência: sênior + subordinada difere de 100%.")
    if data.pl_mm is None:
        data.validations.append("PL atual indisponível na base mensal.")
    if not data.issuances:
        data.validations.append("Histórico de emissões não reconstruído com confiança.")
    elif any(i.date == "N/D" for i in data.issuances):
        data.validations.append("Há emissões sem data confiável; primeira/última emissão podem ser parciais.")
    if data.rdb_status == "Não localizado":
        data.validations.append("RDB não localizado nos documentos extraídos; ausência não prova inexistência.")


def _pct(value: float | None) -> str:
    return "N/D" if value is None else f"{value * 100:.1f}%".replace(".", ",")


def _money(value: float | None) -> str:
    return "N/D" if value is None else f"R$ {fmt_mm(value)} mm"


def build_fidc_analytical_pptx_bytes(funds: Iterable[object], *, portfolio_name: str) -> bytes:
    data = build_slide_data(funds)
    deck = Deck("CRÉDITO ESTRUTURADO | CURADORIA DOCUMENTAL")
    _render_rdb_synthesis(deck, data, portfolio_name)
    for item in data:
        _render_fund_slides(deck, item)
    output = BytesIO(); deck.prs.save(output); return output.getvalue()


def _render_rdb_synthesis(deck: Deck, data: list[FidcSlideData], portfolio_name: str) -> None:
    slide = deck.slide(f"{portfolio_name} | RDB Parati nos regulamentos")
    deck.text(slide, "Leitura contratual por veículo", .62, 1.34, 4.0, .3, size=12, bold=True)
    rows = [["FIDC / classe", "Status", "Leitura"]]
    for item in data:
        rows.append([item.name[:37], item.rdb_status, item.rdb_finding[:105]])
    deck.native_table(slide, rows, .62, 1.72, [3.0, 1.25, 7.8], row_height=.39, header_height=.34, size=8.6, aligns="lll")
    deck.text(slide, "Conclusão documental", .62, 6.15, 2.2, .24, size=10, color=ORANGE, bold=True)
    deck.text(slide, "A DF da Parati descreve RDBs emitidos em favor de FIDCs com finalidade de garantia. Nos regulamentos lidos, o RDB aparece como ativo da classe; o MT FGTS Receivables II detalha a compra direta e a alternativa via CCB Tudo Serviços. A vinculação ampla do RDB como garantia de cada CCB cedida não foi localizada.", 2.55, 6.08, 10.0, .62, size=9.2, color=GRAY_700)
    deck.footer(slide, "Fontes: Fundos.NET/B3, documentos regulatórios indicados por linha; Parati CFI, DF dez/2025. Geração: 14/08/2026. Uso interno.")


def _render_fund_slides(deck: Deck, data: FidcSlideData) -> None:
    chunks = [data.issuances[i:i + 7] for i in range(0, len(data.issuances), 7)] or [[]]
    for page, chunk in enumerate(chunks, 1):
        title = f"{data.name} | estrutura e emissões" + (f" ({page}/{len(chunks)})" if len(chunks) > 1 else "")
        slide = deck.slide(title)
        if len(title) > 58:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text == title:
                    shape.text_frame.paragraphs[0].font.size = Pt(17.5)
        if page == 1:
            snapshot = [
                ["Snapshot", "Valor"], ["PL atual", _money(data.pl_mm)], ["Direitos creditórios / PL", _pct(data.receivables_pct_pl)],
                ["Cotas sênior / PL", _pct(data.senior_pct_pl)], ["Cotas subordinadas / PL", _pct(data.subordinated_pct_pl)],
                ["Volume emitido identificado", _money(data.identified_issuance_volume_mm)], ["Primeira / última emissão", f"{data.first_issuance} / {data.last_issuance}"],
            ]
            deck.native_table(slide, snapshot, .62, 1.43, [2.18, 1.55], row_height=.34, header_height=.34, size=9, aligns="lr")
            deck.text(slide, "Emissões e remuneração", 4.62, 1.40, 4.0, .3, size=11, bold=True)
            x, w = 4.62, 5.25
        else:
            deck.text(slide, "Emissões e remuneração — continuação", .62, 1.40, 5.0, .3, size=11, bold=True)
            x, w = .62, 9.25
        rows = [["Data", "Série / classe", "Volume", "Remuneração"]]
        for issuance in chunk:
            rows.append([issuance.date, issuance.series[:32], _money(issuance.volume_mm).replace("R$ ", ""), issuance.remuneration[:55]])
        if not chunk:
            rows.append(["N/D", "Dados insuficientes para reconstrução", "N/D", "N/D"])
        deck.native_table(slide, rows, x, 1.76, [0.65, 2.05 if page == 1 else 3.0, .9, w - (3.6 if page == 1 else 4.55)], row_height=.43, header_height=.34, size=7.8, aligns="llrl")
        if page == 1:
            deck.text(slide, "Subordinação mínima", .62, 4.37, 2.0, .24, size=10, color=ORANGE, bold=True)
            deck.text(slide, data.minimum_subordination, 2.45, 4.34, 1.9, .36, size=11, bold=True)
            deck.text(slide, "RDB Parati", .62, 4.91, 2.0, .24, size=10, color=ORANGE, bold=True)
            deck.text(slide, f"{data.rdb_status}: {data.rdb_finding}", .62, 5.19, 3.73, .92, size=8.8, color=GRAY_700)
            deck.text(slide, "Perfil de amortização", 10.15, 1.40, 2.3, .3, size=11, bold=True)
            deck.block(slide, 10.15, 1.76, 2.53, 3.67, GRAY_100)
            deck.text(slide, "Dados insuficientes para reconstrução confiável do waterfall.", 10.42, 3.05, 2.0, .7, size=10, color=GRAY_500, bold=True, align=1)
            deck.text(slide, "Controles", 10.15, 5.62, 1.3, .22, size=9, color=ORANGE, bold=True)
            deck.text(slide, "\n".join(f"• {x}" for x in data.validations[:3]) or "• Sem divergência material identificada", 10.15, 5.88, 2.53, .72, size=7.5, color=GRAY_700)
        sources = sorted({i.source for i in chunk if i.source != "N/D"})
        source = f"CVM Informe Mensal {data.competence}; Fundos.NET/B3: " + ("; ".join(sources)[:340] if sources else data.rdb_source)
        deck.footer(slide, source + ". Geração: 14/08/2026. Uso interno.")


def build_fidc_analytical_xlsx_bytes(funds: Iterable[object], *, portfolio_name: str) -> bytes:
    data = build_slide_data(funds)
    wb = Workbook(); wb.remove(wb.active)
    _write_sheet(wb, "Snapshot", [
        ["Carteira", "CNPJ", "FIDC", "Competência", "PL (R$ mm)", "DC / PL", "Sênior / PL", "Subordinada / PL", "Volume emitido (R$ mm)", "Primeira emissão", "Última emissão", "Subordinação mínima"]
    ] + [[portfolio_name, d.cnpj, d.name, d.competence, d.pl_mm, d.receivables_pct_pl, d.senior_pct_pl, d.subordinated_pct_pl, d.identified_issuance_volume_mm, d.first_issuance, d.last_issuance, d.minimum_subordination] for d in data], "tbSnapshot")
    _write_sheet(wb, "Emissões", [["CNPJ", "FIDC", "Data", "Série / classe", "Tipo", "Volume (R$ mm)", "Remuneração", "Amortização", "Fonte"]] + [[d.cnpj, d.name, i.date, i.series, i.quota_class, i.volume_mm, i.remuneration, i.amortization, i.source] for d in data for i in d.issuances], "tbEmissoes")
    evidence = pd.read_csv(RDB_EVIDENCE, dtype=str, keep_default_na=False)
    _write_sheet(wb, "RDB", [list(evidence.columns)] + evidence.values.tolist(), "tbRDB")
    _write_sheet(wb, "Validações", [["CNPJ", "FIDC", "Controle / divergência"]] + [[d.cnpj, d.name, issue] for d in data for issue in (d.validations or ["Sem divergência material identificada"])] , "tbValidacoes")
    ws = wb.create_sheet("Metodologia"); ws.append(["Campo", "Regra"])
    for row in [("Separação", "extração → normalização → cálculos → validação → renderização"), ("PL e composição", "última competência disponível por CNPJ em vehicle_monthly.csv.gz"), ("Emissões", "perfil regulatório curado; deduplicação por data, série, volume e primeira fonte"), ("Waterfall", "somente produzido com cronograma documental quantificável; caso contrário, N/D"), ("RDB", "matriz documental por CNPJ, com página e status explícito/parcial/não localizado")]: ws.append(row)
    _style_sheet(ws); output = BytesIO(); wb.save(output); return output.getvalue()


def _write_sheet(wb: Workbook, title: str, rows: list[list[object]], table_name: str) -> None:
    ws = wb.create_sheet(title)
    for row in rows or [["N/D"]]: ws.append(row)
    _style_sheet(ws)
    if ws.max_row >= 2 and ws.max_column >= 1:
        ref = f"A1:{ws.cell(ws.max_row, ws.max_column).coordinate}"
        table = Table(displayName=table_name, ref=ref); table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True, showFirstColumn=False, showLastColumn=False)
        ws.add_table(table)


def _style_sheet(ws) -> None:
    ws.freeze_panes = "A2"; ws.auto_filter.ref = ws.dimensions
    for cell in ws[1]:
        cell.fill = PatternFill("solid", fgColor="14315C"); cell.font = Font(name="Arial", size=10, bold=True, color="FFFFFF"); cell.alignment = Alignment(vertical="center")
    for column in ws.columns:
        letter = column[0].column_letter; ws.column_dimensions[letter].width = min(55, max(11, max(len(str(c.value or "")) for c in column) + 2))
        for cell in column[1:]: cell.font = Font(name="Arial", size=9); cell.alignment = Alignment(vertical="top", wrap_text=True)


__all__ = ["FidcSlideData", "Issuance", "build_slide_data", "build_fidc_analytical_pptx_bytes", "build_fidc_analytical_xlsx_bytes"]
