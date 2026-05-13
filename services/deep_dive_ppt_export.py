from __future__ import annotations

from datetime import datetime
from io import BytesIO
from zoneinfo import ZoneInfo

import pandas as pd

from services.deep_dive_models import DeepDiveManifest, DeepDiveTableSpec


SLIDE_W = 13.333
SLIDE_H = 7.5
LEFT = 0.34
RIGHT = 0.28
TOP = 0.26
FOOTER_TOP = 7.08
CONTENT_W = SLIDE_W - LEFT - RIGHT
TABLE_TOP = 0.96
TABLE_BOTTOM = 6.88
TABLE_H = TABLE_BOTTOM - TABLE_TOP

BLACK = "1F1F1F"
HEADER = "111827"
ORANGE = "EC7000"
WHITE = "FFFFFF"
SOFT = "F7F7F7"
GRID = "D9DEE5"
MID = "6B7280"
HIGHLIGHT = "FFF2E8"
RED_TEXT = "C8102E"
FONT = "Calibri"


def build_deep_dive_pptx_bytes(
    manifest: DeepDiveManifest,
    tables: list[tuple[DeepDiveTableSpec, pd.DataFrame]],
    *,
    highlighted_column: str | None = None,
    generated_at: datetime | None = None,
) -> bytes:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("Dependência python-pptx não instalada.") from exc

    tz = ZoneInfo("America/Sao_Paulo")
    if generated_at is None:
        generated_at = datetime.now(tz)
    elif generated_at.tzinfo is None:
        generated_at = generated_at.replace(tzinfo=tz)
    else:
        generated_at = generated_at.astimezone(tz)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    def rgb(hex_color: str):  # noqa: ANN202
        value = str(hex_color).strip().lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    def add_text(slide, left, top, width, height, text, *, size, bold=False, color=BLACK, align=PP_ALIGN.LEFT):  # noqa: ANN001, ANN202
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.01)
        tf.margin_right = Inches(0.01)
        tf.margin_top = Inches(0.00)
        tf.margin_bottom = Inches(0.00)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text or "")
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return box

    def add_header(slide, title: str, subtitle: str) -> None:  # noqa: ANN001
        add_text(slide, LEFT, TOP, 8.5, 0.26, manifest.title, size=17, bold=True, color=BLACK)
        line2 = subtitle or manifest.subtitle
        if line2:
            add_text(slide, LEFT, TOP + 0.30, 9.5, 0.18, line2, size=8.2, color=MID)
        source = manifest.source or "Deep Dive offline"
        add_text(slide, 9.25, TOP, 3.55, 0.18, source, size=7.5, color=MID, align=PP_ALIGN.RIGHT)

    def add_footer(slide, page: int, total: int) -> None:  # noqa: ANN001
        sep = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(LEFT), Inches(6.98), Inches(CONTENT_W), Inches(0.005))
        sep.fill.solid()
        sep.fill.fore_color.rgb = rgb(GRID)
        sep.line.fill.background()
        stamp = generated_at.strftime("%d/%m/%Y %H:%M")
        add_text(slide, LEFT, FOOTER_TOP, 8.8, 0.16, f"Fonte: {manifest.source or 'pacote offline'} | Gerado em: {stamp} | {manifest.confidentiality}", size=7, color=MID)
        add_text(slide, 11.20, FOOTER_TOP, 1.55, 0.16, f"Página {page} de {total}", size=7, color=MID, align=PP_ALIGN.RIGHT)

    slide_jobs: list[tuple[DeepDiveTableSpec, pd.DataFrame, int, int]] = []
    for spec, frame in tables:
        normalized = _normalize_table(frame, first_column=spec.first_column)
        col_chunks = _column_chunks(normalized)
        for col_start, col_end in col_chunks:
            chunk_cols = [spec.first_column, *normalized.columns[col_start:col_end].tolist()]
            col_frame = normalized[chunk_cols].copy()
            rows_per_slide = _rows_per_slide(col_frame)
            for row_start in range(0, len(col_frame), rows_per_slide):
                slide_jobs.append((spec, col_frame.iloc[row_start : row_start + rows_per_slide].copy(), row_start, len(col_frame)))
    if not slide_jobs:
        slide_jobs.append((DeepDiveTableSpec(id="empty", title="Deep Dive", source_file=""), pd.DataFrame({"Nome": ["Sem dados"]}), 0, 1))

    total_pages = len(slide_jobs)
    for page, (spec, frame, row_start, row_total) in enumerate(slide_jobs, start=1):
        slide = prs.slides.add_slide(blank)
        suffix = f"{spec.title}"
        if row_total > len(frame):
            suffix = f"{suffix} · linhas {row_start + 1}-{row_start + len(frame)} de {row_total}"
        add_header(slide, manifest.title, suffix)
        add_footer(slide, page, total_pages)
        _add_table(
            slide,
            frame,
            highlighted_column=highlighted_column,
            rgb=rgb,
            Inches=Inches,
            Pt=Pt,
            MSO_ANCHOR=MSO_ANCHOR,
            PP_ALIGN=PP_ALIGN,
        )

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _normalize_table(frame: pd.DataFrame, *, first_column: str) -> pd.DataFrame:
    if frame is None or frame.empty:
        return pd.DataFrame({first_column: ["Sem dados"]})
    output = frame.copy().fillna("—").replace("", "—")
    if first_column not in output.columns:
        output = output.rename(columns={output.columns[0]: first_column})
    return output


def _column_chunks(frame: pd.DataFrame) -> list[tuple[int, int]]:
    other_cols = list(range(1, len(frame.columns)))
    if len(other_cols) <= 5:
        return [(1, len(frame.columns))]
    chunks = []
    for start in range(1, len(frame.columns), 5):
        chunks.append((start, min(start + 5, len(frame.columns))))
    return chunks


def _rows_per_slide(frame: pd.DataFrame) -> int:
    cols = len(frame.columns)
    if cols <= 4:
        return 25
    if cols <= 6:
        return 22
    return 19


def _column_widths(columns: list[str]) -> list[float]:
    total = CONTENT_W
    if len(columns) == 1:
        return [total]
    first = 2.25 if len(columns) <= 5 else 2.05
    other = (total - first) / (len(columns) - 1)
    return [first, *([other] * (len(columns) - 1))]


def _add_table(slide, frame: pd.DataFrame, *, highlighted_column: str | None, rgb, Inches, Pt, MSO_ANCHOR, PP_ALIGN) -> None:  # noqa: ANN001, PLR0913
    rows = len(frame) + 1
    cols = len(frame.columns)
    table = slide.shapes.add_table(rows, cols, Inches(LEFT), Inches(TABLE_TOP), Inches(CONTENT_W), Inches(TABLE_H)).table
    widths = _column_widths(list(frame.columns))
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for row in table.rows:
        row.height = Inches(TABLE_H / rows)

    for col_idx, column in enumerate(frame.columns):
        table.cell(0, col_idx).text = str(column)
    for row_idx, (_, row) in enumerate(frame.iterrows(), start=1):
        for col_idx, column in enumerate(frame.columns):
            table.cell(row_idx, col_idx).text = _clip_cell(row.get(column, "—"))

    for row_idx in range(rows):
        for col_idx in range(cols):
            column = frame.columns[col_idx]
            cell = table.cell(row_idx, col_idx)
            is_header = row_idx == 0
            is_first = col_idx == 0
            is_highlight = bool(highlighted_column) and column == highlighted_column
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.045)
            cell.margin_right = Inches(0.045)
            cell.margin_top = Inches(0.012)
            cell.margin_bottom = Inches(0.012)
            cell.fill.solid()
            if is_header:
                cell.fill.fore_color.rgb = rgb(ORANGE if is_highlight else HEADER)
            elif is_highlight:
                cell.fill.fore_color.rgb = rgb(HIGHLIGHT)
            elif is_first:
                cell.fill.fore_color.rgb = rgb(SOFT)
            else:
                cell.fill.fore_color.rgb = rgb(WHITE if row_idx % 2 else "FBFBFB")
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if is_first else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    text = str(run.text or "")
                    run.font.name = FONT
                    run.font.size = Pt(_font_size(row_count=rows - 1, col_count=cols, is_header=is_header))
                    run.font.bold = is_header or is_first
                    run.font.color.rgb = rgb(WHITE if is_header else RED_TEXT if _looks_relevant(text) else BLACK)
            cell.text_frame.word_wrap = True


def _font_size(*, row_count: int, col_count: int, is_header: bool) -> float:
    if is_header:
        return 8.0 if col_count <= 5 else 7.4
    if row_count >= 23 or col_count >= 6:
        return 7.2
    if row_count >= 19:
        return 7.7
    return 8.2


def _clip_cell(value: object) -> str:
    text = str(value if value is not None else "—").strip() or "—"
    if len(text) <= 110:
        return text
    return text[:107].rstrip() + "..."


def _looks_relevant(text: str) -> bool:
    lowered = str(text or "").lower()
    return any(token in lowered for token in ("diverg", "alerta", "lacuna", "não identific", "ausente", "waiver"))
