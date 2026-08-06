"""Shared 16:9 PowerPoint scaffolding for the BBA decks.

Palette and typography match ``services/industry_ppt_export.py`` so decks built
here sit alongside the existing ones without a visual break: white surface,
orange accent, Arial, 13.333 × 7.5 inches, kicker + title + rule, footer with
source and page number.

Everything rendered is native and editable — no rasterized images.
"""

from __future__ import annotations

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BLACK = "151515"
ORANGE = "E36C0A"
ORANGE_LIGHT = "F8E9DE"
GRAY_900 = "30353A"
GRAY_700 = "5D6369"
GRAY_500 = "8D9399"
GRAY_300 = "D7DADD"
GRAY_200 = "E7E9EB"
GRAY_100 = "F5F6F7"
WHITE = "FFFFFF"

#: Chart emphasis pair.  ``HOUSE_BLUE`` marks the house; ``NEUTRAL`` is the
#: recessive backdrop for every competitor.  Validated for CVD separation
#: (ΔE 33.9 protan) and ≥3:1 contrast against the white chart surface.
HOUSE_BLUE = "14315C"
NEUTRAL = "8D9399"

FONT = "Arial"

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MARGIN_IN = 0.62
CONTENT_WIDTH_IN = 12.05


def fmt_mm(value: float, decimals: int = 1) -> str:
    """Brazilian thousand/decimal separators for a R$ million figure."""

    if pd.isna(value):
        return "—"
    text = f"{value:,.{decimals}f}"
    return text.replace(",", "@").replace(".", ",").replace("@", ".")


def fmt_pct(value: float, decimals: int = 1) -> str:
    if pd.isna(value):
        return "—"
    return f"{value * 100:,.{decimals}f}%".replace(".", ",")


def fmt_pp(value: float, decimals: int = 1) -> str:
    if pd.isna(value):
        return "—"
    sign = "+" if value > 0 else ""
    return f"{sign}{value:,.{decimals}f}".replace(".", ",")


def fmt_rank(value: object) -> str:
    return "—" if value is None or pd.isna(value) else f"{int(value)}º"


#: "No Style, No Grid" — the built-in table style that ships no banding and no
#: borders, so the deck's own styling is the only thing on screen.
_PLAIN_TABLE_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

_A_NAMESPACE = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _use_plain_table_style(table) -> None:
    graphic = table._graphic_frame._element  # noqa: SLF001 - no public accessor
    for properties in graphic.iter(f"{{{_A_NAMESPACE}}}tblPr"):
        for existing in properties.findall(f"{{{_A_NAMESPACE}}}tableStyleId"):
            properties.remove(existing)
        style = properties.makeelement(f"{{{_A_NAMESPACE}}}tableStyleId", {})
        style.text = _PLAIN_TABLE_STYLE
        properties.append(style)


def _set_cell_borders(cell, *, bottom: str | None = None) -> None:
    """Draw a hairline under a table cell; python-pptx exposes no border API."""

    if bottom is None:
        return
    properties = cell._tc.get_or_add_tcPr()  # noqa: SLF001 - no public accessor
    tag = f"{{{_A_NAMESPACE}}}lnB"
    for existing in properties.findall(tag):
        properties.remove(existing)
    line = properties.makeelement(tag, {"w": "6350", "cap": "flat"})
    fill = properties.makeelement(f"{{{_A_NAMESPACE}}}solidFill", {})
    color = properties.makeelement(f"{{{_A_NAMESPACE}}}srgbClr", {"val": bottom})
    fill.append(color)
    line.append(fill)
    properties.append(line)


class Deck:
    """A minimal, opinionated deck builder for the BBA house style."""

    def __init__(self, kicker: str) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_WIDTH_IN)
        self.prs.slide_height = Inches(SLIDE_HEIGHT_IN)
        self.kicker = kicker
        self.page = 0

    @staticmethod
    def rgb(value: str) -> RGBColor:
        return RGBColor.from_string(value)

    def text(
        self,
        slide,
        content: object,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float = 12,
        color: str = GRAY_900,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = valign
        frame.margin_left = frame.margin_right = Inches(0)
        frame.margin_top = frame.margin_bottom = Inches(0)
        paragraph = frame.paragraphs[0]
        paragraph.text = str(content)
        paragraph.alignment = align
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = self.rgb(color)
        return box

    def rule(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        *,
        color: str = GRAY_300,
        height: float = 0.012,
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.rgb(color)
        shape.line.fill.background()
        return shape

    def block(self, slide, x: float, y: float, w: float, h: float, color: str):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.rgb(color)
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def blank(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.rgb(WHITE)
        return slide

    def slide(self, title: str, kicker: str | None = None):
        slide = self.blank()
        self.text(
            slide,
            kicker or self.kicker,
            MARGIN_IN,
            0.28,
            8.0,
            0.22,
            size=10,
            color=ORANGE,
            bold=True,
        )
        self.text(
            slide, title, MARGIN_IN, 0.57, 12.1, 0.48, size=22, color=BLACK, bold=True
        )
        self.rule(slide, MARGIN_IN, 1.15, CONTENT_WIDTH_IN, color=GRAY_300, height=0.018)
        return slide

    def footer(self, slide, source: str) -> None:
        self.page += 1
        self.rule(slide, MARGIN_IN, 6.95, CONTENT_WIDTH_IN, color=GRAY_200, height=0.01)
        self.text(slide, source, MARGIN_IN, 7.02, 11.45, 0.3, size=8, color=GRAY_500)
        self.text(
            slide,
            str(self.page),
            12.25,
            7.01,
            0.42,
            0.22,
            size=8,
            color=GRAY_500,
            align=PP_ALIGN.RIGHT,
        )

    def stat_cards(
        self, slide, cards: list[tuple[str, str, str]], y: float = 1.45
    ) -> float:
        """Render up to four headline figures across the content width."""

        for index, (label, value, note) in enumerate(cards):
            x = MARGIN_IN + index * 3.06
            self.block(slide, x, y, 2.86, 1.16, GRAY_100)
            self.text(slide, label.upper(), x + 0.2, y + 0.15, 2.5, 0.22, size=9, color=GRAY_500, bold=True)
            self.text(slide, value, x + 0.2, y + 0.41, 2.5, 0.42, size=20, color=BLACK, bold=True)
            self.text(slide, note, x + 0.2, y + 0.85, 2.5, 0.22, size=9, color=GRAY_500)
        return y + 1.16

    def native_table(
        self,
        slide,
        rows: list[list[str]],
        x: float,
        y: float,
        widths: list[float],
        *,
        highlight: int | None = None,
        row_height: float = 0.28,
        header_height: float = 0.32,
        size: float = 10,
        aligns: str | None = None,
        align_right_from: int = 2,
        header_fill: str = GRAY_100,
        header_color: str = GRAY_500,
        emphasis_rows: tuple[int, ...] = (),
        emphasis_fill: str = ORANGE,
        emphasis_color: str = WHITE,
    ):
        """Insert a real PowerPoint table — editable, with addable rows/columns.

        Unlike a grid of text boxes, this is a single object the reader can
        select, restyle, sort, or paste into Excel, and into which rows and
        columns can be inserted with the normal Office commands.
        """

        row_count = len(rows)
        column_count = len(rows[0])
        total_height = header_height + (row_count - 1) * row_height
        graphic = slide.shapes.add_table(
            row_count,
            column_count,
            Inches(x),
            Inches(y),
            Inches(sum(widths)),
            Inches(total_height),
        )
        table = graphic.table
        _use_plain_table_style(table)
        table.first_row = True
        table.horz_banding = False

        for index, width in enumerate(widths):
            table.columns[index].width = Inches(width)
        table.rows[0].height = Inches(header_height)
        for index in range(1, row_count):
            table.rows[index].height = Inches(row_height)

        def alignment(index: int):
            if aligns is not None and index < len(aligns):
                return PP_ALIGN.RIGHT if aligns[index] == "r" else PP_ALIGN.LEFT
            return PP_ALIGN.RIGHT if index >= align_right_from else PP_ALIGN.LEFT

        for row_index, values in enumerate(rows):
            is_header = row_index == 0
            is_house = highlight is not None and row_index == highlight + 1
            is_emphasis = row_index in emphasis_rows
            for column_index, value in enumerate(values):
                cell = table.cell(row_index, column_index)
                cell.text = str(value)
                cell.margin_left = Inches(0.06)
                cell.margin_right = Inches(0.06)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid()
                if is_header:
                    fill_color = header_fill
                elif is_emphasis:
                    fill_color = emphasis_fill
                elif is_house:
                    fill_color = ORANGE_LIGHT
                else:
                    fill_color = WHITE
                cell.fill.fore_color.rgb = self.rgb(fill_color)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.alignment = alignment(column_index)
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = FONT
                run.font.size = Pt(size - 1 if is_header else size)
                run.font.bold = is_header or is_house or is_emphasis
                if is_header:
                    text_color = header_color
                elif is_emphasis:
                    text_color = emphasis_color
                elif is_house:
                    text_color = BLACK
                else:
                    text_color = GRAY_900
                run.font.color.rgb = self.rgb(text_color)
                _set_cell_borders(
                    cell,
                    bottom=(GRAY_300 if is_header else GRAY_200),
                )
        return y + total_height

    def save(self, path) -> None:
        self.prs.save(path)


__all__ = [
    "BLACK",
    "CONTENT_WIDTH_IN",
    "Deck",
    "FONT",
    "GRAY_100",
    "GRAY_200",
    "GRAY_300",
    "GRAY_500",
    "GRAY_700",
    "GRAY_900",
    "HOUSE_BLUE",
    "MARGIN_IN",
    "NEUTRAL",
    "ORANGE",
    "ORANGE_LIGHT",
    "SLIDE_HEIGHT_IN",
    "SLIDE_WIDTH_IN",
    "WHITE",
    "fmt_mm",
    "fmt_pct",
    "fmt_pp",
    "fmt_rank",
]
