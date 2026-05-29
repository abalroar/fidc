from __future__ import annotations

from io import BytesIO
from typing import Any

import pandas as pd

from services.cloudwalk_financial_cost import FinancialCostOutputs
from services.cloudwalk_pl_waterfall import CloudwalkPlWaterfall
from services.fidc_model.b3_cdi import B3CdiMonthlyRate


BRAND_ORANGE = "D35714"
INK = "111111"
MUTED = "68727D"
GRID = "DDE3EA"
GREEN = "0F9D58"
RED = "D93025"


def build_cloudwalk_financial_cost_xlsx_bytes(
    outputs: FinancialCostOutputs,
    *,
    pl_waterfall: CloudwalkPlWaterfall | None = None,
    monthly_cdi_rates: tuple[B3CdiMonthlyRate, ...] = (),
) -> bytes:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, Reference
    from openpyxl.styles import Alignment, Font, PatternFill, Side, Border
    from openpyxl.utils import get_column_letter
    from openpyxl.utils.dataframe import dataframe_to_rows

    wb = Workbook()
    ws = wb.active
    ws.title = "Resumo"
    thin = Side(style="thin", color=GRID)
    header_fill = PatternFill("solid", fgColor=INK)
    accent_fill = PatternFill("solid", fgColor=BRAND_ORANGE)

    recommended = _summary_row(outputs.summary_df, "2_programado_bruto_com_amortizacao")
    net = _summary_row(outputs.summary_df, "3_programado_liquido_caixa_lft")
    cash_base = _cash_base(outputs)
    metrics = [
        ("Despesa financeira bruta", recommended.get("despesa_financeira_bruta")),
        ("Gross-up de receita", recommended.get("receita_antecipacao_gross_up_sugerida")),
        ("Aplicação caixa/LFT", cash_base),
        ("Rendimento caixa/LFT", net.get("rendimento_caixa_lft")),
        ("Despesa líquida", net.get("despesa_financeira_liquida")),
        ("Saldo médio remunerado", recommended.get("saldo_base")),
        ("Linhas sem CDI+", recommended.get("linhas_sem_spread")),
    ]
    ws["A1"] = "Cloudwalk - custo financeiro FIDCs"
    ws["A1"].font = Font(size=18, bold=True, color=INK)
    ws["A2"] = "Memória gerencial com CDI mensal composto, amortizações documentais e waterfall de PL."
    ws["A2"].font = Font(size=10, color=MUTED)
    for idx, (label, value) in enumerate(metrics, start=4):
        ws.cell(idx, 1, label)
        ws.cell(idx, 2, _as_number(value))
        ws.cell(idx, 2).number_format = '#,##0'
    for row in ws.iter_rows(min_row=4, max_row=4 + len(metrics) - 1, min_col=1, max_col=2):
        for cell in row:
            cell.border = Border(bottom=thin)
            cell.alignment = Alignment(vertical="center")
    ws.column_dimensions["A"].width = 34
    ws.column_dimensions["B"].width = 18

    cost_by_fund = _cost_by_fund(outputs.line_df)
    _write_df(wb.create_sheet("Preço por Série"), _series_price_frame(outputs.line_df), header_fill)
    _write_df(wb.create_sheet("Memoria Fundo a Fundo"), _memory_frame(outputs.line_df), header_fill)
    _write_df(wb.create_sheet("Mensal"), outputs.monthly_df, header_fill)
    _write_df(wb.create_sheet("CDI mensal"), _cdi_frame(monthly_cdi_rates), header_fill)
    _write_df(wb.create_sheet("Caixa LFT"), _cash_lft_frame(outputs.ime_snapshot_df), header_fill)
    _write_df(wb.create_sheet("Custo por Fundo"), cost_by_fund, header_fill)
    _write_df(wb.create_sheet("Lacunas CDI"), outputs.missing_inputs_df, header_fill)
    if pl_waterfall is not None:
        _write_df(wb.create_sheet("PL Waterfall"), pl_waterfall.by_fund_df, header_fill)
        _write_df(wb.create_sheet("PL Steps"), pl_waterfall.steps_df, header_fill)

    if not cost_by_fund.empty:
        chart_ws = wb["Custo por Fundo"]
        chart = BarChart()
        chart.title = "Custo bruto por fundo"
        chart.y_axis.title = "R$"
        chart.x_axis.title = "FIDC"
        data = Reference(chart_ws, min_col=3, min_row=1, max_row=len(cost_by_fund.index) + 1)
        cats = Reference(chart_ws, min_col=1, min_row=2, max_row=len(cost_by_fund.index) + 1)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.height = 7
        chart.width = 14
        ws.add_chart(chart, "D4")

    monthly = outputs.monthly_df.copy()
    if not monthly.empty:
        monthly_summary = monthly.groupby("mes", as_index=False)["custo_programado_bruto"].sum()
        temp = wb.create_sheet("_chart_mensal")
        temp.sheet_state = "hidden"
        _write_df(temp, monthly_summary, header_fill)
        chart = LineChart()
        chart.title = "Custo bruto mensal"
        chart.y_axis.title = "R$"
        chart.x_axis.title = "Mês"
        data = Reference(temp, min_col=2, min_row=1, max_row=len(monthly_summary.index) + 1)
        cats = Reference(temp, min_col=1, min_row=2, max_row=len(monthly_summary.index) + 1)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.height = 7
        chart.width = 14
        ws.add_chart(chart, "D20")

    for sheet in wb.worksheets:
        sheet.freeze_panes = "A2"
        for column_cells in sheet.columns:
            length = max(len(str(cell.value or "")) for cell in column_cells[:80])
            sheet.column_dimensions[get_column_letter(column_cells[0].column)].width = min(max(length + 2, 10), 42)
    ws.freeze_panes = "A4"
    buffer = BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def build_cloudwalk_financial_cost_pptx_bytes(
    outputs: FinancialCostOutputs,
    *,
    pl_waterfall: CloudwalkPlWaterfall | None = None,
) -> bytes:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    recommended = _summary_row(outputs.summary_df, "2_programado_bruto_com_amortizacao")
    net = _summary_row(outputs.summary_df, "3_programado_liquido_caixa_lft")
    cash_base = _cash_base(outputs)
    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Custo financeiro bruto estimado em R$ %.1f bi" % (_as_number(recommended["despesa_financeira_bruta"]) / 1e9))
    _add_subtitle(
        slide,
        "Despesa programada com amortizações documentais e CDI mensal composto; gross-up gerencial replica a despesa bruta em receita.",
    )
    metrics = [
        ("Despesa bruta", recommended["despesa_financeira_bruta"]),
        ("Custo líquido", net["despesa_financeira_liquida"]),
        ("Saldo médio", recommended["saldo_base"]),
        ("Aplicação caixa/LFT", cash_base),
        ("Rendimento caixa/LFT", net["rendimento_caixa_lft"]),
    ]
    for i, (label, value) in enumerate(metrics):
        _add_metric(slide, Inches(0.55 + i * 2.52), Inches(1.25), label, _fmt_money(value), width=Inches(2.28))

    cost_by_fund = _cost_by_fund(outputs.line_df).head(10)
    if not cost_by_fund.empty:
        chart_data = CategoryChartData()
        chart_data.categories = cost_by_fund["short_name"].tolist()
        chart_data.add_series("Custo bruto", cost_by_fund["custo_programado_bruto"].tolist())
        graphic = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.75),
            Inches(2.35),
            Inches(7.1),
            Inches(4.45),
            chart_data,
        )
        chart = graphic.chart
        chart.has_legend = False
        chart.value_axis.tick_labels.font.size = Pt(8)
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        chart.plots[0].data_labels.number_format = '#,##0'

    _add_table(
        slide,
        _memory_frame(outputs.line_df).head(9)[["short_name", "classe", "spread_cdi_plus_aa", "saldo_medio_programado", "custo_programado_bruto"]],
        Inches(8.2),
        Inches(2.35),
        Inches(4.55),
        Inches(4.45),
    )

    price_frame = _series_price_frame(outputs.line_df)
    if not price_frame.empty:
        slide = prs.slides.add_slide(blank)
        _add_title(slide, "Preço por série dos FIDCs")
        _add_subtitle(slide, "CDI+ a.a. usado no motor, com indicação das séries inputadas manualmente.")
        midpoint = (len(price_frame.index) + 1) // 2
        left_table = price_frame.head(midpoint)
        right_table = price_frame.iloc[midpoint:]
        cols = ["FIDC", "Série", "CDI+ a.a.", "Fonte", "Volume", "Custo"]
        _add_table(slide, left_table[cols], Inches(0.55), Inches(1.35), Inches(6.05), Inches(5.55), max_rows=13)
        if not right_table.empty:
            _add_table(slide, right_table[cols], Inches(6.82), Inches(1.35), Inches(5.95), Inches(5.55), max_rows=13)

    if pl_waterfall is not None:
        slide = prs.slides.add_slide(blank)
        _add_title(slide, "Waterfall de PL dos fundos")
        _add_subtitle(slide, "PL final = PL inicial + captações - resgates - amortizações + accrual/rentabilidade residual.")
        _draw_waterfall(slide, pl_waterfall.steps_df, Inches(0.65), Inches(1.65), Inches(7.25), Inches(4.95))
        _add_table(
            slide,
            pl_waterfall.by_fund_df[["short_name", "pl_inicial", "captacoes", "amortizacoes", "accrual_rentabilidade_residual", "pl_final"]].head(10),
            Inches(8.2),
            Inches(1.65),
            Inches(4.55),
            Inches(4.95),
        )

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _write_df(ws: Any, frame: pd.DataFrame, header_fill: Any) -> None:
    from openpyxl.styles import Alignment, Font
    from openpyxl.utils.dataframe import dataframe_to_rows

    clean = frame.copy() if frame is not None else pd.DataFrame()
    if clean.empty:
        clean = pd.DataFrame({"status": ["Sem dados"]})
    for row in dataframe_to_rows(clean, index=False, header=True):
        ws.append(row)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = Font(color="FFFFFF", bold=True)
        cell.alignment = Alignment(horizontal="center")
    for row in ws.iter_rows(min_row=2):
        for cell in row:
            if isinstance(cell.value, float):
                cell.number_format = '#,##0.00'


def _summary_row(summary: pd.DataFrame, estimativa: str) -> pd.Series:
    return summary.loc[summary["estimativa"].eq(estimativa)].iloc[0]


def _memory_frame(line_df: pd.DataFrame) -> pd.DataFrame:
    frame = line_df[line_df["included_in_cost"].fillna(False).astype(bool)].copy()
    if frame.empty:
        return pd.DataFrame()
    frame["short_name"] = frame["fund_name"].map(_short_name)
    return frame[
        [
            "fund_name",
            "short_name",
            "cnpj",
            "classe",
            "class_macro",
            "issue_date",
            "spread_cdi_plus_aa",
            "saldo_inicio_periodo",
            "saldo_medio_programado",
            "amortizacao_no_periodo",
            "custo_programado_bruto",
            "amortization_convention",
            "warnings",
            "source",
        ]
    ].sort_values("custo_programado_bruto", ascending=False)


def _series_price_frame(line_df: pd.DataFrame) -> pd.DataFrame:
    frame = line_df[line_df["included_in_cost"].fillna(False).astype(bool)].copy()
    if frame.empty:
        return pd.DataFrame(columns=["FIDC", "Série", "Classe", "CDI+ a.a.", "Fonte", "Volume", "Saldo médio", "Custo", "Chave"])
    frame["FIDC"] = frame["fund_name"].map(_short_name)
    frame["Série"] = frame["classe"]
    frame["Classe"] = frame["class_macro"].map(_class_label)
    frame["CDI+ a.a."] = pd.to_numeric(frame["spread_cdi_plus_aa"], errors="coerce").map(_fmt_percent)
    frame["Fonte"] = frame["spread_source"].map(_spread_source_label)
    frame["Volume"] = pd.to_numeric(frame["volume_emitido"], errors="coerce").map(_fmt_money)
    frame["Saldo médio"] = pd.to_numeric(frame["saldo_medio_programado"], errors="coerce").map(_fmt_money)
    frame["Custo"] = pd.to_numeric(frame["custo_programado_bruto"], errors="coerce").map(_fmt_money)
    frame["Chave"] = frame["line_key"]
    return frame[["FIDC", "Série", "Classe", "CDI+ a.a.", "Fonte", "Volume", "Saldo médio", "Custo", "Chave"]].sort_values(
        ["FIDC", "Série"]
    )


def _cost_by_fund(line_df: pd.DataFrame) -> pd.DataFrame:
    frame = line_df[line_df["included_in_cost"].fillna(False).astype(bool)].copy()
    if frame.empty:
        return pd.DataFrame(columns=["short_name", "saldo_medio_programado", "custo_programado_bruto"])
    frame["short_name"] = frame["fund_name"].map(_short_name)
    grouped = (
        frame.groupby("short_name", as_index=False)
        .agg(
            saldo_medio_programado=("saldo_medio_programado", "sum"),
            custo_programado_bruto=("custo_programado_bruto", "sum"),
        )
        .sort_values("custo_programado_bruto", ascending=False)
    )
    return grouped


def _cash_lft_frame(ime_df: pd.DataFrame) -> pd.DataFrame:
    if ime_df.empty:
        return pd.DataFrame(
            columns=[
                "FIDC",
                "CNPJ",
                "Competência",
                "PL total",
                "Caixa",
                "Títulos públicos",
                "Recebíveis",
                "Caixa/Títulos reportado",
                "Proxy PL - recebíveis",
                "Aplicação caixa/LFT",
                "Método",
                "Rendimento estimado",
                "Incluído",
                "Fonte",
            ]
        )
    frame = ime_df.copy()
    frame["FIDC"] = frame["fund_name"].map(_short_name)
    return frame[
        [
            "FIDC",
            "cnpj",
            "competencia",
            "pl_total",
            "caixa",
            "titulos_publicos",
            "recebiveis",
            "cash_like_reported",
            "cash_like_residual_proxy",
            "cash_like_caixa_lft",
            "cash_like_method",
            "rendimento_estimado_caixa_lft",
            "included",
            "source",
        ]
    ].rename(
        columns={
            "cnpj": "CNPJ",
            "competencia": "Competência",
            "pl_total": "PL total",
            "caixa": "Caixa",
            "titulos_publicos": "Títulos públicos",
            "recebiveis": "Recebíveis",
            "cash_like_reported": "Caixa/Títulos reportado",
            "cash_like_residual_proxy": "Proxy PL - recebíveis",
            "cash_like_caixa_lft": "Aplicação caixa/LFT",
            "cash_like_method": "Método",
            "rendimento_estimado_caixa_lft": "Rendimento estimado",
            "included": "Incluído",
            "source": "Fonte",
        }
    ).sort_values(["FIDC", "Competência"])


def _cash_base(outputs: FinancialCostOutputs) -> float:
    frame = outputs.ime_snapshot_df
    if frame.empty:
        return 0.0
    return _sum_numeric(frame.get("cash_like_caixa_lft"))


def _cdi_frame(monthly_cdi_rates: tuple[B3CdiMonthlyRate, ...]) -> pd.DataFrame:
    return pd.DataFrame(
        [
            {
                "mes": item.mes,
                "cdi_mensal": item.cdi_mensal,
                "cdi_aa_equivalente": item.cdi_aa_equivalente,
                "dias_uteis": item.dias_uteis,
                "data_inicio": item.data_inicio.isoformat(),
                "data_fim": item.data_fim.isoformat(),
                "source": item.source,
            }
            for item in monthly_cdi_rates
        ]
    )


def _add_title(slide: Any, text: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.45))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(INK)


def _add_subtitle(slide: Any, text: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.57), Inches(0.82), Inches(11.8), Inches(0.35))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(10.5)
    run.font.color.rgb = RGBColor.from_string(MUTED)


def _add_metric(slide: Any, x: Any, y: Any, label: str, value: str, *, width: Any | None = None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, width or Inches(2.75), Inches(0.72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(247, 248, 250)
    shape.line.color.rgb = RGBColor.from_string(GRID)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = label + "\n"
    r.font.size = Pt(8)
    r.font.color.rgb = RGBColor.from_string(MUTED)
    r2 = p.add_run()
    r2.text = value
    r2.font.size = Pt(15)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor.from_string(INK)


def _add_table(slide: Any, frame: pd.DataFrame, x: Any, y: Any, w: Any, h: Any, *, max_rows: int = 11) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    clean = frame.copy()
    rows = min(len(clean.index) + 1, max_rows)
    cols = len(clean.columns)
    table = slide.shapes.add_table(rows, cols, x, y, w, h).table
    for col_idx, column in enumerate(clean.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(column).replace("_", " ").title()
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(INK)
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(6.5)
    for row_idx, (_, row) in enumerate(clean.head(rows - 1).iterrows(), start=1):
        for col_idx, column in enumerate(clean.columns):
            cell = table.cell(row_idx, col_idx)
            cell.text = _fmt_cell(row.get(column))
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(6.2)
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string(INK)


def _draw_waterfall(slide: Any, steps_df: pd.DataFrame, x: Any, y: Any, w: Any, h: Any) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    values = steps_df["valor"].astype(float).tolist()
    labels = steps_df["etapa"].tolist()
    cumulative = []
    running = 0.0
    for idx, value in enumerate(values):
        if idx == 0:
            running = value
        elif idx == len(values) - 1:
            running = value
        else:
            running += value
        cumulative.append(running)
    max_abs = max([abs(item) for item in cumulative + values] + [1.0])
    baseline = y + h * 0.82
    chart_h = h * 0.68
    bar_w = w / max(len(values), 1) * 0.55
    gap = w / max(len(values), 1)
    running_before = 0.0
    for idx, (label, value) in enumerate(zip(labels, values)):
        cx = x + gap * idx + gap * 0.22
        if idx == 0 or idx == len(values) - 1:
            top_value = value
            base_value = 0.0
            color = INK
        else:
            base_value = running_before
            top_value = running_before + value
            color = GREEN if value >= 0 else RED
        y1 = baseline - (max(base_value, top_value) / max_abs) * chart_h
        y2 = baseline - (min(base_value, top_value) / max_abs) * chart_h
        height = max(y2 - y1, Inches(0.05))
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cx, y1, bar_w, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor.from_string(color)
        rect.line.color.rgb = RGBColor.from_string(color)
        label_box = slide.shapes.add_textbox(cx - Inches(0.12), baseline + Inches(0.08), bar_w + Inches(0.25), Inches(0.4))
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(6.5)
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        value_box = slide.shapes.add_textbox(cx - Inches(0.12), y1 - Inches(0.3), bar_w + Inches(0.25), Inches(0.25))
        value_box.text_frame.text = _fmt_money(value)
        value_box.text_frame.paragraphs[0].font.size = Pt(6.5)
        value_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if idx == 0:
            running_before = value
        elif idx < len(values) - 1:
            running_before += value


def _as_number(value: Any) -> float:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return 0.0
    return 0.0 if pd.isna(number) else number


def _fmt_money(value: Any) -> str:
    number = _as_number(value)
    sign = "-" if number < 0 else ""
    number = abs(number)
    if number >= 1e9:
        return f"{sign}R$ {number / 1e9:.1f} bi"
    if number >= 1e6:
        return f"{sign}R$ {number / 1e6:.0f} mi"
    return f"{sign}R$ {number:,.0f}"


def _fmt_percent(value: Any) -> str:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return ""
    if pd.isna(number):
        return ""
    return f"{number * 100:.2f}%"


def _fmt_cell(value: Any) -> str:
    if isinstance(value, float) or isinstance(value, int):
        if abs(float(value)) > 1000:
            return _fmt_money(value)
        return f"{float(value):.2%}" if abs(float(value)) < 1 else f"{float(value):,.2f}"
    return "" if value is None or pd.isna(value) else str(value)[:60]


def _short_name(name: str) -> str:
    text = str(name or "").upper()
    for token in [
        "CLOUDWALK",
        "FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS",
        "FIDC",
        "SEGMENTO MEIOS DE PAGAMENTO",
        "RESPONSABILIDADE LIMITADA",
    ]:
        text = text.replace(token, " ")
    return " ".join(text.title().split()) or str(name or "")


def _class_label(value: Any) -> str:
    mapping = {"senior": "Sênior", "mezzanino": "Mezanino", "subordinada": "Subordinada"}
    return mapping.get(str(value or "").strip().lower(), str(value or ""))


def _spread_source_label(value: Any) -> str:
    text = str(value or "")
    if text.startswith("override:"):
        return "Manual"
    if text == "curadoria:remuneracao":
        return "Documento"
    if text == "pendente":
        return "Pendente"
    return text


def _sum_numeric(series: Any) -> float:
    if series is None:
        return 0.0
    return float(pd.to_numeric(series, errors="coerce").fillna(0.0).sum())
