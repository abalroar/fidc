from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
import math
from pathlib import Path
from typing import Sequence
from zoneinfo import ZoneInfo

import pandas as pd
from PIL import Image, ImageDraw, ImageFont

from services.export_chart_labels import (
    DEFAULT_LABEL_FONT_SIZE_PT,
    choose_export_label_policy,
    format_export_label,
)
from services.fundonet_dashboard import FundonetDashboardData


ORANGE = "#EC7000"
BLACK = "#1F1F1F"
DARK_GRAY = "#4D4D4D"
MID_GRAY = "#757575"
GRID_GRAY = "#E0E0E0"
CHART_GRID_GRAY = "#E8E8E8"
SOFT_GRAY = "#F7F7F7"
NEGATIVE_RED = "#CC0000"
WHITE = "#ffffff"
SERIES_COLORS = [ORANGE, BLACK, "#757575", "#BDBDBD", "#4D4D4D"]
OVER_PPT_COLORS = [ORANGE, BLACK, "#757575", "#BDBDBD", "#4D4D4D", "#E0E0E0"]
AGING_PPT_COLORS = [
    ORANGE,
    BLACK,
    "#757575",
    "#BDBDBD",
    "#4D4D4D",
    "#E0E0E0",
    "#8C8C8C",
    "#2F2F2F",
    "#A6A6A6",
    "#D0D0D0",
]
COVERAGE_LINE_COLOR = ORANGE
FONT_FAMILY = "Calibri"

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MARGIN_LEFT_IN = 0.45
MARGIN_RIGHT_IN = 0.45
CONTENT_WIDTH_IN = SLIDE_WIDTH_IN - MARGIN_LEFT_IN - MARGIN_RIGHT_IN

TITLE_SIZE = 24
SECTION_SIZE = 11
SUBTITLE_SIZE = 10
BODY_SIZE = 9
LABEL_SIZE = 8
AXIS_SIZE = 9
FOOTER_SIZE = 8
CARD_VALUE_SIZE = 24
SLIDE_RENDER_DPI = 170


@dataclass(frozen=True)
class MoneyScale:
    divisor: float
    suffix: str
    label: str


def build_dashboard_pptx_bytes(
    dashboard: FundonetDashboardData,
    *,
    generated_at: datetime | None = None,
    requested_period_label: str | None = None,
) -> bytes:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("Dependência python-pptx não instalada.") from exc

    sao_paulo_tz = ZoneInfo("America/Sao_Paulo")
    if generated_at is None:
        generated_at = datetime.now(sao_paulo_tz)
    elif generated_at.tzinfo is None:
        generated_at = generated_at.replace(tzinfo=sao_paulo_tz)
    else:
        generated_at = generated_at.astimezone(sao_paulo_tz)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]

    def rgb(hex_color: str) -> RGBColor:
        value = str(hex_color or BLACK).strip().lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    def add_textbox(  # noqa: ANN202
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        *,
        size: int,
        bold: bool = False,
        italic: bool = False,
        color: str = BLACK,
        align=PP_ALIGN.LEFT,
        word_wrap: bool = True,
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = word_wrap
        frame.auto_size = MSO_AUTO_SIZE.NONE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT_FAMILY
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color)
        return box

    def style_shape_border(shape, *, line_color: str = GRID_GRAY, line_width_pt: float = 0.8) -> None:  # noqa: ANN001
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_width_pt)

    def add_panel(slide, left: float, top: float, width: float, height: float):  # noqa: ANN001
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(WHITE)
        style_shape_border(shape)
        return shape

    def add_card(slide, left: float, top: float, width: float, height: float, label: str, value: str, note: str = "") -> None:
        panel = add_panel(slide, left, top, width, height)
        style_shape_border(panel, line_color=GRID_GRAY, line_width_pt=1.0)
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.055),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = rgb(ORANGE)
        accent.line.fill.background()
        value_color = NEGATIVE_RED if str(value).strip().startswith("-") else BLACK
        add_textbox(
            slide,
            left + 0.16,
            top + 0.16,
            width - 0.32,
            0.20,
            label.upper(),
            size=LABEL_SIZE,
            bold=False,
            color=MID_GRAY,
        )
        add_textbox(
            slide,
            left + 0.16,
            top + 0.41,
            width - 0.32,
            0.34,
            value,
            size=CARD_VALUE_SIZE,
            bold=True,
            color=value_color,
        )
        if note:
            add_textbox(
                slide,
                left + 0.16,
                top + 0.82,
                width - 0.32,
                0.18,
                note,
                size=FOOTER_SIZE,
                italic=True,
                color=MID_GRAY,
            )

    def add_empty_state_panel(
        slide,
        *,
        title: str,
        message: str,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        panel = add_panel(slide, left, top, width, height)
        style_shape_border(panel, line_color=GRID_GRAY, line_width_pt=0.7)
        add_textbox(slide, left + 0.14, top + 0.12, width - 0.28, 0.18, title, size=SECTION_SIZE, bold=True, color=BLACK)
        add_textbox(
            slide,
            left + 0.26,
            top + (height / 2) - 0.28,
            width - 0.52,
            0.56,
            message,
            size=12,
            color=MID_GRAY,
            align=PP_ALIGN.CENTER,
        )

    def add_footer(slide, timestamp_text: str, page_number: int | None = None) -> None:  # noqa: ANN001
        separator = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(MARGIN_LEFT_IN),
            Inches(6.98),
            Inches(CONTENT_WIDTH_IN),
            Inches(0.004),
        )
        separator.fill.solid()
        separator.fill.fore_color.rgb = rgb(GRID_GRAY)
        separator.line.fill.background()
        add_textbox(
            slide,
            MARGIN_LEFT_IN,
            7.10,
            8.4,
            0.22,
            f"Fonte: Informe Mensal CVM    |    Gerado em: {timestamp_text}",
            size=FOOTER_SIZE,
            color=MID_GRAY,
        )
        if page_number is not None:
            add_textbox(
                slide,
                11.05,
                7.10,
                1.80,
                0.22,
                f"Página {page_number}",
                size=FOOTER_SIZE,
                color=MID_GRAY,
                align=PP_ALIGN.RIGHT,
            )

    def add_cover_slide(
        *,
        title: str,
        subtitle_text: str,
        scope_text: str,
    ) -> None:
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(BLACK)
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.70),
            Inches(2.35),
            Inches(0.07),
            Inches(1.55),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = rgb(ORANGE)
        accent.line.fill.background()
        add_textbox(slide, 0.95, 2.20, 10.60, 0.66, title, size=36, bold=True, color=WHITE)
        add_textbox(slide, 0.97, 2.94, 10.90, 0.20, scope_text.upper(), size=11, bold=True, color=WHITE)
        add_textbox(slide, 0.97, 3.26, 10.90, 0.20, subtitle_text, size=11, color=WHITE)
        add_textbox(slide, 0.97, 3.56, 10.90, 0.18, "Fonte: Informe Mensal CVM", size=9, color=GRID_GRAY)
        add_textbox(
            slide,
            9.15,
            6.82,
            3.45,
            0.22,
            "Toma Conta | Análise Institucional",
            size=9,
            color=WHITE,
            align=PP_ALIGN.RIGHT,
        )

    def add_divider_slide(*, number: str, title: str, subtitle_text: str) -> None:
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(BLACK)
        add_textbox(slide, 1.02, 2.66, 1.05, 0.34, f"{number} —", size=20, bold=True, color=ORANGE)
        add_textbox(slide, 2.05, 2.54, 8.95, 0.55, title, size=30, bold=True, color=WHITE)
        add_textbox(slide, 2.08, 3.18, 9.10, 0.28, subtitle_text, size=14, color=GRID_GRAY)

    def set_table_style(
        table,
        *,
        header_fill: str = BLACK,
        header_font_size: int = LABEL_SIZE,
        body_font_size: int = BODY_SIZE,
    ) -> None:  # noqa: ANN001
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                if row_idx == 0:
                    cell.fill.fore_color.rgb = rgb(header_fill)
                else:
                    cell.fill.fore_color.rgb = rgb(WHITE if row_idx % 2 else SOFT_GRAY)
                cell.text_frame.word_wrap = False
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
                    for run in paragraph.runs:
                        raw_text = str(run.text or "").strip()
                        is_negative = raw_text.startswith("-") or raw_text.startswith("−")
                        run.font.name = FONT_FAMILY
                        run.font.size = Pt(header_font_size if row_idx == 0 else body_font_size)
                        run.font.bold = row_idx == 0
                        run.font.color.rgb = rgb(WHITE if row_idx == 0 else NEGATIVE_RED if is_negative else BLACK)
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)

    def add_table(  # noqa: ANN202
        slide,
        df: pd.DataFrame,
        *,
        title: str | None,
        left: float,
        top: float,
        width: float,
        height: float,
        col_widths: Sequence[float] | None = None,
        max_rows: int = 30,
    ):
        if title:
            add_textbox(slide, left, top - 0.22, width, 0.18, title, size=SECTION_SIZE, bold=True, color=BLACK)
        frame = df.copy()
        row_height_floor = 0.24
        max_rows_fit = max(1, int(max((height - 0.30), row_height_floor) / row_height_floor) - 1)
        effective_max_rows = min(max_rows, max_rows_fit) if max_rows else max_rows_fit
        if effective_max_rows and len(frame) > effective_max_rows:
            frame = frame.iloc[:effective_max_rows]
        rows = max(len(frame.index), 1)
        cols = max(len(frame.columns), 1)
        table = slide.shapes.add_table(rows + 1, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
        total_rows = rows + 1
        row_height = height / total_rows if total_rows > 0 else height
        for row in table.rows:
            row.height = Inches(row_height)
        body_font_size = BODY_SIZE
        header_font_size = LABEL_SIZE
        if rows >= 8:
            body_font_size = 8
            header_font_size = 8
        if rows >= 12:
            body_font_size = 7
            header_font_size = 7
        # Always scale col_widths so they sum exactly to `width`, preventing horizontal overflow
        if col_widths and len(col_widths) == cols:
            total_cw = sum(col_widths)
            scale = width / total_cw if total_cw > 0 else 1.0
            for idx, col_width in enumerate(col_widths):
                table.columns[idx].width = Inches(col_width * scale)
        else:
            even = width / cols
            for col in table.columns:
                col.width = Inches(even)
        _MAX_CELL = 100
        if frame.empty:
            table.cell(0, 0).text = "Sem dados"
            table.cell(1, 0).text = "-"
            set_table_style(table, header_font_size=header_font_size, body_font_size=body_font_size)
            return table
        for col_idx, column in enumerate(frame.columns):
            table.cell(0, col_idx).text = str(column)
        for row_idx, (_, row) in enumerate(frame.iterrows(), start=1):
            for col_idx, value in enumerate(row):
                cell_text = str(value)
                if len(cell_text) > _MAX_CELL:
                    cell_text = cell_text[:_MAX_CELL] + "…"
                table.cell(row_idx, col_idx).text = cell_text
        set_table_style(table, header_font_size=header_font_size, body_font_size=body_font_size)
        return table

    def add_manual_legend(
        slide,
        items: Sequence[tuple[str, str]],
        *,
        left: float,
        top: float,
        max_width: float,
        font_size: int = LABEL_SIZE,
    ) -> None:
        cursor_left = left
        for label, color in items:
            label_width = max(0.75, min(2.6, 0.18 + (len(label) * 0.075)))
            if cursor_left + label_width + 0.35 > left + max_width:
                top += 0.18
                cursor_left = left
            marker = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cursor_left),
                Inches(top + 0.02),
                Inches(0.18),
                Inches(0.05),
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = rgb(color)
            marker.line.fill.background()
            add_textbox(
                slide,
                cursor_left + 0.22,
                top - 0.02,
                label_width,
                0.14,
                label,
                size=font_size,
                color=DARK_GRAY,
                word_wrap=False,
            )
            cursor_left += label_width + 0.45

    def _data_label_position(position_name: str):  # noqa: ANN202
        mapping = {
            "above": XL_DATA_LABEL_POSITION.ABOVE,
            "outside_end": XL_DATA_LABEL_POSITION.OUTSIDE_END,
            "center": XL_DATA_LABEL_POSITION.CENTER,
            "inside_end": XL_DATA_LABEL_POSITION.INSIDE_END,
        }
        return mapping[position_name]

    def add_chart(  # noqa: ANN202
        slide,
        *,
        title: str | None,
        chart_type,
        categories: Sequence[str],
        series_map: Sequence[tuple[str, Sequence[float]]],
        left: float,
        top: float,
        width: float,
        height: float,
        number_format: str,
        percent_axis: bool = False,
        money_axis: bool = False,
        gap_width: int | None = None,
        overlap: int | None = None,
        label_position: str = "above",
        label_color: str = BLACK,
        label_font_size: int = LABEL_SIZE,
        series_colors: Sequence[str] | None = None,
        title_suffix: str = "",
        value_max: float | None = None,
        value_min: float | None = None,
        show_data_labels: bool = True,
        show_legend: bool = False,
        legend_position: str = "bottom",
    ):
        chart_data = CategoryChartData()
        chart_data.categories = list(categories)
        for series_name, values in series_map:
            chart_data.add_series(series_name, _sanitize_chart_series(values))
        chart_shape = slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data)
        chart = chart_shape.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = f"{title}{title_suffix}"
            title_paragraph = chart.chart_title.text_frame.paragraphs[0]
            title_paragraph.font.name = FONT_FAMILY
            title_paragraph.font.size = Pt(SECTION_SIZE)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = rgb(BLACK)
        try:
            chart.chart_area.format.line.fill.background()
        except Exception:  # noqa: BLE001
            pass
        try:
            chart.chart_area.format.fill.solid()
            chart.chart_area.format.fill.fore_color.rgb = rgb(WHITE)
        except Exception:  # noqa: BLE001
            pass
        try:
            chart.plot_area.format.fill.solid()
            chart.plot_area.format.fill.fore_color.rgb = rgb(WHITE)
        except Exception:  # noqa: BLE001
            pass
        chart.has_legend = show_legend
        if show_legend and chart.legend is not None:
            legend_positions = {
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "right": XL_LEGEND_POSITION.RIGHT,
                "top": XL_LEGEND_POSITION.TOP,
                "left": XL_LEGEND_POSITION.LEFT,
            }
            chart.legend.position = legend_positions.get(legend_position, XL_LEGEND_POSITION.BOTTOM)
            if hasattr(chart.legend, "include_in_layout"):
                chart.legend.include_in_layout = True
            try:
                chart.legend.font.name = FONT_FAMILY
                chart.legend.font.size = Pt(9)
                chart.legend.font.color.rgb = rgb(MID_GRAY)
            except Exception:  # noqa: BLE001
                pass

        plot = chart.plots[0]
        if gap_width is not None and hasattr(plot, "gap_width"):
            plot.gap_width = gap_width
        if overlap is not None and hasattr(plot, "overlap"):
            plot.overlap = overlap

        if chart_type == XL_CHART_TYPE.LINE_MARKERS:
            export_chart_kind = "line" if len(series_map) <= 2 else "multi_line"
        elif chart_type == XL_CHART_TYPE.COLUMN_STACKED:
            export_chart_kind = "stacked_bar"
        else:
            export_chart_kind = "bar"
        export_metric_kind = "money" if money_axis else "general_pct" if percent_axis else "number"
        export_policy = choose_export_label_policy(
            [values for _, values in series_map],
            chart_kind=export_chart_kind,
            metric_kind=export_metric_kind,
        )
        use_chart_level_labels = show_data_labels and export_policy.shows_all_points

        plot.has_data_labels = use_chart_level_labels
        if use_chart_level_labels:
            labels = plot.data_labels
            labels.show_value = True
            labels.number_format = number_format
            labels.position = _data_label_position(label_position)
            labels.font.name = FONT_FAMILY
            labels.font.size = Pt(max(label_font_size, DEFAULT_LABEL_FONT_SIZE_PT))
            labels.font.bold = True
            labels.font.color.rgb = rgb(label_color)

        if hasattr(chart, "category_axis"):
            chart.category_axis.tick_labels.font.name = FONT_FAMILY
            chart.category_axis.tick_labels.font.size = Pt(AXIS_SIZE)
            chart.category_axis.tick_labels.font.color.rgb = rgb(MID_GRAY)
            chart.category_axis.format.line.color.rgb = rgb(MID_GRAY)
        if hasattr(chart, "value_axis"):
            chart.value_axis.tick_labels.font.name = FONT_FAMILY
            chart.value_axis.tick_labels.font.size = Pt(AXIS_SIZE)
            chart.value_axis.tick_labels.font.color.rgb = rgb(MID_GRAY)
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = rgb(CHART_GRID_GRAY)
            chart.value_axis.has_minor_gridlines = False
            chart.value_axis.format.line.color.rgb = rgb(MID_GRAY)
            if value_min is not None:
                chart.value_axis.minimum_scale = value_min
            if value_max is not None:
                chart.value_axis.maximum_scale = value_max
            if percent_axis:
                if value_min is None:
                    chart.value_axis.minimum_scale = 0.0
                if value_max is None:
                    chart.value_axis.maximum_scale = 110.0
            if money_axis:
                if value_min is None:
                    chart.value_axis.minimum_scale = 0.0

        applied_colors = list(series_colors or SERIES_COLORS)
        for idx, series in enumerate(chart.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = rgb(applied_colors[idx % len(applied_colors)])
            series.format.line.color.rgb = rgb(applied_colors[idx % len(applied_colors)])
            if use_chart_level_labels and hasattr(series, "data_labels"):
                series.data_labels.show_value = True
                series.data_labels.number_format = number_format
                series.data_labels.position = _data_label_position(label_position)
                series.data_labels.font.name = FONT_FAMILY
                series.data_labels.font.size = Pt(max(label_font_size, DEFAULT_LABEL_FONT_SIZE_PT))
                series.data_labels.font.bold = True
                series.data_labels.font.color.rgb = rgb(label_color)

        if chart_type == XL_CHART_TYPE.LINE_MARKERS:
            for idx, series in enumerate(chart.series):
                _style_line_series(series, color=applied_colors[idx % len(applied_colors)], width_pt=2.2, marker_size=8)

        if show_data_labels and not use_chart_level_labels:
            for series_idx, point_indices in enumerate(export_policy.indices_by_series):
                if series_idx >= len(chart.series):
                    continue
                _, values = series_map[series_idx]
                for point_idx in point_indices:
                    if point_idx >= len(values):
                        continue
                    try:
                        point = chart.series[series_idx].points[point_idx]
                        label = point.data_label
                        label.position = _data_label_position(label_position)
                        label.has_text_frame = True
                        text_frame = label.text_frame
                        text_frame.clear()
                        run = text_frame.paragraphs[0].add_run()
                        run.text = format_export_label(
                            values[point_idx],
                            metric_kind=export_metric_kind,
                            percent_value=False,
                        )
                        run.font.name = FONT_FAMILY
                        run.font.size = Pt(export_policy.font_size_pt)
                        run.font.bold = True
                        if export_chart_kind in {"line", "multi_line"}:
                            run.font.color.rgb = rgb(applied_colors[series_idx % len(applied_colors)])
                        else:
                            run.font.color.rgb = rgb(label_color)
                    except Exception:  # noqa: BLE001
                        continue

        return chart

    def _first_axis_element(chart, axis_name: str):  # noqa: ANN202
        plot_area = chart._chartSpace.xpath("./c:chart/c:plotArea")
        if not plot_area:
            return None
        matches = plot_area[0].xpath(f"./c:{axis_name}")
        return matches[0] if matches else None

    def _set_axis_hidden(chart, axis_name: str, hidden: bool) -> None:  # noqa: ANN001
        axis = _first_axis_element(chart, axis_name)
        if axis is None:
            return
        delete_nodes = axis.xpath("./c:delete")
        if delete_nodes:
            delete_nodes[0].set("val", "1" if hidden else "0")

    def _set_value_axis_right(chart) -> None:  # noqa: ANN001
        val_axis = _first_axis_element(chart, "valAx")
        cat_axis = _first_axis_element(chart, "catAx")
        if val_axis is not None:
            ax_pos = val_axis.xpath("./c:axPos")
            if ax_pos:
                ax_pos[0].set("val", "r")
        if cat_axis is not None:
            crosses = cat_axis.xpath("./c:crosses")
            if crosses:
                crosses[0].set("val", "max")

    def _style_line_series(
        series,  # noqa: ANN001
        *,
        color: str,
        width_pt: float = 2.5,
        marker_size: int | None = None,
        dashed: bool = False,
        hide_marker: bool = False,
    ) -> None:
        series.format.line.color.rgb = rgb(color)
        series.format.line.width = Pt(width_pt)
        if dashed:
            series.format.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if hide_marker:
            series.marker.style = XL_MARKER_STYLE.NONE
        else:
            series.marker.style = XL_MARKER_STYLE.CIRCLE
            if marker_size is not None:
                series.marker.size = marker_size
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = rgb(color)
            series.marker.format.line.color.rgb = rgb(color)

    def _repel_label_positions(values: Sequence[float], min_gap: float) -> list[float]:
        indexed = sorted(enumerate(values), key=lambda item: item[1])
        if not indexed:
            return []
        adjusted: dict[int, float] = {}
        previous: float | None = None
        for idx, value in indexed:
            candidate = value if previous is None else max(value, previous + min_gap)
            adjusted[idx] = candidate
            previous = candidate
        return [adjusted[idx] for idx in range(len(values))]

    def add_line_end_labels(
        slide,
        *,
        labels: Sequence[str],
        values: Sequence[float],
        colors: Sequence[str],
        left: float,
        top: float,
        width: float,
        height: float,
        axis_max: float,
        axis_min: float = 0.0,
        font_size: int = 11,
        fill_colors: Sequence[str] | None = None,
        text_colors: Sequence[str] | None = None,
        side: str = "right",
    ) -> None:
        numeric_values = [float(value) for value in values]
        plot_top = top + 0.22
        plot_height = max(height - 0.58, 0.1)
        plot_right = left + width - 0.38
        plot_left = left + 0.62
        value_span = max(axis_max - axis_min, 1.0)
        raw_positions = [
            plot_top + plot_height * (1.0 - ((value - axis_min) / value_span))
            for value in numeric_values
        ]
        adjusted_positions = _repel_label_positions(raw_positions, min_gap=0.18)
        for idx, label in enumerate(labels):
            label_width = max(1.12, min(2.35, 0.50 + len(str(label)) * 0.12))
            label_left = max(plot_left, left + 0.08) if side == "left" else min(plot_right, left + width - label_width - 0.10)
            label_top = max(top + 0.02, adjusted_positions[idx] - 0.10)
            fill_color = fill_colors[idx % len(fill_colors)] if fill_colors else None
            if fill_color:
                badge = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    Inches(label_left),
                    Inches(label_top - 0.02),
                    Inches(label_width),
                    Inches(0.30),
                )
                badge.fill.solid()
                badge.fill.fore_color.rgb = rgb(fill_color)
                badge.line.fill.background()
            add_textbox(
                slide,
                label_left,
                label_top,
                label_width,
                0.22,
                label,
                size=font_size,
                bold=True,
                color=(text_colors[idx % len(text_colors)] if text_colors else colors[idx % len(colors)]),
                align=PP_ALIGN.CENTER,
            )

    def add_overlay_combo_credit_chart(
        slide,
        *,
        title: str,
        categories: Sequence[str],
        bar_series_map: Sequence[tuple[str, Sequence[float]]],
        line_series_map: Sequence[tuple[str, Sequence[float]]],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        title_height = 0.22
        legend_top = top + height - 0.18
        chart_top = top + title_height
        chart_height = max(height - title_height - 0.24, 0.8)
        bar_values = [value for _, values in bar_series_map for value in values]
        line_values = [value for _, values in line_series_map for value in values]
        bar_axis_max = _dashboard_percent_axis_upper(bar_values, max_cap=140.0)
        max_line_value = max(
            [float(value) for value in line_values if _to_float(value) is not None],
            default=100.0,
        )
        line_axis_max = max(max_line_value * 1.12, max_line_value + 20.0, 120.0)

        add_textbox(slide, left, top - 0.02, width, 0.18, title, size=SECTION_SIZE, bold=True, color=BLACK)
        add_manual_legend(
            slide,
            [
                ("Inadimplência", ORANGE),
                ("Provisão", BLACK),
                ("Cobertura", COVERAGE_LINE_COLOR),
            ],
            left=left,
            top=legend_top,
            max_width=width,
        )

        add_chart(
            slide,
            title=None,
            chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
            categories=categories,
            series_map=bar_series_map,
            left=left,
            top=chart_top,
            width=width,
            height=chart_height,
            number_format='0.00"%"',
            percent_axis=True,
            gap_width=42,
            label_position="outside_end",
            label_font_size=8,
            show_legend=False,
            value_max=bar_axis_max,
            show_data_labels=True,
            series_colors=[ORANGE, BLACK],
        )

        plot_left = left + 0.78
        plot_right = left + width - 1.02
        plot_top = chart_top + 0.34
        plot_bottom = chart_top + chart_height - 0.56
        plot_height = max(plot_bottom - plot_top, 0.1)
        plot_width = max(plot_right - plot_left, 0.1)
        x_positions = [
            plot_left + (plot_width * ((idx + 0.5) / max(len(categories), 1)))
            for idx in range(len(categories))
        ]

        for tick_idx in range(5):
            tick_value = line_axis_max * (tick_idx / 4)
            y_pos = plot_bottom - (plot_height * (tick_value / line_axis_max))
            add_textbox(
                slide,
                plot_right + 0.08,
                y_pos - 0.06,
                0.72,
                0.12,
                _format_percent(tick_value),
                size=AXIS_SIZE,
                color=COVERAGE_LINE_COLOR,
                align=PP_ALIGN.LEFT,
            )
        add_textbox(
            slide,
            plot_right + 0.02,
            plot_top - 0.18,
            0.92,
            0.16,
            "Cobertura (%)",
            size=LABEL_SIZE,
            bold=True,
            color=COVERAGE_LINE_COLOR,
            align=PP_ALIGN.LEFT,
        )

        coverage_values = []
        parity_values = []
        for series_name, values in line_series_map:
            if series_name == "Cobertura":
                coverage_values = [float(_to_float(value) or 0.0) for value in values]
            elif series_name.startswith("100%"):
                parity_values = [float(_to_float(value) or 0.0) for value in values]

        if parity_values:
            parity_y = plot_bottom - (plot_height * (min(parity_values[0], line_axis_max) / line_axis_max))
            parity_shape = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(plot_left),
                Inches(parity_y),
                Inches(plot_right),
                Inches(parity_y),
            )
            parity_shape.line.color.rgb = rgb(MID_GRAY)
            parity_shape.line.width = Pt(1.1)
            parity_shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        if coverage_values:
            raw_label_positions = []
            marker_centers: list[tuple[float, float]] = []
            for x_pos, raw_value in zip(x_positions, coverage_values, strict=False):
                clamped_value = max(min(raw_value, line_axis_max), 0.0)
                y_pos = plot_bottom - (plot_height * (clamped_value / line_axis_max))
                marker_centers.append((x_pos, y_pos))
                raw_label_positions.append(y_pos - 0.16)

            adjusted_label_positions = _repel_label_positions(raw_label_positions, min_gap=0.18)
            for idx in range(len(marker_centers) - 1):
                x1, y1 = marker_centers[idx]
                x2, y2 = marker_centers[idx + 1]
                segment = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(x1),
                    Inches(y1),
                    Inches(x2),
                    Inches(y2),
                )
                segment.line.color.rgb = rgb(COVERAGE_LINE_COLOR)
                segment.line.width = Pt(2.4)

            marker_w = 0.10
            marker_h = 0.10
            for idx, ((x_pos, y_pos), raw_value) in enumerate(zip(marker_centers, coverage_values, strict=False)):
                marker = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    Inches(x_pos - (marker_w / 2)),
                    Inches(y_pos - (marker_h / 2)),
                    Inches(marker_w),
                    Inches(marker_h),
                )
                marker.fill.solid()
                marker.fill.fore_color.rgb = rgb(COVERAGE_LINE_COLOR)
                marker.line.color.rgb = rgb(COVERAGE_LINE_COLOR)
                label_top = max(plot_top - 0.02, adjusted_label_positions[idx])
                add_textbox(
                    slide,
                    x_pos - 0.24,
                    label_top,
                    0.48,
                    0.14,
                    _format_percent(raw_value),
                    size=DEFAULT_LABEL_FONT_SIZE_PT,
                    bold=True,
                    color=COVERAGE_LINE_COLOR,
                    align=PP_ALIGN.CENTER,
                )

    def add_compounding_waterfall_chart(
        slide,
        *,
        categories: Sequence[str],
        step_values: Sequence[float],
        total_value: float,
        left: float,
        top: float,
        width: float,
        height: float,
        number_format: str,
        value_max: float | None,
        title_suffix: str = "",
        label_font_size: int = 9,
    ) -> None:
        base_values: list[float] = []
        flow_values: list[float] = []
        total_series_values: list[float] = []
        cumulative = 0.0
        for idx, step_value in enumerate(step_values):
            numeric_value = float(step_value)
            base_values.append(cumulative)
            flow_values.append(numeric_value)
            total_series_values.append(0.0)
            cumulative += numeric_value
        base_values.append(0.0)
        flow_values.append(0.0)
        total_series_values.append(float(total_value))
        chart = add_chart(
            slide,
            title=f"Waterfall de vencimento da carteira{title_suffix}",
            chart_type=XL_CHART_TYPE.COLUMN_STACKED,
            categories=list(categories) + ["Total"],
            series_map=[
                ("Base", base_values),
                ("Fluxo", flow_values),
                ("Total", total_series_values),
            ],
            left=left,
            top=top,
            width=width,
            height=height,
            number_format=number_format,
            money_axis=True,
            gap_width=36,
            overlap=100,
            label_position="outside_end",
            label_font_size=label_font_size,
            value_max=value_max,
            show_data_labels=True,
            show_legend=False,
            series_colors=[WHITE, ORANGE, BLACK],
        )
        base_series = chart.series[0]
        base_series.format.fill.background()
        base_series.format.line.fill.background()
        if hasattr(base_series, "data_labels"):
            base_series.data_labels.show_value = False
        for idx in range(1, len(chart.series)):
            series = chart.series[idx]
            series.data_labels.font.bold = True
            series.data_labels.font.size = Pt(max(label_font_size, DEFAULT_LABEL_FONT_SIZE_PT))
            series.data_labels.font.color.rgb = rgb(BLACK if idx == 2 else DARK_GRAY)

    def add_picture_bytes(slide, image_bytes: bytes, *, left: float, top: float, width: float, height: float) -> None:  # noqa: ANN001
        slide.shapes.add_picture(BytesIO(image_bytes), Inches(left), Inches(top), Inches(width), Inches(height))

    timestamp_text = f"{generated_at.strftime('%d/%m/%Y %H:%M')} GMT-3"
    title_fund = _fund_title(dashboard)
    is_portfolio_scope = str(dashboard.fund_info.get("aggregation_scope") or "").strip().lower() == "portfolio"
    scope_label = "Carteira agregada" if is_portfolio_scope else "FIDC"
    data_base_label = _format_competencia(dashboard.latest_competencia)
    requested_period_text = (
        f"  |  Janela solicitada: {requested_period_label}"
        if requested_period_label and requested_period_label != dashboard.fund_info.get("periodo_analisado")
        else ""
    )
    subtitle = (
        f"{data_base_label}"
        f"  ·  {dashboard.fund_info.get('periodo_analisado', 'N/D')}"
        f"{requested_period_text}"
    )

    def add_slide_header(slide, section_title: str) -> None:  # noqa: ANN001
        add_textbox(slide, 0.50, 0.15, 9.80, 0.42, section_title, size=TITLE_SIZE, bold=True, color=BLACK)
        add_textbox(slide, 0.50, 0.60, 10.8, 0.20, title_fund, size=SUBTITLE_SIZE, bold=True, color=ORANGE)
        add_textbox(slide, 0.50, 0.82, 11.1, 0.18, subtitle, size=FOOTER_SIZE, color=MID_GRAY)
        chip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(10.65),
            Inches(0.16),
            Inches(2.20),
            Inches(0.30),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb(BLACK)
        chip.line.fill.background()
        add_textbox(
            slide,
            10.78,
            0.22,
            1.92,
            0.20,
            f"Data-base: {data_base_label}",
            size=10,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    def add_summary_cards(slide) -> None:  # noqa: ANN001
        summary = dashboard.summary
        specs = [
            ("Ativo total", _format_brl_compact(summary.get("ativos_totais")), "APLIC_ATIVO/VL_SOM_APLIC_ATIVO"),
            (
                "DCs totais",
                _format_brl_compact(summary.get("direitos_creditorios") or summary.get("inadimplencia_denominador")),
                "Direitos creditórios reportados",
            ),
            ("PL total", _format_brl_compact(summary.get("pl_total")), "Sênior + mezzanino + subordinada"),
            ("Vencidos", _format_brl_compact(summary.get("inadimplencia_total")), "Créditos vencidos/inadimplentes"),
            ("Cobertura de provisão", _format_percent(summary.get("cobertura_pct")), "Provisão / créditos vencidos"),
            (
                "Subordinação reportada",
                _format_percent(summary.get("subordinacao_pct")),
                "Mezzanino + subordinada residual / PL",
            ),
        ]
        gap = 0.26
        card_width = (CONTENT_WIDTH_IN - (2 * gap)) / 3
        card_height = 1.05
        start_top = 1.18
        row_gap = 0.34
        for idx, (label, value, note) in enumerate(specs):
            row = idx // 3
            col = idx % 3
            add_card(
                slide,
                MARGIN_LEFT_IN + col * (card_width + gap),
                start_top + row * (card_height + row_gap),
                card_width,
                card_height,
                label,
                value,
                note,
            )

    full_width = CONTENT_WIDTH_IN
    split_gap = 0.30
    left_wide_width = 7.45
    right_narrow_width = CONTENT_WIDTH_IN - left_wide_width - split_gap
    right_col_left = MARGIN_LEFT_IN + left_wide_width + split_gap
    top_row_top = 1.02

    page_number = 0

    def next_page() -> int:
        nonlocal page_number
        page_number += 1
        return page_number

    add_cover_slide(
        title=title_fund,
        subtitle_text=f"Data-base {data_base_label} | {dashboard.fund_info.get('periodo_analisado', 'N/D')}",
        scope_text=f"Visão executiva — {scope_label}",
    )
    next_page()
    add_divider_slide(
        number="01",
        title="Visão executiva",
        subtitle_text=(
            "PL, estrutura, crédito, liquidez e prazo agregados"
            if is_portfolio_scope
            else "PL, estrutura, crédito, liquidez e prazo do fundo"
        ),
    )
    next_page()

    # Slide 1 — FIDC summary
    slide = prs.slides.add_slide(blank)
    current_page = next_page()
    add_slide_header(slide, "Resumo da carteira" if is_portfolio_scope else "Resumo do FIDC")
    add_summary_cards(slide)
    add_footer(slide, timestamp_text, current_page)

    # Slide 2 — structure and capital
    slide = prs.slides.add_slide(blank)
    current_page = next_page()
    add_slide_header(slide, "Estrutura e capital")

    sub_df = _sort_competencia_frame(dashboard.subordination_history_df, ascending=True)
    if not sub_df.empty:
        sub_series = [("Subordinação reportada", _series_numeric(sub_df, "subordinacao_pct").fillna(0.0).tolist())]
        sub_chart = add_chart(
            slide,
            title="Subordinação reportada",
            chart_type=XL_CHART_TYPE.LINE_MARKERS,
            categories=_competencia_labels(sub_df["competencia"].tolist()),
            series_map=sub_series,
            left=MARGIN_LEFT_IN,
            top=top_row_top,
            width=full_width,
            height=2.05,
            number_format='0.0"%"',
            percent_axis=True,
            label_position="above",
            show_data_labels=len(sub_df["competencia"].drop_duplicates()) <= 12,
            label_font_size=9,
            value_max=_dashboard_percent_axis_upper(sub_series[0][1], max_cap=80.0),
        )
        _style_line_series(sub_chart.series[0], color=ORANGE, width_pt=2.8, marker_size=10)
        if len(sub_df["competencia"].drop_duplicates()) > 12:
            sub_values = [float(value or 0.0) for value in sub_series[0][1]]
            add_line_end_labels(
                slide,
                labels=[_format_percent(sub_values[-1])],
                values=[sub_values[-1]],
                colors=[ORANGE],
                left=MARGIN_LEFT_IN,
                top=top_row_top,
                width=full_width,
                height=2.05,
                axis_max=_dashboard_percent_axis_upper(sub_series[0][1], max_cap=80.0),
                fill_colors=[WHITE],
                text_colors=[ORANGE],
                side="right",
            )
        add_textbox(
            slide,
            MARGIN_LEFT_IN,
            top_row_top + 2.06,
            full_width,
            0.16,
            "Subordinação reportada = (PL mezzanino + PL subordinada residual) / PL total.",
            size=FOOTER_SIZE,
            color=MID_GRAY,
        )

    quota_values = _quota_pl_value_pivot(dashboard.quota_pl_history_df)
    if not quota_values.empty:
        quota_scale = _money_scale(quota_values.drop(columns=["competencia"]).stack())
        quota_series = [
            (column, _scale_values(quota_values[column], quota_scale).tolist())
            for column in quota_values.columns
            if column != "competencia"
        ]
        add_chart(
            slide,
            title="PL por tipo de cota",
            chart_type=XL_CHART_TYPE.COLUMN_STACKED,
            categories=_competencia_labels(quota_values["competencia"].tolist()),
            series_map=quota_series,
            left=MARGIN_LEFT_IN,
            top=3.32,
            width=full_width,
            height=3.35,
            number_format=_money_label_number_format(quota_scale),
            money_axis=True,
            gap_width=78,
            overlap=100,
            label_position="center",
            label_color=WHITE,
            label_font_size=8,
            series_colors=SERIES_COLORS,
            value_min=0.0,
            value_max=_money_axis_max([value for _, values in quota_series for value in values]),
            show_legend=True,
            legend_position="bottom",
            show_data_labels=True,
        )
    add_footer(slide, timestamp_text, current_page)

    # Slide 3 — returns and term
    slide = prs.slides.add_slide(blank)
    current_page = next_page()
    add_slide_header(slide, "Rentabilidade e prazo")

    selected_labels = _top_return_labels_for_ppt(dashboard)
    return_table_df = _build_return_inline_table_for_ppt(dashboard, selected_labels=selected_labels, months=12)
    if not return_table_df.empty:
        month_count = max(len(return_table_df.columns) - 3, 0)
        add_table(
            slide,
            return_table_df,
            title="Rentabilidade por tipo de cota (% a.m.) — últimos 12 meses",
            left=MARGIN_LEFT_IN,
            top=1.15,
            width=full_width,
            height=1.40,
            col_widths=[2.15] + ([0.58] * month_count) + [0.86, 1.02],
            max_rows=8,
        )

    base100_df = _build_return_base100_for_ppt(dashboard, selected_labels=selected_labels, months=12)
    if not base100_df.empty:
        ordered_base100 = _sort_competencia_frame(base100_df, ascending=True)
        base100_series = [
            (str(series_name), pd.to_numeric(group["valor"], errors="coerce").fillna(0.0).tolist())
            for series_name, group in ordered_base100.groupby("serie", dropna=False)
        ]
        base100_numeric = pd.to_numeric(base100_df["valor"], errors="coerce")
        base100_min = float(base100_numeric.min()) if not base100_numeric.dropna().empty else 100.0
        base100_max = float(base100_numeric.max()) if not base100_numeric.dropna().empty else 100.0
        base100_axis_min = min(95.0, base100_min * 0.98)
        base100_axis_max = max(base100_max * 1.06, 105.0)
        base100_chart = add_chart(
            slide,
            title="Índice acumulado base 100",
            chart_type=XL_CHART_TYPE.LINE_MARKERS,
            categories=_competencia_labels(ordered_base100["competencia"].drop_duplicates().tolist()),
            series_map=base100_series,
            left=MARGIN_LEFT_IN,
            top=2.82,
            width=full_width,
            height=1.40,
            number_format="0.0",
            label_position="above",
            label_font_size=8,
            show_data_labels=False,
            series_colors=SERIES_COLORS,
            value_min=base100_axis_min,
            value_max=base100_axis_max,
            show_legend=True,
            legend_position="bottom",
        )
        for idx, series in enumerate(base100_chart.series):
            _style_line_series(series, color=SERIES_COLORS[idx % len(SERIES_COLORS)], width_pt=2.0, marker_size=7)
        base100_last_values = [values[-1] for _, values in base100_series if values]
        if base100_last_values:
            base100_label_colors = [SERIES_COLORS[idx % len(SERIES_COLORS)] for idx in range(len(base100_last_values))]
            add_line_end_labels(
                slide,
                labels=[_format_decimal(value, decimals=1) for value in base100_last_values],
                values=base100_last_values,
                colors=base100_label_colors,
                left=MARGIN_LEFT_IN,
                top=2.82,
                width=full_width,
                height=1.40,
                axis_min=base100_axis_min,
                axis_max=base100_axis_max,
                font_size=11,
                fill_colors=[WHITE] * len(base100_last_values),
                text_colors=base100_label_colors,
                side="right",
            )

    maturity_df = _latest_maturity_chart_frame(dashboard.maturity_latest_df)
    if not maturity_df.empty:
        maturity_df = maturity_df[pd.to_numeric(maturity_df["valor"], errors="coerce").fillna(0.0) > 0].copy()
    if not maturity_df.empty:
        maturity_scale = _money_scale(pd.to_numeric(maturity_df["valor"], errors="coerce"))
        maturity_values = _scale_values(maturity_df["valor"], maturity_scale).tolist()
        maturity_title = (
            f"Prazo de vencimento dos DCs a vencer em {_format_competencia(dashboard.latest_competencia)}"
        )
        maturity_chart = add_chart(
            slide,
            title=maturity_title,
            chart_type=XL_CHART_TYPE.BAR_CLUSTERED,
            categories=maturity_df["faixa"].astype(str).tolist(),
            series_map=[("Saldo", maturity_values)],
            left=MARGIN_LEFT_IN,
            top=4.20,
            width=left_wide_width,
            height=2.40,
            number_format=_money_number_format(maturity_scale),
            gap_width=48,
            label_position="outside_end",
            label_font_size=8,
            series_colors=[ORANGE],
            value_min=0.0,
            value_max=_money_axis_max(list(maturity_values) + [sum(maturity_values)]),
            show_legend=False,
            show_data_labels=True,
        )
        try:
            maturity_chart.category_axis.reverse_order = True
        except Exception:  # noqa: BLE001
            pass

    duration_df = _sort_competencia_frame(dashboard.duration_history_df, ascending=True)
    if not duration_df.empty:
        duration_series = [("Prazo médio proxy (dias)", _series_numeric(duration_df, "duration_days").fillna(0.0).tolist())]
        duration_values = [float(value or 0.0) for value in duration_series[0][1]]
        sorted_duration_values = sorted(duration_values)
        second_highest = sorted_duration_values[-2] if len(sorted_duration_values) >= 2 else (sorted_duration_values[-1] if sorted_duration_values else 0.0)
        max_duration = sorted_duration_values[-1] if sorted_duration_values else 0.0
        duration_axis_max = max(max(duration_series[0][1], default=30.0) * 1.12, max(duration_series[0][1], default=30.0) + 4.0, 30.0)
        duration_outlier_competencia = None
        duration_axis_note = None
        if max_duration > 0 and second_highest > 0 and max_duration >= second_highest * 3:
            duration_axis_max = max(second_highest * 1.35, second_highest + 8.0, 30.0)
            max_idx = duration_values.index(max_duration)
            duration_outlier_competencia = _format_competencia(str(duration_df.iloc[max_idx]["competencia"]))
            duration_axis_note = (
                f"Escala visual limitada a {int(round(duration_axis_max))} dias; "
                f"pico de {int(round(max_duration))} dias em {duration_outlier_competencia}."
            )
        duration_chart = add_chart(
            slide,
            title="Prazo médio proxy dos recebíveis (dias)",
            chart_type=XL_CHART_TYPE.LINE_MARKERS,
            categories=_competencia_labels(duration_df["competencia"].tolist()),
            series_map=duration_series,
            left=right_col_left,
            top=4.20,
            width=right_narrow_width,
            height=2.40,
            number_format='0.0',
            label_position="above",
            label_font_size=8,
            show_data_labels=False,
            value_max=duration_axis_max,
            value_min=0.0,
        )
        _style_line_series(duration_chart.series[0], color=ORANGE, width_pt=2.2, marker_size=9)
        if duration_values:
            add_line_end_labels(
                slide,
                labels=[f"{_format_decimal(duration_values[-1], decimals=1)} dias"],
                values=[duration_values[-1]],
                colors=[ORANGE],
                left=right_col_left,
                top=4.20,
                width=right_narrow_width,
                height=2.40,
                axis_max=duration_axis_max,
                axis_min=0.0,
                font_size=11,
                fill_colors=[WHITE],
                text_colors=[ORANGE],
                side="right",
            )
        if duration_axis_note is not None:
            add_textbox(
                slide,
                right_col_left,
                6.48,
                right_narrow_width,
                0.18,
                duration_axis_note,
                size=FOOTER_SIZE,
                color=MID_GRAY,
            )
    add_footer(slide, timestamp_text, current_page)

    default_df = _sort_competencia_frame(dashboard.default_history_df, ascending=True)
    credit_categories = _competencia_labels(default_df["competencia"].tolist()) if not default_df.empty else []
    credit_bar_series: list[tuple[str, Sequence[float]]] = []
    credit_line_series: list[tuple[str, Sequence[float]]] = []
    if not default_df.empty:
        credit_bar_series = [
            ("Inadimplência", _series_numeric(default_df, "inadimplencia_pct").fillna(0.0).tolist()),
            ("Provisão", _series_numeric(default_df, "provisao_pct_direitos").fillna(0.0).tolist()),
        ]
        credit_line_series = [
            ("Cobertura", _series_numeric(default_df, "cobertura_pct").fillna(0.0).tolist()),
            ("100% (paridade)", [100.0] * len(credit_categories)),
        ]
    credit_has_values = _series_map_has_nonzero_values(credit_bar_series) or (
        _series_map_has_nonzero_values([credit_line_series[0]]) if credit_line_series else False
    )
    aging_history = _build_aging_history_for_ppt(dashboard)
    aging_columns = [column for column in aging_history.columns if column != "competencia"] if not aging_history.empty else []
    aging_series_map = [(column, aging_history[column].tolist()) for column in aging_columns]
    aging_has_values = _series_map_has_nonzero_values(aging_series_map)

    if credit_has_values or aging_has_values:
        slide = prs.slides.add_slide(blank)
        current_page = next_page()
        add_slide_header(slide, "Crédito e cobertura")
        if credit_has_values and aging_has_values:
            add_overlay_combo_credit_chart(
                slide,
                title="Inadimplência e provisão",
                categories=credit_categories,
                bar_series_map=credit_bar_series,
                line_series_map=credit_line_series,
                left=MARGIN_LEFT_IN,
                top=top_row_top,
                width=full_width,
                height=2.42,
            )
            add_chart(
                slide,
                title="Aging",
                chart_type=XL_CHART_TYPE.COLUMN_STACKED,
                categories=_competencia_labels(aging_history["competencia"].tolist()),
                series_map=aging_series_map,
                left=MARGIN_LEFT_IN,
                top=3.90,
                width=full_width,
                height=2.70,
                number_format='0"%"',
                percent_axis=True,
                gap_width=44,
                overlap=100,
                label_position="center",
                label_color=WHITE,
                label_font_size=8,
                series_colors=AGING_PPT_COLORS,
                value_max=_dashboard_percent_axis_upper([value for _, values in aging_series_map for value in values], max_cap=120.0),
                show_legend=True,
                legend_position="bottom",
                show_data_labels=False,
            )
        elif credit_has_values:
            add_overlay_combo_credit_chart(
                slide,
                title="Inadimplência e provisão",
                categories=credit_categories,
                bar_series_map=credit_bar_series,
                line_series_map=credit_line_series,
                left=MARGIN_LEFT_IN,
                top=top_row_top,
                width=full_width,
                height=3.90,
            )
            add_empty_state_panel(
                slide,
                title="Aging",
                message="Sem saldos vencidos distribuídos por faixa de atraso na janela exibida.",
                left=MARGIN_LEFT_IN,
                top=5.10,
                width=full_width,
                height=1.25,
            )
        elif aging_has_values:
            add_chart(
                slide,
                title="Aging",
                chart_type=XL_CHART_TYPE.COLUMN_STACKED,
                categories=_competencia_labels(aging_history["competencia"].tolist()),
                series_map=aging_series_map,
                left=MARGIN_LEFT_IN,
                top=top_row_top,
                width=full_width,
                height=5.55,
                number_format='0"%"',
                percent_axis=True,
                gap_width=44,
                overlap=100,
                label_position="center",
                label_color=WHITE,
                label_font_size=8,
                series_colors=AGING_PPT_COLORS,
                value_max=_dashboard_percent_axis_upper(_series_values(aging_series_map), max_cap=120.0),
                show_legend=True,
                legend_position="bottom",
                show_data_labels=False,
            )
        add_footer(slide, timestamp_text, current_page)

    over_history = _build_over_aging_history_for_ppt(dashboard)
    if not over_history.empty:
        over_columns = [column for column in over_history.columns if column != "competencia"]
        over_series_map = [(column, over_history[column].tolist()) for column in over_columns]
        if _series_map_has_nonzero_values(over_series_map):
            slide = prs.slides.add_slide(blank)
            current_page = next_page()
            add_slide_header(slide, "Deterioração acumulada")
            over_axis_max = _dashboard_percent_axis_upper(_series_values(over_series_map), max_cap=120.0)
            over_chart = add_chart(
                slide,
                title="Inadimplência Over",
                chart_type=XL_CHART_TYPE.LINE_MARKERS,
                categories=_competencia_labels(over_history["competencia"].tolist()),
                series_map=over_series_map,
                left=MARGIN_LEFT_IN,
                top=top_row_top,
                width=full_width,
                height=5.55,
                number_format='0.0"%"',
                percent_axis=True,
                label_position="above",
                label_font_size=8,
                series_colors=OVER_PPT_COLORS,
                value_max=over_axis_max,
                show_data_labels=False,
                show_legend=True,
                legend_position="bottom",
            )
            for idx, series in enumerate(over_chart.series):
                _style_line_series(
                    series,
                    color=OVER_PPT_COLORS[idx % len(OVER_PPT_COLORS)],
                    width_pt=2.2,
                    marker_size=8,
                )
            add_footer(slide, timestamp_text, current_page)

    total_pages = len(prs.slides)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text or ""
                    if text.startswith("Página ") and " de " not in text:
                        run.text = f"{text} de {total_pages}"

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def build_dashboard_pptx_file(
    dashboard: FundonetDashboardData,
    output_path: Path,
    *,
    generated_at: datetime | None = None,
) -> Path:
    output_path.write_bytes(build_dashboard_pptx_bytes(dashboard, generated_at=generated_at))
    return output_path


def _fund_title(dashboard: FundonetDashboardData) -> str:
    return str(dashboard.fund_info.get("nome_fundo") or dashboard.fund_info.get("nome_classe") or "FIDC selecionado")


def _to_float(value: object) -> float | None:
    try:
        if value is None or value == "":
            return None
        return float(value)
    except (TypeError, ValueError):
        return None


def _format_decimal(value: object, *, decimals: int = 2) -> str:
    numeric = _to_float(value)
    if numeric is None:
        return "N/D"
    return f"{numeric:,.{decimals}f}".replace(",", "X").replace(".", ",").replace("X", ".")


def _format_percent(value: object) -> str:
    numeric = _to_float(value)
    if numeric is None:
        return "N/D"
    return f"{_format_decimal(numeric, decimals=1)}%"


def _format_metric_value(value: object, unit: str) -> str:
    if unit == "R$":
        return _format_brl_compact(value)
    if unit == "%":
        return _format_percent(value)
    return _format_decimal(value)


def _format_brl_compact(value: object) -> str:
    numeric = _to_float(value)
    if numeric is None:
        return "N/D"
    scale = _money_scale(pd.Series([numeric]))
    if scale.divisor == 1.0:
        return f"R$ {_format_decimal(numeric, decimals=2)}"
    decimals = 2 if scale.divisor >= 1_000_000_000 else 1
    return f"R$ {_format_decimal(numeric / scale.divisor, decimals=decimals)} {scale.suffix}"


def _format_competencia(value: object) -> str:
    text = str(value or "").strip()
    if len(text) == 7 and "/" in text:
        month, year = text.split("/", 1)
        month_map = {
            "01": "jan", "02": "fev", "03": "mar", "04": "abr", "05": "mai", "06": "jun",
            "07": "jul", "08": "ago", "09": "set", "10": "out", "11": "nov", "12": "dez",
        }
        return f"{month_map.get(month, month)}-{year[-2:]}"
    return text or "N/D"


def _competencia_labels(values: Sequence[object]) -> list[str]:
    return [_format_competencia(value) for value in values]


def _sort_competencia_frame(frame: pd.DataFrame, *, ascending: bool = True) -> pd.DataFrame:
    if frame.empty:
        return frame.copy()
    output = frame.copy()
    helper_column = "__competencia_sort_dt"
    if "competencia_dt" in output.columns:
        return output.sort_values("competencia_dt", ascending=ascending, kind="stable").reset_index(drop=True)
    if "competencia" in output.columns:
        output[helper_column] = pd.to_datetime("01/" + output["competencia"].astype(str), format="%d/%m/%Y", errors="coerce")
        output = output.sort_values(helper_column, ascending=ascending, kind="stable").drop(columns=[helper_column])
        return output.reset_index(drop=True)
    return output.reset_index(drop=True)


def _safe_pct(numerator: object, denominator: object) -> float | None:
    num = _to_float(numerator)
    den = _to_float(denominator)
    if num is None or den is None or den <= 0:
        return None
    return num / den * 100.0


def _series_numeric(frame: pd.DataFrame, column: str) -> pd.Series:
    if column not in frame.columns:
        return pd.Series(dtype="float64")
    numeric = pd.to_numeric(frame[column], errors="coerce")
    finite_mask = numeric.map(lambda value: math.isfinite(value) if pd.notna(value) else False)
    return numeric.where(finite_mask, pd.NA)


def _percent_axis_max(series_map: Sequence[tuple[str, Sequence[float]]], *, cap: float = 110.0) -> float:
    values: list[float] = []
    for _, series_values in series_map:
        for value in series_values:
            numeric = _to_float(value)
            if numeric is None or not math.isfinite(numeric):
                continue
            values.append(float(numeric))
    if not values:
        return cap
    max_value = max(values)
    if max_value <= 0:
        return cap
    return min(cap, max(max_value * 1.18, max_value + 4.0))


def _money_axis_max(values: Sequence[object]) -> float | None:
    numeric = pd.to_numeric(pd.Series(values), errors="coerce")
    numeric = numeric[numeric.map(lambda value: math.isfinite(value) if pd.notna(value) else False)]
    if numeric.empty:
        return None
    max_value = float(numeric.max())
    if max_value <= 0:
        return None
    return max_value * 1.16


def _money_scale(values: pd.Series) -> MoneyScale:
    numeric = pd.to_numeric(values, errors="coerce")
    numeric = numeric[numeric.map(lambda value: math.isfinite(value) if pd.notna(value) else False)]
    max_abs = float(numeric.abs().max()) if not numeric.empty else 0.0
    if max_abs >= 1_000_000_000_000:
        return MoneyScale(1_000_000_000_000.0, "tri", "R$ tri")
    if max_abs >= 1_000_000_000:
        return MoneyScale(1_000_000_000.0, "bi", "R$ bi")
    if max_abs >= 1_000_000:
        return MoneyScale(1_000_000.0, "mm", "R$ mm")
    return MoneyScale(1.0, "", "R$")


def _scale_values(values: Sequence[object], scale: MoneyScale) -> pd.Series:
    numeric = pd.to_numeric(pd.Series(values), errors="coerce")
    numeric = numeric.where(numeric.map(lambda value: math.isfinite(value) if pd.notna(value) else False), pd.NA).fillna(0.0)
    if scale.divisor <= 0:
        return numeric
    return numeric / scale.divisor


def _sanitize_chart_value(value: object) -> float | None:
    numeric = _to_float(value)
    if numeric is None or not math.isfinite(numeric):
        return None
    return float(numeric)


def _sanitize_chart_series(values: Sequence[object]) -> tuple[float | None, ...]:
    return tuple(_sanitize_chart_value(value) for value in values)


def _money_number_format(scale: MoneyScale) -> str:
    if scale.divisor == 1.0:
        return '#,##0.00'
    return '0.0'


def _money_label_number_format(scale: MoneyScale) -> str:
    if scale.divisor == 1.0:
        return '#,##0'
    if scale.suffix:
        return f'0" {scale.suffix}"'
    return '0'


def _latest_aging_table_frame(frame: pd.DataFrame) -> pd.DataFrame:
    if frame.empty:
        return pd.DataFrame({"Faixa": ["Sem dados"], "Valor": ["-"], "% inadimplência": ["-"], "% DCs": ["-"]})
    output = frame.copy()
    output["Faixa"] = output["faixa"]
    output["Valor"] = output["valor"].map(_format_brl_compact)
    aging_column = "percentual_inadimplencia" if "percentual_inadimplencia" in output.columns else "percentual"
    dc_column = "percentual_direitos_creditorios" if "percentual_direitos_creditorios" in output.columns else None
    output["% inadimplência"] = output[aging_column].map(_format_percent) if aging_column in output.columns else "N/D"
    output["% DCs"] = output[dc_column].map(_format_percent) if dc_column else "N/D"
    return output[["Faixa", "Valor", "% inadimplência", "% DCs"]]


def _build_over_aging_history_for_ppt(dashboard: FundonetDashboardData) -> pd.DataFrame:
    over_history_df = _sort_competencia_frame(dashboard.default_over_history_df, ascending=False)
    if over_history_df.empty:
        return pd.DataFrame(columns=["competencia"])
    ordered_competencias = _ordered_competencias(over_history_df, ascending=False)
    pivot = (
        over_history_df[over_history_df["calculo_status"] == "calculado"]
        .pivot(index="competencia", columns="serie", values="percentual")
        .reset_index()
    )
    ordered_columns = ["competencia", "Over 1", "Over 30", "Over 60", "Over 90", "Over 180", "Over 360"]
    for column in ordered_columns:
        if column not in pivot.columns:
            pivot[column] = 0.0 if column != "competencia" else pivot.get(column)
    pivot = pivot[ordered_columns].copy()
    pivot = pivot.fillna(0.0)
    return _reorder_competencia_pivot(pivot, ordered_competencias)


def _build_aging_history_for_ppt(dashboard: FundonetDashboardData) -> pd.DataFrame:
    aging_history_df = _sort_competencia_frame(dashboard.default_aging_history_df, ascending=False)
    if aging_history_df.empty:
        return pd.DataFrame(columns=["competencia"])
    ordered_competencias = _ordered_competencias(aging_history_df, ascending=False)
    pivot = (
        aging_history_df.pivot_table(
            index="competencia",
            columns="faixa",
            values="percentual_inadimplencia",
            aggfunc="sum",
        )
        .fillna(0.0)
        .reset_index()
    )
    ordered_columns = [
        "competencia",
        "Até 30 dias",
        "31 a 60 dias",
        "61 a 90 dias",
        "91 a 120 dias",
        "121 a 150 dias",
        "151 a 180 dias",
        "181 a 360 dias",
        "361 a 720 dias",
        "721 a 1080 dias",
        "Acima de 1080 dias",
    ]
    for column in ordered_columns:
        if column not in pivot.columns:
            pivot[column] = 0.0 if column != "competencia" else pivot.get(column)
    pivot = pivot[ordered_columns].copy()
    return _reorder_competencia_pivot(pivot, ordered_competencias)


def _latest_over_table_frame(frame: pd.DataFrame) -> pd.DataFrame:
    if frame.empty:
        return pd.DataFrame({"Faixa": ["Sem dados"], "% DCs": ["-"], "Status": ["-"]})
    latest_row = frame.iloc[-1]
    rows = []
    for column in frame.columns:
        if column == "competencia":
            continue
        rows.append(
            {
                "Faixa": column,
                "% DCs": _format_percent(latest_row[column]),
                "Status": "Acumulado",
            }
        )
    return pd.DataFrame(rows)


def _latest_events_table_frame(frame: pd.DataFrame) -> pd.DataFrame:
    if frame.empty:
        return pd.DataFrame({"Evento": ["Sem dados"], "Valor": ["-"], "% PL": ["-"]})
    output = frame.copy()
    output["Evento"] = output["evento"]
    output["Valor"] = output["valor_total"].map(_format_brl_compact)
    output["% PL"] = output["valor_total_pct_pl"].map(_format_percent)
    return output[["Evento", "Valor", "% PL"]]


def _class_display_column(df: pd.DataFrame) -> str:
    if "class_label" in df.columns:
        return "class_label"
    return "label"


def _latest_quota_table_frame(frame: pd.DataFrame, latest_competencia: str) -> pd.DataFrame:
    if frame.empty:
        return pd.DataFrame({"Classe": ["Sem dados"], "Qt. cotas": ["-"], "PL": ["-"]})
    output = frame[frame["competencia"] == latest_competencia].copy()
    if output.empty:
        return pd.DataFrame({"Classe": ["Sem dados"], "Qt. cotas": ["-"], "PL": ["-"]})
    if "aggregation_scope" in output.columns and output["aggregation_scope"].eq("portfolio").all():
        output["Classe"] = output.get("class_macro_label", output[_class_display_column(output)])
        output["PL"] = output["pl"].map(_format_brl_compact)
        output["% do PL"] = output["pl_share_pct"].map(_format_percent)
        return output[["Classe", "PL", "% do PL"]]
    output["Classe"] = output[_class_display_column(output)]
    output["Qt. cotas"] = output["qt_cotas"].map(_format_decimal_0_or_2)
    output["PL"] = output["pl"].map(_format_brl_compact)
    return output[["Classe", "Qt. cotas", "PL"]]


def _latest_maturity_chart_frame(frame: pd.DataFrame) -> pd.DataFrame:
    if frame.empty:
        return pd.DataFrame(columns=["faixa", "valor", "ordem"])
    output = frame.copy()
    output["valor"] = pd.to_numeric(output["valor"], errors="coerce").fillna(0.0)
    output = output[output["faixa"] != "Vencidos"].copy()
    if output.empty:
        return pd.DataFrame(columns=["faixa", "valor", "ordem"])
    if "ordem" in output.columns:
        output = output.sort_values("ordem")
        return output[["faixa", "valor", "ordem"]]
    return output[["faixa", "valor"]]


def _ordered_competencias(frame: pd.DataFrame, *, ascending: bool = True) -> list[str]:
    if frame.empty or "competencia" not in frame.columns:
        return []
    output = _sort_competencia_frame(frame, ascending=ascending)
    output["competencia"] = output["competencia"].astype(str)
    return output["competencia"].dropna().drop_duplicates().tolist()


def _reorder_competencia_pivot(frame: pd.DataFrame, ordered_competencias: Sequence[str]) -> pd.DataFrame:
    if frame.empty or "competencia" not in frame.columns or not ordered_competencias:
        return frame.reset_index(drop=True)
    output = frame.copy()
    output["competencia"] = output["competencia"].astype(str)
    output["competencia"] = pd.Categorical(output["competencia"], categories=list(ordered_competencias), ordered=True)
    output = output.sort_values("competencia").reset_index(drop=True)
    output["competencia"] = output["competencia"].astype(str)
    return output


def _series_map_has_nonzero_values(series_map: Sequence[tuple[str, Sequence[float]]], *, tolerance: float = 1e-9) -> bool:
    for _, values in series_map:
        numeric = pd.to_numeric(pd.Series(list(values)), errors="coerce").fillna(0.0)
        if numeric.abs().gt(tolerance).any():
            return True
    return False


def _series_values(series_map: Sequence[tuple[str, Sequence[float]]]) -> list[float]:
    values: list[float] = []
    for _, series in series_map:
        numeric = pd.to_numeric(pd.Series(list(series)), errors="coerce").fillna(0.0)
        values.extend([float(value) for value in numeric.tolist()])
    return values


def _quota_pl_share_pivot(frame: pd.DataFrame) -> pd.DataFrame:
    if frame.empty:
        return pd.DataFrame(columns=["competencia"])
    output = frame.copy()
    ordered_competencias = _ordered_competencias(output, ascending=False)
    output["percentual"] = pd.to_numeric(output["pl_share_pct"], errors="coerce").fillna(0.0)
    label_column = "class_macro_label" if "class_macro_label" in output.columns else _class_display_column(output)
    pivot = (
        output.pivot_table(index="competencia", columns=label_column, values="percentual", aggfunc="sum")
        .fillna(0.0)
        .reset_index()
    )
    return _reorder_competencia_pivot(pivot, ordered_competencias)


def _quota_pl_value_pivot(frame: pd.DataFrame) -> pd.DataFrame:
    if frame.empty:
        return pd.DataFrame(columns=["competencia"])
    output = frame.copy()
    ordered_competencias = _ordered_competencias(output, ascending=False)
    output["valor"] = pd.to_numeric(output["pl"], errors="coerce").fillna(0.0)
    label_column = "class_macro_label" if "class_macro_label" in output.columns else _class_display_column(output)
    pivot = (
        output.pivot_table(index="competencia", columns=label_column, values="valor", aggfunc="sum")
        .fillna(0.0)
        .reset_index()
    )
    return _reorder_competencia_pivot(pivot, ordered_competencias)


def _top_return_labels_for_ppt(dashboard: FundonetDashboardData, *, limit: int | None = None) -> list[str]:
    return_history_df = dashboard.return_history_df.copy()
    if return_history_df.empty:
        return []
    label_column = _class_display_column(return_history_df)
    labels = return_history_df[label_column].dropna().astype(str).drop_duplicates().tolist()
    if limit is None or len(labels) <= limit:
        return labels
    quota_df = dashboard.quota_pl_history_df.copy()
    if quota_df.empty or dashboard.latest_competencia not in quota_df["competencia"].astype(str).tolist():
        return labels[:limit]
    quota_df = quota_df[quota_df["competencia"] == dashboard.latest_competencia].copy()
    quota_df["pl"] = pd.to_numeric(quota_df["pl"], errors="coerce")
    ordered = (
        quota_df.sort_values("pl", ascending=False)[_class_display_column(quota_df)]
        .dropna()
        .astype(str)
        .drop_duplicates()
        .tolist()
    )
    selected = [label for label in ordered if label in set(labels)]
    remaining = [label for label in labels if label not in set(selected)]
    return (selected + remaining)[:limit]


def _return_history_pivot(dashboard: FundonetDashboardData, *, selected_labels: Sequence[str] | None = None) -> pd.DataFrame:
    return_history_df = dashboard.return_history_df.sort_values("competencia_dt").copy()
    if return_history_df.empty:
        return pd.DataFrame(columns=["competencia"])
    ordered_competencias = _ordered_competencias(return_history_df)
    label_column = _class_display_column(return_history_df)
    if selected_labels:
        return_history_df = return_history_df[return_history_df[label_column].isin(list(selected_labels))].copy()
    if return_history_df.empty:
        return pd.DataFrame(columns=["competencia"])
    return_history_df["retorno_mensal_pct"] = pd.to_numeric(return_history_df["retorno_mensal_pct"], errors="coerce")
    pivot = (
        return_history_df.pivot_table(
            index="competencia",
            columns=label_column,
            values="retorno_mensal_pct",
            aggfunc="last",
        )
        .reset_index()
    )
    ordered_columns = ["competencia"] + [label for label in (selected_labels or []) if label in pivot.columns]
    remaining_columns = [column for column in pivot.columns if column not in ordered_columns]
    pivot = pivot[ordered_columns + remaining_columns].copy()
    return _reorder_competencia_pivot(pivot, ordered_competencias)


def _risk_metrics_table_frame(metrics_df: pd.DataFrame, risk_block: str) -> pd.DataFrame:
    if metrics_df.empty:
        return pd.DataFrame(columns=["Métrica", "Valor", "Leitura"])
    output = metrics_df[metrics_df["risk_block"] == risk_block].copy()
    if output.empty:
        return pd.DataFrame(columns=["Métrica", "Valor", "Leitura"])
    output["Métrica"] = output["label"]
    output["Valor"] = output.apply(
        lambda row: _format_metric_value(row.get("value"), str(row.get("unit") or "")),
        axis=1,
    )
    output["Leitura"] = output["interpretation"].fillna("N/D")
    return output[["Métrica", "Valor", "Leitura"]]


def _format_decimal_0_or_2(value: object) -> str:
    numeric = _to_float(value)
    if numeric is None:
        return "N/D"
    decimals = 0 if float(numeric).is_integer() else 2
    return _format_decimal(numeric, decimals=decimals)


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    color = str(value or BLACK).strip().lstrip("#")
    return (int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def _load_pil_font(size: int, *, bold: bool = False) -> ImageFont.ImageFont:
    font_paths = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for path in font_paths:
        font_path = Path(path)
        if font_path.exists():
            return ImageFont.truetype(str(font_path), size=size)
    return ImageFont.load_default()


def _text_size(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> tuple[int, int]:
    left, top, right, bottom = draw.textbbox((0, 0), text, font=font)
    return right - left, bottom - top


def _chart_canvas(
    *,
    width_px: int,
    height_px: int,
    title: str | None,
    legend_items: Sequence[tuple[str, str]],
) -> tuple[Image.Image, ImageDraw.ImageDraw, dict[str, int], ImageFont.ImageFont, ImageFont.ImageFont, ImageFont.ImageFont]:
    image = Image.new("RGB", (width_px, height_px), _hex_to_rgb(WHITE))
    draw = ImageDraw.Draw(image)
    title_font = _load_pil_font(30, bold=True)
    label_font = _load_pil_font(17, bold=False)
    small_font = _load_pil_font(15, bold=False)
    top_pad = 18
    if title:
        draw.text((26, 18), title, fill=_hex_to_rgb(BLACK), font=title_font)
        top_pad = 60
    legend_rows = 0
    if legend_items:
        max_per_row = max(1, int((width_px - 120) / 210))
        legend_rows = math.ceil(len(legend_items) / max_per_row)
    legend_height = 16 + (legend_rows * 24) if legend_rows else 0
    plot_box = {
        "left": 86,
        "top": top_pad,
        "right": width_px - 34,
        "bottom": height_px - (48 + legend_height),
        "legend_top": height_px - (legend_height if legend_height else 0) - 20,
        "width": width_px,
        "height": height_px,
    }
    if legend_items:
        _draw_image_legend(
            draw,
            legend_items=legend_items,
            left=plot_box["left"],
            top=plot_box["legend_top"],
            max_width=plot_box["right"] - plot_box["left"],
            font=label_font,
        )
    return image, draw, plot_box, title_font, label_font, small_font


def _draw_image_legend(
    draw: ImageDraw.ImageDraw,
    *,
    legend_items: Sequence[tuple[str, str]],
    left: int,
    top: int,
    max_width: int,
    font: ImageFont.ImageFont,
) -> None:
    cursor_x = left
    cursor_y = top
    row_height = 22
    for label, color in legend_items:
        marker_w = 18
        marker_h = 6
        label_w, label_h = _text_size(draw, label, font)
        total_w = marker_w + 8 + label_w + 18
        if cursor_x + total_w > left + max_width:
            cursor_x = left
            cursor_y += row_height
        draw.rounded_rectangle(
            (cursor_x, cursor_y + 5, cursor_x + marker_w, cursor_y + 5 + marker_h),
            radius=2,
            fill=_hex_to_rgb(color),
            outline=_hex_to_rgb(color),
        )
        draw.text((cursor_x + marker_w + 8, cursor_y), label, fill=_hex_to_rgb(DARK_GRAY), font=font)
        cursor_x += total_w


def _draw_dashed_line(
    draw: ImageDraw.ImageDraw,
    *,
    start: tuple[float, float],
    end: tuple[float, float],
    color: str,
    width: int = 2,
    dash: int = 8,
    gap: int = 6,
) -> None:
    x1, y1 = start
    x2, y2 = end
    total_len = math.hypot(x2 - x1, y2 - y1)
    if total_len <= 0:
        return
    dx = (x2 - x1) / total_len
    dy = (y2 - y1) / total_len
    cursor = 0.0
    while cursor < total_len:
        seg_start = cursor
        seg_end = min(cursor + dash, total_len)
        sx = x1 + dx * seg_start
        sy = y1 + dy * seg_start
        ex = x1 + dx * seg_end
        ey = y1 + dy * seg_end
        draw.line((sx, sy, ex, ey), fill=_hex_to_rgb(color), width=width)
        cursor += dash + gap


def _repel_positions(values: Sequence[float], *, min_gap: float, lower: float, upper: float) -> list[float]:
    if not values:
        return []
    indexed = sorted(enumerate(values), key=lambda item: item[1])
    adjusted: dict[int, float] = {}
    previous: float | None = None
    for idx, value in indexed:
        candidate = value if previous is None else max(value, previous + min_gap)
        adjusted[idx] = min(candidate, upper)
        previous = adjusted[idx]
    ordered = [adjusted[idx] for idx in range(len(values))]
    if ordered and ordered[-1] > upper:
        shift = ordered[-1] - upper
        ordered = [max(lower, value - shift) for value in ordered]
    return ordered


def _value_to_y(value: float, *, axis_min: float, axis_max: float, top: int, bottom: int) -> float:
    span = max(axis_max - axis_min, 1e-9)
    return bottom - ((value - axis_min) / span) * (bottom - top)


def _draw_vertical_axis(
    draw: ImageDraw.ImageDraw,
    *,
    left: int,
    top: int,
    bottom: int,
    ticks: Sequence[float],
    axis_min: float,
    axis_max: float,
    formatter,
    label_font: ImageFont.ImageFont,
    label_color: str = MID_GRAY,
    grid_color: str = GRID_GRAY,
    right: int | None = None,
    align_right: bool = False,
) -> None:
    axis_x = right if align_right and right is not None else left
    draw.line((axis_x, top, axis_x, bottom), fill=_hex_to_rgb(GRID_GRAY), width=1)
    grid_left = left if right is None else left
    grid_right = right if right is not None else axis_x
    for tick in ticks:
        y = _value_to_y(float(tick), axis_min=axis_min, axis_max=axis_max, top=top, bottom=bottom)
        draw.line((grid_left, y, grid_right, y), fill=_hex_to_rgb(grid_color), width=1)
        label = formatter(float(tick))
        label_w, label_h = _text_size(draw, label, label_font)
        if align_right and right is not None:
            draw.text((right + 10, y - (label_h / 2)), label, fill=_hex_to_rgb(label_color), font=label_font)
        else:
            draw.text((left - 12 - label_w, y - (label_h / 2)), label, fill=_hex_to_rgb(label_color), font=label_font)


def _draw_x_axis(
    draw: ImageDraw.ImageDraw,
    *,
    categories: Sequence[str],
    x_positions: Sequence[float],
    bottom: int,
    font: ImageFont.ImageFont,
) -> None:
    for category, x in zip(categories, x_positions, strict=False):
        label_w, label_h = _text_size(draw, category, font)
        draw.text((x - (label_w / 2), bottom + 10), category, fill=_hex_to_rgb(DARK_GRAY), font=font)


def _nice_ticks(axis_min: float, axis_max: float, *, steps: int = 4) -> list[float]:
    if axis_max <= axis_min:
        return [axis_min]
    return [axis_min + ((axis_max - axis_min) * idx / steps) for idx in range(steps + 1)]


def _finite_float(value: object) -> float | None:
    numeric = _to_float(value)
    if numeric is None or not math.isfinite(numeric):
        return None
    return float(numeric)


def _finite_or_default(value: object, default: float = 0.0) -> float:
    numeric = _finite_float(value)
    return default if numeric is None else numeric


def _percent_tick_formatter(value: float) -> str:
    return f"{value:.0f}%".replace(".", ",")


def _number_tick_formatter(value: float) -> str:
    return f"{value:.0f}".replace(".", ",")


def _money_tick_formatter(scale: MoneyScale):
    def _formatter(value: float) -> str:
        return _format_decimal(value, decimals=0 if scale.divisor != 1 else 2)

    return _formatter


def _line_chart_png(
    *,
    title: str,
    categories: Sequence[str],
    series_map: Sequence[tuple[str, Sequence[float]]],
    colors: Sequence[str],
    width_px: int,
    height_px: int,
    axis_min: float = 0.0,
    axis_max: float | None = None,
    tick_formatter=None,
    legend: bool = True,
    show_point_labels: bool = False,
    show_end_labels: bool = True,
    reference_lines: Sequence[tuple[float, str, str]] | None = None,
) -> bytes:
    legend_items = list(zip([name for name, _ in series_map], colors[: len(series_map)], strict=False)) if legend else []
    image, draw, box, _, label_font, small_font = _chart_canvas(
        width_px=width_px,
        height_px=height_px,
        title=title,
        legend_items=legend_items,
    )
    values = [numeric for _, series in series_map for value in series if (numeric := _finite_float(value)) is not None]
    if reference_lines:
        values.extend([float(value) for value, _, _ in reference_lines])
    axis_max = axis_max if axis_max is not None else _dashboard_percent_axis_upper(values) if tick_formatter == _percent_tick_formatter else (max(values) * 1.14 if values else 100.0)
    axis_max = max(axis_max, axis_min + 1.0)
    formatter = tick_formatter or _number_tick_formatter
    ticks = _nice_ticks(axis_min, axis_max, steps=4)
    _draw_vertical_axis(
        draw,
        left=box["left"],
        top=box["top"],
        bottom=box["bottom"],
        ticks=ticks,
        axis_min=axis_min,
        axis_max=axis_max,
        formatter=formatter,
        label_font=small_font,
    )
    plot_width = box["right"] - box["left"]
    step = plot_width / max(len(categories) - 1, 1)
    x_positions = [box["left"] + (idx * step) for idx in range(len(categories))]
    draw.line((box["left"], box["bottom"], box["right"], box["bottom"]), fill=_hex_to_rgb(GRID_GRAY), width=1)
    _draw_x_axis(draw, categories=categories, x_positions=x_positions, bottom=box["bottom"], font=small_font)
    if reference_lines:
        for value, label, color in reference_lines:
            y = _value_to_y(float(value), axis_min=axis_min, axis_max=axis_max, top=box["top"], bottom=box["bottom"])
            _draw_dashed_line(
                draw,
                start=(box["left"], y),
                end=(box["right"], y),
                color=color,
                width=2,
                dash=10,
                gap=6,
            )
            label_w, _ = _text_size(draw, label, small_font)
            draw.text((box["left"], y - 18), label, fill=_hex_to_rgb(color), font=small_font)
    last_points: list[tuple[str, float, float, str, str]] = []
    for series_idx, (series_name, series_values) in enumerate(series_map):
        color = colors[series_idx % len(colors)]
        points: list[tuple[float, float, float]] = []
        for x, value in zip(x_positions, series_values, strict=False):
            numeric = _finite_float(value)
            if numeric is None:
                continue
            y = _value_to_y(numeric, axis_min=axis_min, axis_max=axis_max, top=box["top"], bottom=box["bottom"])
            points.append((x, y, numeric))
        if len(points) >= 2:
            draw.line([(x, y) for x, y, _ in points], fill=_hex_to_rgb(color), width=4)
        for point_idx, (x, y, numeric) in enumerate(points):
            draw.ellipse((x - 5, y - 5, x + 5, y + 5), fill=_hex_to_rgb(color), outline=_hex_to_rgb(color))
            if show_point_labels:
                label = formatter(numeric)
                label_w, label_h = _text_size(draw, label, small_font)
                offset = 18 + ((point_idx + series_idx) % 2) * 12
                draw.text((x - (label_w / 2), max(box["top"], y - offset - label_h)), label, fill=_hex_to_rgb(color), font=small_font)
        if points:
            x, y, numeric = points[-1]
            last_points.append((series_name, x, y, color, formatter(numeric)))
    if show_end_labels and last_points:
        y_values = [item[2] for item in last_points]
        adjusted = _repel_positions(y_values, min_gap=20.0, lower=box["top"] + 6, upper=box["bottom"] - 6)
        label_x = box["right"] + 18
        for (series_name, x, y, color, label), label_y in zip(last_points, adjusted, strict=False):
            _draw_dashed_line(
                draw,
                start=(x, y),
                end=(label_x - 8, label_y + 8),
                color=color,
                width=2,
                dash=8,
                gap=5,
            )
            draw.text((label_x, label_y), label, fill=_hex_to_rgb(color), font=label_font)
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _stacked_bar_chart_png(
    *,
    title: str,
    categories: Sequence[str],
    series_map: Sequence[tuple[str, Sequence[float]]],
    colors: Sequence[str],
    width_px: int,
    height_px: int,
    percent_axis: bool = False,
    money_scale: MoneyScale | None = None,
    show_latest_callouts: bool = False,
) -> bytes:
    legend_items = list(zip([name for name, _ in series_map], colors[: len(series_map)], strict=False))
    image, draw, box, _, label_font, small_font = _chart_canvas(
        width_px=width_px,
        height_px=height_px,
        title=title,
        legend_items=legend_items,
    )
    totals = []
    for idx in range(len(categories)):
        totals.append(sum(_finite_or_default(values[idx]) for _, values in series_map))
    axis_min = 0.0
    axis_max = _dashboard_percent_axis_upper(totals, max_cap=140.0) if percent_axis else (max(totals) * 1.15 if totals else 1.0)
    axis_max = max(axis_max, 100.0 if percent_axis else 1.0)
    formatter = _percent_tick_formatter if percent_axis else _money_tick_formatter(money_scale or MoneyScale(1.0, "", "R$"))
    ticks = _nice_ticks(axis_min, axis_max, steps=4)
    _draw_vertical_axis(
        draw,
        left=box["left"],
        top=box["top"],
        bottom=box["bottom"],
        ticks=ticks,
        axis_min=axis_min,
        axis_max=axis_max,
        formatter=formatter,
        label_font=small_font,
    )
    draw.line((box["left"], box["bottom"], box["right"], box["bottom"]), fill=_hex_to_rgb(GRID_GRAY), width=1)
    plot_width = box["right"] - box["left"]
    if len(categories) == 1:
        group_width = min(plot_width * 0.34, 240)
        x_positions = [box["left"] + (plot_width / 2)]
    else:
        group_width = plot_width / max(len(categories), 1)
        x_positions = [box["left"] + ((idx + 0.5) * group_width) for idx in range(len(categories))]
    bar_width = min(58, max(26, int(group_width * 0.52)))
    _draw_x_axis(draw, categories=categories, x_positions=x_positions, bottom=box["bottom"], font=small_font)
    latest_segments: list[tuple[str, float, float, str, str]] = []
    for cat_idx, x in enumerate(x_positions):
        current_top_value = 0.0
        for series_idx, (series_name, series_values) in enumerate(series_map):
            color = colors[series_idx % len(colors)]
            value = _finite_or_default(series_values[cat_idx])
            next_top_value = current_top_value + value
            y_bottom = _value_to_y(current_top_value, axis_min=axis_min, axis_max=axis_max, top=box["top"], bottom=box["bottom"])
            y_top = _value_to_y(next_top_value, axis_min=axis_min, axis_max=axis_max, top=box["top"], bottom=box["bottom"])
            draw.rectangle((x - (bar_width / 2), y_top, x + (bar_width / 2), y_bottom), fill=_hex_to_rgb(color), outline=_hex_to_rgb(WHITE))
            if show_latest_callouts and cat_idx == len(x_positions) - 1 and value > 0:
                label = _format_percent(value) if percent_axis else _format_brl_compact((money_scale.divisor if money_scale else 1.0) * value)
                latest_segments.append((series_name, x + (bar_width / 2), (y_top + y_bottom) / 2, color, label))
            current_top_value = next_top_value
    if show_latest_callouts and latest_segments:
        segment_y = [item[2] for item in latest_segments]
        adjusted = _repel_positions(segment_y, min_gap=18.0, lower=box["top"] + 4, upper=box["bottom"] - 4)
        label_x = box["right"] + 24
        for (series_name, x, y, color, label), label_y in zip(latest_segments, adjusted, strict=False):
            _draw_dashed_line(
                draw,
                start=(x, y),
                end=(label_x - 10, label_y + 8),
                color=color,
                width=2,
                dash=8,
                gap=5,
            )
            draw.text((label_x, label_y), label, fill=_hex_to_rgb(color), font=label_font)
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _grouped_bar_line_chart_png(
    *,
    title: str,
    categories: Sequence[str],
    bar_series_map: Sequence[tuple[str, Sequence[float]]],
    line_series_map: Sequence[tuple[str, Sequence[float]]],
    width_px: int,
    height_px: int,
) -> bytes:
    legend_items = [
        *[(name, SERIES_COLORS[idx % len(SERIES_COLORS)]) for idx, (name, _) in enumerate(bar_series_map)],
        *[(name, COVERAGE_LINE_COLOR if idx == 0 else MID_GRAY) for idx, (name, _) in enumerate(line_series_map)],
    ]
    image, draw, box, _, label_font, small_font = _chart_canvas(
        width_px=width_px,
        height_px=height_px,
        title=title,
        legend_items=legend_items,
    )
    bar_values = [_finite_or_default(value) for _, values in bar_series_map for value in values]
    line_values = [_finite_or_default(value) for _, values in line_series_map for value in values]
    bar_axis_max = _dashboard_percent_axis_upper(bar_values, max_cap=140.0)
    line_axis_max = _dashboard_percent_axis_upper(line_values, max_cap=max(max(line_values, default=100.0) * 1.4, 650.0))
    bar_ticks = _nice_ticks(0.0, bar_axis_max, steps=4)
    line_ticks = _nice_ticks(0.0, line_axis_max, steps=4)
    _draw_vertical_axis(
        draw,
        left=box["left"],
        top=box["top"],
        bottom=box["bottom"],
        ticks=bar_ticks,
        axis_min=0.0,
        axis_max=bar_axis_max,
        formatter=_percent_tick_formatter,
        label_font=small_font,
    )
    _draw_vertical_axis(
        draw,
        left=box["left"],
        right=box["right"],
        top=box["top"],
        bottom=box["bottom"],
        ticks=line_ticks,
        axis_min=0.0,
        axis_max=line_axis_max,
        formatter=_percent_tick_formatter,
        label_font=small_font,
        label_color=COVERAGE_LINE_COLOR,
        grid_color=WHITE,
        align_right=True,
    )
    draw.text((box["left"], box["top"] - 24), "% dos DCs", fill=_hex_to_rgb(MID_GRAY), font=small_font)
    draw.text((box["right"] - 88, box["top"] - 24), "Cobertura (%)", fill=_hex_to_rgb(COVERAGE_LINE_COLOR), font=small_font)
    draw.line((box["left"], box["bottom"], box["right"], box["bottom"]), fill=_hex_to_rgb(GRID_GRAY), width=1)
    plot_width = box["right"] - box["left"]
    if len(categories) == 1:
        group_width = min(plot_width * 0.36, 260)
        x_positions = [box["left"] + (plot_width / 2)]
    else:
        group_width = plot_width / max(len(categories), 1)
        x_positions = [box["left"] + ((idx + 0.5) * group_width) for idx in range(len(categories))]
    inner_width = group_width * 0.62
    bar_width = inner_width / max(len(bar_series_map), 1)
    _draw_x_axis(draw, categories=categories, x_positions=x_positions, bottom=box["bottom"], font=small_font)
    for cat_idx, x in enumerate(x_positions):
        left_edge = x - (inner_width / 2)
        for series_idx, (series_name, series_values) in enumerate(bar_series_map):
            value = _finite_or_default(series_values[cat_idx])
            x0 = left_edge + (series_idx * bar_width)
            x1 = x0 + (bar_width * 0.82)
            y = _value_to_y(value, axis_min=0.0, axis_max=bar_axis_max, top=box["top"], bottom=box["bottom"])
            color = SERIES_COLORS[series_idx % len(SERIES_COLORS)]
            draw.rectangle((x0, y, x1, box["bottom"]), fill=_hex_to_rgb(color), outline=_hex_to_rgb(color))
            label = _format_percent(value)
            label_w, label_h = _text_size(draw, label, small_font)
            draw.text((x0 + ((x1 - x0 - label_w) / 2), max(box["top"], y - 18 - label_h)), label, fill=_hex_to_rgb(BLACK), font=small_font)
    reference_series = {name: values for name, values in line_series_map}
    if "100% (paridade)" in reference_series:
        y = _value_to_y(100.0, axis_min=0.0, axis_max=line_axis_max, top=box["top"], bottom=box["bottom"])
        _draw_dashed_line(
            draw,
            start=(box["left"], y),
            end=(box["right"], y),
            color=MID_GRAY,
            width=2,
            dash=10,
            gap=6,
        )
    coverage_values = None
    for series_name, values in line_series_map:
        if series_name == "Cobertura":
            coverage_values = values
            break
    if coverage_values:
        points: list[tuple[float, float, float]] = []
        for x, raw_value in zip(x_positions, coverage_values, strict=False):
            numeric = _finite_or_default(raw_value)
            y = _value_to_y(numeric, axis_min=0.0, axis_max=line_axis_max, top=box["top"], bottom=box["bottom"])
            points.append((x, y, numeric))
        if len(points) >= 2:
            draw.line([(x, y) for x, y, _ in points], fill=_hex_to_rgb(COVERAGE_LINE_COLOR), width=4)
        for idx, (x, y, numeric) in enumerate(points):
            draw.ellipse((x - 5, y - 5, x + 5, y + 5), fill=_hex_to_rgb(COVERAGE_LINE_COLOR), outline=_hex_to_rgb(COVERAGE_LINE_COLOR))
            label = _format_percent(numeric)
            label_w, label_h = _text_size(draw, label, small_font)
            offset = 22 + (idx % 2) * 12
            draw.text((x - (label_w / 2), max(box["top"], y - offset - label_h)), label, fill=_hex_to_rgb(COVERAGE_LINE_COLOR), font=small_font)
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _waterfall_chart_png(
    *,
    title: str,
    categories: Sequence[str],
    step_values: Sequence[float],
    total_value: float,
    width_px: int,
    height_px: int,
    money_scale: MoneyScale,
) -> bytes:
    image, draw, box, _, _, small_font = _chart_canvas(
        width_px=width_px,
        height_px=height_px,
        title=title,
        legend_items=[("Fluxo", ORANGE), ("Total", BLACK)],
    )
    axis_max = _money_axis_max(list(step_values) + [total_value]) or max(total_value, 1.0)
    ticks = _nice_ticks(0.0, axis_max, steps=4)
    _draw_vertical_axis(
        draw,
        left=box["left"],
        top=box["top"],
        bottom=box["bottom"],
        ticks=ticks,
        axis_min=0.0,
        axis_max=axis_max,
        formatter=_money_tick_formatter(money_scale),
        label_font=small_font,
    )
    draw.line((box["left"], box["bottom"], box["right"], box["bottom"]), fill=_hex_to_rgb(GRID_GRAY), width=1)
    all_categories = list(categories) + ["Total"]
    plot_width = box["right"] - box["left"]
    group_width = plot_width / max(len(all_categories), 1)
    bar_width = min(54, max(24, int(group_width * 0.52)))
    x_positions = [box["left"] + ((idx + 0.5) * group_width) for idx in range(len(all_categories))]
    _draw_x_axis(draw, categories=all_categories, x_positions=x_positions, bottom=box["bottom"], font=small_font)
    cumulative = 0.0
    for idx, (category, value) in enumerate(zip(categories, step_values, strict=False)):
        next_cumulative = cumulative + float(value)
        y0 = _value_to_y(cumulative, axis_min=0.0, axis_max=axis_max, top=box["top"], bottom=box["bottom"])
        y1 = _value_to_y(next_cumulative, axis_min=0.0, axis_max=axis_max, top=box["top"], bottom=box["bottom"])
        x = x_positions[idx]
        draw.rectangle((x - (bar_width / 2), min(y0, y1), x + (bar_width / 2), max(y0, y1)), fill=_hex_to_rgb(ORANGE), outline=_hex_to_rgb(ORANGE))
        label = _format_decimal((float(value)), decimals=1)
        label_w, label_h = _text_size(draw, label, small_font)
        draw.text((x - (label_w / 2), min(y0, y1) - label_h - 10), label, fill=_hex_to_rgb(DARK_GRAY), font=small_font)
        cumulative = next_cumulative
    total_x = x_positions[-1]
    total_y = _value_to_y(total_value, axis_min=0.0, axis_max=axis_max, top=box["top"], bottom=box["bottom"])
    draw.rectangle((total_x - (bar_width / 2), total_y, total_x + (bar_width / 2), box["bottom"]), fill=_hex_to_rgb(BLACK), outline=_hex_to_rgb(BLACK))
    total_label = _format_decimal(total_value, decimals=1)
    total_w, total_h = _text_size(draw, total_label, small_font)
    draw.text((total_x - (total_w / 2), total_y - total_h - 10), total_label, fill=_hex_to_rgb(BLACK), font=small_font)
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _build_return_inline_table_for_ppt(
    dashboard: FundonetDashboardData,
    *,
    selected_labels: Sequence[str] | None = None,
    months: int = 12,
) -> pd.DataFrame:
    history_df = dashboard.return_history_df.sort_values("competencia_dt").copy()
    summary_df = dashboard.return_summary_df.copy()
    if history_df.empty or summary_df.empty:
        return pd.DataFrame(columns=["Classe", "YTD", "12 meses"])
    label_column = _class_display_column(history_df)
    summary_label_column = _class_display_column(summary_df)
    if selected_labels:
        history_df = history_df[history_df[label_column].isin(list(selected_labels))].copy()
        summary_df = summary_df[summary_df[summary_label_column].isin(list(selected_labels))].copy()
    competencias = history_df["competencia"].drop_duplicates().tail(months).tolist()
    display_competencias = list(reversed(competencias))
    history_df = history_df[history_df["competencia"].isin(competencias)].copy()
    if history_df.empty:
        return pd.DataFrame(columns=["Classe", "YTD", "12 meses"])
    pivot = (
        history_df.pivot_table(
            index=label_column,
            columns="competencia",
            values="retorno_mensal_pct",
            aggfunc="last",
        )
        .reindex(columns=display_competencias)
    )
    ordered_labels = [label for label in (selected_labels or []) if label in pivot.index]
    ordered_labels += [label for label in pivot.index.tolist() if label not in set(ordered_labels)]
    pivot = pivot.reindex(ordered_labels).reset_index(drop=True)
    output = pd.DataFrame({"Classe": pd.Series(ordered_labels, dtype="object")})
    for competencia in display_competencias:
        output[_format_competencia(competencia)] = pd.Series(pivot[competencia].tolist(), dtype="object").map(_format_percent)
    summary_lookup = summary_df.set_index(summary_label_column)
    retorno_ano = summary_lookup["retorno_ano_pct"] if "retorno_ano_pct" in summary_lookup.columns else pd.Series(dtype="float64")
    retorno_12m = summary_lookup["retorno_12m_pct"] if "retorno_12m_pct" in summary_lookup.columns else pd.Series(dtype="float64")
    output["YTD"] = output["Classe"].map(lambda label: _format_percent(retorno_ano.get(label)))
    output["12 meses"] = output["Classe"].map(lambda label: _format_percent(retorno_12m.get(label)))
    return output


def _build_return_base100_for_ppt(
    dashboard: FundonetDashboardData,
    *,
    selected_labels: Sequence[str] | None = None,
    months: int = 12,
) -> pd.DataFrame:
    history_df = dashboard.return_history_df.sort_values("competencia_dt").copy()
    if history_df.empty:
        return pd.DataFrame(columns=["competencia", "serie", "valor"])
    label_column = _class_display_column(history_df)
    if selected_labels:
        history_df = history_df[history_df[label_column].isin(list(selected_labels))].copy()
    competencias = history_df["competencia"].drop_duplicates().tail(months).tolist()
    history_df = history_df[history_df["competencia"].isin(competencias)].copy()
    if history_df.empty:
        return pd.DataFrame(columns=["competencia", "serie", "valor"])
    rows: list[dict[str, object]] = []
    for label, group in history_df.groupby(label_column, dropna=False):
        ordered = group.sort_values("competencia_dt").copy()
        current_index = 100.0
        first_value = True
        for _, row in ordered.iterrows():
            monthly_return = pd.to_numeric(pd.Series([row.get("retorno_mensal_pct")]), errors="coerce").iloc[0]
            if first_value:
                current_index = 100.0
                first_value = False
            elif pd.notna(monthly_return):
                current_index *= 1.0 + (float(monthly_return) / 100.0)
            rows.append(
                {
                    "competencia": row.get("competencia"),
                    "competencia_dt": row.get("competencia_dt"),
                    "serie": str(label),
                    "valor": current_index,
                }
            )
    return pd.DataFrame(rows)


def _slide_px_size() -> tuple[int, int]:
    return int(round(SLIDE_WIDTH_IN * SLIDE_RENDER_DPI)), int(round(SLIDE_HEIGHT_IN * SLIDE_RENDER_DPI))


def _in_to_px(value_in: float) -> int:
    return int(round(value_in * SLIDE_RENDER_DPI))


def _new_slide_canvas(
    *,
    section_title: str,
    title_fund: str,
    subtitle: str,
    timestamp_text: str,
) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    width_px, height_px = _slide_px_size()
    image = Image.new("RGB", (width_px, height_px), _hex_to_rgb(WHITE))
    draw = ImageDraw.Draw(image)
    title_font = _load_pil_font(34, bold=True)
    subtitle_font = _load_pil_font(19, bold=True)
    meta_font = _load_pil_font(16, bold=False)
    footer_font = _load_pil_font(12, bold=False)
    draw.text((_in_to_px(0.45), _in_to_px(0.18)), section_title, fill=_hex_to_rgb(BLACK), font=title_font)
    draw.text((_in_to_px(0.45), _in_to_px(0.48)), title_fund, fill=_hex_to_rgb(ORANGE), font=subtitle_font)
    draw.text((_in_to_px(0.45), _in_to_px(0.72)), subtitle, fill=_hex_to_rgb(MID_GRAY), font=meta_font)
    footer = f"Fonte: Informe Mensal - CVM    |    Gerado em: {timestamp_text}"
    footer_w, footer_h = _text_size(draw, footer, footer_font)
    draw.text((_in_to_px(MARGIN_LEFT_IN), _in_to_px(7.08)), footer, fill=_hex_to_rgb(MID_GRAY), font=footer_font)
    return image, draw


def _paste_png(image: Image.Image, png_bytes: bytes, *, left: float, top: float, width: float, height: float) -> None:
    chart = Image.open(BytesIO(png_bytes)).convert("RGB")
    chart = chart.resize((_in_to_px(width), _in_to_px(height)))
    image.paste(chart, (_in_to_px(left), _in_to_px(top)))


def _draw_text(
    draw: ImageDraw.ImageDraw,
    *,
    left: float,
    top: float,
    text: str,
    size: int,
    color: str = BLACK,
    bold: bool = False,
) -> None:
    draw.text((_in_to_px(left), _in_to_px(top)), text, fill=_hex_to_rgb(color), font=_load_pil_font(size, bold=bold))


def _draw_table_on_slide(
    image: Image.Image,
    draw: ImageDraw.ImageDraw,
    df: pd.DataFrame,
    *,
    title: str,
    left: float,
    top: float,
    width: float,
    height: float,
    col_widths: Sequence[float] | None = None,
) -> None:
    _draw_text(draw, left=left, top=top - 0.18, text=title, size=15, bold=True)
    if df.empty:
        return
    left_px = _in_to_px(left)
    top_px = _in_to_px(top)
    width_px = _in_to_px(width)
    height_px = _in_to_px(height)
    rows = len(df.index) + 1
    cols = len(df.columns)
    row_h = max(28, int(height_px / max(rows, 1)))
    if col_widths and len(col_widths) == cols:
        total = sum(col_widths)
        widths = [int(round(width_px * (col_w / total))) for col_w in col_widths]
        widths[-1] += width_px - sum(widths)
    else:
        base = int(width_px / cols)
        widths = [base] * cols
        widths[-1] += width_px - sum(widths)
    header_font = _load_pil_font(11, bold=True)
    body_font = _load_pil_font(10, bold=False)
    x = left_px
    for col_idx, column in enumerate(df.columns):
        cell_w = widths[col_idx]
        draw.rectangle((x, top_px, x + cell_w, top_px + row_h), fill=_hex_to_rgb(BLACK), outline=_hex_to_rgb(GRID_GRAY))
        draw.text((x + 6, top_px + 6), str(column), fill=_hex_to_rgb(WHITE), font=header_font)
        x += cell_w
    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        x = left_px
        y = top_px + row_idx * row_h
        for col_idx, value in enumerate(row):
            cell_w = widths[col_idx]
            draw.rectangle((x, y, x + cell_w, y + row_h), fill=_hex_to_rgb(WHITE), outline=_hex_to_rgb(GRID_GRAY))
            text = str(value)
            while text and _text_size(draw, text, body_font)[0] > cell_w - 12:
                text = text[:-2] + "…"
            draw.text((x + 6, y + 6), text, fill=_hex_to_rgb(BLACK), font=body_font)
            x += cell_w


def _dashboard_percent_axis_upper(values: Sequence[object], *, max_cap: float | None = None) -> float:
    """Round the axis ceiling up to the next clean step above the data maximum.

    Step table (matches spec section 6.4):
      max < 1%   → next 0.5 p.p.   e.g. 0.8 → 1.0
      max < 10%  → next 1 p.p.     e.g. 3.2 → 4.0
      max < 100% → next 5 p.p.     e.g. 42  → 45
      max ≥ 100% → next 10         e.g. 134 → 140
    """
    numeric = pd.to_numeric(pd.Series(values), errors="coerce").dropna()
    if numeric.empty:
        return 4.0
    max_value = float(numeric.max())
    if max_value <= 0:
        upper = 1.0
    else:
        step = 0.5 if max_value < 1.0 else 1.0 if max_value < 10.0 else 5.0 if max_value < 100.0 else 10.0
        upper = (math.floor(max_value / step) + 1) * step
    if max_cap is not None:
        upper = min(upper, max_cap)
    # Minimum axis ceiling: one step above zero so the chart isn't collapsed
    return max(upper, 0.5)


def _render_structure_slide_png(
    *,
    dashboard: FundonetDashboardData,
    section_title: str,
    title_fund: str,
    subtitle: str,
    timestamp_text: str,
) -> bytes:
    image, draw = _new_slide_canvas(section_title=section_title, title_fund=title_fund, subtitle=subtitle, timestamp_text=timestamp_text)
    full_width = CONTENT_WIDTH_IN
    top_row_top = 1.02
    sub_df = _sort_competencia_frame(dashboard.subordination_history_df, ascending=True)
    if not sub_df.empty:
        sub_series = [("Subordinação reportada", _series_numeric(sub_df, "subordinacao_pct").fillna(0.0).tolist())]
        _paste_png(
            image,
            _line_chart_png(
                title="Subordinação reportada",
                categories=_competencia_labels(sub_df["competencia"].tolist()),
                series_map=sub_series,
                colors=[ORANGE],
                width_px=int(full_width * SLIDE_RENDER_DPI),
                height_px=int(2.05 * SLIDE_RENDER_DPI),
                axis_min=0.0,
                axis_max=_dashboard_percent_axis_upper(sub_series[0][1], max_cap=80.0),
                tick_formatter=_percent_tick_formatter,
                show_point_labels=len(sub_df["competencia"].drop_duplicates()) <= 12,
                show_end_labels=len(sub_df["competencia"].drop_duplicates()) > 12,
            ),
            left=MARGIN_LEFT_IN,
            top=top_row_top,
            width=full_width,
            height=2.05,
        )
        _draw_text(
            draw,
            left=MARGIN_LEFT_IN,
            top=top_row_top + 2.06,
            text="Subordinação reportada = (PL mezzanino + PL subordinada residual) / PL total.",
            size=12,
            color=MID_GRAY,
        )
    quota_values = _quota_pl_value_pivot(dashboard.quota_pl_history_df)
    if not quota_values.empty:
        quota_scale = _money_scale(quota_values.drop(columns=["competencia"]).stack())
        quota_series = [(column, _scale_values(quota_values[column], quota_scale).tolist()) for column in quota_values.columns if column != "competencia"]
        _paste_png(
            image,
            _stacked_bar_chart_png(
                title="PL por tipo de cota",
                categories=_competencia_labels(quota_values["competencia"].tolist()),
                series_map=quota_series,
                colors=SERIES_COLORS,
                width_px=int(full_width * SLIDE_RENDER_DPI),
                height_px=int(3.35 * SLIDE_RENDER_DPI),
                percent_axis=False,
                money_scale=quota_scale,
                show_latest_callouts=False,
            ),
            left=MARGIN_LEFT_IN,
            top=3.32,
            width=full_width,
            height=3.35,
        )
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _render_returns_slide_png(
    *,
    dashboard: FundonetDashboardData,
    section_title: str,
    title_fund: str,
    subtitle: str,
    timestamp_text: str,
) -> bytes:
    image, draw = _new_slide_canvas(section_title=section_title, title_fund=title_fund, subtitle=subtitle, timestamp_text=timestamp_text)
    full_width = CONTENT_WIDTH_IN
    left_wide_width = 7.45
    split_gap = 0.30
    right_narrow_width = CONTENT_WIDTH_IN - left_wide_width - split_gap
    right_col_left = MARGIN_LEFT_IN + left_wide_width + split_gap
    selected_labels = _top_return_labels_for_ppt(dashboard)
    return_table_df = _build_return_inline_table_for_ppt(dashboard, selected_labels=selected_labels, months=12)
    if not return_table_df.empty:
        month_count = max(len(return_table_df.columns) - 3, 0)
        _draw_table_on_slide(
            image,
            draw,
            return_table_df,
            title="Rentabilidade por tipo de cota (% a.m.) — últimos 12 meses",
            left=MARGIN_LEFT_IN,
            top=1.18,
            width=full_width,
            height=1.38,
            col_widths=[2.15] + ([0.58] * month_count) + [0.86, 1.02],
        )
    base100_df = _build_return_base100_for_ppt(dashboard, selected_labels=selected_labels, months=12)
    if not base100_df.empty:
        ordered_base100 = _sort_competencia_frame(base100_df, ascending=True)
        base100_series = [
            (str(series_name), pd.to_numeric(group["valor"], errors="coerce").fillna(0.0).tolist())
            for series_name, group in ordered_base100.groupby("serie", dropna=False)
        ]
        base100_numeric = pd.to_numeric(base100_df["valor"], errors="coerce")
        base100_min = float(base100_numeric.min()) if not base100_numeric.dropna().empty else 100.0
        base100_max = float(base100_numeric.max()) if not base100_numeric.dropna().empty else 100.0
        _paste_png(
            image,
            _line_chart_png(
                title="Índice acumulado base 100",
                categories=_competencia_labels(ordered_base100["competencia"].drop_duplicates().tolist()),
                series_map=base100_series,
                colors=SERIES_COLORS,
                width_px=int(full_width * SLIDE_RENDER_DPI),
                height_px=int(1.45 * SLIDE_RENDER_DPI),
                axis_min=min(95.0, base100_min * 0.98),
                axis_max=max(base100_max * 1.06, 105.0),
                tick_formatter=_number_tick_formatter,
                show_point_labels=False,
                show_end_labels=True,
            ),
            left=MARGIN_LEFT_IN,
            top=2.80,
            width=full_width,
            height=1.45,
        )
    maturity_df = _latest_maturity_chart_frame(dashboard.maturity_latest_df)
    if not maturity_df.empty:
        maturity_scale = _money_scale(pd.to_numeric(maturity_df["valor"], errors="coerce"))
        maturity_values = _scale_values(maturity_df["valor"], maturity_scale).tolist()
        _paste_png(
            image,
            _waterfall_chart_png(
                title=f"Prazo de vencimento dos DCs a vencer ({maturity_scale.label})" if maturity_scale.label else "Prazo de vencimento dos DCs a vencer",
                categories=maturity_df["faixa"].astype(str).tolist(),
                step_values=maturity_values,
                total_value=float(sum(maturity_values)),
                width_px=int(left_wide_width * SLIDE_RENDER_DPI),
                height_px=int(2.35 * SLIDE_RENDER_DPI),
                money_scale=maturity_scale,
            ),
            left=MARGIN_LEFT_IN,
            top=4.30,
            width=left_wide_width,
            height=2.35,
        )
    duration_df = _sort_competencia_frame(dashboard.duration_history_df, ascending=True)
    if not duration_df.empty:
        duration_values = _series_numeric(duration_df, "duration_days").fillna(0.0).tolist()
        duration_upper = max(max(duration_values or [0.0]) * 1.12, max(duration_values or [0.0]) + 4.0, 30.0)
        _paste_png(
            image,
            _line_chart_png(
                title="Prazo médio proxy dos recebíveis (dias)",
                categories=_competencia_labels(duration_df["competencia"].tolist()),
                series_map=[("Prazo médio proxy (dias)", duration_values)],
                colors=[ORANGE],
                width_px=int(right_narrow_width * SLIDE_RENDER_DPI),
                height_px=int(2.35 * SLIDE_RENDER_DPI),
                axis_min=0.0,
                axis_max=duration_upper,
                tick_formatter=_number_tick_formatter,
                show_point_labels=False,
                show_end_labels=True,
            ),
            left=right_col_left,
            top=4.30,
            width=right_narrow_width,
            height=2.35,
        )
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _render_credit_slide_png(
    *,
    dashboard: FundonetDashboardData,
    section_title: str,
    title_fund: str,
    subtitle: str,
    timestamp_text: str,
) -> bytes:
    image, _ = _new_slide_canvas(section_title=section_title, title_fund=title_fund, subtitle=subtitle, timestamp_text=timestamp_text)
    full_width = CONTENT_WIDTH_IN
    top_row_top = 1.02
    default_df = _sort_competencia_frame(dashboard.default_history_df, ascending=True)
    if not default_df.empty:
        categories = _competencia_labels(default_df["competencia"].tolist())
        _paste_png(
            image,
            _grouped_bar_line_chart_png(
                title="Inadimplência e provisão",
                categories=categories,
                bar_series_map=[
                    ("Inadimplência", _series_numeric(default_df, "inadimplencia_pct").fillna(0.0).tolist()),
                    ("Provisão", _series_numeric(default_df, "provisao_pct_direitos").fillna(0.0).tolist()),
                ],
                line_series_map=[
                    ("Cobertura", _series_numeric(default_df, "cobertura_pct").fillna(0.0).tolist()),
                    ("100% (paridade)", [100.0] * len(categories)),
                ],
                width_px=int(full_width * SLIDE_RENDER_DPI),
                height_px=int(2.30 * SLIDE_RENDER_DPI),
            ),
            left=MARGIN_LEFT_IN,
            top=top_row_top,
            width=full_width,
            height=2.30,
        )
    aging_history = _build_aging_history_for_ppt(dashboard)
    if not aging_history.empty:
        aging_columns = [column for column in aging_history.columns if column != "competencia"]
        aging_series_map = [(column, aging_history[column].tolist()) for column in aging_columns]
        _paste_png(
            image,
            _stacked_bar_chart_png(
                title="Aging",
                categories=_competencia_labels(aging_history["competencia"].tolist()),
                series_map=aging_series_map,
                colors=AGING_PPT_COLORS,
                width_px=int(full_width * SLIDE_RENDER_DPI),
                height_px=int(2.75 * SLIDE_RENDER_DPI),
                percent_axis=True,
                show_latest_callouts=True,
            ),
            left=MARGIN_LEFT_IN,
            top=3.85,
            width=full_width,
            height=2.75,
        )
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _render_over_slide_png(
    *,
    dashboard: FundonetDashboardData,
    section_title: str,
    title_fund: str,
    subtitle: str,
    timestamp_text: str,
) -> bytes:
    image, _ = _new_slide_canvas(section_title=section_title, title_fund=title_fund, subtitle=subtitle, timestamp_text=timestamp_text)
    full_width = CONTENT_WIDTH_IN
    top_row_top = 1.02
    over_history = _build_over_aging_history_for_ppt(dashboard)
    if not over_history.empty:
        over_columns = [column for column in over_history.columns if column != "competencia"]
        over_series_map = [(column, over_history[column].tolist()) for column in over_columns]
        over_axis_max = _dashboard_percent_axis_upper([value for _, values in over_series_map for value in values], max_cap=120.0)
        _paste_png(
            image,
            _line_chart_png(
                title="Inadimplência Over",
                categories=_competencia_labels(over_history["competencia"].tolist()),
                series_map=over_series_map,
                colors=OVER_PPT_COLORS,
                width_px=int(full_width * SLIDE_RENDER_DPI),
                height_px=int(5.55 * SLIDE_RENDER_DPI),
                axis_min=0.0,
                axis_max=over_axis_max,
                tick_formatter=_percent_tick_formatter,
                show_point_labels=False,
                show_end_labels=True,
            ),
            left=MARGIN_LEFT_IN,
            top=top_row_top,
            width=full_width,
            height=5.55,
        )
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()
