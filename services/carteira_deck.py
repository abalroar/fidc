"""Os slides de subordinação da carteira, e a troca deles no deck padrão.

O deck padrão é um binário publicado: ele chega pronto do bundle e é validado
contra um manifesto.  Editá-lo na origem exigiria republicar o bundle inteiro,
então a substituição acontece no caminho de exportação — sobre a apresentação
já carregada em memória, do mesmo modo que o ranking ANBIMA é acoplado.  O
bundle no disco permanece intacto e continua validando.

Os seis slides antigos são **reescritos no lugar**, e não removidos e
recriados.  A diferença não é estética: o ``next_partname`` do python-pptx
devolve nomes já ocupados quando a numeração das partes deixa de ser contígua,
e remover slides do meio do deck é exatamente o que abre esse buraco — o pacote
sairia com duas partes ``ppt/slides/slideN.xml``.  Acrescentar no fim é seguro,
e é assim que os slides excedentes entram.

Cada slide traz o gráfico à esquerda e, à direita, **uma tabela nativa do
Office** com todos os veículos daquele gráfico.  O gráfico nomeia só os poucos
que importam; a tabela é quem garante que nenhum FIDC fique sem identificação.
Ela é um objeto de tabela de verdade: dá para ordenar, editar valores e colar
no Excel.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pandas as pd

from services.carteira_subordinacao import (
    DEFAULT_DATA_DIR,
    dumbbell_figure,
    figure_png_bytes,
    resolve_portfolio,
    short_fund_name,
)


#: Faixa de slides que a carteira substitui, 1-indexada e inclusiva — é a
#: sequência "RISCO ESTRUTURAL · CARTEIRA I · <categoria>" do deck publicado.
REPLACED_SLIDE_RANGE = (18, 23)
#: Abaixo disso um gráfico de hastes não comunica nada: vira um punhado de
#: pontos soltos.
MIN_FUNDS_PER_CATEGORY = 3
#: A tabela ocupa a metade inferior em duas tranches lado a lado; cada uma
#: leva metade das linhas.  O que passar vai para um slide de continuação da
#: mesma categoria, em vez de encolher a fonte abaixo do legível.
ROWS_PER_TRANCHE = 11
MAX_ROWS_PER_SLIDE = ROWS_PER_TRANCHE * 2

# Geometria da lâmina: o gráfico toma a largura inteira do conteúdo e a
# tabela ocupa a metade inferior, dividida em duas tranches.
MARGIN_IN = 0.62
CHART_TOP_IN = 1.30
CHART_WIDTH_IN = 12.05
CHART_HEIGHT_IN = 3.10
TABLE_TOP_IN = 4.52
TRANCHE_WIDTH_IN = 5.90
TRANCHE_GAP_IN = 0.25
TRANCHE_COLUMNS = (2.30, 0.95, 0.80, 0.80, 1.05)

KICKER = "CARTEIRA 101 · RISCO ESTRUTURAL"
SOURCE = (
    "CVM, Informe Mensal FIDC (competência mais recente de cada fundo) e "
    "regulamentos na FundosNet/B3. Folga = subordinação atual − mínimo exigido; "
    "mínimo estrutural (subordinada + mezanino) quando o regulamento o define."
)

# Bandas da folga, em pontos percentuais.  Só sinalização de risco — sem ícone,
# sem forma, sem decoração.  O vermelho marca exclusivamente o desenquadramento.
FOLGA_ATENCAO = 0.0
FOLGA_CONFORTAVEL = 1.5
FILL_VERDE = "E4F2EC"
FILL_AMARELO = "FBF0D9"
FILL_VERMELHO = "F8DFE1"
TEXTO_VERDE = "17654F"
TEXTO_AMARELO = "8A6410"
TEXTO_VERMELHO = "A32130"
#: Destaque discreto para os maiores PLs — fundo cinza claríssimo, sem borda.
FILL_MATERIAL = "F2F3F4"
TOP_PL_DESTACADOS = 3


def slot_count() -> int:
    first, last = REPLACED_SLIDE_RANGE
    return last - first + 1


def clear_slide(slide) -> None:
    """Esvazia um slide, preservando a parte e a posição dele no deck."""

    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def move_slides(presentation, first: int, count: int, destination: int) -> None:
    """Move ``count`` slides a partir de ``first`` para a posição ``destination``.

    Índices 1-indexados, medidos antes do movimento.  Mexe apenas na ordem
    declarada em ``sldIdLst``: nenhuma parte é criada ou destruída, então a
    numeração do pacote fica intacta.
    """

    id_list = presentation.slides._sldIdLst
    entries = list(id_list)
    moving = entries[first - 1 : first - 1 + count]
    for entry in moving:
        id_list.remove(entry)
    for offset, entry in enumerate(moving):
        id_list.insert(destination - 1 + offset, entry)


def _competence_label(competencia: str) -> str:
    return (
        f"{competencia[5:]}/{competencia[2:4]}"
        if len(competencia) == 7 and competencia[4] == "-"
        else competencia
    )


def order_for_table(frame: pd.DataFrame) -> pd.DataFrame:
    """Materialidade primeiro, risco de desenquadramento em seguida.

    O PL manda: é ele que diz o tamanho do problema.  A folga desempata, de
    modo que, entre fundos de porte parecido, o mais perto do limite sobe.
    """

    return frame.sort_values(
        ["pl_mm", "folga_pp"], ascending=[False, True], na_position="last"
    ).reset_index(drop=True)


#: Cabeçalho da tabela.  O PL entra logo depois do nome: é ele que ordena e é
#: por ele que o leitor mede o tamanho de cada caso.
TABLE_HEADER = ("FIDC", "PL jun-26 (R$ mm)", "Mín. (%)", "Atual (%)", "Folga (p.p.)")
COL_FOLGA = len(TABLE_HEADER) - 1


def _mm(value: object) -> str:
    if pd.isna(value):
        return "—"
    return f"{float(value):,.1f}".replace(",", "@").replace(".", ",").replace("@", ".")


def fund_label(registro) -> str:
    """O nome curto, com a marca de multicedente/multissacado quando houver."""

    nome = short_fund_name(str(registro.fundo), limite=26)
    marca = str(getattr(registro, "multi_flag", "") or "").strip()
    return f"{nome} · {marca}" if marca else nome


def _pct(value: object) -> str:
    return "—" if pd.isna(value) else f"{float(value):,.1f}".replace(".", ",")


def _folga(value: object) -> str:
    if pd.isna(value):
        return "—"
    number = float(value)
    sinal = "+" if number > 0 else ""
    return f"{sinal}{number:,.1f}".replace(".", ",")


def table_rows(frame: pd.DataFrame) -> tuple[list[list[str]], dict, dict]:
    """As linhas da tabela e as cores que sinalizam risco e materialidade."""

    rows: list[list[str]] = [list(TABLE_HEADER)]
    fills: dict[tuple[int, int], str] = {}
    colors: dict[tuple[int, int], str] = {}

    materiais = set(frame.head(TOP_PL_DESTACADOS)["cnpj"])
    for indice, registro in enumerate(frame.itertuples(), start=1):
        rows.append(
            [
                fund_label(registro),
                _mm(registro.pl_mm),
                _pct(registro.referencia_pct),
                _pct(registro.sub_atual_pct),
                _folga(registro.folga_pp),
            ]
        )
        folga = registro.folga_pp
        if pd.notna(folga):
            if folga < FOLGA_ATENCAO:
                fills[(indice, COL_FOLGA)] = FILL_VERMELHO
                colors[(indice, COL_FOLGA)] = TEXTO_VERMELHO
            elif folga <= FOLGA_CONFORTAVEL:
                fills[(indice, COL_FOLGA)] = FILL_AMARELO
                colors[(indice, COL_FOLGA)] = TEXTO_AMARELO
            else:
                fills[(indice, COL_FOLGA)] = FILL_VERDE
                colors[(indice, COL_FOLGA)] = TEXTO_VERDE
        if registro.cnpj in materiais:
            # O destaque de materialidade fica só na coluna do nome, para não
            # competir com a sinalização de risco na coluna da folga.
            fills[(indice, 0)] = FILL_MATERIAL
    return rows, fills, colors


def slide_plans(data_dir: Path = DEFAULT_DATA_DIR) -> list[dict[str, object]]:
    """Um plano por slide: categoria, recorte da tabela e gráfico completo."""

    position = resolve_portfolio(data_dir, somente_ativos=True)
    comparable = position.frame[position.frame["comparavel"]]
    rotulo = _competence_label(position.competencia_base)

    plans: list[dict[str, object]] = []
    counts = comparable["categoria_estrutural"].value_counts()
    for categoria in counts[counts.ge(MIN_FUNDS_PER_CATEGORY)].index:
        subset = order_for_table(
            comparable[comparable["categoria_estrutural"].eq(categoria)]
        )
        breaches = int(subset["abaixo_do_minimo"].fillna(False).sum())
        titulo = (
            f"{categoria} | {breaches} de {len(subset)} abaixo do mínimo"
            if breaches
            else f"{categoria} | os {len(subset)} estão acima do mínimo"
        )
        paginas = max(1, -(-len(subset) // MAX_ROWS_PER_SLIDE))
        for pagina in range(paginas):
            recorte = subset.iloc[
                pagina * MAX_ROWS_PER_SLIDE : (pagina + 1) * MAX_ROWS_PER_SLIDE
            ]
            plans.append(
                {
                    "categoria": categoria,
                    "titulo": titulo if pagina == 0 else f"{titulo} (cont.)",
                    "nota": f"{categoria} · FIDCs ativos · base até {rotulo}",
                    "grafico": subset,
                    "tabela": recorte,
                    "desenha_grafico": pagina == 0,
                }
            )
    return plans


def draw_carteira_slide(deck, slide, plano: dict[str, object]) -> None:
    """Gráfico em largura total no topo; tabela em duas tranches embaixo."""

    from pptx.util import Inches

    deck.compose(slide, str(plano["titulo"]), KICKER)
    grafico = plano["grafico"]

    if plano["desenha_grafico"]:
        figure = dumbbell_figure(
            grafico, nomear_todos=True, figsize=(CHART_WIDTH_IN, CHART_HEIGHT_IN)
        )
        slide.shapes.add_picture(
            BytesIO(figure_png_bytes(figure)),
            Inches(MARGIN_IN),
            Inches(CHART_TOP_IN),
            width=Inches(CHART_WIDTH_IN),
        )

    # Duas tranches lado a lado: com onze linhas cada, a categoria inteira cabe
    # na metade inferior sem apertar a fonte.
    tabela = plano["tabela"]
    for indice, inicio in enumerate(range(0, len(tabela), ROWS_PER_TRANCHE)):
        recorte = tabela.iloc[inicio : inicio + ROWS_PER_TRANCHE]
        if recorte.empty:
            continue
        rows, fills, colors = table_rows(recorte)
        deck.native_table(
            slide,
            rows,
            MARGIN_IN + indice * (TRANCHE_WIDTH_IN + TRANCHE_GAP_IN),
            TABLE_TOP_IN,
            list(TRANCHE_COLUMNS),
            row_height=0.175,
            header_height=0.22,
            size=7.6,
            aligns="lrrrr",
            cell_fills=fills,
            cell_colors=colors,
        )
    deck.footer(slide, SOURCE)


def replace_structural_slides(
    presentation, data_dir: Path = DEFAULT_DATA_DIR
) -> int:
    """Reescreve os slides estruturais com os da carteira e devolve quantos."""

    from services.bba_deck import Deck

    first, _ = REPLACED_SLIDE_RANGE
    slots = slot_count()
    plans = slide_plans(data_dir)
    if not plans:
        return 0

    deck = Deck(KICKER, presentation)
    if len(presentation.slides) < first - 1 + slots:
        raise IndexError(
            f"O deck tem {len(presentation.slides)} slides; a carteira precisa "
            f"dos slots {first}–{first + slots - 1}. O deck publicado mudou de "
            "tamanho."
        )

    extras = len(plans) - slots
    if extras > 0:
        total = len(presentation.slides)
        for _ in range(extras):
            deck.blank()
        move_slides(presentation, total + 1, extras, first + slots)

    deck.page = first - 1
    slides = list(presentation.slides)
    for offset, plano in enumerate(plans):
        slide = slides[first - 1 + offset]
        clear_slide(slide)
        draw_carteira_slide(deck, slide, plano)
    return len(plans)


__all__ = [
    "MARGIN_IN",
    "TRANCHE_COLUMNS",
    "TRANCHE_GAP_IN",
    "TRANCHE_WIDTH_IN",
    "TABLE_TOP_IN",
    "CHART_HEIGHT_IN",
    "CHART_WIDTH_IN",
    "CHART_TOP_IN",
    "KICKER",
    "MAX_ROWS_PER_SLIDE",
    "ROWS_PER_TRANCHE",
    "TABLE_HEADER",
    "MIN_FUNDS_PER_CATEGORY",
    "REPLACED_SLIDE_RANGE",
    "clear_slide",
    "draw_carteira_slide",
    "move_slides",
    "order_for_table",
    "replace_structural_slides",
    "slide_plans",
    "slot_count",
    "table_rows",
]
