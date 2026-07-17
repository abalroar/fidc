"""Native PowerPoint and Excel exports for the FIDC Industry executive pack.

The analytical rules live in :mod:`services.industry_executive_pack`.  This
module is deliberately limited to loading versioned artefacts, formatting
editable Office objects and preserving the audit trail in the exported files.
"""

from __future__ import annotations

import json
import re
import zipfile
from functools import lru_cache
from io import BytesIO
from pathlib import Path
from typing import Iterable, Sequence

import pandas as pd

from services.industry_executive_pack import (
    ANBIMA_FOCUS_BY_TYPE,
    ANBIMA_ND,
    ANBIMA_TYPES,
    HOLDER_BUCKETS,
    IndustryExecutivePack,
    build_industry_executive_pack,
)
from services.industry_intelligence import canonical_provider


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_DATA_DIR = ROOT / "data" / "industry_study"

# Itaú BBA-inspired editorial palette.  Orange is used for emphasis, not as
# decorative chrome; the remaining hierarchy is black/navy and neutral gray.
BLACK = "151515"
NAVY = "151515"
ORANGE = "E36C0A"
ORANGE_LIGHT = "F8E9DE"
GRAY_900 = "30353A"
GRAY_700 = "5D6369"
GRAY_500 = "8D9399"
GRAY_300 = "D7DADD"
GRAY_200 = "E7E9EB"
GRAY_100 = "F5F6F7"
WHITE = "FFFFFF"
RED = "A83C32"
GREEN = "2D6F51"

ROLE_COLUMNS = {
    "administrador": "admin_nome",
    "gestor": "gestor_nome",
    "custodiante": "custodiante_nome",
}
ROLE_LABELS = {
    "administrador": "Administração",
    "gestor": "Gestão",
    "custodiante": "Custódia",
}
TYPE_COLORS = {
    "Fomento Mercantil": ORANGE,
    "Agro, Indústria e Comércio": NAVY,
    "Financeiro": "747A80",
    "Outros": "AEB3B7",
    ANBIMA_ND: "D8DADD",
}
MODEL_COLORS = {
    "Monoestrutura": ORANGE,
    "Administração + Gestão": NAVY,
    "Administração + Custódia": "656B70",
    "Gestão + Custódia": "959A9F",
    "Três prestadores distintos": "C2C6C9",
    "Dados incompletos": "E3E5E7",
}

# The executive series starts in 2015 to preserve the full comparable window
# used in the reference deck.  Closed years always use December; the current
# year is handled separately by ``_annual_history`` using the latest complete
# common competence.
ANNUAL_HISTORY_START_YEAR = 2015

_CHART_AXIS_ID_RE = re.compile(rb'(<c:(?:axId|crossAx)\b[^>]*\bval=")(-\d+)(")')
_PACK_INPUTS = (
    "vehicle_monthly.csv.gz",
    "industry_competence_status.csv",
    "industry_monthly.csv",
    "industry_anbima_classification.csv.gz",
    "industry_large_fund_classification.csv",
)


def _read(data_dir: Path, name: str) -> pd.DataFrame:
    path = data_dir / name
    return pd.read_csv(path, low_memory=False) if path.exists() else pd.DataFrame()


def _read_manifest(data_dir: Path) -> dict:
    path = data_dir / "industry_intelligence_manifest.json"
    return json.loads(path.read_text(encoding="utf-8")) if path.exists() else {}


def _data_signature(data_dir: Path) -> tuple[tuple[str, int, int], ...]:
    signature: list[tuple[str, int, int]] = []
    for name in _PACK_INPUTS:
        path = data_dir / name
        stat = path.stat() if path.exists() else None
        signature.append((name, stat.st_size if stat else -1, stat.st_mtime_ns if stat else -1))
    return tuple(signature)


@lru_cache(maxsize=3)
def _cached_executive_pack(
    data_dir_text: str,
    _signature: tuple[tuple[str, int, int], ...],
) -> IndustryExecutivePack:
    data_dir = Path(data_dir_text)
    return build_industry_executive_pack(
        vehicle_monthly=_read(data_dir, "vehicle_monthly.csv.gz"),
        competence_status=_read(data_dir, "industry_competence_status.csv"),
        industry_monthly=_read(data_dir, "industry_monthly.csv"),
        anbima_classification=_read(data_dir, "industry_anbima_classification.csv.gz"),
        published_classifications=_read(data_dir, "industry_large_fund_classification.csv"),
    )


def _load_executive_pack(data_dir: Path) -> IndustryExecutivePack:
    data_dir = data_dir.resolve()
    return _cached_executive_pack(str(data_dir), _data_signature(data_dir))


def _competence_label(value: object, *, short: bool = False, lower: bool = False) -> str:
    full_names = (
        "Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho",
        "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro",
    )
    short_names = ("Jan", "Fev", "Mar", "Abr", "Mai", "Jun", "Jul", "Ago", "Set", "Out", "Nov", "Dez")
    try:
        period = pd.Period(str(value), freq="M")
    except (TypeError, ValueError):
        return str(value)
    label = f"{(short_names if short else full_names)[period.month - 1]}/{str(period.year)[-2:]}"
    return label.lower() if lower else label


def _date_label(value: object) -> str:
    parsed = pd.to_datetime(value, errors="coerce")
    return str(value) if pd.isna(parsed) else parsed.strftime("%d/%m/%Y")


def _normalize_chart_axis_ids(payload: bytes) -> bytes:
    """Convert signed chart-axis IDs to their OpenXML unsigned representation."""

    def replace_axis_id(match: re.Match[bytes]) -> bytes:
        unsigned_value = int(match.group(2)) % (2**32)
        return match.group(1) + str(unsigned_value).encode("ascii") + match.group(3)

    output = BytesIO()
    with zipfile.ZipFile(BytesIO(payload), "r") as source, zipfile.ZipFile(
        output, "w", zipfile.ZIP_DEFLATED
    ) as target:
        for member in source.infolist():
            data = source.read(member.filename)
            if member.filename.startswith("ppt/charts/chart") and member.filename.endswith(".xml"):
                data = _CHART_AXIS_ID_RE.sub(replace_axis_id, data)
            target.writestr(member, data)
    return output.getvalue()


def _as_float(value: object, default: float = 0.0) -> float:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return default
    return default if pd.isna(number) else number


def _fmt_bi(value: object, decimals: int = 1) -> str:
    return f"R$ {_as_float(value) / 1e9:.{decimals}f} bi".replace(".", ",")


def _fmt_mi(value: object, decimals: int = 0) -> str:
    return f"R$ {_as_float(value) / 1e6:.{decimals}f} mi".replace(".", ",")


def _fmt_pct(value: object, decimals: int = 1) -> str:
    return f"{_as_float(value) * 100:.{decimals}f}%".replace(".", ",")


def _fmt_pp(value: object, decimals: int = 1) -> str:
    number = _as_float(value)
    sign = "+" if number > 0 else ""
    number_text = f"{sign}{number:.{decimals}f}".replace(".", ",")
    return f"{number_text} p.p."


def _short_name(value: object, limit: int = 52) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    replacements = (
        ("FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS", "FIDC"),
        ("FUNDO DE INVESTIMENTO EM DIREITOS CREDITORIOS", "FIDC"),
        ("RESPONSABILIDADE LIMITADA", "RL"),
    )
    for old, new in replacements:
        text = text.replace(old, new)
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "…"


def _annual_history(industry: pd.DataFrame, pack: IndustryExecutivePack) -> pd.DataFrame:
    frame = industry.copy()
    frame["competencia"] = frame["competencia"].astype(str)
    frame = frame[frame["competencia"].le(pack.competences.latest_complete)]
    frame["year"] = frame["competencia"].str[:4].astype(int)
    rows: list[pd.Series] = []
    latest_year = int(pack.competences.latest_complete[:4])
    for year, group in frame.groupby("year"):
        if year < ANNUAL_HISTORY_START_YEAR:
            continue
        december = group[group["competencia"].str.endswith("-12")]
        if year == latest_year:
            selected = group[group["competencia"].eq(pack.competences.latest_complete)]
        else:
            selected = december
        if not selected.empty:
            rows.append(selected.sort_values("competencia").iloc[-1])
    output = pd.DataFrame(rows).sort_values("year") if rows else pd.DataFrame()
    if output.empty:
        return output
    output["period_label"] = output.apply(
        lambda row: str(int(row["year"]))
        if str(row["competencia"]).endswith("-12")
        else _competence_label(row["competencia"], short=True),
        axis=1,
    )
    output["pl_ex_fic"] = pd.to_numeric(output["pl_total"], errors="coerce") - pd.to_numeric(
        output["pl_fic_fidc"], errors="coerce"
    ).fillna(0)
    output["pl_total_growth"] = pd.to_numeric(output["pl_total"], errors="coerce").pct_change()
    output["pl_ex_fic_growth"] = output["pl_ex_fic"].pct_change()
    return output


def _overall_rankings(pack: IndustryExecutivePack, role: str, limit: int = 10) -> pd.DataFrame:
    frames: list[pd.DataFrame] = []
    role_column = ROLE_COLUMNS[role]
    for competence, period in pack.competences.period_by_competence.items():
        frame = pack.fund_monthly[pack.fund_monthly["competencia"].eq(competence)].copy()
        frame["participant"] = frame.get(role_column, "").map(canonical_provider)
        frame["pl"] = pd.to_numeric(frame["pl"], errors="coerce")
        total = frame["pl"].sum(min_count=1)
        grouped = (
            frame[frame["participant"].ne("Não informado")]
            .groupby("participant", as_index=False)["pl"]
            .sum(min_count=1)
            .sort_values(["pl", "participant"], ascending=[False, True])
            .reset_index(drop=True)
        )
        grouped["rank"] = range(1, len(grouped) + 1)
        grouped["share"] = grouped["pl"] / total if total else float("nan")
        grouped["period"] = period
        frames.append(grouped)
    long = pd.concat(frames, ignore_index=True)
    latest_period = pack.competences.period_by_competence[pack.competences.latest_complete]
    leaders = long[long["period"].eq(latest_period)].head(limit)["participant"].tolist()
    rows: list[dict[str, object]] = []
    periods = [pack.competences.period_by_competence[c] for c in pack.competences.ordered]
    for participant in leaders:
        row: dict[str, object] = {"Participante": participant}
        period_labels = {
            periods[0]: ("Rank dez/24", "Share dez/24"),
            periods[1]: ("Rank dez/25", "Share dez/25"),
            periods[2]: (f"Rank {latest_period_label(pack)}", f"Share {latest_period_label(pack)}"),
        }
        for period in periods:
            match = long[(long["participant"].eq(participant)) & (long["period"].eq(period))]
            rank_header, share_header = period_labels[period]
            row[rank_header] = f"#{int(match.iloc[0]['rank'])}" if not match.empty else "—"
            row[share_header] = _fmt_pct(match.iloc[0]["share"]) if not match.empty else "—"
        first = long[(long["participant"].eq(participant)) & (long["period"].eq(periods[0]))]
        last = long[(long["participant"].eq(participant)) & (long["period"].eq(periods[-1]))]
        row["Δ share"] = (
            _fmt_pp((_as_float(last.iloc[0]["share"]) - _as_float(first.iloc[0]["share"])) * 100)
            if not first.empty and not last.empty
            else "—"
        )
        rows.append(row)
    return pd.DataFrame(rows)


def latest_period_label(pack: IndustryExecutivePack) -> str:
    return _competence_label(pack.competences.latest_complete, short=True).lower()


def _focus_ranking_table(pack: IndustryExecutivePack, role: str) -> pd.DataFrame:
    rankings = pack.rankings
    if rankings.empty:
        return pd.DataFrame()
    latest = pack.competences.period_by_competence[pack.competences.latest_complete]
    scoped = rankings[(rankings["role"].eq(role)) & rankings["scope"].eq("foco")].copy()
    latest_rows = scoped[scoped["period"].eq(latest)]
    focus_sizes = (
        latest_rows.groupby(["anbima_tipo", "anbima_foco"])["pl_brl"]
        .sum()
        .sort_values(ascending=False)
    )
    focus_sizes = focus_sizes[focus_sizes.ge(1_000_000_000.0)]
    periods = [pack.competences.period_by_competence[c] for c in pack.competences.ordered]
    rows: list[dict[str, object]] = []
    for anbima_type, focus in focus_sizes.index:
        current = latest_rows[
            latest_rows["anbima_tipo"].eq(anbima_type)
            & latest_rows["anbima_foco"].eq(focus)
        ].sort_values("rank").head(1)
        for item in current.itertuples(index=False):
            row: dict[str, object] = {
                "Tipo / foco": f"{str(anbima_type)[:12]} · {focus}",
                "Participante": item.participant,
            }
            period_labels = {
                periods[0]: ("Rank dez/24", "Share dez/24"),
                periods[1]: ("Rank dez/25", "Share dez/25"),
                periods[2]: (f"Rank {latest_period_label(pack)}", f"Share {latest_period_label(pack)}"),
            }
            for period in periods:
                match = scoped[
                    scoped["anbima_tipo"].eq(anbima_type)
                    & scoped["anbima_foco"].eq(focus)
                    & scoped["participant"].eq(item.participant)
                    & scoped["period"].eq(period)
                ]
                rank_header, share_header = period_labels[period]
                row[rank_header] = f"#{int(match.iloc[0]['rank'])}" if not match.empty else "—"
                row[share_header] = _fmt_pct(match.iloc[0]["share_pl"]) if not match.empty else "—"
            rows.append(row)
    return pd.DataFrame(rows)


def _current_provider_leaders(pack: IndustryExecutivePack) -> dict[str, pd.DataFrame]:
    latest = pack.latest_funds.copy()
    latest["pl"] = pd.to_numeric(latest["pl"], errors="coerce")
    outputs: dict[str, pd.DataFrame] = {}
    canonical: dict[str, pd.Series] = {}
    for role, column in ROLE_COLUMNS.items():
        canonical[role] = latest.get(column, "").map(canonical_provider)
        grouped = (
            latest.assign(participant=canonical[role])
            .loc[lambda frame: frame["participant"].ne("Não informado")]
            .groupby("participant", as_index=False)
            .agg(PL=("pl", "sum"), Fundos=("fund_key", "nunique"))
            .sort_values("PL", ascending=False)
            .head(5)
        )
        grouped["PL"] = grouped["PL"].map(_fmt_bi)
        outputs[role] = grouped.rename(columns={"participant": "Participante"})
    same = (canonical["administrador"] == canonical["gestor"]) & (
        canonical["gestor"] == canonical["custodiante"]
    ) & canonical["administrador"].ne("Não informado")
    integrated = (
        latest.assign(participant=canonical["administrador"])
        .loc[same]
        .groupby("participant", as_index=False)
        .agg(PL=("pl", "sum"), Fundos=("fund_key", "nunique"))
        .sort_values("PL", ascending=False)
        .head(5)
    )
    integrated["PL"] = integrated["PL"].map(_fmt_bi)
    outputs["integrados"] = integrated.rename(columns={"participant": "Participante"})
    return outputs


def _top_current_funds(pack: IndustryExecutivePack, limit: int = 15) -> pd.DataFrame:
    latest = pack.latest_funds.copy()
    latest["pl"] = pd.to_numeric(latest["pl"], errors="coerce")
    latest = latest.sort_values("pl", ascending=False).head(limit)
    return pd.DataFrame(
        {
            "#": range(1, len(latest) + 1),
            "Fundo": latest["denominacao"].map(lambda value: _short_name(value, 45)),
            "PL": latest["pl"].map(_fmt_bi),
            "Tipo ANBIMA": latest["anbima_tipo"].replace({ANBIMA_ND: "N/D*"}),
            "Administrador": latest["admin_nome"].map(canonical_provider),
        }
    )


def _write_sheet(writer: pd.ExcelWriter, sheet: str, frame: pd.DataFrame) -> None:
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    safe_sheet = sheet[:31]
    frame.to_excel(writer, sheet_name=safe_sheet, index=False)
    ws = writer.book[safe_sheet]
    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions
    header_fill = PatternFill("solid", fgColor=NAVY)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = Font(name="Arial", bold=True, color=WHITE)
        cell.alignment = Alignment(vertical="center")
    pct_tokens = ("share", "coverage", "growth", "pct", "ratio")
    money_tokens = ("pl_brl", "volume_brl", "valor")
    for column_index, column_name in enumerate(frame.columns, start=1):
        normalized = str(column_name).lower()
        cells = list(ws.iter_cols(min_col=column_index, max_col=column_index))[0]
        content_width = max(len(str(cell.value or "")) for cell in cells[:300]) + 2
        width = min(max(content_width, 10), 52)
        ws.column_dimensions[get_column_letter(column_index)].width = width
        wrap_long_text = safe_sheet in {"Warnings", "Top 20 Outros", "Conflitos Tab IV"} and (
            "warning" in normalized
            or normalized in {"denominacao", "justificativa/fonte", "tab4_warning"}
        )
        for cell in cells[1:]:
            cell.font = Font(name="Arial", size=10)
            cell.alignment = Alignment(vertical="top", wrap_text=wrap_long_text)
            if wrap_long_text and len(str(cell.value or "")) > 52:
                ws.row_dimensions[cell.row].height = max(ws.row_dimensions[cell.row].height or 15, 42)
            if "_pp_" in normalized or normalized.endswith("_pp"):
                cell.number_format = '0.00 "p.p."'
            elif any(token in normalized for token in pct_tokens):
                cell.number_format = "0.00%"
            elif (
                normalized == "pl"
                or normalized.startswith("pl_")
                or normalized.endswith("_pl")
                or any(token in normalized for token in money_tokens)
            ):
                cell.number_format = 'R$ #,##0.00'


def _build_legacy_industry_xlsx_bytes(data_dir: Path = DEFAULT_DATA_DIR) -> bytes:
    """Build the auditable workbook used alongside the executive deck."""

    pack = _load_executive_pack(data_dir)
    industry = _read(data_dir, "industry_monthly.csv")
    annual_history = _annual_history(industry, pack)
    history_columns = [
        "period_label",
        "competencia",
        "n_veiculos",
        "pl_total",
        "pl_fic_fidc",
        "pl_ex_fic",
        "pl_total_growth",
        "pl_ex_fic_growth",
        "tab4_duplicate_cnpjs",
        "tab4_pl_conflict_cnpjs",
    ]
    annual_history = annual_history[
        [column for column in history_columns if column in annual_history.columns]
    ]
    top_20_editable = pack.top_20_outros.copy()
    top_20_editable["Tipo revisado"] = ""
    top_20_editable["Foco revisado"] = ""
    top_20_editable["Justificativa/Fonte"] = ""
    requested = {
        "PL histórico": annual_history,
        "PL anual": pack.annual_pl,
        "Mix ANBIMA": pack.market_share,
        "Top 20 Outros": top_20_editable,
        "Fila curadoria": pack.curation_queue,
        "Hist cotistas": pack.holder_histogram,
        "Monoestrutura": pack.monostructure_history,
        "Rankings ANBIMA": pack.rankings,
        "Cobertura": pack.coverage,
        "Conflitos Tab IV": pack.source_conflicts,
        "Warnings": pd.DataFrame({"warning": list(pack.warnings)}),
    }
    # Keep legacy audit tabs to avoid breaking analysts' existing workflows.
    raw = {
        "Competências": "industry_competence_status.csv",
        "Indústria mensal": "industry_monthly.csv",
        "Ofertas anual": "industry_offers_annual.csv",
        "Posição Itaú": "industry_competitive_position.csv",
        "Ranking ofertas": "industry_offer_rankings.csv.gz",
        "Cedentes": "industry_originators_annual.csv",
        "Investidores hist": "industry_investor_distribution.csv",
        "Tipos investidor": "industry_investor_types.csv",
        "FIDCs >5bi": "industry_large_fund_classification.csv",
    }
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        for sheet, frame in requested.items():
            _write_sheet(writer, sheet, frame)
        for sheet, filename in raw.items():
            _write_sheet(writer, sheet, _read(data_dir, filename))
        from openpyxl.worksheet.datavalidation import DataValidation

        list_sheet = writer.book.create_sheet("_Listas")
        type_values = [*ANBIMA_TYPES, ANBIMA_ND]
        focus_values = sorted(
            {
                focus
                for focuses in ANBIMA_FOCUS_BY_TYPE.values()
                for focus in focuses
                if focus and focus != ANBIMA_ND
            }
        )
        list_sheet.append(["Tipos", "Focos"])
        for index in range(max(len(type_values), len(focus_values))):
            list_sheet.append(
                [
                    type_values[index] if index < len(type_values) else None,
                    focus_values[index] if index < len(focus_values) else None,
                ]
            )
        list_sheet.sheet_state = "hidden"
        ws = writer.book["Top 20 Outros"]
        headers = {cell.value: cell.column_letter for cell in ws[1]}
        type_validation = DataValidation(
            type="list",
            formula1=f"'_Listas'!$A$2:$A${len(type_values) + 1}",
            allow_blank=True,
        )
        focus_validation = DataValidation(
            type="list",
            formula1=f"'_Listas'!$B$2:$B${len(focus_values) + 1}",
            allow_blank=True,
        )
        ws.add_data_validation(type_validation)
        ws.add_data_validation(focus_validation)
        type_validation.add(f"{headers['Tipo revisado']}2:{headers['Tipo revisado']}{len(top_20_editable) + 1}")
        focus_validation.add(f"{headers['Foco revisado']}2:{headers['Foco revisado']}{len(top_20_editable) + 1}")
    return output.getvalue()


def _build_legacy_industry_pptx_bytes(data_dir: Path = DEFAULT_DATA_DIR) -> bytes:
    """Build a 16:9, native and fully editable executive PowerPoint deck."""

    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    pack = _load_executive_pack(data_dir)
    industry = _read(data_dir, "industry_monthly.csv")
    concentration = _read(data_dir, "concentration_monthly.csv")
    segments = _read(data_dir, "segments_monthly.csv")
    offers = _read(data_dir, "industry_offers_annual.csv")
    competitive = _read(data_dir, "industry_competitive_position.csv")
    originators = _read(data_dir, "industry_originators_annual.csv")
    manifest = _read_manifest(data_dir)
    annual_history = _annual_history(industry, pack)
    latest_complete = pack.competences.latest_complete
    latest_available = pack.competences.latest_available
    latest_label = _competence_label(latest_complete, short=True)
    periods = [pack.competences.period_by_competence[c] for c in pack.competences.ordered]
    offers_as_of = _date_label(manifest.get("as_of_date") or "n/d")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def rgb(hex_color: str) -> RGBColor:
        return RGBColor.from_string(hex_color)

    def add_text(
        slide,
        text: object,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float = 12,
        color: str = GRAY_900,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        italic: bool = False,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = valign
        frame.margin_left = frame.margin_right = Inches(0)
        frame.margin_top = frame.margin_bottom = Inches(0)
        paragraph = frame.paragraphs[0]
        paragraph.text = str(text)
        paragraph.alignment = align
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.italic = italic
        paragraph.font.color.rgb = rgb(color)
        return box

    def line(slide, x: float, y: float, w: float, color: str = GRAY_300, height: float = 0.012):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
        return shape

    def base_slide(title: str, kicker: str = "INDÚSTRIA FIDC"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(WHITE)
        add_text(slide, kicker.upper(), 0.62, 0.28, 4.0, 0.22, size=10, color=ORANGE, bold=True)
        add_text(slide, title, 0.62, 0.57, 12.0, 0.48, size=22, color=BLACK, bold=True)
        line(slide, 0.62, 1.15, 12.05, color=GRAY_300, height=0.018)
        return slide

    def footer(slide, source: str, page: int):
        line(slide, 0.62, 6.95, 12.05, color=GRAY_200, height=0.01)
        add_text(slide, source, 0.62, 7.02, 11.45, 0.24, size=8, color=GRAY_500)
        add_text(slide, str(page), 12.25, 7.01, 0.38, 0.22, size=8, color=GRAY_500, align=PP_ALIGN.RIGHT)

    def add_table(
        slide,
        frame: pd.DataFrame,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        font_size: float = 10,
        widths: Sequence[float] | None = None,
        highlight: str = "",
        header_color: str = BLACK,
    ):
        if frame.empty:
            add_text(slide, "Sem dados disponíveis", x, y, w, h, size=12, color=GRAY_500)
            return None
        table_shape = slide.shapes.add_table(
            len(frame) + 1, len(frame.columns), Inches(x), Inches(y), Inches(w), Inches(h)
        )
        table = table_shape.table
        if widths:
            total_width = sum(widths)
            for index, share in enumerate(widths):
                table.columns[index].width = Inches(w * share / total_width)
        for col, name in enumerate(frame.columns):
            cell = table.cell(0, col)
            cell.text = str(name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(header_color)
        for row_index, values in enumerate(frame.itertuples(index=False), start=1):
            row_highlight = bool(highlight) and any(highlight.lower() in str(value).lower() for value in values)
            for col_index, value in enumerate(values):
                cell = table.cell(row_index, col_index)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(
                    ORANGE_LIGHT if row_highlight else (WHITE if row_index % 2 else GRAY_100)
                )
        for row_index, row in enumerate(table.rows):
            for cell in row.cells:
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.025)
                cell.margin_bottom = Inches(0.025)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = "Arial"
                    paragraph.font.size = Pt(font_size)
                    paragraph.font.bold = row_index == 0
                    paragraph.font.color.rgb = rgb(WHITE if row_index == 0 else GRAY_900)
                    for run in paragraph.runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(font_size)
                        run.font.bold = row_index == 0
                        run.font.color.rgb = rgb(WHITE if row_index == 0 else GRAY_900)
        return table

    def add_chart(
        slide,
        categories: Iterable[object],
        series: Sequence[tuple[str, Iterable[object]]],
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        chart_type,
        colors: Sequence[str] | None = None,
        legend: bool = True,
        value_format: str = "0.0",
        min_scale: float | None = None,
        max_scale: float | None = None,
    ):
        data = CategoryChartData()
        data.categories = [str(value) for value in categories]
        for name, values in series:
            data.add_series(str(name), [_as_float(value) for value in values])
        chart = slide.shapes.add_chart(
            chart_type, Inches(x), Inches(y), Inches(w), Inches(h), data
        ).chart
        chart.has_title = False
        chart.has_legend = legend
        if legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.name = "Arial"
            chart.legend.font.size = Pt(9)
            chart.legend.include_in_layout = False
        chart.value_axis.tick_labels.font.name = "Arial"
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.category_axis.tick_labels.font.name = "Arial"
        chart.category_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.number_format = value_format
        chart.value_axis.major_gridlines.format.line.color.rgb = rgb(GRAY_300)
        chart.value_axis.format.line.color.rgb = rgb(GRAY_500)
        chart.category_axis.format.line.color.rgb = rgb(GRAY_500)
        if min_scale is not None:
            chart.value_axis.minimum_scale = min_scale
        if max_scale is not None:
            chart.value_axis.maximum_scale = max_scale
        palette = list(colors or (ORANGE, NAVY, "777D82", "A8ADB1", "D0D3D5"))
        for index, chart_series in enumerate(chart.series):
            color = rgb(palette[index % len(palette)])
            chart_series.format.fill.solid()
            chart_series.format.fill.fore_color.rgb = color
            chart_series.format.line.color.rgb = color
        return chart

    # 1 — cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BLACK)
    line(slide, 0.7, 0.72, 1.0, color=ORANGE, height=0.06)
    add_text(slide, "Indústria de FIDCs", 0.72, 1.55, 11.5, 0.72, size=44, color=WHITE, bold=True)
    add_text(
        slide,
        "Evolução do mercado, composição, prestadores e estrutura de investidores",
        0.74,
        2.45,
        11.1,
        0.45,
        size=20,
        color="D9DCDE",
    )
    add_text(
        slide,
        f"Base consolidada: {latest_label} · prévia {latest_available} excluída por cobertura incompleta",
        0.74,
        3.22,
        11.2,
        0.34,
        size=15,
        color=ORANGE,
        bold=True,
    )
    add_text(
        slide,
        "Apresentação executiva para Flávio",
        0.74,
        5.93,
        5.0,
        0.3,
        size=13,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "Fontes: CVM Informe Mensal e Ofertas Públicas; ANBIMA Data e Deliberação nº 72.",
        0.74,
        6.48,
        11.4,
        0.3,
        size=10,
        color="B9BDC1",
    )

    # 2 — executive synthesis
    latest_annual = pack.annual_pl.iloc[-1]
    latest_mix = pack.market_share[pack.market_share["competencia"].eq(latest_complete)]
    largest_type = latest_mix[latest_mix["anbima_tipo"].isin(ANBIMA_TYPES)].sort_values(
        "pl_brl", ascending=False
    ).iloc[0]
    mono = pack.monostructure_history[
        pack.monostructure_history["competencia"].eq(latest_complete)
        & pack.monostructure_history["structure_model"].eq("Monoestrutura")
    ].iloc[0]
    holder_cov = pack.holder_coverage.iloc[0]
    slide = base_slide("O mercado cresce, mas sua estrutura exige três leituras distintas", "SÍNTESE EXECUTIVA")
    summary_rows = [
        (
            "01",
            f"{_fmt_bi(latest_annual['pl_ex_fic_brl'], 0)} de PL ex-FIC em {latest_label}",
            f"Alta de {_fmt_pct(latest_annual['pl_ex_fic_growth'])} frente a dez/25; a comparação de 2026 usa a última competência completa.",
        ),
        (
            "02",
            f"{largest_type['anbima_tipo']} lidera o mix, com {_fmt_pct(largest_type['share_ex_fic'])}",
            "O market share mantém N/D separado e marca equivalências CVM como proxy; não transforma dúvida em Outros.",
        ),
        (
            "03",
            f"Monoestrutura representa {_fmt_pct(mono['pl_share_known'])} do PL com prestadores conhecidos",
            "A maior parte do mercado combina plataformas integradas e prestadores especializados — os dois modelos coexistem.",
        ),
        (
            "04",
            f"Histograma cobre {int(holder_cov['included_funds']):,} fundos acima de R$ 200 mi".replace(",", "."),
            f"Cobertura de {_fmt_pct(holder_cov['pl_coverage'])} do PL elegível; quantidade e volume contam histórias diferentes.",
        ),
    ]
    for index, (number, headline, body) in enumerate(summary_rows):
        y = 1.45 + index * 1.30
        add_text(slide, number, 0.72, y, 0.55, 0.32, size=15, color=ORANGE, bold=True)
        add_text(slide, headline, 1.45, y - 0.02, 10.85, 0.36, size=16, color=BLACK, bold=True)
        add_text(slide, body, 1.45, y + 0.43, 10.75, 0.48, size=11.5, color=GRAY_700)
        if index < len(summary_rows) - 1:
            line(slide, 1.45, y + 1.02, 10.9, color=GRAY_200)
    footer(slide, f"Fonte: CVM e ANBIMA; fotografia consolidada em {latest_complete}. Síntese calculada sobre o pack auditável.", 2)

    # 3 — annual PL
    first_history = annual_history.iloc[0]
    last_history = annual_history.iloc[-1]
    slide = base_slide(
        f"O PL ex-FIC avançou de {_fmt_bi(first_history['pl_ex_fic'], 0)} em "
        f"{int(first_history['year'])} para {_fmt_bi(last_history['pl_ex_fic'], 0)} em {latest_label}",
        "EVOLUÇÃO DO PL",
    )
    add_chart(
        slide,
        annual_history["period_label"],
        [
            ("PL bruto", annual_history["pl_total"] / 1e9),
            ("PL ex-FIC", annual_history["pl_ex_fic"] / 1e9),
        ],
        0.68,
        1.45,
        8.3,
        4.9,
        chart_type=XL_CHART_TYPE.LINE,
        colors=("73787D", ORANGE),
        value_format='0 "bi"',
    )
    history_table = pd.DataFrame(
        {
            "Data-base": annual_history["period_label"],
            "PL ex-FIC": annual_history["pl_ex_fic"].map(lambda value: _fmt_bi(value, 0)),
            "Δ": annual_history["pl_ex_fic_growth"].map(
                lambda value: "—" if pd.isna(value) else _fmt_pct(value)
            ),
        }
    ).tail(8)
    add_table(slide, history_table, 9.25, 1.48, 3.4, 4.75, font_size=10, widths=(1.1, 1.4, 0.8))
    footer(slide, f"Fonte: CVM, Informe Mensal FIDC. Dezembro contra dezembro; 2026 = {latest_label}. PL ex-FIC reduz dupla contagem.", 3)

    # 4 — market mix
    slide = base_slide(
        f"{largest_type['anbima_tipo']} lidera o mix; N/D e proxies permanecem visíveis",
        "MIX POR TIPO ANBIMA",
    )
    mix_pivot = pack.market_share.pivot(index="period", columns="anbima_tipo", values="share_ex_fic")
    mix_order = [*ANBIMA_TYPES, ANBIMA_ND]
    mix_pivot = mix_pivot.reindex(index=periods, columns=mix_order, fill_value=0)
    add_chart(
        slide,
        mix_pivot.index,
        [(column if column != ANBIMA_ND else "N/D*", mix_pivot[column]) for column in mix_order],
        0.68,
        1.46,
        7.75,
        4.8,
        chart_type=XL_CHART_TYPE.COLUMN_STACKED_100,
        colors=[TYPE_COLORS[column] for column in mix_order],
        value_format="0%",
        min_scale=0,
        max_scale=1,
    )
    latest_mix_display = latest_mix.sort_values("category_order").copy()
    latest_mix_display["Tipo"] = latest_mix_display["anbima_tipo"].replace({ANBIMA_ND: "N/D*"})
    latest_mix_display["PL"] = latest_mix_display["pl_brl"].map(lambda value: _fmt_bi(value, 0))
    latest_mix_display["Share"] = latest_mix_display["share_ex_fic"].map(_fmt_pct)
    add_table(
        slide,
        latest_mix_display[["Tipo", "PL", "Share"]],
        8.72,
        1.52,
        3.93,
        3.45,
        font_size=10,
        widths=(2.0, 1.0, 0.8),
    )
    latest_coverage = pack.coverage[pack.coverage["competencia"].eq(latest_complete)].iloc[0]
    add_text(
        slide,
        f"Cobertura oficial ANBIMA: {_fmt_pct(latest_coverage['official_anbima_ex_fic_pl_coverage'])}",
        8.77,
        5.12,
        3.7,
        0.28,
        size=11,
        color=BLACK,
        bold=True,
    )
    add_text(
        slide,
        f"Proxy CVM*: {_fmt_pct(latest_coverage['proxy_ex_fic_pl_share'])} · N/D*: {_fmt_pct(latest_coverage['nd_ex_fic_pl_share'])}",
        8.77,
        5.54,
        3.7,
        0.42,
        size=10.5,
        color=RED,
        bold=True,
    )
    add_text(
        slide,
        "* Equivalências proxy exigem validação; N/D não é somado a Outros.",
        8.77,
        6.08,
        3.7,
        0.38,
        size=9.5,
        color=GRAY_700,
    )
    footer(
        slide,
        f"Fonte: ANBIMA Data + CVM. PL ex-FIC em dez/24, dez/25 e {latest_label}. "
        "Classificação ANBIMA dez/25; proxy e N/D sinalizados.",
        4,
    )

    # 5–6 — top 20 Outros, split to keep the names readable
    top_outros = pack.top_20_outros.copy()
    outros_total = _as_float(
        latest_mix.loc[latest_mix["anbima_tipo"].eq("Outros"), "pl_brl"].sum()
    )
    top_outros["Fundo"] = top_outros["denominacao"].map(lambda value: _short_name(value, 56))
    top_outros["Foco ANBIMA / evid."] = top_outros["anbima_foco"]
    top_outros["PL"] = top_outros["pl"].map(_fmt_bi)
    top_outros["% Outros"] = top_outros["pl"].map(
        lambda value: _fmt_pct(_as_float(value) / outros_total if outros_total else 0)
    )
    def top_outros_source(row: pd.Series) -> str:
        label = {
            "oficial_anbima": "ANBIMA oficial",
            "evidencia_publicada": "Evidência publicada",
        }.get(str(row.get("classification_tier")), "Revisar")
        warning = str(row.get("classification_warning") or "")
        if "diverge da evidência documental" in warning:
            return f"{label}**"
        if bool(row.get("classification_requires_warning")):
            return f"{label}*"
        return label

    top_outros["Fonte"] = top_outros.apply(top_outros_source, axis=1)
    top_20_share = top_outros["pl"].sum() / outros_total if outros_total else 0
    official_outros = int(top_outros["classification_tier"].eq("oficial_anbima").sum())
    evidence_outros = int(top_outros["classification_tier"].eq("evidencia_publicada").sum())
    for page, start in ((5, 0), (6, 10)):
        end = start + 10
        slide = base_slide(
            f"Top 20 de Outros concentram {_fmt_pct(top_20_share)} da classe — revisão {start + 1} a {end}",
            "FILA DE RECLASSIFICAÇÃO",
        )
        display = top_outros.iloc[start:end][
            ["outros_rank", "Fundo", "Foco ANBIMA / evid.", "PL", "% Outros", "Fonte"]
        ].rename(columns={"outros_rank": "#"})
        add_table(
            slide,
            display,
            0.63,
            1.43,
            12.05,
            4.92,
            font_size=9.2,
            widths=(0.35, 3.8, 1.35, 0.8, 0.65, 1.15),
            header_color=BLACK,
        )
        add_text(
            slide,
            f"Os 20 fundos estão em Outros segundo o cadastro oficial ANBIMA ({official_outros}) "
            f"ou evidência publicada ({evidence_outros}). * Fotografia ANBIMA de dez/25 aplicada "
            "a mai/26; ** há divergência com evidência documental posterior. Detalhes no XLSX.",
            0.72,
            6.43,
            11.8,
            0.35,
            size=9.7,
            color=GRAY_700,
        )
        footer(slide, f"Fonte: ANBIMA Data e evidência documental publicada; PL CVM em {latest_complete}. Classificação revisada é campo separado no XLSX.", page)

    # 7–9 — overall provider rankings
    for page, role in ((7, "administrador"), (8, "gestor"), (9, "custodiante")):
        ranking = _overall_rankings(pack, role)
        leader = ranking.iloc[0]["Participante"] if not ranking.empty else "N/D"
        slide = base_slide(
            f"{ROLE_LABELS[role]}: {leader} lidera o estoque em {latest_label}",
            f"RANKING GERAL · {ROLE_LABELS[role].upper()}",
        )
        add_table(
            slide,
            ranking,
            0.62,
            1.40,
            12.05,
            5.25,
            font_size=8.8,
            widths=(2.0, 0.65, 0.85, 0.65, 0.85, 0.65, 0.85, 0.8),
            highlight="Itaú",
        )
        historical_provider_note = (
            " *Dez/24 e dez/25 usam cadastro vigente (reconstrução)."
            if role in {"gestor", "custodiante"}
            else ""
        )
        footer(
            slide,
            f"Fonte: CVM, Informe Mensal. Rank e share sobre PL bruto em dez/24, dez/25 e {latest_label}; "
            f"prestadores canônicos por conglomerado.{historical_provider_note}",
            page,
        )

    # 10–12 — all material ANBIMA focuses, one current leader per focus
    for page, role in ((10, "administrador"), (11, "gestor"), (12, "custodiante")):
        focus_table = _focus_ranking_table(pack, role)
        slide = base_slide(
            f"Os líderes mudam materialmente entre os focos ANBIMA — {ROLE_LABELS[role].lower()}",
            "RANKING POR TIPO E FOCO ANBIMA",
        )
        add_table(
            slide,
            focus_table,
            0.58,
            1.36,
            12.18,
            5.45,
            font_size=7.8,
            widths=(2.15, 1.35, 0.55, 0.72, 0.55, 0.72, 0.55, 0.72),
            highlight="Itaú",
        )
        historical_provider_note = (
            " Gestor/custodiante de dez/24 e dez/25 usam o cadastro CVM vigente; reconstrução indicativa.*"
            if role in {"gestor", "custodiante"}
            else ""
        )
        footer(
            slide,
            f"Fonte: ANBIMA Data + CVM. Líder por foco (PL ≥ R$ 1 bi); rank/share em dez/24, dez/25 e {latest_label}. "
            f"Proxies no XLSX.{historical_provider_note}",
            page,
        )

    # 13 — monthly net flows
    slide = base_slide("A captação líquida é volátil mês a mês", "CAPTAÇÃO E FLUXO")
    flow = industry[
        industry["competencia"].astype(str).between("2023-01", latest_complete)
    ].copy()
    flow["label"] = flow["competencia"].map(lambda value: _competence_label(value, short=True))
    add_chart(
        slide,
        flow["label"],
        [("Captação líquida", flow["captacao_liquida"] / 1e9)],
        0.65,
        1.46,
        8.15,
        4.85,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        colors=(ORANGE,),
        legend=False,
        value_format='0 "bi"',
    )
    annual_flows = flow.assign(year=flow["competencia"].str[:4]).groupby("year", as_index=False).agg(
        Captações=("captacoes", "sum"),
        Resgates=("resgates", "sum"),
        Amortizações=("amortizacoes", "sum"),
        Líquida=("captacao_liquida", "sum"),
    )
    for column in ("Captações", "Resgates", "Amortizações", "Líquida"):
        annual_flows[column] = annual_flows[column].map(lambda value: _fmt_bi(value, 0))
    annual_flows = annual_flows.rename(columns={"year": "Ano"}).tail(4)
    add_table(slide, annual_flows, 9.02, 1.50, 3.65, 3.6, font_size=8.8, widths=(0.6, 1, 1, 1, 1))
    add_text(slide, "Captação líquida = captações − resgates − amortizações.", 9.08, 5.45, 3.4, 0.55, size=10.5, color=GRAY_700, bold=True)
    footer(slide, f"Fonte: CVM, Informe Mensal FIDC até {latest_complete}. Junho/26 permanece fora da série por cobertura incompleta.", 13)

    # 14 — accounts and vehicles
    slide = base_slide("O mercado ganhou veículos e contas de cotistas em paralelo", "BASE INVESTIDORA")
    accounts = industry[
        industry["competencia"].astype(str).between("2021-01", latest_complete)
    ].copy()
    accounts["label"] = accounts["competencia"].map(lambda value: _competence_label(value, short=True))
    add_text(slide, "Contas de cotistas (mil)", 0.72, 1.38, 5.7, 0.25, size=12, color=BLACK, bold=True)
    add_chart(
        slide,
        accounts["label"],
        [("Contas", accounts["cotistas_total"] / 1e3)],
        0.66,
        1.72,
        5.9,
        4.55,
        chart_type=XL_CHART_TYPE.LINE,
        colors=(ORANGE,),
        legend=False,
        value_format='0 "mil"',
    )
    add_text(slide, "Veículos reportantes", 6.90, 1.38, 5.7, 0.25, size=12, color=BLACK, bold=True)
    add_chart(
        slide,
        accounts["label"],
        [("Veículos", accounts["n_veiculos"])],
        6.83,
        1.72,
        5.83,
        4.55,
        chart_type=XL_CHART_TYPE.LINE,
        colors=(NAVY,),
        legend=False,
        value_format="0",
    )
    footer(slide, "Fonte: CVM, Informe Mensal FIDC. Contas por classe/série (Tabela X.1); não representam CPFs únicos.", 14)

    # 15 — delinquency
    delinquency = industry[
        industry["competencia"].astype(str).between("2021-01", latest_complete)
    ].copy()
    delinquency["label"] = delinquency["competencia"].map(lambda value: _competence_label(value, short=True))
    latest_delinquency = delinquency.sort_values("competencia").iloc[-1]
    slide = base_slide(
        f"A inadimplência ajustada está em {_fmt_pct(latest_delinquency['inad_pct_ajustada'])}; "
        f"a bruta, em {_fmt_pct(latest_delinquency['inad_pct'])}",
        "QUALIDADE DA CARTEIRA",
    )
    add_chart(
        slide,
        delinquency["label"],
        [
            ("Inadimplência bruta", delinquency["inad_pct"]),
            ("Inadimplência ajustada", delinquency["inad_pct_ajustada"]),
        ],
        0.72,
        1.48,
        11.9,
        4.9,
        chart_type=XL_CHART_TYPE.LINE,
        colors=("6D7276", ORANGE),
        value_format="0.0%",
        min_scale=0,
    )
    add_text(
        slide,
        "Ajustada: limita, veículo a veículo, a inadimplência reportada à própria carteira de direitos creditórios antes da agregação, reduzindo distorções de NPL a valor de face.",
        0.78,
        6.42,
        11.6,
        0.35,
        size=10,
        color=GRAY_700,
    )
    footer(slide, f"Fonte: CVM, Informe Mensal FIDC até {latest_complete}. Percentuais sobre carteira de direitos creditórios.", 15)

    # 16 — administrator concentration, matching the site chart
    concentration = concentration[
        concentration["competencia"].astype(str).between("2021-01", latest_complete)
    ].copy()
    concentration["label"] = concentration["competencia"].map(
        lambda value: _competence_label(value, short=True)
    )
    latest_concentration = concentration.sort_values("competencia").iloc[-1]
    slide = base_slide(
        f"Top 5 administradores concentram {_fmt_pct(latest_concentration['share_top5'])} do PL; "
        f"Top 10, {_fmt_pct(latest_concentration['share_top10'])}",
        "CONCENTRAÇÃO DE ADMINISTRADORES",
    )
    add_chart(
        slide,
        concentration["label"],
        [
            ("Top 10", concentration["share_top10"]),
            ("Top 5", concentration["share_top5"]),
        ],
        0.72,
        1.48,
        11.9,
        4.9,
        chart_type=XL_CHART_TYPE.LINE,
        colors=(ORANGE, NAVY),
        value_format="0%",
        min_scale=0,
        max_scale=1,
    )
    footer(
        slide,
        f"Fonte: CVM, Informe Mensal FIDC até {latest_complete}. Participação sobre o PL administrado; consolidação por CNPJ do administrador.",
        16,
    )

    # 17 — official CVM receivables mix, matching the site chart
    latest_segments = segments[
        segments["competencia"].astype(str).eq(latest_complete)
        & segments["nivel"].astype(str).eq("top")
    ].copy()
    latest_segments["valor"] = pd.to_numeric(latest_segments["valor"], errors="coerce")
    latest_segments = latest_segments[latest_segments["valor"].gt(50_000_000)].sort_values(
        "valor", ascending=True
    )
    slide = base_slide(
        "A carteira combina recebíveis financeiros, comerciais e corporativos",
        "CARTEIRA POR TIPO DE RECEBÍVEL",
    )
    add_chart(
        slide,
        latest_segments["segmento"],
        [("Carteira", latest_segments["valor"] / 1e9)],
        0.72,
        1.43,
        11.85,
        5.05,
        chart_type=XL_CHART_TYPE.BAR_CLUSTERED,
        colors=(ORANGE,),
        legend=False,
        value_format='0 "bi"',
        min_scale=0,
    )
    footer(
        slide,
        f"Fonte: CVM, Informe Mensal FIDC, Tabela II em {latest_complete}. Segmentos com valor superior a R$ 50 milhões; valores em R$ bilhões.",
        17,
    )

    # 18–19 — holder histograms, exact same buckets
    histogram = pack.holder_histogram
    histogram_pivot_count = histogram.pivot(index="cotistas_bucket", columns="anbima_tipo", values="fund_count").reindex(
        index=HOLDER_BUCKETS, columns=[*ANBIMA_TYPES, ANBIMA_ND], fill_value=0
    )
    histogram_pivot_pl = histogram.pivot(index="cotistas_bucket", columns="anbima_tipo", values="pl_brl").reindex(
        index=HOLDER_BUCKETS, columns=[*ANBIMA_TYPES, ANBIMA_ND], fill_value=0
    )
    hist_colors = [TYPE_COLORS[column] for column in [*ANBIMA_TYPES, ANBIMA_ND]]
    slide = base_slide("Fundos acima de R$ 200 mi se distribuem de forma desigual por número de cotistas", "HISTOGRAMA · QUANTIDADE")
    add_chart(
        slide,
        histogram_pivot_count.index,
        [
            (column if column != ANBIMA_ND else "N/D*", histogram_pivot_count[column])
            for column in histogram_pivot_count.columns
        ],
        0.68,
        1.46,
        11.95,
        4.95,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        colors=hist_colors,
        value_format="0",
    )
    footer(slide, f"Fonte: CVM + ANBIMA, {latest_complete}. Eixo X = cotistas (0; 1; 2–3; 4–10; 11–50; 51+); Y = quantidade de fundos.", 18)

    slide = base_slide("Em volume, as faixas de cotistas revelam outra concentração", "HISTOGRAMA · PL")
    add_chart(
        slide,
        histogram_pivot_pl.index,
        [
            (column if column != ANBIMA_ND else "N/D*", histogram_pivot_pl[column] / 1e9)
            for column in histogram_pivot_pl.columns
        ],
        0.68,
        1.46,
        11.95,
        4.95,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        colors=hist_colors,
        value_format='0 "bi"',
    )
    footer(slide, f"Fonte: CVM + ANBIMA, {latest_complete}. Mesmo universo e mesmos buckets do slide anterior; Y = PL em R$ bilhões.", 19)

    # 20 — monostructure history and current decomposition
    mono_history = pack.monostructure_history[
        pack.monostructure_history["structure_model"].eq("Monoestrutura")
    ].set_index("period").reindex(periods)
    slide = base_slide(
        f"Monoestrutura representa {_fmt_pct(mono_history.iloc[-1]['pl_share_known'])} do PL conhecido; "
        "modelos híbridos seguem dominantes",
        "MODELO DE PRESTAÇÃO",
    )
    add_text(slide, "Monoestrutura entre dados conhecidos", 0.72, 1.36, 5.65, 0.25, size=12, color=BLACK, bold=True)
    monostructure_labels = [f"{period}*" for period in periods[:-1]] + [latest_label]
    add_chart(
        slide,
        monostructure_labels,
        [
            ("% fundos", mono_history["fund_share_known"]),
            ("% PL", mono_history["pl_share_known"]),
        ],
        0.66,
        1.72,
        5.75,
        4.45,
        chart_type=XL_CHART_TYPE.LINE,
        colors=(NAVY, ORANGE),
        value_format="0%",
        min_scale=0,
    )
    latest_models = pack.monostructure_history[
        pack.monostructure_history["competencia"].eq(latest_complete)
    ].sort_values("model_order")
    add_text(slide, f"Decomposição do PL em {latest_label}", 6.88, 1.36, 5.6, 0.25, size=12, color=BLACK, bold=True)
    model_display = pd.DataFrame(
        {
            "Modelo": latest_models["structure_model"].str.replace("Administração", "Adm.", regex=False),
            "% PL total": latest_models["pl_share_total"].map(_fmt_pct),
            "Fundos": latest_models["funds"].astype(int),
        }
    )
    add_table(slide, model_display, 6.86, 1.73, 5.77, 4.45, font_size=9.3, widths=(2.3, 0.8, 0.7))
    footer(
        slide,
        f"Fonte: CVM. Mai/26 é fotografia vigente; dez/24 e dez/25* são reconstruções indicativas, pois gestor/custodiante vêm do cadastro atual. Cobertura de PL atual: {_fmt_pct(latest_models.iloc[0]['provider_pl_coverage'])}.",
        20,
    )

    # 21 — integrated and specialist leaders
    slide = base_slide("Plataformas integradas e especialistas coexistem entre os líderes", "LÍDERES POR MODELO")
    leaders = _current_provider_leaders(pack)
    blocks = [
        ("Monoestrutura", "integrados", 0.65, 1.45),
        ("Administração", "administrador", 6.78, 1.45),
        ("Gestão", "gestor", 0.65, 4.15),
        ("Custódia", "custodiante", 6.78, 4.15),
    ]
    for label, key, x, y in blocks:
        add_text(slide, label, x, y, 5.7, 0.25, size=12, color=BLACK, bold=True)
        add_table(
            slide,
            leaders[key],
            x,
            y + 0.34,
            5.88,
            2.05,
            font_size=8.8,
            widths=(2.2, 0.9, 0.65),
            highlight="Itaú",
        )
    footer(slide, f"Fonte: CVM, estoque em {latest_complete}. Monoestrutura = mesmo conglomerado em administração, gestão e custódia.", 21)

    # 22 — offers and origination
    slide = base_slide("Ofertas mantêm ritmo elevado e a originação permanece pulverizada", "OFERTAS E ORIGINAÇÃO")
    add_chart(
        slide,
        offers["period"],
        [
            ("Volume registrado", offers["valid_registered_volume_brl"] / 1e9),
            ("Volume inicial", offers["initial_registered_volume_brl"] / 1e9),
        ],
        0.65,
        1.48,
        6.2,
        4.75,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        colors=(NAVY, ORANGE),
        value_format='0 "bi"',
    )
    latest_originators = originators[
        originators["period"].astype(str).eq(str(offers.iloc[-1]["period"]))
    ].sort_values("rank").head(9)
    originator_display = pd.DataFrame(
        {
            "Cedente / originador": latest_originators["originator_group"],
            "Volume": latest_originators["volume_brl"].map(_fmt_bi),
            "% total": latest_originators["share_of_total"].map(_fmt_pct),
        }
    )
    add_text(slide, "Principais originadores nomináveis — 2026YTD", 7.15, 1.40, 5.2, 0.28, size=12, color=BLACK, bold=True)
    add_table(slide, originator_display, 7.12, 1.78, 5.5, 4.45, font_size=9.2, widths=(2.1, 0.9, 0.75))
    coverage = _as_float(latest_originators["identified_volume_coverage"].max()) if not latest_originators.empty else 0
    footer(slide, f"Fonte: CVM Ofertas Públicas até {offers_as_of}. Ranking nominal cobre {_fmt_pct(coverage)} do volume; não identificado fica fora.", 22)

    # 23 — largest funds
    slide = base_slide("Os maiores FIDCs concentram teses muito diferentes sob o mesmo veículo", "TOP FUNDOS")
    add_table(
        slide,
        _top_current_funds(pack),
        0.60,
        1.34,
        12.13,
        5.50,
        font_size=8.1,
        widths=(0.3, 3.0, 0.7, 1.45, 1.35),
        highlight="Itaú",
    )
    footer(slide, f"Fonte: CVM + ANBIMA, {latest_complete}. N/D* permanece sem reenquadramento automático; PL inclui fundos e classes agregados por CNPJ.", 23)

    # 24 — methodology and sources
    slide = base_slide("A leitura é defensável porque separa dado oficial, evidência e inferência", "METODOLOGIA E FONTES")
    methodology = pd.DataFrame(
        [
            ["PL, fluxos, cotistas, inadimplência", "CVM oficial", f"Dez/dez; 2026 = {latest_label}. Junho/26 é prévia parcial e foi excluído."],
            [
                "Tipo e foco ANBIMA",
                "Oficial / evidência*",
                "Fotografia pública de dez/25 aplicada como ponte cadastral; divergências, proxies e novos fundos permanecem sinalizados.",
            ],
            ["Equivalência CVM", "Proxy*", "Usada apenas quando a taxonomia CVM permite ponte segura; sempre marcada para validação."],
            ["N/D", "Não disponível*", "Nunca é convertido silenciosamente em Outros e permanece visível no mix e no XLSX."],
            ["Ranking de prestadores", "CVM + de-para", "Consolidação por conglomerado; ausência de prestador reduz cobertura, não gera inferência."],
            [
                "Monoestrutura",
                "Atual / histórico proxy*",
                "Mai/26 é fotografia atual; 2024/2025 combinam administrador mensal com gestor/custodiante do cadastro vigente.",
            ],
            ["Top 20 Outros", "Curadoria", "Campos revisados no Excel não sobrescrevem tipo, foco, status ou fonte oficiais."],
        ],
        columns=["Dimensão", "Status", "Regra"],
    )
    add_table(slide, methodology, 0.62, 1.38, 12.05, 4.86, font_size=9.3, widths=(1.65, 1.2, 5.2))
    add_text(
        slide,
        "Fontes primárias: dados.cvm.gov.br (Informe Mensal FIDC e Ofertas Públicas); ANBIMA Data; ANBIMA Deliberação nº 72; CVM/FundosNet para evidência documental.",
        0.72,
        6.43,
        11.7,
        0.38,
        size=9.6,
        color=GRAY_700,
    )
    footer(slide, f"Elaboração: Toma Conta · base consolidada {latest_complete} · workbook anexo preserva tabelas, coberturas e fila de curadoria.", 24)

    output = BytesIO()
    prs.save(output)
    return _normalize_chart_axis_ids(output.getvalue())


def build_industry_pptx_bytes(data_dir: Path = DEFAULT_DATA_DIR) -> bytes:
    """Return the audited 42-slide deck used by the Industry Data surface.

    The reviewed presentation is generated with ``@oai/artifact-tool`` from
    the same versioned payload consumed by the application.  A visually
    different legacy deck is intentionally not used as a silent fallback.
    """

    from services.industry_revision_export import build_revision_pptx_bytes

    return build_revision_pptx_bytes(data_dir)


def build_industry_xlsx_bytes(data_dir: Path = DEFAULT_DATA_DIR) -> bytes:
    """Return the audited workbook paired with the reviewed presentation."""

    from services.industry_revision_export import build_revision_xlsx_bytes

    return build_revision_xlsx_bytes(data_dir)


__all__ = [
    "build_industry_pptx_bytes",
    "build_industry_xlsx_bytes",
    "_normalize_chart_axis_ids",
]
