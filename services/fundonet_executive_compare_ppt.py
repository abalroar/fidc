from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Sequence
from zoneinfo import ZoneInfo

import pandas as pd


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
LEFT_IN = 0.55
RIGHT_IN = 0.40
TOP_IN = 0.45
CONTENT_WIDTH_IN = SLIDE_WIDTH_IN - LEFT_IN - RIGHT_IN

BLACK = "1F1F1F"
ORANGE = "EC7000"
WHITE = "FFFFFF"
MID_GRAY = "757575"
GRID_GRAY = "E0E0E0"
SOFT_GRAY = "F7F7F7"
HIGHLIGHT_FILL = "FFF2E8"
FONT = "Calibri"


def build_executive_comparison_pptx_bytes(
    comparison_df: pd.DataFrame,
    *,
    highlighted_column: str | None = None,
    title: str = "Comparativo Executivo de FIDCs",
    subtitle: str | None = None,
    generated_at: datetime | None = None,
) -> bytes:
    """Export an already-built comparison dataframe to PPTX.

    The input dataframe is the source of truth. This function intentionally does
    not receive dashboards and does not recalculate metrics.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("Dependência python-pptx não instalada.") from exc

    frame = _normalize_frame(comparison_df)
    if generated_at is None:
        generated_at = datetime.now(ZoneInfo("America/Sao_Paulo"))
    elif generated_at.tzinfo is None:
        generated_at = generated_at.replace(tzinfo=ZoneInfo("America/Sao_Paulo"))
    else:
        generated_at = generated_at.astimezone(ZoneInfo("America/Sao_Paulo"))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]

    def rgb(hex_color: str):  # noqa: ANN202
        value = str(hex_color).strip().lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    def add_textbox(slide, left, top, width, height, text, *, size, bold=False, color=BLACK, align=PP_ALIGN.LEFT):  # noqa: ANN001, ANN202
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame_text = box.text_frame
        frame_text.clear()
        frame_text.word_wrap = True
        frame_text.auto_size = MSO_AUTO_SIZE.NONE
        frame_text.margin_left = Inches(0.02)
        frame_text.margin_right = Inches(0.02)
        frame_text.margin_top = Inches(0.01)
        frame_text.margin_bottom = Inches(0.01)
        paragraph = frame_text.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(text or "")
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return box

    def add_footer(slide, page_number: int, total_pages: int) -> None:  # noqa: ANN001
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(LEFT_IN),
            Inches(6.95),
            Inches(CONTENT_WIDTH_IN),
            Inches(0.006),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = rgb(GRID_GRAY)
        line.line.fill.background()
        stamp = generated_at.strftime("%d/%m/%Y %H:%M")
        add_textbox(
            slide,
            LEFT_IN,
            7.08,
            8.6,
            0.18,
            f"Fonte: Informe Mensal CVM | Gerado em: {stamp}",
            size=7.5,
            color=MID_GRAY,
        )
        add_textbox(
            slide,
            10.9,
            7.08,
            1.9,
            0.18,
            f"Página {page_number} de {total_pages}",
            size=7.5,
            color=MID_GRAY,
            align=PP_ALIGN.RIGHT,
        )

    def add_header(slide, page_index: int, total_pages: int) -> None:  # noqa: ANN001
        add_textbox(slide, LEFT_IN, TOP_IN, 8.7, 0.34, title, size=23, bold=True, color=BLACK)
        if subtitle:
            add_textbox(slide, LEFT_IN, TOP_IN + 0.36, 10.8, 0.20, subtitle, size=9, color=MID_GRAY)
        add_footer(slide, page_index, total_pages)

    max_rows = _rows_per_slide(frame)
    chunks = [frame.iloc[idx : idx + max_rows].copy() for idx in range(0, len(frame), max_rows)]
    if not chunks:
        chunks = [frame]
    total_pages = len(chunks)
    for page_index, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(blank)
        add_header(slide, page_index, total_pages)
        _add_comparison_table(
            slide,
            chunk,
            highlighted_column=highlighted_column,
            rgb=rgb,
            left=LEFT_IN,
            top=1.08,
            width=CONTENT_WIDTH_IN,
            height=5.68,
        )

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _normalize_frame(frame: pd.DataFrame) -> pd.DataFrame:
    if frame is None or frame.empty:
        return pd.DataFrame({"Métrica": ["Sem dados disponíveis"]})
    output = frame.copy()
    output = output.fillna("—")
    for column in output.columns:
        output[column] = output[column].map(lambda value: str(value).strip() if str(value).strip() else "—")
    if "Métrica" not in output.columns:
        output.insert(0, "Métrica", [f"Linha {idx + 1}" for idx in range(len(output))])
    return output


def _rows_per_slide(frame: pd.DataFrame) -> int:
    if frame.empty:
        return 1
    cols = max(1, len(frame.columns))
    if cols <= 4:
        return 24
    if cols <= 5:
        return 22
    return 18


def _column_widths(columns: Sequence[str], total_width: float) -> list[float]:
    cols = len(columns)
    if cols <= 1:
        return [total_width]
    metric_width = 2.45 if cols <= 5 else 2.25
    metric_width = min(metric_width, total_width * 0.34)
    other_width = (total_width - metric_width) / (cols - 1)
    return [metric_width, *([other_width] * (cols - 1))]


def _add_comparison_table(  # noqa: PLR0913
    slide,
    frame: pd.DataFrame,
    *,
    highlighted_column: str | None,
    rgb,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    rows = max(1, len(frame.index))
    cols = max(1, len(frame.columns))
    table = slide.shapes.add_table(rows + 1, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    col_widths = _column_widths(list(frame.columns), width)
    for idx, col_width in enumerate(col_widths):
        table.columns[idx].width = Inches(col_width)
    row_height = height / (rows + 1)
    for row in table.rows:
        row.height = Inches(row_height)

    for col_idx, column in enumerate(frame.columns):
        table.cell(0, col_idx).text = str(column)
    for row_idx, (_, row) in enumerate(frame.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = str(value)

    for row_idx in range(rows + 1):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            is_header = row_idx == 0
            is_highlight = bool(highlighted_column) and frame.columns[col_idx] == highlighted_column
            cell.fill.solid()
            if is_header and is_highlight:
                cell.fill.fore_color.rgb = rgb(ORANGE)
            elif is_header:
                cell.fill.fore_color.rgb = rgb(BLACK)
            elif is_highlight:
                cell.fill.fore_color.rgb = rgb(HIGHLIGHT_FILL)
            else:
                cell.fill.fore_color.rgb = rgb(WHITE if row_idx % 2 else SOFT_GRAY)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(_font_size(rows=rows, cols=cols, is_header=is_header))
                    run.font.bold = is_header or col_idx == 0
                    run.font.color.rgb = rgb(WHITE if is_header else BLACK)
            cell.text_frame.word_wrap = True


def _font_size(*, rows: int, cols: int, is_header: bool) -> float:
    if is_header:
        return 8.5 if cols <= 5 else 7.5
    if rows >= 22 or cols >= 6:
        return 7.2
    if rows >= 18:
        return 7.8
    return 8.4
