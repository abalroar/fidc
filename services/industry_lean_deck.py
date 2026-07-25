"""Recorte enxuto do deck executivo publicado.

A versão enxuta não é renderizada de novo: ela seleciona slides do PPTX que já
passou pela validação do bundle. Isso mantém gráficos, paleta e tipografia
idênticos ao deck completo e elimina a possibilidade de as duas versões
divergirem. A seleção é feita pelo texto de cabeçalho de cada slide, e não por
índice, para que a inserção de slides no meio do deck não quebre o recorte.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
import re
import unicodedata

from pptx import Presentation

from services.industry_revision_export import (
    DEFAULT_DATA_DIR,
    build_revision_pptx_bytes,
)


# Ordem do recorte para DCM e book comercial: escala, comparação com renda fixa,
# base investidora, taxonomia, prestadores, bloco de ofertas e as duas páginas de
# conclusão. Cada entrada é o cabeçalho do slide no deck completo.
LEAN_SLIDE_HEADINGS: tuple[str, ...] = (
    "INDÚSTRIA DE FIDCs",
    "ESCALA DA INDÚSTRIA",
    "OFERTAS ENCERRADAS · RENDA FIXA",
    "BASE INVESTIDORA",
    "DISTRIBUIÇÃO POR NÚMERO DE COTISTAS",
    "TAXONOMIA VIGENTE",
    "PRESTADORES · EVOLUÇÃO E RANKING",
    "OFERTAS ENCERRADAS · VOLUME E TICKET",
    "OFERTAS ENCERRADAS · DISTRIBUIÇÃO DO TICKET",
    "OFERTAS · VOLUME E REGIME",
    "TOP 15 · OFERTAS ENCERRADAS",
    "PRINCIPAIS CONCLUSÕES",
    "CONCLUSÕES ESTRATÉGICAS · DCM",
)


class LeanDeckUnavailable(RuntimeError):
    """Levantada quando o deck publicado não contém os slides do recorte."""


def _normalize(value: str) -> str:
    stripped = unicodedata.normalize("NFKD", value)
    without_marks = "".join(char for char in stripped if not unicodedata.combining(char))
    collapsed = re.sub(r"\s+", " ", without_marks).strip().upper()
    return collapsed.replace("·", "-").replace("—", "-").replace("–", "-")


def _slide_heading(slide) -> str:
    """Primeiro texto não vazio do slide, que é sempre o eyebrow do cabeçalho."""

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text:
            return text.splitlines()[0].strip()
    return ""


def _drop_slide(presentation: Presentation, index: int) -> None:
    slide_id_list = presentation.slides._sldIdLst  # noqa: SLF001 - única via de remoção
    entry = list(slide_id_list)[index]
    presentation.part.drop_rel(entry.rId)
    slide_id_list.remove(entry)


def build_lean_pptx_bytes(
    pptx_bytes: bytes,
    headings: tuple[str, ...] = LEAN_SLIDE_HEADINGS,
) -> bytes:
    """Devolve o deck contendo apenas ``headings``, na ordem em que aparecem.

    Falha fechado: se qualquer cabeçalho esperado não existir no deck publicado,
    nada é gerado. O caso mais provável é o bundle ainda não ter sido republicado
    após a inclusão de um slide novo.
    """

    presentation = Presentation(BytesIO(pptx_bytes))
    wanted = {_normalize(heading) for heading in headings}
    found: dict[str, int] = {}
    keep: list[int] = []
    for index, slide in enumerate(presentation.slides):
        key = _normalize(_slide_heading(slide))
        if key in wanted and key not in found:
            found[key] = index
            keep.append(index)

    missing = [heading for heading in headings if _normalize(heading) not in found]
    if missing:
        raise LeanDeckUnavailable(
            "Deck publicado não contém: " + "; ".join(missing)
        )

    total = len(presentation.slides)
    for index in sorted(set(range(total)) - set(keep), reverse=True):
        _drop_slide(presentation, index)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_lean_revision_pptx_bytes(data_dir: Path = DEFAULT_DATA_DIR) -> bytes:
    """Recorte enxuto a partir do bundle publicado e já validado."""

    return build_lean_pptx_bytes(build_revision_pptx_bytes(data_dir))


__all__ = [
    "LEAN_SLIDE_HEADINGS",
    "LeanDeckUnavailable",
    "build_lean_pptx_bytes",
    "build_lean_revision_pptx_bytes",
]
