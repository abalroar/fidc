"""Reordenação de lâminas e renumeração de páginas no deck montado.

Duas operações que só fazem sentido **depois** que todas as seções já entraram:
mover um bloco de lâminas para outro ponto do deck e acertar os números de
página.

Mover é seguro porque o que muda é a ordem dos ``sldId`` no XML — as partes
ficam onde estão.  O caminho perigoso é o oposto: remover uma lâmina do meio e
recriá-la, porque o ``next_partname`` do python-pptx devolve nomes de parte já
ocupados assim que a numeração deixa de ser contígua, e o pacote sai com dois
``ppt/slides/slideN.xml``.

Os números de página não são campos do PowerPoint: são caixas de texto escritas
na montagem.  Quem reordena precisa reescrevê-las, senão o rodapé passa a
descrever a ordem antiga.
"""

from __future__ import annotations

from pptx.util import Inches


#: Onde a caixa do número de página mora no rodapé padrão do deck.
PAGE_BOX_X_IN = 12.25
PAGE_BOX_MIN_TOP_IN = 6.8
PAGE_BOX_TOLERANCE_IN = 0.3


def page_box(slide):
    """A caixinha do número de página, achada pela geometria do rodapé.

    Nem toda lâmina tem uma — a capa e os separadores não são numerados —, e
    nesse caso a função devolve ``None``.
    """

    for forma in slide.shapes:
        if not forma.has_text_frame:
            continue
        texto = forma.text_frame.text.strip()
        if (
            texto.isdigit()
            and forma.top > Inches(PAGE_BOX_MIN_TOP_IN)
            and abs(forma.left - Inches(PAGE_BOX_X_IN)) < Inches(PAGE_BOX_TOLERANCE_IN)
        ):
            return forma
    return None


def set_page_number(slide, numero: int) -> bool:
    caixa = page_box(slide)
    if caixa is None:
        return False
    caixa.text_frame.paragraphs[0].runs[0].text = str(numero)
    return True


def renumber_pages(presentation) -> int:
    """Cada lâmina numerada passa a exibir a sua própria posição.

    Devolve quantas foram reescritas.  As não numeradas continuam sem número —
    inventar um para a capa mudaria o desenho da lâmina, não só o rodapé.
    """

    escritas = 0
    for posicao, slide in enumerate(presentation.slides, start=1):
        escritas += int(set_page_number(slide, posicao))
    return escritas


def move_slides(presentation, primeiro: int, ultimo: int, destino: int) -> None:
    """Move o bloco ``primeiro..ultimo`` para começar em ``destino``.

    Tudo 1-based e inclusivo, e ``destino`` é lido na **numeração atual**: mover
    9–17 para 52 põe o bloco imediatamente antes da lâmina que hoje é a 52ª.
    Como o bloco sai de antes dela, no deck final ele termina uma posição antes
    disso — e a lâmina que era a 52ª continua sendo a 52ª.
    """

    ids = presentation.slides._sldIdLst
    elementos = list(ids)
    total = len(elementos)
    if not (1 <= primeiro <= ultimo <= total):
        raise IndexError(
            f"O bloco {primeiro}–{ultimo} não cabe num deck de {total} lâminas."
        )
    if primeiro <= destino <= ultimo + 1:
        return

    bloco = elementos[primeiro - 1 : ultimo]
    ancora = elementos[destino - 1] if destino <= total else None
    for elemento in bloco:
        ids.remove(elemento)

    restantes = list(ids)
    if ancora is None:
        posicao = len(restantes)
    else:
        posicao = restantes.index(ancora)
    for deslocamento, elemento in enumerate(bloco):
        ids.insert(posicao + deslocamento, elemento)


__all__ = [
    "PAGE_BOX_X_IN",
    "move_slides",
    "page_box",
    "renumber_pages",
    "set_page_number",
]
