"""Somatório FIDCs — deck v3 para comitê de crédito.

Estrutura: 19 slides
  1  Capa
  2  Divisor 01 — Visão consolidada
  3  Consolidado: PL por classe + subordinação + NPL + cobertura
  4  Consolidado: Carteira + crescimento YoY + NPL por severidade + roll rates
  5  Consolidado: Duration por FIDC + cohorts recentes
  6  Consolidado: Roll rates por mês do ano (placeholder se sem dados)
  7  Divisor 02 — Mercado Crédito FIDC
  8  MC FIDC: PL por classe + subordinação + NPL + cobertura
  9  MC FIDC: Carteira + crescimento + NPL severidade + roll rates
  10 MC FIDC: Duration + cohorts
  11 Divisor 03 — Mercado Crédito I Brasil FIDC
  12 MC I Brasil: PL + subordinação + NPL + cobertura
  13 MC I Brasil: Carteira + crescimento + NPL severidade + roll rates
  14 MC I Brasil: Duration + cohorts
  15 Divisor 04 — Mercado Crédito II Brasil FIDC
  16 MC II Brasil: PL + subordinação + NPL + cobertura
  17 MC II Brasil: Carteira + crescimento + NPL severidade + roll rates
  18 MC II Brasil: Duration + cohorts
  19 Síntese — Pontos de atenção e destaques
"""
from __future__ import annotations

from io import BytesIO
from typing import Any

import pandas as pd

from services.export_chart_labels import choose_export_label_policy, format_export_label

# ---------------------------------------------------------------------------
# Paleta de cores — Itaú BBA
# ---------------------------------------------------------------------------
_BLACK = "1F1F1F"
_ORANGE = "EC7000"
_GRAY = "6B6B6B"
_GRAY_LIGHT = "E5E5E5"
_GRAY_MID = "9C9C9C"
_GRAY_PALE = "BDBDBD"
_CARD_BG = "F5F5F5"
_WHITE = "FFFFFF"
_DARK_BG = "1F1F1F"   # capa e divisores
_SUBTITLE_FG = "CCCCCC"

# ---------------------------------------------------------------------------
# Abreviações de mês
# ---------------------------------------------------------------------------
_PT_MONTH = {
    1: "jan", 2: "fev", 3: "mar", 4: "abr", 5: "mai", 6: "jun",
    7: "jul", 8: "ago", 9: "set", 10: "out", 11: "nov", 12: "dez",
}

# ---------------------------------------------------------------------------
# Mapeamento de nomes longos → (header curto, subtítulo)
# ---------------------------------------------------------------------------
_FUND_NAMES: dict[str, tuple[str, str]] = {
    # chaves normalizadas (upper + strip)
    "MERCADO CRÉDITO FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS RESPONSABILIDADE LIMITADA": (
        "Mercado Crédito FIDC",
        "Mercado Crédito Fundo de Investimento em Direitos Creditórios Resp. Ltda.",
    ),
    "MERCADO CRÉDITO I BRASIL FIDC SEGMENTO FINANCEIRO DE RESPONSABILIDADE LIMITADA": (
        "Mercado Crédito I Brasil FIDC",
        "Mercado Crédito I Brasil FIDC Segmento Financeiro Resp. Ltda.",
    ),
    "MERCADO CRÉDITO II BRASIL FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS DE RESPONSABILIDADE LIMITADA": (
        "Mercado Crédito II Brasil FIDC",
        "Mercado Crédito II Brasil Fundo de Investimento em Direitos Creditórios Resp. Ltda.",
    ),
}

# Fallback: prefixo parcial
_FUND_PREFIXES = [
    ("MERCADO CRÉDITO II BRASIL", "Mercado Crédito II Brasil FIDC",
     "Mercado Crédito II Brasil Fundo de Investimento em Direitos Creditórios Resp. Ltda."),
    ("MERCADO CRÉDITO I BRASIL", "Mercado Crédito I Brasil FIDC",
     "Mercado Crédito I Brasil FIDC Segmento Financeiro Resp. Ltda."),
    ("MERCADO CRÉDITO", "Mercado Crédito FIDC",
     "Mercado Crédito Fundo de Investimento em Direitos Creditórios Resp. Ltda."),
]

# ---------------------------------------------------------------------------
# Geometria do slide
# ---------------------------------------------------------------------------
_SW = 13.333   # largura
_SH = 7.5      # altura
_HDR_H = 0.50  # banda header
_KPI_TOP = 0.55
_KPI_H = 0.70
_CTX_TOP = _KPI_TOP + _KPI_H + 0.08   # ~ 1.33"
_FTR_TOP = 7.20
_MX = 0.50     # margem horizontal
_CTX_H = _FTR_TOP - _CTX_TOP          # ~ 5.87"
_GAP_X = 0.15
_GAP_Y = 0.20
_TOTAL_SLIDES = 19


# ===========================================================================
# ENTRY POINT
# ===========================================================================

def build_somatorio_fidcs_pptx_bytes(
    *,
    outputs: Any,
    monitor_outputs: Any,
    research_outputs: Any | None = None,
) -> bytes:
    """Gera o deck v3 completo (19 slides) em formato PPTX."""
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import (
            XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LABEL_POSITION,
            XL_LEGEND_POSITION, XL_MARKER_STYLE,
        )
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("Dependência python-pptx não instalada.") from exc

    prs = Presentation()
    prs.slide_width = Inches(_SW)
    prs.slide_height = Inches(_SH)
    layout = prs.slide_layouts[6]

    deps = dict(
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_DATA_LABEL_POSITION=XL_DATA_LABEL_POSITION,
        XL_LABEL_POSITION=XL_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )

    # Data-base: último mês do consolidado
    con_monthly = _last_12m(getattr(outputs, "consolidated_monthly", pd.DataFrame()))
    data_base = _detect_data_base(con_monthly)

    # Mapa fund_cnpj → DataFrames
    fund_monthly_map = {
        cnpj: _last_12m(frame)
        for cnpj, frame in getattr(outputs, "fund_monthly", {}).items()
    }
    fund_monitor_map = {
        cnpj: _chart_monthly(frame)
        for cnpj, frame in getattr(monitor_outputs, "fund_monitor", {}).items()
    }
    fund_cohort_map = getattr(monitor_outputs, "fund_cohorts", {})
    con_monitor = _chart_monthly(getattr(monitor_outputs, "consolidated_monitor", pd.DataFrame()))
    con_cohorts = getattr(monitor_outputs, "consolidated_cohorts", pd.DataFrame())
    roll_df = getattr(research_outputs, "roll_seasonality", pd.DataFrame()) if research_outputs else pd.DataFrame()

    page = [0]  # mutable counter

    def next_page() -> int:
        page[0] += 1
        return page[0]

    # ---- Slide 1: Capa ----
    _add_cover(prs, layout, data_base, **deps)
    next_page()

    # ---- Slide 2: Divisor 01 ----
    _add_divider(prs, layout, "01", "Visão consolidada",
                 "PL, carteira, risco e duration agregados", **deps)
    next_page()

    # ---- Slides 3–6: Consolidado ----
    _add_base_slide(prs, layout, con_monthly, con_monitor, "Somatório FIDCs — Consolidado",
                    next_page(), _TOTAL_SLIDES, data_base, is_subordination_fund=False, **deps)
    _add_credit_slide(prs, layout, con_monitor, "Somatório FIDCs — Consolidado",
                      next_page(), _TOTAL_SLIDES, data_base, **deps)
    _add_detail_slide(prs, layout, con_monitor, fund_monitor_map, con_cohorts,
                      "Somatório FIDCs — Consolidado",
                      next_page(), _TOTAL_SLIDES, data_base, **deps)
    _add_roll_slide(prs, layout, con_monitor, roll_df, "Somatório FIDCs — Consolidado",
                    next_page(), _TOTAL_SLIDES, data_base, **deps)

    # ---- Per-fund blocks ----
    section_meta = [
        ("02", _FUND_NAMES.get("MERCADO CRÉDITO FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS RESPONSABILIDADE LIMITADA",
                                ("Mercado Crédito FIDC", "Mercado Crédito Fundo de Investimento em Direitos Creditórios Resp. Ltda."))),
        ("03", _FUND_NAMES.get("MERCADO CRÉDITO I BRASIL FIDC SEGMENTO FINANCEIRO DE RESPONSABILIDADE LIMITADA",
                                ("Mercado Crédito I Brasil FIDC", "Mercado Crédito I Brasil FIDC Segmento Financeiro Resp. Ltda."))),
        ("04", _FUND_NAMES.get("MERCADO CRÉDITO II BRASIL FUNDO DE INVESTIMENTO EM DIREITOS CREDITÓRIOS DE RESPONSABILIDADE LIMITADA",
                                ("Mercado Crédito II Brasil FIDC", "Mercado Crédito II Brasil Fundo de Investimento em Direitos Creditórios Resp. Ltda."))),
    ]
    fund_cnpjs = list(fund_monthly_map.keys())

    for idx, (sec_num, (hdr, sub)) in enumerate(section_meta):
        _add_divider(prs, layout, sec_num, hdr, sub, **deps)
        next_page()

        cnpj = fund_cnpjs[idx] if idx < len(fund_cnpjs) else None
        fm = fund_monthly_map.get(cnpj, pd.DataFrame()) if cnpj else pd.DataFrame()
        mon = fund_monitor_map.get(cnpj, pd.DataFrame()) if cnpj else pd.DataFrame()
        coh = fund_cohort_map.get(cnpj, pd.DataFrame()) if cnpj else pd.DataFrame()
        is_mc_fidc = idx == 0  # outlier fix only for first fund

        _add_base_slide(prs, layout, fm, mon, hdr,
                        next_page(), _TOTAL_SLIDES, data_base,
                        is_subordination_fund=is_mc_fidc, **deps)
        _add_credit_slide(prs, layout, mon, hdr,
                          next_page(), _TOTAL_SLIDES, data_base, **deps)
        _add_detail_slide(prs, layout, mon, {}, coh, hdr,
                          next_page(), _TOTAL_SLIDES, data_base, **deps)

    # ---- Slide 19: Síntese ----
    _add_synthesis(prs, layout, next_page(), _TOTAL_SLIDES, data_base, **deps)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ===========================================================================
# STRUCTURAL SLIDES
# ===========================================================================

def _add_cover(prs, layout, data_base: str, *, Inches, Pt, RGBColor, **_kw) -> None:
    slide = prs.slides.add_slide(layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(_DARK_BG, RGBColor)

    # Linha laranja vertical (4pt × 1.5")
    line = slide.shapes.add_shape(1, Inches(0.60), Inches(2.80), Inches(4 / 72), Inches(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(_ORANGE, RGBColor)
    line.line.fill.background()

    # Título principal
    _textbox(slide, "Mercado de FIDCs",
             left=1.0, top=2.50, width=10.0, height=0.80,
             size=44, bold=True, color=_WHITE, Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    # Subtítulo
    subtitle = f"Visão consolidada e por fundo — Data-base {data_base}"
    _textbox(slide, subtitle,
             left=1.0, top=3.50, width=10.5, height=0.45,
             size=18, bold=False, color=_SUBTITLE_FG, Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    # Rodapé canto inferior direito
    _textbox(slide, "Itaú BBA | Análise Setorial de Crédito",
             left=_SW - 5.5, top=_SH - 0.50, width=5.0, height=0.35,
             size=10, bold=False, color=_SUBTITLE_FG,
             align_right=True, Inches=Inches, Pt=Pt, RGBColor=RGBColor)


def _add_divider(
    prs, layout, number: str, title: str, subtitle: str,
    *, Inches, Pt, RGBColor, **_kw
) -> None:
    slide = prs.slides.add_slide(layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(_DARK_BG, RGBColor)

    # Numeração laranja
    _textbox(slide, f"{number} —",
             left=1.0, top=2.60, width=3.0, height=0.45,
             size=24, bold=True, color=_ORANGE, Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    # Título branco
    _textbox(slide, title,
             left=1.0, top=3.10, width=11.0, height=0.60,
             size=32, bold=True, color=_WHITE, Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    # Subtítulo cinza claro
    _textbox(slide, subtitle,
             left=1.0, top=3.80, width=11.0, height=0.40,
             size=14, bold=False, color=_SUBTITLE_FG, Inches=Inches, Pt=Pt, RGBColor=RGBColor)


def _add_synthesis(
    prs, layout, page_num: int, total: int, data_base: str,
    *, Inches, Pt, RGBColor, **_kw
) -> None:
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(_WHITE, RGBColor)

    _slide_header(slide, "Síntese — Pontos de atenção e destaques", data_base,
                  Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _slide_footer(slide, page_num, total, Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    col_top = _HDR_H + 0.15
    col_h = _FTR_TOP - col_top - 0.10
    col_w = 5.50
    col_gap = 0.50
    col_left_1 = _MX
    col_left_2 = _MX + col_w + col_gap

    bullets = ["[INSERIR PONTO 1]", "[INSERIR PONTO 2]", "[INSERIR PONTO 3]"]

    for col_left, hdr_color, hdr_text in [
        (col_left_1, _ORANGE, "Pontos de atenção"),
        (col_left_2, _BLACK, "Destaques positivos"),
    ]:
        # Header da coluna
        _textbox(slide, hdr_text,
                 left=col_left, top=col_top, width=col_w, height=0.35,
                 size=16, bold=True, color=hdr_color,
                 Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        # Bullets placeholder
        bullet_top = col_top + 0.45
        for bullet in bullets:
            _textbox(slide, f"• {bullet}",
                     left=col_left, top=bullet_top, width=col_w, height=0.30,
                     size=14, bold=True, color=_BLACK,
                     Inches=Inches, Pt=Pt, RGBColor=RGBColor)
            bullet_top += 0.40


# ===========================================================================
# CONTENT SLIDES
# ===========================================================================

def _add_base_slide(
    prs, layout,
    monthly_df: pd.DataFrame,
    monitor_df: pd.DataFrame,
    title: str,
    page_num: int, total: int, data_base: str,
    *,
    is_subordination_fund: bool = False,
    Inches, Pt, RGBColor,
    CategoryChartData, XL_CHART_TYPE, XL_DATA_LABEL_POSITION,
    XL_LEGEND_POSITION, **_kw,
) -> None:
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(_WHITE, RGBColor)

    _slide_header(slide, title, data_base, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _slide_footer(slide, page_num, total, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _kpi_strip(slide, monitor_df, monthly_df, Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    df = monthly_df
    cats = _category_labels(df)
    slots = _grid_2x2()

    _chart_pl_stacked(slide, slots[0], df, cats,
                      CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                      XL_CHART_TYPE=XL_CHART_TYPE, XL_LEGEND_POSITION=XL_LEGEND_POSITION,
                      XL_DATA_LABEL_POSITION=XL_DATA_LABEL_POSITION,
                      Inches=Inches, Pt=Pt)
    _chart_percent_line(slide, slots[1],
                        title="% Subordinação Total ex-360",
                        column="subordinacao_total_ex360_pct",
                        df=df, cats=cats, color=_GRAY,
                        clip_outlier=is_subordination_fund,
                        CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                        XL_CHART_TYPE=XL_CHART_TYPE, Inches=Inches, Pt=Pt)
    _chart_percent_line(slide, slots[2],
                        title="NPL Over 90d ex-360 / Carteira",
                        column="npl_over90_ex360_pct",
                        df=df, cats=cats, color=_BLACK, metric_kind="npl_pct",
                        CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                        XL_CHART_TYPE=XL_CHART_TYPE, Inches=Inches, Pt=Pt)
    _chart_percent_line(slide, slots[3],
                        title="Cobertura PDD / NPL Over 90d ex-360",
                        column="pdd_npl_over90_ex360_pct",
                        df=df, cats=cats, color=_ORANGE, metric_kind="coverage_pct",
                        CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                        XL_CHART_TYPE=XL_CHART_TYPE, Inches=Inches, Pt=Pt)


def _add_credit_slide(
    prs, layout,
    monitor_df: pd.DataFrame,
    title: str,
    page_num: int, total: int, data_base: str,
    *, Inches, Pt, RGBColor,
    CategoryChartData, XL_CHART_TYPE, XL_LABEL_POSITION,
    XL_LEGEND_POSITION, XL_MARKER_STYLE, **_kw,
) -> None:
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(_WHITE, RGBColor)

    _slide_header(slide, title, data_base, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _slide_footer(slide, page_num, total, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _kpi_strip(slide, monitor_df, pd.DataFrame(), Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    df = monitor_df
    cats = _category_labels(df)
    slots = _grid_2x2()

    _chart_bar_money(slide, slots[0], title="Carteira ex-360",
                     column="carteira_ex360", series_name="Carteira ex-360",
                     df=df, cats=cats,
                     CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                     XL_CHART_TYPE=XL_CHART_TYPE,
                     XL_LABEL_POSITION=XL_LABEL_POSITION,
                     Inches=Inches, Pt=Pt)
    _chart_multi_line(slide, slots[1], title="Crescimento YoY carteira ex-360",
                      df=df, cats=cats,
                      series_specs=[("Crescimento YoY", "carteira_ex360_yoy_pct", _ORANGE)],
                      y_title="%", value_is_percent=True,
                      CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                      XL_CHART_TYPE=XL_CHART_TYPE, XL_LABEL_POSITION=XL_LABEL_POSITION,
                      XL_LEGEND_POSITION=XL_LEGEND_POSITION, XL_MARKER_STYLE=XL_MARKER_STYLE,
                      Inches=Inches, Pt=Pt)
    _chart_stacked_npl(slide, slots[2], df=df, cats=cats,
                       CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                       XL_CHART_TYPE=XL_CHART_TYPE,
                       XL_LABEL_POSITION=XL_LABEL_POSITION,
                       XL_LEGEND_POSITION=XL_LEGEND_POSITION,
                       Inches=Inches, Pt=Pt)
    _chart_multi_line(slide, slots[3], title="Roll rates",
                      df=df, cats=cats,
                      series_specs=[
                          ("Roll 61-90 M-3", "roll_61_90_m3_pct", _BLACK),
                          ("Roll 91-120 M-4", "roll_91_120_m4_pct", _GRAY),
                          ("Roll 121-150 M-5", "roll_121_150_m5_pct", _GRAY_MID),
                          ("Roll 151-180 M-6", "roll_151_180_m6_pct", _ORANGE),
                      ],
                      y_title="%", value_is_percent=True, metric_kind="roll_pct",
                      CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                      XL_CHART_TYPE=XL_CHART_TYPE, XL_LABEL_POSITION=XL_LABEL_POSITION,
                      XL_LEGEND_POSITION=XL_LEGEND_POSITION, XL_MARKER_STYLE=XL_MARKER_STYLE,
                      Inches=Inches, Pt=Pt)


def _add_detail_slide(
    prs, layout,
    monitor_df: pd.DataFrame,
    fund_monitor_map: dict[str, pd.DataFrame],
    cohort_df: pd.DataFrame,
    title: str,
    page_num: int, total: int, data_base: str,
    *, Inches, Pt, RGBColor,
    CategoryChartData, XL_CHART_TYPE, XL_LABEL_POSITION,
    XL_LEGEND_POSITION, XL_MARKER_STYLE, **_kw,
) -> None:
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(_WHITE, RGBColor)

    _slide_header(slide, title, data_base, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _slide_footer(slide, page_num, total, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _kpi_strip(slide, monitor_df, pd.DataFrame(), Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    slots = _grid_1x2()

    # Duration chart — legenda ABAIXO do eixo X para evitar colisão
    dur_df = _build_duration_frame(monitor_df, fund_monitor_map)
    dur_series = [c for c in dur_df.columns if c not in {"competencia_dt", "competencia"}]
    dur_specs = [(s, _short_name(s), color)
                 for s, color in zip(dur_series, [_BLACK, _ORANGE, _GRAY, _GRAY_MID])]
    _chart_multi_line(slide, slots[0],
                      title="Duration por FIDC",
                      df=dur_df, cats=_category_labels(dur_df),
                      series_specs=dur_specs,
                      y_title="meses", value_is_percent=False, metric_kind="duration",
                      legend_below=True,
                      CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                      XL_CHART_TYPE=XL_CHART_TYPE, XL_LABEL_POSITION=XL_LABEL_POSITION,
                      XL_LEGEND_POSITION=XL_LEGEND_POSITION, XL_MARKER_STYLE=XL_MARKER_STYLE,
                      Inches=Inches, Pt=Pt)
    _chart_cohorts(slide, slots[1], cohort_df=cohort_df,
                   CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                   XL_CHART_TYPE=XL_CHART_TYPE, XL_LABEL_POSITION=XL_LABEL_POSITION,
                   XL_LEGEND_POSITION=XL_LEGEND_POSITION, XL_MARKER_STYLE=XL_MARKER_STYLE,
                   Inches=Inches, Pt=Pt)


def _add_roll_slide(
    prs, layout,
    monitor_df: pd.DataFrame,
    roll_df: pd.DataFrame,
    title: str,
    page_num: int, total: int, data_base: str,
    *, Inches, Pt, RGBColor,
    CategoryChartData, XL_CHART_TYPE, XL_LABEL_POSITION,
    XL_LEGEND_POSITION, XL_MARKER_STYLE, **_kw,
) -> None:
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(_WHITE, RGBColor)

    _slide_header(slide, title, data_base, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _slide_footer(slide, page_num, total, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _kpi_strip(slide, monitor_df, pd.DataFrame(), Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    if roll_df is None or not isinstance(roll_df, pd.DataFrame) or roll_df.empty:
        _empty_placeholder(slide, _grid_2x2()[0],
                           "Roll rates por mês do ano",
                           "Dados de sazonalidade não disponíveis.",
                           Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        return

    slots = _grid_2x2()
    roll_metrics = [
        ("Roll 61-90 por mês do ano", "roll_61_90_m3"),
        ("Roll 91-120 por mês do ano", "roll_91_120_m4"),
        ("Roll 121-150 por mês do ano", "roll_121_150_m5"),
        ("Roll 151-180 por mês do ano", "roll_151_180_m6"),
    ]
    colors = [_BLACK, _GRAY, _GRAY_MID, _GRAY_PALE]
    for slot, (r_title, metric_id) in zip(slots, roll_metrics):
        chart_df = _research_roll_wide(roll_df, metric_id=metric_id)
        series = [c for c in chart_df.columns if c not in {"competencia_dt", "competencia", "month", "ordem"}]
        _chart_multi_line(slide, slot, title=r_title,
                          df=chart_df,
                          cats=chart_df.get("competencia", pd.Series(dtype="str")).astype(str).tolist(),
                          series_specs=[(s, s, colors[i % len(colors)]) for i, s in enumerate(series)],
                          y_title="%", value_is_percent=True, metric_kind="roll_pct",
                          CategoryChartData=CategoryChartData, RGBColor=RGBColor,
                          XL_CHART_TYPE=XL_CHART_TYPE, XL_LABEL_POSITION=XL_LABEL_POSITION,
                          XL_LEGEND_POSITION=XL_LEGEND_POSITION, XL_MARKER_STYLE=XL_MARKER_STYLE,
                          Inches=Inches, Pt=Pt)


# ===========================================================================
# CHROME: HEADER / FOOTER / KPI STRIP
# ===========================================================================

def _slide_header(slide, title: str, data_base: str, *, Inches, Pt, RGBColor) -> None:
    """Faixa superior: título (esq) + chip data-base (dir)."""
    # Título
    _textbox(slide, title,
             left=_MX, top=0.10, width=10.0, height=0.36,
             size=22, bold=True, color=_BLACK,
             Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    # Chip "Data-base: mar/26" — retângulo arredondado #F5F5F5
    chip_w = 2.20
    chip_h = 0.32
    chip_left = _SW - _MX - chip_w
    chip_top = 0.09
    chip = slide.shapes.add_shape(5, Inches(chip_left), Inches(chip_top),
                                   Inches(chip_w), Inches(chip_h))
    chip.fill.solid()
    chip.fill.fore_color.rgb = _rgb(_CARD_BG, RGBColor)
    chip.line.fill.background()
    _shape_text(chip, f"Data-base: {data_base}", size=9, bold=False,
                color=_GRAY, Pt=Pt, RGBColor=RGBColor, align_center=True)


def _slide_footer(slide, page_num: int, total: int, *, Inches, Pt, RGBColor) -> None:
    """Faixa inferior: fonte (esq) + paginação (dir)."""
    y = _FTR_TOP + 0.04
    h = 0.22

    _textbox(slide, "Fonte: CVM Fundos.NET | Elaboração: Toma Conta",
             left=_MX, top=y, width=8.0, height=h,
             size=8, bold=False, color=_GRAY,
             Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    _textbox(slide, f"Página {page_num} de {total}",
             left=_SW - _MX - 2.5, top=y, width=2.5, height=h,
             size=8, bold=False, color=_GRAY,
             align_right=True, Inches=Inches, Pt=Pt, RGBColor=RGBColor)


def _kpi_strip(
    slide, monitor_df: pd.DataFrame, monthly_df: pd.DataFrame,
    *, Inches, Pt, RGBColor,
) -> None:
    """KPI cards v3: fundo #F5F5F5, cantos arredondados, sem borda."""
    cards = _build_kpi_cards(monitor_df, monthly_df)
    if not cards:
        return

    total_w = _SW - _MX * 2
    gap = 0.15
    n = len(cards)
    card_w = (total_w - gap * (n - 1)) / n
    card_h = _KPI_H

    for i, (label, value) in enumerate(cards):
        left = _MX + i * (card_w + gap)
        top = _KPI_TOP

        card = slide.shapes.add_shape(
            5, Inches(left), Inches(top), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(_CARD_BG, RGBColor)
        card.line.fill.background()

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.02)

        # Label (uppercase, 9pt, cinza)
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = label.upper()
        r1.font.name = "Calibri"
        r1.font.size = Pt(9)
        r1.font.bold = False
        r1.font.color.rgb = _rgb(_GRAY, RGBColor)

        # Valor (18pt, bold, preto)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = value
        r2.font.name = "Calibri"
        r2.font.size = Pt(18)
        r2.font.bold = True
        r2.font.color.rgb = _rgb(_BLACK, RGBColor)


# ===========================================================================
# CHART BUILDERS
# ===========================================================================

def _chart_pl_stacked(
    slide, slot, df: pd.DataFrame, cats: list[str],
    *, CategoryChartData, RGBColor, XL_CHART_TYPE,
    XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION, Inches, Pt,
) -> None:
    if df.empty or not cats:
        _empty_placeholder(slide, slot, "PL por classe", "Sem dados.", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        return

    scale_div, scale_lbl = _money_scale(
        pd.concat([
            pd.to_numeric(df.get("pl_senior"), errors="coerce"),
            pd.to_numeric(df.get("pl_subordinada_mezz_ex360"), errors="coerce"),
        ], ignore_index=True)
    )
    senior_vals = [_scaled(v, scale_div) for v in df.get("pl_senior", pd.Series(dtype="float"))]
    subord_vals = [_scaled(v, scale_div) for v in df.get("pl_subordinada_mezz_ex360", pd.Series(dtype="float"))]

    cdata = CategoryChartData()
    cdata.categories = cats
    cdata.add_series("Sênior", senior_vals)
    cdata.add_series("Subordinada + Mezanino ex-360", subord_vals)

    chart = _place_chart(slide, slot, XL_CHART_TYPE.COLUMN_STACKED, cdata, Inches)
    _style_chart(chart, "PL por classe", "Competência", scale_lbl, "#,##0.0", RGBColor, Pt)
    chart.has_legend = True
    _style_legend(chart, XL_LEGEND_POSITION.BOTTOM, RGBColor, Pt)
    if chart.series:
        _series_fill(chart.series[0], _rgb(_BLACK, RGBColor))
    if len(chart.series) > 1:
        _series_fill(chart.series[1], _rgb(_ORANGE, RGBColor))
    _apply_export_labels(
        chart,
        [senior_vals, subord_vals],
        [_WHITE, _BLACK],
        chart_kind="stacked_bar",
        metric_kind="money",
        percent=False,
        decimals=0,
        RGBColor=RGBColor,
        Pt=Pt,
        positions=[
            XL_DATA_LABEL_POSITION.INSIDE_END,
            XL_DATA_LABEL_POSITION.INSIDE_END,
        ],
    )

    # Badge "PL total: R$ x mm"
    if not df.empty:
        latest = df.iloc[-1]
        total = (_num(latest.get("pl_senior")) or 0) + (_num(latest.get("pl_subordinada_mezz_ex360")) or 0)
        if total > 0:
            lft, top, wid, _ = slot
            badge = f"PL total: {_fmt_money_scaled(total, scale_div, scale_lbl)}"
            _textbox(slide, badge,
                     left=lft + wid - 2.40, top=top + 0.28, width=2.30, height=0.22,
                     size=10, bold=True, color=_BLACK,
                     Inches=Inches, Pt=Pt, RGBColor=RGBColor)


def _chart_percent_line(
    slide, slot, title: str, column: str,
    df: pd.DataFrame, cats: list[str], color: str,
    *, CategoryChartData, RGBColor, XL_CHART_TYPE, Inches, Pt,
    clip_outlier: bool = False, metric_kind: str = "general_pct",
) -> None:
    if df.empty or not cats or column not in df.columns:
        _empty_placeholder(slide, slot, title, "Sem dados.", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        return

    raw = pd.to_numeric(df[column], errors="coerce")
    values = [_pct_val(v) for v in raw]

    cdata = CategoryChartData()
    cdata.categories = cats
    cdata.add_series(title, values)

    chart = _place_chart(slide, slot, XL_CHART_TYPE.LINE_MARKERS, cdata, Inches)
    _style_chart(chart, title, "Competência", "%", "0.0%", RGBColor, Pt)
    chart.has_legend = False

    if chart.series:
        _series_line(chart.series[0], _rgb(color, RGBColor))

    # Outlier fix: corta eixo Y em 70% e adiciona anotação
    if clip_outlier and values:
        non_null = [v for v in values if v is not None]
        if non_null and min(non_null) < 0.10:  # outlier detectado
            try:
                chart.value_axis.minimum_scale = 0.70
            except Exception:
                pass
            lft, top, wid, hgt = slot
            _textbox(slide,
                     "Outlier: dado mai/25 inconsistente na CVM",
                     left=lft + wid * 0.55,
                     top=top + hgt * 0.70,
                     width=wid * 0.42, height=0.30,
                     size=9, bold=False, color=_GRAY, italic=True,
                     Inches=Inches, Pt=Pt, RGBColor=RGBColor)

    _apply_export_labels(
        chart,
        [values],
        [color],
        chart_kind="line",
        metric_kind=metric_kind,
        percent=True,
        RGBColor=RGBColor,
        Pt=Pt,
    )


def _chart_bar_money(
    slide, slot, title: str, column: str, series_name: str,
    df: pd.DataFrame, cats: list[str],
    *, CategoryChartData, RGBColor, XL_CHART_TYPE,
    XL_LABEL_POSITION, Inches, Pt,
) -> None:
    if df.empty or not cats or column not in df.columns:
        _empty_placeholder(slide, slot, title, "Sem dados.", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        return

    scale_div, scale_lbl = _money_scale(pd.to_numeric(df[column], errors="coerce"))
    values = [_scaled(v, scale_div) for v in df[column]]

    cdata = CategoryChartData()
    cdata.categories = cats
    cdata.add_series(series_name, values)

    chart = _place_chart(slide, slot, XL_CHART_TYPE.COLUMN_CLUSTERED, cdata, Inches)
    _style_chart(chart, title, "Competência", scale_lbl, "#,##0.0", RGBColor, Pt)
    chart.has_legend = False
    if chart.series:
        _series_fill(chart.series[0], _rgb(_BLACK, RGBColor))
    _apply_export_labels(
        chart,
        [values],
        [_BLACK],
        chart_kind="bar",
        metric_kind="money",
        percent=False,
        decimals=0,
        RGBColor=RGBColor,
        Pt=Pt,
        positions=[XL_LABEL_POSITION.OUTSIDE_END],
    )


def _chart_multi_line(
    slide, slot, title: str,
    df: pd.DataFrame, cats: list[str],
    series_specs: list[tuple[str, str, str]],
    y_title: str,
    *,
    CategoryChartData, RGBColor, XL_CHART_TYPE,
    XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE, Inches, Pt,
    value_is_percent: bool = True,
    metric_kind: str = "general_pct",
    legend_below: bool = False,
) -> None:
    valid = [(nm, col, clr) for nm, col, clr in series_specs
             if col in df.columns or not df.empty]
    if df.empty or not cats or not valid:
        _empty_placeholder(slide, slot, title, "Sem dados.", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        return

    cdata = CategoryChartData()
    cdata.categories = cats
    all_vals: list[list] = []
    for name, col, _ in valid:
        if col in df.columns:
            vals = [_pct_val(v) if value_is_percent else _num(v) for v in df[col]]
        else:
            vals = [None] * len(cats)
        cdata.add_series(name, vals)
        all_vals.append(vals)

    chart = _place_chart(slide, slot, XL_CHART_TYPE.LINE_MARKERS, cdata, Inches)
    fmt = "0.0%" if value_is_percent else "#,##0.0"
    _style_chart(chart, title, "Competência", y_title, fmt, RGBColor, Pt)
    chart.has_legend = len(valid) > 1
    if chart.has_legend:
        pos = XL_LEGEND_POSITION.BOTTOM
        _style_legend(chart, pos, RGBColor, Pt)

    for series, (_, _, clr) in zip(chart.series, valid):
        _series_line(series, _rgb(clr, RGBColor))

    _apply_export_labels(
        chart,
        all_vals,
        [clr for _, _, clr in valid],
        chart_kind="multi_line" if len(valid) > 1 else "line",
        metric_kind=metric_kind,
        percent=value_is_percent,
        RGBColor=RGBColor,
        Pt=Pt,
        positions=[
            XL_LABEL_POSITION.RIGHT,
            XL_LABEL_POSITION.ABOVE,
            XL_LABEL_POSITION.BELOW,
            XL_LABEL_POSITION.LEFT,
        ],
    )


def _chart_stacked_npl(
    slide, slot, df: pd.DataFrame, cats: list[str],
    *, CategoryChartData, RGBColor, XL_CHART_TYPE,
    XL_LABEL_POSITION, XL_LEGEND_POSITION, Inches, Pt,
) -> None:
    if df.empty or not cats:
        _empty_placeholder(slide, slot, "NPL ex-360 por severidade", "Sem dados.", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        return

    v1 = [_pct_val(v) for v in df.get("npl_1_90_pct", pd.Series(dtype="float"))]
    v2 = [_pct_val(v) for v in df.get("npl_91_360_pct", pd.Series(dtype="float"))]

    cdata = CategoryChartData()
    cdata.categories = cats
    cdata.add_series("NPL 1-90d", v1)
    cdata.add_series("NPL 91-360d", v2)

    chart = _place_chart(slide, slot, XL_CHART_TYPE.COLUMN_STACKED, cdata, Inches)
    _style_chart(chart, "NPL ex-360 por severidade", "Competência",
                 "% carteira ex-360", "0.0%", RGBColor, Pt)
    chart.has_legend = True
    _style_legend(chart, XL_LEGEND_POSITION.BOTTOM, RGBColor, Pt)
    if chart.series:
        _series_fill(chart.series[0], _rgb(_BLACK, RGBColor))
    if len(chart.series) > 1:
        _series_fill(chart.series[1], _rgb(_ORANGE, RGBColor))
    _apply_export_labels(
        chart,
        [v1, v2],
        [_WHITE, _BLACK],
        chart_kind="stacked_bar",
        metric_kind="npl_pct",
        percent=True,
        RGBColor=RGBColor,
        Pt=Pt,
        positions=[XL_LABEL_POSITION.CENTER, XL_LABEL_POSITION.CENTER],
    )


def _chart_cohorts(
    slide, slot, cohort_df: pd.DataFrame,
    *, CategoryChartData, RGBColor, XL_CHART_TYPE,
    XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE, Inches, Pt,
) -> None:
    if cohort_df is None or cohort_df.empty:
        _empty_placeholder(slide, slot, "Cohorts recentes", "Sem dados.", Inches=Inches, Pt=Pt, RGBColor=RGBColor)
        return

    df = cohort_df.copy()
    recent = (df[["cohort", "cohort_dt"]].drop_duplicates()
              .sort_values("cohort_dt").tail(6)["cohort"].tolist())
    df = df[df["cohort"].isin(recent)]
    months = ["M1", "M2", "M3", "M4", "M5", "M6"]
    gray_series = [_BLACK, _GRAY, _GRAY_MID, _GRAY_PALE, "383838", "585858"]

    cdata = CategoryChartData()
    cdata.categories = months
    all_vals: list[list] = []
    for coh in recent:
        g = df[df["cohort"].eq(coh)].set_index("mes_ciclo")
        vals = [_pct_val(g.loc[m, "valor_pct"]) if m in g.index else None for m in months]
        cdata.add_series(coh, vals)
        all_vals.append(vals)

    chart = _place_chart(slide, slot, XL_CHART_TYPE.LINE_MARKERS, cdata, Inches)
    _style_chart(chart, "Cohorts recentes", "Mês de maturação",
                 "% saldo a vencer 30d", "0.0%", RGBColor, Pt)
    chart.has_legend = True
    _style_legend(chart, XL_LEGEND_POSITION.BOTTOM, RGBColor, Pt)
    for series, clr in zip(chart.series, gray_series):
        _series_line(series, _rgb(clr, RGBColor))
    _apply_export_labels(
        chart,
        all_vals,
        gray_series,
        chart_kind="cohort",
        metric_kind="cohort_pct",
        percent=True,
        RGBColor=RGBColor,
        Pt=Pt,
        positions=[
            XL_LABEL_POSITION.RIGHT,
            XL_LABEL_POSITION.ABOVE,
            XL_LABEL_POSITION.BELOW,
            XL_LABEL_POSITION.LEFT,
        ],
    )


# ===========================================================================
# GRID LAYOUTS
# ===========================================================================

def _grid_2x2() -> list[tuple[float, float, float, float]]:
    w = (_SW - _MX * 2 - _GAP_X) / 2
    h = (_CTX_H - _GAP_Y) / 2
    return [
        (_MX, _CTX_TOP, w, h),
        (_MX + w + _GAP_X, _CTX_TOP, w, h),
        (_MX, _CTX_TOP + h + _GAP_Y, w, h),
        (_MX + w + _GAP_X, _CTX_TOP + h + _GAP_Y, w, h),
    ]


def _grid_1x2() -> list[tuple[float, float, float, float]]:
    w = (_SW - _MX * 2 - _GAP_X) / 2
    return [
        (_MX, _CTX_TOP, w, _CTX_H),
        (_MX + w + _GAP_X, _CTX_TOP, w, _CTX_H),
    ]


# ===========================================================================
# CHART STYLING HELPERS
# ===========================================================================

def _place_chart(slide, slot, chart_type, chart_data, Inches):
    l, t, w, h = slot
    return slide.shapes.add_chart(
        chart_type, Inches(l), Inches(t), Inches(w), Inches(h), chart_data
    ).chart


def _style_chart(chart, title, x_title, y_title, num_fmt, RGBColor, Pt) -> None:
    chart.has_title = True
    _set_text(chart.chart_title.text_frame, title, Pt(11), True, _BLACK, RGBColor)

    chart.category_axis.has_title = True
    _set_text(chart.category_axis.axis_title.text_frame, x_title, Pt(9), False, _GRAY, RGBColor)
    chart.value_axis.has_title = True
    _set_text(chart.value_axis.axis_title.text_frame, y_title, Pt(9), False, _GRAY, RGBColor)

    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.name = "Calibri"
        axis.tick_labels.font.size = Pt(9)
        axis.tick_labels.font.color.rgb = _rgb(_GRAY, RGBColor)

    chart.value_axis.tick_labels.number_format = num_fmt
    chart.value_axis.has_major_gridlines = True
    try:
        chart.value_axis.major_gridlines.format.line.color.rgb = _rgb(_GRAY_LIGHT, RGBColor)
        chart.value_axis.major_gridlines.format.line.width = 6350  # 0.5pt
    except Exception:
        pass

    # Remove plot area border
    try:
        chart.plot_area.format.line.fill.background()
    except Exception:
        pass

    # Remove vertical gridlines
    try:
        chart.category_axis.has_major_gridlines = False
    except Exception:
        pass


def _style_legend(chart, position, RGBColor, Pt) -> None:
    if chart.legend is None:
        return
    chart.legend.position = position
    chart.legend.include_in_layout = True
    chart.legend.font.name = "Calibri"
    chart.legend.font.size = Pt(9)
    chart.legend.font.color.rgb = _rgb(_GRAY, RGBColor)


def _series_fill(series, color) -> None:
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.color.rgb = color


def _series_line(series, color) -> None:
    series.format.line.color.rgb = color
    series.format.line.width = 15875  # ~1.25pt
    try:
        from pptx.enum.chart import XL_MARKER_STYLE
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = 4
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        series.marker.format.line.color.rgb = color
    except Exception:
        pass


def _apply_export_labels(
    chart,
    series_values: list[list],
    colors: list[str],
    *,
    chart_kind: str,
    metric_kind: str,
    percent: bool,
    RGBColor,
    Pt,
    decimals: int | None = None,
    positions: list | None = None,
) -> None:
    """Apply shared label policy to a python-pptx chart."""

    policy = choose_export_label_policy(
        series_values,
        chart_kind=chart_kind,
        metric_kind=metric_kind,
    )
    if policy.mode == "none":
        return

    for series_idx, point_indices in enumerate(policy.indices_by_series):
        if series_idx >= len(chart.series):
            continue
        series = chart.series[series_idx]
        values = series_values[series_idx]
        for point_idx in point_indices:
            if point_idx >= len(values):
                continue
            try:
                point = series.points[point_idx]
                label = point.data_label
                if positions:
                    label.position = positions[series_idx % len(positions)]
                label.has_text_frame = True
                text_frame = label.text_frame
                text_frame.clear()
                run = text_frame.paragraphs[0].add_run()
                run.text = format_export_label(
                    values[point_idx],
                    metric_kind=metric_kind,
                    percent_value=percent,
                    decimals=decimals,
                )
                run.font.name = "Calibri"
                run.font.size = Pt(policy.font_size_pt)
                run.font.bold = True
                run.font.color.rgb = _rgb(colors[series_idx % len(colors)], RGBColor)
            except Exception:
                continue


def _empty_placeholder(slide, slot, title, msg, *, Inches, Pt, RGBColor) -> None:
    l, t, w, h = slot
    _textbox(slide, title, left=l, top=t + 0.04, width=w, height=0.24,
             size=11, bold=True, color=_BLACK, Inches=Inches, Pt=Pt, RGBColor=RGBColor)
    _textbox(slide, msg, left=l, top=t + h / 2 - 0.1, width=w, height=0.28,
             size=10, bold=False, color=_GRAY, Inches=Inches, Pt=Pt, RGBColor=RGBColor)


# ===========================================================================
# TEXT HELPERS
# ===========================================================================

def _textbox(
    slide, text: str, *, left, top, width, height, size, bold, color,
    italic: bool = False, align_right: bool = False,
    Inches, Pt, RGBColor,
) -> None:
    from pptx.enum.text import PP_ALIGN
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = False
    tf.clear()
    p = tf.paragraphs[0]
    if align_right:
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color, RGBColor)


def _shape_text(shape, text: str, *, size, bold, color, Pt, RGBColor,
                align_center: bool = False) -> None:
    from pptx.enum.text import PP_ALIGN
    tf = shape.text_frame
    tf.word_wrap = False
    tf.clear()
    p = tf.paragraphs[0]
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color, RGBColor)


def _set_text(text_frame, text, size, bold, color, RGBColor) -> None:
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = _rgb(color, RGBColor)


# ===========================================================================
# DATA HELPERS
# ===========================================================================

def _last_12m(df: pd.DataFrame) -> pd.DataFrame:
    """Filtra para os últimos 12 meses da série."""
    if not isinstance(df, pd.DataFrame) or df.empty:
        return df
    out = df.copy()
    out["_dt"] = pd.to_datetime(out.get("competencia_dt", out.get("competencia")), errors="coerce")
    out = out.sort_values("_dt").dropna(subset=["_dt"])
    if len(out) > 12:
        out = out.tail(12)
    return out.drop(columns=["_dt"], errors="ignore").reset_index(drop=True)


def _chart_monthly(df: pd.DataFrame) -> pd.DataFrame:
    if not isinstance(df, pd.DataFrame) or df.empty:
        return pd.DataFrame()
    out = df.copy()
    out["competencia_dt"] = pd.to_datetime(
        out.get("competencia_dt", out.get("competencia")), errors="coerce"
    )
    return out.sort_values("competencia_dt").reset_index(drop=True)


def _category_labels(df: pd.DataFrame) -> list[str]:
    if df.empty:
        return []
    labels = []
    for _, row in df.iterrows():
        ts = pd.to_datetime(row.get("competencia_dt"), errors="coerce")
        if pd.notna(ts):
            labels.append(f"{_PT_MONTH[int(ts.month)]}/{str(int(ts.year))[-2:]}")
        else:
            labels.append(str(row.get("competencia") or ""))
    return labels


def _detect_data_base(df: pd.DataFrame) -> str:
    if not isinstance(df, pd.DataFrame) or df.empty:
        return "mar/26"
    out = _chart_monthly(df)
    if out.empty:
        return "mar/26"
    ts = pd.to_datetime(out.iloc[-1].get("competencia_dt"), errors="coerce")
    if pd.notna(ts):
        return f"{_PT_MONTH[int(ts.month)]}/{str(ts.year)[-2:]}"
    return "mar/26"


def _build_kpi_cards(
    monitor_df: pd.DataFrame, monthly_df: pd.DataFrame
) -> list[tuple[str, str]]:
    """Retorna lista (label, value) para o KPI strip."""
    row = None
    if isinstance(monitor_df, pd.DataFrame) and not monitor_df.empty:
        df = _chart_monthly(monitor_df)
        if not df.empty:
            row = df.iloc[-1]

    if row is None and isinstance(monthly_df, pd.DataFrame) and not monthly_df.empty:
        df = _chart_monthly(monthly_df)
        if not df.empty:
            row = df.iloc[-1]

    if row is None:
        return []

    comp = _fmt_month(row.get("competencia_dt") or row.get("competencia"))
    cart_label = f"Carteira ex-360 · {comp}" if comp else "Carteira ex-360"
    return [
        (cart_label, _fmt_money(_num(row.get("carteira_ex360")))),
        ("NPL Over 1d ex-360", _fmt_pct(_num(row.get("npl_over1_ex360_pct")))),
        ("NPL Over 30d ex-360", _fmt_pct(_num(row.get("npl_over30_ex360_pct")))),
        ("NPL Over 60d ex-360", _fmt_pct(_num(row.get("npl_over60_ex360_pct")))),
        ("NPL Over 90d ex-360", _fmt_pct(_num(row.get("npl_over90_ex360_pct")))),
        ("Duration", _fmt_duration(_num(row.get("duration_months")))),
    ]


def _build_duration_frame(
    consolidated: pd.DataFrame, fund_map: dict[str, pd.DataFrame]
) -> pd.DataFrame:
    rows = []
    if consolidated is not None and not consolidated.empty:
        rows.extend(_duration_rows(consolidated, "Consolidado"))
    for cnpj, frame in fund_map.items():
        name = _fund_display_name(frame, fallback=str(cnpj))
        rows.extend(_duration_rows(frame, name))
    if not rows:
        return pd.DataFrame()
    long = pd.DataFrame(rows)
    wide = (long.pivot_table(index=["competencia_dt", "competencia"],
                             columns="serie", values="duration_months",
                             aggfunc="last")
            .reset_index())
    wide.columns.name = None
    # Rename series cols to short names
    renames = {c: _short_name(c) for c in wide.columns
               if c not in {"competencia_dt", "competencia"}}
    wide = wide.rename(columns=renames)
    return wide.sort_values("competencia_dt").reset_index(drop=True)


def _duration_rows(frame: pd.DataFrame, label: str) -> list[dict]:
    if frame is None or frame.empty:
        return []
    df = _chart_monthly(frame)
    vals = pd.to_numeric(df.get("duration_months"), errors="coerce")
    return [
        {
            "competencia_dt": row.get("competencia_dt"),
            "competencia": row.get("competencia"),
            "serie": label,
            "duration_months": vals.iloc[i],
        }
        for i, (_, row) in enumerate(df.iterrows())
    ]


def _research_roll_wide(roll_df: pd.DataFrame, *, metric_id: str) -> pd.DataFrame:
    if roll_df is None or roll_df.empty:
        return pd.DataFrame()
    mask = (
        roll_df.get("scope", pd.Series(dtype="object")).astype(str).eq("consolidado")
        & roll_df["metric_id"].eq(metric_id)
    )
    df = roll_df[mask].copy()
    if df.empty:
        return pd.DataFrame()
    df["month"] = pd.to_numeric(df.get("month"), errors="coerce")
    df = df.sort_values("month")
    wide = (df.pivot_table(index=["month", "month_label"],
                           columns="series_name", values="value_pct",
                           aggfunc="last")
            .reset_index())
    wide.columns.name = None
    return wide.rename(columns={"month_label": "competencia"}).reset_index(drop=True)


# ===========================================================================
# FORMATTING HELPERS
# ===========================================================================

def _fund_display_name(frame: pd.DataFrame, *, fallback: str) -> str:
    if frame is not None and not frame.empty and "fund_name" in frame.columns:
        raw = str(frame["fund_name"].dropna().iloc[0]) if frame["fund_name"].notna().any() else fallback
    else:
        raw = fallback
    return _short_name(raw)


def _short_name(raw: str) -> str:
    """Converte nome longo para header curto."""
    key = raw.strip().upper()
    if key in _FUND_NAMES:
        return _FUND_NAMES[key][0]
    for prefix, short, _ in _FUND_PREFIXES:
        if key.startswith(prefix):
            return short
    return raw


def _fmt_month(value) -> str:
    ts = pd.to_datetime(value, errors="coerce")
    return f"{_PT_MONTH[int(ts.month)]}/{str(int(ts.year))[-2:]}" if pd.notna(ts) else ""


def _fmt_money(v) -> str:
    if v is None:
        return "N/D"
    scale_div, scale_lbl = _money_scale(pd.Series([v]))
    scaled = v / scale_div
    dec = 0 if abs(scaled) >= 10 else 1
    unit = scale_lbl.replace("R$ ", "")
    return f"R$ {_br(scaled, dec)} {unit}".strip()


def _fmt_pct(v) -> str:
    if v is None:
        return "N/D"
    return f"{_br(v, 1)}%"


def _fmt_duration(v) -> str:
    if v is None:
        return "N/D"
    return f"{_br(v, 1)} meses"


def _fmt_money_scaled(total: float, div: float, lbl: str) -> str:
    scaled = total / div if div else total
    dec = 0 if abs(scaled) >= 10 else 1
    n = _br(scaled, dec)
    if lbl == "R$":
        return f"R$ {n}"
    return f"R$ {n} {lbl.replace('R$ ', '')}"


def _fmt_val(v, *, percent: bool, decimals: int) -> str:
    if v is None:
        return "N/D"
    if percent:
        return f"{_br(v * 100.0, decimals)}%"
    return _br(v, decimals)


def _br(v: float, decimals: int) -> str:
    return f"{v:,.{decimals}f}".replace(",", "_").replace(".", ",").replace("_", ".")


def _pct_val(v) -> float | None:
    n = _num(v)
    return None if n is None else n / 100.0


def _scaled(v, div) -> float | None:
    n = _num(v)
    return None if n is None else n / div


def _num(v) -> float | None:
    parsed = pd.to_numeric(pd.Series([v]), errors="coerce").iloc[0]
    return None if pd.isna(parsed) else float(parsed)


def _last_non_null(values: list) -> int | None:
    for i in range(len(values) - 1, -1, -1):
        if values[i] is not None and not (isinstance(values[i], float) and pd.isna(values[i])):
            return i
    return None


def _money_scale(values: pd.Series) -> tuple[float, str]:
    mx = pd.to_numeric(values, errors="coerce").abs().max()
    if pd.isna(mx):
        mx = 0.0
    if mx >= 1_000_000_000_000:
        return 1_000_000_000.0, "R$ bi"
    if mx >= 1_000_000:
        return 1_000_000.0, "R$ mm"
    if mx >= 1_000:
        return 1_000.0, "R$ mil"
    return 1.0, "R$"


def _rgb(hex_color: str, RGBColor):
    c = str(hex_color).strip().lstrip("#")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
