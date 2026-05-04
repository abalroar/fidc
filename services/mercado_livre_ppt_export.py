from __future__ import annotations

from io import BytesIO
from typing import Any

import pandas as pd

from services.export_chart_labels import choose_export_label_policy, format_export_label
from services.mercado_livre_dashboard import PT_MONTH_ABBR


MELI_BLACK = "000000"
MELI_ORANGE = "E47811"
MELI_DARK_GRAY = "3F3F3F"
MELI_MEDIUM_GRAY = "8C8C8C"


def build_pptx_export_bytes(outputs: Any) -> bytes:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - environment guard
        raise RuntimeError("Dependência python-pptx não instalada.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    scopes: list[tuple[str, pd.DataFrame]] = [("Carteira consolidada", outputs.consolidated_monthly)]
    for cnpj, frame in outputs.fund_monthly.items():
        fund_name = str(frame["fund_name"].iloc[0]) if isinstance(frame, pd.DataFrame) and not frame.empty and "fund_name" in frame.columns else str(cnpj)
        scopes.append((fund_name, frame))

    for scope_name, monthly_df in scopes:
        _add_scope_grid_slide(
            prs=prs,
            layout=blank_layout,
            scope_name=scope_name,
            monthly_df=monthly_df,
            CategoryChartData=CategoryChartData,
            RGBColor=RGBColor,
            XL_CHART_TYPE=XL_CHART_TYPE,
            XL_DATA_LABEL_POSITION=XL_DATA_LABEL_POSITION,
            XL_LEGEND_POSITION=XL_LEGEND_POSITION,
            XL_MARKER_STYLE=XL_MARKER_STYLE,
            Inches=Inches,
            Pt=Pt,
        )

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _add_scope_grid_slide(
    *,
    prs,
    layout,
    scope_name: str,
    monthly_df: pd.DataFrame,
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    Inches,
    Pt,
) -> None:
    df = _chart_monthly(monthly_df)
    categories = _category_labels(df)
    slide = prs.slides.add_slide(layout)
    _style_slide_background(slide, RGBColor)
    _add_scope_header(slide, scope_name, RGBColor, Inches, Pt)
    slots = _grid_2x2_slots()
    if df.empty or not categories:
        for slot, title in zip(
            slots,
            [
                "PL por classe",
                "Subordinação total ex-360",
                "NPL Over 90d ex-360",
                "Cobertura PDD ex-360",
            ],
            strict=False,
        ):
            _add_empty_chart_placeholder(slide, slot, title, "Sem dados para o período selecionado.", RGBColor, Inches, Pt)
        return

    _add_pl_stack_chart(
        slide=slide,
        slot=slots[0],
        df=df,
        categories=categories,
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_DATA_LABEL_POSITION=XL_DATA_LABEL_POSITION,
        XL_LEGEND_POSITION=XL_LEGEND_POSITION,
        Inches=Inches,
        Pt=Pt,
    )
    _add_percent_line_chart(
        slide=slide,
        slot=slots[1],
        title="% Subordinação Total ex-360",
        series_name="% Subordinação Total ex-360",
        column="subordinacao_total_ex360_pct",
        df=df,
        categories=categories,
        color=MELI_DARK_GRAY,
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_DATA_LABEL_POSITION=XL_DATA_LABEL_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    _add_percent_line_chart(
        slide=slide,
        slot=slots[2],
        title="NPL Over 90d ex-360 / Carteira",
        series_name="NPL Over 90d ex-360 / Carteira",
        column="npl_over90_ex360_pct",
        df=df,
        categories=categories,
        color=MELI_BLACK,
        metric_kind="npl_pct",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_DATA_LABEL_POSITION=XL_DATA_LABEL_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )
    _add_percent_line_chart(
        slide=slide,
        slot=slots[3],
        title="Cobertura PDD / NPL Over 90d ex-360",
        series_name="PDD Ex / NPL Over 90d ex-360",
        column="pdd_npl_over90_ex360_pct",
        df=df,
        categories=categories,
        color=MELI_ORANGE,
        metric_kind="coverage_pct",
        CategoryChartData=CategoryChartData,
        RGBColor=RGBColor,
        XL_CHART_TYPE=XL_CHART_TYPE,
        XL_DATA_LABEL_POSITION=XL_DATA_LABEL_POSITION,
        XL_MARKER_STYLE=XL_MARKER_STYLE,
        Inches=Inches,
        Pt=Pt,
    )


def _grid_2x2_slots() -> list[tuple[float, float, float, float]]:
    margin_x = 0.48
    top = 0.72
    gap_x = 0.26
    gap_y = 0.28
    slide_width = 13.333
    slide_height = 7.5
    bottom = 0.24
    width = (slide_width - margin_x * 2 - gap_x) / 2
    height = (slide_height - top - bottom - gap_y) / 2
    return [
        (margin_x, top, width, height),
        (margin_x + width + gap_x, top, width, height),
        (margin_x, top + height + gap_y, width, height),
        (margin_x + width + gap_x, top + height + gap_y, width, height),
    ]


def _add_scope_header(slide, scope_name: str, RGBColor, Inches, Pt) -> None:  # noqa: ANN001
    _add_textbox(slide, scope_name, 0.48, 0.18, 12.35, 0.32, Pt(16), bold=True, RGBColor=RGBColor)


def _add_pl_stack_chart(
    *,
    slide,
    slot: tuple[float, float, float, float],
    df: pd.DataFrame,
    categories: list[str],
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_LEGEND_POSITION,
    Inches,
    Pt,
) -> None:
    scale_divisor, scale_label = _money_scale(
        pd.concat(
            [
                pd.to_numeric(df.get("pl_senior"), errors="coerce"),
                pd.to_numeric(df.get("pl_subordinada_mezz_ex360"), errors="coerce"),
            ],
            ignore_index=True,
        )
    )
    chart_data = CategoryChartData()
    chart_data.categories = categories
    senior_values = [_scaled_value(value, scale_divisor) for value in df.get("pl_senior", pd.Series(dtype="float"))]
    subord_values = [
        _scaled_value(value, scale_divisor)
        for value in df.get("pl_subordinada_mezz_ex360", pd.Series(dtype="float"))
    ]
    chart_data.add_series("Sênior", senior_values)
    chart_data.add_series(
        "Subordinada + Mezanino ex-360",
        subord_values,
    )
    left, top, width, height = slot
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    _style_chart_common(
        chart,
        title="PL por classe",
        x_axis_title="Competência",
        y_axis_title=scale_label,
        number_format="#,##0.0",
        RGBColor=RGBColor,
        Pt=Pt,
    )
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    _style_legend(chart, RGBColor=RGBColor, Pt=Pt)
    if chart.series:
        _set_series_fill(chart.series[0], _rgb(MELI_BLACK, RGBColor))
    if len(chart.series) > 1:
        _set_series_fill(chart.series[1], _rgb(MELI_ORANGE, RGBColor))
    _apply_ppt_export_labels(
        chart,
        [senior_values, subord_values],
        ["FFFFFF", MELI_BLACK],
        chart_kind="stacked_bar",
        metric_kind="money",
        percent=False,
        positions=[XL_DATA_LABEL_POSITION.INSIDE_END, XL_DATA_LABEL_POSITION.INSIDE_END],
        RGBColor=RGBColor,
        Pt=Pt,
    )
    _add_latest_pl_total_badge(slide, slot, df, scale_divisor, scale_label, RGBColor, Inches, Pt)


def _add_percent_line_chart(
    *,
    slide,
    slot: tuple[float, float, float, float],
    title: str,
    series_name: str,
    column: str,
    df: pd.DataFrame,
    categories: list[str],
    color: str,
    CategoryChartData,
    RGBColor,
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_MARKER_STYLE,
    Inches,
    Pt,
    metric_kind: str = "general_pct",
) -> None:
    values = [_percent_value(value) for value in df.get(column, pd.Series(dtype="float"))]
    if not any(value is not None for value in values):
        _add_empty_chart_placeholder(slide, slot, title, "Sem dados calculáveis.", RGBColor, Inches, Pt)
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    left, top, width, height = slot
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    _style_chart_common(
        chart,
        title=title,
        x_axis_title="Competência",
        y_axis_title="%",
        number_format="0.0%",
        RGBColor=RGBColor,
        Pt=Pt,
        show_gridlines=False,
    )
    chart.has_legend = False
    if chart.series:
        _set_series_line(chart.series[0], _rgb(color, RGBColor), XL_MARKER_STYLE.CIRCLE)
        _apply_ppt_export_labels(
            chart,
            [values],
            [color],
            chart_kind="line",
            metric_kind=metric_kind,
            percent=True,
            positions=[XL_DATA_LABEL_POSITION.ABOVE],
            RGBColor=RGBColor,
            Pt=Pt,
        )


def _style_slide_background(slide, RGBColor) -> None:  # noqa: ANN001
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def _add_empty_chart_placeholder(
    slide,
    slot: tuple[float, float, float, float],
    title: str,
    message: str,
    RGBColor,
    Inches,
    Pt,
) -> None:  # noqa: ANN001
    left, top, width, height = slot
    _add_textbox(slide, title, left, top + 0.04, width, 0.22, Pt(11), bold=True, RGBColor=RGBColor)
    _add_textbox(slide, message, left, top + height / 2 - 0.1, width, 0.28, Pt(10), color=MELI_MEDIUM_GRAY, RGBColor=RGBColor)


def _style_chart_common(
    chart,
    *,
    title: str,
    x_axis_title: str,
    y_axis_title: str,
    number_format: str,
    RGBColor,
    Pt,
    show_gridlines: bool = True,
) -> None:  # noqa: ANN001
    _set_chart_title(chart, title, RGBColor, Pt)
    _set_axis_title(chart.category_axis, x_axis_title, RGBColor, Pt)
    _set_axis_title(chart.value_axis, y_axis_title, RGBColor, Pt)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = _rgb(MELI_DARK_GRAY, RGBColor)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.color.rgb = _rgb(MELI_DARK_GRAY, RGBColor)
    chart.value_axis.tick_labels.number_format = number_format
    chart.value_axis.has_major_gridlines = bool(show_gridlines)
    if show_gridlines:
        try:
            chart.value_axis.major_gridlines.format.line.color.rgb = _rgb("E5E5E5", RGBColor)
            chart.value_axis.major_gridlines.format.line.width = 6350
        except AttributeError:
            pass
    try:
        chart.category_axis.format.line.color.rgb = _rgb("E5E5E5", RGBColor)
        chart.value_axis.format.line.color.rgb = _rgb("E5E5E5", RGBColor)
    except AttributeError:
        pass


def _set_chart_title(chart, title: str, RGBColor, Pt) -> None:  # noqa: ANN001
    chart.has_title = True
    _set_text_frame(chart.chart_title.text_frame, title, font_size=Pt(11), bold=True, color=MELI_BLACK, RGBColor=RGBColor)


def _set_axis_title(axis, title: str, RGBColor, Pt) -> None:  # noqa: ANN001
    axis.has_title = True
    _set_text_frame(axis.axis_title.text_frame, title, font_size=Pt(10), bold=False, color=MELI_DARK_GRAY, RGBColor=RGBColor)


def _set_text_frame(text_frame, text: str, *, font_size, bold: bool, color: str, RGBColor) -> None:  # noqa: ANN001
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = _rgb(color, RGBColor)


def _add_textbox(slide, text: str, left: float, top: float, width: float, height: float, font_size, *, bold: bool = False, color: str = MELI_BLACK, RGBColor) -> None:  # noqa: ANN001
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = _rgb(color, RGBColor)


def _style_legend(chart, *, RGBColor, Pt) -> None:  # noqa: ANN001
    if chart.legend is None:
        return
    chart.legend.font.name = "Calibri"
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = _rgb(MELI_DARK_GRAY, RGBColor)


def _add_latest_pl_total_badge(
    slide,
    slot: tuple[float, float, float, float],
    df: pd.DataFrame,
    scale_divisor: float,
    scale_label: str,
    RGBColor,
    Inches,
    Pt,
) -> None:
    if df.empty:
        return
    latest = df.iloc[-1]
    total = (_num(latest.get("pl_senior")) or 0.0) + (_num(latest.get("pl_subordinada_mezz_ex360")) or 0.0)
    if total <= 0:
        return
    left, top, width, _ = slot
    text = f"PL total: {_format_money_scaled(total, scale_divisor, scale_label)}"
    box = slide.shapes.add_textbox(Inches(left + width - 2.38), Inches(top + 0.28), Inches(2.30), Inches(0.24))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = _rgb(MELI_BLACK, RGBColor)


def _format_money_scaled(value: float, divisor: float, scale_label: str) -> str:
    scaled = value / divisor if divisor else value
    decimals = 0 if abs(scaled) >= 10 else 1
    number = _format_number_br(scaled, decimals)
    if scale_label == "R$":
        return f"R$ {number}"
    return f"R$ {number} {scale_label.replace('R$ ', '')}"


def _format_number_br(value: float, decimals: int) -> str:
    return f"{value:,.{decimals}f}".replace(",", "_").replace(".", ",").replace("_", ".")


def _rgb(hex_color: str, RGBColor):  # noqa: ANN001
    clean = str(hex_color).strip().lstrip("#")
    return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))


def _chart_monthly(monthly_df: pd.DataFrame) -> pd.DataFrame:
    if not isinstance(monthly_df, pd.DataFrame):
        return pd.DataFrame()
    df = monthly_df.copy()
    if "competencia_dt" not in df.columns:
        df["competencia_dt"] = pd.to_datetime(df.get("competencia"), errors="coerce")
    else:
        df["competencia_dt"] = pd.to_datetime(df["competencia_dt"], errors="coerce")
    return df.sort_values("competencia_dt").reset_index(drop=True)


def _category_labels(df: pd.DataFrame) -> list[str]:
    labels: list[str] = []
    for _, row in df.iterrows():
        ts = pd.to_datetime(row.get("competencia_dt"), errors="coerce")
        if pd.notna(ts):
            labels.append(f"{PT_MONTH_ABBR[int(ts.month)]}/{str(int(ts.year))[-2:]}")
        else:
            labels.append(str(row.get("competencia") or ""))
    return labels


def _money_scale(values: pd.Series) -> tuple[float, str]:
    max_value = pd.to_numeric(values, errors="coerce").abs().max()
    if pd.isna(max_value):
        max_value = 0.0
    if max_value >= 1_000_000_000_000:
        return 1_000_000_000.0, "R$ bi"
    if max_value >= 1_000_000:
        return 1_000_000.0, "R$ mm"
    if max_value >= 1_000:
        return 1_000.0, "R$ mil"
    return 1.0, "R$"


def _scaled_value(value: object, divisor: float) -> float | None:
    numeric = _num(value)
    return None if numeric is None else numeric / divisor


def _percent_value(value: object) -> float | None:
    numeric = _num(value)
    return None if numeric is None else numeric / 100.0


def _num(value: object) -> float | None:
    parsed = pd.to_numeric(pd.Series([value]), errors="coerce").iloc[0]
    if pd.isna(parsed):
        return None
    return float(parsed)


def _set_series_fill(series, color) -> None:  # noqa: ANN001
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.color.rgb = color


def _set_series_line(series, color, marker_style) -> None:  # noqa: ANN001
    series.format.line.color.rgb = color
    series.format.line.width = 19050
    series.marker.style = marker_style
    series.marker.size = 5
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = color
    series.marker.format.line.color.rgb = color


def _apply_ppt_export_labels(
    chart,
    series_values: list[list[float | None]],
    colors: list[str],
    *,
    chart_kind: str,
    metric_kind: str,
    percent: bool,
    positions: list,
    RGBColor,
    Pt,
) -> None:
    policy = choose_export_label_policy(
        series_values,
        chart_kind=chart_kind,
        metric_kind=metric_kind,
    )
    if policy.mode == "none":
        return
    for series_idx, point_indices in enumerate(policy.indices_by_series):
        if series_idx >= len(chart.series):
            continue
        values = series_values[series_idx]
        for point_idx in point_indices:
            if point_idx >= len(values):
                continue
            try:
                label = chart.series[series_idx].points[point_idx].data_label
                label.position = positions[series_idx % len(positions)]
                label.has_text_frame = True
                text_frame = label.text_frame
                text_frame.clear()
                run = text_frame.paragraphs[0].add_run()
                run.text = format_export_label(values[point_idx], metric_kind=metric_kind, percent_value=percent)
                run.font.name = "Calibri"
                run.font.size = Pt(policy.font_size_pt)
                run.font.bold = True
                run.font.color.rgb = _rgb(colors[series_idx % len(colors)], RGBColor)
            except Exception:
                continue


def _style_data_labels(
    series,
    *,
    position,
    number_format: str,
    font_color,
    fill_color,
    RGBColor,
    Pt,
    font_size=None,
) -> None:  # noqa: ANN001
    labels = series.data_labels
    labels.show_value = True
    labels.number_format = number_format
    labels.number_format_is_linked = False
    labels.position = position
    labels.font.size = font_size or Pt(8)
    labels.font.bold = True
    labels.font.color.rgb = font_color
    if fill_color is not None:
        _set_data_labels_fill(labels, fill_color)


def _set_data_labels_fill(labels, color) -> None:  # noqa: ANN001
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn

    element = labels._element
    existing = element.find(qn("c:spPr"))
    if existing is not None:
        element.remove(existing)
    sp_pr = parse_xml(
        f"""
        <c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"
                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <a:solidFill><a:srgbClr val="{_rgb_to_hex(color)}"/></a:solidFill>
          <a:ln><a:noFill/></a:ln>
        </c:spPr>
        """
    )
    insert_at = 0
    for idx, child in enumerate(element):
        if child.tag in {
            qn("c:dLblPos"),
            qn("c:showLegendKey"),
            qn("c:showVal"),
            qn("c:showCatName"),
            qn("c:showSerName"),
            qn("c:showPercent"),
            qn("c:showBubbleSize"),
            qn("c:showLeaderLines"),
        }:
            insert_at = idx
            break
    else:
        insert_at = len(element)
    element.insert(insert_at, sp_pr)


def _rgb_to_hex(color) -> str:  # noqa: ANN001
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
