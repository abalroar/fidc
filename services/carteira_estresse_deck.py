"""As duas lâminas do teste de estresse: método e nove desenquadramentos.

A primeira lâmina registra as equações. A segunda preserva a leitura de PDD
sobre inadimplência e expõe júnior, mezanino, folga e aporte para os nove fundos
que desenquadram na triagem já validada.

Fora da área visível — à direita do limite da lâmina, no mesmo arquivo — vai a
tabela de apuração: quem não declarou PDD, inadimplência ou as duas.  Ela não
compete com a leitura da lâmina, mas viaja com o arquivo e é editável como
qualquer tabela do Office.
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from services.bba_deck import (
    CONTENT_WIDTH_IN,
    GRAY_100,
    GRAY_500,
    GRAY_700,
    GRAY_900,
    MARGIN_IN,
    ORANGE,
    SLIDE_WIDTH_IN,
    WHITE,
    Deck,
    fmt_mm,
)
from services.carteira_estresse import estressar, nao_reportantes
from services.carteira_apuracao_documental import load_apuracao
from services.carteira_provisao import (
    attach_provisao,
)
from services.carteira_deck import REPLACED_SLIDE_RANGE
from services.carteira_subordinacao import resolve_portfolio, short_fund_name
from services.carteira_validacao_analistas import slide_frame
from services.deck_layout import move_slides, renumber_pages

KICKER = "CARTEIRA 101 · PROVISÃO NÃO RECONHECIDA"

#: A lâmina fecha o bloco de subordinação, logo após os índices de sub v. mínimo.
_, ULTIMO_SLIDE_ESTRUTURAL = REPLACED_SLIDE_RANGE

CHART_TOP_IN = 1.24
CHART_HEIGHT_IN = 2.24
TABLE_TOP_IN = 3.78
TRANCHE_GAP_IN = 0.05
#: Onde o número de página mora no rodapé padrão do deck.
PAGINA_X_IN = 12.25
#: FIDC · cobertura · PL · mínimo · júnior · júnior+mezz · folga · aporte · aporte/PL.
PAUTA_COLUMNS = (4.34, 0.72, 1.05, 0.86, 0.86, 1.02, 0.82, 1.20, 1.18)
TABLE_HEADER = (
    "FIDC",
    "Cob.",
    "PL R$ mm",
    "Sub Mín.",
    "Sub/PL",
    "Sub+Mez/PL",
    "Folga",
    "Aporte R$ mm",
    "Aporte/PL",
)

COLOR_ABAIXO = "C8102E"
COLOR_ACIMA = "17A398"
COLOR_GRID = "E4E6E8"

#: Vermelho de desenquadramento e verde de folga, os mesmos do resto do deck.
FILL_ABAIXO = "F7D5DA"
FILL_ATENCAO = "FCEFCF"

#: A tabela acessória mora à direita da lâmina — presente no arquivo, ausente
#: da projeção.
OFFSLIDE_LEFT_IN = SLIDE_WIDTH_IN + 0.60
APURACAO_COLUMNS = (2.60, 1.45, 0.95, 0.95, 0.95, 9.60)
APURACAO_HEADER = (
    "FIDC",
    "Seção",
    "Carteira R$ mm",
    "Inad. R$ mm",
    "PDD R$ mm",
    "Caso e apuração documental",
)


def _pp(valor: object) -> str:
    if valor is None or pd.isna(valor):
        return "—"
    return f"{float(valor):+.1f}".replace(".", ",")


def _pct(valor: object) -> str:
    if valor is None or pd.isna(valor):
        return "—"
    return f"{float(valor):.0f}%"


def _brl_mm(valor: object) -> str:
    if valor is None or pd.isna(valor):
        return "—"
    return fmt_mm(float(valor) / 1e6, 1)


def _brl_mm_fino(valor: object) -> str:
    """Como ``_brl_mm``, mas sem transformar centavos de milhão em zero.

    Na tabela de apuração convivem R$ 2.966 mm e R$ 0,0002 mm, e arredondar o
    segundo para ``0,0`` faria parecer que o fundo não declarou nada — que é
    exatamente a distinção que a tabela existe para fazer.
    """

    if valor is None or pd.isna(valor):
        return "—"
    numero = float(valor) / 1e6
    if numero == 0:
        return "0"
    if abs(numero) < 0.01:
        return "< 0,01"
    return fmt_mm(numero, 2)


def _dados(data_dir: Path) -> pd.DataFrame:
    frame = attach_provisao(resolve_portfolio(data_dir).frame, data_dir)
    return frame.assign(rotulo=frame["fundo"].map(short_fund_name))


def _linhas_pauta(mesa: pd.DataFrame) -> list[list[str]]:
    return [
        [
            str(linha.rotulo)[:40],
            _pct(linha.cobertura_pct),
            _brl_mm(linha.pl_total_cotas_brl),
            _pct(linha.referencia_pct),
            _pct(float(linha.subordinada_sobre_pl) * 100.0),
            (
                _pct(float(linha.submaismez_sobre_pl) * 100.0)
                if linha.tem_mezanino and pd.notna(linha.submaismez_sobre_pl)
                else "—"
            ),
            _pp(linha.folga_pos_pp),
            _brl_mm(linha.aporte_brl) if linha.aporte_brl and linha.aporte_brl > 0 else "—",
            _pct(float(linha.aporte_sobre_pl) * 100.0),
        ]
        for linha in mesa.itertuples()
    ]


def _pinta(mesa: pd.DataFrame) -> dict[tuple[int, int], str]:
    """Vermelho nas três medidas que registram o desenquadramento."""

    fills: dict[tuple[int, int], str] = {}
    for posicao, linha in enumerate(mesa.itertuples(), start=1):
        for coluna in (6, 7, 8):
            fills[(posicao, coluna)] = FILL_ABAIXO
    return fills


def _tabela_pauta(deck: Deck, slide, mesa: pd.DataFrame) -> None:
    linhas = [list(TABLE_HEADER)] + _linhas_pauta(mesa)
    deck.native_table(
        slide,
        linhas,
        MARGIN_IN,
        TABLE_TOP_IN,
        list(PAUTA_COLUMNS),
        aligns="lrrrrrrrr",
        size=8,
        row_height=0.23,
        header_height=0.27,
        header_fill=GRAY_100,
        header_color=GRAY_500,
        cell_fills=_pinta(mesa),
    )


def _grafico_cobertura(deck: Deck, slide, mesa: pd.DataFrame) -> None:
    """Gráfico nativo para evitar deslocamento de PNG no Office."""

    chart_left = 1.28
    chart_right = 12.58
    chart_top = 1.43
    baseline = 3.05
    plot_height = baseline - chart_top
    for value in (0, 50, 100, 150, 200):
        y = baseline - plot_height * value / 200.0
        if value == 100:
            dash_width = 0.12
            x = chart_left
            while x < chart_right:
                deck.rule(slide, x, y, min(dash_width, chart_right - x), color=GRAY_900, height=0.012)
                x += 0.20
        else:
            deck.rule(slide, chart_left, y, chart_right - chart_left, color=COLOR_GRID, height=0.008)
        deck.text(
            slide,
            "100%" if value == 100 else str(value),
            MARGIN_IN,
            y - 0.08,
            0.52,
            0.16,
            size=6.5,
            color=GRAY_700,
            align=PP_ALIGN.RIGHT,
        )

    count = len(mesa)
    slot = (chart_right - chart_left) / max(count, 1)
    width = min(0.78, slot * 0.64)
    for index, linha in enumerate(mesa.itertuples()):
        value = max(0.0, min(float(linha.cobertura_pct), 200.0))
        height = plot_height * value / 200.0
        x = chart_left + index * slot + (slot - width) / 2
        deck.block(slide, x, baseline - height, width, max(height, 0.012), COLOR_ABAIXO)
        deck.text(
            slide,
            _pct(linha.cobertura_pct),
            x - 0.05,
            max(chart_top, baseline - height - 0.17),
            width + 0.10,
            0.15,
            size=6.2,
            color=GRAY_900,
            align=PP_ALIGN.CENTER,
        )
        deck.text(
            slide,
            str(linha.rotulo)[:20],
            x - 0.18,
            baseline + 0.05,
            width + 0.36,
            0.27,
            size=5.8,
            color=GRAY_700,
            align=PP_ALIGN.CENTER,
        )

    deck.block(slide, 10.42, 1.30, 0.11, 0.08, COLOR_ABAIXO)
    deck.text(slide, "Abaixo de 100%", 10.60, 1.27, 0.88, 0.15, size=6.5, color=GRAY_900)
    deck.block(slide, 11.68, 1.30, 0.11, 0.08, COLOR_ACIMA)
    deck.text(slide, "100% ou mais", 11.86, 1.27, 0.82, 0.15, size=6.5, color=GRAY_900)


def _tabela_apuracao(deck: Deck, slide, pendentes: pd.DataFrame, data_dir: Path) -> None:
    """A lista de quem não reportou, fora da área projetada.

    A última coluna carrega, além do caso, o que a varredura dos documentos do
    próprio fundo encontrou — e o que não encontrou.
    """

    diagnosticos = load_apuracao(data_dir).set_index("cnpj")

    deck.text(
        slide,
        "Fora da lâmina — apuração: fundos sem PDD e/ou sem inadimplência declarada",
        OFFSLIDE_LEFT_IN,
        0.30,
        sum(APURACAO_COLUMNS),
        0.26,
        size=11,
        color=ORANGE,
        bold=True,
    )
    linhas = [list(APURACAO_HEADER)] + [
        [
            str(linha.rotulo)[:44],
            str(linha.categoria_estrutural),
            _brl_mm_fino(linha.carteira_dc),
            _brl_mm_fino(linha.dc_inadimplentes),
            _brl_mm_fino(linha.pdd_brl),
            _caso_com_apuracao(linha, diagnosticos),
        ]
        for linha in pendentes.itertuples()
    ]
    deck.native_table(
        slide,
        linhas,
        OFFSLIDE_LEFT_IN,
        0.62,
        list(APURACAO_COLUMNS),
        aligns="llrrrl",
        size=7,
        row_height=0.19,
        header_height=0.24,
        header_fill=GRAY_100,
        header_color=GRAY_500,
    )


def _caso_com_apuracao(linha, diagnosticos: pd.DataFrame) -> str:
    caso = str(linha.caso).replace(" — apurar", "")
    if linha.cnpj not in diagnosticos.index:
        return caso
    registro = diagnosticos.loc[linha.cnpj]
    partes = [caso, str(registro.get("diagnostico", "") or "")]
    falta = str(registro.get("lacunas", "") or "")
    if falta:
        partes.append(f"falta: {falta}")
    return " · ".join(p for p in partes if p)


def _caixa_de_pagina(slide):
    """A caixinha do número de página, achada pela geometria do rodapé."""

    for forma in slide.shapes:
        if not forma.has_text_frame:
            continue
        texto = forma.text_frame.text.strip()
        if (
            texto.isdigit()
            and forma.top > Inches(6.8)
            and abs(forma.left - Inches(PAGINA_X_IN)) < Inches(0.3)
        ):
            return forma, int(texto)
    return None, None


def _mover_para(presentation, depois_do_slide: int, quantidade: int = 1) -> None:
    """Leva as últimas lâminas para logo depois do ponto informado.

    A lâmina é **acrescentada** e só então reposicionada, porque o
    ``next_partname`` do python-pptx devolve nomes de parte já ocupados quando a
    numeração deixa de ser contígua — inserir no meio abriria esse buraco.  Aqui
    o que muda é apenas a ordem dos ``sldId`` no XML: as partes ficam onde estão,
    e o pacote continua íntegro.

    Depois da mudança de ordem, os números de página de tudo que vem adiante
    andam um.  Eles não são campos do PowerPoint, e sim caixas de texto escritas
    na montagem, então quem insere no meio tem de corrigi-los.
    """

    total = len(presentation.slides._sldIdLst)
    if total <= depois_do_slide or quantidade <= 0:
        return
    primeiro = total - quantidade + 1
    move_slides(presentation, primeiro, total, depois_do_slide + 1)
    renumber_pages(presentation)


def _slide_metodologia(deck: Deck):
    slide = deck.slide("Stress Test | Metodologia")
    deck.block(slide, MARGIN_IN, 1.42, 5.92, 2.05, GRAY_100)
    deck.text(slide, "PERDA RECONHECIDA", 0.78, 1.66, 5.35, 0.22, size=10, color=ORANGE, bold=True)
    deck.text(slide, "Δ = máx(Inadimplência − PDD, 0)", 0.78, 2.05, 5.35, 0.38, size=20, color=GRAY_900, bold=True)
    deck.text(
        slide,
        "A premissa aloca a PDD integralmente aos créditos inadimplentes. O Δ reduz a cota subordinada e o total de cotas.",
        0.78,
        2.63,
        5.25,
        0.56,
        size=11,
        color=GRAY_700,
    )

    deck.block(slide, 6.67, 1.42, 6.13, 2.05, GRAY_100)
    deck.text(slide, "ÍNDICE PÓS-ESTRESSE", 6.92, 1.66, 5.45, 0.22, size=10, color=ORANGE, bold=True)
    deck.text(slide, "Sub pós = (Sub + Mez − Δ) ÷ (Total − Δ)", 6.92, 2.05, 5.55, 0.38, size=18, color=GRAY_900, bold=True)
    deck.text(
        slide,
        "A parcela mezanino entra quando existe na Tabela X.2. O denominador é o total de cotas do Informe Mensal.",
        6.92,
        2.63,
        5.45,
        0.56,
        size=11,
        color=GRAY_700,
    )

    deck.block(slide, MARGIN_IN, 3.74, CONTENT_WIDTH_IN, 2.56, WHITE)
    deck.rule(slide, MARGIN_IN, 3.74, CONTENT_WIDTH_IN, color=ORANGE, height=0.03)
    deck.text(slide, "REENQUADRAMENTO", 0.78, 4.05, 3.1, 0.22, size=10, color=ORANGE, bold=True)
    deck.text(slide, "Folga = Sub pós − Sub Mínima", 0.78, 4.43, 5.2, 0.36, size=18, color=GRAY_900, bold=True)
    deck.text(
        slide,
        "Aporte = máx[0; (m × (Total − Δ) − (Sub + Mez − Δ)) ÷ (1 − m)]",
        0.78,
        5.06,
        11.8,
        0.38,
        size=17,
        color=GRAY_900,
        bold=True,
    )
    deck.text(
        slide,
        "m representa a Sub Mínima usada no stress. O aporte zera a folga negativa e entra simultaneamente no numerador e no denominador.",
        0.78,
        5.63,
        11.8,
        0.42,
        size=10,
        color=GRAY_700,
    )
    deck.footer(
        slide,
        "Metodologia da Carteira 101. Valores monetários em R$; índices em percentual do total de cotas.",
    )
    return slide


def append_stress_slide(presentation, data_dir: Path):
    """Acrescenta as duas lâminas ao fim e reposiciona o bloco por ``sldId``."""

    dados = estressar(_dados(data_dir))
    mesa = slide_frame(data_dir).assign(
        rotulo=lambda frame: frame["fundo"].map(short_fund_name)
    )
    pendentes = nao_reportantes(dados)
    if mesa.empty:
        return presentation

    deck = Deck(KICKER, presentation)
    _slide_metodologia(deck)
    slide = deck.slide("Stress Test | Nove FIDCs c/ cobertura < 100%")

    _grafico_cobertura(deck, slide, mesa)

    # O racional em uma linha, e só ele.
    deck.text(
        slide,
        "Δ = Inadimplência − PDD (o que falta provisionar). Reconhecido, sai da "
        "subordinada e do total de cotas: Sub pós = (Sub + Mez − Δ) ÷ (Total − Δ).",
        MARGIN_IN,
        CHART_TOP_IN + 2.10 + 0.04,
        CONTENT_WIDTH_IN,
        0.22,
        size=9,
        color=GRAY_900,
    )

    _tabela_pauta(deck, slide, mesa)

    deck.text(
        slide,
        "A tabela reúne os casos com folga negativa após os filtros de materialidade da triagem. Sub/PL mostra a cota júnior; Sub+Mez/PL aparece quando há mezanino.",
        MARGIN_IN,
        6.28,
        CONTENT_WIDTH_IN,
        0.22,
        size=8,
        color=GRAY_700,
    )

    _tabela_apuracao(deck, slide, pendentes, data_dir)
    deck.footer(
        slide,
        "CVM, Informe Mensal FIDC — Tabela I (competência mais recente de cada fundo) "
        "e regulamentos na FundosNet/B3. A carteira reportada já é líquida de PDD; a "
        "provisão é tratada como integralmente alocada aos créditos inadimplentes.",
    )
    _mover_para(presentation, ULTIMO_SLIDE_ESTRUTURAL, quantidade=2)
    return presentation


__all__ = ["append_stress_slide"]
