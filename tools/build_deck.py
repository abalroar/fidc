# -*- coding: utf-8 -*-
"""Deck para comitê de crédito, gerado a partir de data/solfacil_claude/*.csv.

Prioriza gráfico e diagrama sobre texto. Todo elemento é shape, tabela ou gráfico
nativo do Office, alinhado a uma grade explícita em EMU. Sem imagens e sem ícones.
"""
import csv, os, sys
from datetime import date
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
DATA = os.path.join(ROOT, "data", "solfacil_claude")
OUTDIR = os.path.join(ROOT, "outputs", "solfacil")
sys.path.insert(0, HERE)
from deck_base import *          # noqa: F401,F403
import deck_base as B
import solfacil_comite as K
import solfacil_estrutura as E
import solfacil_mercado as M
import solfacil_funding as FU

prs = Presentation()
prs.slide_width, prs.slide_height = SL_W, SL_H
VAZIO = prs.slide_layouts[6]
FONTE_GERAL = ("Fontes: escrituras, prospectos, lâminas, comunicados, anúncios de oferta e Termos de Securitização "
               "das emissões; relatório do agente fiduciário da debênture; análise de crédito de 21/08/2026. "
               "Data-base: FIDCs em 31/07/2026; CRIs na última competência por operação.")


def L(nome):
    return ler(DATA, nome)


def grafico(sl, tipo, x, y, w, h, categorias, series, cores, **kw):
    d = CategoryChartData()
    d.categories = categorias
    for nome, vals in series:
        d.add_series(nome, tuple(vals))
    gf = sl.shapes.add_chart(tipo, x, y, w, h, d)
    estilizar_gr(gf.chart, cores, **kw)
    return gf.chart


# ============================================================ 1. Capa
sl = prs.slides.add_slide(VAZIO)
txt(sl, MARGEM, Emu(960120), col_w(8), Emu(228600), "SOLFÁCIL | CRÉDITO ESTRUTURADO",
    tam=10.5, negrito=True, cor=CINZA_MED)
txt(sl, MARGEM, Emu(1280160), col_w(9), Emu(731520),
    "Sete warehouses, seis take-outs, uma debênture", tam=33, negrito=True, cor=PRETO)
txt(sl, MARGEM, Emu(2057400), col_w(8), Emu(457200),
    "Preço por tranche, elegibilidade, cascata, extração de subordinada e descasamento de prazo.\n"
    "Cada número resolve para uma linha de CSV com a fonte documental identificada.",
    tam=13, cor=CINZA_ESC, espaco=3)

tabela(sl, MARGEM, Emu(2880360), col_w(7), Emu(1691640),
       ["O universo", "Nº", "Como fecha"],
       [["Veículos", "14", "7 FIDCs de warehouse, 6 operações de CRI e 1 debênture de 2022"],
        ["Tranches precificadas", "70", "34 cotas de FIDC, 34 séries de CRI e 2 séries de debênture"],
        ["Emitido no programa", "R$ 8,30 bi", "R$ 4.573,9 mi em FIDC, R$ 3.670,6 mi em CRI e R$ 60,0 mi em debênture"],
        ["Retido pelo originador", "R$ 107,0 mi", "Séries subordinadas privadas dos seis CRIs, de 2,49% a 3,00% de cada"]],
       larguras=[0.28, 0.14, 0.58], tam=10.5, altura_linha=Emu(365760), alinh_centro=(1,))

txt(sl, col_x(8), Emu(2880360), col_w(4), Emu(320040), "As três perguntas do comitê", tam=13, negrito=True, cor=PRETO)
for i, (n, t) in enumerate([
        ("1", "A subordinada da originadora pode virar caixa antes do fim? Em que condições?"),
        ("2", "O prazo do ativo cabe no prazo do passivo? Que prazo aprovar na sênior?"),
        ("3", "Quanto se pagou pelos créditos cedidos - e o que disso é verificável?")]):
    y = Emu(3291840 + i * 411480)
    bloco(sl, col_x(8), y, Emu(274320), Emu(274320), n, LARANJA, PRETO, tam=11)
    txt(sl, Emu(col_x(8).emu + 365760), y, Emu(col_w(4).emu - 365760), Emu(365760), t,
        tam=10.5, cor=CINZA_ESC, anchor=MSO_ANCHOR.MIDDLE)
rodape(sl, FONTE_GERAL)

# ============================================================ 2. Mapa do programa
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Estrutura", "Como o dinheiro entra e sai do programa",
          "O FIDC financia a carteira enquanto ela nasce; o CRI compra um pool fechado e alonga o passivo; a debênture antecede os dois.")
etapas = [("ORIGINAÇÃO", "Solfácil e cerca de\n4 mil integradores\n\nCCB pré-fixada de PF e PJ\ncom alienação fiduciária\ndo equipamento", CINZA_ESC),
          ("WAREHOUSE", "FIDCs I a VII\n\nCompram e financiam\ndurante a originação\n\nRevolvência e\nreinvestimento", PETROLEO),
          ("TAKE-OUT", "Kanastra 1ª a 4ª\nVERT 174ª e 177ª\n\nCessão definitiva sem\ncoobrigação; novo\npatrimônio separado", LARANJA),
          ("INVESTIDORES", "Sênior a Subordinado Jr.\n\nSéries públicas ao mercado\ne série privada retida\npela originadora", CINZA_ESC)]
larg, y_cx, h_cx = col_w(3), Emu(1737360), Emu(2011680)
caixas = []
for i, (tit, corpo, cor) in enumerate(etapas):
    x = Emu(int(MARGEM.emu + i * (larg.emu + Emu(228600).emu)))
    if x.emu + larg.emu > SL_W.emu - MARGEM.emu:
        x = Emu(SL_W.emu - MARGEM.emu - larg.emu)
    cx = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_cx, larg, h_cx)
    cx.fill.solid(); cx.fill.fore_color.rgb = BRANCO
    cx.line.color.rgb = CINZA_CLARO; cx.line.width = Pt(1); sem_sombra(cx)
    caixas.append(cx)
    f = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_cx, larg, Emu(365760))
    f.fill.solid(); f.fill.fore_color.rgb = cor
    f.line.color.rgb = cor; sem_sombra(f)
    tf = f.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"{i + 1}.  {tit}"
    r.font.size = Pt(11); r.font.bold = True; r.font.name = FONTE; r.font.color.rgb = BRANCO
    txt(sl, Emu(x.emu + 114300), Emu(y_cx.emu + 502920), Emu(larg.emu - 228600), Emu(1417320),
        corpo, tam=10, cor=CINZA_ESC)
for i in range(3):
    a = caixas[i]
    s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(a.left.emu + a.width.emu + 45720),
                            Emu(y_cx.emu + int(h_cx.emu / 2) - 68580), Emu(137160), Emu(137160))
    s.fill.solid(); s.fill.fore_color.rgb = CINZA_MED; s.line.fill.background(); sem_sombra(s)
leitura(sl, Emu(3931920),
        "Três fundos - II, IV e VI - foram reutilizados em vários take-outs. O programa não cria um warehouse por emissão de CRI.")
txt(sl, MARGEM, Emu(4343400), UTIL_W, Emu(1188720),
    "Onde o comitê deve olhar: a cessão é sem coobrigação, então o risco de crédito migra integralmente do fundo para o "
    "patrimônio separado do CRI. Nenhum CRI tem garantia constituída no próprio título - a garantia real é a alienação "
    "fiduciária do equipamento, contratada uma camada abaixo, na CCB. Para o investidor do CRI, executar garantia significa "
    "retomar sistema fotovoltaico instalado em telhado de terceiro.",
    tam=11, cor=CINZA_ESC, espaco=3)
rodape(sl, "Lâminas de CRI-I, CRI-III e CRI-V, seção Garantias; Termo de Securitização da 177ª emissão, definição de Cessão. " + FONTE_GERAL)

# ============================================================ 3. Timeline consolidada (Gantt)
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Linha do tempo", "Todos os veículos do programa, do primeiro ao último",
          "Barras em escala de tempo real. O FIDC vai da constituição até a data-base; o CRI e a debênture vão da emissão até o vencimento legal.")

T0 = date(2020, 12, 1)
HOJE = date(2026, 7, 31)


def meses(d):
    return (d.year - T0.year) * 12 + (d.month - T0.month)


_GANTT = [
    ("FIDC I", date(2020, 12, 9), HOJE, "FIDC"), ("FIDC II", date(2021, 10, 7), HOJE, "FIDC"),
    ("FIDC IV", date(2022, 6, 23), HOJE, "FIDC"), ("FIDC V", date(2022, 12, 8), HOJE, "FIDC"),
    ("FIDC III", date(2023, 7, 10), HOJE, "FIDC"), ("FIDC VI", date(2024, 11, 6), HOJE, "FIDC"),
    ("FIDC VII", date(2026, 1, 13), HOJE, "FIDC"),
    ("Debênture", date(2022, 2, 18), date(2035, 8, 18), "DEB"),
    ("CRI-I  KAN 1ª", date(2024, 1, 15), date(2036, 1, 15), "CRI"),
    ("CRI-II  KAN 2ª", date(2024, 6, 25), date(2036, 6, 6), "CRI"),
    ("CRI-III  KAN 3ª", date(2025, 5, 28), date(2037, 5, 8), "CRI"),
    ("CRI-IV  KAN 4ª", date(2025, 9, 28), date(2037, 9, 22), "CRI"),
    ("CRI-V  VERT 174ª", date(2026, 5, 20), date(2038, 5, 20), "CRI"),
    ("CRI-VI  VERT 177ª", date(2026, 7, 21), date(2038, 7, 20), "CRI"),
]
_GANTT = list(reversed(_GANTT))          # o gráfico de barras desenha de baixo para cima
cats = [g[0] for g in _GANTT]
offset = [meses(g[1]) for g in _GANTT]
dur_fidc = [meses(g[2]) - meses(g[1]) if g[3] == "FIDC" else 0 for g in _GANTT]
dur_deb = [meses(g[2]) - meses(g[1]) if g[3] == "DEB" else 0 for g in _GANTT]
dur_cri = [meses(g[2]) - meses(g[1]) if g[3] == "CRI" else 0 for g in _GANTT]
gr = grafico(sl, XL_CHART_TYPE.BAR_STACKED, MARGEM, Emu(1417320), col_w(9), Emu(4023360),
             cats, [("", offset), ("FIDC - em funcionamento", dur_fidc),
                    ("Debênture - até o vencimento", dur_deb), ("CRI - até o vencimento legal", dur_cri)],
             [None, PETROLEO, CINZA_MED, LARANJA], unidade_y=None)
gr.value_axis.maximum_scale = 220
gr.value_axis.minimum_scale = 0
gr.value_axis.tick_labels.number_format = '0'
gr.value_axis.tick_labels.number_format_is_linked = False
txt(sl, MARGEM, Emu(5486400), col_w(9), Emu(228600),
    "Eixo em meses desde dezembro de 2020. 0 = dez/2020 · 60 = dez/2025 · 120 = dez/2030 · 180 = dez/2035",
    tam=9, cor=CINZA_MED)

txt(sl, col_x(9), Emu(1417320), col_w(3), Emu(320040), "O que a linha mostra", tam=13, negrito=True, cor=PRETO)
txt(sl, col_x(9), Emu(1783080), col_w(3), Emu(3383280),
    "Os warehouses vieram primeiro e continuam todos abertos: nenhum FIDC foi liquidado.\n\n"
    "Os take-outs começam em janeiro de 2024 e desde então há uma emissão de CRI a cada cinco a "
    "onze meses.\n\n"
    "As duas pontas convivem. O programa não substituiu o warehouse pelo take-out - ele empilhou "
    "os dois, e é isso que faz a exposição consolidada da originadora ser difícil de medir por fora.\n\n"
    "A debênture de 2022 é o veículo mais longo em prazo contratual e o menor em volume.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
leitura(sl, Emu(5760720),
        "Nenhum dos sete FIDCs foi encerrado. O que muda de veículo para veículo é a fase: originando, em revolvência ou em runoff pós-cessão.")
rodape(sl, "FIDC sem data de vencimento publicada no acervo: a barra vai até a data-base de 31/07/2026, não até um vencimento contratual. " + FONTE_GERAL)

# ============================================================ 4. Muro de vencimentos
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Amortizações futuras", "O muro de vencimentos concentra em 2031",
          "Soma do valor nominal das séries de CRI e de debênture por ano de vencimento legal. Os FIDCs não entram: suas cotas não têm vencimento publicado.")
c, r = L("26b_vencimentos_por_ano.csv")
anos = [x[0] for x in r]
vals = [float(x[1]) for x in r]
gr = grafico(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, MARGEM, Emu(1417320), col_w(8), Emu(3383280),
             anos, [("Vencimento legal no ano (R$ mi)", vals)], [LARANJA],
             legenda=False, unidade_y="R$ mi", rotulos=True)

txt(sl, col_x(8), Emu(1417320), col_w(4), Emu(320040), "Leitura de refinanciamento", tam=13, negrito=True, cor=PRETO)
tabela(sl, col_x(8), Emu(1783080), col_w(4), Emu(1005840),
       ["Janela", "R$ mi", "% do total"],
       [["2029 a 2031", "2.506,9", "70,7%"], ["2032 a 2035", "982,6", "27,7%"], ["2036 a 2038", "241,1", "6,8%"]],
       larguras=[0.40, 0.30, 0.30], tam=10, alinh_centro=(1, 2), altura_linha=Emu(274320))
txt(sl, col_x(8), Emu(2926080), col_w(4), Emu(1828800),
    "Setenta por cento do passivo de CRI vence em três anos, entre 2029 e 2031, com pico de "
    "R$ 1,24 bi em 2031.\n\n"
    "Isso não é um muro de refinanciamento no sentido usual: as séries amortizam mensalmente ao longo "
    "da vida e o vencimento legal é o limite, não a expectativa. Mas é o prazo até o qual o investidor "
    "fica exposto se a amortização não ocorrer no ritmo previsto - e a amortização é condicionada a "
    "'caso exista disponibilidade'.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
leitura(sl, Emu(4937760),
        "O vencimento legal é o prazo de exposição de pior caso. A duration divulgada é o cenário-base, e as lâminas só a qualificam para baixo.")
rodape(sl, "Soma por ano de vencimento das 34 séries de CRI e das 2 séries de debênture, a valor nominal de emissão. " + FONTE_GERAL)

# ============================================================ 5. De-para FIDC -> CRI
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Cessões", "Qual fundo alimentou qual operação",
          "Fluxo documentado. Linha cheia é documento primário; linha tracejada vem da análise de crédito, sem Termo de Securitização no acervo.")
FIDCS = [("FIDC I", "Nunca cedeu"), ("FIDC II", "4 CRIs"), ("FIDC III", "Nunca cedeu"),
         ("FIDC IV", "4 CRIs"), ("FIDC V", "Nunca cedeu"), ("FIDC VI", "4 CRIs"), ("FIDC VII", "1 CRI")]
CRIS = [("CRI-I", "KAN 1ª"), ("CRI-II", "KAN 2ª"), ("CRI-III", "KAN 3ª"),
        ("CRI-IV", "KAN 4ª"), ("CRI-V", "VERT 174ª"), ("CRI-VI", "VERT 177ª")]
LIG = {("FIDC II", "CRI-I"): 1, ("FIDC IV", "CRI-I"): 1, ("FIDC II", "CRI-II"): 1, ("FIDC IV", "CRI-II"): 1,
       ("FIDC II", "CRI-III"): 0, ("FIDC IV", "CRI-III"): 0, ("FIDC VI", "CRI-III"): 0,
       ("FIDC II", "CRI-IV"): 0, ("FIDC IV", "CRI-IV"): 0, ("FIDC VI", "CRI-IV"): 0,
       ("FIDC VI", "CRI-V"): 0, ("FIDC VI", "CRI-VI"): 1, ("FIDC VII", "CRI-VI"): 1}

X_F, X_C = MARGEM, col_x(8)
W_BOX, H_BOX = col_w(2), Emu(365760)
Y_F0, Y_C0 = Emu(1508760), Emu(1691640)
PASSO_F, PASSO_C = Emu(457200), Emu(548640)
pos_f, pos_c = {}, {}
for i, (nome, sub) in enumerate(FIDCS):
    y = Emu(Y_F0.emu + i * PASSO_F.emu)
    cedeu = "Nunca" not in sub
    s = bloco(sl, X_F, y, W_BOX, H_BOX, nome, PETROLEO if cedeu else CINZA_CLARO,
              BRANCO if cedeu else CINZA_ESC, tam=10.5)
    txt(sl, Emu(X_F.emu + W_BOX.emu + 45720), y, Emu(640080), H_BOX, sub, tam=8.5,
        cor=CINZA_MED, anchor=MSO_ANCHOR.MIDDLE)
    pos_f[nome] = (Emu(X_F.emu + W_BOX.emu), Emu(y.emu + int(H_BOX.emu / 2)))
for i, (vid, nome) in enumerate(CRIS):
    y = Emu(Y_C0.emu + i * PASSO_C.emu)
    bloco(sl, X_C, y, W_BOX, H_BOX, f"{vid}  {nome}", LARANJA, PRETO, tam=10.5)
    pos_c[vid] = (X_C, Emu(y.emu + int(H_BOX.emu / 2)))
for (f, cri), primario in LIG.items():
    x1, y1 = pos_f[f]; x2, y2 = pos_c[cri]
    cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(x1.emu + 640080), y1, x2, y2)
    cn.line.color.rgb = CINZA_ESC if primario else CINZA_CLARO
    cn.line.width = Pt(1.75 if primario else 1)
    if not primario:
        cn.line.dash_style = 4      # tracejado

txt(sl, col_x(3), Emu(1417320), col_w(4), Emu(228600), "CEDENTES", tam=9, negrito=True, cor=CINZA_MED, alinh=PP_ALIGN.CENTER)
txt(sl, col_x(3), Emu(4663440), col_w(4), Emu(731520),
    "Dois dos sete fundos nunca aparecem como cedentes em nenhum documento: o FIDC I e o FIDC III. "
    "O FIDC V também não - e é o único cujo mandato admite CPR-F e prazo de até 4.760 dias, acima do "
    "teto de 3.845 dias de todos os CRIs.",
    tam=10, cor=CINZA_ESC, espaco=3, alinh=PP_ALIGN.CENTER)
leitura(sl, Emu(5486400),
        "Só três das seis operações têm o cedente nomeado em documento primário: CRI-I e CRI-II pelos Prospectos, CRI-VI pelo Termo de Securitização.")
rodape(sl, "Linha cheia: Prospectos Definitivos da 1ª e 2ª emissões Kanastra e Termo de Securitização da 177ª VERT. Linha tracejada: análise de crédito de 21/08/2026; os Termos das 1ª, 3ª e 4ª Kanastra e da 174ª VERT não estão no acervo.")

# ============================================================ 6. Situação hoje
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Situação corrente", "O que continua vigente hoje",
          "Nenhum veículo foi encerrado. O que muda é a fase de cada um, e é a fase que define o risco corrente.")
c, r = L("27b_fidc_status.csv")
linhas = [[x[0].replace("FIDC-", "FIDC "), x[3], x[4], f"R$ {x[5]} mi", f"R$ {x[6]} mi"] for x in r]
tabela(sl, MARGEM, Emu(1417320), col_w(12), Emu(2103120),
       ["Fundo", "Cedeu para", "Situação hoje", "PL", "Carteira"], linhas,
       larguras=[0.11, 0.28, 0.35, 0.13, 0.13], tam=9.5, altura_linha=Emu(256032), alinh_centro=(3, 4))
c2, r2 = L("27_depara_fidc_cri.csv")
linhas2 = [[x[1], x[2][:7], x[3], x[4], x[5]] for x in r2]
txt(sl, MARGEM, Emu(3657600), col_w(12), Emu(274320), "AS SEIS OPERAÇÕES DE CRI", tam=11, negrito=True, cor=PRETO)
tabela(sl, MARGEM, Emu(3931920), col_w(12), Emu(1554480),
       ["Operação", "Emissão", "Cedentes documentados", "Cessão direta da Solfácil", "Situação"], linhas2,
       larguras=[0.14, 0.11, 0.30, 0.20, 0.25], tam=9.5, altura_linha=Emu(219456), alinh_centro=(1, 3))
leitura(sl, Emu(5623560),
        "O FIDC VI atravessa quatro operações e hoje está em runoff; o FIDC VII entrou só na 177ª e segue em revolvência obrigatória de 12 meses.")
rodape(sl, K.NOTA_DEPARA)

# ============================================================ 7. Pilha de funding
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Pilha de funding", "Quem financiou o quê, ano a ano",
          "Setenta tranches em seis anos. A virada é 2024: até 2023 o funding vem inteiro de FIDC e debênture.")
c, r = L("25_captacao_por_ano.csv")
anos = sorted({x[0] for x in r})
mapa = {(x[0], x[1]): float(x[2]) for x in r}
grafico(sl, XL_CHART_TYPE.COLUMN_STACKED, MARGEM, Emu(1417320), col_w(7), Emu(3200400), anos,
        [("FIDC (warehouse)", [mapa.get((a, "FIDC (warehouse)"), 0.0) for a in anos]),
         ("CRI (take-out)", [mapa.get((a, "CRI (take-out)"), 0.0) for a in anos]),
         ("Debênture", [mapa.get((a, "Debênture"), 0.0) for a in anos])],
        [PETROLEO, LARANJA, CINZA_MED], unidade_y="R$ mi emitidos no ano")
txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(320040), "Os quatro tipos de veículo", tam=13, negrito=True, cor=PRETO)
tabela(sl, col_x(7), Emu(1783080), col_w(5), Emu(1188720),
       ["Instrumento", "Veículos", "Tranches", "Emitido (R$ mi)"],
       [["FIDC (warehouse)", "7", "34", "4.573,9"], ["CRI (take-out)", "6", "34", "3.670,6"],
        ["Debênture", "1", "2", "60,0"]],
       larguras=[0.38, 0.16, 0.17, 0.29], tam=10, alinh_centro=(1, 2, 3), altura_linha=Emu(274320))
txt(sl, col_x(7), Emu(3108960), col_w(5), Emu(1417320),
    "O valor por ano é emissão de tranche, não saldo em aberto: os FIDCs reabrem classes ao longo do "
    "tempo, então a soma das cotas de um mesmo fundo em anos diferentes não é o tamanho dele hoje.\n\n"
    "A comparação entre instrumentos aqui é de volume, não de custo - preço só se compara dentro da "
    "mesma família de indexador.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
leitura(sl, Emu(4937760),
        "De 2024 em diante o CRI responde pela maior parte da captação de cada ano, mas nenhum warehouse foi desmontado no processo.")
rodape(sl, "O gráfico por ano exclui uma tranche de R$ 50,0 mi - a cota subordinada júnior do FIDC IV, cuja data de registro não consta do acervo. Por isso a soma das barras é R$ 8.254,5 mi contra R$ 8.304,5 mi emitidos. " + FONTE_GERAL)

# ============================================================ 8. Custo de captação
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Custo de captação", "O sênior ficou 200 pontos-base mais barato em três anos",
          "Família DI+, a única série longa e comparável sem recorrer à curva de juros. Cada ponto é a taxa contratada na emissão do veículo.")
c, r = L("22_custo_senior_timeline.csv")
di = [x for x in r if x[6] == "DI+"]
gr = grafico(sl, XL_CHART_TYPE.LINE_MARKERS, MARGEM, Emu(1417320), col_w(7), Emu(2925763),
             [f"{x[2]}\n{x[1][:7]}" for x in di],
             [("Spread da camada sênior sobre o DI", [float(x[8]) for x in di])],
             [LARANJA], legenda=False, unidade_y="% a.a. sobre o DI", rotulos=True)
for s in gr.series:
    s.format.line.width = Pt(2.5); s.smooth = False
txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(320040), "As outras famílias, sem misturar", tam=13, negrito=True, cor=PRETO)
outros = [x for x in r if x[6] in ("IPCA+", "Pré")]
tabela(sl, col_x(7), Emu(1783080), col_w(5), Emu(2011680),
       ["Veículo", "Data", "Família", "Taxa"],
       [[x[2], x[1][:7], x[6], x[7].replace(" a.a.", "").replace("IPCA + ", "")] for x in outros],
       larguras=[0.30, 0.22, 0.18, 0.30], tam=9, alinh_centro=(1, 2, 3), altura_linha=Emu(228600))
txt(sl, col_x(7), Emu(3931920), col_w(5), Emu(914400),
    "7,22% em IPCA+ e 14,81% em pré-fixado descrevem custos que só entram na mesma régua depois de "
    "convertidos pela curva DI e pela inflação implícita da data-base. Essa conversão não foi feita: falta o insumo.",
    tam=10.5, cor=CINZA_ESC, espaco=4)
leitura(sl, Emu(4937760),
        "Em julho de 2026 o take-out em CRI capta a DI+1,50% - cinquenta pontos-base abaixo do warehouse que o alimenta, o FIDC VII, a DI+2,00%.")
txt(sl, MARGEM, Emu(5349240), UTIL_W, Emu(548640),
    "Nas camadas subordinadas o movimento é o mesmo: o Mezanino caiu de DI+6,00% na 2ª emissão para DI+5,50% da 4ª em diante, e o "
    "Subordinado de DI+10,00% para DI+8,00% a partir da 174ª. Em CRI-III as duas séries pré pararam exatamente no piso da lâmina - "
    "15,50% e 16,50% - enquanto em CRI-I e CRI-V a perna indexada ao DI futuro prevaleceu.",
    tam=11, cor=CINZA_ESC, espaco=3)
rodape(sl, "A queda de spread é consistente com melhora de percepção de risco, mas o dado público não sustenta causalidade: mudaram também o indexador, a camada e o ciclo de juros. " + FONTE_GERAL)

# ============================================================ 9. Preço de aquisição
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Preço de cessão", "Os dois lados do programa precificam de formas diferentes",
          "O FIDC compra por percentual do saldo contábil. O CRI compra por valor presente descontado a uma taxa - e o que o contrato limita é a taxa, não o percentual.")
bloco(sl, MARGEM, Emu(1417320), col_w(6), Emu(365760), "FIDC — teto em percentual do saldo contábil",
      PETROLEO, BRANCO, tam=11)
tabela(sl, MARGEM, Emu(1828800), col_w(6), Emu(1417320),
       ["Fundo", "Teto de preço", "Quem atesta"],
       [["FIDC I", "100,4% do saldo", "Gestor"], ["FIDC II", "100,5% do saldo", "Gestor"],
        ["FIDC V", "100,5% do saldo", "Gestor"], ["FIDC VI", "101% do saldo", "Gestora"],
        ["FIDC VII", "104% do saldo", "Endossantes e originador"]],
       larguras=[0.24, 0.36, 0.40], tam=9.5, altura_linha=Emu(228600), alinh_centro=(1,), destaque=(4,))
bloco(sl, col_x(6), Emu(1417320), col_w(6), Emu(365760), "CRI — teto em taxa de desconto",
      LARANJA, PRETO, tam=11)
txt(sl, col_x(6), Emu(1874520), col_w(6), Emu(1188720),
    "Preço de Aquisição  =  Σᵢ  Parcela do Direito Creditórioᵢ  ÷  (1 + Taxa de Retorno)^(N/252)\n\n"
    "O valor bruto dessa fórmula é o máximo que pode ser pago. A elegibilidade exige Taxa de Retorno "
    "igual ou superior à Taxa Média Mínima de Retorno — e quanto maior a taxa exigida, menor o preço.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
tabela(sl, col_x(6), Emu(3108960), col_w(6), Emu(731520),
       ["Operação", "Taxa mínima de retorno"],
       [["CRI-I  Kanastra 1ª", "21,5% a.a."], ["CRI-II  Kanastra 2ª", "21,0% a.a."]],
       larguras=[0.55, 0.45], tam=9.5, altura_linha=Emu(228600), alinh_centro=(1,))
txt(sl, MARGEM, Emu(3474720), col_w(6), Emu(731520),
    "O teto subiu ao longo do programa: de 100,4% no fundo I para 104% no fundo VII. No FIDC VII "
    "mudou também quem atesta - deixou de ser a gestora e passou a ser o próprio originador.",
    tam=10.5, cor=CINZA_ESC, espaco=4)
bloco(sl, MARGEM, Emu(4297680), UTIL_W, Emu(548640),
      "Preço efetivamente praticado em cada cessão: não disponível. É evidenciado no respectivo Termo de Cessão, "
      "que não é público em nenhuma das seis operações de CRI.", CINZA_FUNDO, PRETO, tam=11, negrito=True,
      borda=CINZA_CLARO)
leitura(sl, Emu(5029200),
        "A única aquisição com preço observável no acervo é a da debênture: R$ 17,2 mi de valor de aquisição para R$ 17,2 mi de valor nominal — compra ao par.")
rodape(sl, "Fórmula: definição de Preço de Aquisição do Termo de Securitização da 177ª VERT. Taxas mínimas: Prospectos Definitivos da 1ª e 2ª Kanastra. Tetos dos FIDCs: análise de crédito de 21/08/2026.")

# ============================================================ 10. Elegibilidade
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Seleção do pool", "Elegibilidade, critério a critério",
          "Redação literal das lâminas. O CRI compra um recorte mais curto e mais granular do que o mandato do warehouse permite originar.")
tabela(sl, MARGEM, Emu(1417320), col_w(12), Emu(3931920),
       ["Critério", "FIDCs (faixa dos sete)", "CRI-I  jan/2024", "CRI-III  mai/2025", "CRI-V  mai/2026"],
       [["Taxa da CCB", "Fixa", "Pré-fixada", "Pré-fixada", "Pré-fixada"],
        ["Moeda do pagamento", "Reais (VI e VII)", "n/d na lâmina", "n/d na lâmina", "Reais, explícito"],
        ["Cap por devedor", "2% a 20% do patrimônio", "0,10% do Patr. Separado", "0,10% do Patr. Separado",
         "0,15% até 470.600 cotas;\n0,07% a partir de 750.000"],
        ["Cap dos 10 maiores", "1% a 10%", "1% do Patr. Separado", "Não consta", "Não consta"],
        ["WAM do pool", "2.000 a 2.400 dias", "2.000 dias", "2.000 dias", "2.000 dias, sobre o valor presente"],
        ["Prazo máx. por recebível", "até 4.760 dias (FIDC V)", "3.845 dias", "3.845 dias", "3.845 dias corridos"],
        ["Devedor inadimplente", "Vedado (II a VII)", "Vedado na Data de Oferta", "Vedado na Data de Oferta", "Vedado na Data de Oferta"],
        ["Parcela balão final", "n/d", "Não consta", "Não consta", "Vedada expressamente"],
        ["Idade máxima PF", "71 anos (VI e VII)", "71 anos", "71 anos", "71 anos"],
        ["Pessoa jurídica", "2 anos de constituição", "2 anos", "2 anos", "2 anos e enquadrada na CMN 5.118"],
        ["Carência máxima", "180 a 366 dias", "185 dias", "185 dias", "185 dias"],
        ["Valor presente máx. PF / PJ", "R$ 201-500 mil / 500-700 mil", "R$ 350 mil / R$ 600 mil",
         "R$ 350 mil / R$ 700 mil", "R$ 350 mil / R$ 700 mil"],
        ["Seasoning mínimo", "n/d", "Não exigido", "Não exigido", "Não exigido"],
        ["Quem atesta", "Gestora / administrador", "Emissora, com dados dos Cedentes",
         "Emissora, com dados dos Cedentes", "Emissora, com dados da Gestora\ndo Cedente Fundo"]],
       larguras=[0.19, 0.20, 0.20, 0.20, 0.21], tam=8.5, tam_cab=9,
       altura_cab=Emu(274320), altura_linha=Emu(228600), destaque=(2, 7, 12))
leitura(sl, Emu(5486400),
        "Nenhuma das seis operações exige safra performada, MoB mínimo ou histórico de inadimplência do lote: a seleção é documental, não de performance.")
rodape(sl, "Faixa dos FIDCs conforme análise de crédito; colunas de CRI com redação literal das lâminas de 15/01/2024, 22/04/2025 e 17/04/2026. CRI-II, CRI-IV e CRI-VI seguem o mesmo desenho.")

# ============================================================ 11. Ranking de permissividade
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Ranking de risco", "Do mandato mais restritivo ao mais permissivo",
          "Índice de 0 a 100 agregando quatro parâmetros documentados. Mede o quanto cada documento permite, não o que a carteira de fato tem.")
c, r = L("28_ranking_permissividade.csv")
iv, ii, ix = c.index("veiculo_id"), c.index("instrumento"), c.index("indice_permissividade")
validos = [x for x in r if x[ix] != "n/d"]
cats = [f"{x[iv].replace('FIDC-', 'FIDC ').replace('CRI-', 'CRI ')}" for x in validos][::-1]
vals_cri = [float(x[ix]) if x[ii] == "CRI" else 0 for x in validos][::-1]
vals_fidc = [float(x[ix]) if x[ii] == "FIDC" else 0 for x in validos][::-1]
grafico(sl, XL_CHART_TYPE.BAR_STACKED, MARGEM, Emu(1417320), col_w(6), Emu(3931920), cats,
        [("CRI (take-out)", vals_cri), ("FIDC (warehouse)", vals_fidc)],
        [LARANJA, PETROLEO], unidade_y="Índice de permissividade (0 = mais restritivo)", rotulos=True)

txt(sl, col_x(6), Emu(1417320), col_w(6), Emu(274320), "O que entra no índice", tam=13, negrito=True, cor=PRETO)
tabela(sl, col_x(6), Emu(1737360), col_w(6), Emu(1005840),
       ["Eixo", "Restritivo", "Permissivo", "Escala"],
       [["Cap por devedor", "0,07%", "20%", "Logarítmica"],
        ["WAM máximo do pool", "2.000 d", "2.400 d", "Linear"],
        ["Prazo máx. por recebível", "3.836 d", "4.760 d", "Linear"],
        ["Ticket máximo PJ", "R$ 500 mil", "R$ 700 mil", "Linear"]],
       larguras=[0.36, 0.20, 0.20, 0.24], tam=9, alinh_centro=(1, 2, 3), altura_linha=Emu(219456))
txt(sl, col_x(6), Emu(2926080), col_w(6), Emu(2377440),
    "Os extremos são inequívocos. O FIDC IV é o mandato mais permissivo do programa - 20% por devedor, "
    "sem piso de subordinação e hoje sem first loss corrente - e é também o fundo de onde saiu o maior "
    "volume de subordinada: R$ 438,9 mi.\n\n"
    "No outro extremo, os seis CRIs ocupam a metade restritiva, e o FIDC III é o único warehouse que se "
    "posiciona junto deles: 0,10% por devedor, WAM de 2.000 dias e teto de 1% para os dez maiores.\n\n"
    "Ressalva de comparabilidade: CRI-II, CRI-IV e CRI-VI têm três eixos avaliados em vez de quatro, "
    "porque o ticket máximo PJ não consta dos documentos disponíveis. Os FIDCs VI e VII também têm "
    "índice parcial - o cap por devedor não é divulgado.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
leitura(sl, Emu(5486400),
        "O índice ordena mandatos, não desempenho. Um fundo restritivo com carteira ruim continua ruim; o índice diz só o que o documento deixa entrar.")
rodape(sl, "Insumos de 03_Elegibilidade e 04_Concentracao; fórmula e ressalvas em 18_Metodologia. Inferido: agregação de parâmetros documentados, com escala logarítmica no eixo de concentração porque o cap varia 285 vezes entre os extremos.")

# ============================================================ 12. Mismatch de prazo
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Descasamento de prazo", "O ativo é longo; o passivo sênior, muito mais curto",
          "Tudo em meses, na mesma escala. E o descasamento não é uniforme: depende da operação, não do programa.")
grafico(sl, XL_CHART_TYPE.BAR_CLUSTERED, MARGEM, Emu(1417320), col_w(7), Emu(2743200),
        ["Prazo máx. do\nrecebível", "WAM contratual\ndo pool", "Duration sênior\nCRI-III",
         "Duration sênior\nCRI-I", "Duration sênior\nCRI-V"],
        [("Meses", [126.3, 65.7, 59.3, 40.9, 21.6])], [LARANJA],
        legenda=False, unidade_y="Meses", rotulos=True)
txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(274320), "Gap por operação", tam=13, negrito=True, cor=PRETO)
c, r = L("32_mismatch_prazo.csv")
tabela(sl, col_x(7), Emu(1737360), col_w(5), Emu(1005840),
       ["Operação", "Duration", "Venc. legal", "Gap ativo × duration"],
       [[x[0], f"{x[3]} m", f"{x[5]} m", f"{x[8]} m"] for x in r],
       larguras=[0.28, 0.22, 0.22, 0.28], tam=9.5, alinh_centro=(1, 2, 3), altura_linha=Emu(256032))
txt(sl, col_x(7), Emu(2926080), col_w(5), Emu(2103120),
    "A mesma camada tem duration de 1.806 dias em CRI-III e de 659 dias em CRI-V. Tratar a faixa curta "
    "de CRI-V como padrão do programa seria erro material.\n\n"
    "Quem absorve o gap é a estrutura: o recebível longo continua no patrimônio separado depois que a "
    "sênior foi resgatada. O que sustenta o resgate no prazo é a amortização do pool, não o vencimento dele.\n\n"
    "O WAM observado de cada pool não é publicado em nenhuma operação. Só existe o teto contratual.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
leitura(sl, Emu(5257800),
        "O risco que sobra para o investidor da sênior não é de prazo do ativo — é de extensão: se o pool amortizar devagar, o recebimento alonga até o vencimento legal.")
rodape(sl, "Duration aproximada informada nas lâminas, sujeita a redução por amortização extraordinária. CRI-II, CRI-IV e CRI-VI não publicam duration. Conversão a 30,4375 dias por mês.")

# ============================================================ 13. Prazo sugerido de exposição
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Prazo de aprovação", "Dimensionar pelo vencimento legal, não pela duration",
          "A duration divulgada é cenário-base e as lâminas só a qualificam para baixo. O risco de extensão está no próprio contrato.")
c, r = L("32_mismatch_prazo.csv")
tabela(sl, MARGEM, Emu(1417320), col_w(12), Emu(1188720),
       ["Operação", "Camada sênior", "Duration", "Vencimento legal", "Regime muda em", "Prazo sugerido de aprovação"],
       [[x[0], x[1], f"{x[3]} meses", f"{x[4]}  ({x[5]} m)", x[9], x[10]] for x in r],
       larguras=[0.11, 0.22, 0.12, 0.19, 0.19, 0.17], tam=9.5, altura_linha=Emu(274320), alinh_centro=(2,))

txt(sl, MARGEM, Emu(2834640), col_w(12), Emu(274320), "POR QUE NÃO APROVAR PELA DURATION", tam=11, negrito=True, cor=PRETO)
razoes = [
    ("Condicionalidade", "Toda amortização é condicionada a 'caso exista disponibilidade' e ao limite do caixa disponível — redação repetida para todas as séries."),
    ("Saldo Devedor Target", "A série amortiza até um saldo alvo, não até o cronograma. Se o pool render menos, o alvo não é atingido e o saldo permanece."),
    ("Regime sequencial no mês 48", "A partir do mês 48 o regime vira sequencial e as camadas abaixo da sênior param de receber principal — o que acelera a sênior, mas só se houver caixa."),
    ("Gatilho de cobertura", "Desenquadramento das Razões de Cobertura em 2 datas consecutivas ou 4 alternadas em 12 meses aciona Evento de Desalavancagem."),
]
for i, (tit, corpo) in enumerate(razoes):
    y = Emu(3154680 + i * 502920)
    bloco(sl, MARGEM, y, Emu(365760), Emu(365760), str(i + 1), LARANJA, PRETO, tam=11)
    txt(sl, Emu(MARGEM.emu + 457200), y, Emu(1828800), Emu(365760), tit, tam=10.5, negrito=True,
        cor=PRETO, anchor=MSO_ANCHOR.MIDDLE)
    txt(sl, Emu(MARGEM.emu + 2377440), y, Emu(UTIL_W.emu - 2377440), Emu(457200), corpo,
        tam=10, cor=CINZA_ESC, anchor=MSO_ANCHOR.MIDDLE)
leitura(sl, Emu(5303520),
        "Um limite aprovado pela duration venceria antes do papel em qualquer cenário de performance abaixo do previsto.")
rodape(sl, "Leitura de estrutura para o comitê, com base nas lâminas e nos Termos de Securitização da 2ª Kanastra e da 177ª VERT. Não constitui recomendação de investimento. " + FONTE_GERAL)


# ---------------------------------------------------------------- helper de waterfall
def waterfall(sl, x, y, w, h, rotulos_cat, valores, cor_barra, total_label="Patrimônio Separado",
              unidade="R$ mi"):
    """Waterfall nativo: coluna empilhada com série-base invisível."""
    base, val, cats = [], [], [total_label] + rotulos_cat
    total = sum(valores)
    val.append(total); base.append(0.0)
    restante = total
    for v in valores:
        restante -= v
        base.append(restante); val.append(v)
    gr = grafico(sl, XL_CHART_TYPE.COLUMN_STACKED, x, y, w, h, cats,
                 [("", base), (unidade, val)], [None, cor_barra],
                 legenda=False, unidade_y=unidade, rotulos=False)
    gr.plots[0].gap_width = 40
    pt = gr.series[1]
    pt.has_data_labels = True
    return gr


# ============================================================ 14. Waterfall CRI-II
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Ordem de pagamentos", "CRI-II: como o patrimônio separado é consumido por senioridade",
          "Cada degrau só recebe se as Razões de Cobertura de todas as camadas acima estiverem enquadradas na Data de Pagamento.")
waterfall(sl, MARGEM, Emu(1417320), col_w(7), Emu(3017520),
          ["Super Sênior", "Sênior", "Mezanino", "Subordinado", "Subordinado Jr."],
          [487.5, 135.0, 75.0, 30.0, 22.5], LARANJA)
txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(274320), "A trava de cada degrau", tam=13, negrito=True, cor=PRETO)
c, r = L("29b_razoes_de_cobertura.csv")
cri2 = [x for x in r if x[0] == "CRI-II"]
tabela(sl, col_x(7), Emu(1737360), col_w(5), Emu(1005840),
       ["Camada", "Razão de cobertura mínima"],
       [[x[1], f"{x[2]}%"] for x in cri2],
       larguras=[0.55, 0.45], tam=10, alinh_centro=(1,), altura_linha=Emu(228600))
txt(sl, col_x(7), Emu(2926080), col_w(5), Emu(2286000),
    "O regime chamado de pró-rata não paga em paralelo. A 2ª série só recebe se a cobertura Super "
    "Sênior estiver enquadrada; a 3ª exige as duas acima; a 4ª exige três; a 5ª exige as quatro.\n\n"
    "E cada uma recebe até um Saldo Devedor Target, não até o cronograma do Anexo I — o cronograma é "
    "alvo, não promessa.\n\n"
    "No fim da fila, item (bb): todo o valor remanescente do patrimônio separado vai à 5ª série a "
    "título de Prêmio Final. É a série retida pela originadora.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
leitura(sl, Emu(4663440),
        "A 5ª série tem cupom declarado de 11,93% — o menor da emissão — mas recebe todo o resíduo do patrimônio separado. A taxa declarada subestima a economia dela.")
rodape(sl, "Cláusula 6.5.1, itens (a) a (cc), do 2º Aditamento ao Termo de Securitização da 2ª emissão Kanastra, consolidado e registrado na JUCEMG em 10/06/2025.")

# ============================================================ 15. Waterfall CRI-VI
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Ordem de pagamentos", "CRI-VI: mesma espinha, com o derivativo na frente de todos",
          "A 177ª VERT introduz contratos de swap. O ajuste do derivativo é pago antes de qualquer investidor e tem reserva própria.")
waterfall(sl, MARGEM, Emu(1417320), col_w(7), Emu(3017520),
          ["Sênior A", "Sênior B", "Mezanino I", "Mezanino II", "Subordinado Jr."],
          [100.0, 450.0, 51.765, 25.882, 19.412], LARANJA)
txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(274320), "O que muda frente a CRI-II", tam=13, negrito=True, cor=PRETO)
tabela(sl, col_x(7), Emu(1737360), col_w(5), Emu(1188720),
       ["Camada", "Cobertura mínima"],
       [["Sênior (A e B)", "120,48%"], ["Mezanino I", "109,89%"], ["Mezanino II", "105,26%"],
        ["Índice de Atraso de Estoque", "máx. 15%"]],
       larguras=[0.58, 0.42], tam=10, alinh_centro=(1,), altura_linha=Emu(228600))
txt(sl, col_x(7), Emu(3108960), col_w(5), Emu(2103120),
    "Três Razões de Cobertura em vez de quatro: as duas séries sênior são pari passu e a cláusula 7.5.3 "
    "trava a proporção entre elas quando o caixa é insuficiente.\n\n"
    "O ajuste dos contratos de derivativos é o item (c) da cascata — antes da remuneração sênior. Há "
    "ainda uma Reserva de Caixa Derivativo, limitada a 1% do valor presente da parcela pós-fixada.\n\n"
    "E o diferencial de taxa do swap, positivo ou negativo, impacta a remuneração da 5ª série.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
leitura(sl, Emu(4663440),
        "Em CRI-VI a série retida pela originadora absorve o risco de base do hedge: o diferencial do swap recai sobre ela e, se a reserva faltar, há cobertura extraordinária dos titulares da 5ª série.")
rodape(sl, "Cláusulas 7.5.1 e 7.5.2 da Ordem de Alocação; cl. 7.5.3 razão entre séries sênior; cl. 15.14.6 item 7 e cl. 15.14.8 item 3 do Termo de Securitização da 177ª emissão VERT.")

# ============================================================ 16. Waterfall consolidado
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Cascata consolidada", "A pilha de absorção de perda do programa inteiro",
          "Soma das 70 tranches por senioridade. Onde está o primeiro prejuízo e quanto de proteção existe abaixo de cada camada.")
c, r = L("02_series.csv")
iv, ic, im = c.index("veiculo_id"), c.index("camada"), c.index("montante_reportado_deck_Rmi")
GRUPO = {"Super Senior": "Sênior", "Super Sênior": "Sênior", "Super Sênior A": "Sênior",
         "Super Sênior B": "Sênior", "Senior": "Sênior", "Sênior": "Sênior",
         "Sênior A": "Sênior", "Sênior B": "Sênior",
         "Mezanino": "Mezanino", "Mezanino I": "Mezanino", "Mezanino II": "Mezanino",
         "Subordinado": "Subordinado", "Subordinado Jr.": "Subordinado Jr."}
agr = {}
for x in r:
    if x[im] in ("n/d", ""):
        continue
    g = GRUPO.get(x[ic])
    if g:
        agr[g] = agr.get(g, 0.0) + float(x[im])
ordem = ["Sênior", "Mezanino", "Subordinado", "Subordinado Jr."]
vals = [round(agr.get(k, 0.0), 1) for k in ordem]
waterfall(sl, MARGEM, Emu(1417320), col_w(7), Emu(3017520), ordem, vals, LARANJA,
          total_label="Emitido no programa")
total = sum(vals)
prot = [round(100 * sum(vals[i + 1:]) / total, 1) for i in range(len(vals))]
txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(274320), "Proteção abaixo de cada camada", tam=13, negrito=True, cor=PRETO)
tabela(sl, col_x(7), Emu(1737360), col_w(5), Emu(1188720),
       ["Camada", "R$ mi", "% do total", "Proteção abaixo"],
       [[k, f"{v:,.1f}".replace(",", "."), f"{100 * v / total:.1f}%", f"{p:.1f}%"]
        for k, v, p in zip(ordem, vals, prot)],
       larguras=[0.34, 0.22, 0.22, 0.22], tam=9.5, alinh_centro=(1, 2, 3), altura_linha=Emu(228600))
txt(sl, col_x(7), Emu(3108960), col_w(5), Emu(2103120),
    "A leitura consolidada tem um limite importante: somar camadas de veículos diferentes não cria "
    "uma estrutura única. Cada patrimônio separado e cada fundo tem cascata própria e não há "
    "cross-collateral entre eles — a subordinada de um CRI não protege a sênior de outro.\n\n"
    "O número serve para dimensionar quanto do programa a originadora financiou nas camadas que "
    "absorvem perda primeiro, e não como medida de proteção de nenhuma série individual.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
leitura(sl, Emu(4663440),
        "Não há cross-collateral entre veículos: a proteção consolidada é uma medida de programa, não de série. A proteção que importa ao investidor é a do seu próprio patrimônio separado.")
rodape(sl, "Agregação das 70 tranches de 02_Series por camada, a valor nominal de emissão. As classes de cotas de FIDC entram pelo mesmo critério de senioridade. " + FONTE_GERAL)

# ============================================================ 17. Resgate da subordinada
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Extração de subordinada", "Em que condições a subordinada sai antes do fim",
          "A diferença entre FIDC e CRI não é se pode sair, mas o que trava a saída: governança de um lado, cobertura do outro.")
bloco(sl, MARGEM, Emu(1417320), col_w(6), Emu(365760), "FIDC — depende de pedido e de quórum", PETROLEO, BRANCO, tam=11)
c, r = L("30_resgate_subordinada.csv")
fid = [x for x in r if x[1] == "FIDC"]
tabela(sl, MARGEM, Emu(1828800), col_w(6), Emu(1737360),
       ["Fundo", "Quem pede", "Testes exigidos"],
       [[x[0].replace("FIDC-", "FIDC "), x[4] if x[4] != "n/d" else "n/d", x[3][:64]] for x in fid],
       larguras=[0.16, 0.28, 0.56], tam=8.5, altura_linha=Emu(228600), destaque=(3,))
bloco(sl, col_x(6), Emu(1417320), col_w(6), Emu(365760), "CRI — automático, travado por cobertura", LARANJA, PRETO, tam=11)
txt(sl, col_x(6), Emu(1874520), col_w(6), Emu(1691640),
    "Não há pedido, não há quórum, não há deliberação. A série retida recebe em cada Data de Pagamento "
    "se as Razões de Cobertura estiverem enquadradas e o saldo estiver acima do Target.\n\n"
    "Ao final, depois do resgate das quatro séries públicas, ela recebe todo o remanescente do "
    "patrimônio separado a título de Prêmio Final.\n\n"
    "A trava é a cobertura, não a governança.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
tabela(sl, col_x(6), Emu(2926080), col_w(6), Emu(640080),
       ["Operação", "Coberturas exigidas", "Carência documentada"],
       [["CRI-II", "4 razões (159/123/110/105%)", "13º pagamento — 07/07/2025"],
        ["CRI-VI", "3 razões (120,48/109,89/105,26%)", "Reserva de derivativo antes da 5ª série"]],
       larguras=[0.18, 0.42, 0.40], tam=9, altura_linha=Emu(228600))
bloco(sl, MARGEM, Emu(3931920), UTIL_W, Emu(640080),
      "Nos FIDCs VI e VII o quórum para a amortização extraordinária é de 75% dos titulares da classe júnior — "
      "que é a classe da originadora. Quem pede e quem se beneficia são a mesma parte.",
      CINZA_FUNDO, PRETO, tam=11, negrito=True, borda=CINZA_CLARO)
leitura(sl, Emu(4754880),
        "Sim, a subordinada pode funcionar como fonte de caixa da originadora — e nos CRIs isso não depende de nenhuma aprovação, só de a cobertura estar enquadrada na data.")
txt(sl, MARGEM, Emu(5166360), UTIL_W, Emu(548640),
    "A exceção documentada é a debênture: a 2ª série é integralmente subordinada à 1ª por contrato e, em 31/07/2026, "
    "o relatório do agente fiduciário registra R$ 0,00 distribuídos aos seus investidores — consistente com a cláusula 3.3.2 da escritura.",
    tam=11, cor=CINZA_ESC, espaco=3)
rodape(sl, "Cl. 6.5.1 (y), (z) e (bb) do Termo da 2ª Kanastra; cl. 7.5.1 (r), (s) e (u) do Termo da 177ª VERT; cl. 3.3.2 da escritura da debênture. Regras dos FIDCs conforme análise de crédito — os regulamentos não estão no acervo.")

# ============================================================ 18. Histórico de saques
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Histórico", "Já saiu — em todos os sete fundos e em três dos seis CRIs",
          "Principal de mezanino e júnior efetivamente amortizado, conforme os informes mensais. O que não se verifica é se cada saída passou nos testes.")
c, r = L("33_historico_saques_subordinada.csv")
por_v = {}
for x in r:
    if x[1] == "Debênture":
        continue
    por_v.setdefault(x[0], {"Mezanino": 0.0, "Subordinado Jr.": 0.0})
    chave = "Mezanino" if x[2] == "Mezanino" else "Subordinado Jr."
    por_v[x[0]][chave] += float(x[6])
ordem_v = sorted(por_v, key=lambda k: -sum(por_v[k].values()))
cats = [v.replace("FIDC-", "FIDC ").replace("CRI-", "CRI ") for v in ordem_v][::-1]
grafico(sl, XL_CHART_TYPE.BAR_STACKED, MARGEM, Emu(1417320), col_w(7), Emu(3383280), cats,
        [("Mezanino", [por_v[v]["Mezanino"] for v in ordem_v][::-1]),
         ("Subordinado / Júnior", [por_v[v]["Subordinado Jr."] for v in ordem_v][::-1])],
        [CINZA_MED, LARANJA], unidade_y="R$ mi de principal amortizado")
txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(274320), "Os dois casos que pedem explicação", tam=13, negrito=True, cor=PRETO)
txt(sl, col_x(7), Emu(1737360), col_w(5), Emu(1554480),
    "FIDC IV — R$ 438,9 mi de mezanino e júnior amortizados, com um único mês de R$ 151,4 mi. É o maior "
    "volume do programa, saiu do fundo com o mandato mais permissivo, e hoje o fundo está sem first "
    "loss corrente e com PDD de 69,9% da carteira.\n\n"
    "FIDC VII — R$ 7,7 mi de júnior amortizados em julho de 2026, no mesmo mês do take-out da 177ª "
    "emissão e no primeiro mês em que o fundo cedeu carteira.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
bloco(sl, col_x(7), Emu(3383280), col_w(5), Emu(914400),
      "Se cada saída respeitou o regulamento: não disponível. O informe mensal comprova o pagamento, "
      "mas o teste de subordinação e cobertura na data de cada amortização não é publicado.",
      CINZA_FUNDO, PRETO, tam=10.5, negrito=True, borda=CINZA_CLARO)
leitura(sl, Emu(4937760),
        "R$ 1,06 bi de mezanino e júnior já saíram dos sete fundos. A ocorrência comprova extração econômica; a aderência contratual, não.")
txt(sl, MARGEM, Emu(5349240), UTIL_W, Emu(548640),
    "É o item de diligência com maior assimetria entre o que se observa e o que se pode afirmar. O pedido ao originador e aos "
    "administradores é objetivo: os demonstrativos dos índices de subordinação, cobertura e reservas na data de cada amortização "
    "extraordinária, fundo a fundo, desde 2022.",
    tam=11, cor=CINZA_ESC, espaco=3)
rodape(sl, "CVM — Informe Mensal FIDC e Informe Mensal CRI, até 31/07/2026, conforme análise de crédito de 21/08/2026. A debênture não aparece: sua 2ª série não recebeu nada até a data-base.")

# ============================================================ 19. PDD
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Risco de crédito", "Duas curvas de provisão, e por que a razão passa de 100%",
          "A curva posterior reconhece mais perda entre 91 e 180 dias. Reconhecimento mais conservador, não piora do ativo.")
grafico(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, MARGEM, Emu(1417320), col_w(7), Emu(3017520),
        ["Até 15d", "16-30", "31-60", "61-90", "91-120", "121-150", "151-180", "> 180"],
        [("Curva inicial  CRI-I e CRI-II", [0.0, 1.0, 3.0, 10.0, 30.0, 50.0, 70.0, 100.0]),
         ("Curva posterior  CRI-III a CRI-VI", [0.0, 1.5, 5.0, 10.0, 37.0, 58.0, 78.0, 100.0])],
        [CINZA_CLARO, LARANJA], unidade_y="% do valor presente do recebível")
txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(274320), "Por que PDD / >90d passa de 100%", tam=13, negrito=True, cor=PRETO)
txt(sl, col_x(7), Emu(1737360), col_w(5), Emu(1828800),
    "Dois mecanismos contratuais distintos.\n\n"
    "A provisão incide sobre o valor presente do recebível — o saldo que resta da CCB — e não sobre a "
    "parcela vencida. O numerador mede contrato inteiro; o denominador, só as parcelas atrasadas.\n\n"
    "E o Efeito Vagão: se um devedor atrasa em um contrato, todos os contratos dele passam a ser "
    "tratados pelo pior atraso, inclusive os que estão em dia.",
    tam=10.5, cor=CINZA_ESC, espaco=5)
tabela(sl, col_x(7), Emu(3657600), col_w(5), Emu(914400),
       ["Veículo", "PDD/carteira", "PDD / >90d"],
       [["FIDC IV", "69,9%", "550%"], ["FIDC VI", "48,8%", "581%"], ["FIDC VII", "0,2%", "3.725%"]],
       larguras=[0.40, 0.30, 0.30], tam=10, alinh_centro=(1, 2), altura_linha=Emu(228600))
leitura(sl, Emu(4754880),
        "A razão PDD sobre saldo vencido não é cobertura de perda: é efeito de dois denominadores diferentes somado ao arrasto entre contratos do mesmo devedor.")
txt(sl, MARGEM, Emu(5166360), UTIL_W, Emu(548640),
    "A inadimplência por safra que as lâminas descrevem tem outro denominador ainda — o total originado na safra — e separa perda bruta "
    "de perda líquida, sendo a diferença as recuperações após 90 dias. As duas métricas não se comparam.",
    tam=11, cor=CINZA_ESC, espaco=3)
rodape(sl, "Tabelas de PDD por faixa de atraso dos Prospectos Definitivos da 1ª e 2ª emissões; definição literal de Efeito Vagão do mesmo documento. " + FONTE_GERAL)

# ============================================================ 20. Veredito
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Veredito", "Onde o CRI ganha, onde não ganha, e o que não dá para afirmar",
          "Doze dimensões comparadas. A terceira coluna é a que importa para o comitê: o que ainda falta para fechar cada uma.")
c, r = L("15_fidc_vs_cri.csv")
tabela(sl, MARGEM, Emu(1417320), col_w(12), Emu(3931920),
       ["Dimensão", "Vantagem", "O que falta para confirmar"],
       [[x[0], x[3], x[5]] for x in r],
       larguras=[0.24, 0.13, 0.63], tam=8.5, tam_cab=9.5,
       altura_cab=Emu(274320), altura_linha=Emu(274320), alinh_centro=(1,))
leitura(sl, Emu(5486400),
        "Cinco dimensões favorecem o CRI, quatro favorecem o FIDC e três não têm veredito com dado público. Nenhuma conclusão de custo é possível sem a curva de juros.")
rodape(sl, "Detalhe por dimensão, com a evidência de cada linha, na aba 15_FIDC_vs_CRI do workbook. " + FONTE_GERAL)

# ============================================================ 21. Pontos em aberto
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Due diligence", "O que permanece sem resposta nos documentos disponíveis",
          "Dez perguntas em ordem de prioridade, com o documento que as fecharia e o impacto de cada uma na decisão.")
c, r = L("34_pontos_em_aberto.csv")
tabela(sl, MARGEM, Emu(1417320), col_w(12), Emu(3931920),
       ["#", "Pergunta do comitê", "O que falta", "Onde obter", "Impacto na decisão"],
       [[x[0], x[1], x[2], x[3], x[4]] for x in r],
       larguras=[0.04, 0.24, 0.25, 0.20, 0.27], tam=8, tam_cab=9,
       altura_cab=Emu(274320), altura_linha=Emu(274320), alinh_centro=(0,))
leitura(sl, Emu(5486400),
        "Os quatro primeiros bloqueiam, nesta ordem: a economia da cessão, a aderência das saídas de subordinada, a cascata de quatro operações e os critérios literais dos fundos.")
rodape(sl, "Inventário completo de fontes, incluindo as buscas sem resultado, na aba 17_Fontes do workbook. " + FONTE_GERAL)

# ---------------------------------------------------------------- salvar
os.makedirs(OUTDIR, exist_ok=True)
caminho = os.path.join(OUTDIR, "Solfacil_CRI_FIDC_20260822_claude.pptx")
prs.save(caminho)
print(f"Deck salvo: {caminho}")
print(f"Slides: {len(prs.slides._sldIdLst)}")
