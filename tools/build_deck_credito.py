# -*- coding: utf-8 -*-
"""Deck executivo de renovação de crédito - Solfácil FIDC/CRI.

Versão adicional, derivada do deck analítico. Corpo principal enxuto com
conclusão no título; detalhamento documental no appendix. Todos os objetos
são nativos e editáveis.
"""
import csv, os, sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
DATA = os.path.join(ROOT, "data", "solfacil_claude")
OUTDIR = os.path.join(ROOT, "outputs", "solfacil")
sys.path.insert(0, HERE)
from deck_credito_base import *   # noqa: F401,F403
import deck_credito_base as C

prs = Presentation()
prs.slide_width, prs.slide_height = SL_W, SL_H
VAZIO = prs.slide_layouts[6]
DB = "data-base jul/26"


def L(nome):
    return ler(DATA, nome)


def grafico(sl, tipo, x, y, w, h, cats, series, cores, **kw):
    d = CategoryChartData(); d.categories = cats
    for nome, vals in series:
        d.add_series(nome, tuple(vals))
    gf = sl.shapes.add_chart(tipo, x, y, w, h, d)
    estilizar_gr(gf.chart, cores, **kw)
    return gf.chart


# ══════════════════════════════════════════════ 1. Capa e síntese
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "Renovação de crédito · Solfácil · Crédito estruturado")
titulo_conclusivo(
    sl, "Programa depende de reciclagem contínua de warehouses; proteção contratual é\n"
        "relevante, mas subordinada já foi sacada em 7/7 fundos e performance por safra é inconclusiva")
txt(sl, MARGEM, Emu(1234440), col_w(12), Emu(228600),
    "Sete FIDCs, seis operações de CRI e uma debênture · R$ 8.304,5 mi emitidos · " + DB,
    tam=11, cor=CINZA_ESC)

for i, (v, r) in enumerate([("14", "veículos ativos"), ("R$ 1.234,6 mi", "PL somado dos FIDCs"),
                            ("R$ 3.670,6 mi", "take-outs em CRI"), ("R$ 1.076,2 mi", "subordinada já amortizada"),
                            ("26,0%", "subordinação do FIDC VII"), ("+1,0 p.p.", "folga do VII ao piso")]):
    destaque_numero(sl, col_x(2 * i), Emu(1691640), col_w(2), v, r)

rotulo_secao(sl, Emu(2377440), "Síntese de crédito")
caixa_sintese(sl, MARGEM, Emu(2560320), col_w(4), Emu(2011680), "Proteções", [
    ("Granularidade.", "CRI limita concentração individual a 0,07%–0,25% do patrimônio separado e WAM a 2.000 dias."),
    ("Subordinação.", "Piso de 20%–25% em 6/7 fundos; subordinação observada de 20,9% a 36,5%."),
    ("Cascata travada.", "Razões de cobertura de 105% a 159% condicionam cada degrau (CRI-II e CRI-VI)."),
    ("Retenção.", "Série privada de 2,49%–3,00% em 6/6 CRIs, subscrita pela Solfácil."),
    ("PDD contratual.", "Curva posterior provisiona 37%/58%/78% entre 91 e 180 dias."),
], PETROLEO)
caixa_sintese(sl, col_x(4), Emu(2560320), col_w(4), Emu(2011680), "Pontos de atenção", [
    ("Cash leakage.", "R$ 1.076,2 mi de mezanino e júnior sacados; aderência aos testes não verificável."),
    ("FIDC IV.", "Sem first loss e sem piso; PDD/DC de 69,9%; R$ 438,9 mi já sacados."),
    ("Folga mínima.", "FIDC V a +0,9 p.p. e FIDC VII a +1,0 p.p. do piso."),
    ("Basis/hedge.", "Em CRI-VI o diferencial do swap recai sobre a série retida."),
    ("Preço de cessão.", "n/d — evidenciado apenas nos Termos de Cessão, não públicos."),
], LARANJA)
caixa_sintese(sl, col_x(8), Emu(2560320), col_w(4), Emu(2011680), "Condicionantes para aprovação", [
    ("Tape por CCB.", "Sem ele não há conclusão sobre performance por safra."),
    ("Testes de subordinação.", "Demonstrativos na data de cada amortização extraordinária."),
    ("Termos de securitização.", "4 das 6 operações sem termo no acervo."),
    ("Regulamentos dos FIDCs.", "Critérios literais não obtidos; parâmetros vêm de análise secundária."),
    ("Prazo de exposição.", "Aprovar pelo vencimento legal, não pela duration."),
], CINZA_ESC)

txt(sl, MARGEM, Emu(4709160), UTIL_W, Emu(457200),
    "Encaminhamento proposto: renovação condicionada à entrega do tape por CCB, dos demonstrativos de teste "
    "das amortizações subordinadas e dos termos de securitização faltantes. Exposição dimensionada pelo "
    "vencimento legal da série sênior.",
    tam=10.5, cor=PRETO, espaco=3)
fonte(sl, "Fonte: regulamentos e termos de securitização; CVM (cadastro, informes mensais, ofertas); escritura e relatório do agente fiduciário da debênture; " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 1

Universo (14 veículos): 7 FIDCs (CVM, cadastro de fundos), 6 operações de CRI, 1 emissão de debêntures
(Escritura da 1ª emissão da Amazônia Solar Cia. Securitizadora de Créditos Financeiros, 18/02/2022,
CNPJ 43.102.521/0001-62).

R$ 8.304,5 mi emitidos = FIDC R$ 4.573,9 mi + CRI R$ 3.670,6 mi + debênture R$ 60,0 mi. Soma das 70
tranches com montante documentado (02_Series). Valor nominal de emissão, não saldo em aberto.

PL somado dos FIDCs R$ 1.234,6 mi e subordinada amortizada R$ 1.076,2 mi: CVM — Informe Mensal FIDC,
competência 31/07/2026.

Subordinação do FIDC VII de 26,0% e folga de +1,0 p.p. ao piso de 25%: informe mensal 31/07/2026.
Folga = [Sub_NAV − piso × PL] / [1 − piso]. Subordinação observada, não target.

Concentração 0,07%–0,25% e WAM 2.000 dias: critérios de elegibilidade das lâminas de CRI-I (15/01/2024),
CRI-III (22/04/2025) e CRI-V (17/04/2026). São CRITÉRIOS DE ELEGIBILIDADE contratuais, não
características observadas da carteira.

Razões de cobertura 105%–159%: cl. 6.5.1 do 2º Aditamento ao Termo de Securitização da 2ª emissão
Kanastra (registro JUCEMG 12803230, 10/06/2025) e cl. 7.5.1 do Termo da 177ª emissão VERT (20/07/2026).
São TRIGGERS de alocação, não targets de subordinação.

Basis/hedge: cl. 15.14.6 item 7 e cl. 15.14.8 item 3 do Termo da 177ª VERT.

Preço de cessão: fórmula na definição de Preço de Aquisição do Termo da 177ª VERT; o valor praticado é
evidenciado em cada Termo de Cessão, não público. Classificado n/d.

HIERARQUIA DE FONTES: parâmetros de CRI-II e CRI-VI vêm de termo de securitização (Nível 1). Parâmetros
dos FIDCs e das demais operações vêm de análise de crédito secundária de 21/08/2026 consolidando CVM
(Nível 2), por ausência dos regulamentos e termos no acervo — divergência registrada como limitação.""")

# ══════════════════════════════════════════════ 2. Slide A — estrutura de funding
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "A · Estrutura de funding")
titulo_conclusivo(
    sl, "Warehouses são reciclados, não substituídos: 3 fundos alimentaram 4 take-outs cada\n"
        "e nenhum dos 7 FIDCs foi encerrado",
    "Originação Solfácil → FIDC (warehouse) → seleção de CCBs por elegibilidade → CRI (take-out) → investidores")

ETAPAS = [("ORIGINAÇÃO", "Solfácil · ~4.000 integradores\nCCB pré-fixada PF e PJ\nAlienação fiduciária do equipamento", CINZA_ESC),
          ("WAREHOUSE", "FIDCs I–VII · R$ 1.234,6 mi de PL\nRevolvência e reinvestimento\nSubordinação de 20,9% a 36,5%", PETROLEO),
          ("TAKE-OUT", "6 CRIs · R$ 3.670,6 mi\nCessão definitiva sem coobrigação\nNovo patrimônio separado", LARANJA),
          ("INVESTIDORES", "Sênior a Subordinado Jr.\n2,1 a 7,9 mil PF por oferta\nSérie privada retida pela Solfácil", CINZA_ESC)]
LARG, Y_CX, H_CX = col_w(3), Emu(1417320), Emu(1188720)
caixas = []
for i, (t, corpo, cor) in enumerate(ETAPAS):
    x = col_x(3 * i)
    cx = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Y_CX, LARG, H_CX)
    cx.fill.solid(); cx.fill.fore_color.rgb = BRANCO
    cx.line.color.rgb = CINZA_CLARO; cx.line.width = Pt(0.75); sem_sombra(cx)
    caixas.append(cx)
    barra_lateral(sl, x, Y_CX, LARG, Emu(45720), cor)
    txt(sl, Emu(x.emu + 91440), Emu(Y_CX.emu + 128016), Emu(LARG.emu - 182880), Emu(182880),
        f"{i + 1}. {t}", tam=9.5, negrito=True, cor=PRETO)
    txt(sl, Emu(x.emu + 91440), Emu(Y_CX.emu + 365760), Emu(LARG.emu - 182880), Emu(731520),
        corpo, tam=9, cor=CINZA_ESC, espaco=2)
for i in range(3):
    a = caixas[i]
    s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(a.left.emu + a.width.emu + 18288),
                            Emu(Y_CX.emu + int(H_CX.emu / 2) - 54864), Emu(100584), Emu(109728))
    s.fill.solid(); s.fill.fore_color.rgb = CINZA_MED; s.line.fill.background(); sem_sombra(s)

rotulo_secao(sl, Emu(2743200), "Onde o risco permanece e para onde migra")
tab(sl, MARGEM, Emu(2926080), col_w(7), Emu(1737360),
    ["Fundo", "Início", "PL (R$ mi)", "Carteira (R$ mi)", "Subord. obs.", "Cedeu para", "Fase atual"],
    [["I", "dez/20", "83,7", "77,6", "36,5%", "—", "Originando"],
     ["II", "out/21", "94,1", "102,5", "27,1%", "4 CRIs", "Originando"],
     ["III", "jul/23", "141,1", "131,5", "33,2%", "—", "Originando"],
     ["IV", "jun/22", "17,5", "14,6", "0,0%", "4 CRIs", "Sem first loss"],
     ["V", "mar/23", "67,5", "66,7", "20,9%", "—", "Originando"],
     ["VI", "nov/24", "211,1", "147,7", "33,6%", "4 CRIs", "Runoff pós-cessão"],
     ["VII", "fev/26", "619,6", "446,1", "26,0%", "1 CRI", "Revolvência 12m"]],
    [0.07, 0.10, 0.13, 0.15, 0.14, 0.13, 0.28], num_cols=(2, 3, 4), bold_cols=(0,),
    bold_cells=((3, 4), (3, 6), (6, 2)), h_linha=Emu(196596), destaque_linhas=(3,))

txt(sl, col_x(7), Emu(2926080), col_w(5), Emu(182880), "LEITURA DE CRÉDITO", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, col_x(7), Emu(3108960), col_w(5), Emu(1554480),
    "O FIDC VII concentra 50,2% do PL somado dos sete fundos e é o único em revolvência obrigatória — "
    "12 meses com principal suspenso e saldo reinvestido.\n\n"
    "Os fundos I, III e V nunca aparecem como cedentes em nenhum documento do acervo. O programa não "
    "usa todos os warehouses para take-out.\n\n"
    "O FIDC IV é exceção material: subordinação observada de 0,0%, sem piso contratual, e foi de onde "
    "saiu o maior volume de subordinada do programa.",
    tam=9.5, cor=CINZA_ESC, espaco=4)
txt(sl, MARGEM, Emu(4800600), UTIL_W, Emu(365760),
    "Dependência de funding: os take-outs recompram carteira dos warehouses e devolvem caixa, o que sustenta a originação seguinte. "
    "Interromper o canal de CRI transfere a necessidade de funding integralmente para os FIDCs.",
    tam=10, cor=PRETO, espaco=3)
fonte(sl, "Fonte: CVM — Informe Mensal FIDC e cadastro de fundos; prospectos definitivos da 1ª e 2ª emissões Kanastra; termo de securitização da 177ª VERT; " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 2 (A)

PL, carteira e subordinação observada: CVM — Informe Mensal FIDC, competência 31/07/2026, consolidado na
análise de crédito de 21/08/2026. Subordinação observada = (NAV mezanino + NAV júnior) / PL. É medida
OBSERVADA na data-base, distinta do piso contratual (target), que consta do slide D.

Cedentes documentados em fonte Nível 1:
- FIDC II (CNPJ 42.462.306/0001-00) e FIDC IV (CNPJ 44.909.456/0001-44): Prospecto Definitivo da 1ª
  emissão Kanastra, seção 11.1.2 "Representantes de mais de 10% dos Direitos Creditórios Imobiliários
  Cedidos"; e definição de "Cedentes" do Prospecto Definitivo da 2ª emissão.
- FIDC VI (IS Green Solfácil VI FIDC RL, CNPJ 57.028.406/0001-08) e FIDC VII (Solfácil Crédito Pessoal
  VII FIDC RL, CNPJ 63.505.455/0001-89): Termo de Securitização da 177ª emissão VERT, 20/07/2026,
  definição de "Cedentes Fundos".
- Vínculos de CRI-III, CRI-IV e CRI-V com FIDC II/IV/VI: análise de crédito de 21/08/2026 (Nível 6
  consolidando Nível 2). Os termos de securitização dessas operações não estão no acervo — classificado
  como n/l. Registrado como divergência de nível de evidência.

"Nunca cedeu" para I, III e V significa ausência de menção como cedente nos documentos analisados; não
equivale a proibição contratual de ceder.

FIDC VII concentra 50,2% do PL: 619,6 / 1.234,6 = 50,2%. Cálculo próprio sobre o informe mensal.

Revolvência obrigatória de 12 meses do FIDC VII e ausência de piso no FIDC IV: análise de crédito de
21/08/2026; os regulamentos dos fundos não foram localizados (n/l).""")

# ══════════════════════════════════════════════ 3. Slide B — comparativo dos FIDCs
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "B · Comparativo dos warehouses")
titulo_conclusivo(
    sl, "FIDC IV é a exceção material do programa: sem first loss, sem piso e com PDD/carteira de 69,9%;\n"
        "V e VII operam a menos de 1 p.p. do piso de subordinação",
    "Subordinação, PDD e >90d são observados na data-base. Cap, WAM e piso são limites contratuais — não descrevem a carteira efetiva.")
tab(sl, MARGEM, Emu(1371600), col_w(12), Emu(1920240),
    ["FIDC", "Início", "Carteira\n(R$ mi)", "Subord.\nobs.", "Piso", "Folga\n(p.p.)", "PDD /\ncarteira",
     ">90d /\ncarteira", "Cap por\ndevedor", "WAM\nmáx (d)", "Subord.\nsacada (R$ mi)", "Take-out"],
    [["I", "dez/20", "77,6", "36,5%", "25%", "+11,5", "6,8%", "1,3%", "2,00%", "2.135", "130,4", "—"],
     ["II", "out/21", "102,5", "27,1%", "20%", "+7,1", "0,0%", "2,1%", "2,00%", "2.400", "160,8", "4 CRIs"],
     ["III", "jul/23", "131,5", "33,2%", "25%", "+8,2", "5,7%", "1,2%", "0,10%", "2.000", "101,4", "—"],
     ["IV", "jun/22", "14,6", "0,0%", "n.a.", "n.a.", "69,9%", "12,7%", "20,00%", "n/d", "438,9", "4 CRIs"],
     ["V", "mar/23", "66,7", "20,9%", "20%", "+0,9", "9,3%", "1,7%", "2,00%", "n/d", "34,4", "—"],
     ["VI", "nov/24", "147,7", "33,6%", "25%", "+8,6", "48,8%", "8,4%", "n/d", "2.400", "183,2", "4 CRIs"],
     ["VII", "fev/26", "446,1", "26,0%", "25%", "+1,0", "0,2%", "0,0%", "n/d", "2.400", "7,7", "1 CRI"]],
    [0.05, 0.07, 0.08, 0.08, 0.06, 0.07, 0.08, 0.08, 0.09, 0.08, 0.11, 0.09],
    tam=8.5, num_cols=(2, 3, 4, 5, 6, 7, 8, 9, 10), bold_cols=(0,),
    bold_cells=((3, 3), (3, 6), (3, 7), (3, 8), (3, 10), (4, 5), (6, 5), (5, 6), (6, 2)),
    h_linha=Emu(210312), destaque_linhas=(3,))

rotulo_secao(sl, Emu(3383280), "Diferenças materiais")
ITENS = [
    ("FIDC IV", "Subordinação observada de 0,0% e sem piso contratual. PDD/carteira de 69,9% sobre carteira de R$ 14,6 mi. "
                "Cap por devedor de 20% — 200 vezes o do FIDC III. R$ 438,9 mi de subordinada já sacados, o maior volume do programa."),
    ("FIDC V e VII", "Folga ao piso de +0,9 p.p. e +1,0 p.p. Qualquer deterioração relevante do NAV subordinado aproxima do piso. "
                     "O VII responde por 50,2% do PL somado e tem apenas 6 competências de histórico."),
    ("FIDC VI", "PDD/carteira de 48,8% após o take-out de jul/26, que reduziu a carteira de R$ 399,3 mi para R$ 147,7 mi. "
                "O informe não decompõe a variação — efeito de denominador e deterioração não são separáveis."),
    ("FIDC III", "Único warehouse com parâmetros contratuais equivalentes aos dos CRIs: cap de 0,10% e WAM de 2.000 dias."),
]
for i, (t, corpo) in enumerate(ITENS):
    y = Emu(3566160 + i * 493776)
    barra_lateral(sl, MARGEM, y, Emu(36576), Emu(420624), LARANJA if i < 3 else PETROLEO)
    txt(sl, Emu(MARGEM.emu + 109728), y, Emu(1188720), Emu(182880), t, tam=9.5, negrito=True, cor=PRETO)
    txt(sl, Emu(MARGEM.emu + 1371600), y, Emu(UTIL_W.emu - 1371600), Emu(420624), corpo,
        tam=9, cor=CINZA_ESC, espaco=0)
fonte(sl, "Fonte: CVM — Informe Mensal FIDC (subordinação, PDD, >90d, amortização); regulamentos via análise consolidada (cap, WAM, piso); " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 3 (B)

COLUNAS OBSERVADAS (CVM — Informe Mensal FIDC, competência 31/07/2026):
- Carteira: saldo de direitos creditórios.
- Subordinação observada = (NAV mezanino + NAV júnior) / PL. Valores: I 36,46%; II 27,13%; III 33,19%;
  IV 0,00%; V 20,91%; VI 33,63%; VII 26,00%.
- PDD/carteira e >90d/carteira: PDD já reduz o PL; >90d soma parcelas vencidas, não o saldo integral
  da CCB. A razão PDD/>90d supera 100% em 6 dos 7 fundos por dois mecanismos contratuais: a provisão
  incide sobre o valor presente do recebível e o Efeito Vagão arrasta a pior faixa de atraso do devedor
  para todas as CCBs dele (definição literal no Prospecto Definitivo da 1ª emissão Kanastra).
- Subordinada sacada: principal de mezanino e júnior acumulado até 31/07/2026 (FIDC IV até 30/11/2025).

COLUNAS CONTRATUAIS (limites, não observações):
- Piso: TARGET de subordinação do regulamento. FIDC IV sem piso — classificado n.a.
- Folga ao piso = [Sub_NAV − piso × PL] / [1 − piso]. Cálculo próprio.
- Cap por devedor e WAM máximo: CRITÉRIOS DE ELEGIBILIDADE. Fonte: análise de crédito de 21/08/2026
  consolidando regulamentos; os regulamentos vigentes não foram localizados no acervo (n/l). Cap dos
  FIDCs VI e VII não divulgado (n/d). WAM dos FIDCs IV e V não divulgado (n/d).

DIVERGÊNCIA DE FONTES: parâmetros contratuais dos FIDCs seriam Nível 1 (regulamento). Na ausência,
usou-se consolidação secundária. Recomenda-se confirmar contra os regulamentos antes da aprovação.

FIDC VI: PL 30/06/2026 R$ 437,4 mi e carteira R$ 399,3 mi; PL 31/07/2026 R$ 211,1 mi e carteira
R$ 147,7 mi. A emissão da 177ª VERT ocorreu em 21/07/2026. O informe não decompõe a queda entre cessão,
ajuste de valor e distribuição — limitação declarada.

FIDC IV: cap de 20% contra 0,10% do FIDC III = razão de 200x.""")

# ══════════════════════════════════════════════ 4. Slide C — seleção nos take-outs
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "C · Qualidade relativa dos pools cedidos")
titulo_conclusivo(
    sl, "Seleção nos take-outs é contratual, não performada: CRI aperta WAM, concentração e prazo,\n"
        "mas nenhuma operação exige seasoning e não há dado para comprovar cherry-picking",
    "Comparação entre critérios de elegibilidade. Nenhuma coluna descreve carteira observada.")
tab(sl, MARGEM, Emu(1371600), col_w(12), Emu(1737360),
    ["Dimensão", "Warehouses (faixa dos 7)", "CRI-I  jan/24", "CRI-III  mai/25", "CRI-V  mai/26", "Δ contratual"],
    [["Cap por devedor", "0,10% a 20,00%", "0,10%", "0,10%", "0,15% → 0,07%", "CRI aperta"],
     ["Cap dos 10 maiores", "1% a 10%", "1%", "não consta", "não consta", "Afrouxa após CRI-I"],
     ["WAM máx. do pool", "2.000 a 2.400 d", "2.000 d", "2.000 d", "2.000 d (sobre VP)", "CRI aperta"],
     ["Prazo máx. por recebível", "3.836 a 4.760 d", "3.845 d", "3.845 d", "3.845 d", "CRI aperta"],
     ["Ticket máx. PF / PJ", "201–500 / 500–700 mil", "350 / 600 mil", "350 / 700 mil", "350 / 700 mil", "Neutro"],
     ["Adimplência na cessão", "exigida em II–VII", "exigida", "exigida", "exigida", "Equivalente"],
     ["Parcela balão final", "n/d", "não consta", "não consta", "vedada", "CRI aperta em V"],
     ["Idade máx. PF", "71 anos (VI e VII)", "71 anos", "71 anos", "71 anos", "Equivalente"],
     ["Constituição mín. PJ", "2 anos", "2 anos", "2 anos", "2 anos + CMN 5.118", "CRI aperta em V"],
     ["Carência máxima", "180 a 366 d", "185 d", "185 d", "185 d", "CRI aperta"],
     ["Seasoning mínimo", "n/d", "não exigido", "não exigido", "não exigido", "Ausente em 6/6"]],
    [0.19, 0.19, 0.14, 0.14, 0.17, 0.17], tam=8.5, num_cols=(),
    bold_cols=(0,), bold_cells=((0, 4), (6, 4), (8, 4), (10, 1), (10, 2), (10, 3), (10, 4)),
    h_linha=Emu(140208), destaque_linhas=(10,))

rotulo_secao(sl, Emu(3200400), "Respostas objetivas")
QA = [
    ("Os CRIs recebem critérios mais restritivos?", "Sim, em quatro dimensões: concentração individual, WAM do pool, prazo máximo por recebível e carência."),
    ("Em quais dimensões não há diferença?", "Ticket máximo, idade máxima PF e exigência de adimplência na cessão são equivalentes aos warehouses."),
    ("É seleção contratual ou performance superior?", "Apenas contratual. Nenhum documento apresenta NPL, roll rate ou MoB do lote cedido."),
    ("Há evidência de cherry-picking?", "Não há evidência em nenhuma direção. Testar exige o tape por CCB do lote cedido e do que permaneceu no fundo."),
    ("Existe exigência de seasoning?", "Não em nenhuma das seis operações. O crédito pode ser cedido sem histórico de pagamento."),
    ("Os critérios mudaram ao longo das emissões?", "Sim. CRI-V introduziu vedação a parcela balão, enquadramento na CMN 5.118 e cap escalonado por integralização."),
]
for i, (q, a) in enumerate(QA):
    y = Emu(3383280 + (i % 3) * 393192)
    x = MARGEM if i < 3 else col_x(6)
    w = col_w(6)
    txt(sl, x, y, w, Emu(160020), q, tam=9, negrito=True, cor=PRETO)
    txt(sl, x, Emu(y.emu + 155448), w, Emu(228600), a, tam=8.5, cor=CINZA_ESC)

txt(sl, MARGEM, Emu(4663440), UTIL_W, Emu(365760),
    "Implicação de crédito: o pool do CRI é contratualmente mais granular e mais curto que o mandato do warehouse, "
    "mas a ausência de seasoning significa que crédito recém-originado pode ser cedido sem qualquer histórico de pagamento.",
    tam=10, cor=PRETO, espaco=3)
fonte(sl, "Fonte: critérios de elegibilidade das lâminas de CRI-I (15/01/24), CRI-III (22/04/25) e CRI-V (17/04/26); parâmetros dos FIDCs via análise consolidada; " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 4 (C)

Todas as colunas de CRI são CRITÉRIOS DE ELEGIBILIDADE literais, extraídos das lâminas:
- CRI-I: Lâmina da Oferta da 1ª emissão Kanastra, item 3, critérios i a xi. Inclui o cap dos 10 maiores
  devedores em 1% do Patrimônio Separado — único do programa.
- CRI-III: Lâmina de 22/04/2025, critérios i a xi.
- CRI-V: Lâmina de 17/04/2026, critérios i a xiii. Critério (ix) veda parcela final superior às demais
  (balão); critério (xi) e (xiii) exigem enquadramento do PJ na Resolução CMN 5.118; critério (iii)
  escalona o cap: 0,15% até 470.600 cotas integralizadas e 0,07% a partir de 750.000.

CRI-II, CRI-IV e CRI-VI seguem o mesmo desenho; suas lâminas não estão no acervo (n/l). O Termo da 177ª
VERT confirma os parâmetros de CRI-VI por outra via (Nível 1).

"Não consta" para o cap dos 10 maiores em CRI-III e CRI-V: a lâmina lista os critérios de forma
cumulativa e o item não aparece. Não equivale a ausência de limite em outro documento não obtido.

SEASONING: nenhuma das lâminas traz exigência de seasoning, MoB mínimo, safra performada ou histórico
de inadimplência do lote. Verificado nas três lâminas disponíveis. Classificado como ausente, não n/d,
porque a lista de critérios é cumulativa e exaustiva no documento.

CHERRY-PICKING: não testável com o acervo. Exigiria (i) tape por CCB do lote cedido em cada Termo de
Cessão e (ii) tape do que permaneceu no fundo, ambos na mesma data-base. Nenhum dos dois é público.
Declarado como ausência de evidência em qualquer direção — não como ausência do fenômeno.

DISTINÇÃO METODOLÓGICA: este slide compara CRITÉRIOS DE ELEGIBILIDADE (condição para o ativo entrar).
Não compara CONDIÇÕES DE CESSÃO (requisitos do negócio de cessão), COVENANTS do originador nem
características observadas da carteira. A elegibilidade define o universo elegível, não a composição
efetiva do pool adquirido.""")

# ══════════════════════════════════════════════ 5. Slide D — proteções e cash leakage
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "D · Proteção estrutural e extração de caixa")
titulo_conclusivo(
    sl, "Proteção contratual é relevante, mas redutível: subordinada já foi sacada em 7/7 fundos e\n"
        "R$ 1.076,2 mi saíram sem que a aderência aos testes seja verificável",
    "Distinção obrigatória: regra contratual, proteção observada na data-base, mecanismo de redução e evento efetivamente ocorrido.")

rotulo_secao(sl, Emu(1325880), "1 · Mecanismos que reduzem a proteção")
tab(sl, MARGEM, Emu(1508760), col_w(12), Emu(1188720),
    ["Veículo", "Regime corrente", "Vira sequencial em", "Trigger de desalavancagem",
     "Saída da subordinada", "Quem autoriza"],
    [["CRI-II", "Pró-rata condicionado", "mês 48 ou evento",
      "Atraso de estoque >15% em 3 datas; queda de 2 níveis de rating; cobertura desenquadrada em 2 datas consecutivas ou 4 em 12m",
      "Automática, até o Saldo Target", "Ninguém — regra automática"],
     ["CRI-VI", "Pró-rata condicionado", "mês 48 ou evento", "Idem CRI-II",
      "Automática, até o Saldo Target", "Ninguém — regra automática"],
     ["FIDC VI", "Pró-rata condicionado", "eventos; 6 datas aceleram",
      "Subordinação e cobertura pro forma; reserva de MTM; patamares 136,0% / 113,3% / 106,3%",
      "Amortização extraordinária sob testes", "75% da classe júnior"],
     ["FIDC VII", "Revolvência 12m, depois pró-rata", "eventos; 6 datas aceleram",
      "Idem FIDC VI", "Idem FIDC VI, com trava de 3m pós-venda", "75% da classe júnior"]],
    [0.09, 0.16, 0.13, 0.31, 0.17, 0.14], tam=8, num_cols=(),
    bold_cols=(0,), bold_cells=((0, 4), (1, 4), (0, 5), (1, 5), (2, 5), (3, 5)),
    h_linha=Emu(219456))

rotulo_secao(sl, Emu(2834640), "2 · Proteção contratual × observada × extração efetiva")
tab(sl, MARGEM, Emu(3017520), col_w(7), Emu(1554480),
    ["Fundo", "Piso (target)", "Subord. observada", "Folga (p.p.)", "Extração acumulada (R$ mi)"],
    [["I", "25%", "36,5%", "+11,5", "130,4"], ["II", "20%", "27,1%", "+7,1", "160,8"],
     ["III", "25%", "33,2%", "+8,2", "101,4"], ["IV", "n.a.", "0,0%", "n.a.", "438,9"],
     ["V", "20%", "20,9%", "+0,9", "34,4"], ["VI", "25%", "33,6%", "+8,6", "183,2"],
     ["VII", "25%", "26,0%", "+1,0", "7,7"]],
    [0.10, 0.18, 0.22, 0.18, 0.32], tam=8.5, num_cols=(1, 2, 3, 4), bold_cols=(0,),
    bold_cells=((3, 2), (3, 4), (4, 3), (6, 3)), h_linha=Emu(178308), destaque_linhas=(3,))

txt(sl, col_x(7), Emu(2834640), col_w(5), Emu(160020), "CASH LEAKAGE — O QUE JÁ OCORREU", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, col_x(7), Emu(3017520), col_w(5), Emu(1554480),
    "R$ 1.076,2 mi de mezanino e júnior amortizados nos sete fundos, mais R$ 45,4 mi nas séries "
    "subordinadas de três CRIs.\n\n"
    "Nos CRIs a saída não depende de aprovação: a série retida recebe em cada Data de Pagamento se as "
    "Razões de Cobertura estiverem enquadradas, e ao final recebe todo o remanescente do patrimônio "
    "separado como Prêmio Final.\n\n"
    "Nos FIDCs VI e VII o quórum é de 75% da classe júnior — que é a classe da originadora. Quem pede "
    "e quem se beneficia são a mesma parte.",
    tam=9, cor=CINZA_ESC, espaco=4)
txt(sl, MARGEM, Emu(4709160), UTIL_W, Emu(457200),
    "Limitação declarada: o informe mensal comprova o pagamento, mas não publica o teste de subordinação, cobertura e reservas "
    "na data de cada amortização extraordinária. Não há evidência documental suficiente para concluir sobre a aderência contratual "
    "das extrações já ocorridas.",
    tam=10, cor=PRETO, espaco=3)
fonte(sl, "Fonte: 2º aditamento ao TS da 2ª emissão Kanastra (cl. 6.5); TS da 177ª VERT (cl. 7.5); CVM — Informe Mensal FIDC e CRI; " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 5 (D)

REGIME E TRIGGERS (Nível 1):
- CRI-II: 2º Aditamento ao TS da 2ª emissão Kanastra, cl. 6.5.3 a 6.5.6. Pró-rata até o mês 47
  inclusive; sequencial a partir do mês 48 ou por Evento de Desalavancagem. Eventos: (i) Índice de
  Atraso de Estoque desenquadrado em 3 Datas de Verificação consecutivas — o índice não pode superar
  15%; (ii) rebaixamento de 2 níveis do rating dos CRI Super Sênior e/ou Sênior; (iii) não pagamento
  de remuneração ou amortização da 1ª série, cura de 5 dias úteis; (iv) não divulgação do Relatório da
  Emissão; (v) desenquadramento das Razões de Cobertura em 2 Datas de Pagamento consecutivas ou 4
  alternadas em 12 meses. Realavancagem: cl. 6.5.6; se o evento persistir 6 meses consecutivos, a
  operação permanece em sequencial permanentemente salvo aprovação em Assembleia Especial.
- CRI-VI: TS da 177ª VERT, cl. 7.5.5 a 7.5.7, mesma estrutura.

RAZÕES DE COBERTURA (triggers de alocação, não targets de subordinação):
- CRI-II: Super Sênior 159%, Sênior 123%, Mezanino 110%, Subordinada 105%.
- CRI-VI: Sênior 120,48%, Mezanino I 109,89%, Mezanino II 105,26%.
Fórmula: (Valor Presente dos DCI com PDD + Ativo Financeiro) / (Saldo Devedor das camadas até aquele nível).

SAÍDA DA SUBORDINADA:
- CRI-II cl. 6.5.1 (y) e (z): a 5ª série recebe remuneração e amortização se as QUATRO Razões de
  Cobertura estiverem enquadradas e o saldo estiver acima do Target. Item (bb): Prêmio Final — todo o
  remanescente do patrimônio separado à 5ª série após o resgate das quatro séries públicas.
- CRI-VI cl. 7.5.1 (r), (s) e (u): mesma mecânica com três razões.
- FIDCs VI e VII: quórum de 75% da classe júnior, testes de subordinação e cobertura pro forma, reserva
  de MTM e ausência de eventos; patamares 136,0% / 113,3% / 106,3%. Fonte: análise consolidada; os
  regulamentos não foram localizados (n/l).

PISO vs SUBORDINAÇÃO OBSERVADA: piso é TARGET contratual; subordinação observada é medida do informe
mensal em 31/07/2026. Folga = [Sub_NAV − piso × PL] / [1 − piso], cálculo próprio.

EXTRAÇÃO ACUMULADA (eventos observados, Informe Mensal FIDC): I 130,4; II 160,8; III 101,4; IV 438,9
(até 30/11/2025); V 34,4; VI 183,2; VII 7,7. Soma = 1.076,2. CRIs: CRI-I 22,5; CRI-II 7,8 + 6,8;
CRI-III 8,3. Soma = 45,4.

FIDC IV: sem piso contratual (n.a., não n/d) conforme análise consolidada; subordinação observada de
0,00% no informe de 31/07/2026.""")

# ══════════════════════════════════════════════ 6. Slide E — performance e safras
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "E · Performance da carteira")
titulo_conclusivo(
    sl, "Performance por safra é inconclusiva: só há curva por idade do veículo, e o take-out de jul/26\n"
        "altera o denominador do FIDC VI no mesmo mês em que a PDD salta de 16,3% para 48,8%",
    "A curva abaixo é MoB do veículo, não vintage de originação. Efeito de maturação e efeito de safra não são separáveis com o dado disponível.")

MOB = ["M0", "M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8", "M9", "M10",
       "M11", "M12", "M13", "M14", "M15", "M16", "M17", "M18", "M19"]
VI = [0, 0.00005, 0.0005, 0.0005, 0.049, 0.147, 0.163, 0.459, 0.672, 1.008, 1.150,
      1.842, 2.661, 2.767, 3.141, 4.000, 4.948, 5.868, 13.972, 16.259]
VII = [0, 0, 0.001, 0.007, 0.038, 0.078, 0.233, None, None, None, None,
       None, None, None, None, None, None, None, None, None]
gr = grafico(sl, XL_CHART_TYPE.LINE_MARKERS, MARGEM, Emu(1417320), col_w(7), Emu(2560320), MOB,
             [("FIDC VI — PDD / carteira", VI), ("FIDC VII — PDD / carteira", VII)],
             [LARANJA, PETROLEO], unidade_y="% da carteira")
for s in gr.series:
    s.format.line.width = Pt(2.25); s.smooth = False

txt(sl, col_x(7), Emu(1417320), col_w(5), Emu(160020), "O QUE A CURVA MOSTRA E O QUE NÃO MOSTRA", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, col_x(7), Emu(1600200), col_w(5), Emu(1005840),
    "O FIDC VI sai de 0% e chega a 16,3% da carteira em 20 meses de vida do veículo. Em jul/26 marca "
    "48,8% — mas nesse mesmo mês a carteira cai de R$ 399,3 mi para R$ 147,7 mi com a cessão para a "
    "177ª emissão.",
    tam=9, cor=CINZA_ESC, espaco=3)
txt(sl, col_x(7), Emu(2651760), col_w(5), Emu(1188720),
    "Deterioração real e efeito de denominador produzem o mesmo salto e não são distinguíveis sem a "
    "memória contábil da competência.\n\n"
    "O FIDC VII tem 6 competências e PDD de 0,23% — período curto demais para leitura de tendência.",
    tam=9, cor=CINZA_ESC, espaco=3)

rotulo_secao(sl, Emu(4114800), "Inventário de dados de performance")
tab(sl, MARGEM, Emu(4297680), col_w(12), Emu(914400),
    ["Disponível na data-base", "Fonte", "Ausente — bloqueia conclusão", "Classificação"],
    [["PDD/carteira e >90d/carteira por fundo", "Informe Mensal FIDC",
      "Vintage/cohort de originação, MoB por safra de CCB", "n/l"],
     ["Curva de PDD por idade do veículo (VI e VII)", "Informe Mensal FIDC",
      "NPL 30/60/90+, default, write-off, recuperação, CPR", "n/l"],
     ["Curva de PDD contratual por faixa de atraso", "Prospectos e termos",
      "Duration econômica observada, score, renda, LTV, mix PF/PJ efetivo", "n/d"]],
    [0.30, 0.16, 0.38, 0.10], tam=8.5, bold_cols=(), h_linha=Emu(178308))

txt(sl, MARGEM, Emu(5395464), UTIL_W, Emu(457200),
    "Não há evidência documental suficiente para concluir sobre performance relativa das safras, nem para separar efeito de vintage "
    "de efeito de seasoning. A PDD contratual por faixa de atraso é proteção prevista, não perda observada — e a razão PDD/>90d "
    "acima de 100% em 6 dos 7 fundos decorre de o numerador medir o saldo integral do contrato.",
    tam=10, cor=PRETO, espaco=3)
fonte(sl, "Fonte: CVM — Informe Mensal FIDC (PDD, >90d, MoB do veículo); prospectos definitivos (curva contratual de PDD); " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 6 (E)

CURVA MoB: série de PDD/carteira por mês ordinal desde a primeira competência observada de cada fundo,
extraída dos informes mensais e consolidada na análise de crédito de 21/08/2026. FIDC VI: 20 pontos,
de M0 a M19, de 0,00% a 16,26%. FIDC VII: 7 pontos, de M0 a M6, de 0,00% a 0,23%.

DISTINÇÃO CRÍTICA — o eixo é MoB DO VEÍCULO (idade do fundo), não vintage nem cohort de originação das
CCBs. Um fundo em ramp-up tem carteira jovem e PDD baixa por construção; a subida da curva reflete
simultaneamente maturação da carteira e mudança de composição por novas aquisições. Não é possível
isolar efeito de safra (mudança de underwriting) de efeito de seasoning (maturação) com este dado.

EFEITO DE DENOMINADOR — FIDC VI: PL 30/06/2026 R$ 437,4 mi, carteira R$ 399,3 mi; PL 31/07/2026
R$ 211,1 mi, carteira R$ 147,7 mi. Emissão da 177ª VERT em 21/07/2026, tendo o FIDC VI como cedente
(TS da 177ª, definição de Cedentes Fundos). PDD/carteira de 48,8% em 31/07/2026. Uma redução de 63% do
denominador eleva a razão mesmo sem qualquer aumento do numerador. O informe não decompõe a variação
de PL entre cessão, ajuste de valor e distribuição — limitação declarada.

PDD CONTRATUAL (proteção prevista, não perda observada):
- Curva inicial (CRI-I, CRI-II, FIDC IV, FIDC V): 0% até 14d; 1% 15–30; 3% 31–60; 10% 61–90; 30% 91–120;
  50% 121–150; 70% 151–180; 100% acima de 180. Fonte: definição de "Valor Presente dos Direitos
  Creditórios Imobiliários Líquido de PDD", Prospectos Definitivos da 1ª e 2ª emissões Kanastra.
- Curva posterior (CRI-III a CRI-VI, FIDC VI, FIDC VII): 0%/1,5%/5%/10%/37%/58%/78%/100%.

RAZÃO PDD/>90d ACIMA DE 100%: decorre de (i) a provisão incidir sobre o valor presente do recebível —
saldo remanescente da CCB — enquanto o >90d soma apenas parcelas vencidas; e (ii) o Efeito Vagão,
definido nos Prospectos como o arrasto da pior classificação de atraso entre todas as CCBs de um mesmo
devedor, esteja o título vencido ou a vencer. Não é indicador de cobertura de LGD.

INADIMPLÊNCIA POR SAFRA: as lâminas de CRI-III e CRI-V descrevem metodologia por safra — volumes não
liquidados com atraso superior a 90 dias, pelo saldo devedor do contrato, sobre o total originado na
safra, com separação entre perda bruta e perda líquida. A TABELA de resultados está no Prospecto, que
para CRI-V consta do acervo; os valores por safra não foram extraídos para este deck e permanecem como
item de diligência. Denominador distinto da PDD contábil — as duas métricas não se comparam.

FONTES NÃO DISPONÍVEIS: demonstrações financeiras auditadas e relatórios de auditoria (Nível 3) não
constam do acervo; relatórios de rating (Nível 4) constam apenas como notas atribuídas, sem os modelos
e premissas de stress.""")

# ══════════════════════════════════════════════ 7. Custo e dependência de funding
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "Funding")
titulo_conclusivo(
    sl, "Take-out passou a captar abaixo do warehouse: CRI sênior a DI+1,50% contra DI+2,00% do FIDC VII\n"
        "que o alimenta — e o CRI responde por 74,8% do funding emitido desde 2024",
    "Comparação restrita à família DI+, a única com série longa. Pré-fixado e IPCA+ não entram na mesma régua sem a curva de juros da data-base.")
c, r = L("22_custo_senior_timeline.csv")
di = [x for x in r if x[6] == "DI+"]
gr = grafico(sl, XL_CHART_TYPE.LINE_MARKERS, MARGEM, Emu(1417320), col_w(6), Emu(2377440),
             [f"{x[2]}\n{x[1][:7]}" for x in di],
             [("Spread sênior sobre o DI", [float(x[8]) for x in di])], [LARANJA],
             legenda=False, unidade_y="% a.a. sobre o DI", rotulos=True)
for s in gr.series:
    s.format.line.width = Pt(2.5); s.smooth = False

txt(sl, col_x(6), Emu(1417320), col_w(6), Emu(160020), "CAPTAÇÃO POR INSTRUMENTO E ANO (R$ mi)", tam=8.5, negrito=True, cor=CINZA_MED)
c2, r2 = L("25_captacao_por_ano.csv")
anos = sorted({x[0] for x in r2})
mp = {(x[0], x[1]): float(x[2]) for x in r2}
grafico(sl, XL_CHART_TYPE.COLUMN_STACKED, col_x(6), Emu(1600200), col_w(6), Emu(2194560), anos,
        [("FIDC", [mp.get((a, "FIDC (warehouse)"), 0.0) for a in anos]),
         ("CRI", [mp.get((a, "CRI (take-out)"), 0.0) for a in anos]),
         ("Debênture", [mp.get((a, "Debênture"), 0.0) for a in anos])],
        [PETROLEO, LARANJA, CINZA_MED], unidade_y="R$ mi")

rotulo_secao(sl, Emu(3931920), "Dependência estrutural")
tab(sl, MARGEM, Emu(4114800), col_w(12), Emu(731520),
    ["Indicador", "Valor", "Leitura de crédito"],
    [["Emitido em CRI desde jan/24", "R$ 3.670,6 mi", "74,8% de tudo que o programa emitiu no período"],
     ["PL somado dos 7 FIDCs (jul/26)", "R$ 1.234,6 mi", "Os warehouses sozinhos não sustentam o volume originado"],
     ["Spread sênior FIDC VI → VII", "DI+3,50% → DI+2,00%", "Redução de 150 bps em 15 meses no funding de warehouse"],
     ["Spread sênior CRI-VI (jul/26)", "DI+1,50%", "50 bps abaixo do FIDC VII, o warehouse que cede para ele"]],
    [0.28, 0.16, 0.56], tam=9, num_cols=(1,), bold_cols=(1,), h_linha=Emu(160020))
txt(sl, MARGEM, Emu(5029200), UTIL_W, Emu(365760),
    "Implicação: a economia do programa depende da continuidade do canal de securitização. Fechamento do mercado de CRI "
    "transferiria a necessidade de funding para os FIDCs a um custo hoje 50 bps superior, sobre uma base de PL que é um terço do volume já cedido.",
    tam=10, cor=PRETO, espaco=3)
fonte(sl, "Fonte: comunicados de bookbuilding e termos de securitização (taxas contratadas); CVM (registro de cotas); " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 7

SPREADS SÊNIOR, família DI+ (taxa contratada na data de emissão):
- FIDC III, 10/07/2023: CDI + 3,50% a.a.
- FIDC VI, 06/11/2024: DI + 3,50% a.a.
- FIDC VII, 13/01/2026: DI + 2,00% a.a. Rating AA+.br(sf) definitivo em 06/02/2026 (Moody's).
- CRI-VI Sênior A, 21/07/2026: 100% DI + 1,50% a.a., base 252. Fonte Nível 1: TS da 177ª VERT, cl. 6.2.
- CRI-VI Sênior B: 100% DI + 2,00% a.a.

Famílias não comparáveis diretamente: IPCA+ (FIDC I 6,75%; FIDC II 11%→8%→7%; FIDC V 10,00%; debênture
7,22%) e pré-fixado (CRI-I 11,51%; CRI-II 13,1926%; CRI-III 15,50%; CRI-IV 14,2216%; CRI-V 14,8064%).
A conversão para spread sobre DI exigiria a curva DI futura da B3 e a inflação implícita das NTN-B em
cada data-base — insumo não disponível. Classificado n/d, sem estimativa.

74,8%: CRI R$ 3.670,6 mi sobre o total emitido de 2024 a 2026 (CRI 3.670,6 + FIDC 1.235,0 no período) =
74,8%. Cálculo próprio sobre 25_captacao_por_ano.

O gráfico de captação por ano exclui uma tranche de R$ 50,0 mi — cota subordinada júnior do FIDC IV sem
data de registro no acervo. Soma das barras R$ 8.254,5 mi contra R$ 8.304,5 mi emitidos.

HIPÓTESE ANALÍTICA declarada: a leitura de que o fechamento do mercado de CRI transferiria a necessidade
de funding aos FIDCs é inferência de estrutura, não afirmação documentada. Não há no acervo plano de
contingência de funding da companhia.""")

# ══════════════════════════════════════════════ 8. Mismatch e prazo de exposição
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "Prazo")
titulo_conclusivo(
    sl, "Duration varia 2,7x entre operações; aprovar exposição pelo vencimento legal, não pela duration,\n"
        "porque toda amortização é condicionada a disponibilidade de caixa",
    "Duration divulgada é cenário-base e as lâminas só a qualificam para baixo. O risco de extensão está no contrato.")
grafico(sl, XL_CHART_TYPE.BAR_CLUSTERED, MARGEM, Emu(1417320), col_w(6), Emu(2011680),
        ["Prazo máx. do recebível", "WAM contratual do pool", "Duration sênior CRI-III",
         "Duration sênior CRI-I", "Duration sênior CRI-V"],
        [("Meses", [126.3, 65.7, 59.3, 40.9, 21.6])], [LARANJA],
        legenda=False, unidade_y="Meses", rotulos=True)
tab(sl, col_x(6), Emu(1417320), col_w(6), Emu(1005840),
    ["Operação", "Duration", "Venc. legal", "Gap ativo × duration", "Prazo sugerido"],
    [["CRI-I", "40,9 m", "jan/31 (84 m)", "85,4 m", "84 meses"],
     ["CRI-III", "59,3 m", "mai/30 (60 m)", "67,0 m", "60 meses"],
     ["CRI-V", "21,6 m", "mai/31 (60 m)", "104,7 m", "60 meses"]],
    [0.18, 0.16, 0.24, 0.22, 0.20], tam=9, num_cols=(1, 2, 3, 4), bold_cols=(0,),
    bold_cells=((0, 4), (1, 4), (2, 4)), h_linha=Emu(178308))
txt(sl, col_x(6), Emu(2560320), col_w(6), Emu(914400),
    "A mesma camada sênior tem duration de 59,3 meses em CRI-III e 21,6 meses em CRI-V. Tratar a faixa "
    "curta como padrão do programa seria erro material. CRI-II, CRI-IV e CRI-VI não publicam duration.",
    tam=9, cor=CINZA_ESC, espaco=3)

rotulo_secao(sl, Emu(3657600), "Quatro razões contratuais para o risco de extensão")
RAZOES = [("Condicionalidade", "Amortização de toda série condicionada a 'caso exista disponibilidade' e ao limite do caixa disponível."),
          ("Saldo Devedor Target", "A série amortiza até um saldo alvo, não até o cronograma do Anexo I. Pool abaixo do previsto não atinge o alvo."),
          ("Sequencial no mês 48", "Regime vira sequencial; camadas abaixo da sênior param de receber principal, o que acelera a sênior apenas se houver caixa."),
          ("Trigger de cobertura", "Desenquadramento em 2 datas consecutivas ou 4 alternadas em 12 meses aciona Evento de Desalavancagem.")]
for i, (t, corpo) in enumerate(RAZOES):
    x = MARGEM if i % 2 == 0 else col_x(6)
    y = Emu(3840480 + (i // 2) * 585216)
    barra_lateral(sl, x, y, Emu(36576), Emu(457200), LARANJA)
    txt(sl, Emu(x.emu + 109728), y, col_w(6), Emu(160020), t, tam=9.5, negrito=True, cor=PRETO)
    txt(sl, Emu(x.emu + 109728), Emu(y.emu + 173736), Emu(col_w(6).emu - 109728), Emu(310896),
        corpo, tam=8.5, cor=CINZA_ESC)
txt(sl, MARGEM, Emu(5120640), UTIL_W, Emu(320040),
    "Um limite aprovado pela duration venceria antes do papel em qualquer cenário de performance abaixo do previsto. "
    "Leitura de estrutura para o comitê; não constitui recomendação de investimento.",
    tam=10, cor=PRETO, espaco=3)
fonte(sl, "Fonte: lâminas de CRI-I, CRI-III e CRI-V (duration e vencimento); 2º aditamento ao TS da 2ª Kanastra e TS da 177ª VERT (cl. de amortização); " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 8

DURATION (aproximada, informada nas lâminas, em dias corridos, convertida a 30,4375 dias/mês):
- CRI-I: 1ª série 1.246 d (40,9 m); 2ª 1.295; 3ª 1.311; 4ª 1.146. Lâmina de 15/01/2024.
- CRI-III: 1ª e 2ª séries 1.806 d (59,3 m); 3ª 2.900; 4ª e 5ª 3.632. Lâmina de 22/04/2025.
- CRI-V: 1ª 659 d (21,6 m); 2ª 660; 3ª 713; 4ª 690; 5ª 689. Lâmina de 17/04/2026.
As lâminas qualificam a duration apenas para baixo: "observada a possibilidade de amortização
extraordinária, que pode fazer com que a duration seja menor". Não há qualificação para cima.

DURATION ECONÔMICA OBSERVADA: não disponível (n/l). O que consta é duration estimada de emissão. Não há
no acervo série de PU e amortização por série que permita recalcular duration realizada.

WAM contratual 2.000 dias (65,7 m) e prazo máximo por recebível 3.845 dias (126,3 m): critérios de
elegibilidade das lâminas. São LIMITES CONTRATUAIS. O WAM observado de cada pool não é publicado em
nenhuma operação (n/d).

Gap ativo × duration = prazo máximo do recebível (126,3 m) menos duration da sênior. Cálculo próprio.

RISCO DE EXTENSÃO — base contratual:
- Condicionalidade: redação repetida para todas as séries nas lâminas de CRI-III e CRI-V, seção
  Amortização/Juros: "Ressalvadas as hipóteses de Resgate Antecipado Obrigatório, caso exista
  disponibilidade, o saldo do Valor Nominal Unitário será amortizado nas Datas de Pagamento indicadas
  no cronograma do Anexo I, observadas as regras da Ordem de Alocação de Recursos".
- Saldo Devedor Target: TS da 2ª Kanastra cl. 6.5.1 (f); TS da 177ª VERT cl. 7.5.1.
- Mês 48: TS da 2ª Kanastra cl. 6.5.3 e 6.5.4; TS da 177ª VERT cl. 7.5.5 e 7.5.6.
- Trigger de cobertura: TS da 2ª Kanastra cl. 6.5.5 (v).

HIPÓTESE ANALÍTICA: o prazo sugerido de aprovação é igual ao vencimento legal da série sênior. É
recomendação de dimensionamento de limite, não previsão de prazo de recebimento.""")

# ══════════════════════════════════════════════ 9. Conclusão de crédito
sl = prs.slides.add_slide(VAZIO)
kicker(sl, "Conclusão de crédito")
titulo_conclusivo(
    sl, "Renovação condicionada: estrutura contratual sustenta a exposição, mas quatro lacunas de\n"
        "informação impedem conclusão sobre performance, economia da cessão e aderência das extrações",
    "Síntese por natureza de evidência: o que é proteção contratual, o que é observado e o que permanece sem resposta.")
caixa_sintese(sl, MARGEM, Emu(1417320), col_w(4), Emu(2377440), "Proteções", [
    ("Granularidade contratual.", "Cap individual de 0,07%–0,25% do patrimônio separado nos CRIs, contra 0,10%–20,00% nos warehouses."),
    ("WAM e prazo.", "2.000 dias de WAM e 3.845 dias por recebível em 6/6 operações de CRI."),
    ("Subordinação observada.", "20,9% a 36,5% em 6/7 fundos, acima de pisos de 20%–25%."),
    ("Cascata condicionada.", "Razões de cobertura de 105%–159% travam cada degrau; verificação em cada Data de Pagamento."),
    ("Retenção do originador.", "Série privada de 2,49%–3,00% em 6/6 CRIs e 2ª série subordinada na debênture."),
    ("Garantia real no ativo.", "Alienação fiduciária do equipamento contratada na CCB."),
], PETROLEO)
caixa_sintese(sl, col_x(4), Emu(1417320), col_w(4), Emu(2377440), "Pontos de atenção", [
    ("Extração de subordinada.", "R$ 1.076,2 mi sacados em 7/7 fundos; aderência aos testes não verificável."),
    ("FIDC IV.", "Subordinação observada de 0,0%, sem piso, PDD/carteira de 69,9%, R$ 438,9 mi já sacados."),
    ("Folga mínima ao piso.", "FIDC V +0,9 p.p. e FIDC VII +1,0 p.p.; o VII concentra 50,2% do PL."),
    ("Ausência de seasoning.", "Nenhuma das 6 operações exige safra performada ou MoB mínimo."),
    ("Basis/hedge.", "Em CRI-VI o diferencial do swap e a falta de reserva recaem sobre a série retida."),
    ("Recompra e substituição.", "Prevista no Contrato de Cessão; hipóteses e histórico não disponíveis."),
], LARANJA)
caixa_sintese(sl, col_x(8), Emu(1417320), col_w(4), Emu(2377440), "Sem evidência suficiente", [
    ("Performance por safra.", "Sem tape por CCB não se separa vintage de seasoning."),
    ("Preço de cessão.", "Mecanismo documentado; valor praticado só nos Termos de Cessão."),
    ("Cherry-picking.", "Não testável em nenhuma direção com o acervo atual."),
    ("Cascata de 4 CRIs.", "Termos de securitização de CRI-I, III, IV e V não localizados."),
    ("Critérios dos FIDCs.", "Regulamentos vigentes não localizados; parâmetros vêm de fonte secundária."),
    ("Custo all-in.", "Sem curva DI da data-base não há comparação entre famílias de indexador."),
], CINZA_ESC)

rotulo_secao(sl, Emu(3977640), "Encaminhamento")
tab(sl, MARGEM, Emu(4160520), col_w(12), Emu(1097280),
    ["Condição precedente", "Documento a obter", "Fonte", "Nível"],
    [["Concluir sobre performance e underwriting", "Tape por CCB: originação, safra, atraso, pré-pagamento", "Solfácil (originadora e servicer)", "5"],
     ["Validar aderência das extrações subordinadas", "Demonstrativos dos testes na data de cada amortização", "Administradores dos FIDCs", "1"],
     ["Fechar cascata e gatilhos de 4 operações", "Termos de securitização de CRI-I, III, IV e V", "CVM Fundos.NET; Kanastra; VERT", "1"],
     ["Confirmar critérios literais dos warehouses", "Regulamentos vigentes dos sete FIDCs", "CVM Fundos.NET", "1"]],
    [0.30, 0.36, 0.26, 0.08], tam=8.5, num_cols=(3,), bold_cols=(0,), h_linha=Emu(196596))
txt(sl, MARGEM, Emu(5395464), UTIL_W, Emu(320040),
    "Monitoramento sugerido: folga ao piso de subordinação por competência; PDD e >90d sobre carteira; ocorrência de amortização "
    "extraordinária da júnior; enquadramento das razões de cobertura; e decomposição contábil do PL do FIDC VI.",
    tam=10, cor=PRETO, espaco=3)
fonte(sl, "Fonte: consolidação das evidências dos slides A–E. Detalhamento documental no appendix. " + DB + ".")
notas(sl, """RASTREABILIDADE — SLIDE 9

Este slide consolida evidências já rastreadas nos slides 1 a 8. Referências cruzadas:
- Granularidade, WAM e prazo: slide 4 (C), critérios de elegibilidade das lâminas.
- Subordinação observada e pisos: slide 3 (B) e slide 5 (D), Informe Mensal FIDC 31/07/2026.
- Razões de cobertura: slide 5 (D), cl. 6.5.1 do TS da 2ª Kanastra e cl. 7.5.1 do TS da 177ª VERT.
- Retenção: séries privadas de 2,49% a 3,00% em 6/6 CRIs; cl. 3.3.3 da escritura da debênture prevê que
  a 2ª série seja subscrita exclusivamente pela Solfácil.
- Garantia real: lâminas de CRI-I, III e V, seção Garantias — "Não serão constituídas garantias no
  âmbito dos CRI diretamente. Não obstante, os Direitos Creditórios Imobiliários são garantidos por
  alienação fiduciária dos Equipamentos". A garantia é do ativo subjacente, não do título.
- Extração acumulada: slide 5 (D), Informe Mensal FIDC.
- Basis/hedge: cl. 15.14.6 item 7 do TS da 177ª — "Esse diferencial, seja positivo ou negativo, deverá
  impactar a Remuneração dos CRI da Quinta Série"; cl. 15.14.8 item 3 — cobertura extraordinária dos
  titulares da 5ª série se a Reserva de Caixa Derivativos for insuficiente. Reserva limitada a 1% do
  produto do Valor Presente dos DCI pela razão do percentual pós-fixado sobre o saldo total.
- Recompra e substituição: prevista nas lâminas de CRI-I, III e V como "compra e recompra obrigatória
  dos Direitos Creditórios Imobiliários Cedidos na ocorrência de qualquer das hipóteses de recompra das
  Cedentes e/ou Endossante Inicial e/ou Solfácil, nos termos do Contrato de Cessão". Os Contratos de
  Cessão não estão no acervo — hipóteses e histórico classificados n/l.

CLASSIFICAÇÃO DAS LACUNAS: "sem evidência suficiente" reúne itens n/d (não divulgado pelo emissor) e
n/l (deveria existir e não foi localizado). O detalhamento por item está no appendix.

NÍVEL na tabela de encaminhamento refere-se à hierarquia de fontes: 1 = documentação jurídica primária;
5 = informação fornecida pela administração.""")

# ══════════════════════════════════════════════ Divisória do appendix
sl = prs.slides.add_slide(VAZIO)
barra_lateral(sl, Emu(0), Emu(0), Emu(91440), SL_H, LARANJA)
txt(sl, MARGEM, Emu(2926080), col_w(8), Emu(457200), "APPENDIX", tam=28, negrito=True, cor=PRETO)
txt(sl, MARGEM, Emu(3474720), col_w(8), Emu(731520),
    "Detalhamento documental, rastreabilidade e material de suporte. As conclusões de crédito estão no corpo principal, "
    "slides 1 a 9. O appendix preserva a informação necessária para verificação, sem competir com a leitura de decisão.",
    tam=11, cor=CINZA_ESC, espaco=3)
notas(sl, "Divisória. O appendix reúne o detalhamento movido do corpo principal: registro por série, ISIN, "
          "cronologia de emissões, cascatas literais, critérios completos e o inventário de lacunas.")


def app(titulo, subtitulo=None, kick="Appendix"):
    s = prs.slides.add_slide(VAZIO)
    kicker(s, kick)
    txt(s, MARGEM, Y_TITULO, col_w(12), Emu(320040), titulo, tam=15, negrito=True, cor=PRETO)
    if subtitulo:
        txt(s, MARGEM, Emu(731520), col_w(12), Emu(228600), subtitulo, tam=9.5, cor=CINZA_ESC)
    return s


# A1 — timeline
sl = app("A1 · Cronologia das emissões do programa",
         "Registro de cotas de FIDC, emissões de CRI e a debênture, em ordem cronológica")
c, r = L("26_timeline_consolidada.csv")
tab(sl, MARGEM, Emu(1097280), col_w(7), Emu(4754880),
    ["#", "Data", "Veículo", "Evento", "R$ mi"],
    [[x[0], x[1], x[2], x[4][:52], x[5]] for x in r],
    [0.05, 0.13, 0.12, 0.55, 0.15], tam=7.5, num_cols=(0, 4), bold_cols=(2,), h_linha=Emu(146304))
c2, r2 = L("26b_vencimentos_por_ano.csv")
txt(sl, col_x(7), Emu(1097280), col_w(5), Emu(160020), "VENCIMENTOS LEGAIS POR ANO — CRI E DEBÊNTURE", tam=8.5, negrito=True, cor=CINZA_MED)
grafico(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, col_x(7), Emu(1280160), col_w(5), Emu(2194560),
        [x[0] for x in r2], [("R$ mi", [float(x[1]) for x in r2])], [LARANJA],
        legenda=False, unidade_y="R$ mi", rotulos=False)
txt(sl, col_x(7), Emu(3566160), col_w(5), Emu(1188720),
    "70,7% do valor nominal de CRI vence entre 2029 e 2031, com pico de R$ 1.239,4 mi em 2031.\n\n"
    "As séries amortizam mensalmente ao longo da vida; o vencimento legal é o limite de exposição, não a expectativa de recebimento.\n\n"
    "Os FIDCs não entram: as cotas não têm data de vencimento publicada no acervo (n/d).",
    tam=9, cor=CINZA_ESC, espaco=3)
fonte(sl, "Fonte: CVM (registro de cotas e ofertas); prospectos e termos de securitização; escritura da debênture; " + DB + ".")
notas(sl, "Datas de registro de cotas de FIDC conforme CVM e consolidação da análise de 21/08/2026. "
          "Datas de emissão de CRI conforme prospectos e termos: CRI-I 15/01/2024; CRI-II 25/06/2024 "
          "(bookbuilding em 20/06/2024 — eventos distintos); CRI-III 28/05/2025; CRI-IV 28/09/2025 "
          "(bookbuilding 26/09); CRI-V 20/05/2026; CRI-VI 21/07/2026 (TS de 20/07/2026). "
          "Debênture: escritura de 18/02/2022. Vencimentos por ano somam o valor nominal de emissão das "
          "34 séries de CRI e das 2 séries de debênture.")

# A2 — de-para
sl = app("A2 · De-para FIDC → CRI e situação corrente de cada veículo",
         "Vínculo de cessão por operação, com a qualidade da evidência documental")
c, r = L("27_depara_fidc_cri.csv")
EVID = {"CRI-I": "Nível 1 — Prospecto Definitivo nomeia os dois fundos",
        "CRI-II": "Nível 1 — definição de Cedentes do Prospecto",
        "CRI-III": "Nível 6 — lâmina não nomeia fundo; TS não localizado",
        "CRI-IV": "Nível 6 — TS da 4ª emissão não localizado",
        "CRI-V": "Nível 6 — lâmina cita Cedente Fundo sem nomeá-lo",
        "CRI-VI": "Nível 1 — TS define Cedentes Fundos: VI e VII"}
tab(sl, MARGEM, Emu(1097280), col_w(12), Emu(1371600),
    ["Operação", "Emissão", "Cedentes documentados", "Cessão direta", "Situação", "Qualidade da evidência"],
    [[x[1], x[2], x[3], x[4], x[5], EVID[x[0]]] for x in r],
    [0.10, 0.09, 0.17, 0.10, 0.16, 0.38], tam=7.5, bold_cols=(0,), h_linha=Emu(228600))
c2, r2 = L("27b_fidc_status.csv")
txt(sl, MARGEM, Emu(2651760), col_w(12), Emu(160020), "SITUAÇÃO DOS SETE WAREHOUSES", tam=8.5, negrito=True, cor=CINZA_MED)
tab(sl, MARGEM, Emu(2834640), col_w(12), Emu(1554480),
    ["Fundo", "Nome oficial", "CNPJ", "Cedeu para", "Situação hoje", "PL (R$ mi)", "Carteira (R$ mi)"],
    [[x[0].replace("FIDC-", "FIDC "), x[1][:46], x[2], x[3], x[4][:34], x[5], x[6]] for x in r2],
    [0.08, 0.28, 0.14, 0.16, 0.18, 0.08, 0.08], tam=7.5, num_cols=(5, 6), bold_cols=(0,), h_linha=Emu(196596))
txt(sl, MARGEM, Emu(4343400), UTIL_W, Emu(457200),
    "Só CRI-I, CRI-II e CRI-VI têm o cedente nomeado em documento de Nível 1. Os vínculos de CRI-III, CRI-IV e CRI-V "
    "vêm de consolidação secundária — os termos de securitização dessas operações não foram localizados.",
    tam=9.5, cor=PRETO, espaco=3)
fonte(sl, "Fonte: prospectos definitivos da 1ª e 2ª emissões Kanastra; TS da 177ª VERT; CVM — cadastro e informes; " + DB + ".")
notas(sl, "Nível 1: Prospecto Definitivo da 1ª emissão Kanastra, seção 11.1.2, nomeia GREEN SOLFÁCIL II FIDC "
          "(42.462.306/0001-00) e GREEN SOLFÁCIL IV FIDC (44.909.456/0001-44), ambos administrados pelo Banco "
          "Genial (45.246.410/0001-55), como representantes de mais de 10% dos direitos creditórios cedidos. "
          "O Prospecto da 2ª emissão define 'Cedentes' como o FIDC II e o FIDC IV em conjunto. "
          "O TS da 177ª VERT define 'Cedentes Fundos' como IS GREEN SOLFÁCIL VI FIDC RL (57.028.406/0001-08) e "
          "SOLFÁCIL CRÉDITO PESSOAL VII FIDC RL (63.505.455/0001-89), e 'Cedentes' como estes mais a Solfácil "
          "(31.931.053/0001-50). A lâmina de CRI-V confirma a existência de um 'Cedente Fundo' com Gestora e "
          "custodiante próprios sem nomeá-lo — evidência indireta. Vínculos de CRI-III, CRI-IV e CRI-V: análise "
          "de 21/08/2026, classificada como Nível 6 consolidando Nível 2.")

# A3 / A4 — séries de CRI
c, r = L("02_series.csv")
IX = {k: i for i, k in enumerate(c)}
cri = [x for x in r if x[IX["veiculo_id"]].startswith("CRI")]
NOME = {"CRI-I": "KAN 1ª", "CRI-II": "KAN 2ª", "CRI-III": "KAN 3ª",
        "CRI-IV": "KAN 4ª", "CRI-V": "VERT 174ª", "CRI-VI": "VERT 177ª"}
for parte, alvos in enumerate([["CRI-I", "CRI-II", "CRI-III"], ["CRI-IV", "CRI-V", "CRI-VI"]], start=3):
    sub = [x for x in cri if x[IX["veiculo_id"]] in alvos]
    sl = app(f"A{parte} · As 34 séries de CRI — {', '.join(NOME[a] for a in alvos)}",
             "Montante, indexador, taxa contratada, ISIN, vencimento e rating por série")
    tab(sl, MARGEM, Emu(1097280), col_w(12), Emu(4663440),
        ["Operação", "Série", "Camada", "ISIN", "R$ mi", "Indexador", "Taxa contratada",
         "Vencimento", "Rating", "Colocação"],
        [[NOME[x[IX["veiculo_id"]]], x[IX["serie"]], x[IX["camada"]], x[IX["isin"]],
          x[IX["montante_reportado_deck_Rmi"]], x[IX["indexador"]], x[IX["taxa_contratada"]][:34],
          x[IX["data_vencimento"]], x[IX["rating_nota"]][:26], x[IX["colocacao"]]] for x in sub],
        [0.09, 0.05, 0.12, 0.13, 0.07, 0.08, 0.17, 0.10, 0.11, 0.08],
        tam=7.5, num_cols=(4,), bold_cols=(0,), h_linha=Emu(228600))
    fonte(sl, "Fonte: prospectos definitivos, lâminas, comunicados de bookbuilding e termos de securitização das respectivas emissões; " + DB + ".")
    notas(sl, "Montantes: CRI-I e CRI-II conforme anúncio de encerramento e comunicado de bookbuilding; CRI-III "
              "conforme Prospecto Definitivo de 06/06/2025 (R$ 727,5 mi, confirmando exercício integral do lote "
              "adicional de 25% sobre lote base de R$ 582,0 mi); CRI-IV conforme anúncio de início de 29/09/2025; "
              "CRI-V conforme Prospecto Definitivo republicado (R$ 456,481 mi públicos); CRI-VI conforme TS da 177ª, "
              "cl. 5.6.4 e 5.6.6 (R$ 647,059 mi). As séries privadas (Subordinado Jr.) não constam dos anúncios de "
              "encerramento, que registram apenas a oferta pública — por isso os anúncios contam 4, 4 e 6 séries "
              "onde os termos contam 5, 5 e 7. Divergência de perímetro, não de fato.")

# A5 — cotas de FIDC
sl = app("A5 · As 34 classes de cotas dos sete FIDCs",
         "Emissão, montante, remuneração e rating por classe. Dimensão distinta das séries de CRI — não se somam.")
fidc = [x for x in r if x[IX["veiculo_id"]].startswith("FIDC")]
tab(sl, MARGEM, Emu(1097280), col_w(12), Emu(4663440),
    ["Fundo", "Camada", "Classe / emissão", "Registro", "Emitido", "Subscrito", "Remuneração", "Rating"],
    [[x[IX["veiculo_id"]].replace("FIDC-", ""), x[IX["camada"]], x[IX["serie"]][:26],
      x[IX["data_emissao"]], x[IX["montante_ofertado_Rmi"]], x[IX["montante_subscrito_Rmi"]],
      x[IX["taxa_contratada"]][:46], x[IX["rating_nota"]][:24]] for x in fidc],
    [0.06, 0.11, 0.16, 0.09, 0.07, 0.08, 0.28, 0.15],
    tam=7, num_cols=(4, 5), bold_cols=(0,), h_linha=Emu(128016))
fonte(sl, "Fonte: CVM — registro de ofertas de cotas; consolidação da análise de crédito de 21/08/2026; " + DB + ".")
notas(sl, "As 34 classes de cotas de FIDC e as 34 séries de CRI são coincidência numérica sem relação causal. "
          "São dimensões distintas e não devem ser somadas. Emitido = montante registrado; subscrito = "
          "efetivamente colocado. Onde subscrito < emitido, a diferença não foi colocada — não é n/d. "
          "Remuneração de FIDC IV sênior A, B e C não divulgada (n/d). Ratings: Austin Rating para os fundos "
          "I, II, IV, V e VI; Fitch para o III; Moody's para o VII (AA+.br(sf) definitivo em 06/02/2026).")

# A6 — waterfall CRI-II
sl = app("A6 · Ordem de alocação de recursos — CRI-II (2ª emissão Kanastra)",
         "Cascata literal, 24 degraus no regime pró-rata condicionado. Fonte de Nível 1.")
c2, r2 = L("06b_waterfall_degraus.csv")
pro = [x for x in r2 if x[1] == "Pró-rata"][:24]
meta = len(pro) // 2
tab(sl, MARGEM, Emu(1097280), col_w(6), Emu(4297680), ["#", "Item", "Degrau"],
    [[x[2], x[3], x[4][:76]] for x in pro[:meta]],
    [0.06, 0.10, 0.84], tam=7, bold_cols=(0,), h_linha=Emu(146304))
tab(sl, col_x(6), Emu(1097280), col_w(6), Emu(4297680), ["#", "Item", "Degrau"],
    [[x[2], x[3], x[4][:76]] for x in pro[meta:]],
    [0.06, 0.10, 0.84], tam=7, bold_cols=(0,), h_linha=Emu(146304))
fonte(sl, "Fonte: 2º aditamento ao termo de securitização da 2ª emissão Kanastra, cl. 6.5.1, registrado na JUCEMG sob nº 12803230 em 10/06/2025.")
notas(sl, "Cláusula 6.5.1, itens (a) a (cc), do 2º Aditamento ao Termo de Securitização, consolidado e registrado. "
          "A cl. 6.5.2 traz a ordem sequencial, com 21 itens de (a) a (aa): nela desaparecem as condições de "
          "Razão de Cobertura e o Saldo Devedor Target, e cada série é amortizada até 98% do valor nominal antes "
          "de a seguinte receber principal. Regime pró-rata vigora até o mês 47 inclusive (cl. 6.5.3); a partir "
          "do mês 48 ou de Evento de Desalavancagem passa a sequencial (cl. 6.5.4). "
          "Razões de Cobertura: Super Sênior 159%, Sênior 123%, Mezanino 110%, Subordinada 105%. "
          "Índice de Atraso de Estoque limitado a 15%.")

# A7 — waterfall comparado
sl = app("A7 · CRI-II × CRI-VI: as duas cascatas documentadas, degrau a degrau",
         "Mesma espinha estrutural. A divergência material está no tratamento do derivativo.")
c3, r3 = L("29_waterfall_comparado.csv")
tab(sl, MARGEM, Emu(1097280), col_w(12), Emu(2560320),
    ["Degrau", "CRI-II (2ª Kanastra)", "CRI-VI (177ª VERT)", "Divergência"],
    [[x[0], x[2][:60], x[3][:60], x[4][:44]] for x in r3],
    [0.20, 0.28, 0.28, 0.24], tam=7.5, bold_cols=(0,), h_linha=Emu(228600),
    destaque_linhas=(1, 5))
c4, r4 = L("29b_razoes_de_cobertura.csv")
txt(sl, MARGEM, Emu(3840480), col_w(5), Emu(160020), "RAZÕES DE COBERTURA CONTRATUAIS", tam=8.5, negrito=True, cor=CINZA_MED)
tab(sl, MARGEM, Emu(4023360), col_w(5), Emu(1371600), ["Operação", "Camada", "Mínimo"],
    [[x[0], x[1], f"{x[2]}%"] for x in r4], [0.24, 0.46, 0.30], tam=8, num_cols=(2,),
    bold_cols=(0,), h_linha=Emu(164592))
txt(sl, col_x(5), Emu(3840480), col_w(7), Emu(160020), "O QUE A DIVERGÊNCIA SIGNIFICA", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, col_x(5), Emu(4023360), col_w(7), Emu(1371600),
    "CRI-VI introduz contratos de swap para cobrir o descasamento entre o ativo pré-fixado e as séries pós-fixadas. "
    "O ajuste do derivativo é o item (c) da cascata — pago antes da remuneração sênior — e há Reserva de Caixa "
    "Derivativo constituída antes do pagamento à série retida, limitada a 1% do valor presente da parcela pós-fixada.\n\n"
    "O diferencial de taxa, positivo ou negativo, impacta a remuneração da 5ª série. Se a reserva for insuficiente, "
    "há cobertura extraordinária dos titulares dessa mesma série — que é a retida pela originadora.\n\n"
    "Efeito de crédito: o basis risk do hedge está alocado no originador, não nas séries públicas.",
    tam=9, cor=CINZA_ESC, espaco=4)
fonte(sl, "Fonte: cl. 6.5.1 e 6.5.2 do 2º aditamento ao TS da 2ª Kanastra; cl. 7.5.1, 7.5.2, 7.5.3, 15.14.6 e 15.14.8 do TS da 177ª VERT.")
notas(sl, "CRI-VI tem três Razões de Cobertura em vez de quatro porque as séries Sênior A e Sênior B são pari passu. "
          "A cl. 7.5.3 determina que, na falta de caixa, os pagamentos às duas séries sênior sejam feitos na proporção "
          "necessária para manter a razão base entre os respectivos Saldos Devedores Target. "
          "A cl. 7.5.4 prevê Incorporação: remuneração devida e não paga é incorporada ao saldo devedor da respectiva "
          "classe, subclasse ou série — mecanismo de PIK que não existe explicitado no TS da 2ª Kanastra. "
          "Prêmio Final: cl. 6.5.1 (bb) em CRI-II e cl. 7.5.1 (u) em CRI-VI, ambos destinando o remanescente integral "
          "do patrimônio separado à 5ª série após o resgate das séries públicas.")

# A8 — ranking
sl = app("A8 · Índice de permissividade dos mandatos",
         "Agregação de quatro parâmetros contratuais documentados. Mede o que o documento permite, não o que a carteira tem.")
c5, r5 = L("28_ranking_permissividade.csv")
K5 = {k: i for i, k in enumerate(c5)}
tab(sl, MARGEM, Emu(1097280), col_w(7), Emu(3383280),
    ["#", "Veículo", "Tipo", "Cap", "WAM", "Prazo", "Ticket PJ", "Índice", "Eixos"],
    [[x[K5["posicao"]], x[K5["veiculo_id"]], x[K5["instrumento"]], x[K5["cap_individual_pct"]],
      x[K5["wam_max_dias"]], x[K5["prazo_max_dias"]], x[K5["ticket_max_PJ_R"]],
      x[K5["indice_permissividade"]], x[K5["eixos_avaliados"]]] for x in r5],
    [0.06, 0.15, 0.10, 0.12, 0.11, 0.11, 0.15, 0.11, 0.09],
    tam=7.5, num_cols=(0, 3, 4, 5, 6, 7, 8), bold_cols=(1,),
    bold_cells=((12, 7), (12, 3)), h_linha=Emu(219456), destaque_linhas=(12,))
txt(sl, col_x(7), Emu(1097280), col_w(5), Emu(160020), "MÉTODO E RESSALVAS", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, col_x(7), Emu(1280160), col_w(5), Emu(3200400),
    "Cada eixo é normalizado entre o extremo restritivo (0) e o permissivo (100) do universo observado. "
    "O eixo de concentração usa escala logarítmica porque o cap varia de 0,07% a 20,00% — 285 vezes — e "
    "em escala linear todos os CRIs colapsariam em zero.\n\n"
    "O índice é a média simples dos eixos com dado disponível. Veículos com 3 eixos não são estritamente "
    "comparáveis aos com 4: CRI-II, CRI-IV e CRI-VI não publicam ticket máximo PJ; os FIDCs VI e VII não "
    "publicam cap por devedor; o FIDC IV tem apenas o eixo de concentração.\n\n"
    "Classificação: hipótese analítica com método declarado. Não é rating nem medida de perda esperada.",
    tam=9, cor=CINZA_ESC, espaco=4)
fonte(sl, "Fonte: critérios de elegibilidade das lâminas de CRI-I, III e V; TS da 177ª VERT; parâmetros dos FIDCs via análise consolidada de 21/08/2026.")
notas(sl, "Fórmula por eixo: nota = 100 × (valor − mínimo) / (máximo − mínimo), com log10 aplicado a valor, mínimo "
          "e máximo no eixo de concentração. Índice = média das notas disponíveis. Extremos do universo: cap de "
          "0,07% (CRI-V, pool maduro) a 20,00% (FIDC IV); WAM de 2.000 d (CRIs e FIDC III) a 2.400 d (FIDCs II, VI "
          "e VII); prazo de 3.836 d (FIDCs I, VI, VII) a 4.760 d (FIDC V); ticket PJ de R$ 500 mil (FIDC II) a "
          "R$ 700 mil (CRI-III, CRI-V, FIDCs V, VI e VII). "
          "O resultado ordena mandatos e não deve ser lido como ordenação de risco efetivo: um mandato restritivo "
          "com carteira deteriorada continua deteriorado. O FIDC IV aparece em 100 com um único eixo avaliado — "
          "leitura frágil em precisão, robusta em direção.")

# A9 — preço de aquisição
sl = app("A9 · Preço de aquisição dos créditos: mecanismo documentado, valor praticado não disponível",
         "FIDC precifica por percentual do saldo contábil; CRI por valor presente descontado a uma taxa mínima.")
c6, r6 = L("31_preco_de_aquisicao.csv")
tab(sl, MARGEM, Emu(1097280), col_w(12), Emu(2560320),
    ["Veículo", "Tipo", "Mecanismo de preço", "Teto contratual", "Preço praticado", "Quem define"],
    [[x[0], x[1], x[2][:56], x[3][:64], x[4][:34], x[5][:30]] for x in r6],
    [0.09, 0.08, 0.24, 0.28, 0.17, 0.14], tam=7, bold_cols=(0,), h_linha=Emu(196596))
txt(sl, MARGEM, Emu(3840480), col_w(6), Emu(160020), "FÓRMULA DO CRI (NÍVEL 1)", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, MARGEM, Emu(4023360), col_w(6), Emu(1188720),
    "Preço de Aquisição  =  Σᵢ  Parcela do Direito Creditórioᵢ  ÷  (1 + Taxa de Retorno)^(N/252)\n\n"
    "N = dias úteis entre o vencimento de cada parcela vincenda e a Data de Aquisição e Pagamento.\n"
    "O valor bruto da fórmula é o máximo. A elegibilidade exige Taxa de Retorno ≥ Taxa Média Mínima de "
    "Retorno — quanto maior a taxa exigida, menor o preço pago.",
    tam=9, cor=CINZA_ESC, espaco=3)
txt(sl, col_x(6), Emu(3840480), col_w(6), Emu(160020), "EVOLUÇÃO DO TETO NOS WAREHOUSES", tam=8.5, negrito=True, cor=CINZA_MED)
tab(sl, col_x(6), Emu(4023360), col_w(6), Emu(1005840),
    ["Fundo", "Teto", "Quem atesta a elegibilidade"],
    [["FIDC I", "100,4%", "Gestor"], ["FIDC II / V", "100,5%", "Gestor"],
     ["FIDC VI", "101,0%", "Gestora"], ["FIDC VII", "104,0%", "Endossantes e originador"]],
    [0.24, 0.16, 0.60], tam=8.5, num_cols=(1,), bold_cols=(0,), bold_cells=((3, 1), (3, 2)),
    h_linha=Emu(196596), destaque_linhas=(3,))
fonte(sl, "Fonte: definição de Preço de Aquisição do TS da 177ª VERT; prospectos definitivos da 1ª e 2ª Kanastra (taxa mínima); tetos dos FIDCs via análise consolidada.")
notas(sl, "Fórmula literal do TS da 177ª VERT, definição de 'Preço de Aquisição': valor presente das parcelas "
          "vincendas descontadas pela Taxa de Retorno em base 252 dias úteis, evidenciado em cada Termo de Cessão. "
          "Taxa Média Mínima de Retorno documentada: 21,5% a.a. em CRI-I e 21,0% a.a. em CRI-II (Prospectos "
          "Definitivos). Para CRI-III e CRI-V a lâmina exige a comparação com a Taxa Média Mínima sem publicar o "
          "percentual — classificado n/d. "
          "No FIDC VII muda também o responsável por atestar a condição: deixa de ser a gestora e passa a ser "
          "endossantes e originador, simultaneamente ao aumento do teto de 101% para 104%. É mudança de governança "
          "de controle, relevante para a leitura de conflito de interesse. "
          "Preço de Recompra, no TS da 177ª, é definido como o Valor Presente dos DCI na data de recompra. "
          "Debênture: relatório do agente fiduciário registra valor de aquisição acumulado de R$ 17.212.382,98 "
          "igual ao valor nominal total adquirido — aquisição ao par, único preço observável no acervo.")

# A10 — histórico de saques
sl = app("A10 · Histórico de amortização de subordinada por veículo",
         "Eventos efetivamente observados nos informes mensais. A aderência aos testes contratuais não é verificável.")
c7, r7 = L("33_historico_saques_subordinada.csv")
tab(sl, MARGEM, Emu(1097280), col_w(7), Emu(4297680),
    ["Veículo", "Camada", "Primeira", "Última", "Meses", "Total (R$ mi)", "Máx. mensal"],
    [[x[0].replace("FIDC-", "FIDC ").replace("CRI-", "CRI "), x[2][:18], x[3], x[4], x[5], x[6], x[7]] for x in r7],
    [0.14, 0.20, 0.14, 0.14, 0.09, 0.15, 0.14], tam=7.5, num_cols=(4, 5, 6), bold_cols=(0,),
    bold_cells=((6, 5), (7, 5), (7, 6)), h_linha=Emu(219456), destaque_linhas=(6, 7))
txt(sl, col_x(7), Emu(1097280), col_w(5), Emu(160020), "OS DOIS CASOS QUE PEDEM EXPLICAÇÃO", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, col_x(7), Emu(1280160), col_w(5), Emu(1737360),
    "FIDC IV — R$ 438,9 mi de mezanino e júnior, com um mês isolado de R$ 151,4 mi. Maior volume do "
    "programa, saído do fundo com o mandato mais permissivo. Hoje sem first loss e com PDD/carteira de 69,9%.\n\n"
    "FIDC VII — R$ 7,7 mi de júnior em jul/26, no mesmo mês do take-out da 177ª emissão e no primeiro mês "
    "em que o fundo cedeu carteira.",
    tam=9, cor=CINZA_ESC, espaco=4)
txt(sl, col_x(7), Emu(3108960), col_w(5), Emu(160020), "LIMITAÇÃO", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, col_x(7), Emu(3291840), col_w(5), Emu(1188720),
    "O informe mensal comprova o pagamento, mas não publica o teste de subordinação, cobertura e reservas "
    "na data de cada amortização extraordinária.\n\n"
    "Não há evidência documental suficiente para concluir sobre a aderência contratual das extrações "
    "já ocorridas.",
    tam=9, cor=CINZA_ESC, espaco=4)
fonte(sl, "Fonte: CVM — Informe Mensal FIDC e Informe Mensal CRI, até 31/07/2026; escritura e relatório do agente fiduciário da debênture.")
notas(sl, "Valores acumulados de principal de mezanino e classe júnior/subordinada. FIDC IV com última ocorrência "
          "em 30/11/2025 (júnior) e 31/08/2025 (mezanino). Total do programa: R$ 1.076,2 mi nos FIDCs e R$ 45,4 mi "
          "nas séries subordinadas de CRI-I, CRI-II e CRI-III. "
          "Debênture: a 2ª série, subordinada à 1ª pela cl. 3.3.2 da escritura e subscrita exclusivamente pela "
          "Solfácil pela cl. 3.3.3, registra R$ 0,00 distribuídos aos investidores em 31/07/2026 no relatório do "
          "agente fiduciário — consistente com a subordinação integral. É o único caso do programa em que a "
          "subordinada comprovadamente não recebeu. "
          "Distinção: os valores são AMORTIZAÇÃO DE PRINCIPAL observada, não distribuição de resultado nem "
          "pagamento de remuneração.")

# A11 — pontos em aberto e hierarquia de fontes
sl = app("A11 · Lacunas de informação e hierarquia documental aplicada",
         "Classificação por natureza da ausência e nível da fonte que a fecharia.")
c8, r8 = L("34_pontos_em_aberto.csv")
tab(sl, MARGEM, Emu(1097280), col_w(12), Emu(2926080),
    ["#", "Pergunta de crédito", "Documento que fecharia", "Fonte", "Impacto na decisão"],
    [[x[0], x[1][:52], x[2][:56], x[3][:40], x[4][:56]] for x in r8],
    [0.04, 0.22, 0.24, 0.18, 0.32], tam=7, num_cols=(0,), bold_cols=(1,), h_linha=Emu(196596))
txt(sl, MARGEM, Emu(4206240), col_w(6), Emu(160020), "HIERARQUIA DE FONTES APLICADA", tam=8.5, negrito=True, cor=CINZA_MED)
tab(sl, MARGEM, Emu(4389120), col_w(6), Emu(1188720),
    ["Nível", "Fonte", "Uso neste deck"],
    [["1", "Termos de securitização e escrituras", "CRI-II, CRI-VI e debênture"],
     ["1", "Prospectos definitivos e lâminas", "CRI-I, III, IV e V"],
     ["2", "CVM — informes mensais e ofertas", "PL, carteira, PDD, subordinação"],
     ["3", "Demonstrações auditadas", "não disponível (n/l)"],
     ["4", "Relatórios de rating", "notas apenas; modelos n/l"],
     ["6", "Research e análise consolidada", "parâmetros de FIDC, com ressalva"]],
    [0.10, 0.46, 0.44], tam=8, bold_cols=(0,), h_linha=Emu(164592), destaque_linhas=(3, 4))
txt(sl, col_x(6), Emu(4206240), col_w(6), Emu(160020), "CONVENÇÃO DE AUSÊNCIA", tam=8.5, negrito=True, cor=CINZA_MED)
txt(sl, col_x(6), Emu(4389120), col_w(6), Emu(1188720),
    "n/d — não divulgado: o emissor não publica o campo. Exemplos: cap por devedor dos FIDCs VI e VII; "
    "WAM observado dos pools; preço praticado em cada cessão.\n\n"
    "n/l — não localizado: deveria existir e não foi encontrado no acervo analisado. Exemplos: regulamentos "
    "vigentes dos sete FIDCs; termos de securitização de CRI-I, III, IV e V; contratos de cessão; "
    "demonstrações auditadas.\n\n"
    "n.a. — não aplicável. Exemplo: piso de subordinação do FIDC IV, que não possui piso contratual.\n\n"
    "Nenhuma ausência foi convertida em zero.",
    tam=8.5, cor=CINZA_ESC, espaco=3)
fonte(sl, "Fonte: inventário de fontes do estudo, incluindo as buscas registradas sem resultado. " + DB + ".")
notas(sl, "A hierarquia aplicada segue a ordem: (1) documentação jurídica primária da operação; (2) bases "
          "regulatórias oficiais; (3) demonstrações auditadas; (4) relatórios de rating; (5) informações da "
          "companhia; (6) research. "
          "Onde houve conflito, prevaleceu a fonte de maior hierarquia. Divergências registradas: "
          "(a) número de séries de CRI-I, CRI-II e CRI-IV — os anúncios de oferta contam apenas séries públicas "
          "(4, 4 e 6) enquanto os termos e lâminas contam o total incluindo a série privada (5, 5 e 7); "
          "prevaleceu o documento constitutivo. "
          "(b) vencimento da 1ª série da debênture — escritura fixa 18/02/2033 (132 meses) e o relatório do agente "
          "fiduciário exibe 18/08/2035 para as duas séries; prevaleceu a escritura. "
          "(c) montante por série da 177ª — a análise consolidada trazia n/d por série e total de R$ 647,1 mi; o "
          "termo de securitização traz o detalhamento e total de R$ 647,059 mi; prevaleceu o termo. "
          "Buscas registradas sem resultado: FIDC Solfácil VIII ou posterior no cadastro CVM; emissão de CRI "
          "posterior a 31/07/2026; documentos das debêntures SFCL11/21/31/41, dimensionadas em R$ 150 mi pela "
          "análise consolidada mas sem documento no acervo.")

# ---------------------------------------------------------------- salvar
os.makedirs(OUTDIR, exist_ok=True)
caminho = os.path.join(OUTDIR, "Solfacil_Renovacao_Credito_20260824_claude.pptx")
prs.save(caminho)
n = len(prs.slides._sldIdLst)
print(f"Deck salvo: {caminho}")
print(f"Slides: {n}  (corpo principal 1-9, appendix 10-{n})")
