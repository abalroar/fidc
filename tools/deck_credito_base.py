# -*- coding: utf-8 -*-
"""Primitivas do deck de renovação de crédito.

Herda grade, paleta e fontes do deck analítico existente e acrescenta o que o
formato de comitê exige: alinhamento numérico à direita, negrito seletivo por
célula, notas de rodapé curtas e speaker notes com rastreabilidade completa.
"""
import os, sys
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck_base import (LARANJA, LARANJA_ESC, LARANJA_CLA, PETROLEO, PETROLEO_CLA,
                       PRETO, CINZA_ESC, CINZA_MED, CINZA_CLARO, CINZA_FUNDO, BRANCO,
                       FONTE, SL_W, SL_H, MARGEM, UTIL_W, col_x, col_w, ler,
                       sem_sombra, txt, estilizar_gr)

# Faixas verticais do slide executivo
Y_KICKER   = Emu(228600)
Y_TITULO   = Emu(402336)
Y_SUB      = Emu(950976)
Y_CORPO    = Emu(1325880)
Y_RODAPE   = Emu(6355080)


def kicker(sl, texto):
    txt(sl, MARGEM, Y_KICKER, col_w(9), Emu(146304), texto.upper(), tam=8.5,
        negrito=True, cor=CINZA_MED)


def titulo_conclusivo(sl, texto, subtitulo=None):
    """O título carrega a conclusão de crédito do slide."""
    txt(sl, MARGEM, Y_TITULO, col_w(12), Emu(502920), texto, tam=19, negrito=True,
        cor=PRETO, entrelinha=0.92)
    if subtitulo:
        txt(sl, MARGEM, Y_SUB, col_w(12), Emu(283464), subtitulo, tam=11, cor=CINZA_ESC)


def fonte(sl, texto):
    txt(sl, MARGEM, Y_RODAPE, UTIL_W, Emu(182880), texto, tam=7.5, cor=CINZA_MED)


def notas(sl, texto):
    """Rastreabilidade completa: documento, cláusula, data-base, divergência e método."""
    sl.notes_slide.notes_text_frame.text = texto


def rotulo_secao(sl, y, texto):
    txt(sl, MARGEM, y, UTIL_W, Emu(160020), texto.upper(), tam=8.5, negrito=True, cor=CINZA_MED)


def destaque_numero(sl, x, y, w, valor, rotulo, cor_valor=PRETO):
    """Métrica com hierarquia por tamanho, não por cor. Sem caixa nem preenchimento."""
    txt(sl, x, y, w, Emu(320040), valor, tam=20, negrito=True, cor=cor_valor)
    txt(sl, x, Emu(y.emu + 310896), w, Emu(228600), rotulo, tam=8.5, cor=CINZA_MED)


def barra_lateral(sl, x, y, w, h, cor):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = cor
    s.line.fill.background(); sem_sombra(s)
    return s


def tab(sl, x, y, w, h, cabs, linhas, larguras, tam=8.5, tam_cab=8.5,
        num_cols=(), bold_cols=(), bold_cells=(), h_cab=Emu(228600), h_linha=None,
        destaque_linhas=(), cab_cor=CINZA_ESC):
    """Tabela executiva nativa.

    num_cols        colunas alinhadas à direita (números)
    bold_cols       colunas inteiras em negrito
    bold_cells      pares (linha, coluna) em negrito, base 0 nas linhas de dados
    destaque_linhas linhas com fundo neutro, para exceções
    """
    gf = sl.shapes.add_table(len(linhas) + 1, len(cabs), x, y, w, h)
    tb = gf.table
    tb.first_row = False; tb.horz_banding = False
    total = sum(larguras)
    for j, frac in enumerate(larguras):
        tb.columns[j].width = Emu(int(w.emu * frac / total))
    tb.rows[0].height = h_cab
    if h_linha:
        for i in range(1, len(linhas) + 1):
            tb.rows[i].height = h_linha
    for j, c in enumerate(cabs):
        cel = tb.cell(0, j); cel.text = ""
        cel.fill.solid(); cel.fill.fore_color.rgb = cab_cor
        cel.margin_left = cel.margin_right = Emu(45720)
        cel.margin_top = cel.margin_bottom = Emu(18288)
        cel.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cel.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if j in num_cols else PP_ALIGN.LEFT
        r = p.add_run(); r.text = c
        r.font.size = Pt(tam_cab); r.font.bold = True
        r.font.color.rgb = BRANCO; r.font.name = FONTE
    for i, linha in enumerate(linhas, start=1):
        for j, v in enumerate(linha):
            cel = tb.cell(i, j); cel.text = ""
            cel.fill.solid()
            cel.fill.fore_color.rgb = CINZA_FUNDO if (i - 1) in destaque_linhas else BRANCO
            cel.margin_left = cel.margin_right = Emu(45720)
            cel.margin_top = cel.margin_bottom = Emu(13716)
            cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cel.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j in num_cols else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(tam)
            r.font.bold = (j in bold_cols) or ((i - 1, j) in bold_cells)
            r.font.color.rgb = PRETO; r.font.name = FONTE
    return tb


def caixa_sintese(sl, x, y, w, h, titulo, itens, cor_faixa):
    """Bloco de síntese: faixa de título e lista de itens curtos."""
    f = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(256032))
    f.fill.solid(); f.fill.fore_color.rgb = cor_faixa
    f.line.color.rgb = cor_faixa; sem_sombra(f)
    tf = f.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(91440)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = titulo.upper()
    r.font.size = Pt(9); r.font.bold = True; r.font.name = FONTE; r.font.color.rgb = BRANCO
    corpo = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Emu(y.emu + 256032), w, Emu(h.emu - 256032))
    corpo.fill.solid(); corpo.fill.fore_color.rgb = BRANCO
    corpo.line.color.rgb = CINZA_CLARO; corpo.line.width = Pt(0.75); sem_sombra(corpo)
    tfc = corpo.text_frame; tfc.word_wrap = True
    tfc.margin_left = tfc.margin_right = Emu(91440)
    tfc.margin_top = tfc.margin_bottom = Emu(68580)
    tfc.vertical_anchor = MSO_ANCHOR.TOP
    for i, (forte, resto) in enumerate(itens):
        p = tfc.paragraphs[0] if i == 0 else tfc.add_paragraph()
        p.space_after = Pt(5)
        a = p.add_run(); a.text = forte
        a.font.size = Pt(9); a.font.bold = True; a.font.color.rgb = PRETO; a.font.name = FONTE
        b = p.add_run(); b.text = "  " + resto
        b.font.size = Pt(9); b.font.color.rgb = CINZA_ESC; b.font.name = FONTE
    return corpo
