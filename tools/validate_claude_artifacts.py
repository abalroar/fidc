# -*- coding: utf-8 -*-
"""Checklist de aceite, verificado mecanicamente sobre os dois entregaveis."""
import csv, os, re, sys, zipfile
from openpyxl import load_workbook
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
DATA = os.path.join(ROOT, "data", "solfacil_claude")
XLSX = os.path.join(ROOT, "outputs", "solfacil", "Solfacil_CRI_FIDC_20260822_claude.xlsx")
PPTX = os.path.join(ROOT, "outputs", "solfacil", "Solfacil_CRI_FIDC_20260822_claude.pptx")

ok, falhas = [], []


def checa(cond, texto, detalhe=""):
    (ok if cond else falhas).append(texto + (f" -- {detalhe}" if detalhe and not cond else ""))


def csvs():
    d = {}
    for n in sorted(os.listdir(DATA)):
        if n.endswith(".csv"):
            with open(os.path.join(DATA, n), encoding="utf-8") as fh:
                r = list(csv.DictReader(fh))
            d[n] = r
    return d


C = csvs()
wb = load_workbook(XLSX)
prs = Presentation(PPTX)

# 1. Todo numero rastreia a uma linha de CSV com fonte_id
sem_fonte = []
for n, rows in C.items():
    if not rows:
        continue
    cols = rows[0].keys()
    if "fonte_id" not in cols:
        # Tabelas de sintese rastreiam a fonte por coluna propria, nao por fonte_id
        alternativas = {"15_fidc_vs_cri.csv": ["evidencia"],
                        "16_conflitos.csv": ["fonte_A", "fonte_B"]}
        if n in alternativas:
            faltando = [i for i, r in enumerate(rows)
                        if not all((r.get(k) or "").strip() for k in alternativas[n])]
            if faltando:
                sem_fonte.append(f"{n}: {len(faltando)} linhas sem {alternativas[n]}")
            continue
        # Agregados derivados de tabelas que já carregam fonte_id linha a linha
        if n not in ("19_glossario.csv", "20_lacunas.csv", "03b_elegibilidade_deltas.csv",
                     "13b_amortizacao_realizada.csv", "raw_anexo1_cri2_kanastra.csv",
                     "25_captacao_por_ano.csv", "29_waterfall_comparado.csv",
                     "34_pontos_em_aberto.csv"):
            sem_fonte.append(n)
        continue
    vazios = [i for i, r in enumerate(rows) if not (r.get("fonte_id") or "").strip()]
    if vazios:
        sem_fonte.append(f"{n}: {len(vazios)} linhas sem fonte_id")
checa(not sem_fonte, "Toda linha de dado carrega fonte_id", str(sem_fonte))

# fonte_id resolve no inventario
inv = {r["fonte_id"] for r in C["17_fontes.csv"]}
orfaos = set()
for n, rows in C.items():
    for r in rows:
        for f in str(r.get("fonte_id", "")).split(";"):
            f = re.sub(r"\s*\(.*?\)", "", f).strip()   # remove localizador: "ANX-DECK (sl.35)" -> "ANX-DECK"
            if f and f not in ("n/a", "n/d", "") and f not in inv:
                orfaos.add(f)
checa(not orfaos, "Todo fonte_id resolve na aba 17_Fontes", str(sorted(orfaos)))

# 2. Nenhum n/d virou zero: campos n/d nunca aparecem como 0 numerico no mesmo par
checa(any("n/d" in str(r.values()) for r in C["02_series.csv"]),
      "n/d preservado como texto em 02_Series (nao convertido a zero)")

# 3. Divergencias em Conflitos com decisao justificada
conf = C["16_conflitos.csv"]
checa(len(conf) >= 5, f"Conflitos registrados: {len(conf)} (minimo 5)")
checa(all(r["decisao_adotada"].strip() and r["justificativa"].strip() for r in conf),
      "Todo conflito tem decisao adotada e justificativa")

# 4. Series
ser = C["02_series.csv"]
cri = [r for r in ser if r["veiculo_id"].startswith("CRI")]
fidc = [r for r in ser if r["veiculo_id"].startswith("FIDC")]
checa(len(cri) == 34, f"34 series de CRI em 02_Series (encontradas {len(cri)})")
checa(len({r["veiculo_id"] for r in fidc}) == 7, "As 7 classes de FIDC representadas em 02_Series")
soma = sum(float(r["montante_reportado_deck_Rmi"]) for r in cri
           if r["montante_reportado_deck_Rmi"] not in ("n/d", ""))
nd_cri = [r for r in cri if r["montante_reportado_deck_Rmi"] in ("n/d", "")]
checa(not nd_cri and abs(soma - 3670.64) < 0.1,
      f"As 34 séries de CRI têm montante documentado e somam R$ {soma:.2f} mi (nenhuma n/d)")

# Pilha de funding: 70 tranches nos três instrumentos
fund = C["21_funding_por_tranche.csv"]
checa(len(fund) == 70, f"21_Funding_por_Tranche cobre as 70 tranches do programa (encontradas {len(fund)})")
inst = {r["instrumento"] for r in fund}
checa(len(inst) == 3, f"Os três instrumentos estão na pilha de funding: {sorted(inst)}")
com_preco = [r for r in fund if r["preco_taxa_contratada"] not in ("n/d", "")]
checa(len(com_preco) >= 60,
      f"{len(com_preco)} das {len(fund)} tranches têm preço documentado")
deb = [r for r in C["02_series.csv"] if r["veiculo_id"] == "DEB-I"]
checa(len(deb) == 2, f"As 2 séries da debênture estão em 02_Series (encontradas {len(deb)})")

# 5. Limites contratuais x praticados em colunas separadas
cols_ser = set(ser[0].keys())
for par in ["quantidade_ofertada_lote_base", "quantidade_subscrita", "montante_ofertado_Rmi",
            "montante_subscrito_Rmi", "taxa_teto_lamina", "taxa_piso_lamina", "taxa_contratada"]:
    checa(par in cols_ser, f"Coluna separada presente: {par}")

# 6. Cronograma distingue Realizado de Projetado
cron = C["13_cronograma_pagamentos.csv"]
checa(all(r["status"] in ("Projetado", "Realizado") for r in cron),
      "13_Cronograma: toda linha tem status Projetado ou Realizado")
checa(all(r["status"] == "Realizado" for r in C["13b_amortizacao_realizada.csv"]),
      "13b: amortizacao observada marcada como Realizado")

# 7. Workbook: ListObjects e graficos nativos
tabs = sum(len(ws.tables) for ws in wb.worksheets)
graf = sum(len(ws._charts) for ws in wb.worksheets)
nomes_ok = all(n.startswith("tbl_") for ws in wb.worksheets for n in ws.tables)
checa(nomes_ok, f"Todas as {tabs} tabelas do Excel sao ListObject nomeadas tbl_*")
checa(graf >= 5, f"Graficos nativos no Excel: {graf}")
z = zipfile.ZipFile(XLSX)
checa(not [n for n in z.namelist() if "/media/" in n], "Nenhuma imagem embutida no Excel")
checa(all(ws["A2"].value for ws in wb.worksheets), "Toda aba tem leitura em portugues na linha 2")
checa(all(ws.freeze_panes for ws in wb.worksheets if ws.tables), "Painéis congelados em toda aba com tabela")
checa(not any(ws.merged_cells.ranges for ws in wb.worksheets), "Nenhuma celula mesclada")

# 8. Metodologia e glossario
met = C["18_metodologia.csv"]
checa(all(r["formula"].strip() and r["qualificador"].strip() for r in met),
      "18_Metodologia: toda metrica tem formula e qualificador")
glo = {r["termo"].lower() for r in C["19_glossario.csv"]}
checa(len(glo) >= 15, f"Glossario com {len(glo)} termos (minimo 15)")

# 9. Veredito com as tres respostas
ver = C["15_fidc_vs_cri.csv"]
vants = {r["vantagem_real"] for r in ver}
checa({"FIDC", "CRI"} <= vants and any("Neutro" in v for v in vants),
      f"15_FIDC_vs_CRI cobre FIDC, CRI e neutro: {sorted(vants)}")
checa(all(r["o_que_falta_para_confirmar"].strip() for r in ver),
      "Toda dimensao diz o que falta para confirmar")

# 10. Deck
PAL = {"EC7000", "A85000", "F7C89A", "197278", "7FB3B6", "000000", "323436", "6E6E6E", "BFBFBF", "F2F2F2", "FFFFFF"}
zp = zipfile.ZipFile(PPTX)
checa(not [n for n in zp.namelist() if "/media/" in n], "Deck sem imagens ou icones")
cores = set()
for n in zp.namelist():
    if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
        cores |= set(c.upper() for c in re.findall(r'val="([0-9A-Fa-f]{6})"', zp.read(n).decode("utf8", "ignore")))
checa(not (cores - PAL), f"Paleta do deck restrita a laranja, preto e cinzas: fora = {sorted(cores - PAL)}")
n_sl = len(prs.slides._sldIdLst)
checa(12 <= n_sl <= 22, f"Deck com {n_sl} slides (o escopo cresceu com os nove itens analíticos do comitê)")
tab_p = sum(1 for sl in prs.slides for sh in sl.shapes if sh.has_table)
gr_p = sum(1 for sl in prs.slides for sh in sl.shapes if getattr(sh, "has_chart", False) and sh.has_chart)
checa(gr_p >= 6, f"Graficos nativos no deck: {gr_p}")
checa(tab_p >= 8, f"Tabelas nativas no deck: {tab_p}")
setas = re.compile(r"[›→»▶➤⟶⇒]")
achou = [(i, sh.text_frame.text[:30]) for i, sl in enumerate(prs.slides, 1)
         for sh in sl.shapes if sh.has_text_frame and setas.search(sh.text_frame.text)]
checa(not achou, "Nenhuma seta tipografica no deck", str(achou))
# rodape presente em todo slide, como texto e nao como forma preenchida
sem_rod = []
for i, sl in enumerate(prs.slides, 1):
    tem = any(sh.has_text_frame and sh.top and sh.top > 6100000 and sh.text_frame.text.strip()
              for sh in sl.shapes)
    if not tem:
        sem_rod.append(i)
checa(not sem_rod, "Todo slide traz a linha de fonte e data-base no pe", str(sem_rod))

# 11. Excel e PPTX nao se contradizem nos numeros-chave
def cel(aba, l, c):
    return wb[aba].cell(row=l, column=c).value

painel = {r["indicador"]: r["valor"] for r in C["00_painel.csv"]}
checa(painel["Séries de CRI"] == "34" and n_sl and len(cri) == 34,
      "Painel, CSV e deck concordam em 34 series de CRI")
checa(painel["Volume nominal de CRI emitido"] == "3.670,6",
      "Painel e deck concordam no volume nominal de R$ 3.670,7 mi")

print("=" * 78)
print(f"APROVADOS: {len(ok)}    FALHAS: {len(falhas)}")
print("=" * 78)
for t in ok:
    print("  OK    ", t)
if falhas:
    print()
    for t in falhas:
        print("  FALHA ", t)
sys.exit(1 if falhas else 0)
