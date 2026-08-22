# -*- coding: utf-8 -*-
"""Fundação visual do deck: grade em EMU, paleta e primitivas de desenho.

Paleta: laranja Itaú BBA como cor primária, petróleo como segunda cor categórica
(necessária para separar FIDC de CRI nas comparações) e escala de cinzas.
Nenhum elemento decorativo: toda forma colorida carrega dado ou rótulo.
"""
import csv, os
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# ---------------------------------------------------------------- paleta
LARANJA     = RGBColor(0xEC, 0x70, 0x00)   # Itaú BBA - cor primária
LARANJA_ESC = RGBColor(0xA8, 0x50, 0x00)
LARANJA_CLA = RGBColor(0xF7, 0xC8, 0x9A)
PETROLEO    = RGBColor(0x19, 0x72, 0x78)   # segunda cor categórica
PETROLEO_CLA = RGBColor(0x7F, 0xB3, 0xB6)
PRETO       = RGBColor(0x00, 0x00, 0x00)
CINZA_ESC   = RGBColor(0x32, 0x34, 0x36)
CINZA_MED   = RGBColor(0x6E, 0x6E, 0x6E)
CINZA_CLARO = RGBColor(0xBF, 0xBF, 0xBF)
CINZA_FUNDO = RGBColor(0xF2, 0xF2, 0xF2)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
FONTE = "Calibri"

# ---------------------------------------------------------------- grade, em EMU
SL_W, SL_H = Emu(12192000), Emu(6858000)
MARGEM   = Emu(548640)
UTIL_W   = Emu(12192000 - 2 * 548640)
Y_EYEBROW = Emu(228600)
Y_TITULO  = Emu(411480)
Y_SUB     = Emu(795528)
Y_CORPO   = Emu(1417320)
Y_RODAPE  = Emu(6217920)
COLS, GUT = 12, Emu(137160)
_PASSO = (UTIL_W.emu - GUT.emu * (COLS - 1)) / COLS


def col_x(i):
    return Emu(int(MARGEM.emu + i * (_PASSO + GUT.emu)))


def col_w(n):
    return Emu(int(n * _PASSO + (n - 1) * GUT.emu))


def ler(DATA, nome):
    with open(os.path.join(DATA, nome), encoding="utf-8") as fh:
        r = list(csv.reader(fh))
    return r[0], r[1:]


def sem_sombra(shp):
    shp.shadow.inherit = False
    return shp


def txt(sl, x, y, w, h, texto, tam=12, negrito=False, cor=PRETO, alinh=PP_ALIGN.LEFT,
        espaco=0, anchor=MSO_ANCHOR.TOP, italico=False, entrelinha=None):
    cx = sl.shapes.add_textbox(x, y, w, h)
    tf = cx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(texto.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alinh
        p.space_after = Pt(espaco)
        if entrelinha:
            p.line_spacing = entrelinha
        r = p.add_run(); r.text = ln
        r.font.size = Pt(tam); r.font.bold = negrito; r.font.italic = italico
        r.font.color.rgb = cor; r.font.name = FONTE
    return cx


def cabecalho(sl, eyebrow, titulo, subtitulo=None):
    txt(sl, MARGEM, Y_EYEBROW, col_w(9), Emu(160020), eyebrow.upper(), tam=9, negrito=True, cor=CINZA_MED)
    txt(sl, MARGEM, Y_TITULO, col_w(11), Emu(365760), titulo, tam=23, negrito=True, cor=PRETO)
    if subtitulo:
        txt(sl, MARGEM, Y_SUB, col_w(11), Emu(320040), subtitulo, tam=12.5, cor=CINZA_ESC)


def rodape(sl, fonte_txt):
    txt(sl, MARGEM, Y_RODAPE, UTIL_W, Emu(274320), fonte_txt, tam=8.5, cor=CINZA_MED)


def leitura(sl, y, texto, w=None):
    txt(sl, MARGEM, y, w or UTIL_W, Emu(365760), texto, tam=12, negrito=True, cor=PRETO)


def bloco(sl, x, y, w, h, texto, fundo, cor_txt=BRANCO, tam=10.5, negrito=True,
          forma=MSO_SHAPE.RECTANGLE, borda=None):
    s = sl.shapes.add_shape(forma, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fundo
    s.line.color.rgb = borda or fundo; s.line.width = Pt(0.75)
    sem_sombra(s)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.margin_top = tf.margin_bottom = Emu(27432)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = texto
    r.font.size = Pt(tam); r.font.bold = negrito
    r.font.color.rgb = cor_txt; r.font.name = FONTE
    return s


def tabela(sl, x, y, w, h, cabecalhos, linhas, larguras=None, tam=9.5, tam_cab=9.5,
           altura_cab=Emu(320040), altura_linha=None, alinh_centro=(), destaque=()):
    gf = sl.shapes.add_table(len(linhas) + 1, len(cabecalhos), x, y, w, h)
    tb = gf.table
    tb.first_row = False; tb.horz_banding = False
    if larguras:
        total = sum(larguras)
        for j, frac in enumerate(larguras):
            tb.columns[j].width = Emu(int(w.emu * frac / total))
    tb.rows[0].height = altura_cab
    if altura_linha:
        for i in range(1, len(linhas) + 1):
            tb.rows[i].height = altura_linha
    for j, c in enumerate(cabecalhos):
        cel = tb.cell(0, j); cel.text = ""
        cel.fill.solid(); cel.fill.fore_color.rgb = CINZA_ESC
        cel.margin_left = cel.margin_right = Emu(54864)
        cel.margin_top = cel.margin_bottom = Emu(27432)
        cel.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cel.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j in alinh_centro else PP_ALIGN.LEFT
        r = p.add_run(); r.text = c
        r.font.size = Pt(tam_cab); r.font.bold = True
        r.font.color.rgb = BRANCO; r.font.name = FONTE
    for i, linha in enumerate(linhas, start=1):
        for j, v in enumerate(linha):
            cel = tb.cell(i, j); cel.text = ""
            cel.fill.solid()
            cel.fill.fore_color.rgb = CINZA_FUNDO if (i - 1) in destaque else BRANCO
            cel.margin_left = cel.margin_right = Emu(54864)
            cel.margin_top = cel.margin_bottom = Emu(22860)
            cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cel.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j in alinh_centro else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(tam); r.font.bold = (j == 0)
            r.font.color.rgb = PRETO; r.font.name = FONTE
    return tb


def estilizar_gr(gr, cores, legenda=True, tam=9.5, unidade_y=None, rotulos=False):
    gr.has_title = False
    gr.font.size = Pt(tam); gr.font.name = FONTE; gr.font.color.rgb = CINZA_ESC
    if legenda:
        gr.has_legend = True
        gr.legend.position = XL_LEGEND_POSITION.BOTTOM
        gr.legend.include_in_layout = False
        gr.legend.font.size = Pt(tam)
    else:
        gr.has_legend = False
    for s, cor in zip(gr.series, cores):
        if cor is None:            # série de apoio invisível (base do waterfall / offset do Gantt)
            s.format.fill.background()
            s.format.line.fill.background()
            continue
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = cor
        s.format.line.color.rgb = cor
    try:
        gr.value_axis.has_major_gridlines = True
        gr.value_axis.major_gridlines.format.line.color.rgb = CINZA_CLARO
        gr.value_axis.major_gridlines.format.line.width = Pt(0.5)
        gr.value_axis.tick_labels.font.size = Pt(tam - 0.5)
        gr.category_axis.tick_labels.font.size = Pt(tam - 0.5)
        if unidade_y:
            gr.value_axis.axis_title.text_frame.text = unidade_y
            gr.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(tam)
    except Exception:
        pass
    if rotulos:
        gr.plots[0].has_data_labels = True
        gr.plots[0].data_labels.font.size = Pt(tam - 0.5)
        gr.plots[0].data_labels.font.color.rgb = PRETO
    return gr
