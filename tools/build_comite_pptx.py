# -*- coding: utf-8 -*-
"""Os três quadros de comitê em PowerPoint, com tabelas e gráficos nativos.

Referência visual: slides Itaú BBA de captação via FIDC — faixa de cabeçalho em
azul-marinho, acento laranja, tabela compacta, gráfico limpo abaixo do título.
"""
import os, sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
OUTDIR = os.path.join(ROOT, "outputs", "solfacil")
sys.path.insert(0, HERE)
import solfacil_comite_tabelas as T

NAVY = RGBColor(0x1F, 0x38, 0x64)
NAVY2 = RGBColor(0x2E, 0x54, 0x96)
LARANJA = RGBColor(0xEC, 0x70, 0x00)
LARANJA_CLA = RGBColor(0xF7, 0xC8, 0x9A)
ND_FILL = RGBColor(0xFC, 0xE4, 0xD6)
ATR_FILL = RGBColor(0xEA, 0xEF, 0xF7)
PRETO = RGBColor(0, 0, 0)
CINZA_ESC = RGBColor(0x32, 0x34, 0x36)
CINZA_MED = RGBColor(0x6E, 0x6E, 0x6E)
CINZA_CLARO = RGBColor(0xD9, 0xD9, 0xD9)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
FONTE = "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
VAZIO = prs.slide_layouts[6]
MARGEM = Emu(457200)
UTIL = Emu(12192000 - 2 * 457200)


def txt(sl, x, y, w, h, t, tam=10, neg=False, cor=PRETO, al=PP_ALIGN.LEFT, anc=MSO_ANCHOR.TOP, esp=0):
    cx = sl.shapes.add_textbox(x, y, w, h)
    tf = cx.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anc
    for i, ln in enumerate(t.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al; p.space_after = Pt(esp)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(tam); r.font.bold = neg; r.font.color.rgb = cor; r.font.name = FONTE
    return cx


def cabecalho(sl, titulo, sub=None):
    marca = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGEM, Emu(274320), Emu(365760), Emu(228600))
    marca.fill.solid(); marca.fill.fore_color.rgb = NAVY
    marca.line.fill.background(); marca.shadow.inherit = False
    tf = marca.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "IB"
    r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = BRANCO; r.font.name = FONTE
    txt(sl, Emu(MARGEM.emu + 457200), Emu(256032), Emu(10500000), Emu(310896), titulo, tam=17, neg=True)
    if sub:
        txt(sl, Emu(MARGEM.emu + 457200), Emu(594360), Emu(10500000), Emu(228600), sub, tam=9.5, cor=CINZA_ESC)


def rodape(sl, t):
    txt(sl, MARGEM, Emu(6400800), UTIL, Emu(182880), t, tam=7.5, cor=CINZA_MED)


def tabela(sl, x, y, w, h, cols, linhas, larg, tam=8, tam_cab=8, h_cab=Emu(256032),
           h_lin=None, atr_col=True, centro=True):
    gf = sl.shapes.add_table(len(linhas) + 1, len(cols), x, y, w, h)
    tb = gf.table; tb.first_row = False; tb.horz_banding = False
    tot = sum(larg)
    for j, f in enumerate(larg):
        tb.columns[j].width = Emu(int(w.emu * f / tot))
    tb.rows[0].height = h_cab
    if h_lin:
        for i in range(1, len(linhas) + 1):
            tb.rows[i].height = h_lin
    for j, c in enumerate(cols):
        cel = tb.cell(0, j); cel.text = ""
        cel.fill.solid(); cel.fill.fore_color.rgb = NAVY
        cel.margin_left = cel.margin_right = Emu(36576)
        cel.margin_top = cel.margin_bottom = Emu(13716)
        cel.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cel.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if (j > 0 and centro) else PP_ALIGN.LEFT
        r = p.add_run(); r.text = c
        r.font.size = Pt(tam_cab); r.font.bold = True; r.font.color.rgb = BRANCO; r.font.name = FONTE
    for i, linha in enumerate(linhas, start=1):
        for j, v in enumerate(linha):
            cel = tb.cell(i, j); cel.text = ""
            v = str(v)
            cel.fill.solid()
            if j == 0 and atr_col:
                cel.fill.fore_color.rgb = ATR_FILL
            elif v == "n/d":
                cel.fill.fore_color.rgb = ND_FILL
            else:
                cel.fill.fore_color.rgb = BRANCO
            cel.margin_left = cel.margin_right = Emu(36576)
            cel.margin_top = cel.margin_bottom = Emu(9144)
            cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cel.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (j > 0 and centro) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = v
            r.font.size = Pt(tam)
            r.font.bold = (j == 0 and atr_col)
            r.font.italic = (v == "n/d")
            r.font.color.rgb = CINZA_ESC if v == "n/d" else PRETO
            r.font.name = FONTE
    return tb


def grafico(sl, tipo, x, y, w, h, cats, series, cores, unidade=None, legenda=True, rot=False, tam=8):
    d = CategoryChartData(); d.categories = cats
    for n, v in series:
        d.add_series(n, tuple(v))
    gr = sl.shapes.add_chart(tipo, x, y, w, h, d).chart
    gr.has_title = False
    gr.font.size = Pt(tam); gr.font.name = FONTE; gr.font.color.rgb = CINZA_ESC
    if legenda:
        gr.has_legend = True; gr.legend.position = XL_LEGEND_POSITION.BOTTOM
        gr.legend.include_in_layout = False; gr.legend.font.size = Pt(tam)
    else:
        gr.has_legend = False
    for s, c in zip(gr.series, cores):
        s.format.fill.solid(); s.format.fill.fore_color.rgb = c
        s.format.line.color.rgb = c
    try:
        gr.value_axis.has_major_gridlines = True
        gr.value_axis.major_gridlines.format.line.color.rgb = CINZA_CLARO
        gr.value_axis.major_gridlines.format.line.width = Pt(0.5)
        gr.value_axis.tick_labels.font.size = Pt(tam - 0.5)
        gr.category_axis.tick_labels.font.size = Pt(tam - 0.5)
        if unidade:
            gr.value_axis.axis_title.text_frame.text = unidade
            gr.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(tam)
    except Exception:
        pass
    if rot:
        gr.plots[0].has_data_labels = True
        gr.plots[0].data_labels.font.size = Pt(tam - 1)
        gr.plots[0].data_labels.font.color.rgb = PRETO
    return gr

# ═══════════════════════════════ 1 · Quadro 1A — FIDCs
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Solfácil | FIDCs · Perfil dos Veículos",
          "Cadastro e informe mensal da CVM · data-base jul/26 · FIDC VIII em discussão, coluna reservada")
tabela(sl, MARGEM, Emu(1005840), UTIL, Emu(3383280), T.FIDC_COLS, T.FIDC_LINHAS,
       [0.19, 0.10, 0.10, 0.10, 0.10, 0.10, 0.10, 0.10, 0.11], tam=8.5, tam_cab=8.5,
       h_lin=Emu(246888))
txt(sl, MARGEM, Emu(4526280), Emu(5486400), Emu(182880), "COMPOSIÇÃO DO PL POR CLASSE DE COTA (%)",
    tam=8, neg=True, cor=CINZA_MED)
grafico(sl, XL_CHART_TYPE.COLUMN_STACKED, MARGEM, Emu(4709160), Emu(5486400), Emu(1600200),
        [c for c in T.FIDC_COLS[1:8]],
        [("Sênior 1", [63.6, 72.9, 66.8, 100.0, 79.1, 66.4, 74.0]),
         ("Mezanino", [32.5, 14.9, 21.0, 0.0, 13.5, 18.7, 20.8]),
         ("Sub Júnior", [3.9, 12.2, 12.2, 0.0, 7.4, 14.9, 5.2])],
        [NAVY, LARANJA, LARANJA_CLA], unidade="% do PL")
txt(sl, Emu(6217920), Emu(4526280), Emu(5029200), Emu(182880), "PL POR FUNDO (R$ mm)",
    tam=8, neg=True, cor=CINZA_MED)
grafico(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(6217920), Emu(4709160), Emu(5029200), Emu(1600200),
        [c for c in T.FIDC_COLS[1:8]],
        [("PL (R$ mm)", [83.7, 94.1, 141.1, 17.5, 67.5, 211.1, 619.6])],
        [NAVY], unidade="R$ mm", legenda=False, rot=True)
rodape(sl, "Fonte: CVM — cadastro de fundos e Informe Mensal FIDC (31/07/2026); relatórios de rating do acervo. "
           "Células destacadas: rentabilidade YTD por classe é campo do informe mensal, a obter no CVM Fundos.NET.")

# ═══════════════════════════════ 2 · Quadro 1B — CRIs
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Solfácil | CRIs · Perfil das Operações",
          "Abertura por série e por securitizadora · prospectos, comunicados e termos de securitização")
tabela(sl, MARGEM, Emu(1005840), Emu(7315200), Emu(4663440), T.CRI_COLS, T.CRI_LINHAS,
       [0.22, 0.13, 0.13, 0.13, 0.13, 0.13, 0.13], tam=7.5, tam_cab=8, h_lin=Emu(228600))
txt(sl, Emu(8046720), Emu(1005840), Emu(3657600), Emu(182880), "COMPOSIÇÃO DA EMISSÃO POR SÉRIE (%)",
    tam=8, neg=True, cor=CINZA_MED)
grafico(sl, XL_CHART_TYPE.COLUMN_STACKED, Emu(8046720), Emu(1188720), Emu(3657600), Emu(2194560),
        [c.replace("CRI ", "") for c in T.CRI_COLS[1:]],
        [("1ª", [59.7, 65.0, 49.0, 43.3, 22.1, 15.5]), ("2ª", [14.9, 18.0, 16.0, 21.7, 47.9, 69.5]),
         ("3ª", [17.9, 10.0, 18.0, 12.0, 15.0, 8.0]), ("4ª", [5.0, 4.0, 10.0, 6.0, 8.0, 4.0]),
         ("5ª", [2.5, 3.0, 4.0, 10.0, 4.0, 3.0]), ("6ª", [0, 0, 3.0, 4.0, 3.0, 0]),
         ("7ª", [0, 0, 0, 3.0, 0, 0])],
        [NAVY, NAVY2, LARANJA, LARANJA_CLA, CINZA_MED, CINZA_CLARO, RGBColor(0x9D, 0xC3, 0xE6)],
        unidade="% do volume")
txt(sl, Emu(8046720), Emu(3566160), Emu(3657600), Emu(182880), "VOLUME EMITIDO (R$ mm)",
    tam=8, neg=True, cor=CINZA_MED)
grafico(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(8046720), Emu(3749040), Emu(3657600), Emu(1920240),
        [c.replace("CRI ", "") for c in T.CRI_COLS[1:]],
        [("R$ mm", [603.0, 750.0, 750.0, 450.0, 470.6, 647.1])],
        [LARANJA], unidade="R$ mm", legenda=False, rot=True)
rodape(sl, "Fonte: prospectos definitivos, lâminas, comunicados de bookbuilding e termos de securitização. "
           "Células destacadas: rentabilidade YTD por série e over 90 por operação são campos do Informe Mensal CRI, a obter no CVM Fundos.NET.")

# ═══════════════════════════════ 3 · Quadro 2A — emissões FIDC
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Solfácil | FIDCs · Emissões, Remuneração-Alvo e Subordinação Mínima",
          "Volume emitido soma as emissões registradas de cada fundo; não é saldo em aberto")
tabela(sl, MARGEM, Emu(1005840), UTIL, Emu(4297680), T.EMI_FIDC_COLS, T.EMI_FIDC,
       [0.07, 0.08, 0.08, 0.09, 0.30, 0.26, 0.12], tam=7, tam_cab=8, h_lin=Emu(475488),
       atr_col=True, centro=False)
txt(sl, MARGEM, Emu(5486400), UTIL, Emu(640080),
    "Como o índice aparece: nos FIDCs é PISO DE SUBORDINAÇÃO TOTAL, em percentual do PL, somando mezanino e "
    "subordinada júnior — 20% ou 25% conforme o fundo, e ausente no FIDC IV. Não é a mesma métrica usada nos CRIs.",
    tam=9, cor=PRETO, esp=2)
rodape(sl, "Fonte: CVM — registro de ofertas de cotas; suplementos de classe; consolidação de informes. Campos n/d a obter no Fundos.NET.")

# ═══════════════════════════════ 4 · Quadro 2B — emissões CRI
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Solfácil | CRIs · Emissões, Remuneração-Alvo e Cobertura Mínima",
          "Abertura por série e remuneração contratada de cada uma das 34 séries")
tabela(sl, MARGEM, Emu(1005840), UTIL, Emu(3931920), T.EMI_CRI_COLS, T.EMI_CRI,
       [0.12, 0.08, 0.09, 0.28, 0.28, 0.15], tam=7, tam_cab=8, h_lin=Emu(566928),
       atr_col=True, centro=False)
txt(sl, MARGEM, Emu(5120640), UTIL, Emu(914400),
    "Como o índice aparece: nos CRIs não existe piso percentual de subordinação. A proteção é uma RAZÃO DE COBERTURA "
    "por camada — valor presente dos direitos creditórios líquido de PDD, mais ativo financeiro, dividido pelo saldo "
    "devedor daquela camada e de todas acima. São 159%/123%/110%/105% em CRI-II e 120,48%/109,89%/105,26% em CRI-VI. "
    "Um piso de 25% do PL e uma cobertura de 120% do saldo medem coisas diferentes e não se comparam diretamente.",
    tam=9, cor=PRETO, esp=2)
rodape(sl, "Fonte: prospectos definitivos, lâminas e termos de securitização. Subordinação mínima de CRI-I, III, IV e V: termos não localizados no acervo.")

# ═══════════════════════════════ 5 · Quadro 3 — cronograma e runoff
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Solfácil | Cronograma de Amortização e Runoff da Exposição",
          "CRIs somados + debênture, do 2T26 ao 3T38 · vencimento legal documentado; PMT contratual só existe para a 1ª série de CRI-II")

VEIC = ["CRI-I", "CRI-II", "CRI-III", "CRI-IV", "CRI-V", "CRI-VI", "DEB-I"]
ROT = {"CRI-I": "CRI I", "CRI-II": "CRI II", "CRI-III": "CRI III", "CRI-IV": "CRI IV",
       "CRI-V": "CRI V", "CRI-VI": "CRI VI", "DEB-I": "Debênture"}
TOT = {"CRI-I": 603.0, "CRI-II": 750.0, "CRI-III": 750.0, "CRI-IV": 450.0,
       "CRI-V": 470.6, "CRI-VI": 647.059, "DEB-I": 60.0}
saldo = dict(TOT)
tris, venc_tri, pmt_tri, runoff_v = [], [], [], {v: [] for v in VEIC}
for rot, sal_ini, vt, pmt, det in T.RUNOFF:
    for k, v in det.items():
        saldo[k] = max(round(saldo.get(k, 0.0) - v, 3), 0.0)
    tris.append(rot); venc_tri.append(round(vt, 1)); pmt_tri.append(round(pmt, 1))
    for v in VEIC:
        runoff_v[v].append(round(saldo.get(v, 0.0), 1))

txt(sl, MARGEM, Emu(1005840), Emu(7772400), Emu(182880),
    "CRONOGRAMA DE AMORTIZAÇÕES POR TRIMESTRE (R$ mm)", tam=8, neg=True, cor=CINZA_MED)
grafico(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, MARGEM, Emu(1188720), Emu(7772400), Emu(2011680),
        tris, [("Vencimento legal no trimestre", venc_tri),
               ("PMT contratual documentada · CRI-II 1ª série", pmt_tri)],
        [LARANJA, NAVY], unidade="R$ mm")
txt(sl, MARGEM, Emu(3383280), Emu(7772400), Emu(182880),
    "EXPOSIÇÃO TOTAL — SALDO NOMINAL POR INSTRUMENTO (R$ mm)", tam=8, neg=True, cor=CINZA_MED)
grafico(sl, XL_CHART_TYPE.AREA_STACKED, MARGEM, Emu(3566160), Emu(7772400), Emu(2286000),
        tris, [(ROT[v], runoff_v[v]) for v in VEIC],
        [NAVY, NAVY2, LARANJA, LARANJA_CLA, CINZA_MED, CINZA_CLARO, RGBColor(0x9D, 0xC3, 0xE6)],
        unidade="R$ mm")

marcos = [x for x in T.CRONOGRAMA if x[2] != "—"]
txt(sl, Emu(8412480), Emu(1005840), Emu(3322320), Emu(182880),
    "TRIMESTRES COM VENCIMENTO LEGAL", tam=8, neg=True, cor=CINZA_MED)
tabela(sl, Emu(8412480), Emu(1188720), Emu(3322320), Emu(3383280),
       ["Tri", "Vence (R$ mm)", "Saldo após (R$ mm)"],
       [[x[0], x[2], x[4]] for x in marcos],
       [0.22, 0.39, 0.39], tam=7, tam_cab=7.5, h_lin=Emu(146304), atr_col=True)
txt(sl, Emu(8412480), Emu(4663440), Emu(3322320), Emu(1188720),
    "O runoff acima assume que cada série fica integralmente em aberto até o vencimento legal — é limite "
    "superior de exposição.\n\n"
    "Onde o cronograma real existe, a diferença é grande: a 1ª série de CRI-II amortiza R$ 327,8 mm de forma "
    "linear entre o 2T26 e o 2T29, contra a barra única de R$ 487,5 mm que o vencimento legal projeta para o 2T29.",
    tam=8.5, cor=CINZA_ESC, esp=3)
rodape(sl, "Fonte: datas de vencimento das 34 séries de CRI e 2 séries de debênture; Anexo I do 2º aditamento ao TS da 2ª emissão Kanastra. "
           "Os FIDCs não entram: suas cotas não têm data de vencimento publicada. PMT das demais séries de CRI: a obter no Fundos.NET.")

# ═══════════════════════════════ 6 · Fontes e lacunas
sl = prs.slides.add_slide(VAZIO)
cabecalho(sl, "Solfácil | Fontes dos Quadros e Campos a Obter no CVM Fundos.NET",
          "O que sustenta cada linha e o que precisa ser preenchido antes do comitê")
LAC = [
    ("Rentabilidade YTD por classe de cota", "FIDCs I a VII", "Informe Mensal FIDC — campo de rentabilidade da cota", "a obter"),
    ("Rentabilidade YTD por série", "CRIs I a VI", "Informe Mensal CRI", "a obter"),
    ("Over 90 / carteira por operação", "CRIs I a VI", "Informe Mensal CRI", "a obter"),
    ("Saldo devedor atual por série", "CRIs I a VI", "Informe Mensal CRI; B3", "a obter"),
    ("Cronograma de PMT por série", "CRI-I, III, IV, V e VI", "Anexo I dos termos de securitização", "a obter"),
    ("Cronograma de amortização das cotas", "FIDCs I a VII", "Regulamento e suplementos de classe", "a obter"),
    ("Subordinação mínima", "CRI-I, III, IV e V", "Termos de securitização", "a obter"),
    ("Volume e data da cota sub júnior", "FIDCs IV, V, VI e VII", "Registro de ofertas de cotas", "a obter"),
    ("Remuneração das cotas sênior A, B e C", "FIDC IV", "Suplementos de classe", "a obter"),
    ("Dados completos", "FIDC VIII", "Em discussão; sem registro na CVM até a data-base", "reservado"),
    ("PL, composição por cota e over 90", "FIDCs I a VII", "Informe Mensal FIDC, 31/07/2026", "obtido"),
    ("Volume, séries, taxas, vencimentos e ISIN", "CRIs I a VI", "Prospectos, lâminas, comunicados e termos", "obtido"),
    ("Cronograma de PMT", "CRI-II, 1ª série", "Anexo I do 2º aditamento ao TS da 2ª Kanastra", "obtido"),
    ("Razões de cobertura", "CRI-II e CRI-VI", "Termos de securitização", "obtido"),
]
tabela(sl, MARGEM, Emu(1005840), UTIL, Emu(3931920),
       ["Campo", "Veículos", "Onde obter", "Status"], LAC,
       [0.30, 0.20, 0.38, 0.12], tam=8, tam_cab=8.5, h_lin=Emu(256032), atr_col=True, centro=False)
txt(sl, MARGEM, Emu(5120640), UTIL, Emu(731520),
    "Esta compilação não teve acesso à rede: nenhum campo foi consultado diretamente no CVM Fundos.NET. Os campos marcados como "
    "obtidos vieram dos documentos em PDF do acervo — prospectos, lâminas, comunicados, termos de securitização, escritura e "
    "relatório do agente fiduciário — e da consolidação de informes já disponível. Os demais estão destacados nos quadros, "
    "com o valor n/d, para preenchimento direto.",
    tam=9.5, cor=PRETO, esp=3)
rodape(sl, "Compilação de 24/08/2026 · data-base dos FIDCs 31/07/2026 · CRIs na última competência por operação.")

os.makedirs(OUTDIR, exist_ok=True)
caminho = os.path.join(OUTDIR, "Solfacil_Comite_Quadros_20260824_claude.pptx")
prs.save(caminho)
print("Deck salvo:", caminho)
print("Slides:", len(prs.slides._sldIdLst))
