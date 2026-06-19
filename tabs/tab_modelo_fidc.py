from __future__ import annotations

from dataclasses import dataclass, replace
from datetime import date, datetime, timedelta
from html import escape
from io import BytesIO
import re
from zipfile import ZIP_DEFLATED, ZipFile

import altair as alt
from openpyxl.styles import Alignment
import pandas as pd
import streamlit as st

from data_loader import load_model_inputs
from services.fidc_model import (
    AMORTIZATION_MODE_BULLET,
    AMORTIZATION_MODE_LINEAR,
    AMORTIZATION_MODE_NONE,
    AMORTIZATION_MODE_WORKBOOK,
    CREDIT_MODEL_MC3_CARTOES,
    CREDIT_MODEL_MIGRATION,
    CREDIT_MODEL_NPL90,
    INTEREST_PAYMENT_MODE_AFTER_GRACE,
    INTEREST_PAYMENT_MODE_BULLET,
    INTEREST_PAYMENT_MODE_PERIODIC,
    INTERPOLATION_METHOD_FLAT_FORWARD_252,
    INTERPOLATION_METHOD_SPLINE,
    RATE_MODE_POST_CDI,
    RATE_MODE_PRE,
    Premissas,
    annual_252_to_monthly_rate,
    build_flow,
    build_kpis,
    cession_discount_to_monthly_rate,
    monthly_to_annual_252_rate,
    monthly_rate_to_cession_discount,
)
from services.fidc_model.b3_curves import (
    B3CurveError,
    B3CurveSnapshot,
    DEFAULT_TAXASWAP_CURVE_CODE,
    fetch_latest_taxaswap_curve,
    fetch_taxaswap_curve,
)
from services.fidc_model.calendar import (
    B3CalendarError,
    B3CalendarSnapshot,
    build_b3_calendar_snapshot,
    fetch_b3_trading_calendar_html,
    merge_with_b3_market_holidays,
)


SNAPSHOT_CURVE_DATE = date(2025, 2, 25)
CURVE_SOURCE_B3_LATEST = "B3 - último pregão disponível"
CURVE_SOURCE_B3_DATE = "B3 - data escolhida"
CURVE_SOURCE_SNAPSHOT = "Curva local salva"
CURVE_SOURCE_OPTIONS = [CURVE_SOURCE_B3_LATEST, CURVE_SOURCE_B3_DATE, CURVE_SOURCE_SNAPSHOT]
INTERPOLATION_LABEL_B3 = "Flat Forward 252 (metodologia B3)"
INTERPOLATION_LABEL_SPLINE = "Spline"
INTERPOLATION_OPTIONS = [INTERPOLATION_LABEL_B3, INTERPOLATION_LABEL_SPLINE]
CALENDAR_SOURCE_B3_OFFICIAL = "B3 oficial + projeção explícita"
CALENDAR_SOURCE_B3_PROJECTED = "Calendário B3 projetado"
CALENDAR_SOURCE_SNAPSHOT = "Feriados locais salvos"
CALENDAR_SOURCE_OPTIONS = [CALENDAR_SOURCE_B3_OFFICIAL, CALENDAR_SOURCE_B3_PROJECTED, CALENDAR_SOURCE_SNAPSHOT]
DATE_SCHEDULE_WORKBOOK = "Grade semestral padrão"
DATE_SCHEDULE_MONTHLY = "Mensal pelo prazo informado"
PORTFOLIO_MODE_REVOLVING = "Revolvente"
PORTFOLIO_MODE_STATIC = "Carteira estática"
CESSION_INPUT_DISCOUNT = "Taxa de Cessão"
CESSION_INPUT_MONTHLY = "Taxa Mensal (%)"
CREDIT_LABEL_NPL90 = "NPL 90 + cobertura de provisão"
CREDIT_LABEL_MIGRATION = "Migração por faixas de atraso"
CREDIT_LABEL_MC3 = "MC3 Cartões (Over90 + Reneg 100% PDD)"
MODEL_VIEW_GERAL = "Modelo FIDC (geral)"
MODEL_VIEW_MC3 = "FIDC MC3 Cartões"
CREDIT_MODEL_LABELS = {
    CREDIT_LABEL_NPL90: CREDIT_MODEL_NPL90,
    CREDIT_LABEL_MC3: CREDIT_MODEL_MC3_CARTOES,
    CREDIT_LABEL_MIGRATION: CREDIT_MODEL_MIGRATION,
}
DEFAULT_VOLUME_CARTEIRA = 1_000_000_000.0
DEFAULT_TX_CESSAO_AM = 0.14
DEFAULT_CUSTO_ADM_AA = 0.0035
DEFAULT_CUSTO_MIN_MENSAL = 20_000.0
DEFAULT_PERDA_ESPERADA_AM = 0.0
DEFAULT_PERDA_INESPERADA_AM = 0.0
DEFAULT_PERDA_CICLO = 0.40
DEFAULT_NPL90_LAG_MESES = 3
DEFAULT_COBERTURA_NPL90 = 1.0
DEFAULT_RENEGOCIADO_PCT = 0.0
DEFAULT_MATURACAO_OVER90_CAP = 0.40
DEFAULT_LGD = 1.0
DEFAULT_ROLAGEM_ADIMPLENTE_1_30 = 0.0
DEFAULT_ROLAGEM_1_30_31_60 = 0.0
DEFAULT_ROLAGEM_31_60_61_90 = 0.0
DEFAULT_ROLAGEM_61_90_90_PLUS = 0.0
DEFAULT_RECUPERACAO_90_PLUS = 0.0
DEFAULT_WRITEOFF_90_PLUS = 0.0
DEFAULT_AGIO_AQUISICAO = 0.0
DEFAULT_EXCESSO_SPREAD_SENIOR_AM = 0.0
DEFAULT_PROP_SENIOR = 0.70
DEFAULT_PROP_MEZZ = 0.0
DEFAULT_PROP_SUB = 0.30
DEFAULT_TAXA_SENIOR = 0.016
DEFAULT_TAXA_MEZZ = 0.0
DEFAULT_PRAZO_ANOS = 2.0
DEFAULT_QTD_PARCELAS_MEDIA = 7.0
DEFAULT_PRAZO_RECEBIVEIS_MESES = 3.154420901487235
DEFAULT_CARENCIA_PRINCIPAL_MESES = 18.0
DEFAULT_SUBORDINACAO_MINIMA_REINVESTIMENTO = 0.30
MC3_PRESET_VERSION = "mc3-2026-06-19-duration-parcelas-v1"
DEFAULT_CURVE_START_YEAR = 2026
DEFAULT_SELIC_PERPETUAL_YEAR = 2028
DEFAULT_SELIC_AA_2026 = 0.13
DEFAULT_SELIC_AA_2027_ONWARD = 0.12
LABELS_COTAS = {
    "sen": "Cota sênior",
    "mezz": "Cota mezzanino",
    "sub": "Cota subordinada",
}
LABELS_COTAS_ABBR = {
    "sen": "SEN",
    "mezz": "MEZZ",
    "sub": "SUB",
}
# Paleta dos gráficos do modelo: somente tons de preto, laranja e cinza.
COLOR_BLACK = "#1a1a1a"
COLOR_GRAY = "#9c9c9c"
COLOR_ORANGE = "#f28e2b"
COLOR_DARK_ORANGE = "#b35c00"
HELP_CUSTO_ADM_GESTAO = (
    "Custo anual sobre o PL econômico do fundo no início do período; aplica-se o maior entre "
    "o valor mensal composto e o custo mínimo mensal."
)
HELP_CUSTO_MINIMO = (
    "Piso em R$/mês comparado ao custo percentual sobre o PL econômico; o motor usa o maior valor a cada período."
)
AMORTIZATION_LABELS = {
    "Cronograma padrão": AMORTIZATION_MODE_WORKBOOK,
    "Linear após carência": AMORTIZATION_MODE_LINEAR,
    "Bullet no vencimento": AMORTIZATION_MODE_BULLET,
    "Sem amortização programada": AMORTIZATION_MODE_NONE,
}
INTEREST_LABELS = {
    "Pago em todo período": INTEREST_PAYMENT_MODE_PERIODIC,
    "Pago após carência": INTEREST_PAYMENT_MODE_AFTER_GRACE,
    "Bullet no vencimento": INTEREST_PAYMENT_MODE_BULLET,
}


_MODEL_CSS = """
<style>
.fidc-model-header {
    border-bottom: 1px solid #dde3ea;
    margin: 0.15rem 0 1.05rem 0;
    padding-bottom: 1rem;
}

.fidc-model-kicker {
    color: #ff5a00;
    font-size: 0.72rem;
    font-weight: 700;
    letter-spacing: 0.08em;
    line-height: 1.2;
    margin-bottom: 0.35rem;
    text-transform: uppercase;
}

.fidc-model-title {
    color: #283241;
    font-size: 1.65rem;
    font-weight: 700;
    letter-spacing: 0;
    line-height: 1.16;
    margin: 0;
}

.fidc-model-copy {
    color: #6f7a87;
    font-size: 0.96rem;
    line-height: 1.58;
    margin-top: 0.55rem;
    max-width: 58rem;
}

.fidc-model-section-title {
    color: #283241;
    font-size: 1.02rem;
    font-weight: 700;
    margin: 1.15rem 0 0.55rem 0;
}

.fidc-model-kpi-grid {
    display: grid;
    gap: 0.65rem;
    grid-template-columns: repeat(auto-fit, minmax(168px, 1fr));
    margin: 0.4rem 0 0.55rem 0;
}

.fidc-model-kpi-card {
    background: #ffffff;
    border: 1px solid #dde3ea;
    border-radius: 8px;
    border-top: 3px solid #1f77b4;
    min-height: 6.1rem;
    padding: 0.75rem 0.8rem;
}

.fidc-model-kpi-card:nth-child(2n) {
    border-top-color: #ff5a00;
}

.fidc-model-kpi-label {
    align-items: center;
    color: #6f7a87;
    display: flex;
    font-size: 0.72rem;
    font-weight: 700;
    gap: 0.35rem;
    justify-content: space-between;
    letter-spacing: 0.05em;
    line-height: 1.25;
    margin-bottom: 0.35rem;
    text-transform: uppercase;
}

.fidc-model-tooltip {
    align-items: center;
    background: #eef3f8;
    border: 1px solid #cbd6e2;
    border-radius: 999px;
    color: #425166;
    cursor: help;
    display: inline-flex;
    flex: 0 0 auto;
    font-size: 0.68rem;
    font-weight: 700;
    height: 1rem;
    justify-content: center;
    letter-spacing: 0;
    line-height: 1;
    position: relative;
    text-transform: none;
    width: 1rem;
}

.fidc-model-tooltip::after {
    background: #202a36;
    border-radius: 6px;
    box-shadow: 0 8px 22px rgba(32, 42, 54, 0.18);
    color: #ffffff;
    content: attr(data-tooltip);
    font-size: 0.75rem;
    font-weight: 500;
    left: 50%;
    line-height: 1.35;
    max-width: 17rem;
    min-width: 13rem;
    opacity: 0;
    padding: 0.55rem 0.65rem;
    pointer-events: none;
    position: absolute;
    text-align: left;
    text-transform: none;
    top: calc(100% + 0.45rem);
    transform: translate(-50%, -0.15rem);
    transition: opacity 0.12s ease, transform 0.12s ease, visibility 0.12s ease;
    visibility: hidden;
    white-space: normal;
    z-index: 1000;
}

.fidc-model-tooltip:hover::after,
.fidc-model-tooltip:focus::after {
    opacity: 1;
    transform: translate(-50%, 0);
    visibility: visible;
}

.fidc-model-kpi-value {
    color: #202a36;
    font-size: 1.18rem;
    font-weight: 700;
    line-height: 1.18;
}

.fidc-model-kpi-context {
    color: #6f7a87;
    font-size: 0.78rem;
    line-height: 1.35;
    margin-top: 0.35rem;
}

[data-testid="stFormSubmitButton"] button,
[data-testid="stButton"] button[kind="primary"],
[data-testid="stDownloadButton"] button[kind="primary"] {
    background: #1f77b4 !important;
    border-color: #1f77b4 !important;
    color: #ffffff !important;
}

[data-testid="stFormSubmitButton"] button p,
[data-testid="stButton"] button[kind="primary"] p,
[data-testid="stDownloadButton"] button[kind="primary"] p {
    color: #ffffff !important;
}

@media (max-width: 760px) {
    .fidc-model-title {
        font-size: 1.4rem;
    }
}
</style>
"""


@st.cache_data
def _load_inputs(path: str):
    return load_model_inputs(path)


@dataclass(frozen=True)
class _SelectedCurve:
    source_label: str
    curve_code: str
    base_date: date
    curva_du: tuple[float, ...]
    curva_taxa_aa: tuple[float, ...]
    source_url: str
    retrieved_label: str
    content_sha256: str
    point_count: int
    first_du: int | None
    last_du: int | None
    raw_line_count: int | None = None

    @property
    def cache_key(self) -> tuple[str, str, str, str]:
        return (
            self.source_label,
            self.curve_code,
            self.base_date.isoformat(),
            self.content_sha256[:16],
        )


@dataclass(frozen=True)
class _SelectedCalendar:
    source_label: str
    feriados: tuple[date, ...]
    holiday_count: int
    first_holiday: date | None
    last_holiday: date | None
    official_years: tuple[int, ...] = ()
    projected_years: tuple[int, ...] = ()
    source_url: str = ""
    retrieved_label: str = ""
    content_sha256: str = ""

    @property
    def cache_key(self) -> tuple[str, int, str, str, str, str]:
        return (
            self.source_label,
            self.holiday_count,
            self.first_holiday.isoformat() if self.first_holiday else "",
            self.last_holiday.isoformat() if self.last_holiday else "",
            ",".join(str(year) for year in self.official_years),
            ",".join(str(year) for year in self.projected_years),
        )


@dataclass(frozen=True)
class _RevolvencyMetrics:
    portfolio_mode: str
    prazo_total_anos: float
    prazo_medio_recebiveis_meses: float
    giro_estimado: float
    carteira_originada_programatica: float
    reinvestimento_principal_total: float
    reinvestimento_excesso_total: float
    nova_originacao_total: float
    reinvestimento_bloqueado_subordinacao_total: float
    carteira_total_originada: float
    ead_maximo: float
    ead_medio_ponderado: float
    sub_final_sem_inadimplencia: float
    colchao_sem_perdas_sobre_originacao: float | None
    subordinacao_minima_reinvestimento: float
    perda_ciclo_calibrada: float | None
    perda_ciclo_calibrada_anual_equivalente: float | None
    perda_ciclo_calibrada_pos_lgd: float | None
    perda_ciclo_calibrada_excede_limite: bool


@st.cache_data(show_spinner=False, ttl=24 * 60 * 60)
def _load_b3_calendar_snapshot(start_year: int, end_year: int) -> B3CalendarSnapshot:
    html_text, content_hash = fetch_b3_trading_calendar_html()
    datas = [date(start_year, 1, 1), date(end_year, 12, 31)]
    return build_b3_calendar_snapshot(datas, html_text, content_hash=content_hash)


@st.cache_data(show_spinner=False, ttl=6 * 60 * 60)
def _load_b3_curve_for_date(date_iso: str, curve_code: str) -> B3CurveSnapshot:
    return fetch_taxaswap_curve(date.fromisoformat(date_iso), curve_code=curve_code)


@st.cache_data(show_spinner=False, ttl=6 * 60 * 60)
def _load_latest_b3_curve(start_date_iso: str, curve_code: str) -> B3CurveSnapshot:
    return fetch_latest_taxaswap_curve(start_date=date.fromisoformat(start_date_iso), curve_code=curve_code)


def _selected_curve_from_snapshot(inputs) -> _SelectedCurve:
    return _SelectedCurve(
        source_label=CURVE_SOURCE_SNAPSHOT,
        curve_code="PRE",
        base_date=SNAPSHOT_CURVE_DATE,
        curva_du=tuple(float(value) for value in inputs.curva_du),
        curva_taxa_aa=tuple(float(value) for value in inputs.curva_cdi),
        source_url="model_data.json",
        retrieved_label="curva local sem consulta externa",
        content_sha256="local-snapshot",
        point_count=len(inputs.curva_du),
        first_du=int(inputs.curva_du[0]) if inputs.curva_du else None,
        last_du=int(inputs.curva_du[-1]) if inputs.curva_du else None,
        raw_line_count=None,
    )


def _selected_curve_from_b3(snapshot: B3CurveSnapshot, source_label: str) -> _SelectedCurve:
    return _SelectedCurve(
        source_label=source_label,
        curve_code=snapshot.curve_code,
        base_date=snapshot.generated_at,
        curva_du=tuple(snapshot.curva_du),
        curva_taxa_aa=tuple(snapshot.curva_taxa_aa),
        source_url=snapshot.source_url,
        retrieved_label=snapshot.retrieved_at.strftime("%d/%m/%Y %H:%M:%S UTC"),
        content_sha256=snapshot.content_sha256,
        point_count=len(snapshot.points),
        first_du=snapshot.first_du,
        last_du=snapshot.last_du,
        raw_line_count=snapshot.raw_line_count,
    )


def _default_selic_rate_for_year(year: int) -> float:
    return DEFAULT_SELIC_AA_2026 if int(year) == DEFAULT_CURVE_START_YEAR else DEFAULT_SELIC_AA_2027_ONWARD


def _effective_selic_projection_for_dates(
    user_projection: tuple[tuple[int, float], ...],
    simulation_dates: list[datetime],
) -> tuple[tuple[int, float], ...]:
    projection = dict(user_projection)
    if not projection:
        raise ValueError("A curva de SELIC média para caixa precisa ter pelo menos um ano informado.")
    required_years = sorted({dt.year for dt in simulation_dates[1:]})
    effective = dict(projection)
    for year in required_years:
        if year < min(projection):
            effective[year] = projection[min(projection)]
        elif year not in projection:
            previous_years = [projection_year for projection_year in projection if projection_year <= year]
            if not previous_years:
                raise ValueError(f"Curva de SELIC média para caixa sem taxa para {year}.")
            effective[year] = projection[max(previous_years)]
    return tuple(sorted(effective.items()))


def _selic_year_label(year: int, projection_years: list[int] | tuple[int, ...]) -> str:
    if int(year) == max(projection_years):
        return f"{year} em diante"
    return str(year)


def _build_workbook_dates_from_start(template_dates: list[datetime], start: datetime) -> list[datetime]:
    if not template_dates:
        return [start]
    template_start = template_dates[0]
    return [
        _add_months(start, (dt.year - template_start.year) * 12 + (dt.month - template_start.month))
        for dt in template_dates
    ]


def _selected_calendar(
    inputs,
    source_label: str,
    b3_snapshot: B3CalendarSnapshot | None = None,
    datas: list[datetime] | None = None,
) -> _SelectedCalendar:
    flow_dates = datas or inputs.datas
    if source_label == CALENDAR_SOURCE_B3_OFFICIAL:
        if b3_snapshot is None:
            raise B3CalendarError("Snapshot de calendário B3 não foi carregado.")
        holidays = tuple(b3_snapshot.holidays)
        return _SelectedCalendar(
            source_label=source_label,
            feriados=holidays,
            holiday_count=len(holidays),
            first_holiday=holidays[0] if holidays else None,
            last_holiday=holidays[-1] if holidays else None,
            official_years=b3_snapshot.official_years,
            projected_years=b3_snapshot.projected_years,
            source_url=b3_snapshot.source_url,
            retrieved_label=b3_snapshot.retrieved_at.strftime("%d/%m/%Y %H:%M:%S UTC"),
            content_sha256=b3_snapshot.content_sha256,
        )
    if source_label == CALENDAR_SOURCE_B3_PROJECTED:
        holidays = tuple(merge_with_b3_market_holidays(inputs.feriados, flow_dates))
        projected_years = tuple(range(flow_dates[0].year, flow_dates[-1].year + 1)) if flow_dates else ()
    else:
        holidays = tuple(holiday.date() for holiday in inputs.feriados)
        projected_years = ()
    return _SelectedCalendar(
        source_label=source_label,
        feriados=holidays,
        holiday_count=len(holidays),
        first_holiday=holidays[0] if holidays else None,
        last_holiday=holidays[-1] if holidays else None,
        projected_years=projected_years,
    )


def _interpolation_method_from_label(label: str) -> str:
    return INTERPOLATION_METHOD_FLAT_FORWARD_252 if label == INTERPOLATION_LABEL_B3 else INTERPOLATION_METHOD_SPLINE


def _format_number_br(value: float | None, decimals: int = 2) -> str:
    if value is None:
        return "N/D"
    return f"{value:,.{decimals}f}".replace(",", "X").replace(".", ",").replace("X", ".")


def _format_percent(value: float | None, decimals: int = 2) -> str:
    if value is None:
        return "N/D"
    return f"{_format_number_br(value * 100.0, decimals)}%"


def _format_brl(value: float | None) -> str:
    if value is None:
        return "N/D"
    return f"R$ {_format_number_br(value, 2)}"


def _format_input_value(value: float, decimals: int = 2) -> str:
    return _format_number_br(value, decimals)


def _format_brl_input_value(value: float, decimals: int = 2) -> str:
    return f"R$ {_format_input_value(value, decimals)}"


def _format_percent_input_value(value: float, decimals: int = 2) -> str:
    return f"{_format_input_value(value, decimals)}%"


def _parse_br_number(value: str, *, field_name: str) -> float:
    text = str(value or "").strip()
    if not text:
        raise ValueError(f"Informe {field_name}.")
    normalized = text.replace("R$", "").replace("%", "").replace(" ", "")
    if "," in normalized:
        normalized = normalized.replace(".", "").replace(",", ".")
    try:
        return float(normalized)
    except ValueError as exc:
        raise ValueError(f"Valor inválido em {field_name}: {value}") from exc


def _parse_percent_for_visibility(value: str, *, default: float) -> float:
    try:
        return _parse_br_number(value, field_name="percentual") / 100.0
    except ValueError:
        return default


def _principal_wal_from_installments(qtd_parcelas: float) -> float:
    installments = max(int(round(float(qtd_parcelas))), 1)
    return (installments + 1.0) / 2.0


def _installment_macaulay_duration_months(qtd_parcelas: float, monthly_rate: float) -> float:
    installments = max(int(round(float(qtd_parcelas))), 1)
    rate = max(float(monthly_rate), 0.0)
    weighted_pv = 0.0
    total_pv = 0.0
    for month in range(1, installments + 1):
        outstanding_start = (installments - month + 1.0) / installments
        cash_flow = (1.0 / installments) + outstanding_start * rate
        pv = cash_flow / ((1.0 + rate) ** month)
        weighted_pv += month * pv
        total_pv += pv
    return weighted_pv / total_pv if total_pv > 0.0 else _principal_wal_from_installments(installments)


def _monthly_rate_from_discount_and_installment_duration(discount: float, qtd_parcelas: float) -> float:
    rate = cession_discount_to_monthly_rate(discount, _principal_wal_from_installments(qtd_parcelas))
    for _ in range(12):
        duration = _installment_macaulay_duration_months(qtd_parcelas, rate)
        rate = cession_discount_to_monthly_rate(discount, duration)
    return rate


def _safe_installment_duration_preview(qtd_parcelas_text: str, taxa_mensal_text: str) -> float:
    try:
        qtd_parcelas = _parse_br_number(qtd_parcelas_text, field_name="Quantidade média de parcelas")
    except ValueError:
        qtd_parcelas = DEFAULT_QTD_PARCELAS_MEDIA
    try:
        monthly_rate = _parse_br_number(taxa_mensal_text, field_name="Taxa Mensal (%)") / 100.0
    except ValueError:
        monthly_rate = DEFAULT_TX_CESSAO_AM
    return _installment_macaulay_duration_months(qtd_parcelas, monthly_rate)


def _format_raw_input_text(value: str, *, decimals: int, kind: str) -> str:
    parsed = _parse_br_number(value, field_name="valor")
    if kind == "brl":
        return _format_brl_input_value(parsed, decimals)
    if kind == "percent":
        return _format_percent_input_value(parsed, decimals)
    if kind == "number":
        return _format_input_value(parsed, decimals)
    raise ValueError(f"Tipo de input inválido: {kind}")


_INPUT_NORMALIZATION_SPECS = {
    "modelo_volume": (2, "brl"),
    "modelo_tx_cessao_desagio": (2, "percent"),
    "modelo_tx_cessao_mensal": (2, "percent"),
    "modelo_custo_adm_pct": (2, "percent"),
    "modelo_custo_min": (2, "brl"),
    "modelo_perda_esperada_pct": (2, "percent"),
    "modelo_perda_inesperada_pct": (2, "percent"),
    "modelo_perda_ciclo_pct": (2, "percent"),
    "modelo_npl90_lag_meses": (0, "number"),
    "modelo_cobertura_npl90_pct": (1, "percent"),
    "modelo_lgd_pct": (1, "percent"),
    "modelo_roll_adimplente_1_30_pct": (2, "percent"),
    "modelo_roll_1_30_31_60_pct": (2, "percent"),
    "modelo_roll_31_60_61_90_pct": (2, "percent"),
    "modelo_roll_61_90_90_plus_pct": (2, "percent"),
    "modelo_recuperacao_90_plus_pct": (2, "percent"),
    "modelo_writeoff_90_plus_pct": (2, "percent"),
    "modelo_agio_aquisicao_pct": (2, "percent"),
    "modelo_excesso_spread_senior": (2, "percent"),
    "modelo_prop_senior": (1, "percent"),
    "modelo_prop_mezz": (1, "percent"),
    "modelo_prop_sub": (1, "percent"),
    "modelo_taxa_senior": (2, "percent"),
    "modelo_taxa_mezz": (2, "percent"),
    "modelo_taxa_sub": (2, "percent"),
    "modelo_prazo_fidc_anos": (1, "number"),
    "modelo_qtd_parcelas_media": (1, "number"),
    "modelo_prazo_recebiveis_meses": (1, "number"),
    "modelo_prazo_senior_anos": (1, "number"),
    "modelo_prazo_mezz_anos": (1, "number"),
    "modelo_prazo_sub_anos": (1, "number"),
    "modelo_inicio_amort_senior": (0, "number"),
    "modelo_inicio_amort_mezz": (0, "number"),
}


def _normalize_model_input_values() -> None:
    for key, (decimals, kind) in _INPUT_NORMALIZATION_SPECS.items():
        if key not in st.session_state:
            continue
        try:
            st.session_state[key] = _format_raw_input_text(st.session_state[key], decimals=decimals, kind=kind)
        except ValueError:
            continue
    for key in list(st.session_state):
        if not str(key).startswith("modelo_selic_aa_"):
            continue
        try:
            st.session_state[key] = _format_raw_input_text(st.session_state[key], decimals=2, kind="percent")
        except ValueError:
            continue


def _apply_mc3_preset_defaults() -> None:
    if st.session_state.get("modelo_mc3_preset_version") == MC3_PRESET_VERSION:
        return
    preset_values = {
        "modelo_view": MODEL_VIEW_MC3,
        "modelo_taxa_cessao_input_mode": CESSION_INPUT_MONTHLY,
        "modelo_credit_model": CREDIT_LABEL_NPL90,
        "modelo_volume": _format_brl_input_value(DEFAULT_VOLUME_CARTEIRA, 2),
        "modelo_tx_cessao_mensal": _format_percent_input_value(DEFAULT_TX_CESSAO_AM * 100.0, 2),
        "modelo_custo_adm_pct": _format_percent_input_value(DEFAULT_CUSTO_ADM_AA * 100.0, 2),
        "modelo_custo_min": _format_brl_input_value(DEFAULT_CUSTO_MIN_MENSAL, 2),
        "modelo_perda_ciclo_pct": _format_percent_input_value(DEFAULT_PERDA_CICLO * 100.0, 2),
        "modelo_npl90_lag_meses": _format_input_value(DEFAULT_NPL90_LAG_MESES, 0),
        "modelo_cobertura_npl90_pct": _format_percent_input_value(DEFAULT_COBERTURA_NPL90 * 100.0, 1),
        "modelo_lgd_pct": _format_percent_input_value(DEFAULT_LGD * 100.0, 1),
        "modelo_mc3_maturacao_over90_cap": _format_percent_input_value(DEFAULT_MATURACAO_OVER90_CAP * 100.0, 1),
        "modelo_agio_aquisicao_pct": _format_percent_input_value(DEFAULT_AGIO_AQUISICAO * 100.0, 2),
        "modelo_excesso_spread_senior": _format_percent_input_value(DEFAULT_EXCESSO_SPREAD_SENIOR_AM * 100.0, 2),
        "modelo_prop_senior": _format_percent_input_value(DEFAULT_PROP_SENIOR * 100.0, 1),
        "modelo_prop_mezz": _format_percent_input_value(DEFAULT_PROP_MEZZ * 100.0, 1),
        "modelo_prop_sub": _format_percent_input_value(DEFAULT_PROP_SUB * 100.0, 1),
        "modelo_taxa_senior": _format_percent_input_value(DEFAULT_TAXA_SENIOR * 100.0, 2),
        "modelo_taxa_mezz": _format_percent_input_value(DEFAULT_TAXA_MEZZ * 100.0, 2),
        "modelo_taxa_sub": _format_percent_input_value(0.0, 2),
        "modelo_date_schedule": DATE_SCHEDULE_MONTHLY,
        "modelo_prazo_fidc_anos": _format_input_value(DEFAULT_PRAZO_ANOS, 1),
        "modelo_qtd_parcelas_media": _format_input_value(DEFAULT_QTD_PARCELAS_MEDIA, 1),
        "modelo_prazo_recebiveis_meses": _format_input_value(DEFAULT_PRAZO_RECEBIVEIS_MESES, 1),
        "modelo_portfolio_mode": PORTFOLIO_MODE_REVOLVING,
        "modelo_senior_mode": "Pós-fixada: CDI + spread aditivo",
        "modelo_mezz_mode": "Pré-fixada: taxa % a.a.",
        "modelo_sub_mode": "Residual",
        "modelo_amort_senior": "Linear após carência",
        "modelo_amort_mezz": "Sem amortização programada",
        "modelo_juros_senior": "Pago em todo período",
        "modelo_juros_mezz": "Pago em todo período",
        "modelo_prazo_senior_anos": _format_input_value(DEFAULT_PRAZO_ANOS, 1),
        "modelo_prazo_mezz_anos": _format_input_value(DEFAULT_PRAZO_ANOS, 1),
        "modelo_prazo_sub_anos": _format_input_value(DEFAULT_PRAZO_ANOS, 1),
        "modelo_inicio_amort_senior": _format_input_value(DEFAULT_CARENCIA_PRINCIPAL_MESES, 0),
        "modelo_inicio_amort_mezz": _format_input_value(0.0, 0),
    }
    for key, value in preset_values.items():
        st.session_state[key] = value
    st.session_state["modelo_mc3_preset_version"] = MC3_PRESET_VERSION
    for cache_key in ("modelo_fidc_signature", "modelo_fidc_periods", "modelo_fidc_kpis"):
        st.session_state.pop(cache_key, None)


def _add_months(dt: datetime, months: int) -> datetime:
    month_index = dt.month - 1 + months
    year = dt.year + month_index // 12
    month = month_index % 12 + 1
    month_lengths = (31, 29 if _is_leap_year(year) else 28, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31)
    day = min(dt.day, month_lengths[month - 1])
    return dt.replace(year=year, month=month, day=day)


def _is_leap_year(year: int) -> bool:
    return year % 4 == 0 and (year % 100 != 0 or year % 400 == 0)


def _build_monthly_dates(start: datetime, prazo_total_anos: float) -> list[datetime]:
    total_months = max(1, round(prazo_total_anos * 12.0))
    return [_add_months(start, month) for month in range(total_months + 1)]


def _build_simulation_dates(
    inputs,
    schedule_label: str,
    prazo_total_anos: float,
    data_inicial: datetime | None = None,
) -> list[datetime]:
    start = data_inicial or inputs.datas[0]
    if schedule_label == DATE_SCHEDULE_MONTHLY:
        return _build_monthly_dates(start, prazo_total_anos)
    return _build_workbook_dates_from_start(list(inputs.datas), start)


def _effective_cession_discount_after_premium(nominal_discount: float, acquisition_premium: float) -> float:
    return float(nominal_discount) - max(float(acquisition_premium), 0.0)


def _projection_years_for_term(start: datetime, prazo_total_anos: float) -> list[int]:
    return list(range(DEFAULT_CURVE_START_YEAR, DEFAULT_SELIC_PERPETUAL_YEAR + 1))


def _safe_term_years_from_text(value: str, fallback: float = DEFAULT_PRAZO_ANOS) -> float:
    try:
        parsed = _parse_br_number(value, field_name="Prazo total do FIDC (anos)")
    except ValueError:
        return fallback
    return parsed if parsed > 0 else fallback


def _label_for_value(mapping: dict[str, str], value: str) -> str:
    for label, mapped_value in mapping.items():
        if mapped_value == value:
            return label
    return next(iter(mapping))


def _ensure_session_option(key: str, options: list[str], default_index: int = 0) -> str:
    value = st.session_state.get(key)
    if value not in options:
        value = options[default_index]
        st.session_state[key] = value
    return value


def _ensure_session_date(key: str, default_value: date) -> date:
    value = st.session_state.get(key)
    if isinstance(value, datetime):
        value = value.date()
    if not isinstance(value, date):
        value = default_value
        st.session_state[key] = value
    return value


def _build_revolvency_metrics(
    *,
    premissas: Premissas,
    zero_default_results,
    portfolio_mode: str,
    calibrated_loss_cycle: float | None = None,
    calibrated_loss_exceeds_limit: bool = False,
) -> _RevolvencyMetrics:
    prazo_total_anos = float(premissas.prazo_fidc_anos or 0.0)
    prazo_medio_meses = max(float(premissas.prazo_medio_recebiveis_meses), 0.01)
    prazo_principal_meses = max(float(premissas.qtd_parcelas_media or premissas.prazo_medio_recebiveis_meses), 0.01)
    giro_estimado = prazo_total_anos * 12.0 / prazo_medio_meses
    if portfolio_mode == PORTFOLIO_MODE_STATIC:
        carteira_originada_programatica = premissas.volume
    else:
        eligible_months = max(prazo_total_anos * 12.0 - prazo_medio_meses, 0.0)
        carteira_originada_programatica = premissas.volume + premissas.volume * (eligible_months / prazo_principal_meses)
    reinvestimento_principal_total = sum(
        max(float(getattr(period, "reinvestimento_principal", 0.0)), 0.0) for period in zero_default_results[1:]
    )
    reinvestimento_excesso_total = sum(
        max(float(getattr(period, "reinvestimento_excesso", 0.0)), 0.0) for period in zero_default_results[1:]
    )
    nova_originacao_total = sum(
        max(float(getattr(period, "nova_originacao", 0.0)), 0.0) for period in zero_default_results[1:]
    )
    reinvestimento_bloqueado_subordinacao_total = sum(
        max(float(getattr(period, "reinvestimento_bloqueado_subordinacao", 0.0)), 0.0)
        for period in zero_default_results[1:]
    )
    has_motor_origination = any(hasattr(period, "nova_originacao") for period in zero_default_results[1:])
    if zero_default_results and has_motor_origination:
        carteira_total_originada = max(float(premissas.volume), 0.0) + nova_originacao_total
    else:
        carteira_total_originada = carteira_originada_programatica
    sub_final = max(float(zero_default_results[-1].pl_sub_jr), 0.0) if zero_default_results else 0.0
    colchao_sem_perdas = sub_final / carteira_total_originada if carteira_total_originada > 0 else None
    ead_values = [float(getattr(period, "carteira", premissas.volume)) for period in zero_default_results]
    ead_maximo = max(ead_values, default=float(premissas.volume))
    weighted_days = sum(max(float(getattr(period, "delta_dc", 0.0)), 0.0) for period in zero_default_results[1:])
    if weighted_days > 0.0:
        ead_medio = (
            sum(
                float(getattr(period, "carteira", premissas.volume)) * max(float(getattr(period, "delta_dc", 0.0)), 0.0)
                for period in zero_default_results[1:]
            )
            / weighted_days
        )
    else:
        ead_medio = float(premissas.volume)
    calibrated_loss_annual = (
        (1.0 + calibrated_loss_cycle) ** (12.0 / prazo_medio_meses) - 1.0
        if calibrated_loss_cycle is not None and prazo_medio_meses > 0.0
        else None
    )
    calibrated_loss_post_lgd = (
        calibrated_loss_cycle * max(float(getattr(premissas, "lgd", 1.0)), 0.0)
        if calibrated_loss_cycle is not None
        else None
    )
    return _RevolvencyMetrics(
        portfolio_mode=portfolio_mode,
        prazo_total_anos=prazo_total_anos,
        prazo_medio_recebiveis_meses=prazo_medio_meses,
        giro_estimado=giro_estimado,
        carteira_originada_programatica=carteira_originada_programatica,
        reinvestimento_principal_total=reinvestimento_principal_total,
        reinvestimento_excesso_total=reinvestimento_excesso_total,
        nova_originacao_total=nova_originacao_total,
        reinvestimento_bloqueado_subordinacao_total=reinvestimento_bloqueado_subordinacao_total,
        carteira_total_originada=carteira_total_originada,
        ead_maximo=ead_maximo,
        ead_medio_ponderado=ead_medio,
        sub_final_sem_inadimplencia=sub_final,
        colchao_sem_perdas_sobre_originacao=colchao_sem_perdas,
        subordinacao_minima_reinvestimento=max(float(premissas.subordinacao_minima_reinvestimento), 0.0),
        perda_ciclo_calibrada=calibrated_loss_cycle,
        perda_ciclo_calibrada_anual_equivalente=calibrated_loss_annual,
        perda_ciclo_calibrada_pos_lgd=calibrated_loss_post_lgd,
        perda_ciclo_calibrada_excede_limite=calibrated_loss_exceeds_limit,
    )


def _premissas_sem_perdas(premissas: Premissas) -> Premissas:
    return replace(
        premissas,
        inadimplencia=0.0,
        perda_esperada_am=0.0,
        perda_inesperada_am=0.0,
        perda_ciclo=0.0,
        rolagem_adimplente_1_30=0.0,
        rolagem_1_30_31_60=0.0,
        rolagem_31_60_61_90=0.0,
        rolagem_61_90_90_plus=0.0,
        recuperacao_90_plus=0.0,
        writeoff_90_plus=0.0,
    )


def _premissas_perda_ciclo_calibrada(premissas: Premissas, perda_ciclo: float) -> Premissas:
    return replace(
        _premissas_sem_perdas(premissas),
        modelo_credito=CREDIT_MODEL_NPL90,
        perda_ciclo=max(float(perda_ciclo), 0.0),
        recuperacao_90_plus=0.0,
        writeoff_90_plus=0.0,
    )


def _final_sub_for_loss_cycle(
    *,
    datas: list[datetime],
    feriados,
    curva_du,
    curva_taxa_aa,
    premissas: Premissas,
    interpolation_method: str,
    perda_ciclo: float,
) -> float:
    periods = build_flow(
        datas,
        feriados,
        curva_du,
        curva_taxa_aa,
        _premissas_perda_ciclo_calibrada(premissas, perda_ciclo),
        interpolation_method=interpolation_method,
    )
    return float(periods[-1].pl_sub_jr) if periods else 0.0


def _solve_calibrated_loss_cycle(
    *,
    datas: list[datetime],
    feriados,
    curva_du,
    curva_taxa_aa,
    premissas: Premissas,
    interpolation_method: str,
    max_loss_cycle: float = 1.0,
    tolerance: float = 1e-5,
    iterations: int = 36,
) -> tuple[float | None, bool]:
    if not datas:
        return None, False
    base_sub = _final_sub_for_loss_cycle(
        datas=datas,
        feriados=feriados,
        curva_du=curva_du,
        curva_taxa_aa=curva_taxa_aa,
        premissas=premissas,
        interpolation_method=interpolation_method,
        perda_ciclo=0.0,
    )
    if base_sub <= 0.0:
        return 0.0, False

    high_sub = _final_sub_for_loss_cycle(
        datas=datas,
        feriados=feriados,
        curva_du=curva_du,
        curva_taxa_aa=curva_taxa_aa,
        premissas=premissas,
        interpolation_method=interpolation_method,
        perda_ciclo=max_loss_cycle,
    )
    if high_sub > 0.0:
        return None, True

    low = 0.0
    high = max_loss_cycle
    for _ in range(iterations):
        mid = (low + high) / 2.0
        mid_sub = _final_sub_for_loss_cycle(
            datas=datas,
            feriados=feriados,
            curva_du=curva_du,
            curva_taxa_aa=curva_taxa_aa,
            premissas=premissas,
            interpolation_method=interpolation_method,
            perda_ciclo=mid,
        )
        if abs(mid_sub) <= max(abs(base_sub) * tolerance, 1.0):
            return mid, False
        if mid_sub > 0.0:
            low = mid
        else:
            high = mid
    return high, False


def _scheduled_origination_components(
    month_index: float,
    previous_month_index: float,
    premissas: Premissas,
    portfolio_mode: str,
) -> tuple[float, float, float, float]:
    initial_portfolio = max(float(premissas.volume), 0.0)
    if portfolio_mode == PORTFOLIO_MODE_STATIC:
        return initial_portfolio, 0.0, 0.0, initial_portfolio

    prazo_medio_meses = max(float(premissas.prazo_medio_recebiveis_meses), 0.01)
    prazo_total_anos = float(premissas.prazo_fidc_anos or 0.0)
    cutoff_month = max(prazo_total_anos * 12.0 - prazo_medio_meses, 0.0)
    current_capped_month = min(max(float(month_index), 0.0), cutoff_month)
    previous_capped_month = min(max(float(previous_month_index), 0.0), cutoff_month)
    new_origination_period = premissas.volume * max(current_capped_month - previous_capped_month, 0.0) / prazo_medio_meses
    new_origination_cumulative = premissas.volume * current_capped_month / prazo_medio_meses
    return initial_portfolio, new_origination_period, new_origination_cumulative, initial_portfolio + new_origination_cumulative


def _build_time_protection_frame(
    frame: pd.DataFrame,
    *,
    premissas: Premissas,
    portfolio_mode: str,
    scenario_label: str,
) -> pd.DataFrame:
    columns = ["indice", "data", "pl_sub_jr"]
    if "fluxo_remanescente_mezz" in frame.columns:
        columns.append("fluxo_remanescente_mezz")
    if "nova_originacao" in frame.columns:
        columns.append("nova_originacao")
    protection = frame[columns].copy()
    protection["carteira_inicial_considerada"] = max(float(premissas.volume), 0.0)
    previous_indices = protection["indice"].shift(1).fillna(0.0)
    components = [
        _scheduled_origination_components(month_index, previous_month_index, premissas, portfolio_mode)
        for month_index, previous_month_index in zip(protection["indice"], previous_indices)
    ]
    protection["nova_originacao_programatica"] = [values[1] for values in components]
    protection["nova_originacao_programatica_acumulada"] = [values[2] for values in components]
    if "nova_originacao" in protection.columns:
        protection["nova_originacao_motor"] = protection["nova_originacao"].clip(lower=0.0)
        protection["nova_originacao_estimada"] = protection["nova_originacao_motor"]
        protection["nova_originacao_acumulada"] = protection["nova_originacao_motor"].cumsum()
        protection["carteira_originada_acumulada"] = protection["carteira_inicial_considerada"] + protection[
            "nova_originacao_acumulada"
        ]
    else:
        protection["nova_originacao_motor"] = None
        protection["nova_originacao_estimada"] = protection["nova_originacao_programatica"]
        protection["nova_originacao_acumulada"] = protection["nova_originacao_programatica_acumulada"]
        protection["carteira_originada_acumulada"] = [values[3] for values in components]
    protection["prazo_medio_recebiveis_meses"] = float(premissas.prazo_medio_recebiveis_meses)
    if "fluxo_remanescente_mezz" in protection.columns:
        protection["residual_economico_fluxo"] = protection["fluxo_remanescente_mezz"]
    else:
        protection["residual_economico_fluxo"] = protection["pl_sub_jr"].diff().fillna(protection["pl_sub_jr"])
    protection["sub_disponivel"] = protection["pl_sub_jr"].clip(lower=0.0)
    protection["perda_maxima_suportada"] = protection.apply(
        lambda row: (
            row["sub_disponivel"] / row["carteira_originada_acumulada"]
            if row["indice"] > 0 and row["carteira_originada_acumulada"] > 0.0
            else None
        ),
        axis=1,
    )
    protection["serie"] = scenario_label
    protection["valor_pct"] = protection["perda_maxima_suportada"] * 100.0
    protection["valor_formatado"] = protection["perda_maxima_suportada"].map(_format_percent)
    protection["sub_formatada"] = protection["sub_disponivel"].map(_format_brl)
    protection["originada_formatada"] = protection["carteira_originada_acumulada"].map(_format_brl)
    protection["carteira_inicial_formatada"] = protection["carteira_inicial_considerada"].map(_format_brl)
    protection["nova_originacao_formatada"] = protection["nova_originacao_estimada"].map(_format_brl)
    protection["nova_originacao_acumulada_formatada"] = protection["nova_originacao_acumulada"].map(_format_brl)
    protection["nova_originacao_motor_formatada"] = protection["nova_originacao_motor"].map(_format_brl)
    protection["residual_fluxo_formatado"] = protection["residual_economico_fluxo"].map(_format_brl)
    protection["periodo"] = protection["data"].dt.strftime("%d/%m/%Y")
    protection["mes_fidc"] = protection["indice"].map(lambda value: f"Mês {int(value)}")
    return protection.dropna(subset=["perda_maxima_suportada"])


def _text_number_input(
    label: str,
    *,
    default: float,
    key: str,
    decimals: int = 2,
    help_text: str | None = None,
) -> str:
    return st.text_input(label, value=_format_input_value(default, decimals), key=key, help=help_text)


def _text_brl_input(
    label: str,
    *,
    default: float,
    key: str,
    decimals: int = 2,
    help_text: str | None = None,
) -> str:
    return st.text_input(label, value=_format_brl_input_value(default, decimals), key=key, help=help_text)


def _text_percent_input(
    label: str,
    *,
    default: float,
    key: str,
    decimals: int = 2,
    help_text: str | None = None,
) -> str:
    return st.text_input(label, value=_format_percent_input_value(default, decimals), key=key, help=help_text)


def _build_dataframe(results) -> pd.DataFrame:
    frame = pd.DataFrame([result.__dict__ for result in results])
    frame["data"] = pd.to_datetime(frame["data"])
    return frame


def _build_export_dataframe(frame: pd.DataFrame) -> pd.DataFrame:
    export = frame.copy()
    export["data"] = export["data"].dt.strftime("%d/%m/%Y")
    if "perda_carteira_despesa" in export.columns and "inadimplencia_despesa" in export.columns:
        export = export.drop(columns=["inadimplencia_despesa"])
    return export.rename(
        columns={
            "indice": "Índice",
            "data": "Data",
            "dc": "Dias corridos",
            "du": "Dias úteis",
            "delta_dc": "Delta dias corridos",
            "delta_du": "Delta dias úteis",
            "pre_di": "Pre DI",
            "taxa_senior": "Taxa sênior",
            "fra_senior": "FRA sênior",
            "taxa_mezz": "Taxa MEZZ",
            "fra_mezz": "FRA MEZZ",
            "carteira": "Carteira de recebíveis (início do período)",
            "ead_carteira": "EAD / saldo em risco (preço pago)",
            "fluxo_carteira": "Fluxo econômico da carteira",
            "taxa_selic_aa": "Taxa SELIC projetada (% a.a.)",
            "taxa_selic_periodo": "Taxa SELIC do período",
            "saldo_caixa_selic_inicio": "Saldo de caixa aplicado SELIC (início)",
            "principal_para_caixa_selic": "Principal direcionado para caixa SELIC",
            "rendimento_caixa_selic": "Rendimento do caixa SELIC",
            "fluxo_ativos_total": "Fluxo econômico total dos ativos",
            "pl_fidc": "PL econômico do veículo",
            "custos_adm": "Custos administrativos",
            "perda_esperada_despesa": "Provisão prospectiva do ciclo",
            "perda_inesperada_despesa": "Reforço de cobertura ou write-off descoberto",
            "perda_carteira_despesa": "Despesa de provisão/perda",
            "carteira_vencendo": "Carteira vencendo no período",
            "ead_vencendo": "EAD vencendo no período",
            "principal_inadimplente": "Principal inadimplente não recebido",
            "entrada_npl90": "Entrada em NPL 90+",
            "npl90_estoque_inicio": "Estoque NPL 90+ (início)",
            "npl90_estoque_fim": "Estoque NPL 90+ (fim)",
            "provisao_saldo_inicio": "Saldo de provisão (início)",
            "provisao_requerida": "Provisão requerida",
            "despesa_provisao": "Despesa de provisão",
            "provisao_saldo_fim": "Saldo de provisão (fim)",
            "cobertura_npl90": "Cobertura NPL 90+",
            "baixa_credito": "Baixa de crédito/write-off",
            "writeoff_descoberto": "Write-off descoberto pela provisão",
            "recuperacao_credito": "Recuperação de crédito",
            "bucket_adimplente": "Bucket adimplente",
            "bucket_1_30": "Bucket 1-30",
            "bucket_31_60": "Bucket 31-60",
            "bucket_61_90": "Bucket 61-90",
            "bucket_90_plus": "Bucket NPL 90+",
            "resultado_carteira_liquido": "Resultado líquido da carteira",
            "prazo_restante_reinvestimento_meses": "Prazo restante para reinvestimento (meses)",
            "reinvestimento_elegivel": "Reinvestimento elegível",
            "subordinacao_minima_reinvestimento": "Subordinação mínima para reinvestimento",
            "carteira_originada_acumulada": "Carteira originada acumulada",
            "capacidade_reinvestimento_subordinacao": "Caixa elegível para reinvestimento",
            "reinvestimento_bloqueado_subordinacao": "Bloqueio por trava",
            "aporte_subordinacao_minima": "Aporte SUB para trava mínima",
            "principal_recebido_carteira": "Principal recebido da carteira",
            "reinvestimento_principal": "Reinvestimento de principal",
            "reinvestimento_excesso": "Reinvestimento de excesso de caixa",
            "nova_originacao": "Nova originação",
            "carteira_fim": "Carteira ao fim do período",
            "caixa_nao_reinvestido": "Caixa não reinvestido",
            "saldo_caixa_selic_fim": "Saldo de caixa aplicado SELIC (fim)",
            "agio_aquisicao_despesa": "Ágio sobre face informado",
            "preco_pago_fator": "Preço pago / face",
            "tx_cessao_am_input": "Taxa mensal informada",
            "tx_cessao_am_piso": "Piso mensal CDI + spread SEN + excesso",
            "tx_cessao_am_aplicada": "Taxa mensal aplicada",
            "principal_senior": "Principal sênior",
            "juros_senior": "Juros sênior",
            "pmt_senior": "PMT sênior",
            "vp_pmt_senior": "VP PMT sênior",
            "pl_senior": "PL sênior",
            "fluxo_remanescente": "Fluxo remanescente após sênior",
            "principal_mezz": "Principal MEZZ",
            "juros_mezz": "Juros MEZZ",
            "pmt_mezz": "PMT MEZZ",
            "pl_mezz": "PL MEZZ",
            "fluxo_remanescente_mezz": "Fluxo remanescente após MEZZ",
            "principal_sub_jr": "Principal subordinada",
            "juros_sub_jr": "Juros subordinada",
            "pmt_sub_jr": "PMT subordinada",
            "pl_sub_jr": "Saldo residual júnior econômico",
            "subordinacao_pct": "Subordinação econômica",
            "colchao_originada_pct": "Colchão sobre carteira originada",
            "pl_sub_jr_modelo": "Saldo júnior histórico",
            "subordinacao_pct_modelo": "Subordinação histórica",
        }
    )


def _build_display_dataframe(export_frame: pd.DataFrame) -> pd.DataFrame:
    display = export_frame.copy()
    money_tokens = (
        "Carteira",
        "Fluxo",
        "PL",
        "Custos",
        "Perda",
        "Provisão",
        "NPL",
        "Bucket",
        "Baixa",
        "Recuperação",
        "Principal",
        "Juros",
        "PMT",
        "VP",
        "Saldo",
        "Resultado",
        "Reinvestimento",
        "Originação",
        "Caixa",
        "Rendimento",
        "Aporte",
        "EAD",
        "Ágio",
        "Capacidade",
    )
    percent_tokens = ("Pre DI", "Taxa", "FRA", "Subordinação", "Cobertura", "Preço pago / face", "Colchão")
    for column in display.columns:
        if column in {"Índice", "Dias corridos", "Dias úteis", "Delta dias corridos", "Delta dias úteis"}:
            display[column] = display[column].map(lambda value: _format_number_br(float(value), 0) if pd.notna(value) else "N/D")
        elif any(token in column for token in percent_tokens):
            display[column] = display[column].map(lambda value: _format_percent(float(value)) if pd.notna(value) else "N/D")
        elif any(token in column for token in money_tokens):
            display[column] = display[column].map(lambda value: _format_brl(float(value)) if pd.notna(value) else "N/D")
    return display


def _build_committee_timeline_dataframe(display_frame: pd.DataFrame) -> pd.DataFrame:
    committee_columns = [
        "Índice",
        "Data",
        "Carteira de recebíveis (início do período)",
        "Fluxo econômico total dos ativos",
        "PL econômico do veículo",
        "Despesa de provisão/perda",
        "Estoque NPL 90+ (fim)",
        "Provisão requerida",
        "Juros sênior",
        "PL sênior",
        "PL MEZZ",
        "Saldo residual júnior econômico",
        "Subordinação econômica",
        "Carteira originada acumulada",
        "Colchão sobre carteira originada",
        "Caixa elegível para reinvestimento",
        "Bloqueio por trava",
        "Aporte SUB para trava mínima",
        "Caixa não reinvestido",
    ]
    visible_columns = [column for column in committee_columns if column in display_frame.columns]
    return display_frame.loc[:, visible_columns].copy()


def _build_kpi_export_dataframe(kpis) -> pd.DataFrame:
    return pd.DataFrame(
        [
            {
                "Indicador": "Retorno anualizado da classe sênior (econômico)",
                "Valor": kpis.xirr_senior,
                "Definicao": "Taxa interna anual de retorno dos pagamentos projetados da classe sênior.",
            },
            {
                "Indicador": "Retorno anualizado da classe MEZZ (econômico)",
                "Valor": kpis.xirr_mezz,
                "Definicao": "Taxa interna anual de retorno dos pagamentos projetados da classe MEZZ.",
            },
            {
                "Indicador": "Retorno anualizado da júnior residual (econômico)",
                "Valor": kpis.xirr_sub_jr,
                "Definicao": "Não aplicável quando a série de pagamentos da subordinada não tem sinais válidos para calcular retorno interno.",
            },
            {
                "Indicador": "Duration econômica da sênior",
                "Valor": kpis.duration_senior_anos,
                "Definicao": "Prazo médio ponderado, em anos, dos pagamentos descontados da classe sênior.",
            },
            {
                "Indicador": "Pre DI equivalente na duration",
                "Valor": kpis.pre_di_duration,
                "Definicao": "Taxa Pre DI da curva interpolada no ponto correspondente ao duration da sênior.",
            },
            {
                "Indicador": "Excesso de retorno da júnior sobre o Pre DI",
                "Valor": kpis.taxa_retorno_sub_jr_cdi,
                "Definicao": "Excesso de retorno da subordinada sobre a taxa Pre DI no duration, quando a série permite o cálculo.",
            },
        ]
    )


def _build_curve_source_dataframe(
    selected_curve: _SelectedCurve,
    selected_calendar: _SelectedCalendar,
    interpolation_label: str,
) -> pd.DataFrame:
    return pd.DataFrame(
        [
            {"Campo": "Fonte", "Valor": selected_curve.source_label},
            {"Campo": "Curva", "Valor": selected_curve.curve_code},
            {"Campo": "Interpolação", "Valor": interpolation_label},
            {"Campo": "Data-base", "Valor": selected_curve.base_date.strftime("%d/%m/%Y")},
            {"Campo": "Pontos", "Valor": _format_number_br(selected_curve.point_count, 0)},
            {"Campo": "DU inicial", "Valor": selected_curve.first_du if selected_curve.first_du is not None else "N/D"},
            {"Campo": "DU final", "Valor": selected_curve.last_du if selected_curve.last_du is not None else "N/D"},
            {"Campo": "Calendário de dias úteis", "Valor": selected_calendar.source_label},
            {"Campo": "Feriados considerados", "Valor": _format_number_br(selected_calendar.holiday_count, 0)},
            {
                "Campo": "Anos oficiais B3",
                "Valor": ", ".join(str(year) for year in selected_calendar.official_years) or "N/D",
            },
            {
                "Campo": "Anos projetados",
                "Valor": ", ".join(str(year) for year in selected_calendar.projected_years) or "N/D",
            },
            {
                "Campo": "Primeiro feriado",
                "Valor": selected_calendar.first_holiday.strftime("%d/%m/%Y") if selected_calendar.first_holiday else "N/D",
            },
            {
                "Campo": "Último feriado",
                "Valor": selected_calendar.last_holiday.strftime("%d/%m/%Y") if selected_calendar.last_holiday else "N/D",
            },
            {"Campo": "URL calendário", "Valor": selected_calendar.source_url or "N/D"},
            {"Campo": "Calendário baixado em", "Valor": selected_calendar.retrieved_label or "N/D"},
            {"Campo": "SHA-256 calendário", "Valor": selected_calendar.content_sha256 or "N/D"},
            {"Campo": "URL/Origem", "Valor": selected_curve.source_url},
            {"Campo": "Baixado em", "Valor": selected_curve.retrieved_label},
            {"Campo": "SHA-256", "Valor": selected_curve.content_sha256},
        ]
    )


def _build_revolvency_export_dataframe(metrics: _RevolvencyMetrics) -> pd.DataFrame:
    return pd.DataFrame(
        [
            {"Indicador": "Prazo total do FIDC (anos)", "Valor": metrics.prazo_total_anos},
            {"Indicador": "Prazo médio dos recebíveis (meses)", "Valor": metrics.prazo_medio_recebiveis_meses},
            {"Indicador": "Carteira originada programática", "Valor": metrics.carteira_originada_programatica},
            {"Indicador": "Reinvestimento de principal", "Valor": metrics.reinvestimento_principal_total},
            {"Indicador": "Reinvestimento de excesso de spread", "Valor": metrics.reinvestimento_excesso_total},
            {"Indicador": "Nova originação efetiva do motor", "Valor": metrics.nova_originacao_total},
            {"Indicador": "Bloqueio por trava", "Valor": metrics.reinvestimento_bloqueado_subordinacao_total},
            {"Indicador": "Carteira originada efetiva", "Valor": metrics.carteira_total_originada},
            {"Indicador": "SUB final sem perdas", "Valor": metrics.sub_final_sem_inadimplencia},
            {"Indicador": "Subordinação mínima para reinvestimento", "Valor": metrics.subordinacao_minima_reinvestimento},
            {"Indicador": "Colchão sem perdas sobre carteira originada", "Valor": metrics.colchao_sem_perdas_sobre_originacao},
        ]
    )


def _build_time_protection_export_dataframe(protection_frame: pd.DataFrame) -> pd.DataFrame:
    export = protection_frame.copy()
    export["data"] = export["data"].dt.strftime("%d/%m/%Y")
    return export.rename(
        columns={
            "indice": "Mês do FIDC",
            "data": "Data",
            "pl_sub_jr": "SUB residual",
            "carteira_inicial_considerada": "Carteira inicial considerada",
            "nova_originacao_estimada": "Nova originação estimada",
            "nova_originacao_acumulada": "Nova originação acumulada",
            "nova_originacao_motor": "Nova originação econômica do motor",
            "prazo_medio_recebiveis_meses": "Prazo médio usado (meses)",
            "carteira_originada_acumulada": "Carteira originada acumulada",
            "residual_economico_fluxo": "Residual econômico do fluxo",
            "sub_disponivel": "SUB disponível",
            "perda_maxima_suportada": "Colchão de proteção sobre carteira originada",
            "serie": "Cenário",
        }
    )[
        [
            "Mês do FIDC",
            "Data",
            "Cenário",
            "Carteira inicial considerada",
            "Nova originação estimada",
            "Nova originação acumulada",
            "Nova originação econômica do motor",
            "Prazo médio usado (meses)",
            "Carteira originada acumulada",
            "SUB residual",
            "SUB disponível",
            "Residual econômico do fluxo",
            "Colchão de proteção sobre carteira originada",
        ]
    ]


def _build_premissas_summary_dataframe(
    *,
    premissas: Premissas,
    taxa_cessao_input_mode: str,
    tx_cessao_desagio: float,
    tx_cessao_aa_equivalente: float,
    credit_model_label: str,
    portfolio_mode_label: str,
    date_schedule_label: str,
    data_inicial: date,
    selected_curve: _SelectedCurve,
    selected_calendar: _SelectedCalendar,
    interpolation_label: str,
    senior_mode_label: str,
    mezz_mode_label: str,
    sub_mode_label: str,
    senior_amort_label: str,
    mezz_amort_label: str,
    senior_interest_label: str,
    mezz_interest_label: str,
    user_selic_aa_por_ano: tuple[tuple[int, float], ...],
) -> pd.DataFrame:
    rows: list[dict[str, str]] = []

    def add(categoria: str, premissa: str, valor: str, observacao: str = "") -> None:
        rows.append({"Categoria": categoria, "Premissa": premissa, "Valor": valor, "Observação": observacao})

    add("Carteira", "Volume inicial", _format_brl(premissas.volume), "Valor de face dos recebíveis no mês zero.")
    add("Carteira", "Modo da carteira", portfolio_mode_label, "Define se há reciclagem de principal/excesso de spread.")
    add("Carteira", "Prazo total do FIDC", f"{_format_number_br(premissas.prazo_fidc_anos or 0.0, 1)} anos")
    if premissas.qtd_parcelas_media is not None:
        add("Carteira", "Quantidade média de parcelas", f"{_format_number_br(premissas.qtd_parcelas_media, 1)} parcelas")
    add(
        "Carteira",
        "Prazo médio dos recebíveis",
        f"{_format_number_br(premissas.prazo_medio_recebiveis_meses, 2)} meses",
        "Duration econômica Macaulay calculada pela quantidade média de parcelas e taxa mensal da carteira.",
    )
    add("Carteira", "Entrada da taxa da carteira", taxa_cessao_input_mode)
    add("Carteira", "Taxa de cessão equivalente", _format_percent(tx_cessao_desagio), "Deságio no prazo médio.")
    add("Carteira", "Taxa mensal efetiva", _format_percent(premissas.tx_cessao_am))
    add("Carteira", "Taxa anual equivalente base 252", _format_percent(tx_cessao_aa_equivalente))
    add("Carteira", "Ágio sobre face", _format_percent(premissas.agio_aquisicao))
    add("Carteira", "Excesso de spread sobre SEN", _format_percent(premissas.excesso_spread_senior_am), "Piso mensal adicional.")
    add("Custos", "Administração/gestão", _format_percent(premissas.custo_adm_aa), "Percentual anual sobre PL econômico.")
    add("Custos", "Custo mínimo mensal", _format_brl(premissas.custo_min))
    add("Crédito", "Metodologia", credit_model_label)
    add("Crédito", "NPL 90+ esperado por ciclo", _format_percent(premissas.perda_ciclo))
    add("Crédito", "Lag até NPL 90+", f"{premissas.npl90_lag_meses} meses")
    add("Crédito", "Cobertura mínima NPL 90+", _format_percent(premissas.cobertura_minima_npl90))
    add("Crédito", "LGD econômica", _format_percent(premissas.lgd))
    add("Crédito", "Rolagem atual → 1-30", _format_percent(premissas.rolagem_adimplente_1_30))
    add("Crédito", "Rolagem 1-30 → 31-60", _format_percent(premissas.rolagem_1_30_31_60))
    add("Crédito", "Rolagem 31-60 → 61-90", _format_percent(premissas.rolagem_31_60_61_90))
    add("Crédito", "Rolagem 61-90 → 90+", _format_percent(premissas.rolagem_61_90_90_plus))
    add("Crédito", "Recuperação 90+", _format_percent(premissas.recuperacao_90_plus))
    add("Crédito", "Write-off 90+", _format_percent(premissas.writeoff_90_plus))
    add("Crédito", "Renegociado", _format_percent(premissas.renegociado_pct), "Premissa MC3, quando aplicável.")
    add("Crédito", "Teto maturação Over90", _format_percent(premissas.maturacao_over90_cap), "Premissa MC3, quando aplicável.")
    add("Estrutura", "PL SEN", _format_percent(premissas.proporcao_senior))
    add("Estrutura", "PL MEZZ", _format_percent(premissas.proporcao_mezz))
    add("Estrutura", "PL SUB", _format_percent(premissas.proporcao_sub_jr))
    add(
        "Estrutura",
        "Trava mínima de reinvestimento",
        _format_percent(premissas.subordinacao_minima_reinvestimento),
        "Piso estrutural de SUB / PL FIDC; a carteira originada acumulada fica como métrica de colchão econômico.",
    )
    add("Remuneração", "SEN", f"{senior_mode_label}; taxa/spread {_format_percent(premissas.taxa_senior)}")
    add("Remuneração", "MEZZ", f"{mezz_mode_label}; taxa/spread {_format_percent(premissas.taxa_mezz)}")
    add("Remuneração", "SUB", f"{sub_mode_label}; taxa-alvo {_format_percent(premissas.taxa_sub_jr)}")
    add("Waterfall", "Amortização SEN", senior_amort_label)
    add("Waterfall", "Amortização MEZZ", mezz_amort_label)
    add("Waterfall", "Juros SEN", senior_interest_label)
    add("Waterfall", "Juros MEZZ", mezz_interest_label)
    add("Waterfall", "Início amortização SEN", f"Mês {premissas.inicio_amortizacao_senior_meses}")
    add("Waterfall", "Início amortização MEZZ", f"Mês {premissas.inicio_amortizacao_mezz_meses}")
    add("Datas", "Data inicial", data_inicial.strftime("%d/%m/%Y"))
    add("Datas", "Cronograma do fluxo", date_schedule_label)
    add("Curva", "Fonte da curva", selected_curve.source_label)
    add("Curva", "Data-base da curva", selected_curve.base_date.strftime("%d/%m/%Y"))
    add("Curva", "Interpolação", interpolation_label)
    add("Curva", "Calendário de dias úteis", selected_calendar.source_label)
    for year, rate in user_selic_aa_por_ano:
        add("Caixa runoff", f"SELIC média {_selic_year_label(year, [item[0] for item in user_selic_aa_por_ano])}", _format_percent(rate))
    return pd.DataFrame(rows)


def _build_model_dashboard_excel_bytes(
    *,
    export_frame: pd.DataFrame,
    kpi_cards: list[dict[str, str]],
    revolvency_cards: list[dict[str, str]],
    premissas_summary_df: pd.DataFrame,
    memory_df: pd.DataFrame,
    curve_source_df: pd.DataFrame,
    revolvency_export_df: pd.DataFrame,
    protection_export_df: pd.DataFrame,
    balance_chart_df: pd.DataFrame,
    loss_chart_df: pd.DataFrame,
    protection_chart_df: pd.DataFrame,
) -> bytes:
    from openpyxl.chart import AreaChart, LineChart, Reference
    from openpyxl.styles import Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        export_frame.to_excel(writer, index=False, sheet_name="timeline")
        pd.DataFrame(kpi_cards).to_excel(writer, index=False, sheet_name="cards_resumo")
        premissas_summary_df.to_excel(writer, index=False, sheet_name="premissas_resumo")
        memory_df.to_excel(writer, index=False, sheet_name="memoria_calculo")
        curve_source_df.to_excel(writer, index=False, sheet_name="fonte_curva")
        revolvency_export_df.to_excel(writer, index=False, sheet_name="capacidade_perda")
        protection_export_df.to_excel(writer, index=False, sheet_name="protecao_tempo")

        workbook = writer.book
        dashboard = workbook.create_sheet("dashboard", 0)
        chart_data = workbook.create_sheet("graficos_dados")

        fonts = {
            "title": Font(name="Calibri", size=18, bold=True, color="1F1F1F"),
            "section": Font(name="Calibri", size=12, bold=True, color="1F1F1F"),
            "header": Font(name="Calibri", size=9, bold=True, color="FFFFFF"),
            "body": Font(name="Calibri", size=9, color="1F1F1F"),
            "muted": Font(name="Calibri", size=8, color="6B7280"),
            "card_label": Font(name="Calibri", size=8, bold=True, color="6B7280"),
            "card_value": Font(name="Calibri", size=14, bold=True, color="1F1F1F"),
        }
        fills = {
            "header": PatternFill("solid", fgColor="1F1F1F"),
            "orange": PatternFill("solid", fgColor="EC7000"),
            "soft": PatternFill("solid", fgColor="F7F7F7"),
            "white": PatternFill("solid", fgColor="FFFFFF"),
        }
        border = Border(
            left=Side(style="thin", color="E0E0E0"),
            right=Side(style="thin", color="E0E0E0"),
            top=Side(style="thin", color="E0E0E0"),
            bottom=Side(style="thin", color="E0E0E0"),
        )

        _style_dataframe_sheets(workbook, fonts, fills, border)
        _build_dashboard_sheet(
            dashboard,
            kpi_cards=kpi_cards,
            revolvency_cards=revolvency_cards,
            premissas_summary_df=premissas_summary_df,
            fonts=fonts,
            fills=fills,
            border=border,
        )
        chart_ranges = _write_model_chart_sources(chart_data, balance_chart_df, loss_chart_df, protection_chart_df)
        _style_dataframe_sheet(chart_data, fonts, fills, border)
        _add_model_excel_charts(
            dashboard,
            chart_data,
            chart_ranges,
            AreaChart=AreaChart,
            LineChart=LineChart,
            Reference=Reference,
        )
        dashboard.freeze_panes = "A4"
        for col in range(1, 15):
            dashboard.column_dimensions[get_column_letter(col)].width = 12.5

    return output.getvalue()


def _style_dataframe_sheets(workbook, fonts: dict[str, object], fills: dict[str, object], border: object) -> None:
    for worksheet in workbook.worksheets:
        if worksheet.title == "dashboard":
            continue
        _style_dataframe_sheet(worksheet, fonts, fills, border)


def _style_dataframe_sheet(worksheet, fonts: dict[str, object], fills: dict[str, object], border: object) -> None:
    if worksheet.max_row < 1 or (worksheet.max_row == 1 and worksheet.max_column == 1 and worksheet["A1"].value is None):
        return
    worksheet.freeze_panes = "A2"
    worksheet.auto_filter.ref = worksheet.dimensions
    headers = [str(worksheet.cell(1, col).value or "") for col in range(1, worksheet.max_column + 1)]
    for cell in worksheet[1]:
        cell.fill = fills["header"]
        cell.font = fonts["header"]
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border
    money_tokens = ("Carteira", "Fluxo", "PL", "Custos", "Perda", "Provisão", "NPL", "Principal", "Juros", "PMT", "Saldo", "Resultado", "Originação", "Caixa", "EAD", "Valor")
    pct_tokens = ("Taxa", "FRA", "Subordinação", "Cobertura", "Colchão", "Preço pago / face")
    for col_idx, header in enumerate(headers, start=1):
        values = [worksheet.cell(row, col_idx).value for row in range(1, min(worksheet.max_row, 80) + 1)]
        non_empty_values = [value for value in values if value is not None]
        width = min(max(len(str(value)) for value in non_empty_values) + 2, 48) if non_empty_values else 10
        worksheet.column_dimensions[worksheet.cell(1, col_idx).column_letter].width = max(width, 10)
        for row_idx in range(2, worksheet.max_row + 1):
            cell = worksheet.cell(row_idx, col_idx)
            cell.font = fonts["body"]
            cell.border = border
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            if isinstance(cell.value, (int, float)):
                if any(token in header for token in pct_tokens):
                    cell.number_format = "0.00%"
                elif any(token in header for token in money_tokens):
                    cell.number_format = 'R$ #,##0'
                else:
                    cell.number_format = "#,##0.00"


def _build_dashboard_sheet(
    worksheet,
    *,
    kpi_cards: list[dict[str, str]],
    revolvency_cards: list[dict[str, str]],
    premissas_summary_df: pd.DataFrame,
    fonts: dict[str, object],
    fills: dict[str, object],
    border: object,
) -> None:
    worksheet.sheet_view.showGridLines = False
    worksheet.merge_cells("A1:N1")
    worksheet["A1"] = "Modelagem FIDC - Dashboard"
    worksheet["A1"].font = fonts["title"]
    worksheet["A1"].alignment = Alignment(vertical="center")
    worksheet.merge_cells("A2:N2")
    worksheet["A2"] = "Cards, gráficos e premissas resumidas da tela de Modelagem. Abas auxiliares preservam os dados numéricos auditáveis."
    worksheet["A2"].font = fonts["muted"]

    _section_title(worksheet, 4, "Premissas resumidas", fonts)
    premissas_preview = premissas_summary_df.head(10)
    headers = ["Categoria", "Premissa", "Valor", "Observação"]
    for col_idx, header in enumerate(headers, start=1):
        cell = worksheet.cell(5, col_idx, header)
        cell.fill = fills["header"]
        cell.font = fonts["header"]
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border
    for row_offset, row in enumerate(premissas_preview.itertuples(index=False), start=6):
        for col_idx, value in enumerate(row, start=1):
            cell = worksheet.cell(row_offset, col_idx, value)
            cell.fill = fills["soft"] if row_offset % 2 == 0 else fills["white"]
            cell.font = fonts["body"]
            cell.border = border
            cell.alignment = Alignment(vertical="top", wrap_text=True)
    worksheet.merge_cells("E5:N15")
    note = worksheet["E5"]
    note.value = "A planilha `premissas_resumo` contém a lista completa de premissas resumidas. Esta aba mostra uma prévia para leitura executiva."
    note.fill = fills["soft"]
    note.font = fonts["muted"]
    note.alignment = Alignment(vertical="top", wrap_text=True)
    note.border = border

    _section_title(worksheet, 17, "Resumo econômico", fonts)
    _add_cards_grid(worksheet, kpi_cards, start_row=18, fonts=fonts, fills=fills, border=border)
    _section_title(worksheet, 30, "Capacidade de perda e denominadores", fonts)
    _add_cards_grid(worksheet, revolvency_cards, start_row=31, fonts=fonts, fills=fills, border=border)
    _section_title(worksheet, 42, "Evolução do saldo das cotas", fonts)
    _section_title(worksheet, 64, "Perda da carteira", fonts)
    _section_title(worksheet, 64, "Proteção da estrutura", fonts, start_col=8)


def _section_title(worksheet, row: int, text: str, fonts: dict[str, object], *, start_col: int = 1) -> None:
    cell = worksheet.cell(row, start_col, text)
    cell.font = fonts["section"]
    cell.alignment = Alignment(vertical="center")


def _add_cards_grid(
    worksheet,
    cards: list[dict[str, str]],
    *,
    start_row: int,
    fonts: dict[str, object],
    fills: dict[str, object],
    border: object,
) -> None:
    card_width = 4
    card_height = 4
    gap = 1
    for idx, card in enumerate(cards):
        block_row = start_row + (idx // 3) * (card_height + 1)
        block_col = 1 + (idx % 3) * (card_width + gap)
        _add_card_block(worksheet, block_row, block_col, card_width, card, fonts=fonts, fills=fills, border=border)


def _add_card_block(
    worksheet,
    row: int,
    col: int,
    width: int,
    card: dict[str, str],
    *,
    fonts: dict[str, object],
    fills: dict[str, object],
    border: object,
) -> None:
    worksheet.merge_cells(start_row=row, start_column=col, end_row=row, end_column=col + width - 1)
    worksheet.cell(row, col).fill = fills["orange"]
    for offset, value_key, font_key in [(1, "label", "card_label"), (2, "value", "card_value"), (3, "context", "muted")]:
        worksheet.merge_cells(start_row=row + offset, start_column=col, end_row=row + offset, end_column=col + width - 1)
        cell = worksheet.cell(row + offset, col, str(card.get(value_key, "")))
        cell.font = fonts[font_key]
        cell.fill = fills["soft"]
        cell.alignment = Alignment(vertical="center", wrap_text=True)
    for row_idx in range(row, row + 4):
        for col_idx in range(col, col + width):
            worksheet.cell(row_idx, col_idx).border = border


def _write_model_chart_sources(
    worksheet,
    balance_chart_df: pd.DataFrame,
    loss_chart_df: pd.DataFrame,
    protection_chart_df: pd.DataFrame,
) -> dict[str, tuple[int, int, int, int]]:
    balance = _balance_chart_wide(balance_chart_df)
    loss = _long_percent_chart_wide(loss_chart_df)
    protection = _long_percent_chart_wide(protection_chart_df)
    ranges = {}
    ranges["balance"] = _write_dataframe_at(worksheet, balance, start_row=1, start_col=1)
    ranges["loss"] = _write_dataframe_at(worksheet, loss, start_row=ranges["balance"][0] + ranges["balance"][2] + 3, start_col=1)
    ranges["protection"] = _write_dataframe_at(
        worksheet,
        protection,
        start_row=ranges["loss"][0] + ranges["loss"][2] + 3,
        start_col=1,
    )
    return ranges


def _write_dataframe_at(worksheet, frame: pd.DataFrame, *, start_row: int, start_col: int) -> tuple[int, int, int, int]:
    from openpyxl.utils.dataframe import dataframe_to_rows

    for row_offset, row in enumerate(dataframe_to_rows(frame, index=False, header=True), start=0):
        for col_offset, value in enumerate(row, start=0):
            worksheet.cell(start_row + row_offset, start_col + col_offset, value)
    return start_row, start_col, len(frame) + 1, len(frame.columns)


def _balance_chart_wide(chart_df: pd.DataFrame) -> pd.DataFrame:
    if chart_df.empty:
        return pd.DataFrame({"Mês": [], "Data": []})
    wide = (
        chart_df.pivot_table(
            index=["indice", "periodo"],
            columns="classe",
            values="valor_milhoes",
            aggfunc="sum",
            fill_value=0.0,
        )
        .reset_index()
        .rename(columns={"indice": "Mês", "periodo": "Data"})
    )
    for column in [LABELS_COTAS["sen"], LABELS_COTAS["mezz"], LABELS_COTAS["sub"], "Déficit econômico"]:
        if column not in wide.columns:
            wide[column] = 0.0
    return wide[["Mês", "Data", LABELS_COTAS["sen"], LABELS_COTAS["mezz"], LABELS_COTAS["sub"], "Déficit econômico"]]


def _long_percent_chart_wide(chart_df: pd.DataFrame) -> pd.DataFrame:
    if chart_df.empty:
        return pd.DataFrame({"Mês": [], "Data": []})
    wide = (
        chart_df.pivot_table(
            index=["indice", "periodo"],
            columns="serie",
            values="valor",
            aggfunc="sum",
            fill_value=0.0,
        )
        .reset_index()
        .rename(columns={"indice": "Mês", "periodo": "Data"})
    )
    series_columns = [column for column in wide.columns if column not in {"Mês", "Data"}]
    return wide[["Mês", "Data", *series_columns]]


def _add_model_excel_charts(
    dashboard,
    chart_data,
    ranges: dict[str, tuple[int, int, int, int]],
    *,
    AreaChart,
    LineChart,
    Reference,
) -> None:
    _add_area_chart(
        dashboard,
        chart_data,
        ranges["balance"],
        AreaChart=AreaChart,
        Reference=Reference,
        anchor="A43",
        title="Evolução do saldo das cotas",
        y_axis_title="R$ milhões",
        colors=["1A1A1A", "9C9C9C", "F28E2B", "B35C00"],
    )
    _add_line_chart(
        dashboard,
        chart_data,
        ranges["loss"],
        LineChart=LineChart,
        Reference=Reference,
        anchor="A65",
        title="Perda da carteira",
        y_axis_title="Perda da carteira (%)",
        colors=["F28E2B"],
        percent_axis=True,
    )
    _add_line_chart(
        dashboard,
        chart_data,
        ranges["protection"],
        LineChart=LineChart,
        Reference=Reference,
        anchor="H65",
        title="Proteção da estrutura",
        y_axis_title="Proteção da estrutura (%)",
        colors=["1A1A1A", "F28E2B"],
        percent_axis=True,
    )


def _add_area_chart(
    dashboard,
    data_sheet,
    data_range: tuple[int, int, int, int],
    *,
    AreaChart,
    Reference,
    anchor: str,
    title: str,
    y_axis_title: str,
    colors: list[str],
) -> None:
    start_row, start_col, row_count, col_count = data_range
    if row_count <= 1 or col_count <= 2:
        return
    chart = AreaChart()
    chart.title = title
    chart.y_axis.title = y_axis_title
    chart.x_axis.title = "Mês do FIDC"
    chart.grouping = "stacked"
    chart.width = 23
    chart.height = 9
    data = Reference(data_sheet, min_col=start_col + 2, max_col=start_col + col_count - 1, min_row=start_row, max_row=start_row + row_count - 1)
    cats = Reference(data_sheet, min_col=start_col, min_row=start_row + 1, max_row=start_row + row_count - 1)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    for idx, series in enumerate(chart.series):
        color = colors[idx % len(colors)]
        series.graphicalProperties.solidFill = color
        series.graphicalProperties.line.solidFill = color
    chart.legend.position = "b"
    dashboard.add_chart(chart, anchor)


def _add_line_chart(
    dashboard,
    data_sheet,
    data_range: tuple[int, int, int, int],
    *,
    LineChart,
    Reference,
    anchor: str,
    title: str,
    y_axis_title: str,
    colors: list[str],
    percent_axis: bool = False,
) -> None:
    start_row, start_col, row_count, col_count = data_range
    if row_count <= 1 or col_count <= 2:
        return
    chart = LineChart()
    chart.title = title
    chart.y_axis.title = y_axis_title
    chart.x_axis.title = "Mês do FIDC"
    chart.width = 11.2
    chart.height = 8.2
    if percent_axis:
        chart.y_axis.numFmt = "0.0%"
    data = Reference(data_sheet, min_col=start_col + 2, max_col=start_col + col_count - 1, min_row=start_row, max_row=start_row + row_count - 1)
    cats = Reference(data_sheet, min_col=start_col, min_row=start_row + 1, max_row=start_row + row_count - 1)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    for idx, series in enumerate(chart.series):
        color = colors[idx % len(colors)]
        series.graphicalProperties.line.solidFill = color
        series.graphicalProperties.line.width = 25000
        series.marker.symbol = "circle"
        series.marker.size = 5
        series.marker.graphicalProperties.solidFill = color
        series.marker.graphicalProperties.line.solidFill = color
    chart.legend.position = "b"
    dashboard.add_chart(chart, anchor)


def _summary_lookup(premissas_summary_df: pd.DataFrame, premissa: str, default: str = "N/D") -> str:
    if premissas_summary_df.empty or "Premissa" not in premissas_summary_df.columns:
        return default
    rows = premissas_summary_df[premissas_summary_df["Premissa"] == premissa]
    if rows.empty:
        return default
    return str(rows.iloc[0].get("Valor", default) or default)


def _compact_brl_from_text(value: str) -> str:
    try:
        amount = _parse_br_number(value, field_name="valor")
    except ValueError:
        return value
    if abs(amount) >= 1_000_000_000:
        return f"R$ {_format_number_br(amount / 1_000_000_000.0, 1).replace(',0', '')} bi"
    if abs(amount) >= 1_000_000:
        return f"R$ {_format_number_br(amount / 1_000_000.0, 1).replace(',0', '')} mi"
    return _format_brl(amount)


def _short_percent_text(value: str, *, decimals: int = 0) -> str:
    try:
        pct = _parse_br_number(value, field_name="percentual")
    except ValueError:
        return value
    return f"{_format_number_br(pct, decimals)}%"


def _committee_premissas_rows(premissas_summary_df: pd.DataFrame) -> list[tuple[str, str]]:
    return [
        ("Volume inicial", _compact_brl_from_text(_summary_lookup(premissas_summary_df, "Volume inicial"))),
        ("Modo da carteira", _summary_lookup(premissas_summary_df, "Modo da carteira")),
        ("Prazo total do FIDC", _summary_lookup(premissas_summary_df, "Prazo total do FIDC").replace(" anos", " a")),
        ("Prazo médio dos recebíveis", _summary_lookup(premissas_summary_df, "Prazo médio dos recebíveis").replace(" meses", "m")),
        ("Entrada da taxa da carteira", _summary_lookup(premissas_summary_df, "Entrada da taxa da carteira")),
        ("Taxa mensal efetiva", _short_percent_text(_summary_lookup(premissas_summary_df, "Taxa mensal efetiva"), decimals=0)),
        ("Taxa anual equivalente base 252", _summary_lookup(premissas_summary_df, "Taxa anual equivalente base 252")),
        ("Administração/gestão", _summary_lookup(premissas_summary_df, "Administração/gestão")),
        ("Metodologia", "NPL 90 + cobertura de provisão"),
        ("NPL 90+ esperado por ciclo", _summary_lookup(premissas_summary_df, "NPL 90+ esperado por ciclo")),
        ("Lag até NPL 90+", _summary_lookup(premissas_summary_df, "Lag até NPL 90+")),
        ("Cobertura mínima NPL 90+", _short_percent_text(_summary_lookup(premissas_summary_df, "Cobertura mínima NPL 90+"), decimals=0)),
        ("LGD econômica", _short_percent_text(_summary_lookup(premissas_summary_df, "LGD econômica"), decimals=0)),
        ("Teto maturação Over90", _short_percent_text(_summary_lookup(premissas_summary_df, "Teto maturação Over90"), decimals=0)),
        ("PL SEN", _short_percent_text(_summary_lookup(premissas_summary_df, "PL SEN"), decimals=0)),
        ("PL SUB", _short_percent_text(_summary_lookup(premissas_summary_df, "PL SUB"), decimals=0)),
        ("Remuneração Sr.", "CDI +1,60%"),
        ("Amortização SEN", _summary_lookup(premissas_summary_df, "Amortização SEN")),
        ("Juros SEN", _summary_lookup(premissas_summary_df, "Juros SEN")),
        ("Início amortização SEN", _summary_lookup(premissas_summary_df, "Início amortização SEN")),
    ]


def _committee_flow_rows(timeline_frame: pd.DataFrame | None) -> list[tuple[str, str, str]]:
    if timeline_frame is None or timeline_frame.empty:
        return [
            ("Carteira originada", "100,00", "Base do período"),
            ("Receita", "14,00", "~14% a.m."),
            ("Principal vencendo", "14,29", "1/7 da carteira vence"),
            ("PDD", "-13,33", "40% da carteira / 3 meses"),
            ("Custos adm.", "-0,03", "Despesas operacionais"),
            ("Juros sênior", "N/D", "Remuneração da tranche sênior"),
            ("Excesso de Caixa", "N/D", ""),
        ]
    data = timeline_frame[timeline_frame["indice"] > 0].copy()
    if data.empty:
        return _committee_flow_rows(None)
    row = data.iloc[0]
    base = float(timeline_frame.iloc[0].get("carteira", 0.0) or 0.0)
    scale = 100.0 / base if base else 0.0

    def scaled(column: str, sign: float = 1.0) -> str:
        return _format_number_br(float(row.get(column, 0.0) or 0.0) * scale * sign, 2)

    carteira_vencendo = float(row.get("carteira_vencendo", 0.0) or 0.0)
    prazo_medio = base / carteira_vencendo if base and carteira_vencendo else 0.0
    principal_racional = (
        f"1/{_format_number_br(prazo_medio, 1).replace(',0', '')} da carteira vence"
        if prazo_medio > 0.0
        else "Quantidade média de parcelas"
    )
    pdd_value = float(row.get("despesa_provisao", row.get("perda_carteira_despesa", 0.0)) or 0.0)
    pdd_racional = "40% da carteira / 3 meses"
    excesso_caixa = float(row.get("principal_recebido_carteira", 0.0) or 0.0) + float(
        row.get("fluxo_remanescente_mezz", 0.0) or 0.0
    )

    return [
        ("Carteira originada", "100,00", "Base do período"),
        ("Receita", scaled("fluxo_carteira"), "~14% a.m. ajustado por calendário"),
        ("Principal vencendo", scaled("carteira_vencendo"), principal_racional),
        ("PDD", _format_number_br(pdd_value * scale * -1.0, 2), pdd_racional),
        ("Custos adm.", scaled("custos_adm", -1.0), "Despesas operacionais"),
        ("Juros sênior", scaled("juros_senior", -1.0), "Remuneração da tranche sênior"),
        ("Excesso de Caixa", _format_number_br(excesso_caixa * scale, 2), "Reinveste em carteira nova"),
    ]


def _committee_lock_rows(timeline_frame: pd.DataFrame | None) -> pd.DataFrame:
    columns = ["Mês", "Principal recebido", "Excesso reinv.", "Nova originação", "Aporte SUB", "Sub/PL"]
    if timeline_frame is None or timeline_frame.empty:
        return pd.DataFrame([["M1", "N/D", "N/D", "N/D", "N/D", "N/D"]], columns=columns)
    data = timeline_frame[timeline_frame["indice"] > 0].head(8).copy()
    if data.empty:
        return pd.DataFrame([["M1", "N/D", "N/D", "N/D", "N/D", "N/D"]], columns=columns)
    return pd.DataFrame(
        [
            [
                f"M{int(row.get('indice', 0))}",
                _compact_brl_from_text(_format_brl(float(row.get("principal_recebido_carteira", 0.0) or 0.0))),
                _compact_brl_from_text(_format_brl(float(row.get("reinvestimento_excesso", 0.0) or 0.0))),
                _compact_brl_from_text(_format_brl(float(row.get("nova_originacao", 0.0) or 0.0))),
                _compact_brl_from_text(_format_brl(float(row.get("aporte_subordinacao_minima", 0.0) or 0.0))),
                _format_percent(float(row.get("subordinacao_pct", 0.0) or 0.0)),
            ]
            for _, row in data.iterrows()
        ],
        columns=columns,
    )


def _committee_cards(kpi_cards: list[dict[str, str]], revolvency_cards: list[dict[str, str]]) -> list[dict[str, str]]:
    by_label = {card["label"]: card for card in kpi_cards + revolvency_cards}
    selected = [
        by_label.get("Retorno anualizado", {"label": "Retorno anualizado", "value": "N/D", "context": "Cota sênior"}),
        by_label.get("Duration econômica", {"label": "Duration econômica", "value": "N/D", "context": "Cota sênior"}),
        by_label.get("SUB inicial", {"label": "SUB inicial", "value": "N/D", "context": "Colchão subordinado"}),
        by_label.get("Prazo médio recebíveis", {"label": "Prazo médio recebíveis", "value": "N/D", "context": "Prazo de giro"}),
        by_label.get("Carteira originada efetiva", {"label": "Carteira originada efetiva", "value": "N/D", "context": "Base de comparação"}),
        by_label.get("SUB final sem perdas", {"label": "SUB final sem perdas", "value": "N/D", "context": "Colchão acumulado"}),
        by_label.get(
            "Colchão s/ perdas sobre originada",
            {"label": "Colchão s/ perdas sobre originada", "value": "N/D", "context": "SUB / carteira"},
        ),
    ]
    return [
        {
            **card,
            "value": _compact_brl_from_text(card["value"]) if str(card.get("value", "")).startswith("R$") else card["value"],
        }
        for card in selected
    ]


def _build_model_dashboard_pptx_bytes(
    *,
    kpi_cards: list[dict[str, str]],
    revolvency_cards: list[dict[str, str]],
    premissas_summary_df: pd.DataFrame,
    timeline_frame: pd.DataFrame | None = None,
    balance_chart_df: pd.DataFrame,
    loss_chart_df: pd.DataFrame,
    protection_chart_df: pd.DataFrame,
) -> bytes:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - environment guard
        raise RuntimeError("Dependência python-pptx não instalada.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    deps = {
        "CategoryChartData": CategoryChartData,
        "RGBColor": RGBColor,
        "XL_CHART_TYPE": XL_CHART_TYPE,
        "XL_DATA_LABEL_POSITION": XL_DATA_LABEL_POSITION,
        "XL_LEGEND_POSITION": XL_LEGEND_POSITION,
        "XL_MARKER_STYLE": XL_MARKER_STYLE,
        "MSO_AUTO_SHAPE_TYPE": MSO_AUTO_SHAPE_TYPE,
        "MSO_ANCHOR": MSO_ANCHOR,
        "PP_ALIGN": PP_ALIGN,
        "Inches": Inches,
        "Pt": Pt,
    }

    _ppt_add_committee_premissas_slide(prs, layout, premissas_summary_df, timeline_frame, **deps)
    _ppt_add_committee_lock_slide(prs, layout, premissas_summary_df, timeline_frame, **deps)
    _ppt_add_committee_outputs_slide(
        prs,
        layout,
        _committee_cards(kpi_cards, revolvency_cards),
        balance_chart_df,
        protection_chart_df,
        timeline_frame,
        **deps,
    )
    for page, slide in enumerate(prs.slides, start=1):
        _ppt_add_footer(slide, page, len(prs.slides), **deps)

    output = BytesIO()
    prs.save(output)
    return _normalize_pptx_unsigned_chart_ids(output.getvalue())


_PPTX_CHART_AXIS_ID_RE = re.compile(rb'(<c:(?:axId|crossAx)\b[^>]*\bval=")(-?\d+)(")')
_PPTX_CHART_AXID_RE = re.compile(rb'<c:axId\b[^>]*\bval="(-?\d+)"')


def _normalize_pptx_unsigned_chart_ids(pptx_bytes: bytes) -> bytes:
    """Keep chart axis IDs positive and inside Office's signed-32 comfort zone."""
    source = BytesIO(pptx_bytes)
    target = BytesIO()

    with ZipFile(source, "r") as zin, ZipFile(target, "w", compression=ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            payload = zin.read(item.filename)
            if item.filename.startswith("ppt/charts/") and item.filename.endswith(".xml"):
                axis_values = list(dict.fromkeys(_PPTX_CHART_AXID_RE.findall(payload)))
                axis_id_map = {
                    old: str(1_000_000 + chart_index).encode("ascii")
                    for chart_index, old in enumerate(axis_values, start=1)
                }

                def normalize(match: re.Match[bytes]) -> bytes:
                    value = match.group(2)
                    normalized = axis_id_map.get(value, value)
                    return match.group(1) + normalized + match.group(3)

                payload = _PPTX_CHART_AXIS_ID_RE.sub(normalize, payload)
            zout.writestr(item, payload)
    return target.getvalue()


def _ppt_add_committee_header(slide, title: str, **deps) -> None:
    logo = slide.shapes.add_shape(
        deps["MSO_AUTO_SHAPE_TYPE"].ROUNDED_RECTANGLE,
        deps["Inches"](0.50),
        deps["Inches"](0.20),
        deps["Inches"](0.52),
        deps["Inches"](0.36),
    )
    logo.fill.solid()
    logo.fill.fore_color.rgb = deps["RGBColor"].from_string("1F1F1F")
    logo.line.fill.background()
    _ppt_textbox(logo, "itaú", deps["Pt"](10), deps["RGBColor"].from_string("FFFFFF"), bold=True, align=deps["PP_ALIGN"].CENTER)
    bba = slide.shapes.add_textbox(deps["Inches"](1.05), deps["Inches"](0.27), deps["Inches"](0.78), deps["Inches"](0.20))
    _ppt_textbox(bba, "BBA", deps["Pt"](11), deps["RGBColor"].from_string("1F1F1F"), bold=True, align=deps["PP_ALIGN"].LEFT)
    title_box = slide.shapes.add_textbox(deps["Inches"](1.95), deps["Inches"](0.16), deps["Inches"](10.80), deps["Inches"](0.50))
    _ppt_textbox(title_box, title, deps["Pt"](23), deps["RGBColor"].from_string("1F1F1F"), bold=True, align=deps["PP_ALIGN"].LEFT)


def _ppt_add_committee_premissas_slide(prs, layout, premissas_summary_df: pd.DataFrame, timeline_frame: pd.DataFrame | None, **deps) -> None:
    slide = prs.slides.add_slide(layout)
    _ppt_set_background(slide, "FFFFFF", **deps)
    _ppt_add_committee_header(slide, "Modelagem FIDC | Fluxo Econômico e Mecânica Reinvestimento", **deps)
    rows = pd.DataFrame(_committee_premissas_rows(premissas_summary_df), columns=["Premissa", "Valor"])
    _ppt_add_table(
        slide,
        rows,
        x=0.55,
        y=0.82,
        w=5.45,
        h=5.84,
        column_widths=[2.80, 2.65],
        **deps,
    )
    _ppt_add_bullet_block(
        slide,
        "Premissas Principais",
        [
            "Carteira inicial: R$1 bi",
            "Receita Carteira: ~14% a.m.",
            f"Parcelas médias: {_summary_lookup(premissas_summary_df, 'Quantidade média de parcelas', '7,0 parcelas')}",
            f"Prazo médio econômico: {_summary_lookup(premissas_summary_df, 'Prazo médio dos recebíveis', '3,15 meses')}",
            "NPL/LGD assumidos: 40% / 100%",
            "Reinv. de principal e excesso dentro da janela elegível",
        ],
        6.45,
        0.92,
        5.95,
        **deps,
    )
    _ppt_add_bullet_block(
        slide,
        "Fluxo Simplificado",
        [
            "Principal vencendo: 1/7 da carteira por mês.",
            "PDD: 40% da carteira provisionado linearmente até Over90.",
            "Excesso de caixa positivo é reinvestido em nova carteira.",
        ],
        6.45,
        2.78,
        5.95,
        **deps,
    )
    flow = pd.DataFrame(_committee_flow_rows(timeline_frame), columns=["Item", "Valor", "Racional"])
    _ppt_add_table(
        slide,
        flow,
        x=6.50,
        y=4.55,
        w=6.08,
        h=1.82,
        column_widths=[2.15, 1.25, 2.68],
        **deps,
    )


def _ppt_add_committee_lock_slide(prs, layout, premissas_summary_df: pd.DataFrame, timeline_frame: pd.DataFrame | None, **deps) -> None:
    slide = prs.slides.add_slide(layout)
    _ppt_set_background(slide, "FFFFFF", **deps)
    _ppt_add_committee_header(slide, "Modelagem FIDC | Subordinação e Reinvestimento", **deps)
    lock_pct = _summary_lookup(premissas_summary_df, "Trava mínima de reinvestimento", "30,00%")
    _ppt_add_bullet_block(
        slide,
        "Regra de modelagem",
        [
            "Principal recebido e excesso de caixa podem ser reinvestidos integralmente.",
            f"Subordinação mínima estrutural: SUB disponível / PL FIDC >= {lock_pct}.",
            "Se perdas/custos consumirem SUB abaixo do piso, o modelo registra aporte necessário.",
            "SUB / carteira originada acumulada fica como métrica de colchão econômico.",
        ],
        0.62,
        0.95,
        5.60,
        **deps,
    )
    formula_frame = pd.DataFrame(
        [
            ["SUB mínima", "SUB / PL FIDC >= 30%"],
            ["Caixa reinvestível", "principal recebido + max(excesso de caixa, 0)"],
            ["Nova originação", "se elegível: caixa reinvestível; se não: 0"],
            ["Subordinação", "SUB disponível / PL FIDC"],
            ["Aporte SUB", "max((30% x PL - SUB) / 70%, 0)"],
        ],
        columns=["Item", "Fórmula"],
    )
    _ppt_add_table(slide, formula_frame, x=6.55, y=0.98, w=5.95, h=2.20, column_widths=[1.70, 4.25], **deps)
    monthly = _committee_lock_rows(timeline_frame)
    _ppt_add_table(slide, monthly, x=0.62, y=3.60, w=11.88, h=2.30, column_widths=[0.90, 2.25, 2.25, 2.15, 2.15, 2.18], **deps)
    note = slide.shapes.add_textbox(deps["Inches"](0.65), deps["Inches"](6.08), deps["Inches"](11.80), deps["Inches"](0.38))
    _ppt_textbox(
        note,
        "Leitura: reinvestimento troca caixa por carteira e não consome PL; a trava de 30% é testada no PL econômico da estrutura.",
        deps["Pt"](8.4),
        deps["RGBColor"].from_string("6B7280"),
        bold=False,
        align=deps["PP_ALIGN"].LEFT,
    )


def _ppt_add_committee_outputs_slide(
    prs,
    layout,
    cards: list[dict[str, str]],
    balance_chart_df: pd.DataFrame,
    protection_chart_df: pd.DataFrame,
    timeline_frame: pd.DataFrame | None,
    **deps,
) -> None:
    slide = prs.slides.add_slide(layout)
    _ppt_set_background(slide, "FFFFFF", **deps)
    _ppt_add_committee_header(slide, "Modelagem FIDC | Fluxo Econômico e Mecânica Reinvestimento", **deps)
    _ppt_add_committee_kpi_strip(slide, cards, **deps)
    _ppt_add_stacked_area_chart(
        slide,
        _balance_chart_wide(balance_chart_df),
        x=0.55,
        y=1.42,
        w=6.25,
        h=2.70,
        title="Evolução do Saldo das Cotas",
        y_axis_title="R$ milhões",
        colors=["1A1A1A", "9C9C9C", "F28E2B", "B35C00"],
        **deps,
    )
    outputs = _committee_outputs(timeline_frame)
    _ppt_add_bullet_block(slide, "Outputs", outputs, 7.20, 1.48, 5.35, **deps)
    _ppt_add_line_chart(
        slide,
        _long_percent_chart_wide(protection_chart_df),
        x=0.72,
        y=4.62,
        w=11.70,
        h=1.78,
        title="Proteção da estrutura",
        y_axis_title="Proteção da estrutura (%)",
        colors=["1A1A1A", "F28E2B"],
        percent_axis=True,
        **deps,
    )


def _committee_outputs(timeline_frame: pd.DataFrame | None) -> list[str]:
    if timeline_frame is None or timeline_frame.empty:
        return [
            "Principal vence por 1/7 da carteira e PDD é linear até Over90.",
            "Excesso de caixa positivo é reinvestido durante a janela elegível.",
            "Sub / carteira originada é métrica de colchão econômico, não trava de caixa.",
        ]
    last = timeline_frame.iloc[-1]
    total_originated = float(last.get("carteira_originada_acumulada", 0.0) or 0.0)
    final_sub = float(last.get("pl_sub_jr", 0.0) or 0.0)
    colchao = float(last.get("colchao_originada_pct", 0.0) or 0.0)
    subordinacao = float(last.get("subordinacao_pct", 0.0) or 0.0)
    if total_originated <= 0.0:
        initial_portfolio = float(timeline_frame.iloc[0].get("carteira", 0.0) or 0.0)
        new_origination = float(timeline_frame.get("nova_originacao", pd.Series(dtype=float)).clip(lower=0.0).sum() or 0.0)
        total_originated = initial_portfolio + new_origination
    if colchao <= 0.0 and total_originated > 0.0:
        colchao = final_sub / total_originated
    return [
        "PDD: 40% da carteira provisionado em 3 meses até Over90.",
        "Nova originação inclui principal recebido e excesso de caixa positivo.",
        f"Carteira originada efetiva: {_compact_brl_from_text(_format_brl(total_originated))}.",
        f"SUB / PL final: {_format_percent(subordinacao)}; colchão sobre originada: {_format_percent(colchao)}.",
    ]


def _ppt_add_bullet_block(slide, title: str, bullets: list[str], x: float, y: float, w: float, **deps) -> None:
    title_box = slide.shapes.add_textbox(deps["Inches"](x), deps["Inches"](y), deps["Inches"](w), deps["Inches"](0.26))
    _ppt_textbox(title_box, title, deps["Pt"](13.5), deps["RGBColor"].from_string("1F1F1F"), bold=True, align=deps["PP_ALIGN"].LEFT)
    body = slide.shapes.add_textbox(deps["Inches"](x), deps["Inches"](y + 0.38), deps["Inches"](w), deps["Inches"](1.35))
    text = "\n".join(f"•  {item}" for item in bullets)
    _ppt_textbox(body, text, deps["Pt"](10.0), deps["RGBColor"].from_string("2B2B2B"), bold=False, align=deps["PP_ALIGN"].LEFT)


def _ppt_add_committee_kpi_strip(slide, cards: list[dict[str, str]], **deps) -> None:
    left = 0.55
    top = 0.78
    gap = 0.12
    width = (12.25 - gap * 6) / 7
    for idx, card in enumerate(cards[:7]):
        x = left + idx * (width + gap)
        shape = slide.shapes.add_shape(
            deps["MSO_AUTO_SHAPE_TYPE"].ROUNDED_RECTANGLE,
            deps["Inches"](x),
            deps["Inches"](top),
            deps["Inches"](width),
            deps["Inches"](0.50),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = deps["RGBColor"].from_string("F7F7F7")
        shape.line.color.rgb = deps["RGBColor"].from_string("D9D9D9")
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.margin_left = deps["Inches"](0.06)
        text_frame.margin_right = deps["Inches"](0.04)
        text_frame.margin_top = deps["Inches"](0.02)
        text_frame.margin_bottom = deps["Inches"](0.01)
        lines = [
            (str(card.get("label", "")).upper(), 4.8, "6B7280", True),
            (str(card.get("value", "")), 8.0, "1F1F1F", True),
            (str(card.get("context", "")), 4.8, "6B7280", False),
        ]
        for line_idx, (text, size, color, bold) in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if line_idx == 0 else text_frame.add_paragraph()
            paragraph.text = text
            paragraph.font.name = "Calibri"
            paragraph.font.size = deps["Pt"](size)
            paragraph.font.bold = bold
            paragraph.font.color.rgb = deps["RGBColor"].from_string(color)
            paragraph.space_after = deps["Pt"](0)


def _ppt_add_summary_slide(
    prs,
    layout,
    kpi_cards: list[dict[str, str]],
    revolvency_cards: list[dict[str, str]],
    **deps,
) -> None:
    slide = prs.slides.add_slide(layout)
    _ppt_set_background(slide, "FFFFFF", **deps)
    _ppt_add_header(
        slide,
        "Modelagem FIDC — Dashboard",
        "Cards de resumo econômico, capacidade de perda e denominadores da tela de Modelagem.",
        **deps,
    )
    _ppt_add_section_label(slide, "Resumo econômico", 0.55, **deps)
    _ppt_add_card_grid(slide, kpi_cards, x=0.50, y=0.82, w=12.32, h=2.05, columns=4, **deps)
    _ppt_add_section_label(slide, "Capacidade de perda e denominadores", 3.08, **deps)
    _ppt_add_card_grid(slide, revolvency_cards, x=0.50, y=3.35, w=12.32, h=1.05, columns=4, **deps)
    note = slide.shapes.add_textbox(deps["Inches"](0.55), deps["Inches"](4.75), deps["Inches"](12.20), deps["Inches"](1.05))
    _ppt_textbox(
        note,
        "Premissas resumidas aparecem nos slides seguintes; timeline e memória detalhada seguem preservadas no Excel.",
        deps["Pt"](12),
        deps["RGBColor"].from_string("6B7280"),
        bold=False,
        align=deps["PP_ALIGN"].LEFT,
    )


def _ppt_add_premissas_slides(prs, layout, premissas_summary_df: pd.DataFrame, **deps) -> None:
    rows_per_slide = 18
    frame = premissas_summary_df.copy()
    if frame.empty:
        frame = pd.DataFrame([{"Categoria": "Premissas", "Premissa": "N/D", "Valor": "—", "Observação": ""}])
    chunks = [frame.iloc[start : start + rows_per_slide] for start in range(0, len(frame), rows_per_slide)]
    for idx, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(layout)
        _ppt_set_background(slide, "FFFFFF", **deps)
        subtitle = "Resumo das premissas usadas na simulação."
        if len(chunks) > 1:
            subtitle = f"{subtitle} Parte {idx} de {len(chunks)}."
        _ppt_add_header(slide, "Premissas resumidas", subtitle, **deps)
        _ppt_add_table(
            slide,
            chunk[["Categoria", "Premissa", "Valor", "Observação"]],
            x=0.50,
            y=0.88,
            w=12.32,
            h=5.92,
            column_widths=[1.45, 2.55, 2.00, 6.32],
            **deps,
        )


def _ppt_add_balance_slide(prs, layout, balance_chart_df: pd.DataFrame, **deps) -> None:
    slide = prs.slides.add_slide(layout)
    _ppt_set_background(slide, "FFFFFF", **deps)
    _ppt_add_header(
        slide,
        "Evolução do saldo das cotas",
        "Mesma visão exibida no dashboard: SEN, MEZZ, SUB disponível e déficit econômico.",
        **deps,
    )
    wide = _balance_chart_wide(balance_chart_df)
    _ppt_add_stacked_area_chart(
        slide,
        wide,
        x=0.50,
        y=1.00,
        w=12.30,
        h=5.65,
        title="Evolução do saldo das cotas",
        y_axis_title="R$ milhões",
        colors=["1A1A1A", "9C9C9C", "F28E2B", "B35C00"],
        **deps,
    )


def _ppt_add_loss_protection_slide(
    prs,
    layout,
    loss_chart_df: pd.DataFrame,
    protection_chart_df: pd.DataFrame,
    **deps,
) -> None:
    slide = prs.slides.add_slide(layout)
    _ppt_set_background(slide, "FFFFFF", **deps)
    _ppt_add_header(
        slide,
        "Perda da carteira e proteção da estrutura",
        "Mesmas séries exibidas no dashboard, em gráficos editáveis do PowerPoint.",
        **deps,
    )
    _ppt_add_line_chart(
        slide,
        _long_percent_chart_wide(loss_chart_df),
        x=0.50,
        y=1.00,
        w=5.95,
        h=4.65,
        title="Perda da carteira",
        y_axis_title="Perda da carteira (%)",
        colors=["F28E2B"],
        percent_axis=True,
        **deps,
    )
    _ppt_add_line_chart(
        slide,
        _long_percent_chart_wide(protection_chart_df),
        x=6.86,
        y=1.00,
        w=5.95,
        h=4.65,
        title="Proteção da estrutura",
        y_axis_title="Proteção da estrutura (%)",
        colors=["1A1A1A", "F28E2B"],
        percent_axis=True,
        **deps,
    )
    captions = [
        ("Perda da carteira", _chart_definition_caption("loss"), 0.55),
        ("Proteção da estrutura", _chart_definition_caption("protection"), 6.91),
    ]
    for title, text, x in captions:
        box = slide.shapes.add_textbox(deps["Inches"](x), deps["Inches"](5.82), deps["Inches"](5.70), deps["Inches"](0.66))
        _ppt_textbox(
            box,
            f"{title}: {text}",
            deps["Pt"](7.5),
            deps["RGBColor"].from_string("6B7280"),
            bold=False,
            align=deps["PP_ALIGN"].LEFT,
        )


def _ppt_set_background(slide, color: str, **deps) -> None:
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = deps["RGBColor"].from_string(color)


def _ppt_add_header(slide, title: str, subtitle: str, **deps) -> None:
    title_box = slide.shapes.add_textbox(deps["Inches"](0.50), deps["Inches"](0.20), deps["Inches"](8.50), deps["Inches"](0.34))
    _ppt_textbox(title_box, title, deps["Pt"](22), deps["RGBColor"].from_string("1F1F1F"), bold=True, align=deps["PP_ALIGN"].LEFT)
    subtitle_box = slide.shapes.add_textbox(deps["Inches"](0.52), deps["Inches"](0.57), deps["Inches"](11.80), deps["Inches"](0.24))
    _ppt_textbox(subtitle_box, subtitle, deps["Pt"](9), deps["RGBColor"].from_string("6B7280"), bold=False, align=deps["PP_ALIGN"].LEFT)


def _ppt_add_footer(slide, page: int, total_pages: int, **deps) -> None:
    line = slide.shapes.add_shape(deps["MSO_AUTO_SHAPE_TYPE"].RECTANGLE, deps["Inches"](0.50), deps["Inches"](6.94), deps["Inches"](12.33), deps["Inches"](0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = deps["RGBColor"].from_string("E0E0E0")
    line.line.fill.background()
    footer = slide.shapes.add_textbox(deps["Inches"](0.50), deps["Inches"](7.03), deps["Inches"](8.20), deps["Inches"](0.20))
    _ppt_textbox(footer, "Fonte: simulação interna Toma Conta | Modelagem FIDC", deps["Pt"](7.5), deps["RGBColor"].from_string("6B7280"), bold=False, align=deps["PP_ALIGN"].LEFT)
    page_box = slide.shapes.add_textbox(deps["Inches"](10.80), deps["Inches"](7.03), deps["Inches"](2.00), deps["Inches"](0.20))
    _ppt_textbox(page_box, f"Página {page} de {total_pages}", deps["Pt"](7.5), deps["RGBColor"].from_string("6B7280"), bold=False, align=deps["PP_ALIGN"].RIGHT)


def _ppt_add_section_label(slide, text: str, y: float, **deps) -> None:
    box = slide.shapes.add_textbox(deps["Inches"](0.52), deps["Inches"](y), deps["Inches"](6.00), deps["Inches"](0.22))
    _ppt_textbox(box, text, deps["Pt"](11), deps["RGBColor"].from_string("1F1F1F"), bold=True, align=deps["PP_ALIGN"].LEFT)


def _ppt_add_card_grid(
    slide,
    cards: list[dict[str, str]],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    columns: int,
    **deps,
) -> None:
    gap_x = 0.16
    gap_y = 0.14
    rows = max((len(cards) + columns - 1) // columns, 1)
    card_w = (w - gap_x * (columns - 1)) / columns
    card_h = (h - gap_y * (rows - 1)) / rows
    for idx, card in enumerate(cards):
        row = idx // columns
        col = idx % columns
        _ppt_add_card(slide, x + col * (card_w + gap_x), y + row * (card_h + gap_y), card_w, card_h, card, **deps)


def _ppt_add_card(slide, x: float, y: float, w: float, h: float, card: dict[str, str], **deps) -> None:
    shape = slide.shapes.add_shape(
        deps["MSO_AUTO_SHAPE_TYPE"].ROUNDED_RECTANGLE,
        deps["Inches"](x),
        deps["Inches"](y),
        deps["Inches"](w),
        deps["Inches"](h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = deps["RGBColor"].from_string("F7F7F7")
    shape.line.color.rgb = deps["RGBColor"].from_string("E0E0E0")
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = deps["Inches"](0.08)
    text_frame.margin_right = deps["Inches"](0.08)
    text_frame.margin_top = deps["Inches"](0.04)
    text_frame.margin_bottom = deps["Inches"](0.03)
    text_frame.word_wrap = True
    entries = [
        (str(card.get("label", "")).upper(), 6.7, "6B7280", True),
        (str(card.get("value", "")), 12.2, "1F1F1F", True),
        (str(card.get("context", "")), 6.6, "6B7280", False),
    ]
    for idx, (text, size, color, bold) in enumerate(entries):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = text
        paragraph.font.name = "Calibri"
        paragraph.font.size = deps["Pt"](size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = deps["RGBColor"].from_string(color)
        paragraph.alignment = deps["PP_ALIGN"].LEFT
        paragraph.space_after = deps["Pt"](0)


def _ppt_textbox(shape, text: str, size, color, *, bold: bool, align, **_) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Calibri"
    paragraph.font.size = size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _ppt_add_table(
    slide,
    frame: pd.DataFrame,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    column_widths: list[float],
    **deps,
) -> None:
    table_shape = slide.shapes.add_table(len(frame) + 1, len(frame.columns), deps["Inches"](x), deps["Inches"](y), deps["Inches"](w), deps["Inches"](h))
    table = table_shape.table
    for idx, width in enumerate(column_widths):
        table.columns[idx].width = deps["Inches"](width)
    for col_idx, header in enumerate(frame.columns):
        _ppt_style_table_cell(table.cell(0, col_idx), str(header), fill="1F1F1F", color="FFFFFF", bold=True, font_size=7.5, **deps)
    for row_idx, row in enumerate(frame.itertuples(index=False), start=1):
        fill = "F7F7F7" if row_idx % 2 == 0 else "FFFFFF"
        for col_idx, value in enumerate(row):
            _ppt_style_table_cell(table.cell(row_idx, col_idx), str(value or "—"), fill=fill, color="1F1F1F", bold=False, font_size=6.6, **deps)


def _ppt_style_table_cell(
    cell,
    text: str,
    *,
    fill: str,
    color: str,
    bold: bool,
    font_size: float,
    **deps,
) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = deps["RGBColor"].from_string(fill)
    cell.vertical_anchor = deps["MSO_ANCHOR"].MIDDLE
    cell.margin_left = deps["Inches"](0.04)
    cell.margin_right = deps["Inches"](0.04)
    cell.margin_top = deps["Inches"](0.02)
    cell.margin_bottom = deps["Inches"](0.02)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Calibri"
    paragraph.font.size = deps["Pt"](font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = deps["RGBColor"].from_string(color)
    paragraph.alignment = deps["PP_ALIGN"].LEFT


def _ppt_add_stacked_area_chart(
    slide,
    wide: pd.DataFrame,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    y_axis_title: str,
    colors: list[str],
    **deps,
) -> None:
    _ppt_add_native_chart(
        slide,
        wide,
        chart_type=deps["XL_CHART_TYPE"].AREA_STACKED,
        x=x,
        y=y,
        w=w,
        h=h,
        title=title,
        y_axis_title=y_axis_title,
        colors=colors,
        percent_axis=False,
        show_data_labels=False,
        **deps,
    )


def _ppt_add_line_chart(
    slide,
    wide: pd.DataFrame,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    y_axis_title: str,
    colors: list[str],
    percent_axis: bool,
    **deps,
) -> None:
    _ppt_add_native_chart(
        slide,
        wide,
        chart_type=deps["XL_CHART_TYPE"].LINE_MARKERS,
        x=x,
        y=y,
        w=w,
        h=h,
        title=title,
        y_axis_title=y_axis_title,
        colors=colors,
        percent_axis=percent_axis,
        show_data_labels=percent_axis,
        **deps,
    )


def _ppt_add_native_chart(
    slide,
    wide: pd.DataFrame,
    *,
    chart_type,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    y_axis_title: str,
    colors: list[str],
    percent_axis: bool,
    show_data_labels: bool,
    **deps,
) -> None:
    columns = _ppt_chart_series_columns(wide)
    if wide.empty or not columns:
        _ppt_add_chart_no_data(slide, title, x, y, w, **deps)
        return
    chart_data = deps["CategoryChartData"]()
    chart_data.categories = [f"M{int(value)}" if pd.notna(value) else "" for value in wide["Mês"]]
    value_scale = 100.0 if percent_axis else 1.0
    for column in columns:
        chart_data.add_series(
            str(column),
            [float(value or 0.0) * value_scale for value in wide[column].fillna(0.0).tolist()],
        )
    shape = slide.shapes.add_chart(
        chart_type,
        deps["Inches"](x),
        deps["Inches"](y),
        deps["Inches"](w),
        deps["Inches"](h),
        chart_data,
    )
    chart = shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.name = "Calibri"
    chart.chart_title.text_frame.paragraphs[0].font.size = deps["Pt"](11)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = y_axis_title
    chart.value_axis.tick_labels.font.size = deps["Pt"](8)
    chart.category_axis.tick_labels.font.size = deps["Pt"](8)
    if percent_axis:
        chart.value_axis.tick_labels.number_format = "0\\%"
        try:
            chart.value_axis.minimum_scale = 0.0
        except ValueError:
            pass
    chart.has_legend = True
    chart.legend.position = deps["XL_LEGEND_POSITION"].BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = deps["Pt"](8)
    for idx, series in enumerate(chart.series):
        color = colors[idx % len(colors)]
        series.format.line.color.rgb = deps["RGBColor"].from_string(color)
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = deps["RGBColor"].from_string(color)
        if hasattr(series, "marker"):
            series.marker.style = deps["XL_MARKER_STYLE"].CIRCLE
            series.marker.size = 5
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = deps["RGBColor"].from_string(color)
            series.marker.format.line.color.rgb = deps["RGBColor"].from_string(color)
    if show_data_labels:
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            labels = plot.data_labels
            labels.number_format = "0\\%" if percent_axis else "#,##0"
            labels.position = deps["XL_DATA_LABEL_POSITION"].ABOVE
            labels.font.size = deps["Pt"](7)
            labels.font.color.rgb = deps["RGBColor"].from_string("333333")
        except (AttributeError, IndexError, ValueError):
            pass


def _ppt_chart_series_columns(wide: pd.DataFrame) -> list[str]:
    if wide.empty:
        return []
    return [str(column) for column in wide.columns if column not in {"Mês", "Data"}]


def _ppt_add_chart_no_data(slide, title: str, x: float, y: float, w: float, **deps) -> None:
    box = slide.shapes.add_textbox(deps["Inches"](x), deps["Inches"](y), deps["Inches"](w), deps["Inches"](0.45))
    _ppt_textbox(
        box,
        f"{title}: sem dados suficientes.",
        deps["Pt"](10),
        deps["RGBColor"].from_string("6B7280"),
        bold=False,
        align=deps["PP_ALIGN"].LEFT,
    )


def _build_workbook_mechanics_markdown(
    *,
    selected_curve: _SelectedCurve,
    selected_calendar: _SelectedCalendar,
    interpolation_label: str,
    taxa_cessao_input_mode: str,
    data_inicial: date | None = None,
    credit_model_label: str = CREDIT_LABEL_NPL90,
) -> str:
    start_label = data_inicial.strftime("%d/%m/%Y") if data_inicial else "data inicial selecionada"
    return "\n".join(
        [
            "Esta seção descreve a mecânica completa da aba e a base de cálculo usada na simulação.",
            "",
            "### 1. Datas, dias corridos e dias úteis",
            "",
            f"- A data inicial padrão é a data-base da curva carregada. Nesta simulação, a data inicial é `{start_label}`.",
            "- O usuário pode sobrescrever a data inicial; nesse caso, a timeline é deslocada, mas a curva DI/Pré selecionada continua sendo a fonte dos juros.",
            "- Se a data sobrescrita estiver distante da data-base da curva, os resultados devem ser lidos como simulação prospectiva com a estrutura de curva selecionada.",
            "- A simulação roda em grade mensal pelo prazo informado ou na grade padrão deslocada a partir da data inicial.",
            "- `DC` é a quantidade de dias corridos desde a data inicial do fluxo.",
            "- `DU` é a quantidade de dias úteis desde a data inicial, descontando fins de semana e feriados.",
            f"- Calendário usado nesta simulação: `{selected_calendar.source_label}`.",
            f"- Curva usada nesta simulação: `{selected_curve.source_label}`, data-base `{selected_curve.base_date:%d/%m/%Y}`, interpolação `{interpolation_label}`.",
            "- Bases temporais atuais por componente: carteira em `delta_DU / 21`, cotas em `delta_DU / 252`, custos em taxa mensal composta, crédito pelo período mensal da simulação e SELIC projetada em `21 DU` médios por mês.",
            "- A carteira ainda usa 21 dias úteis médios como mês financeiro; cotas e curva DI usam DU efetivos em base 252.",
            "",
            "### 2. Taxa da carteira",
            "",
            "- O usuário pode informar a taxa bruta de duas formas:",
            "- `Taxa de Cessão`: deságio nominal sobre o valor futuro no prazo médio dos recebíveis. Ex.: comprar R$ 100 por R$ 95 significa taxa de cessão nominal de `5,00%` no prazo do recebível.",
            "- `Taxa Mensal (%)`: taxa efetiva bruta cobrada ao mês; o motor converte para o deságio nominal equivalente no prazo médio.",
            "- Se houver ágio sobre face, o FIDC paga mais pelo recebível. O ágio reduz a taxa de cessão efetiva e aumenta o EAD econômico.",
            "",
            "```text",
            "taxa_cessao_efetiva = taxa_cessao_nominal - agio_sobre_face",
            "preco_pago = face * (1 - taxa_cessao_efetiva)",
            "taxa_mensal_efetiva = (1 / (1 - taxa_cessao_efetiva)) ^ (1 / prazo_medio_recebiveis_meses) - 1",
            "taxa_cessao_nominal = 1 - 1 / ((1 + taxa_mensal_bruta) ^ prazo_medio_recebiveis_meses)",
            "```",
            "",
            "- Exemplo: face de `R$ 100`, taxa de cessão nominal de `5,00%`, ágio de `1,00%` e prazo de `1 mês` geram preço pago de `R$ 96`, taxa efetiva de cessão de `4,00%` e taxa mensal efetiva de `4,17% a.m.`.",
            "- Sem ágio, o mesmo recebível comprado por `R$ 95` teria taxa mensal efetiva de `5,26% a.m.`.",
            "- A carteira gera retorno econômico por composição da taxa mensal aplicada no período:",
            "",
            "```text",
            "fluxo_carteira = carteira * ((1 + taxa_cessao_am_aplicada) ^ (delta_DU / 21) - 1)",
            "```",
            "",
            "- Em carteira revolvente, `carteira` é o saldo em aberto usado para juros e perda; ele se mantém reciclando o principal recebido e o excesso de caixa positivo enquanto a nova carteira couber no prazo do FIDC. A subordinação mínima é testada como `SUB / PL FIDC`, não como cap de originação sobre carteira originada.",
            "- A timeline preserva carteira pelo valor de face e mostra EAD/saldo em risco pelo preço pago. A provisão e o NPL usam EAD ajustado quando há ágio sobre face.",
            "",
            f"- Entrada informada pelo usuário nesta simulação: `{taxa_cessao_input_mode}`.",
            "- Em ambos os casos, a aba também calcula o de-para anual em base 252 dias úteis:",
            "",
            "```text",
            "taxa_cessao_aa_252 = (1 + taxa_cessao_am) ^ (252 / 21) - 1",
            "```",
            "",
            "- Se houver excesso de spread sobre a SEN, a taxa mensal aplicada usa um piso:",
            "",
            "```text",
            "taxa_cessao_am_aplicada = max(taxa_informada, CDI + spread_SEN + excesso_spread)",
            "```",
            "",
            "### 3. Custos de administração e gestão",
            "",
            "- O custo percentual é anual, mas o motor calcula uma parcela mensal composta sobre o PL econômico do início do período, com piso mínimo:",
            "",
            "```text",
            "custos_adm = max(PL_inicio * ((1 + custo_adm_aa) ^ (1/12) - 1), custo_minimo_mensal)",
            "```",
            "",
            "- Exemplo: `0,35` na interface significa `0,35% a.a.`; internamente vira `0,0035`.",
            "",
            "### 4. Metodologias de crédito, NPL e provisão",
            "",
            f"- Metodologia selecionada nesta simulação: `{credit_model_label}`.",
            "- Esta seção é a principal para reconciliar a contabilidade econômica do fluxo projetado.",
            "- O motor separa quatro coisas que não devem ser misturadas:",
            "  1. **principal programado para vencer**: parcela da carteira que deveria virar caixa no mês;",
            "  2. **provisão/despesa de PDD**: rubrica prospectiva que reduz o PL econômico;",
            "  3. **entrada em NPL 90+**: maturação do evento de crédito depois do lag informado;",
            "  4. **baixa/write-off**: retirada do crédito da carteira, consumindo provisão já formada quando houver saldo.",
            "- Na metodologia `NPL 90 + cobertura de provisão`, o usuário informa o NPL 90+ esperado por ciclo. O motor reconhece PDD linearmente até o lag e baixa o crédito quando ele entra em NPL 90+.",
            "",
            "```text",
            "carteira_vencendo = carteira_inicio * meses_periodo / qtd_parcelas_media",
            "pdd_prospectiva = carteira_inicio * npl90_esperado_por_ciclo * meses_periodo / lag_npl90_meses",
            "principal_recebido = carteira_vencendo",
            "entrada_npl90_t = npl90_futuro_de_periodos_anteriores_apos_lag",
            "baixa_credito_t = entrada_npl90_t * LGD",
            "estoque_npl90_t = max(estoque_npl90_t-1 + entrada_npl90_t - baixa_credito_t, 0)",
            "provisao_minima = estoque_npl90_t * cobertura_minima * LGD",
            "provisao_prospectiva = pdd_prospectiva * LGD",
            "provisao_pos_writeoff = max(provisao_anterior - baixa_credito_t, 0)",
            "provisao_requerida = max(provisao_pos_writeoff + provisao_prospectiva, provisao_minima)",
            "despesa_provisao = writeoff_descoberto + max(provisao_requerida - provisao_pos_writeoff, 0)",
            "```",
            "",
            "- Leitura prática: com carteira de `R$ 100`, `7` parcelas e `40%` de NPL esperado em `3` meses, vencem `R$ 14,29` de principal e a PDD do mês é `R$ 13,33`.",
            "- O principal programado e a PDD não são a mesma linha: a PDD reduz o PL econômico; o principal recebido pode ser reinvestido enquanto a revolvência estiver elegível.",
            "- Se a provisão anterior for suficiente, o write-off consome provisão e não gera despesa adicional além do reforço necessário para manter a provisão requerida.",
            "- Se a provisão anterior for insuficiente, a parte não coberta aparece como `writeoff_descoberto`, aumentando a despesa de provisão/perda do mês.",
            "- A baixa é imediata no mês em que o crédito entra em NPL 90+. Assim, com `LGD = 100%`, a entrada em NPL 90+ deixa de compor carteira performing no mesmo período.",
            "- A provisão modelada é prospectiva, no estilo `ECL forward-looking`: ela antecipa a perda esperada de ciclo antes do write-off e nunca deixa a provisão abaixo da cobertura mínima do NPL 90+ observado.",
            "- Esta é uma aproximação econômica para simulação e não deve ser lida como regra contábil/regulatória estrita. O objetivo é separar caixa de principal, PDD prospectiva e baixa de crédito.",
            "- A implementação atual não reconhece reversão negativa de provisão quando o estoque NPL 90+ cai. Isso é conservador e pode reduzir a SUB final; a revisão com reversão explícita depende de decisão metodológica posterior.",
            "- Na metodologia `Migração por faixas de atraso`, o usuário informa taxas mensais de rolagem entre buckets:",
            "",
            "```text",
            "adimplente -> 1-30 -> 31-60 -> 61-90 -> NPL 90+",
            "recuperacao_90+ entra como caixa",
            "writeoff_90+ baixa a carteira contra provisao",
            "```",
            "",
            "- A provisão requerida continua sendo função de `estoque NPL 90+ * cobertura mínima * LGD`.",
            "- Esta versão avançada usa rolagens agregadas simples. Curvas MOB por safra podem ser adicionadas depois como refinamento, sem mudar a lógica de cobertura.",
            "- Se a taxa mensal da carteira é `11,00%` e a despesa de provisão do mês é `1,00%` da carteira, a despesa consome cerca de `1 / 11` da receita bruta do mês, mas representa `1,00%` do principal em aberto.",
            "",
            "### 5. Taxas SEN e MEZZ",
            "",
            "- Para cotas pós-fixadas, o motor usa `CDI + spread` em convenção aditiva: o spread informado é somado ao CDI, não multiplicado por ele.",
            "",
            "```text",
            "taxa_classe_aa = CDI/PreDI_interpolado + spread_classe_aa",
            "```",
            "",
            "- Para cotas pré-fixadas, o app usa diretamente a taxa anual informada.",
            "- A taxa de cada período é transformada em FRA pela composição entre o ponto atual e o ponto anterior da curva, em base 252 dias úteis.",
            "- Os juros de cada classe são calculados sobre o PL da classe no início do período:",
            "",
            "```text",
            "juros_classe = pl_classe_inicial * ((1 + fra_classe) ^ (delta_DU / 252) - 1)",
            "```",
            "",
            "### 6. Principal, PMT e waterfall atual",
            "",
            "- O motor tem pagamentos programados de SEN e MEZZ. O pagamento de cada classe é:",
            "",
            "```text",
            "PMT SEN = juros SEN + principal SEN programado",
            "PMT MEZZ = juros MEZZ + principal MEZZ programado",
            "```",
            "",
            "- A SUB é residual: não há principal nem juros programados para SUB nesta versão.",
            "- Os PMTs programados de SEN/MEZZ são calculados mesmo quando o fluxo líquido da carteira fica insuficiente.",
            "- Portanto, a trava de caixa está desligada nesta etapa.",
            "",
            "### 7. O que seria a trava de caixa",
            "",
            "- Com trava de caixa ligada, o modelo não pagaria mais aos cotistas do que o caixa disponível no período.",
            "- A ordem econômica esperada seria:",
            "",
            "```text",
            "1. custos e despesas",
            "2. juros SEN",
            "3. principal SEN",
            "4. juros MEZZ",
            "5. principal MEZZ",
            "6. residual SUB",
            "```",
            "",
            "- Se o caixa acabasse no item 3, por exemplo, MEZZ e SUB não receberiam naquele período.",
            "- Essa trava ainda não está implementada; quando for implementada, deve ficar em premissas avançadas.",
            "",
            "### 8. PL econômico e SUB residual",
            "",
            "- Esta seção mostra como o motor fecha o PL econômico mês a mês.",
            "- A receita da carteira e o rendimento de caixa aumentam o PL. Custos, despesa de provisão/perda e PMTs de SEN/MEZZ reduzem o PL.",
            "- O write-off não é subtraído de novo diretamente do PL quando já foi coberto pela provisão. A baixa reduz o ativo carteira; o efeito de resultado passa pela despesa de provisão/perda e pelo eventual `writeoff_descoberto`.",
            "- Depois de retorno da carteira, rendimento do caixa SELIC, custos, perdas e PMTs, o PL econômico do veículo é:",
            "",
            "```text",
            "fluxo_ativos_total = fluxo_carteira + rendimento_caixa_selic + recuperacao_credito",
            "PL FIDC_t = PL FIDC_t-1 + fluxo_ativos_total - custos - despesa_provisao/perda - PMT SEN - PMT MEZZ",
            "```",
            "",
            "- O saldo econômico de SEN e MEZZ cai conforme o principal programado é amortizado.",
            "- A carteira, por sua vez, evolui pelo saldo inicial, nova originação, principal recebido, baixas e runoff. Por isso a timeline mostra separadamente `PDD`, `Baixa de crédito/write-off`, `Reinvestimento de principal` e `Reinvestimento de excesso de caixa`.",
            "- A SUB residual corrente é:",
            "",
            "```text",
            "SUB residual = PL FIDC - PL SEN - PL MEZZ",
            "```",
            "",
            "- Se a SUB residual fica positiva, ela representa colchão subordinado disponível.",
            "- Se fica negativa, ela não é um saldo de cota para receber; é déficit econômico depois de consumir toda a SUB.",
            "- A timeline detalhada preserva também a coluna residual histórica para conferência, mas os gráficos usam o residual corrente para evitar distorções como percentuais extremos.",
            "",
            "### 9. Revolvência, denominadores e capacidade de perda",
            "",
            "- Quando a carteira é revolvente, o modelo recicla principal recebido e excesso de caixa positivo enquanto o prazo médio econômico dos recebíveis ainda cabe no prazo restante do FIDC:",
            "",
            "```text",
            "giro_estimado = prazo_total_fidc_anos * 12 / prazo_medio_recebiveis_meses",
            "mes_limite_reinvestimento = prazo_total_fidc_meses - prazo_medio_recebiveis_meses",
            "principal_programado = carteira_inicio * meses_periodo / qtd_parcelas_media",
            "principal_recebido = principal_programado",
            "excesso_caixa = max(fluxo_remanescente_apos_MEZZ, 0)",
            "nova_originacao_economica = principal_recebido + excesso_caixa",
            "```",
            "",
            "- Se o mês do FIDC fica depois do mês limite de reinvestimento, a nova originação econômica vira `0` e a carteira começa a amortizar por runoff.",
            "- A compra nova não é reduzida por `SUB / carteira originada acumulada`; reinvestimento troca caixa por carteira e não consome PL.",
            "- Se perdas/custos derrubarem `SUB / PL FIDC` abaixo da subordinação mínima, a timeline registra o aporte subordinado necessário para recompor o piso.",
            "- O denominador principal de carteira originada usa a originação efetiva do motor, isto é, principal reciclado e excesso de caixa reinvestido durante a revolvência.",
            "",
            "```text",
            "carteira_originada_efetiva = volume_inicial + soma(nova_originacao_economica_t)",
            "nova_originacao_economica_t = reinvestimento_principal_t + reinvestimento_excesso_t",
            "```",
            "",
            "- No cenário sem perdas (usado para o colchão econômico), a carteira originada efetiva também incorpora o excesso de caixa positivo gerado pela estrutura.",
            "- Em cenários com perdas, a PDD consome PL e pode reduzir o excesso reinvestível; por isso `carteira_originada_efetiva` pode ficar abaixo do giro programático.",
            "- Quando a carteira é estática, a carteira originada é apenas a compra inicial:",
            "",
            "```text",
            "carteira_originada = volume_inicial",
            "```",
            "",
            "- O colchão econômico sem perdas roda uma simulação paralela com crédito zerado e compara o colchão final positivo da SUB com a carteira total originada:",
            "",
            "```text",
            "colchao_sem_perdas = max(SUB_final_sem_perdas, 0) / carteira_originada",
            "```",
            "",
            "- Ao longo do tempo, a carteira originada acumulada é calculada mês a mês:",
            "",
            "```text",
            "nova_originacao_acumulada = volume_inicial * min(mes_fidc, mes_limite_reinvestimento) / prazo_medio_recebiveis_meses",
            "denominador_mes = volume_inicial + nova_originacao_acumulada",
            "protecao_disponivel_no_mes = SUB_disponivel_no_mes / denominador_mes",
            "```",
            "",
            "- Exemplo: com 7 parcelas e taxa mensal de 14%, a duration econômica fica perto de `3,15 meses`, mas o principal que vence no mês segue a amortização das parcelas: cerca de `1/7` do volume inicial.",
            "- Exemplo: em FIDC de `36 meses` com recebíveis de `12 meses`, a originação nova para quando o fluxo chega perto do mês `24`, porque novos recebíveis de 12 meses já não caberiam no prazo da estrutura.",
            "",
            "### 10. Caixa pós-revolvência e SELIC projetada",
            "",
            "- Enquanto a revolvência é elegível, o modelo recicla principal recebido e excesso de caixa positivo em novos recebíveis.",
            "- Quando o prazo médio dos recebíveis já não cabe mais no prazo restante do FIDC, principal recebido e excesso de caixa deixam de ser reciclados e passam a entrar no saldo de caixa SELIC.",
            "- A taxa SELIC média é uma projeção digitada pelo usuário por ano calendário; nesta etapa ela não vem de fonte externa.",
            "- Esta curva manual remunera o caixa acumulado a cada período: caixa não reinvestido e, depois do mês limite de reinvestimento, também o principal recebido.",
            "- O CDI implícito das cotas pós-fixadas e o Pre DI na duration continuam usando a curva DI/Pré selecionada na fonte B3/local.",
            "- O default é `13,00% a.a.` para 2026, `12,00% a.a.` para 2027 e `12,00% a.a.` para 2028 em diante; o usuário pode sobrescrever os campos exibidos.",
            "- O motor transforma a taxa anual em taxa do período com matemática financeira exponencial e 21 dias úteis médios por mês:",
            "",
            "```text",
            "taxa_selic_periodo = (1 + selic_aa_do_ano) ^ (21 * meses_periodo / 252) - 1",
            "rendimento_caixa_selic = (caixa_selic_inicio + principal_para_caixa_selic) * taxa_selic_periodo",
            "saldo_caixa_selic_fim = caixa_selic_inicio + principal_para_caixa_selic + fluxo_remanescente_apos_MEZZ",
            "```",
            "",
            "- Exemplo: com 7 parcelas, cerca de `1/7` da carteira em aberto vence a cada mês.",
            "- Antes do mês limite de reinvestimento, essa parcela recompra recebíveis; depois do mês limite, ela vai para caixa SELIC.",
            "- O rendimento do caixa SELIC entra no fluxo econômico total dos ativos antes de custos, perdas e pagamentos das cotas.",
            "- O modelo ainda não usa SELIC observada nem curva de mercado para essa projeção; a premissa é manual para manter rastreabilidade.",
            "",
            "### 11. Indicadores do resumo econômico",
            "",
            "- Retorno anualizado SEN: XIRR dos PMTs SEN contra a data de cada fluxo.",
            "- Retorno anualizado MEZZ: XIRR dos PMTs MEZZ contra a data de cada fluxo.",
            "- Retorno anualizado SUB: só aparece nos cards quando a SUB tem série de caixa válida.",
            "- Duration SEN: média ponderada dos pagamentos SEN descontados, usando `DU / 252`.",
            "- Pre DI na duration: taxa interpolada entre os pontos simulados da curva no DU equivalente à duration da SEN.",
            "- SUB inicial: volume inicial multiplicado pela proporção subordinada.",
            "",
            "### 12. Como interpretar os gráficos",
            "",
            "- `Evolução do saldo das cotas`: mostra SEN, MEZZ, SUB disponível e, quando existir, déficit econômico separado.",
            "- `Perda da carteira`: mostra a perda do período, isto é, a despesa de provisão do mês dividida pela carteira do mês.",
            "- `Proteção da estrutura`: mostra subordinação econômica e colchão de proteção.",
            "- `Colchão de proteção` é `SUB disponível / carteira originada acumulada`; ele mede o colchão econômico naquele mês.",
            "- Se a perda do período se sustenta ao longo dos meses enquanto a subordinação disponível cai em direção a zero (gráfico de evolução do saldo das cotas), a estrutura está consumindo o colchão subordinado.",
            "- Se aparece déficit econômico, o cenário já ultrapassou a proteção da SUB dentro da mecânica atual.",
            "",
            "### 13. Limitações atuais",
            "",
            "- Ainda não há trava de caixa ligada no waterfall.",
            "- A migração de crédito é agregada por buckets; ainda não há upload de MOB por safra nem capitalização de juros em atraso.",
            "- Ainda não há amortização customizada por classe via interface avançada.",
            "- Ainda não há fluxo programado para SUB; ela permanece residual nesta versão.",
            "",
            "### 14. Limitações conhecidas em backlog",
            "",
            "- Backlog Fase 2 e Fase 3: itens abaixo estão declarados como fila priorizada, não como comportamento escondido.",
            "- Fase 2: a carteira ainda usa `delta_DU / 21`; uma versão futura pode converter a taxa mensal para taxa anual equivalente e aplicar `delta_DU / 252` explicitamente.",
            "- Fase 2: perda e provisão continuam em lógica mensal agregada; uma versão futura pode desdobrar por DU efetivo ou por safra.",
            "- Fase 2: o Pre DI na duration ainda é interpolado a partir dos pontos simulados; pode evoluir para interpolação direta na curva completa.",
            "- Fase 3: a trava de caixa e gatilhos de avaliação/liquidação ainda não limitam PMTs programados.",
            "- Fase 3: a SELIC de caixa pode sair de input manual para fonte/curva de mercado auditável.",
            "- Direção de viés: sem trava de caixa, cenários de stress podem superestimar pagamentos às cotas; sem safra/MOB, a perda é mais simples e menos granular que um motor de crédito completo.",
        ]
    )


def _build_step_by_step_markdown() -> str:
    return "\n".join(
        [
            "Use esta aba como um livro-caixa econômico simplificado do FIDC. A pergunta principal é: depois de juros da carteira, custos, perdas, pagamentos de SEN/MEZZ e reinvestimentos, quanto sobra de colchão econômico para a SUB?",
            "",
            "### 1. Defina o ativo que o fundo compra",
            "",
            "- **Volume da carteira** é o valor de face dos recebíveis comprados no início da simulação.",
            "- **Taxa de Cessão** é o deságio no prazo médio do recebível. **Taxa Mensal** é a taxa efetiva usada pelo motor. A aba mostra a equivalência entre as duas e a taxa anual em base 252.",
            "- **Ágio sobre face** significa pagar mais pelo recebível. Isso reduz a taxa efetiva da carteira e aumenta o EAD econômico; não aparece como despesa separada no mês zero.",
            "- **Quantidade média de parcelas** gera automaticamente o prazo médio econômico dos recebíveis. Para carteira parcelada, o modelo usa a duration Macaulay do fluxo de parcelas, não o prazo contratual final.",
            "- **Prazo médio dos recebíveis** define quanto da carteira vence por mês. Ex.: com 7 parcelas e taxa mensal de 14%, a duration econômica fica em torno de 3,15 meses; portanto a carteira gira mais rápido do que um prazo contratual de 7 meses.",
            "",
            "### 2. Defina a estrutura de capital",
            "",
            "- **SEN** recebe antes, **MEZZ** recebe depois da SEN, e **SUB** recebe apenas o residual econômico.",
            "- A SUB absorve o impacto econômico das perdas primeiro: se o PL do fundo cai depois de custos, perdas e pagamentos, a sobra residual da SUB diminui.",
            "- As regras de juros e amortização determinam os PMTs de SEN/MEZZ. A SUB não tem pagamento programado nesta etapa.",
            "",
            "### 3. Entenda a contabilidade mensal do fluxo",
            "",
            "- Em cada mês, o motor começa com a carteira e o PL econômico do mês anterior.",
            "- A carteira gera **juros/receita econômica** pela taxa da carteira.",
            "- Uma parte do principal vence. Se esse principal paga normalmente, ele vira caixa e pode ser reinvestido durante a revolvência.",
            "- A PDD de NPL 90+ é reconhecida como despesa prospectiva separada do principal que vence.",
            "- Depois entram custos, despesa de provisão/perda e PMTs de SEN/MEZZ. O que sobra no PL depois de SEN e MEZZ é a SUB residual e pode virar nova carteira enquanto houver janela elegível.",
            "",
            "### 4. Como ler NPL, provisão e write-off",
            "",
            "- Na metodologia **NPL 90 + cobertura**, o usuário informa o NPL 90+ esperado por ciclo e o lag até Over90.",
            "- O modelo provisiona essa perda de forma prospectiva antes de ela aparecer como estoque NPL 90+. É uma visão econômica tipo ECL, não uma regra contábil/regulatória completa.",
            "- Quando o crédito entra em NPL 90+, o motor aplica a LGD. Com LGD de 100%, a entrada em NPL 90+ é baixada integralmente.",
            "- A baixa consome a provisão já formada. Se a provisão for insuficiente, a diferença vira **write-off descoberto** e aumenta a despesa de perda do mês.",
            "- Importante: o write-off reduz o saldo da carteira; ele não é debitado de novo no PL fora da linha de provisão/perda. Isso evita dupla contagem.",
            "",
            "Exemplo simples: com carteira de R$ 100, 7 parcelas e NPL esperado de 40% em 3 meses, vencem R$ 14,29 de principal e a PDD mensal é R$ 13,33. A PDD reduz PL; o principal recebido segue como caixa reinvestível.",
            "",
            "### 5. Entenda a revolvência",
            "",
            "- Enquanto ainda há prazo suficiente para comprar novos recebíveis, o motor recicla o principal recebido e o excesso de caixa positivo.",
            "- **Principal recebido** é o principal que venceu no mês conforme a quantidade média de parcelas.",
            "- **Subordinação mínima** é testada como SUB disponível / PL FIDC; ela não limita a compra de carteira originada pelo denominador acumulado.",
            "- **Excesso de caixa** é o fluxo positivo remanescente após custos, perdas e pagamentos de SEN/MEZZ. Durante a janela elegível, ele compra carteira nova junto com o principal recebido.",
            "- Quando o prazo médio dos recebíveis não cabe mais no prazo restante do FIDC, a carteira entra em runoff: o principal recebido também deixa de comprar nova carteira e passa a render pela SELIC média informada.",
            "",
            "### 6. Entenda os denominadores",
            "",
            "- **Carteira originada programática** é a referência teórica de giro do volume inicial (volume + reposições de principal pelo prazo médio).",
            "- **Carteira originada efetiva** é a base reconciliada com o motor: volume inicial + soma de principal e excesso de caixa reinvestidos mês a mês.",
            "- No cenário sem perdas, a efetiva incorpora o excesso de caixa gerado pela operação. Em cenários com perdas, a PDD consome PL e pode reduzir o excesso disponível para reinvestir.",
            "- **Colchão sem perdas sobre carteira originada** compara a SUB final sem perdas com a carteira originada efetiva; é uma leitura de excess spread e colchão potencial em uma simulação paralela sem perda de crédito.",
            "- O **gráfico de proteção da estrutura** usa o cenário principal com perdas e compara a SUB disponível de cada mês com PL ou carteira originada acumulada. Por isso, ele não precisa bater com o card sem perdas; são duas leituras complementares.",
            "",
            "### 7. Como ler os gráficos e a timeline",
            "",
            "- O gráfico de saldos mostra SEN, MEZZ, SUB residual e déficit econômico ao longo do tempo.",
            "- O gráfico de perda mostra a perda do período: despesa de provisão/perda do mês dividida pela carteira do mês; não é write-off acumulado nem estoque NPL acumulado.",
            "- O gráfico de proteção compara a subordinação econômica com o colchão de proteção sobre a carteira originada acumulada.",
            "- A timeline detalhada é a memória de cálculo: use as colunas de PDD, provisão, baixa, reinvestimento e carteira originada para reconciliar cada mês.",
        ]
    )


def _validate_model_inputs(inputs) -> list[str]:
    errors: list[str] = []
    if not inputs.datas:
        errors.append("datas do fluxo ausentes em model_data.json")
    if not inputs.feriados:
        errors.append("lista de feriados vazia em model_data.json")
    return errors


def _validate_snapshot_curve(inputs) -> list[str]:
    errors: list[str] = []
    if not inputs.curva_du or not inputs.curva_cdi:
        errors.append("curva DI/Pre ausente em model_data.json")
    if len(inputs.curva_du) != len(inputs.curva_cdi):
        errors.append("curva_du e curva_cdi têm tamanhos diferentes")
    return errors


def _curve_comparison_metrics(inputs, selected_curve: _SelectedCurve) -> dict[str, float] | None:
    if inputs is None or selected_curve.source_label == CURVE_SOURCE_SNAPSHOT:
        return None
    if selected_curve.base_date != SNAPSHOT_CURVE_DATE:
        return None

    snapshot_by_du = {int(du): float(rate) for du, rate in zip(inputs.curva_du, inputs.curva_cdi)}
    b3_by_du = {int(du): float(rate) for du, rate in zip(selected_curve.curva_du, selected_curve.curva_taxa_aa)}
    common_du = sorted(set(snapshot_by_du) & set(b3_by_du))
    if not common_du:
        return None

    diffs = [b3_by_du[du] - snapshot_by_du[du] for du in common_du]
    abs_diffs = [abs(value) for value in diffs]
    return {
        "pontos_comuns": float(len(common_du)),
        "diferenca_media_bps": sum(abs_diffs) / len(abs_diffs) * 10_000.0,
        "diferenca_max_bps": max(abs_diffs) * 10_000.0,
    }


def _render_curve_source_info(inputs, selected_curve: _SelectedCurve) -> None:
    st.info(
        "\n".join(
            [
                f"Fonte da curva: {selected_curve.source_label}.",
                f"Curva {selected_curve.curve_code} em % a.a., base 252 dias úteis.",
                f"Data-base: {selected_curve.base_date:%d/%m/%Y}; pontos: {_format_number_br(selected_curve.point_count, 0)}; "
                f"DU inicial/final: {selected_curve.first_du or 'N/D'} / {selected_curve.last_du or 'N/D'}.",
            ]
        )
    )
    if selected_curve.source_label != CURVE_SOURCE_SNAPSHOT:
        st.caption(
            f"Arquivo B3: `{selected_curve.source_url}` | baixado em {selected_curve.retrieved_label} | "
            f"SHA-256: `{selected_curve.content_sha256[:12]}`"
        )
    else:
        st.caption(
            "Curva local usada apenas como referência histórica: dados salvos no repositório "
            f"com data-base {SNAPSHOT_CURVE_DATE:%d/%m/%Y}."
        )

    metrics = _curve_comparison_metrics(inputs, selected_curve)
    if metrics is not None:
        st.caption(
            "Comparação B3 x curva local na mesma data: "
            f"{_format_number_br(metrics['pontos_comuns'], 0)} DUs comuns; "
            f"diferença média absoluta {_format_number_br(metrics['diferenca_media_bps'], 2)} bps; "
            f"diferença máxima absoluta {_format_number_br(metrics['diferenca_max_bps'], 2)} bps."
        )
    elif selected_curve.source_label != CURVE_SOURCE_SNAPSHOT:
        st.caption(
            f"A curva local salva é de {SNAPSHOT_CURVE_DATE:%d/%m/%Y}; por isso a comparação ponto a ponto só é exibida "
            "quando a data B3 selecionada é a mesma."
        )


def _render_calendar_source_info(selected_calendar: _SelectedCalendar) -> None:
    official_years = ", ".join(str(year) for year in selected_calendar.official_years) or "N/D"
    projected_years = ", ".join(str(year) for year in selected_calendar.projected_years) or "N/D"
    st.info(
        "\n".join(
            [
                f"Calendário de dias úteis: {selected_calendar.source_label}.",
                f"Feriados considerados: {_format_number_br(selected_calendar.holiday_count, 0)}.",
                f"Anos oficiais B3: {official_years}; anos projetados: {projected_years}.",
                "Os DUs dos fluxos mensais são recalculados com esse calendário antes da interpolação da curva.",
            ]
        )
    )
    if selected_calendar.source_label == CALENDAR_SOURCE_B3_OFFICIAL:
        st.caption(
            f"Calendário B3 consultado em {selected_calendar.retrieved_label}; "
            f"SHA-256 `{selected_calendar.content_sha256[:12]}`. "
            "Anos ainda não publicados pela B3 são preenchidos por projeção explícita."
        )
    elif selected_calendar.source_label == CALENDAR_SOURCE_B3_PROJECTED:
        st.caption(
            "O calendário projetado combina os feriados locais salvos com regras de mercado para anos futuros "
            "do fluxo, incluindo feriados nacionais, Carnaval, Sexta-feira Santa, Corpus Christi e datas sem pregão "
            "como 24/12 e 31/12."
        )
    else:
        last_holiday = selected_calendar.last_holiday.strftime("%d/%m/%Y") if selected_calendar.last_holiday else "N/D"
        st.warning(
            "O calendário local salvo é útil para comparação histórica, mas a lista termina em "
            f"{last_holiday}."
        )


def _render_curve_source_controls(
    inputs,
    selected_curve: _SelectedCurve | None = None,
    selected_calendar: _SelectedCalendar | None = None,
) -> None:
    with st.expander("Fonte da curva DI/Pré", expanded=False):
        st.selectbox(
            "Origem da curva de juros",
            CURVE_SOURCE_OPTIONS,
            key="modelo_curve_source",
            help=(
                "A curva DI/Pré remunera cotas pós-fixadas e calcula o Pre DI na duration. "
                "Use B3 para dados de mercado ou curva local apenas para comparação histórica."
            ),
        )
        if st.session_state.get("modelo_curve_source") == CURVE_SOURCE_B3_DATE:
            st.date_input(
                "Data-base B3",
                key="modelo_b3_date",
                min_value=date(2000, 1, 1),
                max_value=date.today(),
                help="Data exata; sem fallback silencioso.",
            )
        st.selectbox(
            "Metodologia de interpolação",
            INTERPOLATION_OPTIONS,
            key="modelo_interpolation_label",
            help=(
                "Flat Forward 252 preserva a composição financeira em dias úteis entre vértices. "
                "Spline fica disponível como metodologia alternativa."
            ),
        )
        st.selectbox(
            "Calendário de dias úteis",
            CALENDAR_SOURCE_OPTIONS,
            key="modelo_calendar_source",
            help=(
                "O calendário oficial usa a página pública da B3 para anos publicados e projeção explícita para anos futuros. "
                "Os feriados locais salvos ficam disponíveis para comparação histórica."
            ),
        )
        if selected_curve is not None and selected_calendar is not None:
            st.markdown("##### Detalhes da fonte selecionada")
            _render_curve_source_info(inputs, selected_curve)
            _render_calendar_source_info(selected_calendar)


def _render_selic_projection_info(
    *,
    selic_aa_por_ano: tuple[tuple[int, float], ...],
) -> None:
    with st.expander("SELIC média para caixa em runoff", expanded=False):
        st.caption(
            "Esta premissa remunera apenas o caixa excedente depois que a carteira entra em runoff. "
            "Cotas pós-fixadas e Pre DI na duration continuam usando a curva DI/Pré da B3 ou a fonte selecionada."
        )
        curve_df = pd.DataFrame(
            [
                {
                    "Ano": _selic_year_label(year, [projection_year for projection_year, _ in selic_aa_por_ano]),
                    "SELIC média para caixa (% a.a.)": _format_percent(rate),
                }
                for year, rate in selic_aa_por_ano
            ]
        )
        st.dataframe(curve_df, width="stretch", hide_index=True)


def _rate_mode_from_label(label: str) -> str:
    return RATE_MODE_PRE if label.startswith("Pré") else RATE_MODE_POST_CDI


def _build_balance_area_frame(frame: pd.DataFrame) -> pd.DataFrame:
    chart_frame = frame[["indice", "data", "pl_senior", "pl_mezz", "pl_sub_jr"]].copy()
    chart_frame["pl_sub_available"] = chart_frame["pl_sub_jr"].clip(lower=0.0)
    chart_frame["deficit_economico"] = chart_frame["pl_sub_jr"].where(chart_frame["pl_sub_jr"] < 0.0, 0.0)
    stack_items = [
        ("pl_senior", LABELS_COTAS["sen"], 1),
        ("pl_mezz", LABELS_COTAS["mezz"], 2),
        ("pl_sub_available", LABELS_COTAS["sub"], 3),
    ]
    rows: list[dict[str, object]] = []
    for row in chart_frame.itertuples(index=False):
        base = 0.0
        for column, label, order in stack_items:
            value = max(float(getattr(row, column)), 0.0)
            top = base + value
            rows.append(
                {
                    "indice": row.indice,
                    "data": row.data,
                    "classe": label,
                    "ordem": order,
                    "valor": value,
                    "stack_base": base,
                    "stack_top": top,
                }
            )
            base = top
        deficit = min(float(row.deficit_economico), 0.0)
        rows.append(
            {
                "indice": row.indice,
                "data": row.data,
                "classe": "Déficit econômico",
                "ordem": 4,
                "valor": deficit,
                "stack_base": 0.0,
                "stack_top": deficit,
            }
        )
    long_df = pd.DataFrame(rows)
    long_df["valor_milhoes"] = long_df["valor"] / 1_000_000.0
    long_df["stack_base_milhoes"] = long_df["stack_base"] / 1_000_000.0
    long_df["stack_top_milhoes"] = long_df["stack_top"] / 1_000_000.0
    long_df["valor_formatado"] = long_df["valor"].map(_format_brl)
    long_df["periodo"] = long_df["data"].dt.strftime("%d/%m/%Y")
    long_df["mes_fidc"] = long_df["indice"].map(lambda value: f"Mês {int(value)}")
    return long_df


def _available_subordination_pct(row: pd.Series) -> float | None:
    pl_fidc = row.get("pl_fidc")
    pl_sub_jr = row.get("pl_sub_jr")
    if pd.isna(pl_fidc) or pd.isna(pl_sub_jr) or pl_fidc <= 0:
        return None
    return max(float(pl_sub_jr), 0.0) / float(pl_fidc)


def _build_loss_area_frame(frame: pd.DataFrame) -> pd.DataFrame:
    loss_column = "perda_carteira_despesa" if "perda_carteira_despesa" in frame.columns else "inadimplencia_despesa"
    chart_frame = frame[["indice", "data", "carteira", loss_column]].copy()
    chart_frame = chart_frame.rename(columns={loss_column: "perda_carteira_despesa"})
    chart_frame["perda_periodo_pct"] = chart_frame.apply(
        lambda row: row["perda_carteira_despesa"] / row["carteira"] if row["carteira"] else None,
        axis=1,
    )
    long_df = chart_frame.dropna(subset=["perda_periodo_pct"]).copy()
    long_df["serie"] = "Perda do período"
    long_df["valor"] = long_df["perda_periodo_pct"]
    long_df["formula_tooltip"] = "Numerador: despesa de provisão do mês. Denominador: carteira no início do mês."
    long_df["numerador_formatado"] = long_df["perda_carteira_despesa"].map(_format_brl)
    long_df["denominador_formatado"] = long_df["carteira"].map(_format_brl)
    long_df["valor_pct"] = long_df["valor"] * 100.0
    long_df["valor_formatado"] = long_df["valor"].map(_format_percent)
    long_df["periodo"] = long_df["data"].dt.strftime("%d/%m/%Y")
    long_df["mes_fidc"] = long_df["indice"].map(lambda value: f"Mês {int(value)}")
    return long_df[
        [
            "indice",
            "data",
            "serie",
            "valor",
            "formula_tooltip",
            "numerador_formatado",
            "denominador_formatado",
            "valor_pct",
            "valor_formatado",
            "periodo",
            "mes_fidc",
        ]
    ]


def _build_protection_area_frame(
    frame: pd.DataFrame,
    protection_frame: pd.DataFrame | None = None,
) -> pd.DataFrame:
    chart_frame = frame[["indice", "data", "pl_fidc", "pl_sub_jr"]].copy()
    chart_frame["subordinacao_display"] = chart_frame.apply(_available_subordination_pct, axis=1)
    chart_frame["numerador_formatado"] = chart_frame["pl_sub_jr"].clip(lower=0.0).map(_format_brl)
    chart_frame["denominador_formatado"] = chart_frame["pl_fidc"].map(_format_brl)
    long_df = chart_frame.melt(
        id_vars=["indice", "data", "numerador_formatado", "denominador_formatado"],
        value_vars=["subordinacao_display"],
        var_name="serie",
        value_name="valor",
    ).dropna(subset=["valor"])
    long_df["serie"] = "Subordinação econômica"
    long_df["formula_tooltip"] = "Numerador: SUB disponível, max(PL FIDC - PL SEN - PL MEZZ, 0). Denominador: PL FIDC."
    long_df["valor_pct"] = long_df["valor"] * 100.0
    long_df["valor_formatado"] = long_df["valor"].map(_format_percent)
    long_df["periodo"] = long_df["data"].dt.strftime("%d/%m/%Y")
    long_df["mes_fidc"] = long_df["indice"].map(lambda value: f"Mês {int(value)}")
    if protection_frame is not None and not protection_frame.empty:
        protection_series = protection_frame[
            [
                "indice",
                "data",
                "perda_maxima_suportada",
                "valor_pct",
                "valor_formatado",
                "sub_formatada",
                "originada_formatada",
                "periodo",
                "mes_fidc",
            ]
        ].copy()
        protection_series = protection_series.rename(
            columns={
                "perda_maxima_suportada": "valor",
                "sub_formatada": "numerador_formatado",
                "originada_formatada": "denominador_formatado",
            }
        )
        protection_series["serie"] = "Colchão de proteção"
        protection_series["formula_tooltip"] = (
            "Numerador: SUB disponível no mês. Denominador: carteira originada acumulada (volume inicial + nova originação acumulada)."
        )
        long_df = pd.concat([long_df, protection_series[long_df.columns]], ignore_index=True)
    return long_df


def _protection_ratio_chart(chart_df: pd.DataFrame) -> alt.Chart:
    base = alt.Chart(chart_df).encode(
        x=alt.X("indice:Q", title="Mês do FIDC", axis=alt.Axis(tickMinStep=1)),
        y=alt.Y(
            "valor_pct:Q",
            title="% do denominador",
            axis=alt.Axis(labelExpr="replace(format(datum.value, '.1f'), '.', ',') + '%'"),
        ),
        color=alt.Color(
            "serie:N",
            title="Cenário",
            scale=alt.Scale(range=[COLOR_BLACK, COLOR_ORANGE]),
        ),
        tooltip=[
            alt.Tooltip("mes_fidc:N", title="Mês"),
            alt.Tooltip("periodo:N", title="Período"),
            alt.Tooltip("serie:N", title="Cenário"),
            alt.Tooltip("carteira_inicial_formatada:N", title="Carteira inicial"),
            alt.Tooltip("nova_originacao_formatada:N", title="Nova originação"),
            alt.Tooltip("nova_originacao_acumulada_formatada:N", title="Nova originação acumulada"),
            alt.Tooltip("nova_originacao_motor_formatada:N", title="Originação econômica do motor"),
            alt.Tooltip("residual_fluxo_formatado:N", title="Residual do fluxo"),
            alt.Tooltip("sub_formatada:N", title="SUB disponível"),
            alt.Tooltip("originada_formatada:N", title="Denominador"),
            alt.Tooltip("valor_formatado:N", title="Colchão de proteção"),
        ],
    )
    return (base.mark_line(size=2.4, interpolate="monotone") + base.mark_point(size=42)).properties(height=320)


def _area_money_chart(chart_df: pd.DataFrame) -> alt.Chart:
    color_domain = [LABELS_COTAS["sen"], LABELS_COTAS["mezz"], LABELS_COTAS["sub"], "Déficit econômico"]
    color_range = [COLOR_BLACK, COLOR_GRAY, COLOR_ORANGE, COLOR_DARK_ORANGE]
    x_encoding = alt.X("indice:Q", title="Mês do FIDC", axis=alt.Axis(tickMinStep=1))
    color_encoding = alt.Color(
        "classe:N",
        title="Classe",
        scale=alt.Scale(domain=color_domain, range=color_range),
    )
    tooltip = [
        alt.Tooltip("mes_fidc:N", title="Mês"),
        alt.Tooltip("periodo:N", title="Período"),
        alt.Tooltip("classe:N", title="Classe"),
        alt.Tooltip("valor_formatado:N", title="Valor"),
    ]
    area_base = alt.Chart(chart_df).encode(
        x=x_encoding,
        y=alt.Y(
            "stack_top_milhoes:Q",
            title="R$ milhões",
            axis=alt.Axis(labelExpr="replace(format(datum.value, '.1f'), '.', ',')"),
        ),
        y2=alt.Y2("stack_base_milhoes:Q"),
        color=color_encoding,
        order=alt.Order("ordem:Q", sort="ascending"),
        tooltip=tooltip,
    )
    boundary_base = alt.Chart(chart_df).encode(
        x=x_encoding,
        y=alt.Y(
            "stack_top_milhoes:Q",
            title="R$ milhões",
            axis=alt.Axis(labelExpr="replace(format(datum.value, '.1f'), '.', ',')"),
        ),
        color=color_encoding,
        order=alt.Order("ordem:Q", sort="ascending"),
        tooltip=tooltip,
    )
    area = area_base.mark_area(opacity=0.86, interpolate="monotone")
    boundary = boundary_base.mark_line(size=0.9, opacity=0.72, interpolate="monotone")
    return (area + boundary).properties(height=320)


def _area_percent_chart(
    chart_df: pd.DataFrame,
    *,
    y_title: str,
    color_domain: list[str],
    color_range: list[str],
) -> alt.Chart:
    x_encoding = alt.X("indice:Q", title="Mês do FIDC", axis=alt.Axis(tickMinStep=1))
    tooltip = [
        alt.Tooltip("mes_fidc:N", title="Mês"),
        alt.Tooltip("periodo:N", title="Período"),
        alt.Tooltip("serie:N", title="Série"),
        alt.Tooltip("valor_formatado:N", title="Valor"),
    ]
    if "numerador_formatado" in chart_df.columns:
        tooltip.append(alt.Tooltip("numerador_formatado:N", title="Numerador"))
    if "denominador_formatado" in chart_df.columns:
        tooltip.append(alt.Tooltip("denominador_formatado:N", title="Denominador"))
    if "formula_tooltip" in chart_df.columns:
        tooltip.append(alt.Tooltip("formula_tooltip:N", title="Cálculo"))

    base = alt.Chart(chart_df).encode(
        x=x_encoding,
        y=alt.Y(
            "valor_pct:Q",
            title=y_title,
            axis=alt.Axis(labelExpr="replace(format(datum.value, '.1f'), '.', ',') + '%'"),
        ),
        color=alt.Color(
            "serie:N",
            title="Série",
            scale=alt.Scale(domain=color_domain, range=color_range),
        ),
        tooltip=tooltip,
    )
    return (
        base.mark_area(opacity=0.18, interpolate="monotone")
        + base.mark_line(size=2.2, interpolate="monotone")
    ).properties(height=320)


def _chart_definition_caption(kind: str) -> str:
    if kind == "loss":
        return (
            "Perda do período = despesa de provisão do mês (numerador) dividida pela carteira do mês (denominador). "
            "Em regime estável essa razão tende a perda_ciclo / prazo médio dos recebíveis "
            "(ex.: 40% / 3,15 meses = ~12,7% a.m.); não existe nenhum teto de 25% ou 40% nesta série mensal — "
            "perda_ciclo é uma taxa de perda do ciclo/safra (sobre o que vence em ~prazo médio meses), não uma taxa mensal direta. "
            "O efeito acumulado das perdas sobre o colchão da SUB aparece na evolução do saldo das cotas e no gráfico de proteção da estrutura."
        )
    if kind == "protection":
        return (
            "Este gráfico usa o cenário principal com perdas. Subordinação econômica = SUB disponível no mês, "
            "max(PL FIDC - PL SEN - PL MEZZ, 0) (numerador), "
            "dividida pelo PL FIDC do mês (denominador). "
            "Colchão de proteção = SUB disponível no mês (numerador) dividida pela carteira originada acumulada, "
            "isto é, volume inicial + nova originação acumulada (denominador). O card sem perdas é uma simulação paralela."
        )
    raise ValueError(f"Tipo de legenda de gráfico inválido: {kind}")


def _render_model_header() -> None:
    st.markdown(_MODEL_CSS, unsafe_allow_html=True)
    st.markdown(
        """
        <div class="fidc-model-header">
          <div class="fidc-model-kicker">Simulação econômica</div>
          <h2 class="fidc-model-title">Modelagem</h2>
          <div class="fidc-model-copy">
            Este é um modelo econômico-financeiro para simular cenários de perda máxima em uma
            carteira de FIDC, considerando diferentes níveis de rentabilidade, subordinação e
            perdas de crédito.
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def _model_tooltip_html(tooltip: str) -> str:
    safe_tooltip = escape(tooltip, quote=True)
    return (
        f'<span class="fidc-model-tooltip" title="{safe_tooltip}" '
        f'data-tooltip="{safe_tooltip}" aria-label="{safe_tooltip}" tabindex="0">?</span>'
    )


def _render_model_kpi_cards(kpis, results, *, has_mezz: bool) -> None:
    cards = _model_kpi_cards_data(kpis, results, has_mezz=has_mezz)
    cards_html = "".join(
        (
            '<div class="fidc-model-kpi-card">'
            f'<div class="fidc-model-kpi-label"><span>{escape(card["label"])}</span>'
            f'{_model_tooltip_html(card["tooltip"])}</div>'
            f'<div class="fidc-model-kpi-value">{escape(card["value"])}</div>'
            f'<div class="fidc-model-kpi-context">{escape(card["context"])}</div>'
            "</div>"
        )
        for card in cards
    )
    st.markdown(f'<div class="fidc-model-kpi-grid">{cards_html}</div>', unsafe_allow_html=True)


def _model_kpi_cards_data(kpis, results, *, has_mezz: bool) -> list[dict[str, str]]:
    cards = [
        (
            "Retorno anualizado",
            _format_percent(kpis.xirr_senior),
            LABELS_COTAS["sen"],
            "Taxa interna de retorno anual dos fluxos da SEN, considerando juros, amortizações e datas do fluxo.",
        ),
        (
            "Retorno anualizado",
            _format_percent(kpis.xirr_sub_jr),
            "Júnior residual",
            "TIR anual do residual da SUB; fica N/D quando os fluxos não permitem uma TIR válida.",
        ),
        (
            "Duration econômica",
            f"{_format_number_br(kpis.duration_senior_anos, 2)} anos" if kpis.duration_senior_anos is not None else "N/D",
            LABELS_COTAS["sen"],
            "Prazo médio ponderado dos pagamentos da SEN; fica menor que o prazo final quando há juros ou amortização antes do vencimento.",
        ),
        (
            "Pre DI na duration",
            _format_percent(kpis.pre_di_duration),
            "Curva interpolada",
            "Taxa DI/Pré interpolada no prazo da duration da SEN, usada como referência de mercado.",
        ),
        (
            "SUB inicial",
            _format_brl(results[0].pl_sub_jr),
            "Colchão subordinado",
            "Valor inicial do colchão subordinado que absorve perdas antes de MEZZ e SEN.",
        ),
    ]
    if has_mezz:
        cards.insert(
            1,
            (
                "Retorno anualizado",
                _format_percent(kpis.xirr_mezz),
                LABELS_COTAS["mezz"],
                "Taxa interna de retorno anual dos fluxos da MEZZ, considerando sua posição no waterfall.",
            ),
        )
    if kpis.xirr_sub_jr is not None:
        cards.insert(
            1 + int(has_mezz),
            (
                "Retorno anualizado",
                _format_percent(kpis.xirr_sub_jr),
                "Júnior residual",
                "TIR anual do residual da SUB quando os fluxos permitem uma TIR válida.",
            ),
        )
    return [
        {
            "label": label,
            "value": value,
            "context": context,
            "tooltip": tooltip,
        }
        for label, value, context, tooltip in cards
    ]


def _render_revolvency_cards(metrics: _RevolvencyMetrics) -> None:
    cards = _revolvency_cards_data(metrics)
    cards_html = "".join(
        (
            '<div class="fidc-model-kpi-card">'
            f'<div class="fidc-model-kpi-label"><span>{escape(card["label"])}</span>'
            f'{_model_tooltip_html(card["tooltip"])}</div>'
            f'<div class="fidc-model-kpi-value">{escape(card["value"])}</div>'
            f'<div class="fidc-model-kpi-context">{escape(card["context"])}</div>'
            "</div>"
        )
        for card in cards
    )
    st.markdown(f'<div class="fidc-model-kpi-grid">{cards_html}</div>', unsafe_allow_html=True)


def _revolvency_cards_data(metrics: _RevolvencyMetrics) -> list[dict[str, str]]:
    cards = [
        (
            "Prazo médio recebíveis",
            f"{_format_number_br(metrics.prazo_medio_recebiveis_meses, 2)} meses",
            "Duration econômica",
            "Duration Macaulay calculada a partir da quantidade média de parcelas e da taxa mensal da carteira.",
        ),
        (
            "Carteira originada efetiva",
            _format_brl(metrics.carteira_total_originada),
            "Base de comparação",
            "Volume comprado no motor, somando carteira inicial, principal recebido e excesso de caixa reinvestido.",
        ),
        (
            "Bloqueio por trava",
            _format_brl(metrics.reinvestimento_bloqueado_subordinacao_total),
            "Inativo no preset",
            "O caixa elegível é reinvestido integralmente; a subordinação mínima é testada em SUB / PL FIDC.",
        ),
        (
            "SUB final sem perdas",
            _format_brl(metrics.sub_final_sem_inadimplencia),
            "Cenário paralelo sem perda",
            "Valor econômico residual da SUB no fim do prazo em simulação paralela com perda de crédito zerada.",
        ),
        (
            "Colchão s/ perdas sobre originada",
            _format_percent(metrics.colchao_sem_perdas_sobre_originacao),
            "SUB final sem perdas / carteira originada",
            "Não é a mesma série do gráfico com perdas; mede o excess spread acumulado sem perdas dividido pela carteira originada efetiva.",
        ),
    ]
    return [
        {
            "label": label,
            "value": value,
            "context": context,
            "tooltip": tooltip,
        }
        for label, value, context, tooltip in cards
    ]


def render_tab_modelo_fidc() -> None:
    inputs = _load_inputs("model_data.json")

    _render_model_header()
    source_errors = _validate_model_inputs(inputs)
    if source_errors:
        st.error("Fonte local do modelo incompleta: " + "; ".join(source_errors) + ".")
        return

    _apply_mc3_preset_defaults()

    curve_source_label = _ensure_session_option("modelo_curve_source", CURVE_SOURCE_OPTIONS)
    selected_b3_date = _ensure_session_date("modelo_b3_date", date.today() - timedelta(days=1))
    calendar_source_label = _ensure_session_option("modelo_calendar_source", CALENDAR_SOURCE_OPTIONS, default_index=1)
    interpolation_label = _ensure_session_option("modelo_interpolation_label", INTERPOLATION_OPTIONS)
    interpolation_method = _interpolation_method_from_label(interpolation_label)

    try:
        if curve_source_label == CURVE_SOURCE_SNAPSHOT:
            snapshot_errors = _validate_snapshot_curve(inputs)
            if snapshot_errors:
                st.error("Curva local salva incompleta: " + "; ".join(snapshot_errors) + ".")
                _render_curve_source_controls(inputs)
                return
            selected_curve = _selected_curve_from_snapshot(inputs)
        elif curve_source_label == CURVE_SOURCE_B3_DATE:
            with st.spinner("Consultando curva TaxaSwap na B3..."):
                b3_snapshot = _load_b3_curve_for_date(selected_b3_date.isoformat(), DEFAULT_TAXASWAP_CURVE_CODE)
            selected_curve = _selected_curve_from_b3(b3_snapshot, curve_source_label)
        else:
            with st.spinner("Localizando último pregão TaxaSwap disponível na B3..."):
                b3_snapshot = _load_latest_b3_curve(date.today().isoformat(), DEFAULT_TAXASWAP_CURVE_CODE)
            selected_curve = _selected_curve_from_b3(b3_snapshot, curve_source_label)
    except B3CurveError as exc:
        st.error(f"Não foi possível carregar a curva B3: {exc}")
        st.warning("Selecione a curva local salva apenas se quiser rodar uma comparação histórica sem fonte externa.")
        _render_curve_source_controls(inputs)
        return

    default_tx_cessao_am = DEFAULT_TX_CESSAO_AM
    default_tx_cessao_desagio = monthly_rate_to_cession_discount(
        default_tx_cessao_am,
        DEFAULT_PRAZO_RECEBIVEIS_MESES,
    )
    default_senior_pct = DEFAULT_PROP_SENIOR * 100.0
    default_mezz_pct = DEFAULT_PROP_MEZZ * 100.0
    default_sub_pct = DEFAULT_PROP_SUB * 100.0

    left = st.container()
    with left:
        with st.form("modelo-fidc-premissas"):
            st.markdown("##### Premissas da carteira")
            volume_text = _text_brl_input(
                "Volume da carteira (R$)",
                default=DEFAULT_VOLUME_CARTEIRA,
                key="modelo_volume",
                decimals=2,
                help_text="Valor nominal inicial da carteira em reais, sem escala implícita em milhões ou bilhões.",
            )
            taxa_cessao_input_mode = st.radio(
                "Entrada da taxa da carteira",
                [CESSION_INPUT_DISCOUNT, CESSION_INPUT_MONTHLY],
                index=1,
                horizontal=True,
                key="modelo_taxa_cessao_input_mode",
                help="Escolha entre informar deságio no prazo médio ou taxa efetiva mensal da carteira.",
            )
            if taxa_cessao_input_mode == CESSION_INPUT_DISCOUNT:
                tx_cessao_desagio_text = _text_percent_input(
                    "Taxa de Cessão (%)",
                    default=default_tx_cessao_desagio * 100.0,
                    key="modelo_tx_cessao_desagio",
                    decimals=2,
                    help_text="Deságio sobre o valor futuro no prazo médio dos recebíveis; o motor converte para taxa mensal equivalente.",
                )
                tx_cessao_mensal_text = ""
            else:
                tx_cessao_mensal_text = _text_percent_input(
                    "Taxa Mensal (%)",
                    default=default_tx_cessao_am * 100.0,
                    key="modelo_tx_cessao_mensal",
                    decimals=2,
                    help_text="Taxa efetiva mensal bruta aplicada ao saldo da carteira de recebíveis.",
                )
                tx_cessao_desagio_text = ""

            costs_a, costs_b = st.columns(2)
            with costs_a:
                custo_adm_text = _text_percent_input(
                    "Custo de administração e gestão (% a.a. sobre PL)",
                    default=DEFAULT_CUSTO_ADM_AA * 100.0,
                    key="modelo_custo_adm_pct",
                    decimals=2,
                    help_text=HELP_CUSTO_ADM_GESTAO,
                )
            with costs_b:
                custo_min_text = _text_brl_input(
                    "Custo mínimo de administração e gestão (R$/mês)",
                    default=DEFAULT_CUSTO_MIN_MENSAL,
                    key="modelo_custo_min",
                    decimals=2,
                    help_text=HELP_CUSTO_MINIMO,
                )
            st.markdown("##### Crédito e provisão")
            model_view = st.radio(
                "Sub-aba do modelo",
                [MODEL_VIEW_GERAL, MODEL_VIEW_MC3],
                horizontal=True,
                key="modelo_view",
                help="Use a sub-aba MC3 para abrir o fluxo já focado no modelo de Cartões.",
            )
            mc3_forcado = model_view == MODEL_VIEW_MC3
            credit_model_label = st.selectbox(
                "Metodologia de crédito",
                list(CREDIT_MODEL_LABELS),
                index=list(CREDIT_MODEL_LABELS).index(CREDIT_LABEL_NPL90) if mc3_forcado else 0,
                key="modelo_credit_model",
                disabled=mc3_forcado,
                help="Define se a perda vem de NPL 90+ por ciclo ou de migração mensal entre faixas de atraso.",
            )
            if mc3_forcado:
                st.caption("Sub-aba MC3 ativa: preset de cartões com metodologia NPL 90+ e trava de reinvestimento.")
            common_credit_a, common_credit_b, common_credit_c = st.columns(3)
            with common_credit_a:
                npl90_lag_text = _text_number_input(
                    "Lag até NPL 90+ (meses)",
                    default=DEFAULT_NPL90_LAG_MESES,
                    key="modelo_npl90_lag_meses",
                    decimals=0,
                    help_text="Meses entre a perda esperada do ciclo e sua entrada no estoque NPL 90+.",
                )
            with common_credit_b:
                cobertura_npl90_text = _text_percent_input(
                    "Cobertura mínima NPL 90+ (%)",
                    default=DEFAULT_COBERTURA_NPL90 * 100.0,
                    key="modelo_cobertura_npl90_pct",
                    decimals=1,
                    help_text="Percentual mínimo do estoque NPL 90+ coberto por provisão, multiplicado pela LGD no cálculo.",
                )
            with common_credit_c:
                lgd_text = _text_percent_input(
                    "LGD econômica (%)",
                    default=DEFAULT_LGD * 100.0,
                    key="modelo_lgd_pct",
                    decimals=1,
                    help_text="Percentual do NPL 90+ não recuperado e tratado como perda econômica.",
                )
            if credit_model_label in (CREDIT_LABEL_NPL90, CREDIT_LABEL_MC3):
                perda_ciclo_text = _text_percent_input(
                    "NPL 90+ esperado por ciclo (% dos recebíveis que vencem)",
                    default=DEFAULT_PERDA_CICLO * 100.0,
                    key="modelo_perda_ciclo_pct",
                    decimals=2,
                    help_text="Percentual do principal que vence no período e migra para NPL 90+ após o lag.",
                )
                if credit_model_label == CREDIT_LABEL_MC3:
                    mc3_a, mc3_b = st.columns(2)
                    with mc3_a:
                        renegociado_pct_text = _text_percent_input(
                            "Renegociado (% dos recebíveis que vencem)",
                            default=DEFAULT_RENEGOCIADO_PCT * 100.0,
                            key="modelo_mc3_renegociado_pct",
                            decimals=2,
                            help_text="Parcela renegociada no período que entra integralmente na base de PDD do MC3.",
                        )
                    with mc3_b:
                        maturacao_over90_cap_text = _text_percent_input(
                            "Teto maturação Over90 por safra (%)",
                            default=DEFAULT_MATURACAO_OVER90_CAP * 100.0,
                            key="modelo_mc3_maturacao_over90_cap",
                            decimals=1,
                            help_text="Limite máximo de maturação para Over90 aplicado ao principal que vence no período.",
                        )
                else:
                    renegociado_pct_text = "0,00%"
                    maturacao_over90_cap_text = f"{DEFAULT_MATURACAO_OVER90_CAP * 100:.1f}%".replace(".", ",")
                roll_adimplente_text = roll_1_30_text = roll_31_60_text = roll_61_90_text = "0,00%"
                recuperacao_90_text = writeoff_90_text = "0,00%"
            else:
                perda_ciclo_text = "0,00%"
                renegociado_pct_text = "0,00%"
                maturacao_over90_cap_text = f"{DEFAULT_MATURACAO_OVER90_CAP * 100:.1f}%".replace(".", ",")
                roll_a, roll_b = st.columns(2)
                with roll_a:
                    roll_adimplente_text = _text_percent_input(
                        "Rolagem atual → 1-30 (% a.m.)",
                        default=DEFAULT_ROLAGEM_ADIMPLENTE_1_30 * 100.0,
                        key="modelo_roll_adimplente_1_30_pct",
                        decimals=2,
                        help_text="Percentual mensal da carteira adimplente que entra em atraso de 1 a 30 dias.",
                    )
                    roll_31_60_text = _text_percent_input(
                        "Rolagem 31-60 → 61-90 (% a.m.)",
                        default=DEFAULT_ROLAGEM_31_60_61_90 * 100.0,
                        key="modelo_roll_31_60_61_90_pct",
                        decimals=2,
                        help_text="Percentual mensal do bucket 31-60 que migra para atraso de 61 a 90 dias.",
                    )
                    recuperacao_90_text = _text_percent_input(
                        "Recuperação 90+ (% a.m.)",
                        default=DEFAULT_RECUPERACAO_90_PLUS * 100.0,
                        key="modelo_recuperacao_90_plus_pct",
                        decimals=2,
                        help_text="Percentual mensal do estoque 90+ que é recuperado em caixa.",
                    )
                with roll_b:
                    roll_1_30_text = _text_percent_input(
                        "Rolagem 1-30 → 31-60 (% a.m.)",
                        default=DEFAULT_ROLAGEM_1_30_31_60 * 100.0,
                        key="modelo_roll_1_30_31_60_pct",
                        decimals=2,
                        help_text="Percentual mensal do bucket 1-30 que migra para atraso de 31 a 60 dias.",
                    )
                    roll_61_90_text = _text_percent_input(
                        "Rolagem 61-90 → 90+ (% a.m.)",
                        default=DEFAULT_ROLAGEM_61_90_90_PLUS * 100.0,
                        key="modelo_roll_61_90_90_plus_pct",
                        decimals=2,
                        help_text="Percentual mensal do bucket 61-90 que migra para NPL 90+.",
                    )
                    writeoff_90_text = _text_percent_input(
                        "Write-off 90+ (% a.m.)",
                        default=DEFAULT_WRITEOFF_90_PLUS * 100.0,
                        key="modelo_writeoff_90_plus_pct",
                        decimals=2,
                        help_text="Percentual mensal do estoque 90+ baixado contra provisão.",
                    )
            enhancement_a, enhancement_b = st.columns(2)
            with enhancement_a:
                agio_aquisicao_text = _text_percent_input(
                    "Ágio (% sobre face)",
                    default=DEFAULT_AGIO_AQUISICAO * 100.0,
                    key="modelo_agio_aquisicao_pct",
                    decimals=2,
                    help_text="Prêmio pago sobre o valor de face; reduz a taxa de cessão efetiva e aumenta o EAD.",
                )
            with enhancement_b:
                excesso_spread_base = st.radio(
                    "Base do excesso de spread",
                    ["% a.m.", "% a.a. base 252 dias úteis"],
                    horizontal=True,
                    help="Escolha a periodicidade do piso adicional exigido sobre a remuneração da SEN.",
                )
                excesso_spread_text = _text_percent_input(
                    "Excesso de spread sobre SEN",
                    default=DEFAULT_EXCESSO_SPREAD_SENIOR_AM * 100.0,
                    key="modelo_excesso_spread_senior",
                    decimals=2,
                    help_text="Piso adicional ao CDI + spread da SEN; a taxa da carteira não fica abaixo desse CDI+.",
                )

            st.markdown("##### Estrutura de PL")
            prop_a, prop_b, prop_c = st.columns(3)
            with prop_a:
                senior_pct_text = _text_percent_input(
                    "PL sênior/SEN (%)",
                    default=default_senior_pct,
                    key="modelo_prop_senior",
                    decimals=1,
                    help_text="Percentual do PL total alocado à cota SEN; SEN, MEZZ e SUB devem somar 100%.",
                )
            with prop_b:
                mezz_pct_text = _text_percent_input(
                    "PL mezanino/MEZZ (%)",
                    default=default_mezz_pct,
                    key="modelo_prop_mezz",
                    decimals=1,
                    help_text="Percentual do PL total alocado à cota MEZZ; use 0,0% se não houver mezanino.",
                )
            with prop_c:
                sub_pct_text = _text_percent_input(
                    "PL subordinado/SUB (%)",
                    default=default_sub_pct,
                    key="modelo_prop_sub",
                    decimals=1,
                    help_text="Percentual do PL total alocado à SUB, que absorve perdas antes das demais cotas.",
                )
            has_mezz = _parse_percent_for_visibility(mezz_pct_text, default=DEFAULT_PROP_MEZZ) > 0.000001
            st.caption("As proporções SEN, MEZZ e SUB devem somar 100,00%.")
            st.caption(
                "Trava de reinvestimento do FIDC padrão: nova originação só é permitida enquanto "
                f"SUB disponível / carteira originada acumulada permanecer em pelo menos "
                f"{_format_percent(DEFAULT_SUBORDINACAO_MINIMA_REINVESTIMENTO)}."
            )

            st.markdown("##### Remuneração das cotas")
            senior_mode_label = st.selectbox(
                "Remuneração cota sênior/SEN",
                ["Pós-fixada: CDI + spread aditivo", "Pré-fixada: taxa % a.a."],
                help=(
                    "Pós-fixada soma o spread ao CDI; por exemplo, 1,35% significa CDI + 1,35% a.a., "
                    "não 101,35% do CDI."
                ),
            )
            senior_rate_label = (
                "Spread aditivo CDI+ da cota SEN (% a.a.)"
                if _rate_mode_from_label(senior_mode_label) == RATE_MODE_POST_CDI
                else "Taxa pré-fixada cota SEN (% a.a.)"
            )
            senior_rate_text = _text_percent_input(
                senior_rate_label,
                default=DEFAULT_TAXA_SENIOR * 100.0,
                key="modelo_taxa_senior",
                decimals=2,
                help_text=(
                    "Spread somado ao CDI: 1,35% significa CDI + 1,35% a.a., não 101,35% do CDI."
                    if _rate_mode_from_label(senior_mode_label) == RATE_MODE_POST_CDI
                    else "Informe a taxa anual pré-fixada da cota SEN."
                ),
            )

            if has_mezz:
                mezz_mode_label = st.selectbox(
                    "Remuneração cota mezanino/MEZZ",
                    ["Pós-fixada: CDI + spread aditivo", "Pré-fixada: taxa % a.a."],
                    help=(
                        "Pós-fixada soma o spread ao CDI; por exemplo, 5,00% significa CDI + 5,00% a.a., "
                        "não 105,00% do CDI."
                    ),
                )
                mezz_rate_label = (
                    "Spread aditivo CDI+ da cota MEZZ (% a.a.)"
                    if _rate_mode_from_label(mezz_mode_label) == RATE_MODE_POST_CDI
                    else "Taxa pré-fixada cota MEZZ (% a.a.)"
                )
                mezz_rate_text = _text_percent_input(
                    mezz_rate_label,
                    default=DEFAULT_TAXA_MEZZ * 100.0,
                    key="modelo_taxa_mezz",
                    decimals=2,
                    help_text=(
                        "Spread somado ao CDI: 5,00% significa CDI + 5,00% a.a., não 105,00% do CDI."
                        if _rate_mode_from_label(mezz_mode_label) == RATE_MODE_POST_CDI
                        else "Informe a taxa anual pré-fixada da cota MEZZ."
                    ),
                )
            else:
                st.caption("Sem cota MEZZ: remuneração, prazo, amortização e juros da MEZZ serão ignorados.")
                mezz_mode_label = "Pré-fixada: taxa % a.a."
                mezz_rate_label = "Taxa pré-fixada cota MEZZ (% a.a.)"
                mezz_rate_text = "0,00"

            sub_mode_label = st.selectbox(
                "Remuneração cota subordinada/SUB",
                [
                    "Residual",
                    "Pós-fixada: CDI + spread aditivo (taxa-alvo)",
                    "Pré-fixada: taxa % a.a. (taxa-alvo)",
                ],
                help="A SUB absorve o residual econômico; taxas-alvo ficam explícitas sem criar pagamento programado.",
            )
            sub_rate_text = "0,00"
            if not sub_mode_label.startswith("Residual"):
                sub_rate_text = _text_percent_input(
                        (
                            "Spread aditivo CDI+ alvo da cota SUB (% a.a.)"
                            if sub_mode_label.startswith("Pós")
                            else "Taxa-alvo pré-fixada cota SUB (% a.a.)"
                        ),
                    default=0.0,
                    key="modelo_taxa_sub",
                    decimals=2,
                    help_text=(
                        "Spread-alvo somado ao CDI: 1,35% significa CDI + 1,35% a.a., não 101,35% do CDI."
                        if sub_mode_label.startswith("Pós")
                        else "Informe a taxa-alvo anual pré-fixada da SUB."
                    ),
                )

            with st.expander("Premissas avançadas de prazo, revolvência e waterfall", expanded=False):
                st.markdown("##### Prazo e originação")
                data_inicial_fidc = st.date_input(
                    "Data inicial do FIDC",
                    value=selected_curve.base_date,
                    key="modelo_data_inicial",
                    help=(
                        "Por padrão usa o último pregão da curva carregada; sobrescrever mantém a curva selecionada "
                        "e desloca a timeline do fluxo."
                    ),
                )
                date_schedule_label = st.selectbox(
                    "Cronograma do fluxo",
                    [DATE_SCHEDULE_WORKBOOK, DATE_SCHEDULE_MONTHLY],
                    index=1,
                    key="modelo_date_schedule",
                    help=(
                        "A grade semestral padrão mantém datas espaçadas por semestre. "
                        "O modo mensal gera novas competências até o prazo total informado."
                    ),
                )
                term_a, term_b, term_c, term_d = st.columns(4)
                with term_a:
                    prazo_fidc_text = _text_number_input(
                        "Prazo total do FIDC (anos)",
                        default=DEFAULT_PRAZO_ANOS,
                        key="modelo_prazo_fidc_anos",
                        decimals=1,
                        help_text="Prazo total da simulação em anos, usado para definir o último mês do fluxo.",
                    )
                with term_b:
                    qtd_parcelas_text = _text_number_input(
                        "Qtd. média de parcelas",
                        default=DEFAULT_QTD_PARCELAS_MEDIA,
                        key="modelo_qtd_parcelas_media",
                        decimals=1,
                        help_text="Número médio de parcelas da carteira. O prazo médio econômico é calculado automaticamente.",
                    )
                with term_c:
                    prazo_recebiveis_preview = _safe_installment_duration_preview(
                        qtd_parcelas_text,
                        tx_cessao_mensal_text or _format_percent_input_value(DEFAULT_TX_CESSAO_AM * 100.0, 2),
                    )
                    prazo_recebiveis_text = _format_input_value(prazo_recebiveis_preview, 2)
                    st.metric("Prazo médio econômico", f"{_format_number_br(prazo_recebiveis_preview, 2)} meses")
                    st.caption(
                        "Duration Macaulay do fluxo parcelado: principal fixo + juros sobre saldo, "
                        "descontado pela taxa mensal da carteira."
                    )
                with term_d:
                    portfolio_mode_label = st.selectbox(
                        "Originação da carteira",
                        [PORTFOLIO_MODE_REVOLVING, PORTFOLIO_MODE_STATIC],
                        key="modelo_portfolio_mode",
                        help=(
                            "No modo revolvente, o principal recebido e o excesso de caixa recompram recebíveis "
                            "enquanto o prazo médio couber no prazo restante do FIDC."
                        ),
                    )

                st.markdown("##### Caixa pós-revolvência")
                st.caption(
                    "Quando a carteira deixa de comprar novos recebíveis, o principal recebido entra em caixa "
                    "e passa a render pela SELIC média anual informada abaixo."
                )
                selic_years = _projection_years_for_term(
                    datetime.combine(data_inicial_fidc, datetime.min.time()),
                    _safe_term_years_from_text(prazo_fidc_text),
                )
                selic_text_by_year: dict[int, str] = {}
                for row_start in range(0, len(selic_years), 3):
                    selic_columns = st.columns(min(3, len(selic_years) - row_start))
                    for column, year in zip(selic_columns, selic_years[row_start : row_start + 3]):
                        with column:
                            selic_text_by_year[year] = _text_percent_input(
                                f"SELIC média {_selic_year_label(year, selic_years)} (% a.a.)",
                                default=_default_selic_rate_for_year(year) * 100.0,
                                key=f"modelo_selic_aa_{year}",
                                decimals=2,
                                help_text=(
                                    "Informe a SELIC média anual usada apenas para remunerar caixa em runoff; "
                                    "a última taxa é perpetuada para anos posteriores."
                                ),
                            )

                st.markdown("##### Cotas SEN e MEZZ")
                st.caption(
                    "O cronograma padrão preserva a regra semestral original do motor. "
                    "Os demais modos alteram os pagamentos programados no motor."
                )
                if has_mezz:
                    senior_cfg_a, senior_cfg_b = st.columns(2)
                else:
                    senior_cfg_a = st.container()
                    senior_cfg_b = None
                with senior_cfg_a:
                    prazo_senior_text = _text_number_input(
                        "Prazo cota SEN (anos)",
                        default=DEFAULT_PRAZO_ANOS,
                        key="modelo_prazo_senior_anos",
                        decimals=1,
                        help_text="Prazo legal/econômico da cota SEN, em anos, usado no cronograma de principal.",
                    )
                    senior_amort_label = st.selectbox(
                        "Amortização principal SEN",
                        list(AMORTIZATION_LABELS),
                        index=1,
                        key="modelo_amort_senior",
                        help="Define quando e como o principal da SEN é amortizado no waterfall.",
                    )
                    senior_start_text = _text_number_input(
                        "Carência principal SEN (meses)",
                        default=DEFAULT_CARENCIA_PRINCIPAL_MESES,
                        key="modelo_inicio_amort_senior",
                        decimals=0,
                        help_text="Número de meses antes do início da amortização de principal da SEN.",
                    )
                if has_mezz and senior_cfg_b is not None:
                    with senior_cfg_b:
                        prazo_mezz_text = _text_number_input(
                            "Prazo cota MEZZ (anos)",
                            default=DEFAULT_PRAZO_ANOS,
                            key="modelo_prazo_mezz_anos",
                            decimals=1,
                            help_text="Prazo legal/econômico da cota MEZZ, em anos, usado no cronograma de principal.",
                        )
                        mezz_amort_label = st.selectbox(
                            "Amortização principal MEZZ",
                            list(AMORTIZATION_LABELS),
                            index=1,
                            key="modelo_amort_mezz",
                            help="Define quando e como o principal da MEZZ é amortizado no waterfall.",
                        )
                        mezz_start_text = _text_number_input(
                            "Carência principal MEZZ (meses)",
                            default=DEFAULT_CARENCIA_PRINCIPAL_MESES,
                            key="modelo_inicio_amort_mezz",
                            decimals=0,
                            help_text="Número de meses antes do início da amortização de principal da MEZZ.",
                        )
                else:
                    prazo_mezz_text = prazo_fidc_text
                    mezz_amort_label = "Sem amortização programada"
                    mezz_start_text = "0"

                interest_columns = st.columns(3 if has_mezz else 2)
                interest_a = interest_columns[0]
                interest_b = interest_columns[1] if has_mezz else None
                interest_c = interest_columns[2] if has_mezz else interest_columns[1]
                with interest_a:
                    senior_interest_label = st.selectbox(
                        "Pagamento de juros SEN",
                        list(INTEREST_LABELS),
                        key="modelo_juros_senior",
                        help="Define se os juros da SEN são pagos periodicamente, após carência ou no vencimento.",
                    )
                if has_mezz and interest_b is not None:
                    with interest_b:
                        mezz_interest_label = st.selectbox(
                            "Pagamento de juros MEZZ",
                            list(INTEREST_LABELS),
                            key="modelo_juros_mezz",
                            help="Define se os juros da MEZZ são pagos periodicamente, após carência ou no vencimento.",
                        )
                else:
                    mezz_interest_label = "Pago em todo período"
                with interest_c:
                    prazo_sub_text = _text_number_input(
                        "Prazo cota SUB (anos)",
                        default=DEFAULT_PRAZO_ANOS,
                        key="modelo_prazo_sub_anos",
                        decimals=1,
                        help_text="Prazo informativo da SUB; nesta versão a SUB segue como residual econômico.",
                    )
                st.caption("A SUB segue residual nesta etapa; o prazo da SUB entra na documentação e na análise econômica.")
            submitted = st.form_submit_button(
                "Rodar simulação",
                width="stretch",
                on_click=_normalize_model_input_values,
            )

    try:
        volume = _parse_br_number(volume_text, field_name="Volume da carteira (R$)")
        prazo_fidc_anos = _parse_br_number(prazo_fidc_text, field_name="Prazo total do FIDC (anos)")
        qtd_parcelas_media = _parse_br_number(qtd_parcelas_text, field_name="Quantidade média de parcelas")
        agio_aquisicao = _parse_br_number(
            agio_aquisicao_text,
            field_name="Ágio (% sobre face)",
        ) / 100.0
        if taxa_cessao_input_mode == CESSION_INPUT_DISCOUNT:
            tx_cessao_desagio_nominal = _parse_br_number(tx_cessao_desagio_text, field_name="Taxa de Cessão (%)") / 100.0
            tx_cessao_am_nominal = _monthly_rate_from_discount_and_installment_duration(
                tx_cessao_desagio_nominal,
                qtd_parcelas_media,
            )
        else:
            tx_cessao_am_nominal = _parse_br_number(tx_cessao_mensal_text, field_name="Taxa Mensal (%)") / 100.0
        prazo_medio_recebiveis_meses = _installment_macaulay_duration_months(qtd_parcelas_media, tx_cessao_am_nominal)
        prazo_recebiveis_text = _format_input_value(prazo_medio_recebiveis_meses, 2)
        if taxa_cessao_input_mode == CESSION_INPUT_MONTHLY:
            tx_cessao_desagio_nominal = monthly_rate_to_cession_discount(
                tx_cessao_am_nominal,
                prazo_medio_recebiveis_meses,
            )
        tx_cessao_desagio = _effective_cession_discount_after_premium(tx_cessao_desagio_nominal, agio_aquisicao)
        tx_cessao_am = cession_discount_to_monthly_rate(tx_cessao_desagio, prazo_medio_recebiveis_meses)
        tx_cessao_aa_equivalente = monthly_to_annual_252_rate(tx_cessao_am)
        custo_adm_aa = _parse_br_number(custo_adm_text, field_name="Custo de administração e gestão (% a.a. sobre PL)") / 100.0
        custo_min = _parse_br_number(custo_min_text, field_name="Custo mínimo de administração e gestão (R$/mês)")
        modelo_credito = CREDIT_MODEL_LABELS[credit_model_label]
        perda_ciclo = _parse_br_number(
            perda_ciclo_text,
            field_name="NPL 90+ esperado por ciclo (% dos recebíveis que vencem)",
        ) / 100.0
        npl90_lag_meses = int(round(_parse_br_number(npl90_lag_text, field_name="Lag até NPL 90+ (meses)")))
        cobertura_minima_npl90 = _parse_br_number(
            cobertura_npl90_text,
            field_name="Cobertura mínima NPL 90+ (%)",
        ) / 100.0
        lgd = _parse_br_number(lgd_text, field_name="LGD econômica (%)") / 100.0
        rolagem_adimplente_1_30 = _parse_br_number(
            roll_adimplente_text,
            field_name="Rolagem atual → 1-30 (% a.m.)",
        ) / 100.0
        rolagem_1_30_31_60 = _parse_br_number(
            roll_1_30_text,
            field_name="Rolagem 1-30 → 31-60 (% a.m.)",
        ) / 100.0
        rolagem_31_60_61_90 = _parse_br_number(
            roll_31_60_text,
            field_name="Rolagem 31-60 → 61-90 (% a.m.)",
        ) / 100.0
        rolagem_61_90_90_plus = _parse_br_number(
            roll_61_90_text,
            field_name="Rolagem 61-90 → 90+ (% a.m.)",
        ) / 100.0
        recuperacao_90_plus = _parse_br_number(
            recuperacao_90_text,
            field_name="Recuperação 90+ (% a.m.)",
        ) / 100.0
        writeoff_90_plus = _parse_br_number(
            writeoff_90_text,
            field_name="Write-off 90+ (% a.m.)",
        ) / 100.0
        renegociado_pct = _parse_br_number(
            renegociado_pct_text,
            field_name="Renegociado (% dos recebíveis que vencem)",
        ) / 100.0
        maturacao_over90_cap = _parse_br_number(
            maturacao_over90_cap_text,
            field_name="Teto maturação Over90 por safra (%)",
        ) / 100.0
        if excesso_spread_base == "% a.m.":
            excesso_spread_senior_am = _parse_br_number(
                excesso_spread_text,
                field_name="Excesso de spread sobre SEN (% a.m.)",
            ) / 100.0
            excesso_spread_senior_aa = monthly_to_annual_252_rate(excesso_spread_senior_am)
        else:
            excesso_spread_senior_aa = _parse_br_number(
                excesso_spread_text,
                field_name="Excesso de spread sobre SEN (% a.a. base 252)",
            ) / 100.0
            excesso_spread_senior_am = annual_252_to_monthly_rate(excesso_spread_senior_aa)
        proporcao_senior = _parse_br_number(senior_pct_text, field_name="PL sênior/SEN (%)") / 100.0
        proporcao_mezz = _parse_br_number(mezz_pct_text, field_name="PL mezanino/MEZZ (%)") / 100.0
        proporcao_sub = _parse_br_number(sub_pct_text, field_name="PL subordinado/SUB (%)") / 100.0
        taxa_senior = _parse_br_number(senior_rate_text, field_name=senior_rate_label) / 100.0
        taxa_mezz = _parse_br_number(mezz_rate_text, field_name=mezz_rate_label) / 100.0
        taxa_sub = _parse_br_number(sub_rate_text, field_name="Taxa-alvo cota SUB (% a.a.)") / 100.0
        user_selic_aa_por_ano = tuple(
            (
                year,
                _parse_br_number(
                    text,
                    field_name=f"SELIC média {_selic_year_label(year, list(selic_text_by_year))} (% a.a.)",
                )
                / 100.0,
            )
            for year, text in sorted(selic_text_by_year.items())
        )
        prazo_senior_anos = _parse_br_number(prazo_senior_text, field_name="Prazo cota SEN (anos)")
        prazo_mezz_anos = _parse_br_number(prazo_mezz_text, field_name="Prazo cota MEZZ (anos)")
        prazo_sub_anos = _parse_br_number(prazo_sub_text, field_name="Prazo cota SUB (anos)")
        inicio_amortizacao_senior_meses = int(
            round(_parse_br_number(senior_start_text, field_name="Carência principal SEN (meses)"))
        )
        inicio_amortizacao_mezz_meses = int(
            round(_parse_br_number(mezz_start_text, field_name="Carência principal MEZZ (meses)"))
        )
    except ValueError as exc:
        st.error(str(exc))
        return

    if volume <= 0:
        st.error("O volume da carteira deve ser maior que zero.")
        return
    if prazo_fidc_anos <= 0 or prazo_medio_recebiveis_meses <= 0 or qtd_parcelas_media <= 0:
        st.error("O prazo total do FIDC, a quantidade média de parcelas e o prazo médio calculado devem ser maiores que zero.")
        return
    if min(prazo_senior_anos, prazo_mezz_anos, prazo_sub_anos) <= 0:
        st.error("Os prazos das cotas devem ser maiores que zero.")
        return
    if min(inicio_amortizacao_senior_meses, inicio_amortizacao_mezz_meses) < 0:
        st.error("A carência de principal não pode ser negativa.")
        return
    credit_percent_values = [
        perda_ciclo,
        cobertura_minima_npl90,
        lgd,
        rolagem_adimplente_1_30,
        rolagem_1_30_31_60,
        rolagem_31_60_61_90,
        rolagem_61_90_90_plus,
        recuperacao_90_plus,
        writeoff_90_plus,
    ]
    if min(credit_percent_values) < 0 or npl90_lag_meses < 0:
        st.error("Premissas de crédito, provisão, LGD e rolagem não podem ser negativas.")
        return
    if max(perda_ciclo, lgd, rolagem_adimplente_1_30, rolagem_1_30_31_60, rolagem_31_60_61_90, rolagem_61_90_90_plus, recuperacao_90_plus, writeoff_90_plus, renegociado_pct, maturacao_over90_cap) > 1.0:
        st.error("Percentuais de perda, LGD, rolagem, recuperação e write-off devem ficar entre 0,00% e 100,00%.")
        return
    if agio_aquisicao < 0 or excesso_spread_senior_am < 0:
        st.error("Ágio sobre face e excesso de spread não podem ser negativos.")
        return
    if any(rate < 0 for _, rate in user_selic_aa_por_ano):
        st.error("As taxas SELIC projetadas não podem ser negativas.")
        return

    prop_total = proporcao_senior + proporcao_mezz + proporcao_sub
    if abs(prop_total - 1.0) > 0.0001:
        st.error(f"As proporções de PL precisam somar 100,00%. Soma atual: {_format_percent(prop_total)}.")
        return

    if isinstance(data_inicial_fidc, datetime):
        data_inicial_date = data_inicial_fidc.date()
    else:
        data_inicial_date = data_inicial_fidc
    data_inicial_dt = datetime.combine(data_inicial_date, datetime.min.time())
    simulation_dates = _build_simulation_dates(inputs, date_schedule_label, prazo_fidc_anos, data_inicial_dt)
    if data_inicial_date != selected_curve.base_date:
        st.warning(
            "A data inicial do FIDC foi sobrescrita para "
            f"{data_inicial_date:%d/%m/%Y}, mas a curva DI/Pré carregada tem data-base "
            f"{selected_curve.base_date:%d/%m/%Y}. A simulação segue com essa curva e desloca a timeline."
        )
    try:
        effective_selic_aa_por_ano = _effective_selic_projection_for_dates(user_selic_aa_por_ano, simulation_dates)
    except ValueError as exc:
        st.error(str(exc))
        return

    premissas = Premissas(
        volume=volume,
        tx_cessao_am=tx_cessao_am,
        tx_cessao_cdi_aa=inputs.premissas.get("Tx Cessão (CDI+ %aa)"),
        custo_adm_aa=custo_adm_aa,
        custo_min=custo_min,
        inadimplencia=0.0,
        proporcao_senior=proporcao_senior,
        taxa_senior=taxa_senior,
        proporcao_mezz=proporcao_mezz,
        taxa_mezz=taxa_mezz,
        proporcao_subordinada=proporcao_sub,
        taxa_sub_jr=taxa_sub,
        tipo_taxa_senior=_rate_mode_from_label(senior_mode_label),
        tipo_taxa_mezz=_rate_mode_from_label(mezz_mode_label),
        tipo_taxa_sub_jr=(
            "residual"
            if sub_mode_label.startswith("Residual")
            else RATE_MODE_PRE
            if sub_mode_label.startswith("Pré")
            else RATE_MODE_POST_CDI
        ),
        prazo_fidc_anos=prazo_fidc_anos,
        prazo_medio_recebiveis_meses=prazo_medio_recebiveis_meses,
        qtd_parcelas_media=qtd_parcelas_media,
        carteira_revolvente=portfolio_mode_label == PORTFOLIO_MODE_REVOLVING,
        prazo_senior_anos=prazo_senior_anos,
        prazo_mezz_anos=prazo_mezz_anos,
        prazo_sub_jr_anos=prazo_sub_anos,
        amortizacao_senior=AMORTIZATION_LABELS[senior_amort_label],
        amortizacao_mezz=AMORTIZATION_LABELS[mezz_amort_label],
        juros_senior=INTEREST_LABELS[senior_interest_label],
        juros_mezz=INTEREST_LABELS[mezz_interest_label],
        inicio_amortizacao_senior_meses=inicio_amortizacao_senior_meses,
        inicio_amortizacao_mezz_meses=inicio_amortizacao_mezz_meses,
        perda_esperada_am=0.0,
        perda_inesperada_am=0.0,
        modelo_credito=modelo_credito,
        perda_ciclo=perda_ciclo,
        npl90_lag_meses=npl90_lag_meses,
        cobertura_minima_npl90=cobertura_minima_npl90,
        lgd=lgd,
        rolagem_adimplente_1_30=rolagem_adimplente_1_30,
        rolagem_1_30_31_60=rolagem_1_30_31_60,
        rolagem_31_60_61_90=rolagem_31_60_61_90,
        rolagem_61_90_90_plus=rolagem_61_90_90_plus,
        recuperacao_90_plus=recuperacao_90_plus,
        writeoff_90_plus=writeoff_90_plus,
        renegociado_pct=renegociado_pct if modelo_credito == CREDIT_MODEL_MC3_CARTOES else 0.0,
        maturacao_over90_cap=maturacao_over90_cap if modelo_credito == CREDIT_MODEL_MC3_CARTOES else DEFAULT_MATURACAO_OVER90_CAP,
        agio_aquisicao=agio_aquisicao,
        excesso_spread_senior_am=excesso_spread_senior_am,
        subordinacao_minima_reinvestimento=DEFAULT_SUBORDINACAO_MINIMA_REINVESTIMENTO,
        selic_aa_por_ano=effective_selic_aa_por_ano,
    )

    if agio_aquisicao > 0.0:
        st.caption(
            "Equivalência da taxa da carteira: "
            f"Taxa de Cessão nominal {_format_percent(tx_cessao_desagio_nominal)} menos ágio "
            f"{_format_percent(agio_aquisicao)} = taxa efetiva {_format_percent(tx_cessao_desagio)} "
            f"no prazo médio de {_format_number_br(prazo_medio_recebiveis_meses, 1)} meses | "
            f"Taxa Mensal efetiva {_format_percent(tx_cessao_am)} | "
            f"Taxa anual base 252 {_format_percent(tx_cessao_aa_equivalente)}."
        )
    else:
        st.caption(
            "Equivalência da taxa da carteira: "
            f"Taxa de Cessão {_format_percent(tx_cessao_desagio)} no prazo médio de "
            f"{_format_number_br(prazo_medio_recebiveis_meses, 1)} meses | "
            f"Taxa Mensal {_format_percent(tx_cessao_am)} | "
            f"Taxa anual base 252 {_format_percent(tx_cessao_aa_equivalente)}."
        )
    if modelo_credito == CREDIT_MODEL_NPL90:
        st.caption(
            "Crédito e provisão: "
            f"NPL 90+ esperado {_format_percent(perda_ciclo)} dos recebíveis que vencem; "
            f"lag {npl90_lag_meses} meses; cobertura mínima {_format_percent(cobertura_minima_npl90)} "
            f"do NPL 90+ com LGD {_format_percent(lgd)}."
        )
    elif modelo_credito == CREDIT_MODEL_MC3_CARTOES:
        st.caption(
            "Crédito e provisão (MC3): "
            f"Over90 esperado {_format_percent(perda_ciclo)} dos recebíveis que vencem; "
            f"Reneg {_format_percent(renegociado_pct)} com PDD de 100%; "
            f"teto de maturação Over90 {_format_percent(maturacao_over90_cap)} por safra."
        )
    else:
        st.caption(
            "Crédito e provisão: migração mensal por faixas de atraso até NPL 90+, "
            f"cobertura mínima {_format_percent(cobertura_minima_npl90)} e LGD {_format_percent(lgd)}. "
            "Recuperações entram como caixa; write-offs baixam a carteira contra provisão."
        )
    st.caption(
        "Sofisticações da carteira: "
        f"ágio sobre face {_format_percent(agio_aquisicao)}; "
        f"excesso de spread SEN {_format_percent(excesso_spread_senior_am)} a.m. "
        f"({_format_percent(excesso_spread_senior_aa)} a.a. base 252). "
        "Por ser um piso, a taxa aplicada usa o maior valor entre a taxa informada e CDI + spread da SEN + excesso."
    )
    st.caption(
        "Caixa pós-revolvência: "
        + "; ".join(
            f"{_selic_year_label(year, [projection_year for projection_year, _ in user_selic_aa_por_ano])}: {_format_percent(rate)} a.a."
            for year, rate in user_selic_aa_por_ano
        )
        + ". Esta curva remunera apenas o caixa excedente quando a carteira entra em runoff."
    )

    try:
        if calendar_source_label == CALENDAR_SOURCE_B3_OFFICIAL:
            with st.spinner("Consultando calendário de negociação da B3..."):
                b3_calendar_snapshot = _load_b3_calendar_snapshot(simulation_dates[0].year, simulation_dates[-1].year)
            selected_calendar = _selected_calendar(
                inputs,
                calendar_source_label,
                b3_calendar_snapshot,
                datas=simulation_dates,
            )
        else:
            selected_calendar = _selected_calendar(inputs, calendar_source_label, datas=simulation_dates)
    except B3CalendarError as exc:
        st.error(f"Não foi possível carregar o calendário B3: {exc}")
        st.warning("Selecione o calendário projetado se quiser rodar a simulação sem consultar a página pública da B3.")
        _render_curve_source_controls(inputs, selected_curve)
        return

    if not sub_mode_label.startswith("Residual"):
        st.warning(
            "A taxa da SUB está registrada como taxa-alvo informativa. A versão atual não possui pagamento "
            "programado da SUB; por isso, o waterfall mantém a SUB como residual."
        )

    simulation_signature = (
        selected_curve.cache_key,
        selected_calendar.cache_key,
        interpolation_method,
        tuple(dt.isoformat() for dt in simulation_dates),
        premissas,
    )
    if (
        not submitted
        and st.session_state.get("modelo_fidc_signature") == simulation_signature
        and "modelo_fidc_periods" in st.session_state
    ):
        results = st.session_state["modelo_fidc_periods"]
        kpis = st.session_state["modelo_fidc_kpis"]
    else:
        results = build_flow(
            simulation_dates,
            selected_calendar.feriados,
            selected_curve.curva_du,
            selected_curve.curva_taxa_aa,
            premissas,
            interpolation_method=interpolation_method,
        )
        kpis = build_kpis(results)
        st.session_state["modelo_fidc_signature"] = simulation_signature
        st.session_state["modelo_fidc_periods"] = results
        st.session_state["modelo_fidc_kpis"] = kpis

    if not results:
        st.info("Sem datas suficientes para montar o fluxo.")
        return

    zero_default_results = build_flow(
        simulation_dates,
        selected_calendar.feriados,
        selected_curve.curva_du,
        selected_curve.curva_taxa_aa,
        _premissas_sem_perdas(premissas),
        interpolation_method=interpolation_method,
    )
    revolvency_metrics = _build_revolvency_metrics(
        premissas=premissas,
        zero_default_results=zero_default_results,
        portfolio_mode=portfolio_mode_label,
    )

    frame = _build_dataframe(results)
    zero_default_frame = _build_dataframe(zero_default_results)
    protection_frame = _build_time_protection_frame(
        frame,
        premissas=premissas,
        portfolio_mode=portfolio_mode_label,
        scenario_label="Cenário com perdas",
    )
    zero_protection_frame = _build_time_protection_frame(
        zero_default_frame,
        premissas=premissas,
        portfolio_mode=portfolio_mode_label,
        scenario_label="Cenário sem perdas",
    )
    protection_chart_frame = pd.concat([protection_frame, zero_protection_frame], ignore_index=True)
    export_frame = _build_export_dataframe(frame)
    display_frame = _build_display_dataframe(export_frame)
    committee_timeline_frame = _build_committee_timeline_dataframe(display_frame)
    balance_chart_df = _build_balance_area_frame(frame)
    loss_chart_df = _build_loss_area_frame(frame)
    protection_display_chart_df = _build_protection_area_frame(frame, protection_frame)

    st.markdown('<div class="fidc-model-section-title">Resumo econômico</div>', unsafe_allow_html=True)
    _render_model_kpi_cards(kpis, results, has_mezz=proporcao_mezz > 0.000001)

    st.markdown('<div class="fidc-model-section-title">Capacidade de perda e denominadores</div>', unsafe_allow_html=True)
    _render_revolvency_cards(revolvency_metrics)

    st.markdown('<div class="fidc-model-section-title">Evolução do saldo das cotas</div>', unsafe_allow_html=True)
    st.altair_chart(_area_money_chart(balance_chart_df), width="stretch")

    chart_left, chart_right = st.columns(2)
    with chart_left:
        st.markdown('<div class="fidc-model-section-title">Perda da carteira</div>', unsafe_allow_html=True)
        st.caption(_chart_definition_caption("loss"))
        st.altair_chart(
            _area_percent_chart(
                loss_chart_df,
                y_title="Perda da carteira (%)",
                color_domain=["Perda do período"],
                color_range=[COLOR_ORANGE],
            ),
            width="stretch",
        )
    with chart_right:
        st.markdown('<div class="fidc-model-section-title">Proteção da estrutura</div>', unsafe_allow_html=True)
        st.caption(_chart_definition_caption("protection"))
        st.altair_chart(
            _area_percent_chart(
                protection_display_chart_df,
                y_title="Proteção da estrutura (%)",
                color_domain=["Subordinação econômica", "Colchão de proteção"],
                color_range=[COLOR_BLACK, COLOR_ORANGE],
            ),
            width="stretch",
        )

    memory_df = pd.DataFrame(
        [
            {
                "Indicador": "Data inicial do FIDC",
                "Fórmula": "default = data-base da curva carregada; usuário pode sobrescrever",
                "Observação": f"Data inicial selecionada: {data_inicial_date:%d/%m/%Y}; curva DI/Pré data-base: {selected_curve.base_date:%d/%m/%Y}.",
            },
            {
                "Indicador": "Fluxo econômico da carteira",
                "Fórmula": "carteira * ((1 + tx_cessao_am_aplicada) ^ (delta_du / 21) - 1)",
                "Observação": "Em carteira revolvente, a base de carteira evolui com principal reciclado e excesso de caixa positivo enquanto houver prazo para nova originação.",
            },
            {
                "Indicador": "Rendimento do caixa SELIC",
                "Fórmula": "rendimento_selic = (caixa_selic_inicio + principal_para_caixa_selic) * ((1 + selic_aa) ^ (21 * meses_periodo / 252) - 1)",
                "Observação": "O caixa SELIC recebe apenas caixa não reinvestido; quando a carteira não pode mais comprar recebíveis, o principal que vence também passa a compor esse caixa e a render pela SELIC média anual informada pelo usuário.",
            },
            {
                "Indicador": "Fluxo econômico total dos ativos",
                "Fórmula": "fluxo_ativos_total = fluxo_carteira + rendimento_caixa_selic + recuperacao_credito",
                "Observação": "Este é o fluxo usado antes de custos, provisão/perda e pagamentos das cotas; a SELIC remunera o saldo de caixa acumulado que não foi reinvestido em carteira nova.",
            },
            {
                "Indicador": "De-para da Taxa de Cessão",
                "Fórmula": "tx_cessao_efetiva = tx_cessao_nominal - agio; tx_am = (1 / (1 - tx_cessao_efetiva)) ^ (1 / prazo_medio_recebiveis_meses) - 1",
                "Observação": "Ex.: face R$ 100, deságio 5,00% e ágio 1,00% implicam preço pago R$ 96, taxa efetiva de cessão 4,00% e taxa mensal menor.",
            },
            {
                "Indicador": "Conversão anual base 252",
                "Fórmula": "tx_cessao_aa_252 = (1 + tx_cessao_am) ^ (252 / 21) - 1",
                "Observação": "Conversão financeira efetiva, com 21 dias úteis médios por mês e 252 dias úteis por ano.",
            },
            {
                "Indicador": "Ágio sobre face",
                "Fórmula": "preço pago / face = 1 - taxa_cessao_efetiva; EAD = carteira_face * preço_pago_face",
                "Observação": "O ágio não é debitado como despesa inicial; ele reduz a taxa efetiva da carteira e aumenta a base de exposição em caso de perda.",
            },
            {
                "Indicador": "Piso de taxa da carteira",
                "Fórmula": "tx_cessao_am_aplicada = max(tx_informada, CDI + spread_SEN + excesso_spread)",
                "Observação": "O spread da SEN é aditivo: ele é somado ao CDI antes de comparar o piso com a taxa informada.",
            },
            {
                "Indicador": "Custo adm/gestão",
                "Fórmula": "max(PL_inicio * ((1 + custo_adm_aa) ^ (1/12) - 1), custo_min)",
                "Observação": "O usuário informa percentual em base 100. Ex.: 0,35 significa 0,35% a.a.; custo mínimo é R$/mês.",
            },
            {
                "Indicador": "Metodologia de crédito",
                "Fórmula": credit_model_label,
                "Observação": "A provisão é ECL-style/prospectiva: a aba separa principal recebido, despesa de PDD, estoque NPL 90+, recuperação e write-off para não misturar caixa com estoque.",
            },
            {
                "Indicador": "Carteira vencendo no período",
                "Fórmula": "carteira_inicio * meses_periodo / qtd_parcelas_media",
                "Observação": "É a parcela programada para vencer conforme a amortização das parcelas; no MC3, 7 parcelas implicam cerca de 1/7 da carteira por mês.",
            },
            {
                "Indicador": "PDD prospectiva",
                "Fórmula": "pdd = carteira_inicio * perda_ciclo * meses_periodo / lag_npl90_meses",
                "Observação": "Com carteira R$100, perda 40% e lag 3 meses, a PDD mensal é R$13,33. A PDD reduz PL econômico, mas não é dedução automática do principal que vence.",
            },
            {
                "Indicador": "Entrada e estoque NPL 90+",
                "Fórmula": "entrada_npl90_t = pdd_prospectiva_de_periodos_anteriores_apos_lag; estoque_npl90_t = max(estoque_t-1 + entrada_npl90_t - baixa_credito_t, 0)",
                "Observação": "Na metodologia NPL 90 + cobertura, a entrada usa a perda de ciclo depois do lag e é baixada imediatamente pelo LGD informado.",
            },
            {
                "Indicador": "Provisão requerida",
                "Fórmula": "max(max(provisao_anterior - baixa_credito, 0) + provisao_prospectiva, estoque_npl90 * cobertura_minima * LGD)",
                "Observação": "A baixa consome provisão já constituída; reforços adicionais só entram quando a provisão prospectiva ou mínima exige.",
            },
            {
                "Indicador": "Despesa de provisão/perda",
                "Fórmula": "writeoff_descoberto + max(provisao_requerida - max(provisao_anterior - baixa_credito, 0), 0)",
                "Observação": "A despesa reduz o PL econômico; a baixa reduz carteira quando o crédito matura em NPL 90+.",
            },
            {
                "Indicador": "Resultado líquido da carteira",
                "Fórmula": "fluxo_carteira + recuperacao_credito - despesa_provisao",
                "Observação": "A despesa de provisão/perda passa pelo resultado; o write-off coberto pela provisão reduz a carteira, mas não é debitado de novo no PL.",
            },
            {
                "Indicador": "Elegibilidade de reinvestimento",
                "Fórmula": "mês_fidc <= prazo_total_fidc_meses - prazo_medio_recebiveis_meses",
                "Observação": "Quando o prazo médio já não cabe no prazo restante do FIDC, nova originação fica zerada e a carteira entra em runoff.",
            },
            {
                "Indicador": "Nova originação",
                "Fórmula": "se elegível: principal_recebido + max(fluxo_remanescente_apos_MEZZ, 0); se não elegível: 0",
                "Observação": "Captura o reinvestimento integral do caixa elegível. A subordinação mínima é testada em SUB / PL FIDC e não reduz a nova originação por carteira originada acumulada.",
            },
            {
                "Indicador": "Juros sênior/MEZZ",
                "Fórmula": "pós: CDI/PreDI + spread aditivo; pré: taxa anual informada",
                "Observação": (
                    "Spread aditivo significa CDI + spread, não percentual do CDI. "
                    f"O FRA de cada período usa base 252 DU. Fonte selecionada: {selected_curve.source_label} "
                    f"({selected_curve.base_date:%d/%m/%Y}); interpolação: {interpolation_label}."
                ),
            },
            {
                "Indicador": "Dias úteis do fluxo",
                "Fórmula": "DU = dias úteis entre a data inicial e a data do fluxo, descontando fins de semana e feriados",
                "Observação": f"Calendário selecionado: {selected_calendar.source_label}.",
            },
            {
                "Indicador": "Proporções de PL",
                "Fórmula": "PL sênior + PL MEZZ + PL SUB = 100%",
                "Observação": "A SUB agora é uma premissa explícita na interface; o motor bloqueia soma diferente de 100%.",
            },
            {
                "Indicador": "Carteira originada programática",
                "Fórmula": "revolvente: volume + volume * max(prazo_total_meses - duration_meses, 0) / qtd_parcelas_media; estática: volume",
                "Observação": "Referência de giro teórico do volume inicial usando a quantidade média de parcelas para principal vencendo.",
            },
            {
                "Indicador": "Carteira originada efetiva",
                "Fórmula": "volume_inicial + soma(nova_originacao_economica_t)",
                "Observação": "Denominador reconciliado com o motor; inclui reinvestimento de principal e excesso de caixa que passou pela janela de prazo.",
            },
            {
                "Indicador": "Subordinação mínima estrutural",
                "Fórmula": "SUB_disponivel / PL_FIDC >= subordinacao_minima",
                "Observação": "Reinvestimento troca caixa por carteira e não consome PL. Quando perdas/custos quebram o piso, o motor registra aporte subordinado necessário.",
            },
            {
                "Indicador": "Aporte SUB para trava mínima",
                "Fórmula": "aporte_subordinacao_minima = max((subordinacao_minima * PL_FIDC - SUB_disponivel) / (1 - subordinacao_minima), 0)",
                "Observação": "Quando perdas ou custos consumirem a SUB abaixo do piso, o motor registra o suporte subordinado necessário para que a condição mínima não seja desrespeitada.",
            },
            {
                "Indicador": "Colchão de proteção no tempo",
                "Fórmula": "SUB disponível no mês / (carteira inicial + nova originação acumulada)",
                "Observação": "A tabela de proteção detalha carteira inicial, nova originação, prazo médio, denominador, SUB e residual do fluxo por mês.",
            },
            {
                "Indicador": "Colchão sem perdas sobre originado",
                "Fórmula": "max(SUB final com perdas 0%, 0) / carteira originada efetiva",
                "Observação": "Esta é uma simulação paralela sem perda de crédito para medir excess spread e colchão econômico antes das perdas, reconciliada com as originações do motor.",
            },
            {
                "Indicador": "Júnior residual / subordinação econômica",
                "Fórmula": "Residual econômico = PL econômico do veículo - PL sênior - PL MEZZ; subordinação disponível = max(residual, 0) / PL positivo",
                "Observação": (
                    "A versão atual não remunera a SUB programaticamente; taxas da SUB ficam como premissa-alvo informativa. "
                    "A timeline preserva a coluna histórica para conferência, mas os gráficos usam o residual corrente."
                ),
            },
        ]
    )
    with st.expander("Memória de cálculo", expanded=False):
        st.dataframe(memory_df, width="stretch", hide_index=True)

    st.markdown('<div class="fidc-model-section-title">Timeline de comitê</div>', unsafe_allow_html=True)
    st.dataframe(committee_timeline_frame, width="stretch", hide_index=True)

    csv = export_frame.to_csv(index=False).encode("utf-8")
    premissas_summary_df = _build_premissas_summary_dataframe(
        premissas=premissas,
        taxa_cessao_input_mode=taxa_cessao_input_mode,
        tx_cessao_desagio=tx_cessao_desagio,
        tx_cessao_aa_equivalente=tx_cessao_aa_equivalente,
        credit_model_label=credit_model_label,
        portfolio_mode_label=portfolio_mode_label,
        date_schedule_label=date_schedule_label,
        data_inicial=data_inicial_date,
        selected_curve=selected_curve,
        selected_calendar=selected_calendar,
        interpolation_label=interpolation_label,
        senior_mode_label=senior_mode_label,
        mezz_mode_label=mezz_mode_label,
        sub_mode_label=sub_mode_label,
        senior_amort_label=senior_amort_label,
        mezz_amort_label=mezz_amort_label,
        senior_interest_label=senior_interest_label,
        mezz_interest_label=mezz_interest_label,
        user_selic_aa_por_ano=tuple(user_selic_aa_por_ano),
    )
    kpi_cards = _model_kpi_cards_data(kpis, results, has_mezz=proporcao_mezz > 0.000001)
    revolvency_cards = _revolvency_cards_data(revolvency_metrics)
    excel_bytes = _build_model_dashboard_excel_bytes(
        export_frame=export_frame,
        kpi_cards=kpi_cards,
        revolvency_cards=revolvency_cards,
        premissas_summary_df=premissas_summary_df,
        memory_df=memory_df,
        curve_source_df=_build_curve_source_dataframe(selected_curve, selected_calendar, interpolation_label),
        revolvency_export_df=_build_revolvency_export_dataframe(revolvency_metrics),
        protection_export_df=_build_time_protection_export_dataframe(protection_chart_frame),
        balance_chart_df=balance_chart_df,
        loss_chart_df=loss_chart_df,
        protection_chart_df=protection_display_chart_df,
    )
    pptx_bytes = _build_model_dashboard_pptx_bytes(
        kpi_cards=kpi_cards,
        revolvency_cards=revolvency_cards,
        premissas_summary_df=premissas_summary_df,
        timeline_frame=frame,
        balance_chart_df=balance_chart_df,
        loss_chart_df=loss_chart_df,
        protection_chart_df=protection_display_chart_df,
    )
    st.download_button(
        "Exportar deck de comitê (PPTX)",
        data=pptx_bytes,
        file_name="modelo_fidc_dashboard.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        width="stretch",
    )
    with st.expander("Dados do modelo para diligência", expanded=False):
        st.download_button(
            "Baixar timeline CSV",
            data=csv,
            file_name="modelo_fidc_timeline.csv",
            mime="text/csv",
            width="stretch",
        )
        st.download_button(
            "Baixar dashboard Excel",
            data=excel_bytes,
            file_name="modelo_fidc_dashboard.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            width="stretch",
        )

    st.markdown('<div class="fidc-model-section-title">Fontes de juros</div>', unsafe_allow_html=True)
    _render_curve_source_controls(inputs, selected_curve, selected_calendar)
    _render_selic_projection_info(selic_aa_por_ano=user_selic_aa_por_ano)

    with st.expander("Metodologia e mecânica do modelo", expanded=False):
        st.markdown(_build_step_by_step_markdown())
        st.markdown("---")
        st.markdown(
            _build_workbook_mechanics_markdown(
                selected_curve=selected_curve,
                selected_calendar=selected_calendar,
                interpolation_label=interpolation_label,
                taxa_cessao_input_mode=taxa_cessao_input_mode,
                data_inicial=data_inicial_date,
                credit_model_label=credit_model_label,
            )
        )
