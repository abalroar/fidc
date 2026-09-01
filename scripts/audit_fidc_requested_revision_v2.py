"""Audit the 33-slide v2 director release before packaging."""

from __future__ import annotations

import hashlib
import json
from pathlib import Path
import re
from zipfile import ZipFile

import pandas as pd
from pptx import Presentation


HERE = Path(__file__).resolve()
if HERE.parent.name == "scripts":
    ROOT = HERE.parents[1]
    OUT = ROOT / "outputs/revisao_diretoria_20260901_v2"
else:
    OUT = HERE.parents[1]
    ROOT = OUT
BASE = OUT / "bases"
FULL = OUT / "Industria_FIDC_Completa_Revisada_20260901_v2.pptx"
COMPACT = OUT / "FIDC_Revisao_Diretoria_20260901_v2.pptx"

ANBIMA_TITLES = (
    "Posição competitiva do Itaú BBA",
    "Sumário executivo",
    "Originação de renda fixa",
    "Originação — market share",
    "Distribuição de renda fixa",
    "Distribuição — market share",
    "Visão por produto: mercado, posição e participação do Itaú BBA",
    "Em FIDC o Itaú BBA lidera com 45,7%",
    "As maiores operações do período",
    "Top FIDCs Middle",
    "Top FIDCs Middle — nota técnica de fontes",
    "Como interpretar o ranking ANBIMA",
    "Premissas, limitações e fontes",
)
FORBIDDEN_TITLES = (
    "Fomento Mercantil: crescimento marginal",
    "Risco estrutural · Carteira I",
    "As 15 maiores ofertas de 2025",
    "Top 15 · Histórico",
)
PROVIDER_COLORS = {
    "Itau": "FF5500",
    "Kanastra": "7030A0",
    "QI Tech": "2456D6",
    "BTG Pactual": "1D4080",
    "Oliveira Trust": "7A1F3D",
    "Bradesco": "73787D",
    "Daycoval": "BEC2C5",
    "Genial": "6EC5E9",
    "Tercon": "8D9399",
    "CBSF ou REAG": "73C6A1",
    "Finaxis": "5B6065",
    "BRL Trust": "454A4F",
    "Hemera": "30353A",
}


def sha(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def slide_text(slide) -> str:  # noqa: ANN001
    return " ".join(
        shape.text.strip().replace("\n", " ")
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text.strip()
    )


def has_native_notes(presentation: Presentation) -> bool:
    return any(
        slide.has_notes_slide and "[Sources]" in slide.notes_slide.notes_text_frame.text
        for slide in presentation.slides
    )


def main() -> None:
    stock = pd.read_csv(BASE / "saldo_cenarios.csv")
    pf = pd.read_csv(BASE / "pfpj_24_incluidos.csv", dtype={"cnpj_fundo": str})
    decisions = pd.read_csv(BASE / "pfpj_26_decisoes.csv", dtype={"cnpj_fundo": str})
    issue = pd.read_csv(BASE / "emissoes_por_categoria.csv")
    issue_audit = pd.read_csv(BASE / "emissoes_auditoria.csv")
    latest = stock[
        (stock.competencia == "2026-06")
        & (stock.cenario == "sem_tapso_petrobras")
    ]

    full = Presentation(FULL)
    compact = Presentation(COMPACT)
    texts = [slide_text(slide) for slide in full.slides]
    folded = [text.casefold() for text in texts]

    provider_slide = full.slides[27]
    provider_xml = "".join(
        rel.target_part.blob.decode("utf-8")
        for rel in provider_slide.part.rels.values()
        if rel.reltype.endswith("/chart")
    )

    anbima = [full.slides[index] for index in range(13, 26)]
    anbima_structure = [
        (
            sum(bool(getattr(shape, "has_chart", False)) for shape in slide.shapes),
            sum(bool(getattr(shape, "has_table", False)) for shape in slide.shapes),
            sum(shape.shape_type == 13 for shape in slide.shapes),
        )
        for slide in anbima
    ]

    with ZipFile(FULL) as archive:
        chart_xml = " ".join(
            archive.read(name).decode("utf-8")
            for name in archive.namelist()
            if re.fullmatch(r"ppt/charts/chart\d+\.xml", name)
        )

    compact_text = " ".join(slide_text(slide) for slide in compact.slides)
    checks = {
        "pfpj_24_fundos": len(pf) == 24,
        "decisoes_26": len(decisions) == 26,
        "exclusoes_11_26": set(
            decisions.loc[~decisions.incluir_pfpj.astype(bool), "ordem"].astype(int)
        )
        == {11, 26},
        "pfpj_pl": abs(pf.pl_brl_base.sum() - 7_102_640_442.08) < 0.01,
        "stock_total": abs(latest.pl_brl.sum() - 718_610_541_429.85) < 0.02,
        "financeiro": abs(
            latest.loc[latest.categoria == "Financeiro", "pl_brl"].iloc[0]
            - 275_294_471_988.64
        )
        < 0.02,
        "stock_share": abs(latest.share.sum() - 1) < 1e-9,
        "issue_share": issue.groupby("period_key").share.sum().sub(1).abs().max()
        < 1e-9,
        "issue_pfpj_1s26": abs(
            issue.loc[
                (issue.period_key == "jun26")
                & (issue.categoria == "Multicarteira Pulverizado PF/PJ"),
                "volume_brl",
            ].iloc[0]
            - 2_603_000_000.07
        )
        < 0.01,
        "issue_fic_1s26": abs(
            issue_audit.loc[issue_audit.period_key == "jun26", "fic_excluido_brl"].iloc[0]
            - 2_813_379_549.06
        )
        < 0.01,
        "full_slide_count_33": len(full.slides) == 33,
        "compact_slide_count_3": len(compact.slides) == 3,
        "removed_blocks_absent": not any(
            title.casefold() in text
            for title in FORBIDDEN_TITLES
            for text in folded
        ),
        "anbima_titles_and_order": all(
            expected.casefold() in actual.casefold()
            for expected, actual in zip(ANBIMA_TITLES, texts[13:26], strict=True)
        ),
        "anbima_native_tables": anbima_structure[2][1] == 2
        and anbima_structure[4][1] == 2
        and anbima_structure[6][1] == 1
        and anbima_structure[7][1] == 1,
        "anbima_native_charts": anbima_structure[3][0] == 2
        and anbima_structure[5][0] == 2,
        "anbima_no_rasterized_slides": all(pictures == 0 for _, _, pictures in anbima_structure),
        "provider_slide_position": "Volume por prestador" in texts[27],
        "provider_palette_complete": all(
            color in provider_xml for color in PROVIDER_COLORS.values()
        ),
        "no_negative_chart_axis_ids": not re.search(
            r"<(?:c:)?(?:axId|crossAx)[^>]+val=\"-\d+\"", chart_xml
        ),
        "required_methodology_full": all(
            phrase.casefold() in " ".join(texts).casefold()
            for phrase in (
                "Itaú e Kanastra separados",
                "Sólido e BizCapital",
                "Top1",
                "taxonomia congelada",
            )
        ),
        "required_methodology_compact": all(
            phrase.casefold() in compact_text.casefold()
            for phrase in (
                "Itaú e Kanastra separados",
                "Sólido e BizCapital",
                "Top1",
                "taxonomia congelada",
            )
        ),
        "sources_notes_full": has_native_notes(full),
        "sources_notes_compact": has_native_notes(compact),
        "overflow_full": "Test passed"
        in (OUT / "qa/overflow_apresentacao_completa.txt").read_text(),
        "overflow_compact": "Test passed"
        in (OUT / "qa/overflow_laminas.txt").read_text(),
    }
    checks = {key: bool(value) for key, value in checks.items()}
    if not all(checks.values()):
        raise SystemExit({key: value for key, value in checks.items() if not value})

    result = {
        "status": "pass",
        "checks": checks,
        "slide_count": len(full.slides),
        "anbima_slide_range": [14, 26],
        "provider_slide": 28,
        "provider_palette": PROVIDER_COLORS,
        "full_sha256": sha(FULL),
        "compact_sha256": sha(COMPACT),
        "pfpj_pl_brl": float(pf.pl_brl_base.sum()),
        "stock_total_sem_tapso_petrobras_brl": float(latest.pl_brl.sum()),
        "financeiro_sem_tapso_petrobras_apos_pfpj_brl": float(
            latest.loc[latest.categoria == "Financeiro", "pl_brl"].iloc[0]
        ),
        "issuance_pfpj_jan_jun_2026_brl": float(
            issue.loc[
                (issue.period_key == "jun26")
                & (issue.categoria == "Multicarteira Pulverizado PF/PJ"),
                "volume_brl",
            ].iloc[0]
        ),
    }
    target = OUT / "qa/auditoria_entrega.json"
    target.write_text(json.dumps(result, ensure_ascii=False, indent=2))
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
