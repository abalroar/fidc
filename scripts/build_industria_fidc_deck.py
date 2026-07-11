"""Gera o entregável executivo da Indústria de FIDCs (v2).

Foco: evolução da indústria. A partir das bases já materializadas em
``data/industry_study`` (Toma Conta FIDCs — derivadas do Informe Mensal FIDC +
Cadastro da CVM), consolidando gestores/administradores por conglomerado
(de-para verificado por CNPJ em ``config/conglomerados_fidc.json``), produz:

1. ``outputs/Industria_FIDC_<competencia>.xlsx`` — abas com gráficos nativos.
2. ``outputs/Industria_FIDC_<competencia>.pptx`` — deck executivo com gráficos nativos.

Cada gráfico carrega a fonte exata. Uso:
    python scripts/build_industria_fidc_deck.py
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

import numpy as np
import pandas as pd

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from services.conglomerados import resolve, depara_table  # noqa: E402

# paleta executiva
NAVY = "1F2A44"; NAVY2 = "34405C"; ORANGE = "E8703A"; ORANGE2 = "F4A97F"
GRAY = "8A94A6"; GRAYL = "D9DEE7"; WHITE = "FFFFFF"; INK = "1A1A1A"
STACK = [ORANGE, NAVY, "5B7DB1", "9AA7BF", ORANGE2, GRAY, "C0703A", "6E7B94", GRAYL, "B5C2D6"]

# fontes por gráfico (proveniência exata)
FONTES = {
    "pl": "Toma Conta FIDCs — aba Indústria (série mensal, base ex-FIC-FIDC). Origem: Informe Mensal FIDC/CVM. NÃO usa o Power BI/dashboard ANBIMA (dados presos no visual).",
    "composicao": "Toma Conta FIDCs — dimensão 'segmento' (série mensal). Taxonomia interna alinhada às classes ANBIMA. Origem: Informe Mensal FIDC/CVM.",
    "definicoes": "ANBIMA — Deliberação nº 72 (Diretriz de Classificação de FIDC).",
    "gestor": "Toma Conta FIDCs — dimensão 'gestor', consolidado por conglomerado (de-para por CNPJ). Origem: Cadastro CVM (gestor) + Informe Mensal (PL).",
    "admin": "Toma Conta FIDCs — dimensão 'administrador', consolidado por conglomerado (de-para por CNPJ). Origem: Informe Mensal FIDC/CVM.",
    "controle": "Toma Conta FIDCs — dimensão 'gestor' + taxonomia de controle (Ligada a banco / Independente Grande / Independente; definição do rodapé do deck Itaú BBA).",
    "depara": "De-para próprio, verificado por CNPJ contra industry_fund_snapshot (mai/2026). Config: config/conglomerados_fidc.json.",
    "top25": "Toma Conta FIDCs — snapshot por fundo (PL, gestor consolidado, leituras de segmento/oferta/subordinação quando há regulamento parseado). Origem: Informe Mensal FIDC/CVM + leituras regulatórias do Toma Conta.",
    "cotistas": "Informe Mensal FIDC/CVM — Tabela X (nº de cotistas) via Toma Conta. Filtro PL > R$ 200 mi.",
    "investidor": "Informe Mensal FIDC/CVM — nº de cotistas por segmento de investidor, via Toma Conta.",
    "emissoes": "Toma Conta FIDCs — variação anual de PL (ex-FIC) e captação líquida (dimensão mensal). Métrica preferida: variação de PL. Validação de magnitude: ANBIMA (imprensa, 2024). NÃO é o campo 'emissões encerradas'.",
    "partes": "Toma Conta FIDCs — leitura de regulamentos (cedentes/sacados nomeados, materialidade por PL de fundo). Base: reports/fidc_clean_named_parties.",
    "sumario": "Síntese das dimensões do Toma Conta (gestor/admin/custodiante/segmento) consolidadas por conglomerado.",
    "delta_papel": "Toma Conta FIDCs — dimensão mensal por papel, consolidada por CNPJ. Share = PL do grupo / PL total do papel na competência. Δrank = posição em dez/2023 menos posição em mai/2026. PL de estoque (inclui troca de mandato).",
    "originacao": "Toma Conta FIDCs — fundos por ANO DE 1ª OFERTA (snapshot: first_offer_year). Mede ORIGINAÇÃO (fundo novo), expurgando troca de gestor/adm de fundos já existentes. Cobertura de 1ª oferta: 1.407 de 4.219 fundos.",
    "captacao": "Toma Conta FIDCs — captação líquida = captações − resgates no mês (Informe Mensal FIDC/CVM), somada por ano. Complementa a variação de PL.",
    "top25det": "Informe Mensal Estruturado (CVM), mai/2026 + leituras de regulamento (Toma Conta). Subordinação = PL classe subordinada / PL total (Tab X_2). Cedente = TAB_I2A12 (até 9, com %). Cotistas por investidor = Tab X_1_1.",
}

ITAU_NOTA = ("Itaú = Intrag (adm) + Itaú Asset + Itaú Unibanco + Kinea (gestão) + Itaú Unibanco (custódia); "
             "números por papel, não somáveis. Itaúna Capital NÃO é Itaú. PL da indústria: ex-FIC R$ 880 bi (bruto CVM R$ 959 bi).")

PERIODOS = ["2022-12", "2023-12", "2024-12", "2025-12"]  # + competência atual anexada em runtime


# ---------------------------------------------------------------------------
# Camada de dados
# ---------------------------------------------------------------------------
def _num(s):
    return pd.to_numeric(s, errors="coerce")


def load(industry_dir: Path) -> dict:
    snap = pd.read_csv(industry_dir / "industry_fund_snapshot.csv.gz")
    snap["pl"] = _num(snap["pl"]).fillna(0.0)
    snap["cotistas"] = _num(snap.get("cotistas"))
    dm = pd.read_csv(industry_dir / "industry_dimension_monthly.csv.gz")
    cot_tipo = pd.read_csv(industry_dir / "cotistas_tipo_monthly.csv")
    cot_tipo["n_cotistas"] = _num(cot_tipo["n_cotistas"]).fillna(0.0)
    comp = str(snap["competencia"].dropna().iloc[0])
    # mapa nome -> (grupo, tipo_controle) a partir do snapshot (para consolidar séries por nome)
    name2grp = {}
    for ncol, ccol in [("gestor_nome", "gestor_cnpj"), ("admin_nome", "admin_cnpj")]:
        for _, r in snap[[ncol, ccol]].dropna().iterrows():
            name2grp.setdefault(str(r[ncol]), resolve(r[ncol], r[ccol]))
    return {"snap": snap, "dm": dm, "cot_tipo": cot_tipo, "comp": comp, "name2grp": name2grp}


def _grp(name, name2grp):
    if name in name2grp:
        return name2grp[name]
    return resolve(name, None)


# ---- agregações ----------------------------------------------------------
def pl_industria_anual(dm, comp):
    fic = dm[dm["dimension_id"] == "fic_fidc"]
    net = fic[fic["dimension_value"] == "FIDC direto"].groupby("competencia")["pl_brl"].sum() / 1e9
    gross = fic.groupby("competencia")["pl_brl"].sum() / 1e9
    def yend(s):
        out = {}
        for y in [str(a) for a in range(2015, 2027)]:
            c = [x for x in s.index if x.startswith(y)]
            if c:
                out[y] = float(s[sorted(c)[-1]])
        return out
    net_y, gross_y = yend(net), yend(gross)
    anos = sorted(net_y)
    return pd.DataFrame({"ano": anos,
                         "pl_liquido": [round(net_y[a]) for a in anos],
                         "pl_bruto": [round(gross_y.get(a, np.nan)) for a in anos]})


def composicao_anual(dm):
    seg = dm[dm["dimension_id"] == "segmento"].copy()
    seg["ano"] = seg["competencia"].str.slice(0, 4)
    # último mês de cada ano
    last = seg.groupby("ano")["competencia"].transform("max")
    d = seg[seg["competencia"] == last]
    piv = d.pivot_table(index="ano", columns="dimension_value", values="pl_brl", aggfunc="sum").fillna(0)
    piv = piv[[c for c in piv.columns if piv[c].sum() > 0]]
    # ordena colunas por tamanho
    piv = piv[piv.sum().sort_values(ascending=False).index]
    pct = piv.div(piv.sum(axis=1), axis=0) * 100
    anos = [a for a in pct.index if a >= "2020"]
    return pct.loc[anos].round(1)


def ranking_consolidado(dm, dim, name2grp, comp, top=12):
    d = dm[dm["dimension_id"] == dim].copy()
    d["grupo"] = d["dimension_value"].map(lambda n: _grp(n, name2grp)[0])
    cur = d[d["competencia"] == comp].groupby("grupo")["pl_brl"].sum() / 1e9
    out = cur.sort_values(ascending=False).head(top).round(1).rename("pl_bi").reset_index()
    out.insert(0, "#", range(1, len(out) + 1))
    return out


def evolucao_consolidada(dm, dim, name2grp, periodos, top=12):
    d = dm[dm["dimension_id"] == dim].copy()
    d["grupo"] = d["dimension_value"].map(lambda n: _grp(n, name2grp)[0])
    frames = {}
    for p in periodos:
        frames[p] = d[d["competencia"] == p].groupby("grupo")["pl_brl"].sum() / 1e9
    res = pd.DataFrame(frames).fillna(0.0)
    ult = periodos[-1]
    res = res.sort_values(ult, ascending=False).head(top)
    # ranking em cada período (sobre todos os grupos, não só top)
    full = pd.DataFrame({p: d[d["competencia"] == p].groupby("grupo")["pl_brl"].sum() for p in periodos}).fillna(0.0)
    rank_ini = full[periodos[0]].rank(ascending=False, method="min")
    rank_fim = full[periodos[-1]].rank(ascending=False, method="min")
    res = res.round(1)
    res["Δ PL (bi)"] = (res[periodos[-1]] - res[periodos[0]]).round(1)
    res["Δ rank"] = [int(rank_ini.get(g, np.nan) - rank_fim.get(g, np.nan)) if g in rank_fim else 0 for g in res.index]
    res = res.reset_index().rename(columns={"grupo": "Grupo"})
    return res


def tipo_controle_anual(dm, name2grp, periodos):
    d = dm[dm["dimension_id"] == "gestor"].copy()
    d["tc"] = d["dimension_value"].map(lambda n: _grp(n, name2grp)[1])
    frames = {}
    for p in periodos:
        frames[p] = d[d["competencia"] == p].groupby("tc")["pl_brl"].sum() / 1e9
    res = pd.DataFrame(frames).fillna(0.0)
    order = ["Independente", "Ligada a banco", "Independente Grande"]
    res = res.reindex([o for o in order if o in res.index])
    return res.round(1)


def _tipo_recebivel(row):
    """Deriva o tipo de recebível a partir das leituras do Toma Conta + nome."""
    nome = str(row.get("nome_exibicao", "")).upper()
    sub = str(row.get("subsegmento_estrategia", "") or "")
    seg = str(row.get("segmento_estrategia", "") or "")
    segp = str(row.get("segmento_principal", "") or "")
    bits = []
    for val in (sub, seg):
        if val and val not in ("Não classificado", "Sem classificação", "Sem oferta CVM mapeada", "nan"):
            bits.append(val)
    if not bits and segp and segp != "nan":
        bits.append(segp)
    # heurística de nome
    if "NAO PADRONIZAD" in nome.replace("Ã", "A").replace("Ú", "U") or " NP" in f" {nome}" or nome.endswith(" NP"):
        bits.append("Não padronizado (NP)")
    for kw, lab in [("VEICULO", "Veículos"), ("AUTO", "Veículos/Auto"), ("CONSIGNAD", "Consignado"),
                    ("CARTAO", "Cartões"), ("CARD", "Cartões"), ("AGRO", "Agro"),
                    ("JUDICIAL", "Ações judiciais"), ("PRECATOR", "Precatórios"),
                    ("IMOB", "Imobiliário"), ("ENERGIA", "Energia/Infra")]:
        if kw in nome.replace("Ç", "C").replace("Õ", "O").replace("Ã", "A"):
            bits.append(lab)
    # dedup preservando ordem
    seen, out = set(), []
    for b in bits:
        if b.lower() not in seen:
            seen.add(b.lower()); out.append(b)
    return " · ".join(out[:3]) if out else "—"


def top25_detalhado(snap, name2grp, top=25):
    t = snap.sort_values("pl", ascending=False).head(top).copy()
    t["PL (R$ bi)"] = (t["pl"] / 1e9).round(2)
    t["Gestor"] = t.apply(lambda r: _grp(r["gestor_nome"], name2grp)[0], axis=1)
    t["Tipo de recebível"] = t.apply(_tipo_recebivel, axis=1)
    sub = _num(t.get("sub_min_pct_median"))
    t["Sub. mín."] = ["" if pd.isna(x) else f"{x:.0f}%" for x in sub]
    def _ult(r):
        for c in ("latest_regulamento_date", "document_latest_date"):
            v = r.get(c)
            if isinstance(v, str) and len(v) >= 7:
                return v[:7]
        fo = r.get("first_offer_year")
        return "" if pd.isna(fo) else str(int(fo))
    t["Últ. leitura"] = t.apply(_ult, axis=1)
    def _ced(r):
        v = r.get("cedentes_top")
        return "" if (pd.isna(v) or str(v) == "nan") else str(v).split("|")[0].strip()[:26]
    t["Cedente princ."] = t.apply(_ced, axis=1)
    t.insert(0, "#", range(1, len(t) + 1))
    t["Fundo"] = t["nome_exibicao"].str.slice(0, 52)
    return t[["#", "Fundo", "Tipo de recebível", "Gestor", "PL (R$ bi)", "Sub. mín.", "Últ. leitura", "Cedente princ."]]


def cotistas_dist(snap, corte=0.2):
    big = snap[snap["pl"] > corte * 1e9].copy()
    order = ["Até 2 cotistas", "3 cotistas", "4 cotistas", "5 ou mais"]
    def b(n):
        if pd.isna(n): return None
        n = int(n)
        return "Até 2 cotistas" if n <= 2 else ("3 cotistas" if n == 3 else ("4 cotistas" if n == 4 else "5 ou mais"))
    big["b"] = big["cotistas"].map(b)
    agg = big[big["b"].isin(order)].groupby("b").agg(n_fundos=("cnpj_fundo", "count"),
                                                     pl_bi=("pl", lambda s: round(s.sum() / 1e9))).reindex(order).reset_index()
    return agg, {"n": int(len(big)), "pl": round(big["pl"].sum() / 1e9)}


def cot_investidor(cot_tipo):
    comp = sorted(cot_tipo["competencia"].astype(str).unique())[-1]
    d = cot_tipo[cot_tipo["competencia"].astype(str) == comp]
    out = d.groupby("tipo_cotista")["n_cotistas"].sum().sort_values(ascending=False).head(12)
    return out.rename("n").reset_index(), comp


def emissoes_anual(dm):
    fic = dm[dm["dimension_id"] == "fic_fidc"]
    net = fic[fic["dimension_value"] == "FIDC direto"].groupby("competencia")["pl_brl"].sum() / 1e9
    def yend(s):
        out = {}
        for y in [str(a) for a in range(2019, 2027)]:
            c = [x for x in s.index if x.startswith(y)]
            if c: out[y] = float(s[sorted(c)[-1]])
        return pd.Series(out)
    pl = yend(net)
    var = pl.diff()
    capt = dm[dm["dimension_id"] == "segmento"].groupby("competencia")["captacao_liquida_brl"].sum() / 1e9
    capt.index = pd.to_datetime(capt.index + "-01")
    capt_y = capt.groupby(capt.index.year).sum()
    out = pd.DataFrame({"ano": [int(a) for a in var.index],
                        "var_pl": var.round(0).values,
                        "capt_liq": [round(capt_y.get(int(a), np.nan)) for a in var.index]})
    return out[out["ano"] >= 2020].reset_index(drop=True)


def cedentes_sacados(reports_dir: Path, top=10):
    f = reports_dir / "fidc_clean_named_parties_20260609.csv"
    if not f.exists():
        return None, None
    p = pd.read_csv(f)
    p["mat_bi"] = _num(p["materiality_brl"]) / 1e9
    p = p[~p["party_name"].astype(str).str.match(r"^\d", na=False)]  # tira CNPJs crus
    def tab(tp):
        d = p[p["participant_type"] == tp].sort_values("mat_bi", ascending=False).head(top)
        out = d[["party_name", "setor_n1", "funds", "mat_bi"]].copy()
        out.columns = ["Parte", "Setor", "# Fundos", "Materialidade (R$ bi)"]
        out["Materialidade (R$ bi)"] = out["Materialidade (R$ bi)"].round(2)
        return out.reset_index(drop=True)
    return tab("cedente_originador"), tab("sacado_devedor")


def sumario_executivo(dm, snap, n2g, comp):
    """Números-chave para o slide de sumário."""
    def dimtop(dim):
        d = dm[dm["dimension_id"] == dim].copy()
        d["g"] = d["dimension_value"].map(lambda n: _grp(n, n2g)[0])
        return d[d["competencia"] == comp].groupby("g")["pl_brl"].sum().sort_values(ascending=False) / 1e9
    plserie = pl_industria_anual(dm, comp)
    pl_fim = float(plserie["pl_liquido"].iloc[-1])
    row23 = plserie[plserie["ano"] == "2023"]["pl_liquido"]
    pl_ini = float(row23.iloc[0]) if len(row23) else float(plserie["pl_liquido"].iloc[0])
    adm = dimtop("admin"); ges = dimtop("gestor")
    tc = tipo_controle_anual(dm, n2g, [comp])[comp]
    itau_adm = adm.get("Itau", 0.0)
    btg_adm = adm.get("BTG Pactual", 0.0); qi_adm = adm.get("QI Tech", 0.0)
    return {
        "pl_ini": pl_ini, "pl_fim": pl_fim,
        "ges_top": ges.index[0], "ges_top_v": ges.iloc[0],
        "adm_btg": btg_adm, "adm_qi": qi_adm, "adm_itau": itau_adm,
        "indep_pct": tc.get("Independente", 0) / tc.sum() * 100 if tc.sum() else 0,
        "indep_v": tc.get("Independente", 0),
    }


def deltas_por_papel(dm, dim, name2grp, periodos, top=9):
    """Trajetória de market share % em cada período + Δ share e Δ ranking, por papel.

    ``periodos`` deve ser a lista de competências a exibir (ex.: 2023-12, 2024-12,
    2025-12, competência atual). Mostra o share em cada uma para ver de onde a casa
    saiu e para onde foi.
    """
    d = dm[dm["dimension_id"] == dim].copy()
    d["g"] = d["dimension_value"].map(lambda n: _grp(n, name2grp)[0])
    pl = {p: d[d["competencia"] == p].groupby("g")["pl_brl"].sum() for p in periodos}
    df = pd.DataFrame(pl).fillna(0.0)
    share = {p: (df[p] / df[p].sum() * 100 if df[p].sum() else df[p] * 0) for p in periodos}
    p0, p1 = periodos[0], periodos[-1]
    rank0 = df[p0].rank(ascending=False, method="min")
    rank1 = df[p1].rank(ascending=False, method="min")
    order = df.sort_values(p1, ascending=False).head(top).index

    def lbl(p):
        return "Sh " + (p[:4] if p.endswith("-12") else p[:7]) + " %"

    cols = {"Grupo": list(order), "PL (R$ bi)": [round(df.loc[g, p1] / 1e9, 1) for g in order]}
    for p in periodos:
        cols[lbl(p)] = [round(share[p].get(g, 0), 1) for g in order]
    cols["Δ share"] = [round(share[p1].get(g, 0) - share[p0].get(g, 0), 1) for g in order]
    cols["Δ rank"] = [int(rank0.get(g, 0) - rank1.get(g, 0)) for g in order]
    return pd.DataFrame(cols).reset_index(drop=True)


def originacao_novos(snap, name2grp, anos=(2025, 2026), dim="gestor_nome", cnpjcol="gestor_cnpj", top=8):
    """Ranking de originação: FIDCs nascidos no ano (1ª oferta), por grupo consolidado."""
    s = snap.copy()
    s["foy"] = pd.to_numeric(s["first_offer_year"], errors="coerce")
    res = {}
    meta = {}
    for y in anos:
        nv = s[s["foy"] == y].copy()
        nv["g"] = nv.apply(lambda r: _grp(r[dim], name2grp) if False else _grp(r[dim], name2grp), axis=1)
        # usa resolve por nome+cnpj
        nv["g"] = nv.apply(lambda r: resolve(r[dim], r[cnpjcol])[0], axis=1)
        res[y] = nv.groupby("g")["pl"].sum().sort_values(ascending=False) / 1e9
        meta[y] = {"n": len(nv), "pl": round(nv["pl"].sum() / 1e9)}
    return res, meta


def captacao_liquida_anual(dm):
    capt = dm[dm["dimension_id"] == "segmento"].groupby("competencia")["captacao_liquida_brl"].sum() / 1e9
    capt.index = pd.to_datetime(capt.index + "-01")
    y = capt.groupby(capt.index.year).sum()
    out = pd.DataFrame({"ano": [int(a) for a in y.index], "capt_liq": y.round(0).values})
    return out[out["ano"] >= 2020].reset_index(drop=True)


DEFINICOES_ANBIMA = [
    ("Fomento Mercantil", "Carteira pulverizada de recebíveis originados/vendidos por diversos cedentes que antecipam recursos (duplicatas, notas promissórias, cheques). Veículo de factoring, cooperativas de crédito e assessoria financeira. (≈ multicedente/multissacado)"),
    ("Financeiro", "Recebíveis originados por instituições financeiras: crédito consignado, crédito pessoal, financiamento de veículos/leasing, crédito imobiliário e multicarteira financeiro."),
    ("Agro, Indústria e Comércio", "Recebíveis originados no setor real (empresas): infraestrutura (energia, telecom, saneamento), recebíveis comerciais (duplicatas, carnês, faturas de cartão) e agronegócio."),
    ("Outros", "Inclui Recuperação/Non-Performing Loans (direitos vencidos e inadimplentes), ações judiciais/precatórios e demais focos não enquadrados nas classes acima."),
    ("Multicedente / Multissacado", "Não é classe formal ANBIMA; é característica estrutural (vários cedentes e/ou vários sacados). Na prática, a maioria cai em 'Fomento Mercantil' ou 'Outros'."),
]


# ---------------------------------------------------------------------------
# EXCEL
# ---------------------------------------------------------------------------
def build_excel(D, path: Path):
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.chart.label import DataLabelList
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    comp = D["comp"]; dm = D["dm"]; snap = D["snap"]; n2g = D["name2grp"]
    periodos = PERIODOS + [comp]
    delta_per = [x for x in PERIODOS if x >= "2023"] + [comp]
    wb = Workbook()
    navy = PatternFill("solid", fgColor=NAVY); thin = Side(style="thin", color=GRAYL)
    bd = Border(left=thin, right=thin, top=thin, bottom=thin)
    tf = Font(name="Calibri", size=14, bold=True, color=WHITE)
    hf = Font(name="Calibri", size=10, bold=True, color=WHITE)

    def banner(ws, text, span=6):
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=span)
        c = ws.cell(1, 1, text); c.fill = navy; c.font = tf
        c.alignment = Alignment(horizontal="left", vertical="center", indent=1)
        ws.row_dimensions[1].height = 24

    def note(ws, row, text, span=8):
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=span)
        c = ws.cell(row, 1, "Fonte: " + text); c.font = Font(name="Calibri", size=8, italic=True, color=GRAY)
        c.alignment = Alignment(horizontal="left", vertical="center", indent=1)

    def table(ws, df, r0, numfmt="#,##0.0"):
        for j, col in enumerate(df.columns, 1):
            c = ws.cell(r0, j, str(col)); c.fill = navy; c.font = hf; c.border = bd
            c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for i, (_, row) in enumerate(df.iterrows(), r0 + 1):
            for j, col in enumerate(df.columns, 1):
                v = row[col]
                cell = ws.cell(i, j, None if (isinstance(v, float) and pd.isna(v)) else v); cell.border = bd
                if isinstance(v, (int, float, np.integer, np.floating)) and not pd.isna(v):
                    cell.number_format = "#,##0" if float(v).is_integer() else numfmt
                    cell.alignment = Alignment(horizontal="center")
        return r0 + len(df) + 1

    def barchart(ws, r0, ncol_val, nrow, cats_col, anchor, color, horiz=False, fmt="#,##0"):
        ch = BarChart(); ch.type = "bar" if horiz else "col"; ch.style = 10
        ch.height = 9.5; ch.width = 20; ch.legend = None
        data = Reference(ws, min_col=ncol_val, min_row=r0, max_row=r0 + nrow)
        cats = Reference(ws, min_col=cats_col, min_row=r0 + 1, max_row=r0 + nrow)
        ch.add_data(data, titles_from_data=True); ch.set_categories(cats)
        ch.series[0].graphicalProperties.solidFill = color
        ch.dLbls = DataLabelList(); ch.dLbls.showVal = True; ch.dLbls.numFmt = fmt
        ws.add_chart(ch, anchor)

    # 1. Evolução PL
    ws = wb.active; ws.title = "Evolucao PL"; banner(ws, f"Evolução do PL da indústria de FIDCs (ex-FIC) — R$ bi · {comp}")
    pl = pl_industria_anual(dm, comp)[["ano", "pl_liquido"]].rename(columns={"ano": "Ano", "pl_liquido": "PL líquido (R$ bi)"})
    end = table(ws, pl, 3, "#,##0")
    barchart(ws, 3, 2, len(pl), 1, "D3", ORANGE, horiz=False)
    note(ws, end + 12, FONTES["pl"]); ws.column_dimensions["A"].width = 10; ws.column_dimensions["B"].width = 16

    # 2. Composição por segmento
    ws = wb.create_sheet("Composicao"); banner(ws, "Composição do PL por segmento (%) — evolução anual")
    comp_df = composicao_anual(dm).reset_index().rename(columns={"ano": "Ano"})
    table(ws, comp_df, 3, "#,##0.0")
    from openpyxl.chart import BarChart as BC, Reference as Rf
    ch = BC(); ch.type = "col"; ch.grouping = "percentStacked"; ch.overlap = 100; ch.style = 10
    ch.height = 9.5; ch.width = 22
    data = Rf(ws, min_col=2, max_col=1 + (comp_df.shape[1] - 1), min_row=3, max_row=3 + len(comp_df))
    cats = Rf(ws, min_col=1, min_row=4, max_row=3 + len(comp_df))
    ch.add_data(data, titles_from_data=True); ch.set_categories(cats)
    for k, s in enumerate(ch.series):
        s.graphicalProperties.solidFill = STACK[k % len(STACK)]
    ws.add_chart(ch, "A" + str(6 + len(comp_df)))
    note(ws, 20 + len(comp_df), FONTES["composicao"])

    # 3. Definições ANBIMA
    ws = wb.create_sheet("Definicoes ANBIMA"); banner(ws, "Classes ANBIMA de FIDC — definições oficiais", 2)
    dfd = pd.DataFrame(DEFINICOES_ANBIMA, columns=["Classe", "Definição (ANBIMA — Deliberação nº 72)"])
    table(ws, dfd, 3); ws.column_dimensions["A"].width = 26; ws.column_dimensions["B"].width = 120
    for i in range(len(dfd)):
        ws.cell(4 + i, 2).alignment = Alignment(wrap_text=True, vertical="top")
        ws.row_dimensions[4 + i].height = 60
    note(ws, 5 + len(dfd), FONTES["definicoes"])

    # 4. Ranking gestor
    ws = wb.create_sheet("Ranking Gestor"); banner(ws, f"Ranking de Gestores (consolidado) — {comp} · Top 12 · R$ bi")
    rg = ranking_consolidado(dm, "gestor", n2g, comp, 12).rename(columns={"grupo": "Grupo", "pl_bi": "PL (R$ bi)"})
    table(ws, rg, 3); barchart(ws, 3, 3, len(rg), 2, "F3", NAVY, horiz=True, fmt="#,##0.0")
    note(ws, 6 + len(rg), FONTES["gestor"]); ws.column_dimensions["B"].width = 30

    # 5. Deltas por papel (gestor/admin/custódia) com share e ranking
    ws = wb.create_sheet("Delta por Papel"); banner(ws, "Share e mudança de posição por papel (dez/23 → atual)", 6)
    r0 = 3
    for dim, lab in [("gestor", "GESTOR"), ("admin", "ADMINISTRADOR"), ("custodiante", "CUSTODIANTE")]:
        ws.cell(r0, 1, lab).font = Font(name="Calibri", bold=True, color=NAVY)
        dl = deltas_por_papel(dm, dim, n2g, delta_per, top=10)
        r0 = table(ws, dl, r0 + 1, "#,##0.0") + 1
    note(ws, r0, FONTES["delta_papel"]); ws.column_dimensions["A"].width = 28

    # 6. Ranking admin
    ws = wb.create_sheet("Ranking Admin"); banner(ws, f"Ranking de Administradores (consolidado) — {comp} · Top 12 · R$ bi")
    ra = ranking_consolidado(dm, "admin", n2g, comp, 12).rename(columns={"grupo": "Grupo", "pl_bi": "PL (R$ bi)"})
    table(ws, ra, 3); barchart(ws, 3, 3, len(ra), 2, "F3", ORANGE, horiz=True, fmt="#,##0.0")
    note(ws, 6 + len(ra), FONTES["admin"]); ws.column_dimensions["B"].width = 30

    # 7. Originação de FIDCs novos + captação líquida
    ws = wb.create_sheet("Originacao Novos"); banner(ws, "Originação — FIDCs novos por ano de 1ª oferta (R$ bi)", 4)
    orig, ometa = originacao_novos(snap, n2g, (2025, 2026))
    r0 = 3
    for y in (2025, 2026):
        ws.cell(r0, 1, f"FIDCs nascidos em {y}: {ometa[y]['n']} fundos · R$ {ometa[y]['pl']:.0f} bi").font = Font(name="Calibri", bold=True, color=NAVY)
        od = orig[y].head(10).round(1).rename("PL originado (R$ bi)").reset_index().rename(columns={"g": "Gestor (consolidado)", "index": "Gestor (consolidado)"})
        od.columns = ["Gestor (consolidado)", "PL originado (R$ bi)"]
        r0 = table(ws, od, r0 + 1, "#,##0.0") + 1
    note(ws, r0, FONTES["originacao"]); ws.column_dimensions["A"].width = 40

    ws = wb.create_sheet("Captacao Liquida"); banner(ws, "Captação líquida anual (captações − resgates) — R$ bi")
    cap = captacao_liquida_anual(dm).rename(columns={"ano": "Ano", "capt_liq": "Captação líquida (R$ bi)"})
    table(ws, cap, 3, "#,##0"); barchart(ws, 3, 2, len(cap), 1, "D3", ORANGE, horiz=False)
    note(ws, 6 + len(cap), FONTES["captacao"]); ws.column_dimensions["A"].width = 8

    # 8. Tipo de controle
    ws = wb.create_sheet("Tipo de Controle"); banner(ws, "PL por tipo de controle do gestor — evolução (R$ bi)")
    tc = tipo_controle_anual(dm, n2g, periodos).reset_index().rename(columns={"tc": "Tipo de controle"})
    table(ws, tc, 3); note(ws, 5 + len(tc), FONTES["controle"]); ws.column_dimensions["A"].width = 22

    # 9. De-para
    ws = wb.create_sheet("De-para Conglomerados"); banner(ws, "De-para de conglomerados (auditável) — chave: CNPJ", 4)
    dep = depara_table()
    table(ws, dep, 3); note(ws, 5 + len(dep), FONTES["depara"])
    ws.column_dimensions["A"].width = 20; ws.column_dimensions["B"].width = 20
    ws.column_dimensions["C"].width = 18; ws.column_dimensions["D"].width = 46

    # 10. Top 25 — detalhe do dossiê (IME/CVM) se disponível, senão o resumo
    ws = wb.create_sheet("Top 25 FIDCs"); banner(ws, f"Top 25 FIDCs por PL — {comp} (anatomia documental)", 12)
    dossie_csv = Path("reports") / "top25_fidc_dossie.csv"
    if dossie_csv.exists():
        dd = pd.read_csv(dossie_csv)
        cols = ["#", "fundo", "pl_bi", "gestor", "segmento", "n_classes_series", "subordinacao_%",
                "n_cedentes_identificados", "maior_cedente_%", "maior_cedente", "pct_vencido_inad",
                "pct_judicial", "sacado_perform_%", "n_cotistas", "top_investidores", "n_cotistas_institucionais",
                "ano_1a_oferta"]
        cols = [c for c in cols if c in dd.columns]
        table(ws, dd[cols], 3, "#,##0.0")
        for j, w in enumerate([4, 46, 8, 18, 16, 8, 9, 8, 9, 16, 9, 8, 9, 8, 34, 8, 9]):
            ws.column_dimensions[chr(65 + j) if j < 26 else "A" + chr(65 + j - 26)].width = w
        note(ws, 5 + len(dd), FONTES["top25det"])
    else:
        t25 = top25_detalhado(snap, n2g, 25)
        table(ws, t25, 3, "#,##0.00")
        note(ws, 5 + len(t25), FONTES["top25"])

    # 11. Cotistas
    ws = wb.create_sheet("Cotistas"); banner(ws, "Número de cotistas — fundos > R$ 200 mi de PL")
    cot, meta = cotistas_dist(snap)
    ws.cell(2, 1, f"Universo: {meta['n']} fundos / R$ {meta['pl']:.0f} bi de PL")
    cotd = cot.rename(columns={"b": "Faixa", "n_fundos": "# Fundos", "pl_bi": "PL (R$ bi)"})
    table(ws, cotd, 4, "#,##0"); barchart(ws, 4, 2, len(cotd), 1, "F4", ORANGE, horiz=False)
    note(ws, 8 + len(cotd), FONTES["cotistas"]); ws.column_dimensions["A"].width = 16

    # 12. Segmento investidor
    ws = wb.create_sheet("Segmento Investidor"); inv, ci = cot_investidor(D["cot_tipo"])
    banner(ws, f"Nº de cotistas por segmento de investidor — {ci}")
    invd = inv.rename(columns={"tipo_cotista": "Segmento", "n": "# Cotistas"})
    table(ws, invd, 3, "#,##0"); barchart(ws, 3, 2, len(invd), 1, "D3", NAVY, horiz=True)
    note(ws, 6 + len(invd), FONTES["investidor"]); ws.column_dimensions["A"].width = 30

    # 13. Emissões
    ws = wb.create_sheet("Emissoes"); banner(ws, "Emissões (variação anual de PL, ex-FIC) — R$ bi")
    em = emissoes_anual(dm).rename(columns={"ano": "Ano", "var_pl": "Variação de PL (R$ bi)", "capt_liq": "Captação líquida (R$ bi)"})
    table(ws, em, 3, "#,##0"); barchart(ws, 3, 2, len(em), 1, "F3", ORANGE, horiz=False)
    note(ws, 6 + len(em), FONTES["emissoes"]); ws.column_dimensions["A"].width = 8

    # 13b. Cedentes & Sacados
    ced, sac = cedentes_sacados(Path("reports"))
    if ced is not None:
        ws = wb.create_sheet("Cedentes e Sacados"); banner(ws, "Cedentes e sacados relevantes (leitura de regulamentos)", 4)
        ws.cell(3, 1, "Maiores CEDENTES (originadores)").font = Font(name="Calibri", bold=True, color=NAVY)
        e1 = table(ws, ced, 4)
        ws.cell(e1 + 1, 1, "Maiores SACADOS (devedores)").font = Font(name="Calibri", bold=True, color=NAVY)
        e2 = table(ws, sac, e1 + 2)
        note(ws, e2 + 1, FONTES["partes"])
        ws.column_dimensions["A"].width = 40; ws.column_dimensions["B"].width = 30

    # 14. Fontes
    ws = wb.create_sheet("Fontes"); banner(ws, "Fonte exata por gráfico", 2)
    fdf = pd.DataFrame([(k, v) for k, v in FONTES.items()], columns=["Gráfico", "Fonte"])
    table(ws, fdf, 3); ws.column_dimensions["A"].width = 16; ws.column_dimensions["B"].width = 130
    for i in range(len(fdf)):
        ws.cell(4 + i, 2).alignment = Alignment(wrap_text=True, vertical="top"); ws.row_dimensions[4 + i].height = 42

    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(path)
    print(f"[ok] Excel: {path}")


# ---------------------------------------------------------------------------
# POWERPOINT
# ---------------------------------------------------------------------------
def build_pptx(D, path: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

    comp = D["comp"]; dm = D["dm"]; snap = D["snap"]; n2g = D["name2grp"]
    periodos = PERIODOS + [comp]
    delta_per = [x for x in PERIODOS if x >= "2023"] + [comp]
    rgb = lambda h: RGBColor.from_string(h)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height

    def rect(s, x, y, w, h, color):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(color); sp.line.fill.background(); sp.shadow.inherit = False
        return sp

    def text(s, x, y, w, h, t, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        first = True
        for line in str(t).split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.alignment = align; r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color); r.font.name = "Calibri"
        return tb

    def header(s, kicker, title):
        rect(s, 0, 0, SW, Inches(1.12), NAVY); rect(s, 0, Inches(1.12), SW, Emu(45720), ORANGE)
        if kicker:
            text(s, Inches(0.5), Inches(0.13), Inches(12), Inches(0.3), kicker.upper(), 11, True, ORANGE2)
        text(s, Inches(0.5), Inches(0.4), Inches(12.4), Inches(0.66), title, 22, True, WHITE)

    def src(s, key):
        text(s, Inches(0.5), SH - Inches(0.5), Inches(12.4), Inches(0.45), "Fonte: " + FONTES[key], 8, False, GRAY)

    def style(chart, colors, fmt="#,##0", legend=False, stacked=False, lblsize=10):
        chart.has_title = False
        try:
            plot = chart.plots[0]; plot.has_data_labels = True
            dl = plot.data_labels; dl.number_format = fmt; dl.number_format_is_linked = False
            dl.font.size = Pt(lblsize); dl.font.name = "Calibri"
            if not stacked:
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass
        chart.has_legend = legend
        if legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM; chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
        if stacked:
            for k, s__ in enumerate(chart.series):
                s__.format.fill.solid(); s__.format.fill.fore_color.rgb = rgb(colors[k % len(colors)])
        else:
            ser = chart.series[0]
            try:
                for i, pt in enumerate(ser.points):
                    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = rgb(colors[i % len(colors)])
            except Exception:
                ser.format.fill.solid(); ser.format.fill.fore_color.rgb = rgb(colors[0])
        for ax in ("category_axis", "value_axis"):
            try:
                a = getattr(chart, ax); a.tick_labels.font.size = Pt(9); a.tick_labels.font.name = "Calibri"
            except Exception:
                pass

    def table(s, df, x, y, w, h, widths=None, fs=9, headfs=None):
        rows, cols = df.shape[0] + 1, df.shape[1]
        t = s.shapes.add_table(rows, cols, x, y, w, h).table
        for j, col in enumerate(df.columns):
            c = t.cell(0, j); c.text = str(col); c.fill.solid(); c.fill.fore_color.rgb = rgb(NAVY)
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.runs[0]; r.font.size = Pt(headfs or fs); r.font.bold = True; r.font.color.rgb = rgb(WHITE); r.font.name = "Calibri"
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = c.margin_bottom = Emu(9144)
        for i in range(df.shape[0]):
            for j, col in enumerate(df.columns):
                v = df.iloc[i, j]
                if isinstance(v, float):
                    txt = "" if pd.isna(v) else (f"{v:,.2f}" if abs(v) < 100 else f"{v:,.0f}")
                elif isinstance(v, (int, np.integer)):
                    txt = f"{int(v):+d}" if str(col).lower().startswith("δ rank") else f"{int(v):,}"
                else:
                    txt = "" if v is None else str(v)
                c = t.cell(i + 1, j); c.text = txt
                c.fill.solid(); c.fill.fore_color.rgb = rgb("F3F5F8" if i % 2 == 0 else WHITE)
                p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if isinstance(v, str) else PP_ALIGN.CENTER
                r = p.runs[0] if p.runs else p.add_run(); r.font.size = Pt(fs); r.font.name = "Calibri"; r.font.color.rgb = rgb(INK)
                c.vertical_anchor = MSO_ANCHOR.MIDDLE; c.margin_top = c.margin_bottom = Emu(4572)
        if widths:
            tot = sum(widths)
            for j, cw in enumerate(widths):
                t.columns[j].width = int(w * cw / tot)
        return t

    def chart(s, kind, x, y, w, h, cats, series, colors, fmt="#,##0", stacked=False, legend=False, lblsize=10):
        cd = CategoryChartData(); cd.categories = cats
        for nm, vals in series:
            cd.add_series(nm, vals)
        gf = s.shapes.add_chart(kind, x, y, w, h, cd)
        style(gf.chart, colors, fmt, legend=legend, stacked=stacked, lblsize=lblsize)
        return gf.chart

    # 0 Capa
    s = prs.slides.add_slide(blank); rect(s, 0, 0, SW, SH, NAVY); rect(s, 0, Inches(3.05), SW, Inches(0.08), ORANGE)
    text(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.0), "Indústria de FIDCs", 44, True, WHITE)
    text(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.7), "Evolução da indústria, prestadores de serviço e concentração", 20, False, ORANGE2)
    text(s, Inches(0.8), Inches(4.15), Inches(11.7), Inches(0.5), f"Competência {comp}  ·  consolidação por conglomerado", 15, True, WHITE)
    text(s, Inches(0.8), SH - Inches(0.9), Inches(11.7), Inches(0.5),
         "Toma Conta FIDCs (Informe Mensal FIDC/CVM + Cadastro). Classes: ANBIMA (Delib. 72). Sem Power BI.", 11, False, GRAYL)

    # 0b Sumário executivo
    s = prs.slides.add_slide(blank); header(s, "Leitura de mercado", "Sumário executivo — o que importa para o Comitê")
    sm = sumario_executivo(dm, snap, n2g, comp)
    cards = [
        ("O mercado dobrou", f"R$ {sm['pl_ini']:,.0f} (2023) → {sm['pl_fim']:,.0f} bi",
         "Indústria de FIDC dobrou em ~2,5 anos (ex-FIC), puxada por crédito estruturado."),
        ("Independentes dominam a gestão", f"{sm['indep_pct']:,.0f}% do PL (R$ {sm['indep_v']:,.0f} bi)",
         "A gestão migrou para fora dos bancos de varejo; independentes lideram."),
        ("A briga é administração/custódia", f"BTG {sm['adm_btg']:,.0f} · QI Tech {sm['adm_qi']:,.0f} · Itaú {sm['adm_itau']:,.0f} bi",
         "Itaú administra ~4–5x menos que BTG e QI Tech — o gap a endereçar."),
        ("Crédito colado a pagamentos", "Maior sacado: Mercado Pago (R$ 5,7 bi)",
         "Grandes sacados são plataformas de pagamento/varejo (Mercado Pago, SafraPay, Casas Bahia)."),
    ]
    y = 1.5
    for titulo, num, desc in cards:
        rect(s, Inches(0.5), Inches(y), Inches(0.14), Inches(1.15), ORANGE)
        text(s, Inches(0.8), Inches(y + 0.02), Inches(4.6), Inches(1.15), titulo, 15, True, NAVY, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(5.3), Inches(y + 0.02), Inches(3.1), Inches(1.15), num, 15, True, ORANGE, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(8.5), Inches(y + 0.02), Inches(4.5), Inches(1.15), desc, 11, False, INK, anchor=MSO_ANCHOR.MIDDLE)
        y += 1.28
    src(s, "sumario")

    # 1 Evolução PL
    s = prs.slides.add_slide(blank); header(s, "Crescimento da indústria", "Evolução do PL da indústria de FIDCs (base líquida, ex-FIC)")
    pl = pl_industria_anual(dm, comp)
    chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.45), Inches(8.7), Inches(5.2),
          [str(a) for a in pl["ano"]], [("PL líquido (R$ bi)", [float(v) for v in pl["pl_liquido"]])], [ORANGE])
    v0, v1 = float(pl["pl_liquido"].iloc[0]), float(pl["pl_liquido"].iloc[-1])
    g = (v1 / v0 - 1) * 100 if v0 else 0
    rect(s, Inches(9.5), Inches(1.9), Inches(3.4), Inches(2.1), "F3F5F8")
    text(s, Inches(9.7), Inches(2.05), Inches(3.0), Inches(0.4), "PL atual (ex-FIC)", 12, True, NAVY)
    text(s, Inches(9.7), Inches(2.45), Inches(3.0), Inches(0.7), f"R$ {v1:,.0f} bi", 30, True, ORANGE)
    text(s, Inches(9.7), Inches(3.25), Inches(3.0), Inches(0.6), f"{pl['ano'].iloc[0]}→{pl['ano'].iloc[-1]}: +{g:,.0f}%", 14, True, NAVY)
    text(s, Inches(9.5), Inches(4.2), Inches(3.5), Inches(2.2),
         "• Série líquida (exclui FIC-FIDC) para evitar dupla contagem\n• Ciclo puxado por avanços regulatórios, novos recebíveis e busca por funding", 11, False, INK)
    src(s, "pl")

    # 2 Composição por segmento
    s = prs.slides.add_slide(blank); header(s, "Composição", "Composição do PL por segmento de atuação (%)")
    cdf = composicao_anual(dm)
    cols = list(cdf.columns)[:6]
    series = [(c, [float(cdf.loc[a, c]) for a in cdf.index]) for c in cols]
    chart(s, XL_CHART_TYPE.COLUMN_STACKED_100, Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.1),
          list(cdf.index), series, STACK, fmt="0", stacked=True, legend=True, lblsize=8)
    src(s, "composicao")

    # 3 Definições ANBIMA
    s = prs.slides.add_slide(blank); header(s, "Referencial ANBIMA", "Classes ANBIMA de FIDC — definições oficiais")
    y = 1.5
    for classe, desc in DEFINICOES_ANBIMA:
        rect(s, Inches(0.5), Inches(y), Inches(3.1), Inches(1.02), NAVY)
        text(s, Inches(0.62), Inches(y + 0.06), Inches(2.9), Inches(0.9), classe, 13, True, WHITE, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(3.8), Inches(y - 0.02), Inches(9.0), Inches(1.05), desc, 11, False, INK, anchor=MSO_ANCHOR.MIDDLE)
        y += 1.08
    src(s, "definicoes")

    # 4 Ranking gestor
    s = prs.slides.add_slide(blank); header(s, "Prestadores — gestão", f"Ranking de Gestores (consolidado por conglomerado) — {comp}")
    rg = ranking_consolidado(dm, "gestor", n2g, comp, 12)
    labels = [("Itaú (Itaú Asset+Unibanco+Kinea)" if g == "Itau" else g)[:32] for g in rg["grupo"]][::-1]
    chart(s, XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.45), Inches(12.3), Inches(4.7),
          labels, [("PL (R$ bi)", [float(v) for v in rg["pl_bi"]][::-1])], [NAVY], fmt="#,##0.0")
    text(s, Inches(0.5), Inches(6.35), Inches(12.4), Inches(0.6), ITAU_NOTA, 8.5, False, GRAY)
    src(s, "gestor")

    # 5-7 Deltas por papel (gestor / administrador / custodiante) com share e ranking
    role_meta = [("gestor", "Gestores", "gestão"), ("admin", "Administradores", "administração"),
                 ("custodiante", "Custodiantes", "custódia")]
    for dim, titulo, papel in role_meta:
        s = prs.slides.add_slide(blank)
        header(s, "Movimento de mercado — " + papel, f"{titulo}: share e mudança de posição (dez/23 → {comp})")
        dl = deltas_por_papel(dm, dim, n2g, delta_per, top=9)
        table(s, dl, Inches(0.35), Inches(1.5), Inches(12.6), Inches(4.9),
              widths=[2.6, 1.0, 0.95, 0.95, 0.95, 0.95, 0.95, 0.85], fs=10, headfs=9)
        # destaque textual
        ganhou = dl.sort_values("Δ share", ascending=False).iloc[0]
        perdeu = dl.sort_values("Δ share").iloc[0]
        gsh = float(ganhou["Δ share"]); psh = float(perdeu["Δ share"])
        text(s, Inches(0.5), Inches(6.5), Inches(12.4), Inches(0.5),
             f"▲ Ganhou share: {ganhou['Grupo']} ({gsh:+.1f} pp)    "
             f"▼ Perdeu: {perdeu['Grupo']} ({psh:+.1f} pp)", 12, True, NAVY)
        src(s, "delta_papel")

    # 7b Originação — FIDCs novos 2025/2026
    s = prs.slides.add_slide(blank); header(s, "Originação (não rouba-monte)", "Quem origina FIDCs novos — por ano de 1ª oferta")
    orig, ometa = originacao_novos(snap, n2g, (2025, 2026))
    for col, y in ((0, 2025), (1, 2026)):
        x = Inches(0.5 + col * 6.6)
        text(s, x, Inches(1.3), Inches(6), Inches(0.4),
             f"FIDCs nascidos em {y}: {ometa[y]['n']} fundos · R$ {ometa[y]['pl']:.0f} bi", 13, True, NAVY)
        od = orig[y].head(7)
        cats = [g[:26] for g in od.index][::-1]; vals = [float(v) for v in od.values][::-1]
        chart(s, XL_CHART_TYPE.BAR_CLUSTERED, x, Inches(1.7), Inches(6.2), Inches(4.7),
              cats, [("PL originado (R$ bi)", vals)], [ORANGE if col == 0 else NAVY], fmt="#,##0.0", lblsize=9)
    src(s, "originacao")

    # 8 Tipo de controle
    s = prs.slides.add_slide(blank); header(s, "Independentes x bancos", "PL por tipo de controle do gestor — evolução (R$ bi)")
    tc = tipo_controle_anual(dm, n2g, periodos)
    series = [(idx, [float(tc.loc[idx, p]) for p in periodos]) for idx in tc.index]
    chart(s, XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(1.45), Inches(8.7), Inches(5.1),
          [p[:7] for p in periodos], series, [ORANGE, NAVY, GRAY], fmt="#,##0", stacked=True, legend=True, lblsize=9)
    cur = tc[periodos[-1]]; tot = cur.sum()
    yy = 1.9
    for idx in tc.index:
        text(s, Inches(9.5), Inches(yy), Inches(3.6), Inches(0.6),
             f"{idx}: R$ {cur[idx]:,.0f} bi ({cur[idx]/tot*100:,.0f}%)", 13, True,
             ORANGE if idx == "Independente" else (NAVY if idx == "Ligada a banco" else GRAY))
        yy += 0.7
    text(s, Inches(9.5), Inches(yy + 0.1), Inches(3.7), Inches(2.0),
         "Definição de controle conforme rodapé do deck Itaú BBA. 'Independente Grande' = Oliveira Trust, BRL Trust, BR Trust.", 10, False, GRAY)
    src(s, "controle")

    # 9 De-para
    s = prs.slides.add_slide(blank); header(s, "Metodologia", "De-para de conglomerados (auditável) — chave: CNPJ")
    dep = depara_table()
    dep_show = dep.groupby(["grupo", "tipo_controle"]).agg(
        entidades=("entidade_original", lambda x: " · ".join(sorted(set(x)))[:70]),
        n=("cnpj", "count")).reset_index().rename(columns={"grupo": "Grupo", "tipo_controle": "Controle", "entidades": "Entidades (agrupadas)", "n": "#"})
    dep_show = dep_show.sort_values(["Controle", "Grupo"]).reset_index(drop=True)
    table(s, dep_show, Inches(0.35), Inches(1.4), Inches(12.6), Inches(5.4),
          widths=[1.5, 1.4, 5.5, 0.5], fs=8.5, headfs=9)
    src(s, "depara")

    # 10 Top 25
    s = prs.slides.add_slide(blank); header(s, "Maiores veículos", f"Top 25 FIDCs por PL — {comp} (detalhe)")
    t25 = top25_detalhado(snap, n2g, 25)
    t25["Fundo"] = t25["Fundo"].str.slice(0, 40); t25["Gestor"] = t25["Gestor"].str.slice(0, 22)
    t25["Tipo de recebível"] = t25["Tipo de recebível"].str.slice(0, 30); t25["Cedente princ."] = t25["Cedente princ."].str.slice(0, 20)
    half = 13
    cols_keep = ["#", "Fundo", "Tipo de recebível", "Gestor", "PL (R$ bi)", "Sub. mín."]
    w = [0.35, 2.7, 1.9, 1.5, 0.8, 0.7]
    table(s, t25[cols_keep].iloc[:half].reset_index(drop=True), Inches(0.25), Inches(1.45), Inches(6.45), Inches(5.5), widths=w, fs=7.5, headfs=7.5)
    table(s, t25[cols_keep].iloc[half:].reset_index(drop=True), Inches(6.85), Inches(1.45), Inches(6.45), Inches(5.5), widths=w, fs=7.5, headfs=7.5)
    src(s, "top25")

    # 10b Cedentes & Sacados
    ced, sac = cedentes_sacados(Path("reports"))
    if ced is not None:
        s = prs.slides.add_slide(blank); header(s, "Originação e risco", "Cedentes e sacados relevantes (leitura de regulamentos)")
        text(s, Inches(0.5), Inches(1.25), Inches(6), Inches(0.3), "Maiores CEDENTES (originadores)", 13, True, NAVY)
        text(s, Inches(6.9), Inches(1.25), Inches(6), Inches(0.3), "Maiores SACADOS (devedores)", 13, True, NAVY)
        cw = [3.2, 1.9, 0.7, 1.2]
        cc = ced.copy(); cc["Parte"] = cc["Parte"].str.slice(0, 30); cc["Setor"] = cc["Setor"].str.slice(0, 18)
        ss = sac.copy(); ss["Parte"] = ss["Parte"].str.slice(0, 30); ss["Setor"] = ss["Setor"].str.slice(0, 18)
        table(s, cc, Inches(0.4), Inches(1.6), Inches(6.3), Inches(4.6), widths=cw, fs=8.5, headfs=8.5)
        table(s, ss, Inches(6.85), Inches(1.6), Inches(6.3), Inches(4.6), widths=cw, fs=8.5, headfs=8.5)
        text(s, Inches(0.5), Inches(6.35), Inches(12.4), Inches(0.4),
             "Sacados são majoritariamente plataformas de pagamento/varejo — o crédito está colado a meios de pagamento.", 11, True, ORANGE)
        src(s, "partes")

    # 11 Cotistas
    s = prs.slides.add_slide(blank); header(s, "Estrutura de cotistas", "Número de cotistas — fundos > R$ 200 mi de PL")
    cot, meta = cotistas_dist(snap)
    text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
         f"{meta['n']} fundos · R$ {meta['pl']:.0f} bi de PL · estruturas mono/poucos cotistas são veículos de crédito institucionais", 12, True, NAVY)
    chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.1), Inches(2.0), Inches(7.4), Inches(4.7),
          list(cot["b"]), [("# Fundos", [int(v) for v in cot["n_fundos"]])], [ORANGE, ORANGE2, GRAY, NAVY])
    cotd = cot.rename(columns={"b": "Faixa", "n_fundos": "# Fundos", "pl_bi": "PL (R$ bi)"})
    table(s, cotd, Inches(8.8), Inches(2.4), Inches(4.1), Inches(2.5), widths=[1.6, 1.0, 1.2], fs=11)
    src(s, "cotistas")

    # 12 Segmento investidor
    s = prs.slides.add_slide(blank); inv, ci = cot_investidor(D["cot_tipo"])
    header(s, "Base investidora", f"Nº de cotistas por segmento de investidor — {ci}")
    chart(s, XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.1),
          [n[:30] for n in inv["tipo_cotista"]][::-1], [("# Cotistas", [float(v) for v in inv["n"]][::-1])], [NAVY], fmt="#,##0")
    src(s, "investidor")

    # 13 Emissões
    s = prs.slides.add_slide(blank); header(s, "Fluxo primário", "Emissões — variação anual de PL (ex-FIC), R$ bi")
    em = emissoes_anual(dm)
    chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.5), Inches(8.6), Inches(5.0),
          [str(a) for a in em["ano"]], [("Variação de PL (R$ bi)", [float(v) for v in em["var_pl"]])], [ORANGE])
    text(s, Inches(9.4), Inches(1.9), Inches(3.6), Inches(4.5),
         "Métrica: variação anual de PL (ex-FIC), como na aba Indústria do Toma Conta.\n\n"
         "Não é o campo 'emissões encerradas' (status ANBIMA), que consideramos menos crível.\n\n"
         "Anos completos (jan–dez). 2026 é parcial (até " + comp[-2:] + "/" + comp[:4] + ").\n\n"
         "Validação: ANBIMA reportou FIDCs entre as categorias que mais cresceram em 2024, "
         "consistente com +R$ " + f"{float(em.loc[em['ano']==2024,'var_pl'].iloc[0]):,.0f}" + " bi.", 11, False, INK)
    src(s, "emissoes")

    # 13b Captação líquida
    s = prs.slides.add_slide(blank); header(s, "Fluxo primário", "Captação líquida anual (captações − resgates) — R$ bi")
    cap = captacao_liquida_anual(dm)
    chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.5), Inches(8.6), Inches(5.0),
          [str(a) for a in cap["ano"]], [("Captação líquida (R$ bi)", [float(v) for v in cap["capt_liq"]])], [ORANGE])
    c24 = float(cap.loc[cap["ano"] == 2024, "capt_liq"].iloc[0]) if (cap["ano"] == 2024).any() else 0
    c25 = float(cap.loc[cap["ano"] == 2025, "capt_liq"].iloc[0]) if (cap["ano"] == 2025).any() else 0
    text(s, Inches(9.4), Inches(1.9), Inches(3.6), Inches(4.6),
         f"Captação líquida = captações − resgates no mês (IME/CVM), somada por ano.\n\n"
         f"2024: R$ {c24:,.0f} bi · 2025: R$ {c25:,.0f} bi.\n\n"
         "Complementa a variação de PL (que também capta valorização de cota). "
         "ANBIMA divulga FIDCs entre os líderes de captação em 2024–2025.", 11, False, INK)
    src(s, "captacao")

    # 14 Fontes
    s = prs.slides.add_slide(blank); header(s, "Transparência", "Fonte exata por gráfico")
    fdf = pd.DataFrame([(k.upper(), v) for k, v in FONTES.items()], columns=["Gráfico", "Fonte"])
    table(s, fdf, Inches(0.4), Inches(1.4), Inches(12.6), Inches(5.5), widths=[1.4, 7.0], fs=8)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    print(f"[ok] PPTX: {path}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--industry-dir", type=Path, default=Path("data/industry_study"))
    ap.add_argument("--out", type=Path, default=Path("outputs"))
    args = ap.parse_args()
    D = load(args.industry_dir)
    comp = D["comp"]
    build_excel(D, args.out / f"Industria_FIDC_{comp}.xlsx")
    build_pptx(D, args.out / f"Industria_FIDC_{comp}.pptx")
    print(f"[ok] competência {comp}")


if __name__ == "__main__":
    main()
