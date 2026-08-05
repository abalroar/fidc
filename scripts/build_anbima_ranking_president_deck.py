"""Standalone deck: ANBIMA fixed-income origination ranking, Itaú BBA position.

Reads only the artefacts materialized by
``scripts/build_anbima_fixed_income_ranking.py`` and renders a short, native
and fully editable 16:9 PowerPoint for the Itaú BBA president.

Scope is the consolidated fixed-income ranking (ANBIMA Tipo 1), accumulated in
the year through the cut-off date.  It deliberately does not touch the industry
deck contract in ``services/industry_ppt_export.py``.

    python scripts/build_anbima_ranking_president_deck.py
"""

from __future__ import annotations

import argparse
from datetime import date
import json
from pathlib import Path
import sys

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

DEFAULT_DATA_DIR = Path("data/industry_study")
DEFAULT_OUTPUT = Path(
    "outputs/anbima_ranking_rf_presidente/"
    "ANBIMA_Ranking_Renda_Fixa_Itau_BBA_1S26.pptx"
)

# Same palette and typography as services/industry_ppt_export.py so this deck
# sits alongside the existing BBA decks without a visual break.
BLACK = "151515"
ORANGE = "E36C0A"
ORANGE_LIGHT = "F8E9DE"
GRAY_900 = "30353A"
GRAY_700 = "5D6369"
GRAY_500 = "8D9399"
GRAY_300 = "D7DADD"
GRAY_200 = "E7E9EB"
GRAY_100 = "F5F6F7"
WHITE = "FFFFFF"

FONT = "Arial"
HOUSE = "ITAU BBA"
KICKER = "RANKING ANBIMA · RENDA FIXA"

SOURCE_LINE = (
    "Fonte: ANBIMA, Ranking de Renda Fixa e Híbridos — Originação (Valor), "
    "Tipo 1: Renda Fixa Consolidado, acumulado 2026 · referência Junho/2026"
)

#: ANBIMA publishes participant names in upper case.  Acronyms must survive the
#: conversion to display case, so a plain ``str.title()`` is not usable.
DISPLAY_NAMES: dict[str, str] = {
    "BRADESCO BBI": "Bradesco BBI",
    "ITAU BBA": "Itaú BBA",
    "BTG PACTUAL": "BTG Pactual",
    "SANTANDER": "Santander",
    "XP INVESTIMENTOS": "XP Investimentos",
    "UBS BB": "UBS BB",
    "CEF": "Caixa Econômica Federal",
    "SAFRA": "Safra",
    "ABC BRASIL": "ABC Brasil",
    "VOTORANTIM": "Votorantim",
    "CITIGROUP": "Citigroup",
    "BNDES": "BNDES",
    "BR PARTNERS": "BR Partners",
    "BB-BI": "BB-BI",
    "DAYCOVAL": "Daycoval",
    "BNP PARIBAS": "BNP Paribas",
    "BOCOM BBM": "Bocom BBM",
    "INTER": "Inter",
    "M7 IB": "M7 IB",
    "JP MORGAN": "J.P. Morgan",
}

#: Published sub-rankings that decompose the consolidated fixed-income number.
SUBDIVISIONS: tuple[tuple[str, str], ...] = (
    ("1.1", "Renda fixa — curto prazo"),
    ("1.2", "Renda fixa — longo prazo"),
    ("1.3", "Securitização"),
    ("1.3.1", "Securitização · FIDC"),
    ("1.3.2", "Securitização · CRI"),
    ("1.3.3", "Securitização · CRA"),
)

METHODOLOGY: tuple[tuple[str, str], ...] = (
    (
        "O ranking credita todos os coordenadores, não só o líder",
        "Cada coordenador e coordenador contratado recebe a fatia que lhe cabe "
        "na operação. O líder apenas reporta a operação à ANBIMA.",
    ),
    (
        "O rateio é contratual",
        "Garantia firme: proporção da garantia definida em contrato. Melhores "
        "esforços: proporção do fee de coordenação e/ou estruturação.",
    ),
    (
        "O mês de referência é o do anúncio de encerramento",
        "Não é a data de registro nem a de emissão. Operações fora do prazo de "
        "envio escorregam para o ranking do trimestre seguinte.",
    ),
    (
        "É um ranking declaratório",
        "Operação cujo formulário-padrão não foi enviado não entra. Boa parte "
        "do mercado liderado por administradores e DTVMs fica de fora.",
    ),
    (
        "Operações de empresas ligadas saem do Tipo 1",
        "Coordenador com participação de 10% ou mais na emissora, cedente ou "
        "originadora vai para o Tipo 3, apurado à parte.",
    ),
    (
        "Perímetro do Tipo 1",
        "Debêntures simples, notas promissórias, notas comerciais, CPR-F e "
        "securitização (FIDC, CRI, CRA, CR). FII, FIAGRO e FIP-IE ficam nas "
        "operações híbridas e não entram.",
    ),
)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--data-dir", type=Path, default=DEFAULT_DATA_DIR)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument(
        "--house",
        default=HOUSE,
        help="Participante em destaque no ranking (default: ITAU BBA).",
    )
    return parser.parse_args()


def _bi(value: float) -> str:
    return f"{value / 1e9:,.1f}".replace(",", "@").replace(".", ",").replace("@", ".")


def _pct(value: float, decimals: int = 2) -> str:
    if pd.isna(value):
        return "—"
    return f"{value * 100:,.{decimals}f}%".replace(".", ",")


def _rank(value: object) -> str:
    return "—" if pd.isna(value) else f"{int(value)}º"


def _display_name(value: object) -> str:
    label = str(value).strip()
    return DISPLAY_NAMES.get(label.upper(), label)


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.page = 0

    @staticmethod
    def _rgb(value: str) -> RGBColor:
        return RGBColor.from_string(value)

    def text(
        self,
        slide,
        content: object,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float = 12,
        color: str = GRAY_900,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = valign
        frame.margin_left = frame.margin_right = Inches(0)
        frame.margin_top = frame.margin_bottom = Inches(0)
        paragraph = frame.paragraphs[0]
        paragraph.text = str(content)
        paragraph.alignment = align
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = self._rgb(color)
        return box

    def rule(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        *,
        color: str = GRAY_300,
        height: float = 0.012,
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(color)
        shape.line.fill.background()
        return shape

    def block(
        self, slide, x: float, y: float, w: float, h: float, color: str
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(color)
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def slide(self, title: str, kicker: str = KICKER):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._rgb(WHITE)
        self.text(slide, kicker, 0.62, 0.28, 8.0, 0.22, size=10, color=ORANGE, bold=True)
        self.text(slide, title, 0.62, 0.57, 12.1, 0.48, size=22, color=BLACK, bold=True)
        self.rule(slide, 0.62, 1.15, 12.05, color=GRAY_300, height=0.018)
        return slide

    def footer(self, slide, source: str) -> None:
        self.page += 1
        self.rule(slide, 0.62, 6.95, 12.05, color=GRAY_200, height=0.01)
        self.text(slide, source, 0.62, 7.02, 11.45, 0.3, size=8, color=GRAY_500)
        self.text(
            slide,
            str(self.page),
            12.25,
            7.01,
            0.42,
            0.22,
            size=8,
            color=GRAY_500,
            align=PP_ALIGN.RIGHT,
        )

    def table(
        self,
        slide,
        rows: list[list[str]],
        x: float,
        y: float,
        widths: list[float],
        *,
        highlight: int | None = None,
        row_height: float = 0.34,
        size: float = 11,
        align_right_from: int = 2,
    ) -> float:
        """Render a light, borderless table; returns the bottom y coordinate."""

        header, *body = rows
        cursor = y
        for index, label in enumerate(header):
            self.text(
                slide,
                label,
                x + sum(widths[:index]),
                cursor,
                widths[index],
                row_height,
                size=size - 1,
                color=GRAY_500,
                bold=True,
                align=PP_ALIGN.RIGHT if index >= align_right_from else PP_ALIGN.LEFT,
            )
        cursor += row_height * 0.85
        self.rule(slide, x, cursor, sum(widths), color=GRAY_300, height=0.012)
        cursor += 0.10

        for position, row in enumerate(body):
            is_house = highlight is not None and position == highlight
            if is_house:
                self.block(
                    slide,
                    x - 0.12,
                    cursor - 0.05,
                    sum(widths) + 0.24,
                    row_height,
                    ORANGE_LIGHT,
                )
            for index, value in enumerate(row):
                self.text(
                    slide,
                    value,
                    x + sum(widths[:index]),
                    cursor,
                    widths[index],
                    row_height,
                    size=size,
                    color=BLACK if is_house else GRAY_900,
                    bold=is_house,
                    align=(
                        PP_ALIGN.RIGHT if index >= align_right_from else PP_ALIGN.LEFT
                    ),
                )
            cursor += row_height
            if position < len(body) - 1:
                self.rule(slide, x, cursor - 0.05, sum(widths), color=GRAY_200, height=0.008)
        return cursor


def _load(data_dir: Path) -> tuple[pd.DataFrame, dict]:
    official = pd.read_csv(data_dir / "anbima_rf_ranking_official.csv")
    manifest = json.loads(
        (data_dir / "anbima_rf_ranking_manifest.json").read_text(encoding="utf-8")
    )
    return official, manifest


def _accumulated(official: pd.DataFrame, code: str) -> pd.DataFrame:
    return official[
        official["measure"].eq("originacao_valor")
        & official["window"].eq("acumulado_ano")
        & official["ranking_code"].eq(code)
    ]


def build(data_dir: Path, output: Path, house: str) -> Path:
    official, manifest = _load(data_dir)
    consolidated = _accumulated(official, "1")
    if consolidated.empty:
        raise SystemExit(
            "anbima_rf_ranking_official.csv sem o bloco Tipo 1; rode "
            "scripts/build_anbima_fixed_income_ranking.py antes."
        )

    universe = float(consolidated["value_brl_or_count"].sum())
    ranked = consolidated.sort_values("rank")
    participants = int(consolidated["participant"].nunique())
    house_row = ranked[ranked["participant"].eq(house)]
    if house_row.empty:
        raise SystemExit(f"Participante {house!r} ausente no ranking consolidado.")
    house_row = house_row.iloc[0]
    cutoff = manifest.get("period", {}).get("end", "")
    cutoff_label = (
        date.fromisoformat(cutoff).strftime("%d/%m/%Y") if cutoff else "n/d"
    )

    deck = Deck()

    # ---------------------------------------------------------------- capa
    cover = deck.prs.slides.add_slide(deck.prs.slide_layouts[6])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = deck._rgb(WHITE)
    deck.block(cover, 0.0, 0.0, 0.22, 7.5, ORANGE)
    deck.text(cover, KICKER, 0.92, 2.28, 8.0, 0.26, size=11, color=ORANGE, bold=True)
    deck.text(
        cover,
        "Posição do Itaú BBA na originação de renda fixa",
        0.92,
        2.72,
        11.4,
        0.7,
        size=30,
        color=BLACK,
        bold=True,
    )
    deck.text(
        cover,
        "Primeiro semestre de 2026 · Ranking ANBIMA de Renda Fixa e Híbridos",
        0.92,
        3.52,
        11.4,
        0.34,
        size=14,
        color=GRAY_700,
    )
    deck.rule(cover, 0.92, 4.08, 3.2, color=GRAY_300, height=0.018)
    deck.text(
        cover,
        f"Data-corte {cutoff_label} · referência Junho/2026 · publicado em 27/07/2026",
        0.92,
        4.32,
        11.4,
        0.3,
        size=11,
        color=GRAY_500,
    )
    deck.text(
        cover,
        "Análise Setorial de Crédito — Itaú BBA",
        0.92,
        6.62,
        6.0,
        0.28,
        size=10,
        color=GRAY_500,
    )

    # ------------------------------------------------- resultado do 1S26
    slide = deck.slide("O Itaú BBA encerrou o 1S26 em 2º lugar, com 22,9% do mercado apurado")

    for index, (label, value, note) in enumerate(
        (
            ("Volume originado", f"R$ {_bi(float(house_row['value_brl_or_count']))} bi", "1S26"),
            ("Participação", _pct(float(house_row["share"]), 1), "do mercado apurado"),
            ("Posição", _rank(house_row["rank"]), f"entre {participants} coordenadores"),
            ("Mercado apurado", f"R$ {_bi(universe)} bi", "universo do ranking"),
        )
    ):
        x = 0.62 + index * 3.06
        deck.block(slide, x, 1.45, 2.86, 1.16, GRAY_100)
        deck.text(slide, label.upper(), x + 0.2, 1.6, 2.5, 0.22, size=9, color=GRAY_500, bold=True)
        deck.text(slide, value, x + 0.2, 1.86, 2.5, 0.42, size=20, color=BLACK, bold=True)
        deck.text(slide, note, x + 0.2, 2.3, 2.5, 0.22, size=9, color=GRAY_500)

    top = ranked.head(10)
    rows = [["#", "Coordenador", "Volume (R$ bi)", "Part."]]
    highlight = None
    for position, (_, row) in enumerate(top.iterrows()):
        if row["participant"] == house:
            highlight = position
        rows.append(
            [
                _rank(row["rank"]),
                _display_name(row["participant"]),
                _bi(float(row["value_brl_or_count"])),
                _pct(float(row["share"]), 1),
            ]
        )
    deck.table(
        slide,
        rows,
        0.62,
        2.88,
        [0.7, 6.2, 2.6, 2.55],
        highlight=highlight,
        row_height=0.33,
    )
    deck.footer(slide, SOURCE_LINE)

    # -------------------------------------------- decomposição do share
    slide = deck.slide("A liderança do Itaú BBA está em securitização; a perda foi em dívida corporativa")

    rows = [["Subdivisão do Tipo 1", "Mercado (R$ bi)", "Itaú BBA (R$ bi)", "Part.", "Pos."]]
    for code, label in SUBDIVISIONS:
        block = _accumulated(official, code)
        if block.empty:
            continue
        block_universe = float(block["value_brl_or_count"].sum())
        entry = block[block["participant"].eq(house)]
        if entry.empty:
            volume, share, position = 0.0, float("nan"), float("nan")
        else:
            volume = float(entry["value_brl_or_count"].iloc[0])
            share = float(entry["share"].iloc[0])
            position = entry["rank"].iloc[0]
        rows.append(
            [
                label,
                _bi(block_universe),
                _bi(volume),
                _pct(share, 1),
                _rank(position),
            ]
        )
    bottom = deck.table(
        slide,
        rows,
        0.62,
        1.5,
        [5.0, 2.3, 2.3, 1.4, 1.05],
        row_height=0.46,
        align_right_from=1,
    )

    deck.block(slide, 0.62, bottom + 0.28, 12.05, 1.38, GRAY_100)
    deck.text(
        slide,
        "LEITURA",
        0.92,
        bottom + 0.45,
        3.0,
        0.22,
        size=9,
        color=ORANGE,
        bold=True,
    )
    deck.text(
        slide,
        "O bloco de longo prazo — debêntures e notas comerciais — responde por 79% do mercado "
        "apurado e é onde o Itaú BBA aparece em 2º, com 20,9%. Em securitização, que responde "
        "por 16% do mercado, o banco é 1º com 37,6%. A troca de liderança no consolidado do "
        "semestre veio da dívida corporativa, não da securitização.",
        0.92,
        bottom + 0.72,
        11.45,
        0.85,
        size=12,
        color=GRAY_900,
    )
    deck.footer(
        slide,
        "Fonte: ANBIMA, Ranking de Renda Fixa e Híbridos — Originação (Valor), "
        "subdivisões Tipo 1.1, 1.2 e 1.3, acumulado 2026 · referência Junho/2026",
    )

    # ------------------------------------------------------- metodologia
    slide = deck.slide("Como a ANBIMA apura o ranking")
    cursor = 1.5
    for index, (title, body) in enumerate(METHODOLOGY):
        column = index % 2
        if column == 0 and index:
            cursor += 1.72
        x = 0.62 + column * 6.2
        deck.block(slide, x, cursor, 0.045, 1.06, ORANGE)
        deck.text(slide, title, x + 0.28, cursor, 5.5, 0.5, size=13, color=BLACK, bold=True)
        deck.text(slide, body, x + 0.28, cursor + 0.46, 5.5, 0.95, size=11, color=GRAY_700)
    deck.footer(
        slide,
        "Fonte: ANBIMA, Metodologia do Ranking de Renda Fixa e Híbridos (fev/2026), "
        "capítulos II a VII",
    )

    # ------------------------------------------------ fontes e ressalvas
    slide = deck.slide("Fontes, reprodutibilidade e ressalvas")

    sources = manifest.get("sources", {})
    rows = [["Insumo", "Arquivo", "SHA-256 (12)"]]
    for label, key in (
        ("Ranking publicado", "ranking_workbook"),
        ("Anexo de encerramento", "annex_workbook"),
        ("Ofertas públicas CVM", "cvm_archive"),
    ):
        entry = sources.get(key, {})
        # The CVM archive may be supplied from anywhere on disk; name it by its
        # canonical download URL so the deck never shows a local path.
        reference = str(entry.get("url") or entry.get("path") or "")
        rows.append(
            [
                label,
                Path(reference).name or "—",
                str(entry.get("sha256", ""))[:12] or "—",
            ]
        )
    bottom = deck.table(
        slide, rows, 0.62, 1.5, [3.1, 6.4, 2.5], row_height=0.38, align_right_from=3
    )

    notes = (
        (
            "Reprodução auditada",
            "Somar o anexo operação a operação reconstrói o ranking publicado ao centavo: "
            "divergência máxima de R$ 0,0000076 em originação e distribuição.",
        ),
        (
            "Entidade jurídica",
            "A ANBIMA consolida o grupo sob o rótulo único ITAU BBA e não segrega Itaú BBA "
            "Assessoria Financeira S.A. do Banco Itaú BBA S.A. Na base CVM do 1S26, todas as "
            "ofertas lideradas pelo grupo saíram sob a Assessoria Financeira.",
        ),
        (
            "Universo do ranking",
            "O ranking cobre 71% do volume de renda fixa registrado na CVM no 1S26. A diferença "
            "vem de operações de empresas ligadas e de ofertas cujos formulários não foram "
            "enviados à ANBIMA.",
        ),
    )
    cursor = bottom + 0.3
    for title, body in notes:
        deck.text(slide, title, 0.62, cursor, 3.1, 0.3, size=11, color=ORANGE, bold=True)
        deck.text(slide, body, 3.85, cursor, 8.82, 0.72, size=11, color=GRAY_900)
        cursor += 0.86

    deck.footer(
        slide,
        "Reproduzível por scripts/build_anbima_fixed_income_ranking.py · "
        "data.anbima.com.br/publicacoes/ranking-de-renda-fixa-e-hibridos",
    )

    output.parent.mkdir(parents=True, exist_ok=True)
    deck.prs.save(output)
    return output


def main() -> None:
    args = parse_args()
    path = build(args.data_dir, args.output, args.house)
    print(f"deck: {path}")


if __name__ == "__main__":
    main()
