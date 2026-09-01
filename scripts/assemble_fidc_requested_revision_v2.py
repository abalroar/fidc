"""Assemble the 33-slide director deck while preserving native OOXML objects.

The edited 39-slide deck remains the data-bearing source. This assembler drops
the ranges the director no longer needs, appends the official 13-slide ANBIMA
ranking block, moves that block after the current-offers section, and renumbers
the visible page labels. Charts, tables, embedded workbooks, and notes remain
native.
"""

from __future__ import annotations

import argparse
from io import BytesIO
from pathlib import Path
import re
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation

from services.pptx_merge import merge_pptx_bytes


REMOVED_SOURCE_SLIDES = frozenset((*range(10, 24), *range(28, 33)))
EXPECTED_SOURCE_SLIDES = 39
EXPECTED_ANBIMA_SLIDES = 13
EXPECTED_FINAL_SLIDES = 33
FIRST_TAIL_POSITION = 13
ANBIMA_UNIFIED_SLIDES = frozenset((*range(32, 43), 52, 53))


def _slide_text(slide) -> str:  # noqa: ANN001
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text.strip()
    )


def _drop_slide(presentation: Presentation, index: int) -> None:
    slide_id = presentation.slides._sldIdLst[index]
    relationship_id = slide_id.rId
    presentation.part.drop_rel(relationship_id)
    presentation.slides._sldIdLst.remove(slide_id)


def _filtered_source(payload: bytes) -> bytes:
    presentation = Presentation(BytesIO(payload))
    if len(presentation.slides) != EXPECTED_SOURCE_SLIDES:
        raise ValueError(
            f"Deck-fonte deveria ter {EXPECTED_SOURCE_SLIDES} slides; "
            f"recebeu {len(presentation.slides)}"
        )
    for slide_number in sorted(REMOVED_SOURCE_SLIDES, reverse=True):
        _drop_slide(presentation, slide_number - 1)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _anbima_appendix(payload: bytes) -> bytes:
    presentation = Presentation(BytesIO(payload))
    if len(presentation.slides) < max(ANBIMA_UNIFIED_SLIDES):
        raise ValueError("Versão unificada não contém o bloco ANBIMA esperado")
    for slide_number in range(len(presentation.slides), 0, -1):
        if slide_number not in ANBIMA_UNIFIED_SLIDES:
            _drop_slide(presentation, slide_number - 1)
    if len(presentation.slides) != EXPECTED_ANBIMA_SLIDES:
        raise ValueError("Recorte ANBIMA não contém 13 slides")
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _move_anbima_before_tail(payload: bytes) -> bytes:
    presentation = Presentation(BytesIO(payload))
    if len(presentation.slides) != EXPECTED_FINAL_SLIDES:
        raise ValueError(
            f"Deck montado deveria ter {EXPECTED_FINAL_SLIDES} slides; "
            f"recebeu {len(presentation.slides)}"
        )
    slide_ids = list(presentation.slides._sldIdLst)
    base_count = EXPECTED_FINAL_SLIDES - EXPECTED_ANBIMA_SLIDES
    desired = (
        slide_ids[:FIRST_TAIL_POSITION]
        + slide_ids[base_count:]
        + slide_ids[FIRST_TAIL_POSITION:base_count]
    )
    for slide_id in list(presentation.slides._sldIdLst):
        presentation.slides._sldIdLst.remove(slide_id)
    for slide_id in desired:
        presentation.slides._sldIdLst.append(slide_id)

    for page_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if (
                re.fullmatch(r"\d+", text)
                and shape.left >= 11.8 * 914_400
                and shape.top >= 6.8 * 914_400
            ):
                shape.text_frame.paragraphs[0].runs[0].text = str(page_number)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _normalize_unsigned_chart_axis_ids(payload: bytes) -> bytes:
    """Repair legacy signed axis IDs so strict OOXML readers accept the deck."""
    output = BytesIO()
    with ZipFile(BytesIO(payload), "r") as source, ZipFile(
        output, "w", compression=ZIP_DEFLATED
    ) as target:
        for info in source.infolist():
            content = source.read(info.filename)
            if info.filename.startswith("ppt/charts/") and info.filename.endswith(".xml"):
                root = etree.fromstring(content)
                changed = False
                for node in root.xpath(
                    './/*[local-name()="axId" or local-name()="crossAx"][@val]'
                ):
                    value = int(node.get("val"))
                    if value < 0:
                        node.set("val", str(value + 2**32))
                        changed = True
                if changed:
                    content = etree.tostring(
                        root, xml_declaration=True, encoding="UTF-8", standalone=True
                    )
            target.writestr(info, content)
    return output.getvalue()


def _validate(payload: bytes) -> None:
    presentation = Presentation(BytesIO(payload))
    if len(presentation.slides) != EXPECTED_FINAL_SLIDES:
        raise ValueError("Quantidade final de slides inválida")
    texts = [_slide_text(slide).casefold() for slide in presentation.slides]
    forbidden = (
        "fomento mercantil: crescimento marginal",
        "risco estrutural · carteira i",
        "as 15 maiores ofertas de 2025",
        "top 15 · histórico",
    )
    if any(token in text for token in forbidden for text in texts):
        raise ValueError("Um slide solicitado para remoção permanece no deck")
    required = (
        "visão por produto: mercado, posição e participação do itaú bba",
        "em fidc o itaú bba lidera com 45,7%",
        "como interpretar o ranking anbima",
        "premissas, limitações e fontes",
    )
    for token in required:
        if not any(token in text for text in texts):
            raise ValueError(f"Slide ANBIMA ausente: {token}")
    if "emissões crescem 15%" not in texts[9]:
        raise ValueError("Bloco de ofertas correntes fora da posição esperada")
    if "posição competitiva do itaú bba" not in texts[13]:
        raise ValueError("Bloco ANBIMA fora da posição esperada")
    if "o que muda a leitura do mercado" not in texts[26]:
        raise ValueError("Fechamento executivo fora da posição esperada")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source", required=True, type=Path)
    parser.add_argument("--anbima", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()

    source = args.source.read_bytes()
    anbima = _anbima_appendix(args.anbima.read_bytes())
    filtered = _filtered_source(source)
    merged = merge_pptx_bytes(filtered, anbima)
    assembled = _normalize_unsigned_chart_axis_ids(_move_anbima_before_tail(merged))
    _validate(assembled)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_bytes(assembled)
    print(args.output)


if __name__ == "__main__":
    main()
